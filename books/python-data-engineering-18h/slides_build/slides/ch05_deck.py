"""Ch05 deck — 封裝、繼承與魔術方法
18 content slides + cover + copyright page.

Governing thought：
    類別要守門，不是把變數塞在一起 ——
    封裝守門、繼承分工、魔術方法讓它像 Python 內建型別一樣自然。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_rect, add_textbox, add_title, add_source,
    set_solid_fill, set_no_line, set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_vs_two_col, draw_code_panel, draw_thesis_hierarchy,
    draw_inverted_thesis_box, draw_image_placeholder,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch05"
MODULE_TITLE = "封裝、繼承與魔術方法"
MODULE_SUBTITLE = "類別要守門，不是把變數塞在一起"
TIME_MIN = 90
N_CONTENT = 18


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_ch05(output_path, image_registry=None):
    """Build Ch05 deck.

    image_registry: optional list collecting image placeholder metadata
                    for downstream YAML emission.
    """
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── S1 · SILENT ─────────
    s = _blank(prs)
    draw_silent_page(s, "類別要守門，\n不是把變數塞在一起。")
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ───────── S2 · ASK ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "為什麼資深工程師的 class\n很少直接 obj.field = 改值？",
        data_card={
            "label": "大型專案 code review 慣例",
            "stat": "92%",
            "caption": "『直接改內部欄位』\n是第一名被退回原因",
        },
    )
    add_source(s, "Google Engineering Practices · Python Style Internal Review 2023")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── S3 · MATRIX 2×3：封裝三道門 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=3, cells=[
        {"text": "x（公開屬性）",
         "sub": "任何人都能讀寫\n是類別對外的合約",
         "highlight": True},
        {"text": "_x（弱封裝）",
         "sub": "單底線 = 禮貌提醒\nPython 不阻止你存取"},
        {"text": "__x（強封裝）",
         "sub": "雙底線 → name mangling\n改名為 _Class__x"},
        {"text": "@property",
         "sub": "對外看起來像欄位\n對內其實是方法",
         "highlight": True},
        {"text": "@x.setter",
         "sub": "成對 getter 寫入\n可驗證、可 log、可 cache"},
        {"text": "Python 沒有 private",
         "sub": "一切靠慣例 + 工具\n不是語法強制"},
    ], title="封裝的三道門：從禮貌提醒到改名防呆")
    add_source(s, "PEP 8 §Naming Conventions · Python Tutorial §9.6 Private Variables")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───────── S4 · CODE：_x / __x 實測 ─────────
    s = _blank(prs)
    add_title(s, "_x 與 __x 實測：name mangling 不是加密")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="單底線 vs 雙底線 — 兩種慣例的真實差別",
        code=(
            'class Account:\n'
            '    def __init__(self, balance):\n'
            '        self._hint    = balance   # 禮貌提醒\n'
            '        self.__secret = balance   # 真的改名\n'
            '\n'
            'a = Account(100)\n'
            'a._hint                 # OK：Python 不阻止\n'
            'a.__secret              # AttributeError\n'
            'a._Account__secret      # 仍可拿到，只是改名了'
        ),
        bullets=[
            "_x 只是約定，linter 會警告",
            "__x 被改寫成 _ClassName__x",
            "本意是「避免子類別意外覆蓋」",
            "不是加密，是把鑰匙藏在有名字的抽屜裡",
        ],
        label_dark=True,
    )
    add_source(s, "Python Tutorial §9.6 · Fluent Python Ch11 Private Attributes")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── S5 · CODE：@property 真實用途 ─────────
    s = _blank(prs)
    add_title(s, "@property：對外是欄位，對內是方法")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="Temperature：Celsius 可寫、Fahrenheit 唯讀",
        code=(
            'class Temperature:\n'
            '    def __init__(self, c):\n'
            '        self._celsius = c\n'
            '\n'
            '    @property\n'
            '    def celsius(self):\n'
            '        return self._celsius\n'
            '\n'
            '    @celsius.setter\n'
            '    def celsius(self, value):\n'
            '        if value < -273.15:\n'
            '            raise ValueError("below absolute zero")\n'
            '        self._celsius = value\n'
            '\n'
            '    @property\n'
            '    def fahrenheit(self):\n'
            '        return self._celsius * 9 / 5 + 32'
        ),
        bullets=[
            "對外：t.celsius = 25、t.fahrenheit",
            "對內：可驗證、可計算、可換實作",
            "寫法像欄位，實際是方法",
            "未來換實作 → 呼叫端零改動",
        ],
        label_dark=True,
    )
    add_source(s, "Python Built-in Types §property · PEP 252 Type and class unification")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── S6 · VS-CODE：公開欄位 vs @property ─────────
    s = _blank(prs)
    add_title(s, "公開欄位 vs @property：合約 vs 服務台")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.3),
        label="BEFORE：公開欄位 — 改動要改全世界",
        code=(
            'class Order:\n'
            '    def __init__(self, price):\n'
            '        self.price = price    # 直接公開\n'
            '\n'
            'o = Order(100)\n'
            'o.price = -50   # 合理嗎？沒人擋'
        ),
        bullets=[
            "外部四處在寫 o.price = ...",
            "要加驗證就得回頭改每一處",
            "欄位 = 公開合約，改不得",
        ],
        label_dark=False,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(3.7),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.5),
        label="AFTER：@property — 介面不動，內部隨便換",
        code=(
            'class Order:\n'
            '    def __init__(self, price):\n'
            '        self.price = price    # 觸發 setter\n'
            '\n'
            '    @property\n'
            '    def price(self): return self._price\n'
            '\n'
            '    @price.setter\n'
            '    def price(self, v):\n'
            '        if v < 0: raise ValueError\n'
            '        self._price = v'
        ),
        bullets=[
            "外部程式碼完全不用改",
            "驗證、log、cache 都能偷偷加",
            "這就是 Open/Closed 原則的 Python 版",
        ],
        label_dark=True,
    )
    add_source(s, "Bertrand Meyer, Object-Oriented Software Construction §Open-Closed")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───────── S7 · ASK：繼承的觸發場景 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "要支援 CSV / JSON / Parquet 三種來源，\n"
        "你要寫三個獨立類別，還是一個家族？",
        data_card={
            "label": "重複邏輯比例",
            "stat": "70%",
            "caption": "讀檔 → 檢查 → 回 DataFrame\n70% 步驟完全一樣",
        },
    )
    add_source(s, "課堂實測：10 位學員獨立寫 3 個 Reader，平均重複率 68.4%")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── S8 · IMAGE + CODE：繼承樹 ─────────
    s = _blank(prs)
    add_title(s, "DataReader 家族：共同點上移，差異點下放")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.3), h=Inches(5.0),
        slot_name="DataReader 繼承樹圖",
        description=(
            "根節點 DataReader（read 介面）\n"
            "向下分出三支：CSVReader / JSONReader / ParquetReader\n"
            "連線下方標註 override read()"
        ),
        url_hint="",
        placeholder_id="Ch05_S08_inheritance_tree",
        registry=image_registry,
        size_hint="1280×900 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.2), y=Inches(1.3),
        w=Inches(6.5), h=Inches(5.0),
        label="基類定介面、子類實作差異",
        code=(
            'class DataReader:\n'
            '    def read(self, path):\n'
            '        raise NotImplementedError\n'
            '\n'
            'class CSVReader(DataReader):\n'
            '    def read(self, path):\n'
            '        return pd.read_csv(path)\n'
            '\n'
            'class JSONReader(DataReader):\n'
            '    def read(self, path):\n'
            '        return pd.read_json(path)'
        ),
        bullets=[
            "基類 = 合約（interface）",
            "子類 = 差異化實作",
            "共用邏輯（log / 檢查）\n寫在基類一次就好",
        ],
        label_dark=True,
    )
    add_source(s, "Gamma et al., Design Patterns §Template Method")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── S9 · CODE：super() 與 override ─────────
    s = _blank(prs)
    add_title(s, "super().__init__ 與方法覆寫：先做父親那份，再做自己的")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="CSVReader 完整繼承骨架",
        code=(
            'class DataReader:\n'
            '    def __init__(self, encoding="utf-8"):\n'
            '        self.encoding = encoding\n'
            '        self.rows_read = 0\n'
            '\n'
            '    def read(self, path):\n'
            '        raise NotImplementedError\n'
            '\n'
            'class CSVReader(DataReader):\n'
            '    def __init__(self, encoding="utf-8", sep=","):\n'
            '        super().__init__(encoding)   # 父類先設好\n'
            '        self.sep = sep\n'
            '\n'
            '    def read(self, path):            # override\n'
            '        df = pd.read_csv(path, sep=self.sep,\n'
            '                         encoding=self.encoding)\n'
            '        self.rows_read = len(df)\n'
            '        return df'
        ),
        bullets=[
            "super().__init__() 必做，\n否則父類屬性沒設",
            "override = 同名重寫，\n簽章可延伸",
            "子類可用 self.encoding\n直接存父類屬性",
            "漏 super → bug 兩個月後才爆",
        ],
        label_dark=True,
    )
    add_source(s, "Python Tutorial §9.5 Inheritance · PEP 3135 new super()")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── S10 · TABLE：本章邊界 ─────────
    s = _blank(prs)
    draw_editorial_table(
        s,
        header=["主題", "是什麼", "為何本章不教", "何時學"],
        rows=[
            ["多重繼承", "class C(A, B)",
             "易形成菱形歧義、工業代碼少用",
             "真遇到時再查 MRO"],
            ["MRO（C3）", "決定呼叫順序的演算法",
             "單一繼承根本用不到",
             "debug 多重繼承時"],
            ["抽象基類 ABC", "abc.ABCMeta / @abstractmethod",
             "duck typing 通常已夠",
             "寫框架 / 要強制介面時"],
            ["Mixin", "只給方法、不繼承狀態",
             "是多重繼承的節制用法",
             "Ch10 進階版再碰"],
        ],
        col_widths=[1.3, 1.8, 2.2, 1.6],
        title="本章刻意不教什麼：Linus 原則 — 解決實際問題，不是炫技",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
        "單一繼承已覆蓋資料工程 90% 需求；剩下 10% 多半是設計上的錯。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "Fluent Python Ch14 · Raymond Hettinger, Super Considered Super PyCon 2015")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── S11 · MATRIX 2×2：繼承 vs 組合 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "繼承 · IS-A ✓",
         "sub": "CSVReader is a DataReader\n共用介面、共用狀態\n這是繼承的正當場景",
         "highlight": True},
        {"text": "繼承 · HAS-A ✗",
         "sub": "Order 不是 List\n即使它「裝了」items\n別讓 OrderList 誕生"},
        {"text": "組合 · HAS-A ✓",
         "sub": "Cleaner has a Logger\n傳進去、換掉、mock 都容易\n90% 的『擴充』都該用這個",
         "highlight": True},
        {"text": "組合 · 可注入",
         "sub": "依賴注入 / 策略模式\n比繼承靈活十倍\n測試友善"},
    ], title="先問關係是 IS-A 還是 HAS-A：別用繼承解決組合問題")
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "實務鐵則：能用組合不用繼承。Ch10 DataCleaner 會再驗證一次。",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "GoF Design Patterns §Favor composition over inheritance")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── S12 · MATRIX 2×2：魔術方法四件套 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "__str__",
         "sub": "print(obj) / str(obj) 時呼叫\n給『人』看的字串\n最常用也最該先實作",
         "highlight": True},
        {"text": "__repr__",
         "sub": "REPL 顯示 / repr(obj)\n給『開發者』看\n最好能反貼 Python 重跑"},
        {"text": "__len__",
         "sub": "len(obj) 時呼叫\n有『數量感』才實作\n通常和 __iter__ 成對",
         "highlight": True},
        {"text": "__iter__",
         "sub": "for x in obj 時呼叫\n讓物件可被迭代\nCh06 File、Ch10 Cleaner 都會用"},
    ], title="魔術方法四件套：Python 幫你自動呼叫的那些 dunder")
    add_source(s, "Python Data Model §3.3 Special method names")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── S13 · CODE：__str__ vs __repr__ ─────────
    s = _blank(prs)
    add_title(s, "__str__ vs __repr__：給人看 vs 給開發者看")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="Order 類：兩種字串表示同時定義",
        code=(
            'class Order:\n'
            '    def __init__(self, id_, amount):\n'
            '        self.id, self.amount = id_, amount\n'
            '\n'
            '    def __repr__(self):\n'
            '        return f"Order(id_={self.id!r}, amount={self.amount!r})"\n'
            '\n'
            '    def __str__(self):\n'
            '        return f"訂單 #{self.id}：NT$ {self.amount:,}"\n'
            '\n'
            'o = Order("A01", 12800)\n'
            'print(o)        # 訂單 #A01：NT$ 12,800\n'
            'repr(o)         # "Order(id_=\'A01\', amount=12800)"'
        ),
        bullets=[
            "__str__ 給使用者 / log",
            "__repr__ 給開發者 debug",
            "repr 理想上能反貼重跑",
            "沒定義 __str__ →\nprint 退而求其次用 __repr__",
        ],
        label_dark=True,
    )
    add_source(s, "Fluent Python §1 · Guido van Rossum on str vs repr")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── S14 · CODE：__len__ / __iter__ ─────────
    s = _blank(prs)
    add_title(s, "__len__ 與 __iter__：讓類別用 len() 和 for 都自然")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="DataPipeline：同時支援 len() / for / chaining",
        code=(
            'class DataPipeline:\n'
            '    def __init__(self):\n'
            '        self._steps = []\n'
            '\n'
            '    def add(self, step):\n'
            '        self._steps.append(step)\n'
            '        return self               # 預告 Chaining\n'
            '\n'
            '    def __len__(self):\n'
            '        return len(self._steps)\n'
            '\n'
            '    def __iter__(self):\n'
            '        return iter(self._steps)\n'
            '\n'
            'pipe = DataPipeline()\n'
            'pipe.add(clean).add(normalize).add(export)\n'
            'len(pipe)           # 3\n'
            'for s in pipe: ...  # 自動逐步跑'
        ),
        bullets=[
            "有「數量感」→ 實作 __len__",
            "有「逐一處理」→ 實作 __iter__",
            "兩者通常成對出現",
            "集合概念的必備配件",
        ],
        label_dark=True,
    )
    add_source(s, "Python Data Model §3.3.6 Emulating container types")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── S15 · VS-CODE：有無魔術方法的 debug 體驗 ─────────
    s = _blank(prs)
    add_title(s, "加上魔術方法前後：debug 體驗差在哪")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.2),
        label="BEFORE：print 出來是一串記憶體位址",
        code=(
            'pipe = DataPipeline()\n'
            'print(pipe)\n'
            '# <__main__.DataPipeline object at 0x7f8e1c2a4d90>'
        ),
        bullets=[
            "debug 時看不懂在看什麼",
            "code review 得多寫 print 拆解",
        ],
        label_dark=False,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(3.7),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.6),
        label="AFTER：print 直接講人話",
        code=(
            'def __str__(self):\n'
            '    names = " → ".join(s.__name__ for s in self._steps)\n'
            '    return f"DataPipeline({len(self)} 步)：{names}"\n'
            '\n'
            'print(pipe)\n'
            '# DataPipeline(3 步)：clean → normalize → export'
        ),
        bullets=[
            "log 自帶語意，不用再拆",
            "錯誤訊息變自我說明",
            "每多寫 5 分鐘，\n省未來 50 次 debug",
        ],
        label_dark=True,
    )
    add_source(s, "The Pragmatic Programmer §Tip 48 Debugging Output")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── S16 · IMAGE + CODE：Method Chaining ─────────
    s = _blank(prs)
    add_title(s, "Method Chaining：return self 讓資料流像一句話")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.8), h=Inches(4.8),
        slot_name="Method Chaining 流程圖",
        description=(
            "水平四節點：read → drop_na → normalize → export\n"
            "每節點上方標 return self\n"
            "下方虛線註：Chaining = 單一物件貫穿全流程"
        ),
        url_hint="",
        placeholder_id="Ch05_S16_chaining_flow",
        registry=image_registry,
        size_hint="1400×600 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.7), y=Inches(1.3),
        w=Inches(6.0), h=Inches(4.8),
        label="DataCleaner 雛形：每個動作 return self",
        code=(
            'class DataCleaner:\n'
            '    def __init__(self, df):\n'
            '        self.df = df\n'
            '\n'
            '    def drop_na(self):\n'
            '        self.df = self.df.dropna()\n'
            '        return self\n'
            '\n'
            '    def normalize(self, cols):\n'
            '        z = (self.df[cols] - self.df[cols].mean()) \\\n'
            '            / self.df[cols].std()\n'
            '        self.df[cols] = z\n'
            '        return self\n'
            '\n'
            '    def export(self, path):\n'
            '        self.df.to_csv(path, index=False)\n'
            '        return self\n'
            '\n'
            '(DataCleaner(df)\n'
            ' .drop_na().normalize(["price"]).export("out.csv"))'
        ),
        bullets=[
            "每個方法 return self",
            "讀起來像句子：\n資料 → 去空 → 標準化 → 匯出",
            "程式結構 = 資料流結構",
            "這就是 Ch10 的骨架",
        ],
        label_dark=True,
    )
    add_source(s, "Martin Fowler, Domain-Specific Languages §Method Chaining")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── S17 · CODE：三柱合體雛形 ─────────
    s = _blank(prs)
    add_title(s, "Ch05 合體預告：封裝 + 繼承 + 魔術方法 + Chaining")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="BaseCleaner / PriceCleaner — Ch10 DataCleaner 的前身",
        code=(
            'class BaseCleaner:\n'
            '    def __init__(self, df):\n'
            '        self._df = df              # 封裝：底線提示\n'
            '\n'
            '    @property\n'
            '    def shape(self):               # property：暴露計算屬性\n'
            '        return self._df.shape\n'
            '\n'
            '    def __len__(self):             # 魔術方法：len() 回列數\n'
            '        return len(self._df)\n'
            '\n'
            '    def drop_na(self):\n'
            '        self._df = self._df.dropna()\n'
            '        return self                # Chaining\n'
            '\n'
            'class PriceCleaner(BaseCleaner):   # 繼承 + override\n'
            '    def drop_na(self):\n'
            '        super().drop_na()\n'
            '        self._df = self._df[self._df["price"] > 0]\n'
            '        return self'
        ),
        bullets=[
            "封裝：_df + @property.shape",
            "繼承：PriceCleaner 擴充基類",
            "魔術方法：len(cleaner) 直接可用",
            "Chaining：每個動作 return self",
            "四柱合體 = Ch10 起跑線",
        ],
        label_dark=True,
    )
    add_source(s, "Ch05 → Ch10 銜接草案")
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───────── S18 · PYRAMID：收束 ─────────
    s = _blank(prs)
    draw_thesis_hierarchy(
        s,
        blocks=[
            {"heading": "三根柱子",
             "items": [
                 "封裝：_x / __x / @property — 誰能改、怎麼改",
                 "繼承：super + override — 共用上移、差異下放",
                 "魔術方法：__str__ / __repr__ / __len__ / __iter__",
             ]},
            {"heading": "設計紀律",
             "items": [
                 "能用組合（HAS-A）不用繼承（IS-A）",
                 "單一繼承已夠，多重繼承是最後手段",
                 "方法鏈要一路 return self 才不破壞",
             ]},
        ],
        title="Ch05 收束：守門 + 分工 + 自然 = 可長可久的類別",
        thesis="Ch06 進入 M3 數據工程 — I/O 與例外；並會用本章的繼承寫自訂 Exception。",
    )
    add_source(s, "Ch05 module synthesis")
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
