"""Ch07 deck — NumPy 數值運算與向量化
19 content slides + cover + copyright page.

Governing thought：
    NumPy 不是讓你算快一點，
    而是讓你用 shape × broadcasting 的語言描述計算問題 ——
    從「我寫個 for 迴圈」直接跳到「向量化思維」的分水嶺。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_code_panel, draw_thesis_hierarchy, draw_image_placeholder,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch07"
MODULE_TITLE = "NumPy 數值運算與向量化"
MODULE_SUBTITLE = "shape × broadcasting × axis — 向量化思維的分水嶺"
TIME_MIN = 120
N_CONTENT = 19


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_ch07(output_path, image_registry=None):
    """Build Ch07 deck.

    image_registry: optional list collecting image placeholder metadata
                    for downstream YAML emission.
    """
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── S1 · SILENT ─────────
    s = _blank(prs)
    draw_silent_page(s, "別再寫 for 算數字。\n換一種語言想這件事。")
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ───────── S2 · ASK ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "同樣加總一億個數字，\nlist vs ndarray 差多少？",
        data_card={
            "label": "1e8 筆 float64 加總實測",
            "stat": "80×",
            "caption": "Python list comprehension vs np.sum\n平均時間比",
        },
    )
    add_source(s, "NumPy Benchmarks · Intel MKL 實測 · Author local i7-12700")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── S3 · MATRIX 2×3：list vs ndarray ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=3, cells=[
        {"text": "list 混型別",
         "sub": "一個 list 可裝 int / str / obj\n每格都是 PyObject\n多一層 pointer 包裝",
         "highlight": True},
        {"text": "list 記憶體散落",
         "sub": "pointer 四散於 heap\nCPU cache 無法預測\n每次讀都得重查位址"},
        {"text": "list 無向量化",
         "sub": "Python level for 迴圈\n每一輪都進直譯器\n天花板在 CPython"},
        {"text": "ndarray 統一 dtype",
         "sub": "整塊相同型別\nfloat64 就是 float64\n沒有型別檢查開銷",
         "highlight": True},
        {"text": "ndarray 連續記憶體",
         "sub": "一整條、一個 malloc\nCPU cache line 友善\n預讀就是預測"},
        {"text": "ndarray 向量化 + SIMD",
         "sub": "底層 C 實作\nCPU 一次處理 4/8 筆\n這就是 80× 的來源"},
    ], title="為什麼 list 不夠用：ndarray 的六個差別")
    add_source(s, "NumPy User Guide §The N-dimensional array · CPython Implementation Details")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───────── S4 · CODE：shape / dtype / ndim ─────────
    s = _blank(prs)
    add_title(s, "ndarray 三件套：shape / dtype / ndim")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="描述一個 ndarray，只需三個屬性",
        code=(
            'import numpy as np\n'
            '\n'
            'a = np.array([[1, 2, 3],\n'
            '              [4, 5, 6]])\n'
            '\n'
            'a.shape   # (2, 3)        每個維度的長度\n'
            'a.dtype   # dtype(\'int64\') 所有元素的統一型別\n'
            'a.ndim    # 2              維度數\n'
            'a.size    # 6              總元素數 = 乘積\n'
            '\n'
            '# 看到一個 ndarray，先 print 這三個就對了\n'
            'print(a.shape, a.dtype, a.ndim)'
        ),
        bullets=[
            "shape：tuple，每個維度長度",
            "dtype：全陣列共用一個型別",
            "ndim = len(shape)",
            "99% 的 bug\n先 print 這三個就破案",
        ],
        label_dark=True,
    )
    add_source(s, "NumPy Reference §ndarray attributes")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── S5 · IMAGE + CODE：連續記憶體 ─────────
    s = _blank(prs)
    add_title(s, "連續記憶體：ndarray 跑得快的物理原因")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.8), h=Inches(5.0),
        slot_name="ndarray 連續記憶體佈局",
        description=(
            "左：Python list 的 pointer 散落 heap\n"
            "右：ndarray 一整條連續 float64 格子\n"
            "CPU cache line 能一次載入整塊"
        ),
        url_hint="",
        placeholder_id="Ch07_S05_memory_layout",
        registry=image_registry,
        size_hint="1400×800 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.7), y=Inches(1.3),
        w=Inches(6.0), h=Inches(5.0),
        label="記憶體三兄弟：itemsize / nbytes / strides",
        code=(
            'a = np.zeros(1_000_000, dtype=np.float64)\n'
            '\n'
            'a.itemsize   # 8      每格 8 bytes\n'
            'a.nbytes     # 8_000_000  = 8 MB\n'
            'a.strides    # (8,)   跨一格要走 8 bytes\n'
            '\n'
            '# 同樣 100 萬筆：\n'
            '#   list:    ~56 MB（物件 + pointer）\n'
            '#   ndarray: ~8 MB（純資料）'
        ),
        bullets=[
            "連續記憶體 → CPU cache 預讀",
            "單一 dtype → SIMD 一次多筆",
            "無 PyObject 包裝\n= 純資料壓得很扁",
            "呼應 Ch1：\nCPU/RAM 的現實投影",
        ],
        label_dark=True,
    )
    add_source(s, "NumPy Internals · CPython list object size measurement")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── S6 · CODE：六種建立方式 ─────────
    s = _blank(prs)
    add_title(s, "建立 ndarray 的六種常用方式")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="涵蓋 95% 建立需求的六個 API",
        code=(
            'np.array([1, 2, 3])              # 從 list 轉過來\n'
            'np.zeros((3, 4))                 # 全零佔位\n'
            'np.ones((2, 3), dtype=np.int32)  # 全一 + 指定 dtype\n'
            'np.arange(0, 10, 2)              # 0,2,4,6,8\n'
            'np.linspace(0, 1, 5)             # 0,0.25,...,1 等距 5 份\n'
            'np.random.randn(3, 3)            # 常態分布隨機矩陣'
        ),
        bullets=[
            "array：你有現成資料",
            "zeros/ones：要先開空間\n再填值",
            "arange：整數步進",
            "linspace：固定段數\n不用算步長",
            "randn：ML 常用\n生成模擬資料",
        ],
        label_dark=True,
    )
    add_source(s, "NumPy Reference §Array creation routines")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───────── S7 · CODE：多維切片 ─────────
    s = _blank(prs)
    add_title(s, "多維切片：arr[行, 列] 的心智模型")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="逗號切割維度 — 不是 arr[行][列]",
        code=(
            'a = np.arange(12).reshape(3, 4)\n'
            '# [[ 0  1  2  3]\n'
            '#  [ 4  5  6  7]\n'
            '#  [ 8  9 10 11]]\n'
            '\n'
            'a[1, 2]        # 6            單一元素\n'
            'a[:, 0]        # [0, 4, 8]    第 0 欄\n'
            'a[1, :]        # [4, 5, 6, 7] 第 1 列\n'
            'a[1:3, ::2]    # [[4,6],[8,10]]  列 1-2、欄每隔一格\n'
            'a[..., 0]      # 省略前面維度 = 最後一維取 0'
        ),
        bullets=[
            "逗號切割每個維度",
            "slice 語法跟 list 一致",
            "省略的維度補 :",
            "... (Ellipsis) 代表\n剩下所有維度",
        ],
        label_dark=True,
    )
    add_source(s, "NumPy User Guide §Indexing on ndarrays")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── S8 · CODE：布林索引 ─────────
    s = _blank(prs)
    add_title(s, "布林索引：用條件過濾資料")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="條件 → bool array → 取元素",
        code=(
            'a = np.array([-2, -1, 0, 1, 2, 3, 4])\n'
            '\n'
            'mask = a > 0              # bool array\n'
            '# [False, False, False, True, True, True, True]\n'
            '\n'
            'a[mask]                   # [1, 2, 3, 4]\n'
            'a[a > 0]                  # 同上，一行寫法\n'
            '\n'
            '# 多條件：用 & | ~，不是 and / or / not\n'
            'a[(a > 0) & (a < 3)]      # [1, 2]\n'
            'a[~(a > 0)]               # [-2, -1, 0]\n'
            '\n'
            '# 也能寫入：把負數歸零\n'
            'a[a < 0] = 0'
        ),
        bullets=[
            "mask 與資料同 shape",
            "回傳一維、順序保留",
            "& | ~ 才是 element-wise\nand/or 會炸",
            "資料工程每天用最多次",
        ],
        label_dark=True,
    )
    add_source(s, "NumPy User Guide §Boolean array indexing")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── S9 · CODE：Fancy Indexing ─────────
    s = _blank(prs)
    add_title(s, "Fancy Indexing：用整數陣列當索引")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="依清單抽取 — 可不連續、可重排、可重複",
        code=(
            'a = np.array([10, 20, 30, 40, 50])\n'
            '\n'
            'a[[0, 2, 4]]      # [10, 30, 50]\n'
            'a[[4, 0, 2]]      # [50, 10, 30]  順序隨索引\n'
            'a[[0, 0, 1]]      # [10, 10, 20]  可重複\n'
            '\n'
            '# 二維：配對 (row, col)\n'
            'm = np.arange(12).reshape(3, 4)\n'
            'm[[0, 1, 2], [1, 2, 3]]\n'
            '#   → [m[0,1], m[1,2], m[2,3]] = [1, 6, 11]'
        ),
        bullets=[
            "整數陣列或 list 當索引",
            "可抽不連續 / 重排 / 重複",
            "二維是配對，不是網格",
            "回傳永遠是 copy\n（下一張詳談）",
        ],
        label_dark=True,
    )
    add_source(s, "NumPy User Guide §Integer array indexing")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── S10 · VS-CODE：view vs copy ─────────
    s = _blank(prs)
    add_title(s, "切片是 view、Fancy 是 copy：改一邊會動到另一邊嗎？")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.3),
        label="切片 = view：共享記憶體，改 view 就改原 array",
        code=(
            'a = np.arange(6)          # [0 1 2 3 4 5]\n'
            'b = a[1:4]                # view\n'
            'b[0] = 99\n'
            'a                         # [0 99 2 3 4 5]  ← 原 array 被改\n'
            'b.base is a               # True — 共享記憶體'
        ),
        bullets=[
            "切片回傳的是 view",
            "省記憶體、但偷改原資料",
            "真要獨立：b = a[1:4].copy()",
        ],
        label_dark=False,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(3.7),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.6),
        label="Fancy / Boolean = copy：改它不影響原 array",
        code=(
            'a = np.arange(6)\n'
            'c = a[[1, 3, 5]]          # fancy → copy\n'
            'c[0] = 99\n'
            'a                         # [0 1 2 3 4 5]  ← 原 array 不變\n'
            'c.base is None            # True — 獨立記憶體\n'
            '\n'
            'd = a[a > 2]              # boolean → 也是 copy'
        ),
        bullets=[
            "Fancy / Boolean 都是 copy",
            "改 copy 永遠不動原 array",
            "判斷法則：\n結果 shape 可預測 → view\n內容由 index 決定 → copy",
        ],
        label_dark=True,
    )
    add_source(s, "NumPy User Guide §Copies and views")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── S11 · ASK：Broadcasting 大哉問 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "兩個形狀不同的 ndarray，\n為什麼能直接相乘？",
        data_card={
            "label": "Broadcasting 規則",
            "stat": "3 步",
            "caption": "右對齊 → 缺維補 1\n→ 大小 1 可拉伸",
        },
    )
    add_source(s, "NumPy User Guide §Broadcasting")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── S12 · IMAGE + CODE：Broadcasting 三步規則 ─────────
    s = _blank(prs)
    add_title(s, "Broadcasting 三步規則：右對齊、補 1、拉伸")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.8), h=Inches(5.0),
        slot_name="Broadcasting 規則示意圖",
        description=(
            "A: (2,1) + B: (3,) 三步對齊過程\n"
            "Step 1 右對齊 → Step 2 缺維補 1\n"
            "→ Step 3 大小 1 沿軸拉伸 → (2,3)"
        ),
        url_hint="",
        placeholder_id="Ch07_S12_broadcasting_rule",
        registry=image_registry,
        size_hint="1400×700 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.7), y=Inches(1.3),
        w=Inches(6.0), h=Inches(5.0),
        label="三步規則的 shape 演算",
        code=(
            'A = np.array([[1], [2]])     # shape (2, 1)\n'
            'B = np.array([10, 20, 30])   # shape (3,)\n'
            '\n'
            '# Step 1 右對齊：\n'
            '#   A: (2, 1)\n'
            '#   B: (   3)  →  補成 (1, 3)\n'
            '# Step 2 缺維補 1\n'
            '# Step 3 大小 1 沿該軸拉伸\n'
            '#   A → (2, 3)\n'
            '#   B → (2, 3)\n'
            '\n'
            '(A + B).shape                # (2, 3) ✓\n'
            '\n'
            '# 不符合會炸：\n'
            'np.zeros((2,3)) + np.zeros((2,4))\n'
            '# ValueError: could not be broadcast'
        ),
        bullets=[
            "從右往左對齊",
            "缺的維度補 1",
            "大小 1 的維度\n可沿該軸拉伸",
            "不符合 → ValueError",
        ],
        label_dark=True,
    )
    add_source(s, "NumPy User Guide §General broadcasting rules")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── S13 · CODE：經典範例 ─────────
    s = _blank(prs)
    add_title(s, "Broadcasting 經典範例：prices × discount")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="三個商品 × 兩個折扣方案 — 一行算六個價格",
        code=(
            'prices   = np.array([100, 200, 300])       # shape (3,)\n'
            'discount = np.array([[0.9], [0.8]])        # shape (2, 1)\n'
            '\n'
            'result = prices * discount                 # shape (2, 3)\n'
            '# [[ 90. 180. 270.]   ← 9 折方案\n'
            '#  [ 80. 160. 240.]]  ← 8 折方案\n'
            '\n'
            '# 沒有 for 迴圈、沒有 reshape\n'
            '# shape (2,1) × (3,) → 自動對齊為 (2,3)'
        ),
        bullets=[
            "(2,1) × (3,) → (2,3)",
            "不用 reshape",
            "不用 for",
            "資料工程的日常：\n「一份資料 × 多種參數」",
        ],
        label_dark=True,
    )
    add_source(s, "NumPy User Guide §Broadcasting — Examples")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── S14 · CODE：z-score 實戰 ─────────
    s = _blank(prs)
    add_title(s, "實務場景：用 broadcasting 做 z-score 標準化")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="ML pipeline 裡 90% 都有這一行",
        code=(
            'X = np.random.randn(100, 5)          # 100 筆樣本 × 5 特徵\n'
            '\n'
            'mu    = X.mean(axis=0)               # shape (5,)\n'
            'sigma = X.std(axis=0)                # shape (5,)\n'
            '\n'
            'z = (X - mu) / sigma                 # shape (100, 5)\n'
            '#    (100, 5) - (5,)  → broadcasting 成功\n'
            '#    (100, 5) / (5,)  → 同理\n'
            '\n'
            '# 驗證：標準化後每欄均值 ≈ 0、標準差 ≈ 1\n'
            'z.mean(axis=0)   # ≈ [0, 0, 0, 0, 0]\n'
            'z.std(axis=0)    # ≈ [1, 1, 1, 1, 1]'
        ),
        bullets=[
            "(100,5) - (5,) 會沿列拉伸",
            "一行完成 preprocessing",
            "Ch08 會換成\nsklearn.StandardScaler",
            "底層就是這兩個運算子",
        ],
        label_dark=True,
    )
    add_source(s, "scikit-learn Preprocessing · ISLR §6 Standardization")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── S15 · VS-CODE：for vs 向量化 ─────────
    s = _blank(prs)
    add_title(s, "for 迴圈 vs 向量化：效能差距從哪裡來")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.3),
        label="BEFORE：純 Python for 迴圈 — 4.2 秒",
        code=(
            'data = list(range(10_000_000))\n'
            '\n'
            'total = 0\n'
            'for x in data:        # 一千萬次進直譯器\n'
            '    total += x * x\n'
            '# 耗時：約 4.2 秒'
        ),
        bullets=[
            "每輪進 CPython 直譯器",
            "每個 int 都是 PyObject",
            "沒有 SIMD",
        ],
        label_dark=False,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(3.7),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.6),
        label="AFTER：NumPy 向量化 — 0.05 秒（80× 快）",
        code=(
            'a = np.arange(10_000_000)\n'
            '\n'
            'total = (a ** 2).sum()   # 一句話 → 交給底層 C\n'
            '# 耗時：約 0.05 秒\n'
            '\n'
            '# 同樣一件事，用「shape 語言」描述 →\n'
            '# CPython 把它交給 C / BLAS / SIMD'
        ),
        bullets=[
            "底層 C 迴圈：無直譯開銷",
            "SIMD：一次 4–8 筆",
            "連續記憶體：cache 友善",
            "差距不是優化\n是換了一種計算方式",
        ],
        label_dark=True,
    )
    add_source(s, "NumPy Performance Tips · Intel SIMD intrinsics")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── S16 · TABLE：統計函式 ─────────
    s = _blank(prs)
    draw_editorial_table(
        s,
        header=["函式", "做什麼", "常用 axis", "Pandas 對應"],
        rows=[
            ["np.mean", "算術平均", "axis=0（每欄）/ axis=1（每列）", ".mean()"],
            ["np.std", "標準差", "axis=0 / axis=1", ".std()"],
            ["np.min / np.max", "極值", "axis=0 / axis=1", ".min() / .max()"],
            ["np.argmax / argmin", "極值的索引位置", "axis=0 / axis=1", ".idxmax() / .idxmin()"],
            ["np.percentile", "分位數（如 q=95）", "axis=0 / axis=1", ".quantile()"],
            ["np.cumsum", "累加（沿某軸）", "axis=0 / axis=1", ".cumsum()"],
        ],
        col_widths=[1.6, 2.2, 2.6, 1.8],
        title="NumPy 內建統計函式：一行搞定聚合",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.4),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
        "這六個函式蓋掉 EDA 階段 80% 的描述統計；Ch08 Pandas 用同名 API 直接接手。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "NumPy Reference §Statistics · pandas.DataFrame API")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── S17 · IMAGE + CODE：axis 方向圖 ─────────
    s = _blank(prs)
    add_title(s, "axis 方向圖：axis=0 沿列、axis=1 沿欄")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.6), h=Inches(5.0),
        slot_name="axis 方向示意",
        description=(
            "3×4 矩陣：axis=0 垂直箭頭穿欄 → 每欄聚合\n"
            "axis=1 水平箭頭穿列 → 每列聚合\n"
            "目的：破除「axis=0 是列」的常見誤解"
        ),
        url_hint="",
        placeholder_id="Ch07_S17_axis_diagram",
        registry=image_registry,
        size_hint="1200×900 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.5), y=Inches(1.3),
        w=Inches(6.2), h=Inches(5.0),
        label="axis=k 就是「把第 k 維吃掉」",
        code=(
            'arr = np.arange(12).reshape(3, 4)\n'
            '# [[ 0  1  2  3]\n'
            '#  [ 4  5  6  7]\n'
            '#  [ 8  9 10 11]]\n'
            '\n'
            'arr.mean(axis=0)\n'
            '# shape (4,)   每欄一個均值\n'
            '# [4., 5., 6., 7.]\n'
            '\n'
            'arr.mean(axis=1)\n'
            '# shape (3,)   每列一個均值\n'
            '# [1.5, 5.5, 9.5]'
        ),
        bullets=[
            "axis=0：第 0 維被摺掉\n→ 輸出 shape 去掉第 0 維",
            "axis=1：第 1 維被摺掉\n→ 輸出 shape 去掉第 1 維",
            "記熟輸出 shape\n就不會搞錯方向",
        ],
        label_dark=True,
    )
    add_source(s, "NumPy User Guide §Reductions — axis argument")
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───────── S18 · VS-CODE：axis=0 vs axis=1 ─────────
    s = _blank(prs)
    add_title(s, "axis=0 vs axis=1：同一個 mean 算出兩種答案")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.3),
        label="axis=0：沿列方向 → 每欄一個值（最常用）",
        code=(
            'scores = np.array([[80, 70, 90],\n'
            '                   [85, 75, 95],\n'
            '                   [90, 80, 100]])   # 3 生 × 3 科\n'
            '\n'
            'scores.mean(axis=0)\n'
            '# [85., 75., 95.]   每科平均 — shape (3,)'
        ),
        bullets=[
            "第 0 維被摺掉",
            "輸出 shape: (3 科,)",
            "EDA 最常用：每欄描述統計",
        ],
        label_dark=False,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(3.7),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.6),
        label="axis=1：沿欄方向 → 每列一個值",
        code=(
            'scores.mean(axis=1)\n'
            '# [80., 85., 90.]   每生平均 — shape (3,)\n'
            '\n'
            '# Debug 法則：看不懂 axis？先 print shape\n'
            'print(scores.shape)            # (3, 3)\n'
            'print(scores.mean(axis=0).shape)   # (3,) 欄\n'
            'print(scores.mean(axis=1).shape)   # (3,) 列'
        ),
        bullets=[
            "第 1 維被摺掉",
            "輸出 shape: (3 生,)",
            "遇到 axis bug\n先 print shape\n半秒破案",
        ],
        label_dark=True,
    )
    add_source(s, "NumPy User Guide §Axis — common pitfalls")
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ───────── S19 · PYRAMID：收束 ─────────
    s = _blank(prs)
    draw_thesis_hierarchy(
        s,
        blocks=[
            {"heading": "四件必帶走",
             "items": [
                 "shape：每個維度長度，debug 第一招",
                 "dtype：統一型別是效能的物理基礎",
                 "broadcasting：右對齊 / 補 1 / 拉伸三步",
                 "axis：axis=k 就是把第 k 維吃掉",
             ]},
            {"heading": "為什麼快 80 倍",
             "items": [
                 "連續記憶體：CPU cache 可預讀（Ch1 RAM 層級）",
                 "SIMD 指令：CPU 一次處理 4–8 筆",
                 "底層 C 實作：無 CPython 直譯開銷",
                 "無 PyObject 包裝：純數字沒穿外套",
             ]},
        ],
        title="Ch07 收束：shape × broadcasting × axis = 向量化思維",
        thesis="Ch08 Pandas 的 Series/DataFrame 底層就是 ndarray —— 本章四件套直接在表格上複用。",
    )
    add_source(s, "Ch07 module synthesis · Ch08 銜接草案")
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
