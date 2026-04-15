"""Ch06 deck — 檔案 I/O 與例外處理
17 content slides + cover + copyright page.

Governing thought：
    資料能進、錯誤能擋、大檔吃得下 ——
    這是數據工程的第一道門。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_rect, add_textbox, add_title, add_source,
    set_solid_fill, set_no_line, set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_vs_two_col, draw_code_panel, draw_thesis_hierarchy,
    draw_inverted_thesis_box, draw_image_placeholder,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch06"
MODULE_TITLE = "檔案 I/O 與例外處理"
MODULE_SUBTITLE = "資料能進、錯誤能擋、大檔吃得下"
TIME_MIN = 90
N_CONTENT = 17


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_ch06(output_path, image_registry=None):
    """Build Ch06 deck.

    image_registry: optional list collecting image placeholder metadata
                    for downstream YAML emission.
    """
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── S1 · SILENT ─────────
    s = _blank(prs)
    draw_silent_page(s, "資料能進、錯誤能擋 ——\n這是數據工程的第一道門。")
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ───────── S2 · ASK ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "為什麼生產環境 80% 的 bug，\n不是邏輯錯，是檔案長得跟你想的不一樣？",
        data_card={
            "label": "Airbnb Data Platform Post-mortem 2022",
            "stat": "82%",
            "caption": "資料管線故障源自\n格式 / 缺值 / 編碼問題",
        },
    )
    add_source(s, "Airbnb Engineering Blog · Data Platform Reliability Report 2022")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── S3 · MATRIX 2×3：pathlib vs os.path ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=3, cells=[
        {"text": "os.path.join(...)",
         "sub": "字串拼接\n跨平台要自己管分隔符"},
        {"text": "Path('data') / 'a.csv'",
         "sub": "/ 運算子重載\nOS 自動選對分隔符",
         "highlight": True},
        {"text": "os.listdir(d)",
         "sub": "回傳字串 list\n要再去 join 才能用"},
        {"text": "d.iterdir() / d.glob()",
         "sub": "回傳 Path 物件\n可直接繼續操作",
         "highlight": True},
        {"text": "os.path.splitext(p)[1]",
         "sub": "元組索引\n忘記是 [0] 還 [1] 常出錯"},
        {"text": "p.suffix / p.stem / p.stat()",
         "sub": "屬性化、可讀\n方法鏈一路到底",
         "highlight": True},
    ], title="pathlib 取代 os.path：跨平台六組對照")
    add_source(s, "Python 官方文件 pathlib §pathlib · PEP 428 The pathlib module")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───────── S4 · CODE：pathlib 實戰 ─────────
    s = _blank(prs)
    add_title(s, "pathlib 實戰：遍歷 / 讀入 / 建目錄三件組")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="三個最常用 pattern，佔日常 I/O 九成",
        code=(
            'from pathlib import Path\n'
            '\n'
            'data_dir = Path("data")\n'
            'for csv in data_dir.glob("*.csv"):\n'
            '    print(csv.name, csv.stat().st_size)\n'
            '\n'
            '# 小檔整顆進來\n'
            'text = (data_dir / "config.json").read_text(encoding="utf-8")\n'
            '\n'
            '# 確認路徑並建資料夾\n'
            'out = Path("output") / "2026" / "report.csv"\n'
            'out.parent.mkdir(parents=True, exist_ok=True)'
        ),
        bullets=[
            "glob('*.csv') 內建 pattern",
            "read_text / write_text\n小檔一行搞定",
            "mkdir(parents, exist_ok)\n是建資料夾鐵三角",
            "回傳物件可繼續鏈接操作",
        ],
        label_dark=True,
    )
    add_source(s, "Python Docs §pathlib · Real Python: Python pathlib tutorial")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── S5 · CODE：CSV 讀寫 ─────────
    s = _blank(prs)
    add_title(s, "CSV 讀寫：csv.DictReader 讓欄位有名有姓")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="純 Python csv 模組 — Pandas 讀不進來時的備胎",
        code=(
            'import csv\n'
            '\n'
            '# 讀：每列是 dict，欄位有名字\n'
            'with open("orders.csv", encoding="utf-8", newline="") as f:\n'
            '    reader = csv.DictReader(f)\n'
            '    for row in reader:\n'
            '        print(row["order_id"], row["amount"])\n'
            '\n'
            '# 寫：先給 header，再一列列寫\n'
            'rows = [{"id": 1, "name": "王小明"},\n'
            '        {"id": 2, "name": "陳小華"}]\n'
            'with open("out.csv", "w", encoding="utf-8", newline="") as f:\n'
            '    w = csv.DictWriter(f, fieldnames=["id", "name"])\n'
            '    w.writeheader()\n'
            '    w.writerows(rows)'
        ),
        bullets=[
            "newline='' 必加\n否則 Windows 多出空行",
            "DictReader > reader\n欄位名不會錯位",
            "中文必設 encoding='utf-8'",
            "Ch08 Pandas 會包掉\n但底層就是這個",
        ],
        label_dark=True,
    )
    add_source(s, "Python Docs §csv — CSV File Reading and Writing")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── S6 · CODE：JSON 中文陷阱 ─────────
    s = _blank(prs)
    add_title(s, "JSON 讀寫：ensure_ascii=False 是中文世界的鐵律")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="第一次寫完 JSON 變成 \\u53f0\\u5317 的那個坑",
        code=(
            'import json\n'
            '\n'
            '# 讀\n'
            'with open("config.json", encoding="utf-8") as f:\n'
            '    cfg = json.load(f)          # dict / list\n'
            '\n'
            '# 寫：中文要看得見\n'
            'data = {"city": "台北", "count": 1203}\n'
            '\n'
            'json.dumps(data)\n'
            '# \'{"city": "\\u53f0\\u5317", "count": 1203}\'   ← 預設跳脫\n'
            '\n'
            'json.dumps(data, ensure_ascii=False, indent=2)\n'
            '# \'{\\n  "city": "台北",\\n  "count": 1203\\n}\'  ← 正常'
        ),
        bullets=[
            "ensure_ascii=False\n幾乎永遠要加",
            "indent=2 給人看；\n機器傳輸時省略",
            "json.load(f) vs json.loads(s)\n差一個 s",
            "JSON 沒 date 型別\n日期要自己轉字串",
        ],
        label_dark=True,
    )
    add_source(s, "Python Docs §json — JSON encoder and decoder · RFC 8259")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───────── S7 · VS-CODE：沒 with vs 有 with ─────────
    s = _blank(prs)
    add_title(s, "沒 with vs 有 with：context manager 是 I/O 的止血閥")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.3),
        label="BEFORE：手動 close — 例外會讓它永遠不跑",
        code=(
            'f = open("data.csv", encoding="utf-8")\n'
            'rows = f.readlines()\n'
            'process(rows)            # 如果這裡炸了…\n'
            'f.close()                # 這行永遠不會跑'
        ),
        bullets=[
            "例外跳過 close\n檔案 handle 洩漏",
            "累積到 ulimit\n程式真的卡死",
            "靠 OS 收 = 設計壞掉",
        ],
        label_dark=False,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(3.7),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.5),
        label="AFTER：with — 離開區塊自動收尾，不論成敗",
        code=(
            'with open("data.csv", encoding="utf-8") as f:\n'
            '    rows = f.readlines()\n'
            '    process(rows)        # 這裡炸也沒關係\n'
            '# 離開 with 區塊自動 f.close()'
        ),
        bullets=[
            "__enter__ / __exit__\n保證收尾（Ch5 dunder）",
            "等價於隱式 try / finally",
            "DB 連線、檔案、lock\n全靠它",
            "Python「不寫 try 也安全」\n的主力機制",
        ],
        label_dark=True,
    )
    add_source(s, "PEP 343 The 'with' Statement · Python Docs §contextlib")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── S8 · IMAGE + CODE：四段骨架 ─────────
    s = _blank(prs)
    add_title(s, "try / except / else / finally：四段骨架各司其職")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.3), h=Inches(5.0),
        slot_name="try-except-else-finally 流程圖",
        description=(
            "垂直四區塊流程：try → 成功 else / 失敗 except\n"
            "兩岔匯流後一定進 finally\n"
            "except 區塊以深綠實心突顯攔截語意"
        ),
        url_hint="",
        placeholder_id="Ch06_S08_tryexcept_flow",
        registry=image_registry,
        size_hint="1280×900 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.2), y=Inches(1.3),
        w=Inches(6.5), h=Inches(5.0),
        label="完整四段骨架：each block has one job",
        code=(
            'try:\n'
            '    data = load("orders.csv")   # 可能炸\n'
            'except FileNotFoundError as e:\n'
            '    logger.warning(f"不存在：{e}")\n'
            '    data = []\n'
            'except ValueError as e:\n'
            '    logger.error(f"格式錯：{e}")\n'
            '    raise                        # 重新拋\n'
            'else:\n'
            '    # try 沒炸才跑\n'
            '    logger.info(f"讀入 {len(data)} 列")\n'
            'finally:\n'
            '    # 不論成敗都跑\n'
            '    close_db()'
        ),
        bullets=[
            "try：只包『會炸那行』\n別包太廣",
            "except：依型別分流",
            "else：try 沒炸才執行\n區分讀檔與用資料",
            "finally：釋放資源\n的最後機會",
        ],
        label_dark=True,
    )
    add_source(s, "Python Tutorial §8 Errors and Exceptions")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── S9 · MATRIX 2×2：四種最常見例外 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "FileNotFoundError",
         "sub": "檔案被搬走 / 路徑打錯\n觸發：open('not_there.csv')",
         "highlight": True},
        {"text": "KeyError",
         "sub": "dict / JSON 欄位不存在\n觸發：row['unknown_column']"},
        {"text": "ValueError",
         "sub": "型別對、值不合法\n觸發：int('abc')、float('N/A')",
         "highlight": True},
        {"text": "TypeError",
         "sub": "型別根本錯\n觸發：'3' + 5、len(None)"},
    ], title="四種最常見的例外：90% 資料管線的錯都在這四格")
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "背起來。except 時指定型別，不要用 bare except 一網打盡。",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "Python Docs §Built-in Exceptions · exception hierarchy")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── S10 · VS-CODE：反模式 ─────────
    s = _blank(prs)
    add_title(s, "except Exception: pass 是反模式：吞錯等於放火")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.3),
        label="BEFORE：Bare except — 把火吞進肚子",
        code=(
            'try:\n'
            '    data = pd.read_csv(path)\n'
            '    result = heavy_process(data)\n'
            '    save_to_db(result)\n'
            'except Exception:\n'
            '    pass              # 什麼都不做，假裝沒事'
        ),
        bullets=[
            "KeyboardInterrupt 也被吞",
            "bug 三週後才爆\n沒人找得到",
            "log 一行都沒留",
            "面試被退件的寫法",
        ],
        label_dark=False,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(3.7),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.5),
        label="AFTER：指定型別 + 留證據",
        code=(
            'try:\n'
            '    data = pd.read_csv(path)\n'
            'except FileNotFoundError:\n'
            '    logger.warning(f"跳過不存在檔案：{path}")\n'
            '    return None\n'
            'except pd.errors.ParserError as e:\n'
            '    logger.error(f"CSV 格式錯：{path} — {e}")\n'
            '    raise'
        ),
        bullets=[
            "只 catch 預期會發生的錯",
            "其他錯讓它炸\n讓 stack trace 幫你找 bug",
            "每個 except\n都要有 log 或 raise",
            "『吞錯』是團隊信用\n的頭號殺手",
        ],
        label_dark=True,
    )
    add_source(s, "PEP 8 §Programming Recommendations · Effective Python Item 87")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── S11 · CODE：自訂 Exception（呼應 Ch5） ─────────
    s = _blank(prs)
    add_title(s, "自訂 Exception：用 Ch5 的繼承語法，說出你的語意")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="Exception 就是普通 class — Ch5 的繼承今天派上用場",
        code=(
            'class DataValidationError(Exception):\n'
            '    """資料驗證失敗的語意錯誤"""\n'
            '    def __init__(self, column: str, reason: str):\n'
            '        self.column = column\n'
            '        self.reason = reason\n'
            '        super().__init__(f"[{column}] {reason}")\n'
            '\n'
            'class NegativeAgeError(DataValidationError):\n'
            '    def __init__(self, value):\n'
            '        super().__init__("age", f"age 不可為負，收到 {value}")\n'
            '\n'
            '# 使用\n'
            'if df["age"].min() < 0:\n'
            '    raise NegativeAgeError(df["age"].min())\n'
            '\n'
            '# 上層接\n'
            'try:\n'
            '    validate(df)\n'
            'except DataValidationError as e:\n'
            '    logger.error(f"欄位 {e.column}：{e.reason}")'
        ),
        bullets=[
            "Exception = 普通 class\n沿用 Ch5 繼承語法",
            "自訂型別讓 except\n可精準分流",
            "多攜帶欄位（column/reason）\n給 log 抓",
            "子類化再細分\n= 管線可讀性大躍進",
        ],
        label_dark=True,
    )
    add_source(s, "Python Tutorial §8.5 User-defined Exceptions · Ch05 繼承銜接")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── S12 · CODE：exception chaining ─────────
    s = _blank(prs)
    add_title(s, "raise ... from e：保留兇手指紋的正確姿勢")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="重新包裝錯誤時，別砍掉原始 traceback",
        code=(
            'class PipelineError(Exception):\n'
            '    pass\n'
            '\n'
            'def load_orders(path):\n'
            '    try:\n'
            '        return pd.read_csv(path)\n'
            '    except FileNotFoundError as e:\n'
            '        raise PipelineError(\n'
            '            f"orders 管線無法啟動：{path}"\n'
            '        ) from e\n'
            '        # ^^^^^^ 保留原始 traceback\n'
            '\n'
            '# 錯誤訊息長這樣：\n'
            '#   FileNotFoundError: ... \'orders.csv\'\n'
            '#   The above exception was the direct cause of\n'
            '#   the following exception:\n'
            '#   PipelineError: orders 管線無法啟動：orders.csv'
        ),
        bullets=[
            "raise X from e\n→ 兩層線索都保留",
            "raise X from None\n→ 刻意隱藏（少用）",
            "什麼都不寫\n→ 預設 implicit chain",
            "debug 時兩層 traceback\n都看得到",
        ],
        label_dark=True,
    )
    add_source(s, "PEP 3134 Exception Chaining and Embedded Tracebacks")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── S13 · ASK：大檔問題 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "50GB 的 CSV，你 16GB 的 RAM ——\n怎麼讀進來？",
        data_card={
            "label": "Ch1 OOM 回顧",
            "stat": "3×",
            "caption": "RAM 只夠檔案大小 1/3\n直接 read_csv 必 OOM",
        },
    )
    add_source(s, "Ch01 OS 與 RAM 章節 · Python Memory Management")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── S14 · IMAGE + CODE：Pandas chunking ─────────
    s = _blank(prs)
    add_title(s, "Chunking：把大檔切成 RAM 吃得下的小口")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.8), h=Inches(4.8),
        slot_name="Chunking 記憶體示意圖",
        description=(
            "上：50GB 大檔 vs 16GB RAM 視覺對比\n"
            "下：切成 5 個 chunk，sliding window 滑過\n"
            "底部註：讀入 → 處理 → 釋放 → 讀下一塊"
        ),
        url_hint="",
        placeholder_id="Ch06_S14_chunking_memory",
        registry=image_registry,
        size_hint="1400×700 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.7), y=Inches(1.3),
        w=Inches(6.0), h=Inches(4.8),
        label="Pandas chunking：一次只吃 10 萬列",
        code=(
            'import pandas as pd\n'
            '\n'
            'total_rows = 0\n'
            'total_amount = 0\n'
            '\n'
            'for chunk in pd.read_csv(\n'
            '        "huge_orders.csv",\n'
            '        chunksize=100_000,\n'
            '        dtype={"amount": "float32"}):\n'
            '    chunk = chunk[chunk["status"] == "paid"]\n'
            '    total_rows += len(chunk)\n'
            '    total_amount += chunk["amount"].sum()\n'
            '    # chunk 出區塊就被 GC\n'
            '\n'
            'print(f"{total_rows:,} 列 / ${total_amount:,.0f}")'
        ),
        bullets=[
            "chunksize 是 iterator",
            "每輪只載 10 萬列\nRAM 平穩",
            "常見 10 萬 ~ 100 萬列",
            "可累加（sum/count）\n不可全域排序",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §IO tools — Iterating through files chunk by chunk")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── S15 · CODE：純 Python 逐行讀 ─────────
    s = _blank(prs)
    add_title(s, "純 Python 逐行讀：檔案物件天生是 generator（呼應 Ch2）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="for line in f: 就是 streaming，不需要任何 import",
        code=(
            '# 純 Python 逐行 — 適用任何 text 檔\n'
            'total = 0\n'
            'with open("huge.log", encoding="utf-8") as f:\n'
            '    for line in f:                # 這就是 generator\n'
            '        if "ERROR" in line:\n'
            '            total += 1\n'
            'print(total)\n'
            '\n'
            '# 對比：錯誤寫法\n'
            'with open("huge.log") as f:\n'
            '    lines = f.readlines()         # 一次載全部 → RAM 爆\n'
            '    for line in lines: ...\n'
            '\n'
            '# 進階：寫成自己的 generator（呼應 Ch2）\n'
            'def tail_errors(path):\n'
            '    with open(path, encoding="utf-8") as f:\n'
            '        for line in f:\n'
            '            if "ERROR" in line:\n'
            '                yield line.strip()'
        ),
        bullets=[
            "for line in f:\n是懶載入",
            "readlines() 是眼殘寫法\nRAM 必爆",
            "自製 generator\n可 chain / filter / 複用",
            "這就是 Ch2 generator\n的生產場景",
        ],
        label_dark=True,
    )
    add_source(s, "Python Docs §io · Fluent Python Ch14 Iterables and Iterators")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── S16 · TABLE：本章邊界 ─────────
    s = _blank(prs)
    draw_editorial_table(
        s,
        header=["主題", "是什麼", "為何本章不教", "何時學"],
        rows=[
            ["Dask", "工業級 parallel DataFrame",
             "chunking 已夠、學習成本高",
             "真的要跨機器計算時"],
            ["Polars streaming", "Rust 寫的高效替代",
             "API 與 pandas 不同\n學習分岔成本",
             "Ch08 後想進階時"],
            ["asyncio / aiofiles", "非同步 I/O",
             "資料工程瓶頸多在 CPU\n非 I/O 併發",
             "寫 web backend 時"],
            ["pickle / struct", "二進位序列化",
             "現場 95% 用 CSV/JSON/Parquet",
             "寫框架 / 通訊協定時"],
            ["contextlib 自製 with", "自製 context manager",
             "Ch10 管線會自然碰到",
             "Ch10 之後"],
        ],
        col_widths=[1.4, 2.0, 2.2, 1.4],
        title="本章邊界：Linus 原則 — 解決真問題，不炫技",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.8),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
        "80% 場景直接可用，剩下 20% 遇到再查 —— 基本功夠強就好。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "課程設計紀律 · Linus Torvalds 實用主義")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── S17 · PYRAMID：收束 ─────────
    s = _blank(prs)
    draw_thesis_hierarchy(
        s,
        blocks=[
            {"heading": "三層骨架",
             "items": [
                 "進得來：pathlib + with + csv/json 讀寫",
                 "擋得住：try/except/else/finally + 自訂 Exception + raise from",
                 "吃得下：pd.read_csv(chunksize=) + 檔案 generator 逐行",
             ]},
            {"heading": "心法",
             "items": [
                 "只 catch 你預期的錯，其他讓它炸",
                 "自訂 Exception 讓 log 有結構（Ch5 繼承再次派上用場）",
                 "Streaming > batch：能 chunk 就 chunk",
             ]},
        ],
        title="Ch06 收束：進得來 + 擋得住 + 吃得下 = 可上線的 I/O 層",
        thesis="Ch07 進入 NumPy — 資料進得來、擋得住、吃得下之後，我們開始把它算快。",
    )
    add_source(s, "Ch06 module synthesis · Ch07 NumPy 銜接")
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
