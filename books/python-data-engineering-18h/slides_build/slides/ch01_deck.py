"""Ch01 deck — 計算機概論前導與 OS 角色. 17 content slides + cover + copyright."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_rect, add_textbox, add_title, add_source,
    set_solid_fill, set_no_line, set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_pyramid_stack, draw_inverted_thesis_box,
    draw_vs_two_col, draw_risk_mitigation, draw_grid,
    draw_image_placeholder,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch01"
MODULE_TITLE = "計算機概論前導與 OS 角色"
MODULE_SUBTITLE = "Python 進階數據工程與分析 · M1 系統前導與 Python 機制"
TIME_MIN = 90
N_CONTENT = 17


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_ch01(output_path, image_registry=None):
    """Build Ch01 deck.

    image_registry: optional list to collect image placeholder metadata
                    for downstream YAML emission.
    """
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE, TIME_MIN, N_CONTENT)

    # ---------- S1 SILENT ----------
    s = _blank(prs)
    draw_silent_page(s, '你寫的每一行 Python，\n都是向 OS 借硬體的一張申請單。')
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ---------- S2 ASK ----------
    s = _blank(prs)
    draw_ask_page(
        s,
        '為什麼同一份 CSV，pandas 讀 3 秒、純 Python 要讀 40 秒？',
        data_card={
            'label': '同檔 CSV 載入基準',
            'stat': '13.5×',
            'caption': 'pandas 3.1s vs csv.reader 41.7s\n1000 萬列 × 6 欄 / i7 + SSD',
        },
    )
    add_source(s, '本課實測 2025 Q4 · CPython 3.11 / pandas 2.2')
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ---------- S3 FLOW · 三代語言 ----------
    s = _blank(prs)
    draw_flow_chain(s, nodes=[
        {'label': '機器碼', 'sub': '1940s', 'caption': '01010101 · CPU 直讀'},
        {'label': '組合語言', 'sub': '1950s', 'caption': 'MOV / ADD · 1:1 指令'},
        {'label': '高階語言', 'sub': '1960s–', 'caption': 'Python / C / Java',
         'highlight': True},
    ], title='三代語言 = 三次抽象交易：換可讀性、付效能代價',
       y=3.2)
    add_textbox(s, T.MARGIN_X, Inches(5.4),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
                '抽象越高、離硬體越遠、開發越快、效能代價越透明。',
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, '計算機史標準教材整理')
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ---------- S4 VS · 編譯 vs 直譯 ----------
    s = _blank(prs)
    draw_vs_two_col(s,
        left_title='編譯式（C / Rust / Go）',
        right_title='直譯式（Python / JS / Ruby）',
        left_items=[
            '一次翻譯、多次執行',
            '啟動慢、執行快',
            '錯誤在編譯期暴露',
            '跨平台需重新編譯',
            '適合：系統層、高效能服務',
        ],
        right_items=[
            '逐行翻譯、逐行執行',
            '啟動快、執行慢',
            '錯誤在執行期才爆',
            '同一份 .py 跨平台直跑',
            '適合：資料科學、快速原型',
        ],
        title='編譯 vs 直譯：一個用啟動時間換執行速度，一個反過來',
        summary='Python 用直譯拿開發速度，再用 C 後端把執行速度偷回來——這就是「黏合語言」。',
        delta='70% C 後端',
    )
    add_source(s, 'PyPI 下載 TOP 100 資料套件抽樣 2025')
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ---------- S5 MATRIX · Python 黏合語言 ----------
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=3, cells=[
        {'text': '表格運算層',
         'sub': 'pandas / Polars\nPython 介面 + C/Rust 核心'},
        {'text': '數值運算層',
         'sub': 'NumPy / SciPy\nPython 介面 + C/Fortran BLAS'},
        {'text': '深度學習層',
         'sub': 'PyTorch / TensorFlow\nPython 介面 + C++/CUDA'},
        {'text': '爬蟲解析層',
         'sub': 'lxml / orjson\nPython 介面 + C 核心'},
        {'text': '資料庫驅動',
         'sub': 'psycopg2 / PyMySQL\nPython 介面 + C 核心'},
        {'text': '你寫 Python，跑的是 C。\nNumPy 比純 Python 快 10–100×。',
         'highlight': True},
    ], title='Python 的雙層結構：上層是你，下層是 C')
    add_source(s, '各套件官方文件 2025 Q1')
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ---------- S6 MATRIX · Bit / Byte 階梯 ----------
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=3, cells=[
        {'text': 'Bit',
         'sub': '0 或 1，1 個電位\n最小資訊單位'},
        {'text': 'Byte',
         'sub': '8 Bit = 256 種狀態\n一個 ASCII 字元'},
        {'text': 'KB / MB',
         'sub': '文件、一張圖、\n幾秒音訊'},
        {'text': 'GB',
         'sub': '一部電影、\n百萬筆資料表'},
        {'text': 'TB / PB',
         'sub': '公司資料湖、\n搜尋引擎 index'},
        {'text': '估算資料量 → 選工具\n→ 估 RAM → 避免 OOM。',
         'highlight': True},
    ], title='資料量的階梯：每一級 1024 倍，你要有量級感')
    add_source(s, '本課整理')
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ---------- S7 CHART-like · 1 億筆估算 ----------
    s = _blank(prs)
    add_title(s, '估算的黃金公式：列數 × 欄數 × 每格 Bytes')
    # 等式區
    add_textbox(s, T.MARGIN_X, Inches(1.5),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
                '1e8 列  ×  1 欄  ×  8 Bytes  =  8e8 Bytes',
                font_size=Pt(26), color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER)
    add_textbox(s, T.MARGIN_X, Inches(2.2),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
                '≈ 763 MiB  ≈  800 MB',
                font_size=Pt(22), color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)

    # 量級錨點三格（像小 matrix）
    anchors = [
        ('800 MB', '1 欄 × 1e8 列', '單欄輕鬆'),
        ('8 GB', '10 欄 × 1e8 列', '剛好塞滿筆電 RAM'),
        ('16 GB', '20 欄 × 1e8 列', 'OOM 風險高'),
    ]
    col_w = Inches(3.8)
    gap = Inches(0.3)
    total_anchor_w = col_w * 3 + gap * 2
    start_x = (T.SLIDE_W - total_anchor_w) / 2
    anchor_y = Inches(3.4)
    anchor_h = Inches(1.8)
    for i, (big, mid, sub) in enumerate(anchors):
        x = start_x + i * (col_w + gap)
        rect = add_rect(s, x, anchor_y, col_w, anchor_h)
        if i == 2:
            set_solid_fill(rect, T.PRIMARY)
            set_no_line(rect)
            c_big = T.WHITE
            c_mid = T.LIGHT_GRAY
            c_sub = T.LIGHT_GRAY
        else:
            set_no_fill(rect)
            set_line(rect, T.PRIMARY, 1.0)
            c_big = T.PRIMARY
            c_mid = T.CHARCOAL
            c_sub = T.GRAY_MID
        add_textbox(s, x, anchor_y + Inches(0.2), col_w, Inches(0.6),
                    big, font_size=Pt(28), color=c_big, bold=True,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, x, anchor_y + Inches(0.9), col_w, Inches(0.4),
                    mid, font_size=T.FONT_BODY, color=c_mid,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, x, anchor_y + Inches(1.3), col_w, Inches(0.4),
                    sub, font_size=T.FONT_SMALL, color=c_sub, italic=True,
                    align=PP_ALIGN.CENTER)

    draw_inverted_thesis_box(
        s, '這個公式你會在本課至少用到五次——能估才能選工具。',
        y=5.6, width=11.0)
    add_source(s, 'float64 = 8 Bytes；本課教學估算')
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ---------- S8 PHOTO · 馮紐曼架構 ----------
    s = _blank(prs)
    add_title(s, '馮紐曼架構（1945）——80 年不變的三分法')
    photo_w = Inches(10.0)
    photo_h = Inches(4.2)
    photo_x = (T.SLIDE_W - photo_w) / 2
    photo_y = Inches(1.3)
    draw_image_placeholder(
        s, photo_x, photo_y, photo_w, photo_h,
        slot_name='馮紐曼架構示意圖',
        description='von Neumann architecture：CPU（ALU + Control Unit）· Memory · I/O 三區，以 Data/Address/Control Bus 串接',
        url_hint='https://en.wikipedia.org/wiki/Von_Neumann_architecture',
        size_hint='2000×900 px',
        placeholder_id='Ch01_S08_vonneumann',
        registry=image_registry,
    )
    add_textbox(s, T.MARGIN_X, Inches(5.8),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
                '運算、暫存、進出——三個區的分工，是今天所有效能問題的根。',
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, 'von Neumann, First Draft of a Report on the EDVAC, 1945')
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ---------- S9 VS + PHOTO · CPU vs GPU ----------
    s = _blank(prs)
    add_title(s, 'CPU vs GPU：少數強者 vs 多數平庸者')

    # 左右兩張 CPU / GPU 照片佔位
    photo_top = Inches(1.3)
    photo_h9 = Inches(2.2)
    left_x = Inches(1.2)
    right_x = Inches(7.3)
    photo_col_w = Inches(4.8)

    draw_image_placeholder(
        s, left_x, photo_top, photo_col_w, photo_h9,
        slot_name='CPU 晶片特寫',
        description='現代 x86/ARM CPU 裸片或封裝特寫；強調少而強的核心佈局',
        url_hint='https://www.intel.com/content/www/us/en/newsroom/',
        size_hint='1400×700 px',
        placeholder_id='Ch01_S09_cpu_photo',
        registry=image_registry,
    )
    draw_image_placeholder(
        s, right_x, photo_top, photo_col_w, photo_h9,
        slot_name='GPU 晶片特寫',
        description='現代 GPU 裸片或板卡特寫；強調多而密集的平行單元佈局',
        url_hint='https://nvidianews.nvidia.com/multimedia',
        size_hint='1400×700 px',
        placeholder_id='Ch01_S09_gpu_photo',
        registry=image_registry,
    )

    # 左右條目欄
    list_top = photo_top + photo_h9 + Inches(0.2)
    list_h = Inches(2.4)
    # Left column (CPU)
    lrect = add_rect(s, left_x, list_top, photo_col_w, list_h)
    set_no_fill(lrect)
    set_line(lrect, T.PRIMARY, 1.0)
    add_textbox(s, left_x, list_top, photo_col_w, Inches(0.4),
                'CPU · 少數強者',
                font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, left_x + Inches(0.2), list_top + Inches(0.45),
                photo_col_w - Inches(0.4), list_h - Inches(0.5),
                '• 8–16 核、時脈 3–5 GHz\n'
                '• 擅長：複雜分支、多工切換\n'
                '• 典型場景：OS、Web、一般程式\n'
                '• 設計哲學：少而強',
                font_size=T.FONT_CAPTION, color=T.CHARCOAL,
                line_spacing=1.4)

    # Right column (GPU)
    rrect = add_rect(s, right_x, list_top, photo_col_w, list_h)
    set_no_fill(rrect)
    set_line(rrect, T.PRIMARY, 1.0)
    add_textbox(s, right_x, list_top, photo_col_w, Inches(0.4),
                'GPU · 多數平庸者',
                font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, right_x + Inches(0.2), list_top + Inches(0.45),
                photo_col_w - Inches(0.4), list_h - Inches(0.5),
                '• 數千 CUDA 核、時脈 1–2 GHz\n'
                '• 擅長：同構大量平行\n'
                '• 典型場景：矩陣、渲染、神經網路\n'
                '• 設計哲學：多而密集',
                font_size=T.FONT_CAPTION, color=T.CHARCOAL,
                line_spacing=1.4)

    # 中間 VS
    add_textbox(s, Inches(6.0), list_top + list_h / 2 - Inches(0.3),
                Inches(1.3), Inches(0.6),
                'VS', font_size=Pt(26), color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_source(s, 'NVIDIA H100 / Intel Xeon 規格書 · 2024')
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ---------- S10 FLOW · 儲存金字塔 ----------
    s = _blank(prs)
    draw_flow_chain(s, nodes=[
        {'label': 'Register', 'sub': '<1 ns', 'caption': 'B 量級'},
        {'label': 'Cache', 'sub': '1–10 ns', 'caption': 'KB–MB'},
        {'label': 'RAM', 'sub': '~100 ns', 'caption': 'GB · 資料主戰場',
         'highlight': True},
        {'label': 'SSD', 'sub': '~100 µs', 'caption': '100 GB – TB'},
        {'label': 'HDD', 'sub': '~10 ms', 'caption': 'TB 量級'},
    ], title='儲存金字塔：上下相差 10^7 倍', y=3.0)
    draw_inverted_thesis_box(
        s, '資料搬到 RAM 就別再下去——這是所有效能優化的第一條口訣。',
        y=5.8, width=11.0)
    add_source(s, 'Jeff Dean, Latency Numbers Every Programmer Should Know')
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ---------- S11 ASK · OS 是守門人 ----------
    s = _blank(prs)
    draw_ask_page(
        s,
        '你的 Python 程式從來不能直接摸硬體——\n所有請求都要經過 OS 這位守門人。',
        data_card={
            'label': 'System Call',
            'stat': '每一次',
            'caption': 'open() / read() / fork() / mmap()\n背後都是一次 syscall',
        },
    )
    add_source(s, 'POSIX.1-2017 / Linux man-pages')
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ---------- S12 PYRAMID · OS 三大職責 ----------
    s = _blank(prs)
    draw_pyramid_stack(s,
        layers=[
            {'name': 'Scheduler', 'caption': '誰先跑、跑多久、何時切換\nProcess / Thread / GIL'},
            {'name': 'File System', 'caption': '檔案組織、路徑解析、權限\npathlib / mount / inode'},
            {'name': 'Memory Manager', 'caption': '誰的 RAM、誰進 Swap\nOOM Killer / Page Table'},
            {'name': 'Hardware', 'caption': 'CPU · RAM · Disk · Network'},
        ],
        cross_cuts=['權限控管', 'Interrupt'],
        thesis='不懂 OS，你的 Python 程式永遠只是在黑盒上跳舞。',
        title='OS 三大職責：記憶體、檔案、行程——數據工程的三座隱形山',
    )
    add_source(s, 'Operating System Concepts, Silberschatz 10e')
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ---------- S13 RISK-MITIGATION · OOM / Swap ----------
    s = _blank(prs)
    draw_risk_mitigation(s,
        risks=[
            '一次 pd.read_csv 大檔：整包進 RAM',
            'groupby 後 pivot：中介表格指數膨脹',
            'DataFrame 鏈式 .copy()：同資料在 RAM 複製 5 份',
            'Swap 啟動：效能掉 100–1000×，系統半死',
        ],
        mitigations=[
            '用 chunksize= / Polars / Dask 分塊讀',
            '先 filter 再 groupby；observed=True',
            '用 inplace=True 或明確 view；斬斷 copy 鏈',
            '監看 free -h / psutil.virtual_memory()，提早退場',
        ],
        title='OOM 與 Swap：能估算、能預防的必然，不是意外',
        summary='反射動作：寫完 pd.read_csv(...) 下一秒問自己「它幾 GB、我有幾 GB」——這救你 80% 的 OOM。',
        risks_title='Risk（OOM 來源）',
        miti_title='Mitigation（你的對策）',
    )
    add_source(s, '本課 OOM 事故盤點 2025')
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ---------- S14 MATRIX · pathlib ----------
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {'text': '壞寫法：字串拼路徑',
         'sub': 'path = folder + "\\\\" + name\n\nWindows 硬編、Linux 爆炸\n部署到雲端第一個炸'},
        {'text': '好寫法：pathlib',
         'sub': 'path = Path(folder) / name\n\nOS 自動選分隔符\n語意清晰、零成本'},
        {'text': '壞寫法：相對路徑',
         'sub': 'open("data/2024/q1/sales.csv")\n\n看你從哪個目錄執行\n同事 clone 下來就壞'},
        {'text': '用 pathlib 不是品味問題，\n是跨平台可移植性的最低保險。',
         'highlight': True},
    ], title='檔案系統跨平台小坑：Windows \\ vs POSIX / ')
    add_source(s, 'PEP 428 · Python 3.4+ pathlib')
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ---------- S15 VS · I/O vs CPU Bound ----------
    s = _blank(prs)
    draw_vs_two_col(s,
        left_title='I/O Bound · 瓶頸在等',
        right_title='CPU Bound · 瓶頸在算',
        left_items=[
            '症狀：CPU 很閒、在等待',
            '典型：爬蟲、API、讀大檔、資料庫',
            '對策：async / asyncio / threading',
            'GIL 不影響 I/O，等待可重疊',
            '診斷：watch CPU% 很低',
        ],
        right_items=[
            '症狀：CPU 100%、程式跑不動',
            '典型：特徵工程、矩陣、模型訓練',
            '對策：multiprocessing / 向量化 / GPU',
            'GIL 鎖住 thread，需開新進程',
            '診斷：watch CPU% 撞滿',
        ],
        title='I/O Bound vs CPU Bound：瓶頸在等還是在算？',
        summary='診斷清楚再選工具——選錯，開 100 條 thread 也跑不快。',
        delta='選錯 = 白做工',
    )
    add_source(s, 'Python 官方文件 · concurrent.futures / asyncio 2025')
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ---------- S16 FLOW · read_csv 一生 ----------
    s = _blank(prs)
    draw_flow_chain(s, nodes=[
        {'label': '你', 'sub': 'Step 1', 'caption': 'df = pd.read_csv(...)'},
        {'label': 'CPython', 'sub': 'Step 2', 'caption': '翻成 Bytecode'},
        {'label': 'pandas', 'sub': 'Step 3', 'caption': 'C 後端發 syscall'},
        {'label': 'OS + HW', 'sub': 'Step 4', 'caption': '定位 → 讀磁碟 → RAM',
         'highlight': True},
        {'label': 'CPU', 'sub': 'Step 5', 'caption': '建構 DataFrame'},
        {'label': 'df', 'sub': 'Step 6', 'caption': '物件回到你手上'},
    ], title='一行 pd.read_csv(\'data.csv\') 的完整旅程',
       y=3.0)
    draw_inverted_thesis_box(
        s, '你寫的一行，牽動五個角色——這就是本章要給你的全景圖。',
        y=5.8, width=11.5)
    add_source(s, '本課整理 · 對應 Ch2/Ch5/Ch6/Ch7/Ch9')
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ---------- S17 SILENT closing ----------
    s = _blank(prs)
    draw_silent_page(s,
        '讀懂這張地圖後，\n後面每一章的效能章節，\n都只是在這張地圖上標註座標。')
    add_footer(s, MODULE_CODE, 17, N_CONTENT, dark_bg=True)

    # Copyright
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
