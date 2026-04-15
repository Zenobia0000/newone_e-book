"""Ch09 deck — Matplotlib 探索性視覺化
17 content slides + cover + copyright page (M4 · 1.5 hr).

Governing thought:
    圖不是漂亮，是問題的顯影 ——
    選對四種 EDA 圖（折線／散佈／直方／Box），
    用 Figure/Axes 雙層結構 + subplots，
    讓資料自己開口說話。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_rect, add_textbox, add_title, add_source,
    set_solid_fill, set_no_line, set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_vs_two_col, draw_code_panel, draw_thesis_hierarchy,
    draw_inverted_thesis_box, draw_image_placeholder,
    draw_flow_chain,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch09"
MODULE_TITLE = "Matplotlib 探索性視覺化"
MODULE_SUBTITLE = "Figure / Axes + 四種 EDA 圖 — 讓資料自己說故事"
TIME_MIN = 90
N_CONTENT = 17


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_ch09(output_path, image_registry=None):
    """Build Ch09 deck.

    image_registry: optional list collecting image placeholder metadata
                    for downstream YAML emission.
    """
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── S1 · SILENT ─────────
    s = _blank(prs)
    draw_silent_page(s, "圖不是漂亮，\n是問題的顯影。")
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ───────── S2 · ASK ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "為什麼大家 plot 完\n還要查文件半小時？",
        data_card={
            "label": "Matplotlib 官方 Cheatsheet 下載量（2024）",
            "stat": "4 M+",
            "caption": "「用過就忘」是 Matplotlib 的常態 —\n因為大家背 API，沒建心智模型",
        },
    )
    add_source(s, "Matplotlib Cheatsheets repo stats 2024 · matplotlib/cheatsheets")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── S3 · MATRIX 2×2：EDA 四問題 → 四圖 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "分布如何？",
         "sub": "→ 直方圖 ax.hist\n看鐘形 / 偏態 / 多峰\n單欄第一眼",
         "highlight": True},
        {"text": "有趨勢嗎？",
         "sub": "→ 折線圖 ax.plot\n時間或有序序列\n斜率 = 訊號",
         "highlight": True},
        {"text": "兩變數有關係嗎？",
         "sub": "→ 散佈圖 ax.scatter\n兩連續變數\n關係 + 離群一次看"},
        {"text": "有離群 / 分組差嗎？",
         "sub": "→ Box plot ax.boxplot\n多組分布比較\n空間效率王"},
    ], title="EDA 要回答的四個問題 → 四種必備圖")
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "看圖前先問：「我在回答什麼問題？」— 問題決定圖，不是品味決定圖。",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "Tukey, Exploratory Data Analysis 1977 · Matplotlib User Guide §Plot types")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───────── S4 · IMAGE：Figure/Axes 雙層結構 ─────────
    s = _blank(prs)
    add_title(s, "Figure / Axes 雙層結構：畫布 + 座標系")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(7.8), h=Inches(5.0),
        slot_name="Figure / Axes 雙層結構圖",
        description=(
            "外層大框 Figure（整張畫布、存檔單位）\n"
            "內層一至多個 Axes（座標系）\n"
            "title / xlabel / ylabel / spines / legend\n"
            "全部屬於 Axes"
        ),
        url_hint="https://matplotlib.org/stable/users/explain/figure/figure_intro.html",
        placeholder_id="Ch09_S04_figure_axes_structure",
        registry=image_registry,
        size_hint="1400×900 px",
    )
    # Right-side bullet column
    bullets = (
        "• Figure = 整張畫布\n"
        "   可 savefig 存檔的最小單位\n\n"
        "• Axes = 畫布上一個座標系\n"
        "   title / xlabel / ylabel / legend\n"
        "   都是 Axes 的屬性\n\n"
        "• 一個 Figure 可放多個 Axes\n"
        "   → subplots 就在做這件事\n\n"
        "• 搞懂這層 → 所有 API 各就各位"
    )
    add_textbox(
        s, Inches(8.7), Inches(1.3),
        Inches(4.0), Inches(5.0),
        bullets,
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    add_source(s, "Matplotlib User Guide §Introduction · §Parts of a Figure")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── S5 · VS-CODE：pyplot vs OO ─────────
    s = _blank(prs)
    add_title(s, "狀態式 pyplot vs 物件導向 OO：為何業界推薦 OO")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.3),
        label="BEFORE：pyplot 狀態式 — 依賴「當前 Figure / Axes」全域狀態",
        code=(
            'import matplotlib.pyplot as plt\n'
            'plt.plot(x, y)\n'
            'plt.title("Sales")   # 作用對象？「當前的」Axes\n'
            'plt.xlabel("Month")  # 多圖切換時容易畫錯地方'
        ),
        bullets=[
            "上手快、範例短",
            "多圖、多 Axes 易混",
            "全域狀態不可預測",
        ],
        label_dark=False,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(3.7),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.5),
        label="AFTER：OO 物件導向 — 明確操作對象（業界 99% 用這寫法）",
        code=(
            'import matplotlib.pyplot as plt\n'
            'fig, ax = plt.subplots(figsize=(8, 5))\n'
            'ax.plot(x, y)\n'
            'ax.set_title("Sales")     # 明確：這個 Axes 的標題\n'
            'ax.set_xlabel("Month")    # 明確：這個 Axes 的 x 軸\n'
            'ax.set_ylabel("Revenue")'
        ),
        bullets=[
            "操作對象明確",
            "可預測、可組合",
            "多圖擴展零心智負擔",
            "Linus 觀點：明確 > 隱式",
        ],
        label_dark=True,
    )
    add_source(s, "Matplotlib docs §The Lifecycle of a Plot · §Coding styles")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── S6 · CODE：所有圖的共同骨架 ─────────
    s = _blank(prs)
    add_title(s, "fig, ax = plt.subplots()：所有圖的共同骨架")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="背下這 5 行 — 之後所有圖只換中間一行",
        code=(
            'import matplotlib.pyplot as plt\n'
            '\n'
            '# 共同骨架（5 行）\n'
            'fig, ax = plt.subplots(figsize=(8, 5))\n'
            'ax.plot(x, y)                       # ← 只有這行會換\n'
            'ax.set_title("月度銷售")\n'
            'ax.set_xlabel("月份"); ax.set_ylabel("金額 (千元)")\n'
            'plt.show()\n'
            '\n'
            '# 換成其他圖種 — 只改中間那一行\n'
            'ax.hist(df["score"], bins=20)       # 直方\n'
            'ax.scatter(df["h"], df["w"], alpha=0.6)   # 散佈\n'
            'ax.boxplot([g1, g2, g3], labels=["A","B","C"])   # box'
        ),
        bullets=[
            "figsize 單位 = 英吋",
            "subplots() 預設 1×1\n回傳 (Figure, Axes)",
            "plt.show() 只有 script 需要\nNotebook 自動顯示",
            "同骨架吃掉 95% 日常繪圖",
        ],
        label_dark=True,
    )
    add_source(s, "Matplotlib Quick start guide · §Coding styles")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───────── S7 · SILENT：選對圖 ─────────
    s = _blank(prs)
    draw_silent_page(s, "選對圖，\n資料就會自己說故事。")
    add_footer(s, MODULE_CODE, 7, N_CONTENT, dark_bg=True)

    # ───────── S8 · IMAGE + CODE：折線圖 ─────────
    s = _blank(prs)
    add_title(s, "折線圖 ax.plot：趨勢與時間序列")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.5), h=Inches(5.0),
        slot_name="折線圖範例（月度銷售）",
        description=(
            "深綠折線 + 空心圓 marker\n"
            "x 軸 = 月份、y 軸 = 金額\n"
            "可加第二條 2023 比較線"
        ),
        url_hint="",
        placeholder_id="Ch09_S08_lineplot_sample",
        registry=image_registry,
        size_hint="1200×720 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.4), y=Inches(1.3),
        w=Inches(6.3), h=Inches(5.0),
        label="趨勢 / 時間序列 — ax.plot 是直覺反應",
        code=(
            'fig, ax = plt.subplots(figsize=(8, 5))\n'
            '\n'
            '# 單條線\n'
            'ax.plot(df["month"], df["sales"],\n'
            '        marker="o", color="#1B5E3F")\n'
            '\n'
            '# 多條線 → label + legend\n'
            'ax.plot(m, s_2023, label="2023", color="#808080")\n'
            'ax.plot(m, s_2024, label="2024", color="#1B5E3F")\n'
            'ax.legend(loc="upper left", frameon=False)\n'
            '\n'
            'ax.set_title("月度銷售趨勢")\n'
            'ax.set_xlabel("月份"); ax.set_ylabel("金額 (千元)")'
        ),
        bullets=[
            "時間 / 有序序列用 plot",
            "多條線要 label + legend",
            "marker 幫離散點看清楚",
            "x 軸過密 → fig.autofmt_xdate()",
        ],
        label_dark=True,
    )
    add_source(s, "Matplotlib gallery §Line plots · §Multiple subplots")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── S9 · IMAGE + CODE：散佈圖 ─────────
    s = _blank(prs)
    add_title(s, "散佈圖 ax.scatter：兩變數關係與離群偵測")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.5), h=Inches(5.0),
        slot_name="散佈圖範例（身高 vs 體重）",
        description=(
            "半透明深綠圓點（alpha=0.6）\n"
            "右上一顆明顯離群 → 箭頭標 Outlier\n"
            "關係 + 離群兩個問題一張圖解答"
        ),
        url_hint="",
        placeholder_id="Ch09_S09_scatter_sample",
        registry=image_registry,
        size_hint="1200×720 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.4), y=Inches(1.3),
        w=Inches(6.3), h=Inches(5.0),
        label="兩連續變數關係 — scatter 一眼看完",
        code=(
            'fig, ax = plt.subplots(figsize=(8, 5))\n'
            '\n'
            '# 基本散佈圖\n'
            'ax.scatter(df["height"], df["weight"],\n'
            '           alpha=0.6, color="#1B5E3F")\n'
            '\n'
            '# 第三維用顏色（類別 → c + cmap）\n'
            'ax.scatter(df["height"], df["weight"],\n'
            '           c=df["gender_code"],   # 0/1 類別\n'
            '           cmap="viridis", alpha=0.7)\n'
            '\n'
            'ax.set_title("身高 vs 體重")\n'
            'ax.set_xlabel("身高 (cm)"); ax.set_ylabel("體重 (kg)")'
        ),
        bullets=[
            "alpha 降透明 → 看密度",
            "重疊嚴重而不用 alpha\n= 看起來只有一坨",
            "c + cmap 做第三維著色",
            "離群點從 scatter 一眼看出",
        ],
        label_dark=True,
    )
    add_source(s, "Matplotlib gallery §Scatter plots · McKinney, PDA 3e §9.2")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── S10 · IMAGE + CODE：直方圖 ─────────
    s = _blank(prs)
    add_title(s, "直方圖 ax.hist：單變數分布的第一眼")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.5), h=Inches(5.0),
        slot_name="直方圖範例（考試成績）",
        description=(
            "深綠柱體 + 白色 edgecolor\n"
            "bins=20、右偏分布\n"
            "集中在 70~85 分"
        ),
        url_hint="",
        placeholder_id="Ch09_S10_hist_sample",
        registry=image_registry,
        size_hint="1200×720 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.4), y=Inches(1.3),
        w=Inches(6.3), h=Inches(5.0),
        label="單變數分布 — hist 看鐘形 / 偏態 / 多峰",
        code=(
            'fig, ax = plt.subplots(figsize=(8, 5))\n'
            '\n'
            '# 基本直方\n'
            'ax.hist(df["score"], bins=20,\n'
            '        color="#1B5E3F", edgecolor="white")\n'
            '\n'
            '# 比較兩組分布 → alpha 疊圖\n'
            'ax.hist(g1, bins=20, alpha=0.5, label="A組")\n'
            'ax.hist(g2, bins=20, alpha=0.5, label="B組")\n'
            'ax.legend()\n'
            '\n'
            '# density=True → 機率密度（面積=1）\n'
            'ax.hist(x, bins=30, density=True)'
        ),
        bullets=[
            "bins 是唯一關鍵參數",
            "預設 10 常常不夠\n先試 20~30",
            "bins 太少 → 失真\nbins 太多 → 鋸齒",
            "比較分布 alpha=0.5 疊圖",
        ],
        label_dark=True,
    )
    add_source(s, "Matplotlib docs §hist · Tukey EDA 1977 §Stem-and-leaf")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── S11 · IMAGE + CODE：Box plot ─────────
    s = _blank(prs)
    add_title(s, "Box plot ax.boxplot：分組比較與四分位數")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.5), h=Inches(5.0),
        slot_name="Box plot 範例（三城市房價）",
        description=(
            "三個並排 box：Taipei / Taichung / Kaohsiung\n"
            "顯示 median 線、IQR box、whisker、離群點\n"
            "y 軸 = 房價（百萬）"
        ),
        url_hint="",
        placeholder_id="Ch09_S11_box_sample",
        registry=image_registry,
        size_hint="1200×720 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.4), y=Inches(1.3),
        w=Inches(6.3), h=Inches(5.0),
        label="多組分布比較 — box plot 空間效率王",
        code=(
            'fig, ax = plt.subplots(figsize=(8, 5))\n'
            '\n'
            'cities = ["Taipei", "Taichung", "Kaohsiung"]\n'
            'data = [df[df.city == c]["price"] for c in cities]\n'
            '\n'
            'ax.boxplot(data, labels=cities)\n'
            'ax.set_title("三都房價分布比較")\n'
            'ax.set_ylabel("房價 (百萬)")\n'
            '\n'
            '# 橫式更利閱讀長標籤\n'
            'ax.boxplot(data, labels=cities, vert=False)\n'
            '\n'
            '# box 元素解讀：\n'
            '#   box   = IQR (Q1~Q3)\n'
            '#   線    = 中位數\n'
            '#   whisker = 1.5 × IQR\n'
            '#   點    = 離群'
        ),
        bullets=[
            "分組比較 / 離群偵測",
            "N 個分布一張圖比\nN 張 hist 更有用",
            "vert=False 轉橫式\n適合長標籤",
            "兩組用 box 還好\n三組以上價值才體現",
        ],
        label_dark=True,
    )
    add_source(s, "Tukey, Exploratory Data Analysis 1977 §Box-and-Whisker · Matplotlib docs §boxplot")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── S12 · TABLE：四圖速查 ─────────
    s = _blank(prs)
    draw_editorial_table(
        s,
        header=["圖種", "何時用", "API", "常見誤用"],
        rows=[
            ["折線 plot",
             "時間 / 有序序列",
             "ax.plot(x, y, marker, label)",
             "類別資料硬套折線 → 錯覺趨勢"],
            ["散佈 scatter",
             "兩連續變數關係、離群",
             "ax.scatter(x, y, alpha, c, cmap)",
             "大量重疊不用 alpha → 看成一坨"],
            ["直方 hist",
             "單變數分布（鐘形／偏態）",
             "ax.hist(x, bins=N, density)",
             "bins 沒調 → 分布失真"],
            ["Box plot",
             "多組分布比較、離群偵測",
             "ax.boxplot(data, labels, vert)",
             "只看兩組時 hist 其實更直覺"],
        ],
        col_widths=[1.4, 2.4, 3.0, 3.0],
        title="四圖速查：何時選誰、常見誤用、解讀重點",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.4),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "選圖 = 選問題，不是選品味。這四張吃掉日常 EDA 的 90%。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "Matplotlib gallery overview · Claus Wilke, Fundamentals of Data Visualization")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── S13 · CODE：客製化三件套 ─────────
    s = _blank(prs)
    add_title(s, "客製化三件套：標題 / 軸線 / 圖例 / 配色")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="三件套吃掉 80% 客製化需求",
        code=(
            'fig, ax = plt.subplots(figsize=(8, 5))\n'
            '\n'
            '# 1. 標題 + 軸線\n'
            'ax.set_title("月度銷售", fontsize=14, pad=12, loc="left")\n'
            'ax.set_xlabel("月份"); ax.set_ylabel("金額 (千元)")\n'
            'ax.grid(True, linestyle="--", alpha=0.4)   # 輕量網格\n'
            'ax.spines[["top", "right"]].set_visible(False)   # 編輯風格\n'
            '\n'
            '# 2. 圖例（label 要在 plot 時就指定）\n'
            'ax.plot(m, s_2023, label="2023", color="C7")     # 具名色\n'
            'ax.plot(m, s_2024, label="2024", color="#1B5E3F")   # hex\n'
            'ax.legend(loc="upper left", frameon=False)\n'
            '\n'
            '# 3. 配色\n'
            '#   類別色：C0~C9（預設循環） 或 hex\n'
            '#   連續值：cmap="viridis" / "plasma" / "cividis"\n'
            'ax.scatter(x, y, c=z, cmap="viridis")'
        ),
        bullets=[
            "永遠用 ax.set_*()\n不要混 plt.*()",
            "label 要在 plot 時指定\nlegend 才有內容",
            "spines 藏掉 top/right\n= 編輯級乾淨風",
            "類別色用 C0~C9\n連續值才用 cmap",
        ],
        label_dark=True,
    )
    add_source(s, "Matplotlib docs §Customizing · §Choosing colormaps")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── S14 · CODE：中文字型 + savefig ─────────
    s = _blank(prs)
    add_title(s, "中文字型 + savefig：把圖變成可交付成品")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="CJK 新手的第一個坑 — 負號變方塊、中文顯示 tofu",
        code=(
            'import matplotlib.pyplot as plt\n'
            '\n'
            '# ── 放在專案啟動處（一次設定、全專案受益）\n'
            'plt.rcParams["font.family"] = "Noto Sans CJK TC"\n'
            'plt.rcParams["axes.unicode_minus"] = False   # 負號修正\n'
            '\n'
            '# Windows 可改用 Microsoft JhengHei\n'
            '# Mac 可改用 PingFang TC / Heiti TC\n'
            '\n'
            'fig, ax = plt.subplots(figsize=(8, 5))\n'
            'ax.plot([-1, 0, 1], [3, -2, 5])\n'
            'ax.set_title("負值也不再變方塊")\n'
            '\n'
            '# ── 存檔（交付標配）\n'
            'fig.savefig("out.png", dpi=300, bbox_inches="tight")\n'
            'fig.savefig("out.pdf", bbox_inches="tight")   # 向量\n'
            'fig.savefig("out.svg", bbox_inches="tight")   # 向量'
        ),
        bullets=[
            "axes.unicode_minus=False\n擋掉負號方塊",
            "字型一次設定、全圖受益",
            "dpi=300 + bbox='tight'\n= 交付預設",
            "PDF/SVG 同指令產生向量檔",
        ],
        label_dark=True,
    )
    add_source(s, "Matplotlib FAQ §Working with text · §Saving figures · Google Noto Fonts")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── S15 · IMAGE + CODE：subplots 2×2 ─────────
    s = _blank(prs)
    add_title(s, "subplots(2, 2)：一次看四張，EDA 儀表板")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.5), h=Inches(5.0),
        slot_name="subplots 2×2 EDA 儀表板",
        description=(
            "一頁四圖：\n"
            "(左上) hist - 分布\n"
            "(右上) plot - 趨勢\n"
            "(左下) scatter - 相關\n"
            "(右下) box - 離群"
        ),
        url_hint="",
        placeholder_id="Ch09_S15_subplots_2x2",
        registry=image_registry,
        size_hint="1600×1000 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.4), y=Inches(1.3),
        w=Inches(6.3), h=Inches(5.0),
        label="2×2 = EDA 儀表板，一次回答四個問題",
        code=(
            'fig, axes = plt.subplots(2, 2, figsize=(10, 8))\n'
            '\n'
            '# axes 是 2D numpy array\n'
            'axes[0, 0].hist(df["score"], bins=20)\n'
            'axes[0, 0].set_title("分布")\n'
            '\n'
            'axes[0, 1].plot(df["month"], df["sales"])\n'
            'axes[0, 1].set_title("趨勢")\n'
            '\n'
            'axes[1, 0].scatter(df["h"], df["w"], alpha=0.6)\n'
            'axes[1, 0].set_title("相關")\n'
            '\n'
            'cities = ["Taipei", "Taichung", "Kaohsiung"]\n'
            'data = [df[df.city == c]["price"] for c in cities]\n'
            'axes[1, 1].boxplot(data, labels=cities)\n'
            'axes[1, 1].set_title("離群")\n'
            '\n'
            'fig.suptitle("EDA 儀表板", fontsize=16)\n'
            'fig.tight_layout()'
        ),
        bullets=[
            "axes 用 [i, j] 取子圖",
            "tight_layout 解決\n標題 / 軸重疊",
            "每子圖可獨立 set_title",
            "真實 EDA 幾乎\n都以 2×2 開場",
        ],
        label_dark=True,
    )
    add_source(s, "Matplotlib docs §Arranging multiple Axes in a Figure · tight_layout guide")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── S16 · FLOW：EDA 四圖工作流 ─────────
    s = _blank(prs)
    draw_flow_chain(
        s,
        nodes=[
            {"label": "分布", "sub": "hist", "caption": "每欄先看形狀",
             "highlight": True},
            {"label": "趨勢", "sub": "plot", "caption": "時間欄畫 plot"},
            {"label": "相關", "sub": "scatter", "caption": "兩欄配對 scatter",
             "highlight": True},
            {"label": "離群", "sub": "box", "caption": "分組 box 找異常"},
        ],
        title="EDA 四圖工作流：分布 → 趨勢 → 相關 → 離群",
        y=2.8,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.6),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
        "四個問題 → 四張圖 → 一張 2×2 subplots 封裝。\n"
        "這不是教條，是拿到新資料 30 分鐘內一定跑完的例行公事。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER, line_spacing=1.4,
    )
    add_source(s, "Ch09 module synthesis · Tukey EDA workflow")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── S17 · PYRAMID：收束 ─────────
    s = _blank(prs)
    draw_thesis_hierarchy(
        s,
        blocks=[
            {"heading": "四圖速查（一輩子夠用）",
             "items": [
                 "分布 → ax.hist(x, bins=20)",
                 "趨勢 → ax.plot(x, y, label=...)",
                 "相關 → ax.scatter(x, y, alpha=0.6)",
                 "離群 → ax.boxplot(data, labels=...)",
             ]},
            {"heading": "今天該帶走的四條紀律",
             "items": [
                 "先建 fig, ax — 再畫",
                 "用 OO API（ax.set_*）不要混 pyplot",
                 "中文字型 + unicode_minus 專案啟動時設定",
                 "savefig dpi=300 + bbox_inches='tight' 是交付預設",
             ]},
        ],
        title="Ch09 收束：四圖 + OO 心智模型",
        thesis="Ch10 把 Ch08 清洗與 Ch09 視覺化封裝進 DataCleaner — 端到端的資料工程物件。",
    )
    add_source(s, "Ch09 module synthesis · 銜接 Ch10 OOP × Pandas 整合")
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
