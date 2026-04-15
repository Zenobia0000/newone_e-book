"""Ch10 deck — OOP × Pandas 整合實戰
17 content slides + cover + copyright page (M4 · 1.5 hr · 最終收斂章).

Governing thought:
    十章到這裡，要做的不是再學一個新 API，
    是把前九章每一項能力擰成一條可執行、可擴充、可交付的資料管線 ——
    DataCleaner 是這門課最後的「一個類別」，也是你下一份工作的第一個 commit。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_rect, add_textbox, add_title, add_source,
    set_solid_fill, set_no_line, set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_vs_two_col, draw_code_panel, draw_thesis_hierarchy,
    draw_inverted_thesis_box, draw_image_placeholder,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch10"
MODULE_TITLE = "OOP × Pandas 整合實戰"
MODULE_SUBTITLE = "把九章學過的能力，組成一台會跑的車"
TIME_MIN = 90
N_CONTENT = 17


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_ch10(output_path, image_registry=None):
    """Build Ch10 deck — course-closing synthesis chapter.

    image_registry: optional list collecting image placeholder metadata
                    for downstream YAML emission.
    """
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── S1 · SILENT ─────────
    s = _blank(prs)
    draw_silent_page(s, "這一章不教新語法。\n這一章把你學過的九章，組成一台車。")
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ───────── S2 · ASK ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "如果今天就被要求交出一個『可以跑』的資料管線，\n"
        "你會從 import pandas 開始，還是從 class 開始？",
        data_card={
            "label": "2025 業界 Code Review 觀察",
            "stat": "2 種",
            "caption": "『腳本型』30 行搞定，『管線型』可重用可擴充\n"
                       "— 新人與資深的分水嶺就這一題",
        },
    )
    add_source(s, "內部 code review 聚合 · 依 Linus 實用主義原則歸納")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── S3 · MATRIX 2×2：狀態 vs 動作 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "狀態（Instance Attribute）",
         "sub": "self.path / self.df\n物件生命週期中一直在的東西\n資料的「所有權」在物件手上",
         "highlight": True},
        {"text": "動作（Method）",
         "sub": "validate / clean / apply / eda / export\n做完就回傳 self\n命名用動詞而非名詞",
         "highlight": True},
        {"text": "可鏈式（return self）",
         "sub": "每個動作回傳自己\n呼叫方串成一條 pipeline\n比流程圖還清楚"},
        {"text": "應 raise（自訂 Exception）",
         "sub": "錯資料不吞、不 print\n走 raise 語意\n讓上游精準捕捉"},
    ], title="狀態 vs 動作：DataCleaner 的職責拆分")
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "OOP 第一題永遠不是「要繼承誰」— 是「這個物件有什麼狀態、會做什麼動作」。",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "OOP SOLID §單一職責原則 · Ch4 類別設計回顧")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───────── S4 · IMAGE + CODE：DataCleaner 介面草圖 ─────────
    s = _blank(prs)
    add_title(s, "DataCleaner 介面草圖：六個方法、一條鏈")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.5), h=Inches(5.0),
        slot_name="DataCleaner 類別結構圖",
        description=(
            "中央 class DataCleaner 方塊，列出 6 個方法。\n"
            "左側輸入 raw.csv，右上輸出 eda.png，右下輸出 clean.csv。\n"
            "下方灰虛線標『狀態 vs 動作』兩群"
        ),
        url_hint="",
        placeholder_id="Ch10_S04_datacleaner_flow",
        registry=image_registry,
        size_hint="1400×800 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.4), y=Inches(1.3),
        w=Inches(6.3), h=Inches(5.0),
        label="寫 code 之前先寫介面 — 資深工程師的肌肉記憶",
        code=(
            'class DataCleaner:\n'
            '    def __init__(self, data_path: str):\n'
            '        ...\n'
            '    def validate(self) -> "DataCleaner":\n'
            '        ...\n'
            '    def clean_missing_values(\n'
            '        self, strategy: str = "drop"\n'
            '    ) -> "DataCleaner":\n'
            '        ...\n'
            '    def apply_custom_transform(\n'
            '        self, column: str, func\n'
            '    ) -> "DataCleaner":\n'
            '        ...\n'
            '    def generate_eda_report(\n'
            '        self, out_path: str\n'
            '    ) -> "DataCleaner":\n'
            '        ...\n'
            '    def export_data(\n'
            '        self, out_path: str\n'
            '    ) -> "DataCleaner":\n'
            '        ...'
        ),
        bullets=[
            "每個方法回傳型別\n都是 'DataCleaner'",
            "看型別就知道能接",
            "命名用動詞不用名詞",
            "型別註解 = 未來的自己\n讀得懂的文件",
        ],
        label_dark=True,
    )
    add_source(s, "PEP 484 §Type hints · Ch4 類別設計 + Ch5 Method Chaining")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── S5 · CODE：自訂 Exception ─────────
    s = _blank(prs)
    add_title(s, "自訂 Exception：DataValidationError（呼應 Ch6）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="一個專案的第一個 class，常常就是它的自訂 Exception",
        code=(
            'class DataValidationError(Exception):\n'
            '    """資料驗證失敗時拋出"""\n'
            '    def __init__(self, column: str, reason: str):\n'
            '        self.column = column\n'
            '        self.reason = reason\n'
            '        super().__init__(f"[{column}] {reason}")\n'
            '\n'
            '# 使用方：上游精準攔截\n'
            'try:\n'
            '    cleaner.validate()\n'
            'except DataValidationError as e:\n'
            '    logger.error(f"欄位 {e.column} 驗證失敗：{e.reason}")\n'
            '    alert_ops(column=e.column)   # 結構化資料，不用解析字串\n'
            '    raise'
        ),
        bullets=[
            "不要 raise Exception 的泛型",
            "不要 print 再 return —\n訊息會被吞",
            "Exception 帶結構化資料\n（column / reason）",
            "上游用 except XxxError\n精準攔截，不誤殺",
        ],
        label_dark=True,
    )
    add_source(s, "Python Docs §Built-in Exceptions §User-defined · Ch6 自訂例外回顧")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── S6 · SILENT ─────────
    s = _blank(prs)
    draw_silent_page(
        s,
        "好介面的特徵：\n看名字就知道能做什麼、看型別就知道能怎麼串。"
    )
    add_footer(s, MODULE_CODE, 6, N_CONTENT, dark_bg=True)

    # ───────── S7 · CODE：__init__ ─────────
    s = _blank(prs)
    add_title(s, "__init__：fail fast + pathlib（呼應 Ch4 + Ch6）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="建構子要麼給你一個合法物件，要麼炸 — 不要給半成品",
        code=(
            'import pandas as pd\n'
            'from pathlib import Path\n'
            '\n'
            'class DataCleaner:\n'
            '    def __init__(self, data_path: str):\n'
            '        self.path = Path(data_path)\n'
            '        if not self.path.exists():\n'
            '            raise FileNotFoundError(\n'
            '                f"找不到檔案：{self.path}"\n'
            '            )\n'
            '        self.df = pd.read_csv(self.path)\n'
            '        print(f"已載入 {self.path.name}，shape={self.df.shape}")'
        ),
        bullets=[
            "fail fast：錯就早炸",
            "pathlib 取代 os.path —\n物件導向的路徑",
            "self.df 是狀態 —\n物件的本體",
            "建構子直接讀檔，\n不做懶載入（除非必要）",
        ],
        label_dark=True,
    )
    add_source(s, "Python Docs §pathlib · Ch4 __init__ + Ch6 I/O 錯誤處理")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── S8 · CODE：validate ─────────
    s = _blank(prs)
    add_title(s, "validate：自訂 Exception 在真實管線長這樣（呼應 Ch6）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="驗證是一個動作，不是副作用 — 過了就 return self，沒過就 raise",
        code=(
            'def validate(self) -> "DataCleaner":\n'
            '    if self.df.empty:\n'
            '        raise DataValidationError("df", "資料為空")\n'
            '\n'
            '    if "revenue" in self.df.columns and \\\n'
            '            (self.df["revenue"] < 0).any():\n'
            '        raise DataValidationError("revenue", "出現負值")\n'
            '\n'
            '    # 這裡可繼續加 schema check / dtype check / 值域 check\n'
            '    return self    # ← Method Chaining 的入場券'
        ),
        bullets=[
            "驗證 = 檢查 + 決策",
            "絕對不能偷偷改資料",
            "raise 帶欄名 + 理由 =\n未來 debug 的救命繩",
            "return self →\n下一個動作可直接 .clean()",
        ],
        label_dark=True,
    )
    add_source(s, "Ch6 Exception 語意 · Fail-fast 設計原則")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── S9 · CODE：clean_missing_values ─────────
    s = _blank(prs)
    add_title(s, "clean_missing_values：dropna / fillna 策略參數化（呼應 Ch8）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="用字串派發策略 — 窮人版 Strategy Pattern，但通常就夠用",
        code=(
            'def clean_missing_values(\n'
            '    self, strategy: str = "drop"\n'
            ') -> "DataCleaner":\n'
            '    if strategy == "drop":\n'
            '        self.df = self.df.dropna()\n'
            '    elif strategy == "mean":\n'
            '        self.df = self.df.fillna(\n'
            '            self.df.mean(numeric_only=True)\n'
            '        )\n'
            '    # elif strategy == "median": ...\n'
            '    # elif strategy == "ffill": ...\n'
            '    else:\n'
            '        raise DataValidationError(\n'
            '            "strategy", f"未知策略 {strategy!r}"\n'
            '        )\n'
            '    return self'
        ),
        bullets=[
            "策略以 str 參數切換",
            "比 if-tree 巢狀乾淨",
            "Ch8 四種策略都可擴進來\n（drop / mean / ffill / interp）",
            "真的複雜了 —\n再考慮 Strategy Pattern",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Missing data · Ch8 四種策略回顧")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── S10 · CODE：apply_custom_transform ─────────
    s = _blank(prs)
    add_title(s, "apply_custom_transform：lambda 找到主場（呼應 Ch3 + Ch8）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="讓使用者傳任意函式就能擴充 — 這是 Ch3 first-class function 的真實主場",
        code=(
            'def apply_custom_transform(\n'
            '    self, column: str, func\n'
            ') -> "DataCleaner":\n'
            '    if column not in self.df.columns:\n'
            '        raise DataValidationError(\n'
            '            column, "欄位不存在"\n'
            '        )\n'
            '    self.df[f"{column}_transformed"] = (\n'
            '        self.df[column].apply(func)\n'
            '    )\n'
            '    return self\n'
            '\n'
            '# 呼叫方：lambda 在這裡終於有家\n'
            'cleaner.apply_custom_transform(\n'
            '    "revenue", lambda x: x * 1.05\n'
            ')'
        ),
        bullets=[
            "func 是 first-class\ncitizen（Ch3）",
            "lambda 在 pandas 找到主場\n（Ch8 呼應）",
            "新增欄位而非覆寫 —\n可追溯、可 diff",
            "失敗用自訂 Exception —\n不讓 KeyError 漏上來",
        ],
        label_dark=True,
    )
    add_source(s, "Ch3 lambda / Ch8 apply 設計回顧 · 高階函式模式")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── S11 · IMAGE + CODE：generate_eda_report ─────────
    s = _blank(prs)
    add_title(s, "generate_eda_report：把 Matplotlib 包進類別（呼應 Ch9）")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.5), h=Inches(5.0),
        slot_name="Method Chaining 視覺化",
        description=(
            "六個圓角矩形水平連鏈：\n"
            "__init__ → validate → clean → apply → eda → export\n"
            "每個方塊下方標『return self』\n"
            "鏈末端標註『一條 pipeline』"
        ),
        url_hint="",
        placeholder_id="Ch10_S11_method_chaining",
        registry=image_registry,
        size_hint="1600×600 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.4), y=Inches(1.3),
        w=Inches(6.3), h=Inches(5.0),
        label="EDA 變成 pipeline 的一環 — 呼叫方一行就拿到報表",
        code=(
            'import matplotlib.pyplot as plt\n'
            '\n'
            'def generate_eda_report(\n'
            '    self, out_path: str = "eda.png"\n'
            ') -> "DataCleaner":\n'
            '    numeric = self.df.select_dtypes("number")\n'
            '    fig, axes = plt.subplots(2, 2,\n'
            '                             figsize=(10, 8))\n'
            '    for ax, col in zip(axes.flat,\n'
            '                       numeric.columns[:4]):\n'
            '        ax.hist(numeric[col].dropna(), bins=30)\n'
            '        ax.set_title(col)\n'
            '    fig.tight_layout()\n'
            '    fig.savefig(out_path, dpi=150)\n'
            '    plt.close(fig)   # 關圖是禮貌\n'
            '    return self'
        ),
        bullets=[
            "Ch9 的 2×2 subplot\n骨架直接拿來用",
            "plt.close 防止\n記憶體洩漏",
            "dpi=150：出版級\n品質下限",
            "return self —\n下游繼續串 export",
        ],
        label_dark=True,
    )
    add_source(s, "Matplotlib User Guide §Figure · Ch9 EDA 工作流回顧")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── S12 · CODE：export_data + 完整鏈 ─────────
    s = _blank(prs)
    add_title(s, "export_data + 完整 Method Chaining（呼應 Ch5）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="這段 code 讓你明白為何 Ch5 花時間講 return self — 現在回報你了",
        code=(
            'def export_data(self, out_path: str) -> "DataCleaner":\n'
            '    self.df.to_csv(out_path, index=False)\n'
            '    print(f"匯出 {len(self.df)} 列到 {out_path}")\n'
            '    return self\n'
            '\n'
            '# ── 完整 pipeline，一段 code 講清整個故事 ──\n'
            '(DataCleaner("raw.csv")\n'
            '    .validate()\n'
            '    .clean_missing_values(strategy="mean")\n'
            '    .apply_custom_transform("revenue",\n'
            '                            lambda x: x * 1.05)\n'
            '    .generate_eda_report("report.png")\n'
            '    .export_data("clean.csv"))'
        ),
        bullets=[
            "外層括號包住整條鏈",
            "一個動作一行 —\n每行讀起來像動詞短句",
            "比流程圖還清楚的文件",
            "這就是 Ch5 return self\n+ 魔術方法的實踐成果",
        ],
        label_dark=True,
    )
    add_source(s, "Ch5 Method Chaining 設計回顧 · Fluent Interface pattern")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── S13 · TABLE：可改進方向 ─────────
    s = _blank(prs)
    draw_editorial_table(
        s,
        header=["改進方向", "為什麼", "下一步 API / 關鍵字"],
        rows=[
            ["加入 logging",
             "print 不能關、logging 有等級",
             "import logging / getLogger(__name__)"],
            ["從 config 載入",
             "參數寫死就不可重用",
             "@classmethod from_config(cls, path)"],
            ["單元測試",
             "return self 很好測，tmp_path 隔離副作用",
             "pytest + tmp_path fixture"],
            ["切出 Strategy Pattern",
             "策略太多時 if-tree 會爆炸",
             "abstract class CleanStrategy"],
            ["資料型別自動推斷",
             "pandas 讀 CSV 常把日期當字串",
             "pd.read_csv(parse_dates=...) / dtype="],
        ],
        title="可改進方向：這就是期末作業的加分清單",
        col_widths=[1.3, 2.3, 2.4],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "做兩個以上交出來 — 你的 GitHub 就有一個亮點；五個全做，你就是資深候選人。",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "期末作業指引 · 作業 rubric 對應 §四")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── S14 · PYRAMID 能力盤點（用 MATRIX 2x3 + thesis 呈現五層） ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=3, cells=[
        {"text": "1. 系統直覺",
         "sub": "OS / RAM / I/O 三關鍵字\n看 code 知道在跟誰要資源\n來源：Ch1"},
        {"text": "2. Python 進階",
         "sub": "Generator / Lambda /\nComprehension / 自訂函式\n來源：Ch2–Ch3"},
        {"text": "3. OOP",
         "sub": "類別、繼承、魔術方法\nMethod Chaining\n來源：Ch4–Ch5",
         "highlight": True},
        {"text": "4. 資料工程",
         "sub": "I/O、Exception\nNumPy 向量化、Pandas 七步\n來源：Ch6–Ch8",
         "highlight": True},
        {"text": "5. 視覺化與整合",
         "sub": "Matplotlib EDA 工作流\nOOP × Pandas 封裝\n來源：Ch9–Ch10",
         "highlight": True},
        {"text": "堆疊而非平行",
         "sub": "底層越穩、上層越能長\n今天你能寫 DataCleaner\n是因為前四層都穩了"},
    ], title="18 小時能力盤點：五層能力堆疊（由下而上）")
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "少一層，上面就不穩 — 這是為什麼我們花時間在 Ch1 講 OS 與 I/O。",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "課程架構總覽 · Ch1–Ch10 能力地圖合成")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── S15 · IMAGE + TABLE：五條學習路徑 ─────────
    s = _blank(prs)
    add_title(s, "下一站：五條路線（SQL / sklearn / Polars+DuckDB / Airflow / 雲端）")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=Inches(5.2), h=Inches(5.0),
        slot_name="五條學習路徑生態圖",
        description=(
            "中央『你現在的位置（DataCleaner）』\n"
            "五條放射狀箭頭指向五個目的地\n"
            "SQL / sklearn / Polars+DuckDB /\nAirflow / 雲端"
        ),
        url_hint="",
        placeholder_id="Ch10_S15_learning_ecosystem",
        registry=image_registry,
        size_hint="1600×1000 px",
    )
    # Right-side: compact table
    tx = Inches(6.1)
    tw = T.SLIDE_W - tx - T.MARGIN_X
    hdr_h = Inches(0.4)
    row_h = Inches(0.78)
    # header band
    hdr = add_rect(s, tx, Inches(1.2), tw, hdr_h)
    set_solid_fill(hdr, T.PRIMARY)
    set_no_line(hdr)
    add_textbox(
        s, tx + Inches(0.15), Inches(1.2), tw - Inches(0.3), hdr_h,
        "方向　|　推薦起點　|　一句話",
        font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    paths = [
        ("資料庫整合", "SQLAlchemy", "把 CSV 換成資料庫"),
        ("機器學習", "scikit-learn", "把清洗好的資料餵給模型"),
        ("大規模資料", "Polars / DuckDB", "當 Pandas 變慢時的下一站"),
        ("自動化排程", "Airflow / Prefect / n8n", "讓你的管線每天自動跑"),
        ("雲端部署", "AWS S3+Lambda / GCP BQ", "把腳本搬上雲"),
    ]
    cy = Inches(1.2) + hdr_h
    for i, (d, api, one) in enumerate(paths):
        if i % 2 == 1:
            band = add_rect(s, tx, cy, tw, row_h)
            set_solid_fill(band, T.TABLE_ALT)
            set_no_line(band)
        add_textbox(
            s, tx + Inches(0.15), cy + Inches(0.05),
            tw - Inches(0.3), Inches(0.3),
            f"{d}　·　{api}",
            font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
        )
        add_textbox(
            s, tx + Inches(0.15), cy + Inches(0.38),
            tw - Inches(0.3), Inches(0.35),
            one,
            font_size=T.FONT_SMALL, color=T.CHARCOAL,
        )
        cy += row_h
    add_source(s, "業界資料工程學習路徑地圖 · 2025 版 · 依 Ch10 收斂思路整理")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── S16 · VS-LIST：Linus 三句忠告 ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="常見迷思（別這樣想）",
        right_title="Linus 式紀律（這樣才長久）",
        left_items=[
            "「跑得起來就好」",
            "「先寫再說、能 work 就優化」",
            "「什麼都 OOP 最優雅」",
        ],
        right_items=[
            "跑得起來不算數 — 能讓三個月後的自己看懂才算數",
            "最強的優化是「不要做」— 在向量化之前先問能不能跳過這步",
            "OOP 不是萬靈丹 — 30 行的腳本不需要類別",
        ],
        title="Linus 三句結業忠告：給三個月後的自己",
        summary="寫 code 是寫給未來的自己 — 好品味先於聰明。",
    )
    add_source(s, "Linus Torvalds 公開訪談合輯 · 實用主義原則整理")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── S17 · SILENT 結業 ─────────
    s = _blank(prs)
    draw_silent_page(
        s,
        "18 小時結束了。\n從現在起，程式碼是寫給三個月後的自己看的。"
    )
    add_footer(s, MODULE_CODE, 17, N_CONTENT, dark_bg=True)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
