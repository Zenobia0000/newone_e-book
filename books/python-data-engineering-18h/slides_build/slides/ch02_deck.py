"""Ch02 deck — Python 核心快速複習與資料結構深化.

17 content slides + cover + copyright page.
Composed strictly from primitives in slides_build/primitives.py and
branding.py, palette restricted to 黑 / 灰 / 深綠 (#1B5E3F).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_rect, add_textbox, add_title, add_source,
    set_solid_fill, set_no_line, set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_vs_two_col, draw_inverted_thesis_box,
    draw_code_panel, draw_image_placeholder, draw_emphasis_pill,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch02"
MODULE_TITLE = "Python 核心快速複習與資料結構深化"
MODULE_SUBTITLE = "M1 系統前導與 Python 機制 · 第二章"
TIME_MIN = 90
N_CONTENT = 17


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_ch02(output_path, image_registry=None):
    """Build Ch02 deck.

    image_registry: optional list collecting image placeholder metadata.
    """
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # S1 — SILENT hook
    s = _blank(prs)
    draw_silent_page(
        s,
        "10GB 檔案不會讓你的 RAM 爆炸——\n只要你選對迭代方式。",
    )
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # S2 — ASK: naive vs generator
    s = _blank(prs)
    draw_ask_page(
        s,
        "你要一次讀完 10GB log 找出所有 ERROR 行，\n需要多少 RAM？",
        data_card={
            "label": "Naive vs Generator",
            "stat": "~200×",
            "caption": "naive readlines() 吃 10GB+；generator 逐行 < 50MB",
        },
    )
    add_source(s, "本課程實測 Python 3.11 / 16GB 機器")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # S3 — MATRIX 2x3: 型別快速校準
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=3, cells=[
        {"text": "數字型 int / float",
         "sub": "整數、浮點數\n算術、計數、度量"},
        {"text": "文字 str",
         "sub": "不可變字串\n以 index 存取，以 slice 切片"},
        {"text": "布林與空值 bool / None",
         "sub": "True / False / None\n控制流與「無資料」訊號"},
        {"text": "序列容器 list / tuple",
         "sub": "依序排列\nlist 可變、tuple 不可變"},
        {"text": "鍵值容器 dict",
         "sub": "key → value\nO(1) 平均查找"},
        {"text": "純值 = 一個事實；容器 = 一群事實的結構。",
         "highlight": True},
    ], title="Python 六種內建型別：純值負責表達，容器負責組織")
    add_source(s, "Python 官方文件 docs.python.org/3/library/stdtypes.html")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # S4 — FLOW: name → object → value
    s = _blank(prs)
    draw_flow_chain(s, nodes=[
        {"label": "name", "caption": "便利貼（變數名）"},
        {"label": "object", "caption": "記憶體物件", "highlight": True},
        {"label": "value", "caption": "物件內的實際資料"},
    ], title="變數不是盒子，是便利貼——Python 的三層心智模型",
       y=3.0)
    draw_inverted_thesis_box(
        s,
        "a = b 不是複製，是貼第二張便利貼到同一個東西上。",
        y=5.9, width=10.0,
    )
    add_source(s, "CPython data model 官方說明")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # S5 — TABLE: 四容器選用對照
    s = _blank(prs)
    draw_editorial_table(s,
        header=["容器", "何時用", "關鍵操作", "時間複雜度"],
        rows=[
            ["List",  "順序資料、允許重複",
             "append / index / slice",
             "O(1) append · O(n) search"],
            ["Tuple", "固定欄位、不可變",
             "解構、可當 dict key",
             "同 List，但不可變"],
            ["Dict",  "key-value 查找",
             "get / items / in",
             "O(1) 平均查找"],
            ["Set",   "去重、集合運算",
             "add / & | - / in",
             "O(1) 平均查找"],
        ],
        col_widths=[0.9, 1.6, 1.6, 1.5],
        title="四種容器，四個不同任務——選錯會付兩種代價：效能與正確性",
    )
    add_textbox(s, T.MARGIN_X, Inches(5.5),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
                "要順序選 List；要唯一選 Set；要查找選 Dict；要保證不被改選 Tuple。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "Python Time Complexity wiki · Sunny Data Science 2026 整理")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # S6 — MATRIX 2x2: 可變/不可變 × 能否當 key
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "可變 · 不能當 key",
         "sub": "list · dict · set\n內容會變，hash 會失效"},
        {"text": "不可變 · 能當 key",
         "sub": "int · str · tuple · frozenset\nhashable，適合作索引",
         "highlight": True},
        {"text": "可變 · 只當值",
         "sub": "list / dict / set 放在 value 位置\n最常見用法"},
        {"text": "不可變 · 只當值",
         "sub": "int · str 放在 value 位置\n如計數器、標籤"},
    ], title="hashable = 不可變 = 能當 dict/set 的 key")
    add_textbox(s, T.MARGIN_X, Inches(6.4),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
                "橫軸：可變性 可變 ↔ 不可變   ·   縱軸：用途 作 key ↔ 作 value",
                font_size=T.FONT_SOURCE, color=T.GRAY_MID,
                align=PP_ALIGN.CENTER)
    add_source(s, "Python Language Reference §3.1 Objects, values and types")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # S7 — CODE: [[]] * 3 踩坑
    s = _blank(prs)
    add_title(s, "`[[]] * 3` 不是建三個 list——是一個 list 被指向三次")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.4),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.4),
        label="陷阱示範 · 經典踩坑",
        code=(
            "# 你以為：建立三個獨立的空 list\n"
            "a = [[]] * 3\n"
            "a[0].append(1)\n"
            "\n"
            "print(a)\n"
            "# 結果：[[1], [1], [1]]   ← 三格全中\n"
            "\n"
            "# 正解：用 list comprehension\n"
            "a = [[] for _ in range(3)]\n"
            "a[0].append(1)\n"
            "print(a)   # [[1], [], []]"
        ),
        bullets=[
            "`*` 複製的是參照，不是 object",
            "三張便利貼貼在同一個 list 上",
            "一改全改是必然結果，不是 bug",
            "延伸：numpy.zeros / pandas 同源陷阱",
            "心法：要 n 個獨立容器，一律用 comprehension",
        ],
    )
    add_source(s, "Python FAQ · 'How do I make a list of lists?'")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # S8 — VS: 淺拷貝 vs 深拷貝
    s = _blank(prs)
    draw_vs_two_col(s,
        left_title="淺拷貝（shallow）",
        right_title="深拷貝（deep）",
        left_items=[
            "`list.copy()` / `copy.copy()` / `x[:]` / `list(x)`",
            "只複製外殼，內層仍共用",
            "巢狀 list 改內層 → 污染原本",
            "適用：扁平結構、成本敏感路徑",
            "踩坑：pandas copy / numpy view 同系列",
        ],
        right_items=[
            "`copy.deepcopy(x)`",
            "遞迴複製到葉子節點",
            "改任何一層都不影響原本",
            "適用：巢狀 / 含 mutable 成員的物件",
            "代價：時間與記憶體多一份",
        ],
        title="淺拷貝複製外殼，深拷貝連骨頭一起——巢狀只有後者安全",
        summary="扁平用淺、巢狀用深；不確定時先 deepcopy 保命，效能問題之後再優化。",
        delta="只差一層",
    )
    add_source(s, "Python 標準函式庫 copy module 文件")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # S9 — MATRIX 2x2: 30 秒容器決策樹
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "Q1：需要順序且允許重複？",
         "sub": "→ 選 List\n典型：時序資料、步驟紀錄"},
        {"text": "Q2：需要唯一性 / 集合運算？",
         "sub": "→ 選 Set\n典型：去重、交集、差集"},
        {"text": "Q3：要用 key 快速查找？",
         "sub": "→ 選 Dict\n典型：ID 對應、設定表"},
        {"text": "Q4：要保證欄位不被改 / 能當 key？",
         "sub": "→ 選 Tuple（或 @dataclass(frozen=True)）\n典型：座標、固定 schema",
         "highlight": True},
    ], title="30 秒容器選用決策：照順序問這四個問題")
    add_textbox(s, T.MARGIN_X, Inches(6.4),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
                "實務分布（本課程 code review 樣本）：List 80% · Dict 15% · Set 4% · Tuple 1%",
                font_size=T.FONT_SOURCE, color=T.GRAY_MID,
                align=PP_ALIGN.CENTER)
    add_source(s, "本課程 2023-2025 企業訓練 code review 樣本 n=120 份")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # S10 — ASK: 打開 Iterator 主題
    s = _blank(prs)
    draw_ask_page(
        s,
        "`for line in file:` 這一行，\nPython 底下到底發生了什麼？",
        data_card={
            "label": "for 的真實身份",
            "stat": "糖",
            "caption": "for 是語法糖，底下是 iter() + next() + StopIteration 的契約",
        },
    )
    add_source(s, "Python Language Reference §6.2.9 Yield expressions")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # S11 — FLOW: Iterator 協議
    s = _blank(prs)
    draw_flow_chain(s, nodes=[
        {"label": "Iterable", "caption": "有 __iter__"},
        {"label": "iter(x)",  "caption": "呼叫 __iter__"},
        {"label": "Iterator", "caption": "有 __next__"},
        {"label": "next(it)", "caption": "取下一個 / StopIteration",
         "highlight": True},
    ], title="Iterator 協議只有兩個方法——這是 Python 迴圈的契約",
       y=3.2)
    draw_inverted_thesis_box(
        s,
        "for 迴圈 = 不斷呼叫 next()，直到收到 StopIteration。",
        y=5.9, width=10.0,
    )
    add_source(s, "PEP 234 · Iterators")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # S12 — CODE: yield 的暫停/續行
    s = _blank(prs)
    add_title(s, "yield 不是 return——它讓函式暫停、記住現場、下次從這裡繼續")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.4),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.4),
        label="yield 示範 · generator function",
        code=(
            "def counter():\n"
            "    print('start')\n"
            "    yield 1          # ← 暫停在這裡\n"
            "    print('resume')\n"
            "    yield 2\n"
            "    print('end')\n"
            "\n"
            "g = counter()        # 此時什麼都沒印\n"
            "next(g)   # start  → 1\n"
            "next(g)   # resume → 2\n"
            "next(g)   # end    → StopIteration"
        ),
        bullets=[
            "return 一去不回；yield 可反覆回來",
            "呼叫 g = counter() 不執行本體",
            "每次 yield 凍住「局部變數 + 指令位置」",
            "記憶體只裝「當下那一幀」",
            "把函式想成連續劇，yield 是本週結尾",
        ],
    )
    add_source(s, "PEP 255 · Simple Generators")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # S13 — VS: List Comp vs Generator Expression
    s = _blank(prs)
    draw_vs_two_col(s,
        left_title="List Comprehension  [ ]",
        right_title="Generator Expression  ( )",
        left_items=[
            "`[x*2 for x in range(10**8)]`",
            "一次算完 1 億個，全部存 RAM",
            "記憶體：~ 3.2 GB",
            "適用：小資料、要多次重讀",
            "特徵：有長度、可索引、可 len()",
        ],
        right_items=[
            "`(x*2 for x in range(10**8))`",
            "只算當下一個，下一個才繼續",
            "記憶體：~ 200 bytes",
            "適用：大資料、只讀一遍",
            "特徵：惰性、只能迭代一次、無 len()",
        ],
        title="括號不同、命運不同：List 立刻算完全部，Generator 要一個算一個",
        summary="要存結果用 [ ]；要串管線用 ( )。同樣的語義，記憶體差兩個數量級。",
        delta="~ 200× RAM",
    )
    add_source(s, "本課程實測 Python 3.11 · tracemalloc 觀測")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # S14 — IMAGE placeholder: Generator 記憶體示意
    s = _blank(prs)
    add_title(s, "Generator 只在 RAM 中留一格——真正的資料還在硬碟上")
    draw_image_placeholder(
        s,
        x=Inches(1.2), y=Inches(1.4),
        w=Inches(10.9), h=Inches(4.2),
        slot_name="Generator 記憶體運作示意",
        description=("左：10GB 檔案（硬碟符號）  →  "
                     "中：generator 輸送帶（僅一行在 RAM + 暫停指標）  →  "
                     "右：for 迴圈消費者"),
        url_hint="原創示意圖；可用 draw.io / Excalidraw 繪製後截圖",
        size_hint="1400×800 px",
        placeholder_id="Ch02_S14_generator_memory",
        registry=image_registry,
    )
    add_textbox(s, T.MARGIN_X, Inches(5.8),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.7),
                "你處理的不是資料本身，是「下一筆資料怎麼來」的指令。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER, line_spacing=1.3)
    add_source(s, "本課程原創示意")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # S15 — CODE: 10GB log → ERROR 行實戰
    s = _blank(prs)
    add_title(s, "10GB log 讀 ERROR 行——三行 generator 解決")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.4),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.4),
        label="實戰 · 大檔 streaming",
        code=(
            "def read_errors(path):\n"
            "    with open(path, encoding='utf-8') as f:\n"
            "        for line in f:              # file 本身就是 iterator\n"
            "            if 'ERROR' in line:\n"
            "                yield line.rstrip()\n"
            "\n"
            "# 使用端：完全不 care 檔案多大\n"
            "for err in read_errors('/var/log/app.log'):\n"
            "    send_to_alert(err)"
        ),
        bullets=[
            "open() 回傳的 file 是天生 iterator",
            "for line in f：O(1) 記憶體逐行",
            "yield 讓呼叫端也「要一個給一個」",
            "組合後 10GB 流過，RAM < 50 MB",
            "要過濾 / 轉換 / 去重，再串下一個 generator",
        ],
    )
    add_source(s, "本課程實戰範例 · Ch06 檔案 I/O 會再深化")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # S16 — MATRIX 2x3: 三個常用 Generator 模式
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=3, cells=[
        {"text": "① 逐行讀",
         "sub": "for line in open(path):\n    yield line"},
        {"text": "② 分塊讀",
         "sub": "while chunk := f.read(1<<20):\n    yield chunk"},
        {"text": "③ 管線串接",
         "sub": "map / filter / 轉換串成\n一條 iterator chain"},
        {"text": "文字 log · JSONL",
         "sub": "逐行流處理\n典型 ETL 入口"},
        {"text": "二進位 · 大 CSV",
         "sub": "固定記憶體窗口\n速度可控"},
        {"text": "惰性運算，耗時只在真正消費時發生。",
         "highlight": True},
    ], title="三個你會反覆使用的 Generator 模式——吃下 80% 大檔情境")
    add_source(s, "本課整理 · Ch06 / Ch07 / Ch08 將反覆呼應此三式")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # S17 — SILENT closing
    s = _blank(prs)
    draw_silent_page(
        s,
        "會選容器，代表你能寫對；\n會用 generator，代表你能寫大。",
    )
    add_footer(s, MODULE_CODE, 17, N_CONTENT, dark_bg=True)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
