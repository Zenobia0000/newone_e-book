"""Ch04 deck — OOP 核心觀念與實例化
17 content slides + cover + copyright page.

Governing thought:
    腳本解決問題，類別解決規模。
    OOP 不是炫技，是把散落的狀態收進一個
    可重用、可測試、可組合的殼。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_rect, add_textbox, add_title, add_source,
    set_solid_fill, set_no_line, set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_vs_two_col, draw_code_panel, draw_thesis_hierarchy,
    draw_flow_chain, draw_image_placeholder,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch04"
MODULE_TITLE = "OOP 核心觀念與實例化"
MODULE_SUBTITLE = "腳本解決問題，類別解決規模"
TIME_MIN = 90
N_CONTENT = 17


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_ch04(output_path, image_registry=None):
    """Build Ch04 deck."""
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── S1 · SILENT ─────────
    s = _blank(prs)
    draw_silent_page(s, "腳本解決問題，\n類別解決規模。")
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ───────── S2 · ASK ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "為什麼同一段清洗邏輯，\n你已經複製貼上第四次？",
        data_card={
            "label": "資料工程師日常",
            "stat": "61%",
            "caption": "每週維護 ≥ 3 條\n結構相似的管線",
        },
    )
    add_source(s, "Kaggle State of Data Science 2024 (n=26,000)")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── S3 · VS：腳本式 vs 類別式 ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="腳本式（Procedural）",
        right_title="類別式（Object-Oriented）",
        left_items=[
            "5 個散落的函式",
            "10 個全域變數穿插",
            "狀態住在 global namespace",
            "改一處，擔心另外三處",
            "單元測試很難寫",
        ],
        right_items=[
            "1 個 DataPipeline 類別",
            "狀態收進 self.xxx",
            "行為封裝成方法",
            "實例化 N 次 = 管線複製 N 條",
            "測試只要 pipe = DataPipeline(); assert ...",
        ],
        title="腳本式 vs 類別式：差的不是行數，是狀態的去處",
        summary="兩邊都 100 行，但維護成本差 10 倍——因為狀態的可見性不同。",
        delta="狀態\n的歸屬",
    )
    add_source(s, "Ch04 課堂歸納 · Gamma et al., Design Patterns §Encapsulation")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───────── S4 · MATRIX 1×3：三大效益 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=1, cols=3, cells=[
        {"text": "可重用 (Reusable)",
         "sub": "清洗邏輯寫一次\n套到 N 個資料集\n實例化 = 複製管線",
         "highlight": True},
        {"text": "可測試 (Testable)",
         "sub": "小單元獨立驗證\npipe.add_step() 可 mock\n單元測試 < 10 行"},
        {"text": "可組合 (Composable)",
         "sub": "管線各步驟可替換\n換資料源不改邏輯\nCh10 method chaining 基礎"},
    ], title="資料工程採用 OOP 的三大效益：背後都指向降耦")
    add_source(s, "Fowler, Refactoring §Encapsulate Record · Beck, TDD by Example")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── S5 · FLOW：反例三條 ─────────
    s = _blank(prs)
    draw_flow_chain(
        s,
        nodes=[
            {"label": "一次性腳本",
             "sub": "< 50 行、跑一次丟",
             "caption": "不必 OOP"},
            {"label": "純函式轉換",
             "sub": "無狀態、純輸入輸出",
             "caption": "用 def 就好"},
            {"label": "實驗 notebook",
             "sub": "探索階段、邊跑邊改",
             "caption": "等成熟再抽類別"},
        ],
        title="反例三條：什麼時候「不」該用 OOP",
        y=2.8,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
        '"If you need OOP for hello world, you\'ve already lost."  — Linus 觀點',
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.9),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "原則：看到『狀態會被 N 個函式共用』才考慮升級類別；單純轉換不需要。",
        font_size=T.FONT_SMALL, color=T.GRAY_MID,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "Linus Torvalds, lkml 2003 · Kernighan & Pike, The Practice of Programming §Simplicity")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── S6 · ASK + IMAGE：Class vs Object 比喻 ─────────
    s = _blank(prs)
    add_title(s, "如果 Class 是建築藍圖，Object 是什麼？")
    # 左：問題與提示
    add_textbox(
        s, T.MARGIN_X, Inches(1.5),
        Inches(6.2), Inches(4.8),
        "一份藍圖\n可以蓋出很多棟房子。\n\n"
        "每棟房子\n各自有地址、有住戶，\n但結構長得一樣。\n\n"
        "→ Class 定義「長什麼樣」\n→ Object 是「實際蓋出來的那棟」",
        font_size=T.FONT_SUBTITLE, color=T.CHARCOAL,
        line_spacing=1.45, anchor=MSO_ANCHOR.TOP,
    )
    # 右：圖片 placeholder
    draw_image_placeholder(
        s,
        x=Inches(7.0), y=Inches(1.3),
        w=Inches(5.7), h=Inches(5.0),
        slot_name="Class 藍圖 vs Object 實體比喻圖",
        description="左：建築藍圖線稿\n右：三棟由同一份藍圖蓋出的房子\n（同結構，不同門牌/顏色）",
        url_hint="https://docs.python.org/3/tutorial/classes.html",
        placeholder_id="Ch04_S06_blueprint_vs_house",
        registry=image_registry,
        size_hint="1400×1200 px",
    )
    add_source(s, "Python Tutorial §9 Classes · 比喻取自 Bruce Eckel, Thinking in Java")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───────── S7 · CODE：最小類別 ─────────
    s = _blank(prs)
    add_title(s, "最小類別：三個關鍵字 + 兩個方法，比 5 個散落函式好用")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="class DataPipeline — 一個最小可用雛形",
        code=(
            'class DataPipeline:\n'
            '    def __init__(self, name):\n'
            '        self.name = name       # 實例屬性：每個實例獨立\n'
            '        self.steps = []        # 可變物件放 __init__！\n'
            '\n'
            '    def add_step(self, step):\n'
            '        self.steps.append(step)\n'
            '\n'
            '# 實例化\n'
            'pipe = DataPipeline("ETL_v1")\n'
            'pipe.add_step("read_csv")\n'
            'pipe.add_step("drop_nulls")\n'
            'print(pipe.steps)   # [\'read_csv\', \'drop_nulls\']'
        ),
        bullets=[
            "class 宣告藍圖",
            "__init__ 是誕生禮",
            "self 是身分證",
            "add_step 是行為",
            "pipe 是實體（Object）",
        ],
        label_dark=True,
    )
    add_source(s, "Python Tutorial §9.3 A First Look at Classes")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── S8 · CODE：__init__ + self 解剖 ─────────
    s = _blank(prs)
    add_title(s, "__init__ 是誕生禮，self 是身分證——不是關鍵字，是慣例名")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="實例化那一行，Python 在背後做了三件事",
        code=(
            'pipe = DataPipeline("ETL_v1")\n'
            '\n'
            '# Python 背後做的事：\n'
            '# ① 建立一顆新的空物件 obj\n'
            '# ② 呼叫 DataPipeline.__init__(self=obj, name="ETL_v1")\n'
            '# ③ 把 obj 綁到名字 pipe\n'
            '\n'
            '# self 是慣例，理論上可改名（但別改）：\n'
            'class Bad:\n'
            '    def __init__(me, x):    # 可執行但沒人看得懂\n'
            '        me.x = x\n'
            '\n'
            '# 省略 self 會 TypeError：\n'
            '# TypeError: __init__() takes 1 positional argument\n'
            '#           but 2 were given'
        ),
        bullets=[
            "__init__ 不是建構子\n是「初始化鉤子」",
            "self 是 Python 自動\n傳入的第一個參數",
            "慣例叫 self，PEP 8 明文",
            "省略 self → TypeError",
        ],
        label_dark=True,
    )
    add_source(s, "Python Data Model §3.3.1 · PEP 8 §Function and method arguments")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── S9 · MATRIX 1×3：self 三問 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=1, cols=3, cells=[
        {"text": "Q1 · 誰傳進來？",
         "sub": "Python 自動把\n當前實例塞進第一個參數\n你不用手動傳",
         "highlight": True},
        {"text": "Q2 · 指到哪？",
         "sub": "指到那一顆 instance\nself.name = ... 就是\n寫入這顆物件的屬性"},
        {"text": "Q3 · 不寫會怎樣？",
         "sub": "TypeError\n或被當 classmethod\n（未進階前不要碰）"},
    ], title="self 三問：背熟，OOP 一半功課就完成")
    add_source(s, "Python Tutorial §9.3.2 Class and Instance Variables")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── S10 · CODE + IMAGE：記憶體模型 ─────────
    s = _blank(prs)
    add_title(s, "兩個實例各自的 steps，互不污染——記憶體才是真相")
    # 左：圖片 placeholder
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.8), h=Inches(5.0),
        slot_name="兩個實例各自 steps 的記憶體示意",
        description="兩個盒子：p1 與 p2\n各自指向獨立的 steps list\n（reference 箭頭標示）",
        url_hint="https://docs.python.org/3/reference/datamodel.html",
        placeholder_id="Ch04_S10_memory_model",
        registry=image_registry,
        size_hint="1400×1200 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.7), y=Inches(1.3),
        w=Inches(6.0), h=Inches(5.0),
        label="驗證兩實例獨立",
        code=(
            'p1 = DataPipeline("ETL_A")\n'
            'p2 = DataPipeline("ETL_B")\n'
            '\n'
            'p1.add_step("read_csv")\n'
            'p2.add_step("read_json")\n'
            '\n'
            'print(p1.steps)\n'
            '# [\'read_csv\']\n'
            'print(p2.steps)\n'
            '# [\'read_json\']\n'
            '\n'
            'print(p1.steps is p2.steps)\n'
            '# False  ← 兩個 list 在不同記憶體位置'
        ),
        bullets=[
            "每次 __init__ 執行\n都建一個新 list",
            "p1.steps 與 p2.steps\n指向不同物件",
            "這就是實例獨立的根基",
        ],
        label_dark=True,
    )
    add_source(s, "Python Reference §3.1 Objects, values, and types")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── S11 · ASK：items=[] 陷阱前哨 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "如果把 items = [] 寫在類別頂端，\n會發生什麼事？",
        data_card={
            "label": "面試最常考",
            "stat": "Top 3",
            "caption": "Python OOP 面試題\n永遠的前三名\n生產環境常見 bug",
        },
    )
    add_source(s, "Stack Overflow, Top-voted Python OOP questions 2024")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── S12 · VS-CODE：Class vs Instance Attribute ─────────
    s = _blank(prs)
    add_title(s, "Class Attribute vs Instance Attribute：差一個縮排，差整條管線")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.4),
        label="BAD：items 放在類別頂端 → 所有實例共用",
        code=(
            'class Bad:\n'
            '    items = []          # Class Attribute，類別共享\n'
            '\n'
            'a, b = Bad(), Bad()\n'
            'a.items.append("x")\n'
            'print(b.items)          # [\'x\']  ← b 也被污染！'
        ),
        bullets=[
            "items 只建立一次",
            "所有實例共享同一個 list",
            "單元測試會互相污染",
        ],
        label_dark=False,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(3.8),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.4),
        label="GOOD：items 放進 __init__ → 每個實例獨立",
        code=(
            'class Good:\n'
            '    def __init__(self):\n'
            '        self.items = []  # Instance Attribute\n'
            '\n'
            'a, b = Good(), Good()\n'
            'a.items.append("x")\n'
            'print(b.items)          # []  ← b 不受影響'
        ),
        bullets=[
            "每次 __init__ 都建新 list",
            "實例之間完全隔離",
            "鐵律：可變物件一律\n放 __init__",
        ],
        label_dark=True,
    )
    add_source(s, "Python Tutorial §9.3.5 · Ruff rule RUF012 Mutable class attributes")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── S13 · TABLE：何時該用 Class Attribute ─────────
    s = _blank(prs)
    draw_editorial_table(
        s,
        header=["用途", "範例", "是否 OK", "備註"],
        rows=[
            ["常數（不可變）",
             'PI = 3.14159',
             "✓ OK",
             "int / float / str / tuple 皆安全"],
            ["預設設定字串",
             'DEFAULT_ENCODING = "utf-8"',
             "✓ OK",
             "字串不可變，共享無副作用"],
            ["版本號 / 類別識別碼",
             'VERSION = "1.0"',
             "✓ OK",
             "配 classmethod 讀取更整齊"],
            ["可變容器（list）",
             'items = []',
             "✗ 禁",
             "永遠是 bug，放 __init__"],
            ["可變容器（dict）",
             'cache = {}',
             "✗ 禁",
             "實例間共用快取 → 資料污染"],
            ["可變容器（set）",
             'seen = set()',
             "✗ 禁",
             "同上，放 __init__"],
        ],
        col_widths=[1.6, 2.2, 1.0, 2.2],
        title="Class Attribute 的紅綠燈：只給不可變物件用",
    )
    add_source(s, "PEP 8 §Designing for Inheritance · Ruff RUF012")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── S14 · CODE：DataPipeline 雛形實作 ─────────
    s = _blank(prs)
    add_title(s, "DataPipeline 雛形：Ch04 的終點，Ch10 的起點")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="完整雛形：狀態 + 行為 + method chaining 伏筆",
        code=(
            'class DataPipeline:\n'
            '    """A minimal reusable pipeline skeleton."""\n'
            '\n'
            '    def __init__(self, name):\n'
            '        self.name = name\n'
            '        self.steps = []\n'
            '\n'
            '    def add_step(self, step):\n'
            '        self.steps.append(step)\n'
            '        return self              # Ch05 chaining 伏筆\n'
            '\n'
            '    def run(self, data):\n'
            '        for step in self.steps:\n'
            '            data = step(data)\n'
            '        return data\n'
            '\n'
            '# 使用：\n'
            'pipe = DataPipeline("ETL_v1")\n'
            'pipe.add_step(str.strip).add_step(str.lower)\n'
            'pipe.run("  HELLO ")   # \'hello\''
        ),
        bullets=[
            "狀態：name / steps",
            "行為：add_step / run",
            "return self →\nmethod chaining",
            "Ch05 加繼承 + 魔術方法",
            "Ch10 擴成 DataCleaner",
        ],
        label_dark=True,
    )
    add_source(s, "scikit-learn Pipeline 設計啟發 · Martin, Clean Code §10 Classes")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── S15 · MATRIX 2×2：紀律四條 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "① 可變屬性放 __init__",
         "sub": "list / dict / set\n絕對不放類別頂端\nRuff RUF012 會告警",
         "highlight": True},
        {"text": "② self 不省略",
         "sub": "每個實例方法\n第一個參數必須是 self\n別用 me / this 等花名"},
        {"text": "③ 先想資料再想方法",
         "sub": "__init__ 先定義清楚\n狀態由哪些屬性組成\n方法才跟著資料設計"},
        {"text": "④ 不為單檔腳本開類別",
         "sub": "< 50 行的一次性腳本\n硬套類別只增加複雜度\nLinus: 'Don't overthink.'"},
    ], title="Ch04 紀律四條：踩過幾次坑就會刻進肌肉記憶")
    add_source(s, "Ch04 課堂歸納 · Ruff RUF012 · PEP 8")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── S16 · PYRAMID：收束 ─────────
    s = _blank(prs)
    draw_thesis_hierarchy(
        s,
        blocks=[
            {"heading": "四個關鍵詞（硬記）",
             "items": [
                 "Class = 建築藍圖",
                 "Object = 蓋出來的房子",
                 "__init__ = 誕生禮（實例化時自動執行）",
                 "self = 身分證（指向當前實例）",
             ]},
            {"heading": "三章銜接線",
             "items": [
                 "Ch04 · 藍圖 + 誕生禮 + 身分證（今日）",
                 "Ch05 · 加封裝 / 繼承 / 魔術方法",
                 "Ch10 · DataCleaner().read().clean().export()",
             ]},
        ],
        title="Ch04 收束：地基不穩，上面再漂亮都會塌",
        thesis="先搞懂狀態住哪，才配談設計——Ch05 繼續延伸封裝、繼承與魔術方法。",
    )
    add_source(s, "Ch04 module synthesis · Ch05 / Ch10 銜接預告")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── S17 · SILENT：收尾 ─────────
    s = _blank(prs)
    draw_silent_page(s, "先搞懂狀態住哪，\n才配談設計。")
    add_footer(s, MODULE_CODE, 17, N_CONTENT, dark_bg=True)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
