"""Ch08 deck — Pandas 資料清洗與特徵工程
23 content slides + cover + copyright page (本課程最長章, M3 · 3.5 hr).

Governing thought:
    Pandas 不是 API 字典，是資料工程師的工作台 ——
    讀→看→篩→清→轉→聚→併，七步流水線吃掉 80% 的真實工作時間。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_rect, add_textbox, add_title, add_source,
    set_solid_fill, set_no_line, set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_vs_two_col, draw_code_panel, draw_thesis_hierarchy,
    draw_inverted_thesis_box, draw_image_placeholder,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch08"
MODULE_TITLE = "Pandas 資料清洗與特徵工程"
MODULE_SUBTITLE = "讀→看→篩→清→轉→聚→併：資料工程師的工作台"
TIME_MIN = 210
N_CONTENT = 23


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_ch08(output_path, image_registry=None):
    """Build Ch08 deck.

    image_registry: optional list collecting image placeholder metadata
                    for downstream YAML emission.
    """
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── S1 · SILENT ─────────
    s = _blank(prs)
    draw_silent_page(s, "Pandas 不是 API 字典，\n是資料工程師的工作台。")
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ───────── S2 · ASK ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "如果 80% 的工作時間花在清洗，\n為什麼大家還是邊查文件邊寫？",
        data_card={
            "label": "Kaggle 2024 State of ML & DS Survey",
            "stat": "80%",
            "caption": "受訪資料工作者選 pandas 為日常工具\n且坦承『大多 API 用過就忘』",
        },
    )
    add_source(s, "Kaggle, State of Machine Learning and Data Science 2024 §Tools")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── S3 · MATRIX 2×3：Series / DataFrame / Index 三角 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=3, cells=[
        {"text": "Series",
         "sub": "一維帶索引陣列\nNumPy + Index 標籤\n單欄資料的化身",
         "highlight": True},
        {"text": "DataFrame",
         "sub": "二維表\n多個 Series 共用同一 Index\n類似 SQL Table"},
        {"text": "Index",
         "sub": "資料的「名字」\n不是欄位、不是 row number\n是對齊與合併的根據"},
        {"text": "對齊（alignment）",
         "sub": "兩個 Series 相加\n自動依 Index 對齊\n缺位 → NaN",
         "highlight": True},
        {"text": "Index 不可變",
         "sub": "建立後不能 in-place 改值\n要改用 reindex / rename"},
        {"text": "set/reset_index",
         "sub": "在「Index」與「欄位」\n之間互換的兩道門"},
    ], title="Series / DataFrame / Index：pandas 的三角心智模型")
    add_source(s, "pandas User Guide §Intro to data structures · Wes McKinney, Python for Data Analysis 3e Ch5")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───────── S4 · CODE：建立 + 6 個觀察方法 ─────────
    s = _blank(prs)
    add_title(s, "建立 DataFrame 的三條路 + 看資料的 6 個方法")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="拿到任何新資料：先建立、再用這 6 個方法把它看穿",
        code=(
            'import pandas as pd\n'
            '\n'
            '# 三種建立路徑\n'
            'df1 = pd.DataFrame({"name": ["A","B"], "age": [20, 30]})\n'
            'df2 = pd.DataFrame([{"name":"A","age":20}, {"name":"B","age":30}])\n'
            'df3 = pd.read_csv("orders.csv")\n'
            '\n'
            '# 觀察 6 件套（每次拿到新資料先跑一遍）\n'
            'df3.head()        # 看前 5 列：長相\n'
            'df3.tail()        # 看末 5 列：尾巴乾淨嗎\n'
            'df3.info()        # 型別 + 缺失：debug 第一站\n'
            'df3.describe()    # 數值欄統計：分佈如何\n'
            'df3.dtypes        # 各欄型別：物件型最常出包\n'
            'df3.shape         # (rows, cols)：規模感'
        ),
        bullets=[
            "別急著 query — 先 info()",
            "object dtype 八成是字串\n或混型別 → 麻煩源頭",
            "describe 只看數值欄\n類別欄要用 value_counts",
            "shape 小於預期 → 讀檔出錯",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Essential basic functionality · PyData best practices")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── S5 · IMAGE + CODE：Index 對齊 ─────────
    s = _blank(prs)
    add_title(s, "Index 不是欄位：標籤對齊才是 pandas 的超能力")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.5), h=Inches(5.0),
        slot_name="Index 對齊示意圖",
        description=(
            "兩個 Series s1=[A,B,C] s2=[B,C,D]（不同 Index）\n"
            "中間 + 號 → 結果 [A=NaN, B=合, C=合, D=NaN]\n"
            "強調「對齊靠名字、不是位置」"
        ),
        url_hint="",
        placeholder_id="Ch08_S05_index_alignment",
        registry=image_registry,
        size_hint="1280×720 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.4), y=Inches(1.3),
        w=Inches(6.3), h=Inches(5.0),
        label="標籤對齊：pandas 的隱藏外掛",
        code=(
            'import pandas as pd\n'
            '\n'
            's1 = pd.Series([10, 20, 30],\n'
            '               index=["A","B","C"])\n'
            's2 = pd.Series([1, 2, 3],\n'
            '               index=["B","C","D"])\n'
            '\n'
            's1 + s2\n'
            '# A     NaN  ← 只在 s1\n'
            '# B    21.0\n'
            '# C    32.0\n'
            '# D     NaN  ← 只在 s2'
        ),
        bullets=[
            "對齊靠 Index 名字，\n不是位置",
            "Excel 沒這個能力",
            "merge / join 的底層原理",
            "理解這張圖 = 理解 50% pandas",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Intro §Alignment · McKinney, Python for Data Analysis 3e §5.2")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── S6 · MATRIX 2×2：四把鑰匙 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "df.loc[label]",
         "sub": "用 Index / 欄名取\n包前包後（與 Python 不同）\nSQL 直覺、可讀性高",
         "highlight": True},
        {"text": "df.iloc[pos]",
         "sub": "用整數位置取\n包前不包後（Pythonic）\n資料順序敏感"},
        {"text": "df[df.col > x]",
         "sub": "布林遮罩\n條件式篩選、最 Pythonic\n多條件用 & | 加括號",
         "highlight": True},
        {"text": "df.query(\"col > x\")",
         "sub": "字串條件、可讀性最高\n可用 @var 帶外部變數\n複雜邏輯不會一坨括號"},
    ], title="取資料的四把鑰匙：loc / iloc / 布林 / query")
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "四把鑰匙各有主場，沒有最強只有最對 — 用標籤就 loc、用位置就 iloc、邏輯複雜就 query。",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "pandas User Guide §Indexing and selecting data")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───────── S7 · VS-CODE：loc vs iloc ─────────
    s = _blank(prs)
    add_title(s, "loc vs iloc：用標籤還是用位置 — bug 第一名是混用")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.3),
        label="iloc — 用整數位置（包前不包後，像 Python list）",
        code=(
            'df.iloc[0:3, 0:2]\n'
            '# 取「位置 0、1、2」三列、「位置 0、1」兩欄\n'
            '# end 不含 → 與 list[0:3] 同邏輯'
        ),
        bullets=[
            "資料順序敏感",
            "排序後同樣寫法\n結果就變了",
            "適合 ETL / 隨機抽樣",
        ],
        label_dark=False,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(3.7),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.5),
        label="loc — 用標籤（包前包後，像 SQL）",
        code=(
            'df.loc["A":"C", "name":"age"]\n'
            '# 取 Index 為「A、B、C」三列、欄位「name 到 age」\n'
            '# end 含 → 與 SQL BETWEEN 一致\n'
            '\n'
            'df.loc[df["age"] > 30, ["name","city"]]\n'
            '# 也可放布林條件 + 欄位選取 → 同時做篩列+選欄'
        ),
        bullets=[
            "業務需求 90% 用 loc",
            "可放布林、標籤、切片",
            "iloc 的 0:3 與 loc 的 0:3\n結果不同 → 切莫混用",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Indexing — Different choices for indexing")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── S8 · CODE：query + Copy-on-Write ─────────
    s = _blank(prs)
    add_title(s, "query 與 Copy-on-Write：pandas 2.x 的兩個要點")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="多條件篩選：從一坨括號 → query 字串",
        code=(
            '# 傳統布林：括號滿天飛\n'
            'df[(df["age"] > 30) & (df["city"] == "Taipei")\n'
            '   & (df["score"].between(60, 90))]\n'
            '\n'
            '# 改用 query：可讀性大躍進\n'
            'min_age = 30\n'
            'df.query("age > @min_age and city == \\"Taipei\\" "\n'
            '         "and 60 <= score <= 90")\n'
            '\n'
            '# pandas 2.x 預設 Copy-on-Write\n'
            'pd.options.mode.copy_on_write = True   # 預設已 True\n'
            '\n'
            '# 鏈式賦值：將被棄用，正確寫法用 .loc 一次到位\n'
            'df.loc[df["age"] > 30, "tag"] = "senior"   # ✓'
        ),
        bullets=[
            "query 用 @var 帶外部變數",
            "字串條件易組合、易測試",
            "Copy-on-Write =\n「改一份不影響另一份」",
            "鏈式賦值 → SettingWithCopyWarning\n直接 .loc 一句寫完",
        ],
        label_dark=True,
    )
    add_source(s, "pandas 2.x Release Notes §Copy-on-Write · pandas docs §Indexing.query")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── S9 · ASK：缺失值的真相 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "缺失值出現了 ——\n你第一個想到 dropna()，\n還是先問為什麼缺？",
        data_card={
            "label": "真實專案缺失值來源拆解",
            "stat": "3 種",
            "caption": "「沒填」、「編碼問題」、「計算產生」\n處理策略截然不同",
        },
    )
    add_source(s, "Hadley Wickham, Tidy Data §missing data taxonomy 2014")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── S10 · IMAGE + CODE：缺失值決策樹 ─────────
    s = _blank(prs)
    add_title(s, "缺失值處理決策樹：先偵測、再分類、後對症下藥")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.5), h=Inches(5.0),
        slot_name="缺失值決策樹",
        description=(
            "根節點「為什麼缺？」分三支：\n"
            "(1) 沒填 → dropna 或 fillna\n"
            "(2) 編碼問題 → na_values 在讀檔時擋下\n"
            "(3) 計算產生 → 通常保留或 interpolate"
        ),
        url_hint="",
        placeholder_id="Ch08_S10_missing_decision_tree",
        registry=image_registry,
        size_hint="1400×900 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.4), y=Inches(1.3),
        w=Inches(6.3), h=Inches(5.0),
        label="先偵測 → 分類 → 對症下藥",
        code=(
            '# 1. 偵測：每欄缺幾筆\n'
            'df.isna().sum()\n'
            'df.isna().mean()       # 缺失比例\n'
            '\n'
            '# 2. 編碼問題 → 讀檔時擋下\n'
            'df = pd.read_csv("data.csv",\n'
            '                 na_values=["-", "N/A", "?", ""])\n'
            '\n'
            '# 3. 沒填 → 看資料量決定\n'
            'df = df.dropna(subset=["price"])      # 關鍵欄\n'
            'df["age"] = df["age"].fillna(df["age"].median())\n'
            '\n'
            '# 4. 計算產生（時間平移、除以 0）\n'
            'df["pct"] = df["pct"].interpolate()'
        ),
        bullets=[
            "決策樹照走\n通過 review 比較容易",
            "編碼問題在讀檔時就擋，\n別等清洗階段才補救",
            "fillna(mean) 改變分佈，\nfillna(median) 比較穩",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Working with missing data · McKinney 3e §7.1")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── S11 · TABLE：四種策略 ─────────
    s = _blank(prs)
    draw_editorial_table(
        s,
        header=["策略", "API", "何時用", "風險"],
        rows=[
            ["整列刪除", "dropna()",
             "缺失少且資料量充足",
             "可能丟掉「缺本身就是訊號」"],
            ["整欄刪除", "dropna(axis=1, thresh=...)",
             "整欄高比例缺失（> 70%）",
             "損失維度、未來無法用該特徵"],
            ["填補定值", "fillna(0 / mean / median / mode)",
             "缺失有合理替代值",
             "改變分佈、平均/中位數彼此差很多"],
            ["向前 / 內插", "ffill / bfill / interpolate()",
             "時間序列、有時序鄰近性",
             "假設「鄰近=相似」可能不成立"],
        ],
        col_widths=[1.2, 2.4, 2.0, 2.4],
        title="缺失值四種策略：何時 drop / fillna / ffill / interpolate",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.4),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "沒有預設答案 — 看資料、看下游、看 stakeholder。能在讀檔擋下的就在讀檔擋下。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "pandas docs §Missing data · Kaggle Learn — Data Cleaning §Handling Missing Values")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── S12 · SILENT：Lambda 找到主場 ─────────
    s = _blank(prs)
    draw_silent_page(s, "Ch3 學過的 lambda，\n今天在 Pandas 找到主場。")
    add_footer(s, MODULE_CODE, 12, N_CONTENT, dark_bg=True)

    # ───────── S13 · CODE：apply 三層級 ─────────
    s = _blank(prs)
    add_title(s, "apply 三個層級：Series / DataFrame / element-wise（map）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="lambda + apply 是 pandas 最常見的搭配",
        code=(
            '# 1. Series.apply — 對每個元素\n'
            'df["price_with_tax"] = df["price"].apply(lambda x: x * 1.05)\n'
            '\n'
            '# 2. DataFrame.apply(axis=0) — 對每欄（預設）\n'
            'df[["a","b"]].apply(lambda col: col.max() - col.min())\n'
            '\n'
            '# 3. DataFrame.apply(axis=1) — 對每列\n'
            'df["full_name"] = df.apply(\n'
            '    lambda row: f"{row[\'first\']} {row[\'last\']}", axis=1)\n'
            '\n'
            '# 4. Series.map — 字典映射（取代 deprecated 的 applymap）\n'
            'df["grade_label"] = df["grade"].map(\n'
            '    {"A": "優", "B": "可", "C": "待加強"})'
        ),
        bullets=[
            "Ch3 lambda 在這找到舞台",
            "axis=0：對每欄做（預設）\naxis=1：對每列做",
            "row 級 apply 是新手最常\n用錯 axis 之處",
            "字典映射用 map 比 apply 更直覺",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Function application · Ch3 Lambda 章節銜接")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── S14 · VS-CODE：apply vs 向量化 ─────────
    s = _blank(prs)
    add_title(s, "apply(lambda) vs 向量化：方便 vs 速度的真實落差")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.3),
        label="BEFORE：apply 版 — 1,000 萬列要 ~10 秒",
        code=(
            'df["tax"] = df["price"].apply(lambda x: x * 1.05)\n'
            '# 逐元素呼叫 Python function\n'
            '# 直譯器迴圈 + lambda 封裝開銷 → 慢'
        ),
        bullets=[
            "方便、可讀",
            "但每筆都進 Python 直譯器",
            "資料量大就是災難",
        ],
        label_dark=False,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(3.7),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.5),
        label="AFTER：向量化 — 同樣 1,000 萬列只要 ~0.05 秒（≈ 200×）",
        code=(
            'df["tax"] = df["price"] * 1.05\n'
            '# Series 算術 → 走 NumPy C 層\n'
            '# 一句話、無 Python 迴圈\n'
            '\n'
            '# 字串 / 時間也有專屬 accessor\n'
            'df["clean"] = df["name"].str.strip().str.lower()\n'
            'df["year"]  = df["date"].dt.year'
        ),
        bullets=[
            "能向量化就不要 apply",
            "apply 是 fallback，\n不是 first choice",
            ".str / .dt 是兩個高頻 accessor",
            "Ch7 NumPy 的向量化\n在這裡發揮作用",
        ],
        label_dark=True,
    )
    add_source(s, "pandas docs §Enhancing performance · McKinney 3e §7.4 vectorized string functions")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── S15 · CODE：特徵工程實戰 ─────────
    s = _blank(prs)
    add_title(s, "特徵工程實戰：日期拆解 + 分組標記 + 字串清洗（90% 場景）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="三招組合，覆蓋日常 90% 特徵工程",
        code=(
            '# 招式 1：日期拆解（.dt accessor）\n'
            'df["date"] = pd.to_datetime(df["date"])\n'
            'df["year"]      = df["date"].dt.year\n'
            'df["month"]     = df["date"].dt.month\n'
            'df["dayofweek"] = df["date"].dt.dayofweek      # 0=Mon\n'
            'df["is_weekend"] = df["dayofweek"].isin([5, 6])\n'
            '\n'
            '# 招式 2：分組標記（apply + lambda）\n'
            'df["grade"] = df["score"].apply(\n'
            '    lambda x: "A" if x >= 90 else "B" if x >= 60 else "C")\n'
            '\n'
            '# 招式 3：字串清洗（.str accessor）\n'
            'df["email"] = df["email"].str.strip().str.lower()\n'
            'df["domain"] = df["email"].str.split("@").str[1]'
        ),
        bullets=[
            ".dt / .str 兩大高頻 accessor",
            "能用 accessor 就不要 apply",
            "特徵工程的本質：\n把隱藏訊號顯現出來",
            "is_weekend、domain 這類\n衍生特徵下游模型很愛吃",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Working with text data · §Time series / date functionality")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── S16 · IMAGE + CODE：Split-Apply-Combine ─────────
    s = _blank(prs)
    add_title(s, "groupby 心智模型：Split → Apply → Combine（背概念，不背 API）")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.5), h=Inches(5.0),
        slot_name="Split-Apply-Combine 流程圖",
        description=(
            "完整 DataFrame → 依 city 切 3 組\n"
            "（Taipei / Taichung / Kaohsiung）\n"
            "→ 每組求 mean(revenue)\n"
            "→ 合併回單一結果表"
        ),
        url_hint="",
        placeholder_id="Ch08_S16_split_apply_combine",
        registry=image_registry,
        size_hint="1400×800 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.4), y=Inches(1.3),
        w=Inches(6.3), h=Inches(5.0),
        label="一句話 = split + apply + combine 三步驟",
        code=(
            '# 一句話完成三步驟\n'
            'df.groupby("city")["revenue"].mean()\n'
            '\n'
            '# 拆開看就是：\n'
            '#   split   = 依 city 切組\n'
            '#   apply   = 每組對 revenue 算 mean\n'
            '#   combine = 拼回單一結果（每 city 一列）\n'
            '\n'
            '# Wilkinson (1956) 的概念\n'
            '# 今天還是 pandas / dplyr / SQL\n'
            '# 三大資料工具的核心'
        ),
        bullets=[
            "1956 年的概念，到今天還在用",
            "心智模型清晰後\n語法只是裝飾",
            "split 是 hash partition\napply 是 vectorised\ncombine 是 concat",
        ],
        label_dark=True,
    )
    add_source(s, "Wickham, The Split-Apply-Combine Strategy 2011 · pandas User Guide §Group by")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── S17 · CODE：groupby + agg 完整 ─────────
    s = _blank(prs)
    add_title(s, "groupby + agg：單鍵、多鍵、多欄、自訂聚合一次到位")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="agg 是進階入口 — 業務報表 90% 用這套",
        code=(
            '# 1. 單欄單聚合\n'
            'df.groupby("city")["revenue"].sum()\n'
            '\n'
            '# 2. 多鍵分組（產生 MultiIndex）\n'
            'df.groupby(["city", "product"])["revenue"].sum()\n'
            '\n'
            '# 3. 多欄不同聚合（業務報表最愛）\n'
            'df.groupby("city").agg({\n'
            '    "revenue": "sum",\n'
            '    "qty":     "mean",\n'
            '    "order_id": "count",\n'
            '})\n'
            '\n'
            '# 4. 自訂聚合（lambda 上場）\n'
            'df.groupby("city")["revenue"].agg(\n'
            '    range_=lambda g: g.max() - g.min(),\n'
            '    p90=lambda g: g.quantile(0.9))'
        ),
        bullets=[
            "字典版 agg：欄位 → 函式",
            "lambda 接整個 group（Series）",
            "命名聚合（named aggregation）\n讓結果欄名乾淨",
            "多鍵 → MultiIndex\n常用 .reset_index() 攤平",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Group by — Aggregation · McKinney 3e §10.2")
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───────── S18 · CODE：transform vs agg + pivot_table ─────────
    s = _blank(prs)
    add_title(s, "transform vs agg + pivot_table：保 shape vs 縮減 vs 交叉表")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="同樣是 groupby — transform 保 shape、agg 縮減、pivot_table 變寬",
        code=(
            '# agg：縮減 — 每組變一個值（rows = 組數）\n'
            'df.groupby("city")["revenue"].agg("mean")\n'
            '\n'
            '# transform：保持原 shape — 每列補上組內統計\n'
            'df["city_avg"] = (\n'
            '    df.groupby("city")["revenue"].transform("mean"))\n'
            '# 結果：每筆訂單都帶上「該城市平均」一欄\n'
            '# → 立刻可算 (revenue - city_avg) 衡量相對表現\n'
            '\n'
            '# pivot_table：wide-format 交叉表（Excel 樞紐分析的 pandas 版）\n'
            'df.pivot_table(\n'
            '    index="city", columns="month",\n'
            '    values="revenue", aggfunc="sum",\n'
            '    fill_value=0, margins=True)'
        ),
        bullets=[
            "transform 是初學者最忽略的方法",
            "「給每筆訂單標城市平均」\n業務上每天遇到",
            "pivot_table 是 wide format\n適合人看；agg 是 long format\n適合下游程式吃",
            "margins=True → 加總列/欄",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Group by — Transformation · §Reshaping pivot tables")
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ───────── S19 · MATRIX 2×2：三種合併 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "concat",
         "sub": "沿軸堆疊、不依賴鍵\naxis=0 加列、axis=1 加欄\nETL 拼批次資料最常用",
         "highlight": True},
        {"text": "merge",
         "sub": "依鍵值合併（類似 SQL JOIN）\non= / left_on= / right_on=\n業務 90% 場景",
         "highlight": True},
        {"text": "join",
         "sub": "簡化版 merge\n預設用 Index 合併\n兩表都已 set_index 時最快"},
        {"text": "陷阱 ：_x / _y 重複欄",
         "sub": "兩表同名非鍵欄會自動加後綴\n用 suffixes=('_l','_r')\n或先 rename 一邊"},
    ], title="concat / merge / join：三種合併、各自場景")
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "工業上 90% 用 merge，concat 是 ETL 拼批次，join 是 Index 已對齊時的捷徑。",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "pandas User Guide §Merge, join, concatenate, and compare")
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───────── S20 · IMAGE + CODE：merge 四種 how ─────────
    s = _blank(prs)
    add_title(s, "merge 四種 how：inner / left / right / outer 一張圖說完")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.5), h=Inches(5.0),
        slot_name="merge how Venn 圖",
        description=(
            "四個 Venn 圖水平排列：\n"
            "(1) inner — 交集塗色\n"
            "(2) left — 左圓全塗\n"
            "(3) right — 右圓全塗\n"
            "(4) outer — 整體塗色"
        ),
        url_hint="https://pandas.pydata.org/docs/user_guide/merging.html",
        placeholder_id="Ch08_S20_merge_how_diagram",
        registry=image_registry,
        size_hint="1280×720 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.4), y=Inches(1.3),
        w=Inches(6.3), h=Inches(5.0),
        label="四種 how — 對應 SQL JOIN，預設 inner",
        code=(
            'orders.merge(users, on="user_id",\n'
            '             how="inner")     # 預設：取交集\n'
            'orders.merge(users, on="user_id",\n'
            '             how="left")      # 保留左表全部\n'
            'orders.merge(users, on="user_id",\n'
            '             how="right")     # 保留右表全部\n'
            'orders.merge(users, on="user_id",\n'
            '             how="outer",     # 全保留，缺的 NaN\n'
            '             indicator=True)  # 加 _merge 欄位\n'
            '\n'
            '# debug 黃金公式\n'
            'before = len(orders)\n'
            'after  = len(orders.merge(users, on="user_id"))\n'
            'assert after == before, "鍵不唯一 → 暴漲"'
        ),
        bullets=[
            "預設 how='inner' — 取交集",
            "left 最常用（保留主表）",
            "indicator=True →\n_merge 欄是 debug 神器",
            "merge 完一定 len() 比對\n多了 = 鍵重複\n少了 = 該 left 沒 left",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Database-style DataFrame or named Series joining/merging")
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ───────── S21 · CODE：to_datetime + DatetimeIndex ─────────
    s = _blank(prs)
    add_title(s, "to_datetime + DatetimeIndex：字串變時間就解鎖切片語法")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="先把日期變成「真正的時間」，pandas 才會幫你做事",
        code=(
            '# 1. 字串 → datetime（記得指定 format 加速）\n'
            'df["date"] = pd.to_datetime(\n'
            '    df["date"], format="%Y-%m-%d", errors="coerce")\n'
            '\n'
            '# 2. 設成 Index → 解鎖 partial string indexing\n'
            'df = df.set_index("date").sort_index()\n'
            '\n'
            '# 3. 切片如行雲流水\n'
            'df["2024-01"]                # 整個 1 月\n'
            'df["2024-01":"2024-03"]      # Q1\n'
            'df.loc["2024-Q1"]            # 同上，季表示法\n'
            '\n'
            '# 4. .dt accessor 拆出年月日\n'
            'df_reset = df.reset_index()\n'
            'df_reset["weekday"] = df_reset["date"].dt.day_name()'
        ),
        bullets=[
            "errors='coerce'：壞值轉 NaT，\n不會整個炸掉",
            "format= 對大檔加速 100×",
            "DatetimeIndex 切片 =\n隱藏外掛",
            "排序好的 Index 才能 partial slicing",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Time series / date functionality §Indexing")
    add_footer(s, MODULE_CODE, 21, N_CONTENT)

    # ───────── S22 · CODE：resample + rolling ─────────
    s = _blank(prs)
    add_title(s, "resample 與 rolling：時間版的 groupby + 移動平均")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="resample = 時間軸 groupby；rolling = 滑動窗口統計",
        code=(
            '# 前置：df 已是 DatetimeIndex\n'
            '\n'
            '# 1. resample — 時間版 groupby\n'
            'df["revenue"].resample("D").sum()      # 每日\n'
            'df["revenue"].resample("W").mean()     # 每週\n'
            'df["revenue"].resample("M").sum()      # 每月\n'
            'df["revenue"].resample("Q").sum()      # 每季\n'
            '\n'
            '# 2. rolling — 移動平均（最常見 7 / 30 天）\n'
            'df["ma7"]  = df["revenue"].rolling(window=7).mean()\n'
            'df["ma30"] = df["revenue"].rolling(window=30).mean()\n'
            '\n'
            '# 3. 兩者組合：先 resample 到日，再算 7 日 MA\n'
            'daily = df["revenue"].resample("D").sum()\n'
            'daily_ma7 = daily.rolling(7).mean()'
        ),
        bullets=[
            "resample 一句話完成\n「每日 → 每週」聚合",
            "rolling 是「最近 N 天」窗口",
            "兩者組合是時間序列\n分析的基本盤",
            "資料若有缺日 → resample\n會自動填入 NaN",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Resampling · §Rolling windows · McKinney 3e §11.6")
    add_footer(s, MODULE_CODE, 22, N_CONTENT)

    # ───────── S23 · PYRAMID：收束 ─────────
    s = _blank(prs)
    draw_thesis_hierarchy(
        s,
        blocks=[
            {"heading": "七步流水線（背下這個就贏一半）",
             "items": [
                 "讀（read_csv / parquet）→ 看（info/describe/dtypes/shape）",
                 "篩（loc / iloc / query）→ 清（dropna / fillna / interpolate）",
                 "轉（apply+lambda / .dt / .str / map）",
                 "聚（groupby+agg / transform / pivot_table）→ 併（concat / merge / join）",
             ]},
            {"heading": "今天該帶走的四條紀律",
             "items": [
                 "Index 是「名字」，不是欄位、不是位置",
                 "能向量化就別 apply（200× 落差）",
                 "dropna 不是預設答案 — 先問為什麼缺",
                 "merge 完一定用 len() 比對前後",
             ]},
        ],
        title="Ch08 收束：七步流水線 — 讀→看→篩→清→轉→聚→併",
        thesis="Ch09 進入 Matplotlib — 把這條流水線的產物畫出來給人看；Ch10 把整條流水線封裝成 DataCleaner。",
    )
    add_source(s, "Ch08 module synthesis · 銜接 Ch09 視覺化 / Ch10 OOP×Pandas 整合")
    add_footer(s, MODULE_CODE, 23, N_CONTENT)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
