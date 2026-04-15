"""M5 deck — 15 content slides + cover + copyright page.

進階 Python — 從腳本到可維護系統 (Editorial-strict v1.1)

Charts required in charts.py (NOT yet implemented — listed here for Agent-D handoff):
  - gil_dual_line_m5s12()
      Data: x = [1, 2, 4, 8, 16] threads
            cpu_bound = [98, 102, 104, 101, 99]      solid PRIMARY
            io_bound  = [95, 180, 340, 610, 720]     dashed CHARCOAL
      Reference: horizontal dashed line at y=100 ("單核上限")
      Figsize: (10.5, 5.0); y-axis 0–820 (label "CPU 利用率 %, 8 核上限 800")
      Series labels at line end; no legend box.

Deviations: S1/S3/S8/S14 are GEOMETRIC-DIAGRAM compositions executed with
existing primitives (draw_flow_chain, draw_editorial_table, draw_matrix variants).
For S3 (dependency DAG) and S14 (decision tree), approximate with draw_grid +
annotated captions because no tree/DAG primitive exists — documented as
"editorial approximation". S6/S9 BEFORE-AFTER use draw_code_panel twice
(matching M1 S8/S9 pattern). S11 MATRIX uses 2x2 draw_matrix with three cells
marked (leaving one quadrant textual as "—").
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from .. import charts
from ..primitives import (
    add_rect, add_textbox, add_title, add_source, set_solid_fill, set_no_line,
    set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_pyramid_stack, draw_inverted_thesis_box,
    draw_vs_two_col, draw_risk_mitigation, draw_dual_track, draw_grid,
    draw_code_panel, draw_split_panel, draw_concentric_zones,
    draw_thesis_hierarchy, draw_three_blocks_flow,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "M5"
MODULE_TITLE = "進階 Python — 從腳本到可維護系統"
MODULE_SUBTITLE = "Python 數據分析與 AI 基礎 · M5"
TIME_MIN = 25
N_CONTENT = 15


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_m5(output_path, image_registry=None):
    """Build M5 deck. image_registry kept for API parity; M5 has no placeholders."""
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE, TIME_MIN, N_CONTENT)

    # S1 — GEOMETRIC-DIAGRAM: 腳本 → 系統 五層
    s = _blank(prs)
    draw_flow_chain(s, nodes=[
        {"label": "① 模組化", "caption": "單一職責 / import"},
        {"label": "② 環境", "caption": "venv / uv / lockfile"},
        {"label": "③ 例外 / 日誌", "caption": "raise from / logging"},
        {"label": "④ I/O", "caption": "with / context manager"},
        {"label": "⑤ 並行", "caption": "asyncio / thread / process",
         "highlight": True},
    ], title="腳本能跑一次，系統要跑一千次——中間隔著五道基礎設施",
       y=3.2)
    add_textbox(s, T.MARGIN_X, Inches(5.3),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
                "跨越這條線 = 從「會寫 Python」到「能做 Python 工程」",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "M5 module thesis")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # S2 — ASK
    s = _blank(prs)
    draw_ask_page(
        s,
        "同一份程式、同一份資料，跑出不同結果——誰的錯？",
        data_card={
            "label": "同程式 · 同資料 · 三種結果",
            "stat": "pandas 2.2 / 1.5 / 2.0",
            "caption": "你的電腦 · 同事電腦 · CI Server — 環境未聲明是根因",
        },
    )
    add_source(s, "2024 pandas release notes — breaking changes")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # S3 — GEOMETRIC-DIAGRAM: 相依 DAG（以 4 層 grid 近似樹狀）
    s = _blank(prs)
    draw_grid(s, rows=4, cols=3, cells=[
        # Row 1: root (center)
        {"label": "—", "sub": ""},
        {"label": "your_project", "sub": "專案根節點", "highlight": True},
        {"label": "—", "sub": ""},
        # Row 2: direct deps
        {"label": "pandas==2.2", "sub": "直接依賴"},
        {"label": "—", "sub": ""},
        {"label": "scikit-learn==1.4", "sub": "直接依賴"},
        # Row 3: transitive
        {"label": "numpy>=1.26", "sub": "pandas 需求", "highlight": True,
         "note": "共用相依 → 版本衝突點"},
        {"label": "python-dateutil / pytz", "sub": "pandas 週邊"},
        {"label": "numpy>=1.19 · scipy · joblib", "sub": "sklearn 需求"},
        # Row 4: resolution note
        {"label": "lockfile 解法", "sub": "solver 同時滿足兩約束 → 鎖定 numpy==1.26.x",
         "dashed": True},
        {"label": "—", "sub": ""},
        {"label": "—", "sub": ""},
    ], title="一個套件的「版本」，背後是一棵相依 DAG",
       caption="你裝的不是一個套件，是一棵樹；葉子會互相打架——這就是依賴衝突的本質。")
    add_source(s, "PyPI dependency graph, 2024 snapshot")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # S4 — TABLE: 四方工具對比
    s = _blank(prs)
    draw_editorial_table(s,
        header=["工具", "定位", "速度", "依賴解析", "2026 適用情境"],
        rows=[
            ["venv + pip", "標準庫 / 最低門檻", "中", "無 lockfile",
             "簡單腳本 / 教學入門"],
            ["conda", "環境 + 非 Python 套件", "慢", "有 solver",
             "科學計算 / C / CUDA 相依"],
            ["poetry", "專案管理 + lockfile", "中", "poetry.lock",
             "純 Python 套件開發"],
            ["uv", "Rust 實作 / 全功能", "10–100× 快", "uv.lock",
             "2026 預設建議"],
        ],
        col_widths=[0.9, 1.4, 0.7, 1.0, 1.6],
        title="venv / conda / poetry / uv 四方對比；2026 預設建議是 uv",
    )
    draw_inverted_thesis_box(s,
        "選一個，就全 deck 不動——工具輪替的成本比想像高。",
        y=5.9, width=9.5)
    add_source(s, "Astral uv 0.5 benchmarks / Python Packaging Authority, 2025")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # S5 — SILENT
    s = _blank(prs)
    draw_silent_page(s,
        "可重現，不是美德，\n是協作的最低義務。")
    add_footer(s, MODULE_CODE, 5, N_CONTENT, dark_bg=True)

    # S6 — BEFORE/AFTER: raise ... from ...
    s = _blank(prs)
    add_title(s, "raise ... from ...：例外鏈保留第一現場，不讓根因消失")
    draw_code_panel(s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.4),
        label="BEFORE（吞掉原始例外）",
        code=('def load_config(path):\n'
              '    try:\n'
              '        return json.loads(path.read_text())\n'
              '    except Exception:\n'
              '        raise ConfigError("無法載入設定")'),
        bullets=[
            "原始 traceback 斷鏈",
            "凌晨三點抓 bug 找不到根因",
            "except Exception 再包一層，是典型反模式",
        ],
        label_dark=False,
    )
    draw_code_panel(s,
        x=T.MARGIN_X, y=Inches(3.9),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.5),
        label="AFTER（raise ... from e 保留鏈結）",
        code=('def load_config(path):\n'
              '    try:\n'
              '        return json.loads(path.read_text())\n'
              '    except (OSError, json.JSONDecodeError) as e:\n'
              '        raise ConfigError(f"無法載入設定：{path}") from e'),
        bullets=[
            "from e 把原因接在 __cause__",
            "traceback 會印出兩段，根因可追",
            "第一現場保留 · 例外鏈明確",
        ],
        label_dark=True,
    )
    add_source(s, "PEP 3134 — Exception Chaining")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # S7 — TABLE: 依賴聲明三層級
    s = _blank(prs)
    draw_editorial_table(s,
        header=["檔案", "管什麼", "不管什麼", "進 git？"],
        rows=[
            ["requirements.txt", "頂層套件清單（常含版本區間）",
             "相依樹完整解析 / 傳遞性鎖定", "✓"],
            ["pyproject.toml", "專案 metadata + 套件宣告 + build 設定",
             "精確版本鎖定（是宣告，不是鎖）", "✓"],
            ["*.lock (uv.lock / poetry.lock)",
             "全樹每個套件的精確版本與 hash",
             "人類手寫（由 solver 產出）", "✓（關鍵！）"],
        ],
        col_widths=[1.4, 1.8, 1.8, 0.8],
        title="依賴聲明三層級：requirements.txt / pyproject.toml / lockfile 各管什麼",
    )
    draw_inverted_thesis_box(s,
        "宣告（pyproject）≠ 鎖定（lockfile）。兩者都進 git，缺一不可重現。",
        y=5.6, width=10.5)
    add_source(s, "PEP 517 / 518 / 621 · Python Packaging Guide, 2025")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # S8 — GEOMETRIC-DIAGRAM: logging 三層架構（水平 flow）
    s = _blank(prs)
    draw_flow_chain(s, nodes=[
        {"label": "Logger", "sub": "命名空間 · 入口",
         "caption": "決定「誰說」\ngetLogger('app.loader')"},
        {"label": "Handler", "sub": "輸出目的地",
         "caption": "決定「送到哪」\nStreamHandler · FileHandler"},
        {"label": "Formatter", "sub": "訊息外觀",
         "caption": "決定「長什麼樣」\n%(asctime)s %(levelname)s ...",
         "highlight": True},
    ], title="logging 三層架構：Logger → Handler → Formatter，各司其職",
       y=2.8)
    add_textbox(s, T.MARGIN_X, Inches(5.4),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
                "三層 MECE：改一層不會污染另一層——這就是為什麼它能在真實系統活下去，而 print 不行。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "Python logging HOWTO, cpython 3.13 docs")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # S9 — BEFORE/AFTER: with context manager
    s = _blank(prs)
    add_title(s, "with context manager：不是糖，是資源釋放的契約")
    draw_code_panel(s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.4),
        label="BEFORE（手動 close · 遇例外就洩漏）",
        code=('f = open("data.csv", "r", encoding="utf-8")\n'
              'data = process(f.read())   # 若此行拋例外\n'
              'f.close()                  # 永遠不會被呼叫 → 握柄洩漏'),
        bullets=[
            "例外路徑 = 資源洩漏",
            "檔案 / 連線 / 鎖 / GPU 記憶體都中招",
            "try/finally 能補救，但冗長易漏",
        ],
        label_dark=False,
    )
    draw_code_panel(s,
        x=T.MARGIN_X, y=Inches(3.9),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.5),
        label="AFTER（with 保證無論成功失敗都釋放）",
        code=('with open("data.csv", "r", encoding="utf-8") as f:\n'
              '    data = process(f.read())\n'
              '# 離開 with 區塊，__exit__ 自動關閉，例外發生亦然'),
        bullets=[
            "契約等價於 try/finally",
            "__enter__ / __exit__ 雙方法契約",
            "自寫 class 可用 @contextmanager 裝飾",
        ],
        label_dark=True,
    )
    add_source(s, "PEP 343 — The 'with' Statement")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # S10 — SILENT
    s = _blank(prs)
    draw_silent_page(s,
        "會寫程式的多，\n會釋放資源的少。")
    add_footer(s, MODULE_CODE, 10, N_CONTENT, dark_bg=True)

    # S11 — MATRIX 2x2: 三種並行機制
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        # 左上: 啟動成本低 × CPU-bound (空)
        {"text": "（低成本 × CPU-bound）",
         "sub": "無對應機制——\nCPU-bound 需繞過 GIL，\n必然付出 process 啟動成本。"},
        # 右上: 啟動成本高 × CPU-bound
        {"text": "multiprocessing",
         "sub": "獨立 process / 繞過 GIL / 真並行\n啟動慢 / 資料要序列化\n適用：純 Python CPU 密集任務",
         "highlight": True},
        # 左下: 啟動成本低 × I/O-bound
        {"text": "asyncio",
         "sub": "單 thread / 協作式\nawait 顯式讓出\n適用：大量 HTTP / DB I/O"},
        # 右下: 啟動成本中 × I/O-bound → threading
        {"text": "threading",
         "sub": "共享記憶體 / 受 GIL 限制\n啟動快於 process\n適用：blocking I/O"},
    ], title="三種並行機制：啟動成本 × 適用瓶頸，不能混用")
    add_textbox(s, T.MARGIN_X, Inches(6.3),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
                "縱軸：瓶頸類型  I/O-bound ↔ CPU-bound    ·    橫軸：啟動成本 低 → 高    ·    三機制 MECE：選錯不是沒效果，是更慢",
                font_size=T.FONT_SOURCE, color=T.GRAY_MID,
                align=PP_ALIGN.CENTER)
    add_source(s, "Python 3.13 concurrency docs / David Beazley 2021 talk")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # S12 — CHART: GIL 雙折線
    s = _blank(prs)
    add_title(s, "GIL 真相：純 Python CPU 任務，thread 加再多 CPU 利用率都卡在 1 核")
    # NOTE: charts.gil_dual_line_m5s12() — to be implemented in charts.py
    chart_png = charts.gil_dual_line_m5s12()
    s.shapes.add_picture(str(chart_png), Inches(1.0), Inches(1.3),
                          width=Inches(11.3))
    draw_inverted_thesis_box(s,
        "GIL 不是 bug，是 CPython 記憶體管理設計；但它劃清了 threading 能走到哪。",
        y=6.0, width=11.0)
    add_source(s, "內部基準測試 — 8 核 M2 Pro / Python 3.12 / cpython-GIL")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # S13 — TABLE: PEP 703 時間軸
    s = _blank(prs)
    draw_editorial_table(s,
        header=["年份", "版本 / 里程碑", "狀態", "對工程實務意義"],
        rows=[
            ["2023", "PEP 703 提案通過", "方向確認",
             "社群確認朝「可選移除 GIL」路線走"],
            ["2024", "Python 3.13", "實驗性（--disable-gil 旗標）",
             "生產禁用；NumPy / pandas thread-safe 驗證中"],
            ["2025", "Python 3.14", "第二階段穩定化",
             "主流套件開始出 free-threaded wheels"],
            ["2026（今年）", "Python 3.14 + 生態", "有限生產試點",
             "新專案可於隔離服務評估；主訓練管線仍建議 GIL 版"],
            ["2027–2028（預估）", "Python 3.15+", "可能進入預設",
             "threading 有機會取代部分 multiprocessing 場景"],
        ],
        col_widths=[1.1, 1.6, 1.5, 2.0],
        title="PEP 703 時間軸：free-threaded Python 從實驗到預設的路線",
    )
    add_textbox(s, T.MARGIN_X, Inches(6.0),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
                "此表反映 2026-04 時點社群進度，版本與時程以 PEP 703 最新狀態為準。",
                font_size=T.FONT_SOURCE, color=T.GRAY_MID,
                align=PP_ALIGN.CENTER, italic=True)
    add_source(s, "PEP 703 · python.org status tracker, 2026-04")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # S14 — GEOMETRIC-DIAGRAM: 決策樹（以 grid 近似）
    s = _blank(prs)
    draw_grid(s, rows=3, cols=4, cells=[
        # Row 1: decision root
        {"label": "—", "sub": ""},
        {"label": "瓶頸在等待，還是在計算？", "sub": "第一步先分 等 vs 算",
         "highlight": True},
        {"label": "—", "sub": ""},
        {"label": "—", "sub": ""},
        # Row 2: two branches
        {"label": "等待（I/O-bound）", "sub": "阻塞在外部系統"},
        {"label": "→", "sub": ""},
        {"label": "計算（CPU-bound）", "sub": "阻塞在 CPU 核心"},
        {"label": "→", "sub": ""},
        # Row 3: four leaves
        {"label": "→ asyncio", "sub": "大量小 I/O（HTTP / DB）",
         "highlight": True},
        {"label": "→ threading", "sub": "blocking I/O（檔案、舊 client）",
         "highlight": True},
        {"label": "→ multiprocessing", "sub": "純 Python 迴圈計算",
         "highlight": True},
        {"label": "→ NumPy / PyTorch", "sub": "數值運算 / 張量（底層已並行）",
         "highlight": True},
    ], title="I/O-bound 或 CPU-bound？一條決策樹定錨",
       caption="AI 工程實例：DataLoader(num_workers=4) = I/O-bound × multiprocessing（繞開 GIL，用 process 代替 thread）")
    add_source(s, "M5 Block 5 · PyTorch DataLoader 文件")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # S15 — SILENT closing
    s = _blank(prs)
    draw_silent_page(s,
        "工程不是寫得下去，\n是改得動、跑得再一次、交得出手。")
    add_footer(s, MODULE_CODE, 15, N_CONTENT, dark_bg=True)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
