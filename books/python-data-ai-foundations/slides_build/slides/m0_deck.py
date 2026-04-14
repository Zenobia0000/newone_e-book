"""M0 deck — 15 content slides + cover + copyright page."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from .. import charts
from ..primitives import (
    add_rect, add_textbox, add_title, add_source, set_solid_fill, set_no_line,
    set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_pyramid_stack, draw_inverted_thesis_box,
    draw_vs_two_col, draw_risk_mitigation, draw_dual_track, draw_grid,
    draw_photo_placeholder_triptych, draw_image_placeholder,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "M0"
MODULE_TITLE = "開場：Python 與 AI 系統全景"
MODULE_SUBTITLE = "Python 數據分析與 AI 基礎 · 開場模組"
TIME_MIN = 25
N_CONTENT = 15
COURSE_HOURS_LABEL = "24 小時"  # 7 內容模組 M1-M7 + M0 開場，總課程時長 24h


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_m0(output_path, image_registry=None):
    """Build M0 deck.

    image_registry: optional list to collect image placeholder metadata
                    for downstream YAML emission.
    """
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE, TIME_MIN, N_CONTENT)

    # S1 — SILENT
    s = _blank(prs)
    draw_silent_page(s, "Data is the raw material;\nPython is the operating language.")
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # S2 — ASK
    s = _blank(prs)
    draw_ask_page(
        s,
        f"2026 年，你為什麼要把 {COURSE_HOURS_LABEL}押在 Python 上？",
        data_card={
            "label": "2024–2025 語言使用率",
            "stat": "#1",
            "caption": "Python 連續兩年於 AI / ML / DS 使用率第一",
        },
    )
    add_source(s, "Stack Overflow Developer Survey 2025")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # S3 — CHART
    s = _blank(prs)
    add_title(s, "Python 於 AI / DS / DE 的使用率已從熱潮進入基礎設施鎖定")
    chart_png = charts.line_chart_s3()
    s.shapes.add_picture(str(chart_png), Inches(0.8), Inches(1.2),
                          width=Inches(11.7))
    add_source(s, "JetBrains/PSF Developer Ecosystem 2024 · Stack Overflow 2025 · Anaconda 2024")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # S4 — MATRIX 2×3
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=3, cells=[
        {"text": "資料操作語言", "sub": "pandas / NumPy / Polars\n表格、陣列、管線"},
        {"text": "ML/DL 建模語言", "sub": "scikit-learn / PyTorch / TensorFlow\n訓練、推論、評估"},
        {"text": "大資料橋接語言", "sub": "PySpark / Dask / Ray\n分散式、批次、串流"},
        {"text": "服務部署語言", "sub": "FastAPI / MLflow / Docker\nAPI、版本、容器"},
        {"text": "黏合膠水語言", "sub": "C / C++ / CUDA 後端呼叫\n高效能核心的 Python 介面"},
        {"text": "一份語法，五種現場。", "highlight": True},
    ], title="Python 的五重身份：一個語言、五種用途、一份契約")
    add_source(s, "本課整理自 JetBrains Python Ecosystem 2024")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # S5 — GEOMETRIC: 5-node flow chain + 1 branch
    s = _blank(prs)
    draw_flow_chain(s, nodes=[
        {"label": "Jupyter", "caption": "互動環境"},
        {"label": "NumPy", "caption": "數值核心"},
        {"label": "pandas", "caption": "表格資料"},
        {"label": "scikit-learn", "caption": "傳統 ML"},
        {"label": "PyTorch", "caption": "深度學習", "highlight": True},
    ], title="生態鏈不是並列清單，是有依賴順序的單向流",
       y=3.4,
       branch={"from_index": 2, "label": "PySpark",
               "sub": "分散式延伸（本課僅概念預覽）", "above": False})
    add_source(s, "本課整理")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # S6 — BEFORE/AFTER
    s = _blank(prs)
    add_title(s, "不學 Python 的代價：職缺市場 18 個月的板塊移動")
    chart_png = charts.before_after_bars_s6()
    s.shapes.add_picture(str(chart_png), Inches(1.0), Inches(1.2),
                          width=Inches(11.3))
    add_source(s, "LinkedIn Talent Insights + 104 人力銀行 AI 職缺 n=2,140, 2023 Q2 vs 2025 Q4")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # S7 — PYRAMID 4-layer + cross cuts + inverted thesis
    s = _blank(prs)
    draw_pyramid_stack(s, layers=[
        {"name": "Infra", "caption": "Docker / K8s / 雲平台"},
        {"name": "Runtime", "caption": "FastAPI / ONNX / TorchServe"},
        {"name": "Code", "caption": "sklearn / PyTorch / HuggingFace"},
        {"name": "Data", "caption": "pandas / SQL / Spark"},
    ], cross_cuts=["治理", "評估"],
       thesis="AI product = data + code + runtime + infra；evaluation 與 governance 是貫穿的縱骨。",
       title="AI 系統不是「一個模型」，是四層 + 兩道橫切面")
    add_source(s, "本課改寫自 Google ML Ops whitepaper 2024")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # S8 — DUAL TRACK
    s = _blank(prs)
    draw_dual_track(
        s,
        track_a_label="資料 + AI 能力",
        track_b_label="軟體 + 系統能力",
        a_nodes=["資料素養", "pandas 管線", "EDA 與視覺化", "ML 評估", "DL 入門"],
        b_nodes=["Python 工程", "venv / uv", "Git + Notebook", "Lint / Test", "FastAPI 雛形"],
        bridges=[(1, 1, "W1"), (2, 2, "W2"), (4, 4, "W3")],
        terminal_label="可交付的\nAI 工程基礎",
        title="兩條能力主線並行推進，每三小時交會一次",
    )
    add_source(s, "本課程主幹設計 2025 版")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # S9 — TABLE 3-col
    s = _blank(prs)
    draw_editorial_table(s,
        header=["能力", "具體產出", "驗收標準"],
        rows=[
            ["資料清洗管線", "clean_pipeline.py + Jupyter notebook",
             "從 raw CSV 到可分析 DataFrame 全自動"],
            ["EDA 報告", "Markdown + 圖表資產",
             "能回答三個可驗證的業務假設"],
            ["ML 基線模型", "sklearn Pipeline + 評估指標",
             "F1 / AUC 附信賴區間"],
            ["環境與版本控管", "pyproject.toml + .venv + Git repo",
             "他人 clone + uv sync 可重現"],
            ["AI 系統直覺", "四層架構圖 + 技術選型理由",
             "白板上能解釋 data/code/runtime/infra"],
        ],
        col_widths=[1.0, 1.6, 1.6],
        title=f"{COURSE_HOURS_LABEL}下課時，你能交付這五件事",
    )
    add_textbox(s, T.MARGIN_X, Inches(5.6),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "這不是承諾成為專家；是承諾你有獨立往下走的能力。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL,
                align=PP_ALIGN.CENTER, bold=True)
    add_source(s, "課程驗收規範 v2025")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # S10 — GEOMETRIC grid 4×2 (對齊實際 M0-M7 課程結構，共 24h)
    s = _blank(prs)
    draw_grid(s, rows=2, cols=4, cells=[
        {"label": "M0 · 2h", "sub": "開場與全景"},
        {"label": "M1 · 3h", "sub": "Python 基礎與資料思維"},
        {"label": "M2 · 3h", "sub": "OOP 與程式抽象"},
        {"label": "M3 · 4h", "sub": "NumPy 與 pandas", "highlight": True,
         "note": "里程碑 1：資料管線可交付"},
        {"label": "M4 · 4h", "sub": "EDA、視覺化與統計直覺"},
        {"label": "M5 · 3h", "sub": "進階 Python", "highlight": True,
         "note": "里程碑 2：可重現系統可交付"},
        {"label": "M6 · 3h", "sub": "計算機組織與 OS"},
        {"label": "M7 · 2h", "sub": "ML / DL / 學習路徑"},
    ], title=f"{COURSE_HOURS_LABEL}切成八塊，每塊有入口、有出口、有里程碑",
       caption="基礎期 M0–M3 · 12h  →  擴展期 M4–M5 · 7h  →  路徑期 M6–M7 · 5h  （總計 24 小時）")
    add_source(s, "課程大綱 v2026")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # S11 — MATRIX 2x2 quadrants
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "Data Engineer 路徑",
         "sub": "傳產 × 工程\n延伸：Airflow / Spark / dbt / cloud"},
        {"text": "ML / AI Engineer 路徑（本課後段主線）",
         "sub": "科技 × 工程\n延伸：PyTorch / MLOps / 分散式訓練",
         "highlight": True},
        {"text": "BI / Data Analyst 路徑",
         "sub": "傳產 × 分析\n延伸：SQL / dbt / Tableau / Power BI"},
        {"text": "Data Scientist 路徑",
         "sub": "科技 × 分析\n延伸：統計 / 實驗設計 / MLflow"},
    ], title="學完本課後的四條岔路：你該走哪一條？")
    # Axis labels
    add_textbox(s, T.MARGIN_X, Inches(6.4),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
                "縱軸：個人偏好 工程系統 ↔ 資料分析   ·   橫軸：產業類型 傳產 ↔ AI/科技",
                font_size=T.FONT_SOURCE, color=T.GRAY_MID,
                align=PP_ALIGN.CENTER)
    add_source(s, "本課整理自 2025 AI/DS 職能框架")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # S12 — RISK-MITIGATION
    s = _blank(prs)
    draw_risk_mitigation(s,
        risks=[
            "版本漂移：pandas 2.0 / NumPy 2.0 / PyTorch 2.0 同時升級",
            "生態爆炸：教程停在 2019 年，語法仍對但工法已變",
            "表面學會：能跑不等於能活",
        ],
        mitigations=[
            "鎖版本：pyproject.toml 強制課程級鎖定",
            "抓 changelog：每模組附「這一版變了什麼」",
            "工作坊占比 ≥ 30%：跑 → 活 → 長久",
        ],
        title="學 Python 的兩道真正風險，與對應的課程設計",
        summary="風險不靠口號解；靠版本、資源、練習量三件事。",
    )
    add_source(s, "本課風險盤點 2025")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # S13 — VS
    s = _blank(prs)
    draw_vs_two_col(s,
        left_title="語法學會者",
        right_title="系統交付者",
        left_items=[
            "能寫：for / def / import",
            "停在：單檔 < 200 行",
            "交付：notebook 片段",
            "市場位：初階資料分析助理",
            "參考薪資：NT$ 45–60k / 月（n=340）",
        ],
        right_items=[
            "能寫：Pipeline / Package / API",
            "活過：50+ 檔案多模組",
            "交付：可部署的服務",
            "市場位：ML / Data Engineer",
            "參考薪資：NT$ 150–300k / 月（n=210）",
        ],
        title="語法學會者 vs 系統交付者：同樣「會 Python」，薪資差 3–5 倍",
        summary="同一個工具，兩個結局。關鍵不在會 Python，在會不會把 Python 用到系統層。",
        delta="薪資差 3–5×",
    )
    add_source(s, "104 人力銀行 + CakeResume 2025 Q1 抽樣薪資範圍")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # S14 — PHOTO triptych via image_placeholder primitive (SOP §7)
    s = _blank(prs)
    add_title(s, "本課程的真實棲地：你終將打開這三份官方文件")
    s14_slots = [
        ("M0_S14_pandas", "pandas 官方文件",
         "DataFrame.groupby 首屏截圖（瀏覽器含 URL 列）",
         "https://pandas.pydata.org/docs/reference/groupby.html",
         "1200×1680 px", "使用層"),
        ("M0_S14_sklearn", "scikit-learn Pipeline",
         "sklearn.pipeline.Pipeline 類別文件首屏",
         "https://scikit-learn.org/stable/modules/generated/sklearn.pipeline.Pipeline.html",
         "1200×1680 px", "核心 API"),
        ("M0_S14_pytorch", "PyTorch nn.Module",
         "torch.nn.Module 文件首屏",
         "https://pytorch.org/docs/stable/generated/torch.nn.Module.html",
         "1200×1680 px", "概念預覽"),
    ]
    triptych_top = Inches(1.3)
    triptych_h = Inches(3.6)
    gap = Inches(0.2)
    total_w = T.SLIDE_W - 2 * T.MARGIN_X
    col_w = (total_w - gap * (len(s14_slots) - 1)) / len(s14_slots)
    for i, (pid, name, desc, url, size, depth) in enumerate(s14_slots):
        x = T.MARGIN_X + i * (col_w + gap)
        draw_image_placeholder(
            s, x, triptych_top, col_w, triptych_h,
            slot_name=name, description=desc, url_hint=url,
            size_hint=size, placeholder_id=pid,
            registry=image_registry,
        )
        # Depth tag underneath each placeholder
        add_textbox(s, x, triptych_top + triptych_h + Inches(0.15),
                    col_w, Inches(0.3),
                    f"本課抵達深度：{depth}",
                    font_size=T.FONT_SMALL, color=T.PRIMARY, bold=True,
                    align=PP_ALIGN.CENTER)
    # Thesis
    add_textbox(s, T.MARGIN_X, Inches(5.5),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.7),
                "專業能力 = 讀得懂官方文件 + 能複製貼上後改到合適 + 知道為什麼要改。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER, line_spacing=1.3)
    add_source(s, "pandas.pydata.org / scikit-learn.org / pytorch.org, 2025 Q1")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # S15 — SILENT closing
    s = _blank(prs)
    draw_silent_page(s, "你不是在學一個工具；\n你在取得一張通往 AI 現場的作業權。")
    add_footer(s, MODULE_CODE, 15, N_CONTENT, dark_bg=True)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
