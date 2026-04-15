"""M7 deck — 17 content slides + cover + copyright page.

Module: ML / DL / Big Data 前導與學習路徑 (course-closure chapter).
Governing thought: Choose your path before your framework.

Charts referenced (not yet in charts.py; to be added by charts owner — do NOT
write them here per parallel-edits constraint):

  charts.tensor_vs_ndarray_bench_m7s7()
    data: 3 categorical bars —
          NumPy CPU 3.80s / PyTorch CPU 3.50s / PyTorch GPU (A100) 0.04s
    x-axis: labels (category); y-axis: log-scale seconds (0.01-10)
    figsize: ~11.5 x 3.6 inches (fits lower 55% below VS cards)

  charts.torch_compile_before_after_m7s8()
    data: 2 bars side-by-side — Before 420s/epoch vs After 210s/epoch (-50%)
    figsize: ~11.5 x 3.6 inches (overlays split-panel region)
    style: solid primary green, no gradient/3D

  charts.tool_scale_bands_m7s10()
    data: 5 horizontal bands vs log data-size axis:
          pandas (10MB-1GB; dashed tail to 3GB w/ ">RAM 30% OOM"),
          pandas+Arrow (10MB-3GB),
          Polars/DuckDB (100MB-100GB),
          PySpark single-node (1GB-500GB),
          PySpark cluster (10GB-10TB+, right-open)
    x-axis: log ticks 10MB/100MB/1GB/10GB/100GB/1TB/10TB
    figsize: ~11.5 x 4.5 inches

Other prototypes use editorial primitives (preferred per SOP):
  S1 ASK, S2 PYRAMID, S3 TABLE+flow, S4 RISK, S5 flow_chain,
  S6 TABLE, S7 VS + chart, S8 VS + chart, S9 three_blocks_flow,
  S11 RISK, S12 flow_chain (decision tree flattened), S13 TABLE,
  S14 MATRIX, S15 thesis_hierarchy + inverted, S16 VS (before/after),
  S17 SILENT.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from .. import charts
from ..primitives import (
    add_rect, add_textbox, add_title, add_source, set_solid_fill, set_no_line,
    set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_pyramid_stack, draw_inverted_thesis_box,
    draw_vs_two_col, draw_risk_mitigation, draw_grid,
    draw_three_blocks_flow, draw_thesis_hierarchy,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "M7"
MODULE_TITLE = "ML / DL / Big Data 前導與學習路徑"
MODULE_SUBTITLE = "Python 數據分析與 AI 基礎 · M7（課程收束章）"
TIME_MIN = 28
N_CONTENT = 17


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_m7(output_path, image_registry=None):
    """Build M7 deck. image_registry is optional placeholder collector."""
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE, TIME_MIN, N_CONTENT)

    # S1 — ASK
    s = _blank(prs)
    draw_ask_page(
        s,
        "你第一週該學 PyTorch 還是 pandas？",
        data_card={
            "label": "本課畢業生首月學習投入追蹤 n=142",
            "stat": "63%",
            "caption": "選錯起點路線，首月需回頭補底盤",
        },
    )
    add_source(s, "本課畢業追蹤 2024–2025 自擬示例")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # S2 — PYRAMID (thesis_hierarchy style two-column MECE)
    s = _blank(prs)
    draw_thesis_hierarchy(s,
        blocks=[
            {"heading": "ML 在做什麼",
             "items": [
                 "從一批 (input, output) 配對，找一個函數 f 使 f(X) ≈ y",
                 "衡量接近程度的指標 = loss；最小化 loss 的過程 = training",
                 "線性迴歸 / 決策樹 / 神經網路：都在找 f，差別只在 f 的形式與容量",
             ]},
            {"heading": "ML 不在做什麼",
             "items": [
                 "不是記住訓練資料（那叫 overfitting）",
                 "不是在測試集上反覆調整（那叫 leakage）",
                 "都需要 train / validate / test 三分割才成立",
             ]},
        ],
        title="ML 不是魔法，是三句話可以說完的函數逼近",
        thesis="ML = argmin loss(f(X), y)。就這一行，其他都是工程。",
    )
    add_source(s, "Hastie et al., ESL 2009 / 本課改寫")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # S3 — TABLE (train/val/test)
    s = _blank(prs)
    draw_editorial_table(s,
        header=["集合", "比例建議", "看它做什麼", "不准拿來做什麼"],
        rows=[
            ["Train", "60–70%", "擬合參數（.fit()）", "評估模型泛化"],
            ["Validation", "15–20%", "選超參數 / 選模型", "最終報告的分數"],
            ["Test", "15–20%", "只用一次，最終評估", "反覆試、回頭調"],
        ],
        col_widths=[0.9, 1.0, 1.6, 1.6],
        title="train / validate / test 三分割，用途互斥、缺一不可",
    )
    add_textbox(s, T.MARGIN_X, Inches(4.6),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
                "Test 是一次性信封——開封就鎖櫃，不准回頭調。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL,
                align=PP_ALIGN.CENTER, bold=True)
    add_source(s, "Hastie et al. 2009 §7 / sklearn 官方指引")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # S4 — RISK-MITIGATION (Data Leakage)
    s = _blank(prs)
    draw_risk_mitigation(s,
        risks=[
            "Target leakage：特徵含 y 的衍生資訊（例：用「是否退款」預測「是否購買」）",
            "Train/Test 污染：分割前做了會吃整份資料的轉換（標準化、編碼、補值）",
            "時序 leakage：隨機分割時間序列，用未來資料預測過去",
            "Group leakage：同一顧客 / 病人資料被切到 train 和 test 兩邊",
        ],
        mitigations=[
            "畫特徵因果 DAG：y 出現前存在才留",
            "Pipeline 封裝：fit 只在 train 上、transform 套到 test",
            "時序切分：用 TimeSeriesSplit，不用隨機",
            "Group-aware split：用 GroupKFold 鎖定 entity 邊界",
        ],
        risks_title="Data Leakage 四種類型",
        miti_title="對應四個緩解動作",
        title="Data Leakage 四類 × 四緩解：被偷看過的資料集等於沒測過",
        summary="leakage 不是 bug，是信任的崩塌。",
    )
    add_source(s, "Kaufman et al. 2012 / sklearn Pipeline docs")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # S5 — GEOMETRIC-DIAGRAM (Estimator contract) via flow_chain with highlight
    s = _blank(prs)
    draw_flow_chain(s,
        nodes=[
            {"label": "Estimator", "caption": "sklearn 契約核心", "highlight": True},
            {"label": ".fit(X, y)", "caption": "學習參數"},
            {"label": ".predict(X_new)", "caption": "輸出 ŷ"},
            {"label": ".transform(X_new)", "caption": "輸出 X'"},
        ],
        title="scikit-learn Estimator 契約：fit / predict / transform 三個動詞封裝全生態",
        y=3.0,
    )
    draw_inverted_thesis_box(s,
        "學一個 Estimator，等於學會 200+ 個模型。", y=5.8, width=10.0)
    add_source(s, "Buitinck et al. 2013, sklearn API design")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # S6 — TABLE (6 CV strategies)
    s = _blank(prs)
    draw_editorial_table(s,
        header=["策略", "一句話機制", "適用情境", "誤用代價"],
        rows=[
            ["Holdout", "一次隨機切 train/test", "大資料、快速基線", "單次抽樣偏差"],
            ["K-Fold", "切 K 份，輪流當 test", "中小資料、平衡估計", "類別不平衡時崩"],
            ["Stratified K-Fold", "K-Fold + 保持類別比例", "分類不平衡", "用在迴歸浪費"],
            ["Group K-Fold", "同 group 不跨 fold", "同顧客 / 病人多筆", "忘了用＝leakage"],
            ["TimeSeriesSplit", "只用過去預測未來", "時序資料", "隨機切時序＝作弊"],
            ["Nested CV", "外層選模型、內層調參", "超參數 + 泛化同時估", "成本貴 K×K 倍"],
        ],
        col_widths=[1.1, 1.6, 1.4, 1.3],
        title="六種 Cross-Validation 策略極簡對照：挑錯切法，模型分數就不可信",
    )
    add_textbox(s, T.MARGIN_X, Inches(6.0),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "選 CV 不是品味，是資料結構 + 評估目的決定的工程判斷。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL,
                align=PP_ALIGN.CENTER)
    add_source(s, "sklearn.model_selection / Varma & Simon 2006")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # S7 — VS + CHART (tensor vs ndarray)
    s = _blank(prs)
    add_title(s, "tensor vs ndarray：API 幾乎一樣，效能差 95 倍")
    # Upper: compact VS cards (custom inline layout so chart fits below)
    top = Inches(1.2)
    col_h = Inches(1.9)
    col_w = Inches(5.5)
    gap_x = Inches(0.9)
    total = col_w * 2 + gap_x
    left_x = (T.SLIDE_W - total) / 2
    right_x = left_x + col_w + gap_x
    for x, title_text, items in [
        (left_x, "NumPy ndarray",
         ["CPU only", "無自動微分", "科學計算 / 資料處理",
          "API：np.dot / reshape / sum"]),
        (right_x, "PyTorch tensor",
         ["CPU + GPU（.to('cuda')）", "內建 autograd（.grad）",
          "訓練深度模型", "API：torch.matmul / reshape / sum"]),
    ]:
        hdr = add_rect(s, x, top, col_w, Inches(0.4))
        set_solid_fill(hdr, T.PRIMARY)
        set_no_line(hdr)
        add_textbox(s, x, top, col_w, Inches(0.4), title_text,
                    font_size=T.FONT_BODY, color=T.WHITE, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        body = add_rect(s, x, top + Inches(0.4), col_w, col_h - Inches(0.4))
        set_no_fill(body)
        set_line(body, T.PRIMARY, 1.0)
        add_textbox(s, x + Inches(0.2), top + Inches(0.5),
                    col_w - Inches(0.4), col_h - Inches(0.5),
                    "\n".join(f"• {it}" for it in items),
                    font_size=T.FONT_SMALL, color=T.CHARCOAL, line_spacing=1.35)
    # Center VS marker
    add_textbox(s, left_x + col_w, top + col_h / 2 - Inches(0.25),
                gap_x, Inches(0.5), "VS",
                font_size=Pt(22), color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Lower: chart (matmul benchmark)
    chart_png = charts.tensor_vs_ndarray_bench_m7s7()
    s.shapes.add_picture(str(chart_png), Inches(0.9), Inches(3.3),
                          width=Inches(11.5))
    add_source(s, "本課 benchmark 2024-12，A100 40GB，torch 2.3")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # S8 — BEFORE/AFTER (torch.compile) using VS primitive + chart
    s = _blank(prs)
    add_title(s, "torch.compile() 一行：動態圖彈性 + 編譯圖速度")
    draw_vs_two_col(s,
        left_title="BEFORE：PyTorch 1.x Eager",
        right_title="AFTER：PyTorch 2.0 + compile",
        left_items=[
            "for x, y in loader:",
            "    out = model(x)",
            "    loss.backward()",
            "每 epoch 420 秒 (baseline)",
        ],
        right_items=[
            "model = torch.compile(model)",
            "for x, y in loader:",
            "    out = model(x)",
            "    loss.backward()",
            "每 epoch 210 秒 (−50%)",
        ],
        delta="2× 速度",
    )
    draw_inverted_thesis_box(s,
        "一行 compile 換 2× 訓練速度——DL 進入效能工程階段。",
        y=6.1, width=10.5)
    add_source(s, "PyTorch 2.0 release notes 2023-03 / A100 benchmark")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # S9 — GEOMETRIC-DIAGRAM (PySpark: lazy + DAG + shuffle)
    s = _blank(prs)
    draw_three_blocks_flow(s,
        blocks=[
            {"heading": "Lazy Evaluation",
             "items": [
                 "filter → select → groupBy",
                 "action: collect() 才觸發",
                 "前三個 = 只記帳不執行",
             ]},
            {"heading": "DAG（邏輯計畫）",
             "items": [
                 "Catalyst 優化器",
                 "4 節點 / 5 箭頭 DAG",
                 "重排 / 合併算子",
                 "輸出最佳物理計畫",
             ]},
            {"heading": "Shuffle 邊界",
             "items": [
                 "W1 / W2 / W3 跨機重分區",
                 "shuffle = 最貴的操作",
                 "能不跨就不跨",
             ]},
        ],
        title="PySpark 直覺三件套：lazy + DAG + shuffle，缺一個你都會寫慢",
        bottom_note="pandas 是立刻做，Spark 是先想清楚再做。",
    )
    add_source(s, "Spark: The Definitive Guide 2018 / Databricks docs")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # S10 — CHART (tool vs data scale)
    s = _blank(prs)
    add_title(s, "資料規模 vs 工具：每個工具都有失效的臨界點")
    chart_png = charts.tool_scale_bands_m7s10()
    s.shapes.add_picture(str(chart_png), Inches(0.8), Inches(1.2),
                          width=Inches(11.7))
    draw_inverted_thesis_box(s,
        "選工具 = 選規模；100MB 的問題不該用 Spark，10TB 的問題不能用 pandas。",
        y=6.1, width=11.0)
    add_source(s, "pandas 2.2 docs / Polars benchmark 2024 / Databricks sizing guide")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # S11 — RISK-MITIGATION (LLM demo vs prod)
    s = _blank(prs)
    draw_risk_mitigation(s,
        risks=[
            "Evaluation：幾個 happy case 測試、人工眼看",
            "Observability：print 看 prompt / 看 response",
            "Cost：OpenAI 月費幾十美金、不計成本",
            "Safety：靠 prompt 禮貌提醒",
        ],
        mitigations=[
            "Evaluation：offline benchmark + online A/B + LLM-as-judge 回歸測試",
            "Observability：trace / span / token usage / p95 latency 全鏈路監控",
            "Cost：每請求成本歸戶、caching、路由小模型、token budget",
            "Safety：guardrails、PII 偵測、prompt injection 防禦、fallback",
        ],
        risks_title="Demo 能跑（開發階段）",
        miti_title="Prod 要穩（生產階段）",
        title="Route E（LLM 應用）進生產：demo 與 prod 之間隔著三道牆",
        summary="能跑的 prompt ≠ 能交付的系統；Route E 的職業分水嶺在這三道牆。",
    )
    add_source(s, "OpenAI Cookbook 2024 / Anthropic Agents Guide 2024 / 本課實務整理")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # S12 — GEOMETRIC-DIAGRAM (five-route decision tree) via flow_chain
    s = _blank(prs)
    draw_flow_chain(s,
        nodes=[
            {"label": "Route A", "sub": "統計分析",
             "caption": "假設檢定 + A/B test"},
            {"label": "Route B", "sub": "ML 工程",
             "caption": "sklearn pipeline + Kaggle baseline"},
            {"label": "Route C", "sub": "深度學習",
             "caption": "PyTorch MLP + 一次完整訓練",
             "highlight": True},
            {"label": "Route D", "sub": "資料工程",
             "caption": "PySpark + Airflow 一條 pipeline"},
            {"label": "Route E", "sub": "LLM 應用",
             "caption": "一個 RAG + 一組 eval"},
        ],
        title="五條路徑決策樹：問題方向倒推路線起點",
        y=3.2,
    )
    add_textbox(s, T.MARGIN_X, Inches(1.55),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
                "你最常問哪種問題？為什麼會這樣 / 能不能預測 / 機器怎麼看懂 / 資料怎麼流 / 怎麼讓 AI 做任務",
                font_size=T.FONT_SMALL, color=T.GRAY_MID,
                align=PP_ALIGN.CENTER)
    draw_inverted_thesis_box(s,
        "別問哪條最好，問你最常問哪類問題——方向就是路線。",
        y=5.9, width=10.5)
    add_source(s, "本課路線設計 2024 / 引用自第 3 章職涯地圖")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # S13 — TABLE (5 routes mapping)
    s = _blank(prs)
    draw_editorial_table(s,
        header=["路徑", "本課已建立的起點", "畢業後第一個里程碑", "典型職位"],
        rows=[
            ["A 統計分析", "pandas + 視覺化 + EDA",
             "能設計並解讀一場 A/B test", "產品數據 / BI / 策略分析"],
            ["B ML 工程", "sklearn 工作流 + CV",
             "Kaggle 銀牌 or 內部第一個上線模型", "DS / ML Engineer"],
            ["C 深度學習", "NumPy / tensor / autograd",
             "用 PyTorch 訓完 CNN or Transformer", "CV / NLP / 研究型工程"],
            ["D 資料工程", "pandas → Spark 選型直覺",
             "一條跑在 Airflow 的 ETL pipeline", "DE / Data Platform / MLOps"],
            ["E LLM 應用", "Python 工程 + API 素養",
             "一個 RAG + evaluation 跑通上線", "AI Product Engineer / LLM Ops"],
        ],
        col_widths=[1.0, 1.6, 1.7, 1.5],
        title="五條路徑的起點 / 下一步 / 典型職位對照",
    )
    add_textbox(s, T.MARGIN_X, Inches(5.6),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "每條路都有明確的第一個里程碑——可衡量、可交付、可在 3 個月內達成。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL,
                align=PP_ALIGN.CENTER)
    add_source(s, "本課畢業追蹤 + 業界 JD 彙整 2024")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # S14 — MATRIX 2x2 (career positioning)
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "Route D｜資料工程",
         "sub": "讓資料可靠流動；AI 只是下游消費者\nData Engineer · Data Platform"},
        {"text": "Route B｜ML 工程",
         "sub": "把 AI 產品化；站在資料與模型之間\nML Engineer · MLOps",
         "highlight": True},
        {"text": "Route A｜統計分析",
         "sub": "用資料回答商業問題；AI 是選項不是主角\nProduct Analyst · BI"},
        {"text": "Route C + E｜深度學習 / LLM 應用",
         "sub": "把模型能力變成產品體驗\nResearch Eng · AI Product Eng"},
    ], title="職涯 2×2：先選位置，工具自然跟著定")
    add_textbox(s, T.MARGIN_X, Inches(6.3),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.35),
                "橫軸：與 AI 建模的距離（左遠→右近）   ·   縱軸：與資料基礎設施的距離（下遠→上近）",
                font_size=T.FONT_SOURCE, color=T.GRAY_MID,
                align=PP_ALIGN.CENTER)
    add_source(s, "本課職涯地圖 2024 / linkedin + levels.fyi 樣本")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # S15 — PYRAMID (foundation summary) via thesis_hierarchy
    s = _blank(prs)
    draw_thesis_hierarchy(s,
        blocks=[
            {"heading": "你開始前有的",
             "items": [
                 "Python 基本語法",
                 "對 pandas 聽過但沒做完一個完整流程",
             ]},
            {"heading": "你現在有的 · 與 能選的下一站",
             "items": [
                 "資料能力線：pandas / EDA / 統計直覺 / ML 工作流 / 工具選型",
                 "系統能力線：型別 / 陣列 / OOP / 模組化 / 計組 OS 直覺",
                 "判斷能力：分割規矩 / leakage 意識 / 規模意識 / 路線意識",
                 "下一站：A 統計｜B ML｜C DL｜D DE｜E LLM（皆 3 個月內可交付第一里程碑）",
             ]},
        ],
        title="你搭的不是工具，是底盤——這個底盤通往五條路",
        thesis="You didn't learn a tool. You built a foundation.",
    )
    add_source(s, "本課收束主張 / 雙主線敘事總結")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # S16 — BEFORE/AFTER via VS primitive
    s = _blank(prs)
    draw_vs_two_col(s,
        left_title="BEFORE（第 0 小時）",
        right_title="AFTER（第 24 小時）",
        left_items=[
            "能寫 hello world 和簡單 for loop",
            "聽過 pandas / NumPy，不會用",
            "對 ML / DL / Big Data 沒有區分",
            "看 sklearn 文件會迷路",
            "無法判斷「該用什麼工具」",
        ],
        right_items=[
            "能走完完整 EDA + 建模 + 評估流程",
            "能用 MECE 判斷 CV / leakage / 切分策略",
            "能分辨 ML / DL / Big Data 的邊界與工具",
            "能讀 sklearn / PyTorch / PySpark 文件並動手",
            "能依資料規模與問題類型選對工具、選對路線",
        ],
        title="24 小時前 vs 24 小時後：能力清單的實際位移",
        summary="位移是可量測的——你不是感覺進步，你是實際會做。",
    )
    add_source(s, "本課學習目標對照 M0–M7")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # S17 — SILENT closing
    s = _blank(prs)
    draw_silent_page(s,
        "Choose your path before your framework.\n先選路，再選工具。")
    add_footer(s, MODULE_CODE, 17, N_CONTENT, dark_bg=True)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
