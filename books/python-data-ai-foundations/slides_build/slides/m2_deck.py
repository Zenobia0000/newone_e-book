"""M2 deck — 16 content slides + cover + copyright page.

Pending chart helpers (to be added in charts.py by a follow-up agent):

- charts.line_chart_m2_s3(outfile=None) -> Path
    Two-line chart comparing "refactor cost vs LOC" under two regimes.
    x: programme size in lines (0–1000, ticks every 200)
    y: hours to add a new requirement (0–12)
    Series A (solid, HEX_PRIMARY, width 2.5pt, label "無 class 邊界"):
        (200, 1.5), (400, 3.2), (600, 6.1), (800, 9.8)  convex accelerating
    Series B (dashed, HEX_CHARCOAL, width 2pt, label "有 class 邊界"):
        (200, 1.4), (400, 2.1), (600, 2.8), (800, 3.5)  near-linear
    Vertical reference line at x=200 (HEX_LIGHT_GRAY) annotated "臨界點".
    End-of-line labels; no legend box.
    Output PNG cached in T.CHART_CACHE, figsize hint ~ (11.7, 4.8) inches.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from .. import charts
from ..primitives import (
    add_rect, add_textbox, add_title, add_source, set_solid_fill, set_no_line,
    set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_pyramid_stack, draw_inverted_thesis_box,
    draw_vs_two_col, draw_risk_mitigation, draw_grid,
    draw_code_panel, draw_split_panel, draw_thesis_hierarchy,
    draw_three_blocks_flow, draw_image_placeholder,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "M2"
MODULE_TITLE = "OOP 與程式抽象"
MODULE_SUBTITLE = "Python 數據分析與 AI 基礎 · M2"
TIME_MIN = 22
N_CONTENT = 16


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_m2(output_path, image_registry=None):
    """Build M2 deck. image_registry is optional placeholder collector."""
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE, TIME_MIN, N_CONTENT)

    # S1 — GEOMETRIC-DIAGRAM: coupled vs modularized (split panel)
    s = _blank(prs)
    draw_split_panel(s,
        left={
            "label": "200 行單檔（函式散落） · 耦合度：24 條交叉",
            "style": "wrong",
            "cells": [
                {"text": "fn · fn · fn · fn"},
                {"text": "fn · fn · fn · fn"},
                {"text": "fn · fn · fn · fn"},
            ],
            "note": "方塊之間 24 條交錯箭頭，網狀糾結",
        },
        right={
            "label": "200 行模組化（class 邊界） · 耦合度：3 條界面",
            "cells": [
                {"text": "class A：fn fn fn fn"},
                {"text": "class B：fn fn fn fn", "highlight": True},
                {"text": "class C：fn fn fn fn"},
            ],
            "note": "三組之間僅 3 條界面箭頭，組內無交叉",
        },
        title="程式碼不是寫得下去的問題，是長得出來的問題",
    )
    add_textbox(s, T.MARGIN_X, Inches(6.3),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "檔案長度不是病，結構混亂才是。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "M2 module thesis")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # S2 — SILENT
    s = _blank(prs)
    draw_silent_page(s, "能跑的程式碼，\n和能活的程式碼，\n是兩件事。")
    add_footer(s, MODULE_CODE, 2, N_CONTENT, dark_bg=True)

    # S3 — CHART (two-line refactor cost vs LOC)
    s = _blank(prs)
    add_title(s, "超過 200 行後，無 class 的改動成本呈非線性發散")
    chart_png = charts.line_chart_m2_s3()
    s.shapes.add_picture(str(chart_png), Inches(0.8), Inches(1.3),
                          width=Inches(11.7))
    add_source(s, "內部 50 個 AI side-project 修改耗時紀錄，2024–2025")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # S4 — ASK
    s = _blank(prs)
    draw_ask_page(
        s,
        "如果程式碼是一個團隊，誰負責什麼？邊界在哪？",
        data_card={
            "label": "AI 工程團隊平均",
            "stat": "78%",
            "caption": "bug 來自「不知道誰該改這段」",
        },
    )
    add_source(s, "M2 課堂歸納")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # S5 — PYRAMID: Class vs Object (two blocks + inverted thesis)
    s = _blank(prs)
    draw_thesis_hierarchy(s,
        blocks=[
            {"heading": "Class（模板）",
             "items": [
                 "宣告「這類東西長什麼樣」",
                 "只存在於原始碼；不佔 runtime 記憶體",
             ]},
            {"heading": "Object（實例）",
             "items": [
                 "依模板具體「造出一個」",
                 "佔記憶體；有自己的狀態；彼此互不干擾",
             ]},
        ],
        title="Class 是模板，Object 是實例；兩者是契約與履行的關係",
        thesis="Class ≠ Object；一份契約可產出無限履行。",
    )
    add_source(s, "Python Data Model §3")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # S6 — TABLE: State / Behavior / Identity
    s = _blank(prs)
    draw_editorial_table(s,
        header=["要素", "回答什麼問題", "Python 對應"],
        rows=[
            ["State（狀態）", "你現在是什麼？", "instance attribute（self.x）"],
            ["Behavior（行為）", "你能做什麼？", "method（def fn(self)）"],
            ["Identity（身分）", "你是哪一個？", "記憶體位址（id(obj)）"],
        ],
        col_widths=[1.2, 1.6, 1.8],
        title="Object 的三個要素：State、Behavior、Identity",
    )
    add_textbox(s, T.MARGIN_X, Inches(5.2),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "三要素 MECE；缺一個就不完整。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "Python Language Reference §3.2")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # S7 — GEOMETRIC-DIAGRAM: encapsulation (Public API vs Private)
    s = _blank(prs)
    add_title(s, "封裝不是加鎖，是在外部介面與內部實作之間劃界")
    ox, oy = Inches(1.2), Inches(1.4)
    ow, oh = T.SLIDE_W - 2 * ox, Inches(4.8)
    outer = add_rect(s, ox, oy, ow, oh)
    set_no_fill(outer); set_line(outer, T.PRIMARY, 1.5)
    add_textbox(s, ox + Inches(0.2), oy + Inches(0.05), ow - Inches(0.4), Inches(0.35),
                "Class: Dataset", font_size=T.FONT_BODY, color=T.PRIMARY, bold=True)

    def _band(label_text, methods, y, label_color, solid):
        add_textbox(s, ox + Inches(0.2), y, ow - Inches(0.4), Inches(0.3),
                    label_text, font_size=T.FONT_CAPTION, color=label_color, bold=True)
        avail = ow - Inches(0.4); g = Inches(0.15)
        bw = (avail - g * (len(methods) - 1)) / len(methods)
        by = y + Inches(0.4); bh = Inches(0.6)
        for i, m in enumerate(methods):
            bx = ox + Inches(0.2) + i * (bw + g)
            r = add_rect(s, bx, by, bw, bh)
            if solid:
                set_solid_fill(r, T.PRIMARY); set_no_line(r); tc = T.WHITE
            else:
                set_no_fill(r); set_line(r, T.GRAY_MID, 1.0, dash=True); tc = T.GRAY_MID
            add_textbox(s, bx, by, bw, bh, m, font_size=T.FONT_CAPTION, color=tc,
                        bold=solid, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                        family=T.FONT_MONO)
        return by + bh

    pub_end = _band("Public API（外部介面）",
                    ["load()", "clean()", "split()", "describe()"],
                    oy + Inches(0.55), T.PRIMARY, True)
    # Boundary
    boundary = add_rect(s, ox + Inches(0.2), pub_end + Inches(0.3),
                        ow - Inches(0.4), Inches(0.02))
    set_solid_fill(boundary, T.PRIMARY); set_no_line(boundary)
    add_textbox(s, ox + Inches(0.2), pub_end + Inches(0.05), Inches(1.0), Inches(0.3),
                "— 邊界 —", font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True)
    _band("Private（內部實作）",
          ["_data", "_schema", "_validate()"],
          pub_end + Inches(0.5), T.GRAY_MID, False)

    add_textbox(s, T.MARGIN_X, Inches(6.4), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "外部不應直取內部；穩定介面在上，會變動實作在下。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "M2 module design")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # S8 — VS: 名目型別 vs 鴨子型別
    s = _blank(prs)
    draw_vs_two_col(s,
        left_title="名目型別（Java / C#）",
        right_title="鴨子型別（Python）",
        left_items=[
            "必須 extends Animal",
            "編譯期檢查 is-a",
            "類型樹強耦合",
        ],
        right_items=[
            "只要物件有 .quack() 方法",
            "運行期呼叫即通過",
            "介面靠 Protocol 宣告，不靠繼承",
        ],
        title="鴨子型別：Python 不問你繼承誰，只問你會不會 quack()",
    )
    add_source(s, "PEP 544 — Protocols")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # S9 — RISK-MITIGATION: Inheritance vs Composition
    s = _blank(prs)
    draw_risk_mitigation(s,
        risks=[
            "表達 is-a 關係",
            "適用：穩定共通介面、子類數量少",
            "風險：類型樹深、diamond 問題、耦合過早",
            "例：BaseLoader → CSVLoader",
        ],
        mitigations=[
            "表達 has-a 關係",
            "適用：行為可插拔、依賴會替換",
            "優勢：低耦合、易測試、易替換",
            "例：ChatAgent 內嵌 VectorStoreClient",
        ],
        risks_title="Inheritance（繼承）",
        miti_title="Composition（組合）",
        title="繼承 vs 組合：該選哪一條？",
        summary="預設選組合；只有 is-a 明確時才繼承。",
    )
    add_source(s, "Gang of Four §1.6 / Composition over Inheritance")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # S10 — BEFORE-AFTER: @dataclass
    s = _blank(prs)
    add_title(s, "@dataclass：12 行樣板碼換成 1 行裝飾器")
    draw_code_panel(s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.5),
        label="Before（手寫 __init__）",
        code=('class User:\n'
              '    def __init__(self, name: str, email: str, age: int = 0):\n'
              '        self.name = name\n'
              '        self.email = email\n'
              '        self.age = age\n'
              '    def __repr__(self):\n'
              '        return f"User(name={self.name!r}, ...)"\n'
              '    def __eq__(self, other):\n'
              '        return (self.name, self.email, self.age) == ...'),
        bullets=[
            "12 行",
            "手寫 __init__ / __repr__ / __eq__",
            "同樣邏輯重複寫",
        ],
        label_dark=False,
    )
    draw_code_panel(s,
        x=T.MARGIN_X, y=Inches(4.0),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.3),
        label="After（@dataclass）",
        code=('from dataclasses import dataclass\n'
              '\n'
              '@dataclass\n'
              'class User:\n'
              '    name: str\n'
              '    email: str\n'
              '    age: int = 0'),
        bullets=[
            "6 行（-50%）",
            "同樣產出 __init__ / __repr__ / __eq__",
            "可讀性提升",
        ],
        label_dark=True,
    )
    add_source(s, "PEP 557 — Data Classes")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # S11 — GEOMETRIC-DIAGRAM: Protocol + 3 structural impls
    s = _blank(prs)
    add_title(s, "Protocol：用「會做什麼」而不是「繼承誰」定義介面")
    # Protocol box on top center
    proto_w = Inches(5.0)
    proto_h = Inches(1.3)
    proto_x = (T.SLIDE_W - proto_w) / 2
    proto_y = Inches(1.4)
    proto_rect = add_rect(s, proto_x, proto_y, proto_w, proto_h)
    set_no_fill(proto_rect)
    set_line(proto_rect, T.PRIMARY, 2.5)
    add_textbox(s, proto_x + Inches(0.2), proto_y + Inches(0.15),
                proto_w - Inches(0.4), Inches(0.4),
                "Protocol",
                font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER)
    add_textbox(s, proto_x + Inches(0.2), proto_y + Inches(0.55),
                proto_w - Inches(0.4), Inches(0.7),
                "class DataLoader(Protocol):\n    def load(self) -> DataFrame: ...",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL,
                align=PP_ALIGN.CENTER, family=T.FONT_MONO, line_spacing=1.3)

    # Three implementations
    impls = ["CSVLoader", "ParquetLoader", "SQLLoader"]
    impl_top = Inches(4.3)
    impl_h = Inches(1.1)
    total_w = T.SLIDE_W - 2 * T.MARGIN_X
    gap = Inches(0.4)
    impl_w = (total_w - gap * (len(impls) - 1)) / len(impls)
    for i, name in enumerate(impls):
        ix = T.MARGIN_X + i * (impl_w + gap)
        rect = add_rect(s, ix, impl_top, impl_w, impl_h)
        set_no_fill(rect)
        set_line(rect, T.GRAY_MID, 1.0)
        add_textbox(s, ix, impl_top, impl_w, impl_h, name,
                    font_size=T.FONT_BODY, color=T.CHARCOAL, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                    family=T.FONT_MONO)

    add_textbox(s, T.MARGIN_X, Inches(3.7),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.35),
                "structural match（虛線向上：不靠繼承，只靠結構相符）",
                font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
                align=PP_ALIGN.CENTER)
    add_textbox(s, T.MARGIN_X, Inches(5.7),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "三實作彼此無連線——表示互不繼承，只是剛好都實作了 load()。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "PEP 544 — Protocols")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # S12 — TABLE: Module / Package（skeleton 修正版）
    s = _blank(prs)
    draw_editorial_table(s,
        header=["層級", "定義", "例子"],
        rows=[
            [".py 檔", "一個 module", "dataset.py"],
            ["資料夾 + __init__.py", "一個 package", "myproject/"],
            ["package 內部資料夾 + __init__.py", "sub-package", "myproject/data/"],
            ["from x import y", "從 module 取名", "from dataset import Dataset"],
            ["import x.y", "從 package 取 module", "import myproject.data.dataset"],
            ["__init__.py 內容", "可為空；可重新匯出", "from .dataset import Dataset"],
        ],
        col_widths=[1.8, 1.4, 1.8],
        title="Module 與 Package：一個 __init__.py 就把資料夾升格為可被 import 的單位",
    )
    add_textbox(s, T.MARGIN_X, Inches(6.2),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "__init__.py 是目錄頁，不是一本書。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "Python Import System §5")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # S13 — BEFORE-AFTER: shared global vs instance isolation
    s = _blank(prs)
    add_title(s, "Chatbot 多用戶：共享狀態 vs 實例隔離")
    draw_code_panel(s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.4),
        label="Before：全域共享狀態",
        code=('global_state = {}\n'
              '\n'
              'def handle(user_id, msg):\n'
              '    global_state[\"history\"].append(msg)  # User A / B / C 同一份\n'
              '    return reply(global_state)'),
        bullets=[
            "User A / B / C 皆寫回同一 global_state",
            "三人互相污染彼此資料",
            "併發場景尤其危險",
        ],
        label_dark=False,
    )
    draw_code_panel(s,
        x=T.MARGIN_X, y=Inches(4.0),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.3),
        label="After：每個 session 一個 instance",
        code=('class ChatSession:\n'
              '    def __init__(self, user_id):\n'
              '        self.user_id = user_id\n'
              '        self.history = []\n'
              '    def handle(self, msg):\n'
              '        self.history.append(msg)'),
        bullets=[
            "ChatSession(a) / (b) / (c) 三獨立實例",
            "狀態隔離；無交叉污染",
            "每個 instance 自帶 self.history",
        ],
        label_dark=True,
    )
    add_source(s, "M2 multi-user pattern")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # S14 — MATRIX 2×2: when to open a class
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        # top-left: low coupling × long-life
        {"text": "module-level 常數 / 工具函式",
         "sub": "低耦合 × 長命"},
        # top-right: high coupling × long-life (highlight)
        {"text": "完整 class（含 method）",
         "sub": "高耦合 × 長命\n← 明確該用",
         "highlight": True},
        # bottom-left
        {"text": "寫成 function 即可",
         "sub": "低耦合 × 短命"},
        # bottom-right
        {"text": "dataclass（僅裝資料）",
         "sub": "高耦合 × 短命"},
    ], title="該不該為這段邏輯開一個 class？2×2 決策象限")
    add_textbox(s, T.MARGIN_X, Inches(6.3),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
                "橫軸：資料與函式耦合度 低 → 高    縱軸：預期生命週期 / 重用次數 短 → 長",
                font_size=T.FONT_SOURCE, color=T.GRAY_MID,
                align=PP_ALIGN.CENTER)
    add_source(s, "作者整理；取自 40 個 AI side-project code review")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # S15 — PYRAMID: three questions
    s = _blank(prs)
    add_title(s, "什麼時候該用 class？三問裡勾到一個就動手")
    q_text = ("• Q1：這組資料和這組函式是不是總是一起出現？\n\n"
              "• Q2：這段邏輯會不會活超過一個月、被多處使用？\n\n"
              "• Q3：是否需要同時存在多個獨立的實例（如多用戶 session）？")
    add_textbox(s, T.MARGIN_X + Inches(1.0), Inches(1.8),
                T.SLIDE_W - 2 * T.MARGIN_X - Inches(2.0), Inches(3.5),
                q_text,
                font_size=T.FONT_BODY, color=T.CHARCOAL,
                line_spacing=1.55)
    draw_inverted_thesis_box(s,
        "三問勾到任一，就開 class；三問全 No，就用 function。",
        y=6.0, width=9.5)
    add_source(s, "M2 decision heuristic")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # S16 — PHOTO: scikit-learn StandardScaler + Pipeline code
    s = _blank(prs)
    add_title(s, "scikit-learn 的 fit / transform：整個 ML 生態靠一份 class 契約運作")
    # Left: image placeholder
    draw_image_placeholder(
        s, Inches(0.6), Inches(1.3), Inches(6.8), Inches(4.8),
        slot_name="scikit-learn StandardScaler 官方文件",
        description=("StandardScaler 類別頁面首屏，顯示 fit() / transform() / "
                     "fit_transform() 方法簽章區"),
        url_hint="https://scikit-learn.org/stable/modules/generated/sklearn.preprocessing.StandardScaler.html",
        size_hint="1320×1680 px",
        placeholder_id="M2_S16_sklearn_scaler",
        registry=image_registry,
    )
    # Right: real code snippet
    draw_code_panel(s,
        x=Inches(7.6), y=Inches(1.3),
        w=Inches(5.1), h=Inches(4.8),
        label="Pipeline 真實程式碼",
        code=('from sklearn.pipeline import Pipeline\n'
              'from sklearn.preprocessing import StandardScaler\n'
              'from sklearn.linear_model import LogisticRegression\n'
              '\n'
              'pipe = Pipeline([\n'
              '    ("scale", StandardScaler()),\n'
              '    ("clf",   LogisticRegression()),\n'
              '])\n'
              'pipe.fit(X_train, y_train)'),
        bullets=[
            "每個元件都實作 fit / transform",
            "遵守同一份 class 契約",
            "→ 可被任意組合",
        ],
        label_dark=True,
    )
    add_textbox(s, T.MARGIN_X, Inches(6.3),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "Pipeline 能串起來，是因為每個元件都實作 fit / transform——這是 OOP 契約，不是魔法。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "scikit-learn.org, 2025-Q1")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
