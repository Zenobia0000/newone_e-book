"""M1 deck — 16 content slides + cover + copyright page."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from .. import charts
from ..primitives import (
    add_rect, add_textbox, add_title, add_source, set_solid_fill, set_no_line,
    set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_pyramid_stack, draw_inverted_thesis_box,
    draw_vs_two_col, draw_risk_mitigation, draw_dual_track, draw_grid,
    draw_code_panel, draw_split_panel, draw_concentric_zones,
    draw_thesis_hierarchy, draw_three_blocks_flow,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "M1"
MODULE_TITLE = "Python 基礎與資料思維"
MODULE_SUBTITLE = "把資料變可信：語法是手段，可信是目的"
TIME_MIN = 28
N_CONTENT = 16


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_m1(output_path, image_registry=None):
    """Build M1 deck. image_registry is optional placeholder collector."""
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE, TIME_MIN, N_CONTENT)

    # S1 — CHART stacked bar
    s = _blank(prs)
    add_title(s, "資料分析師 90% 的時間不在畫圖，在擦地板")
    chart_png = charts.stacked_bar_m1s1()
    s.shapes.add_picture(str(chart_png), Inches(0.8), Inches(2.0),
                          width=Inches(11.7))
    add_textbox(s, T.MARGIN_X, Inches(5.0),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "畫圖表只佔 6%；90% 的工時在前三段。",
                font_size=T.FONT_BODY, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "CrowdFlower Data Science Report 2022; Anaconda State of Data Science 2023")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # S2 — SILENT
    s = _blank(prs)
    draw_silent_page(s, "你不是在學 Python，\n你是在學怎麼把資料變得可信。")
    add_footer(s, MODULE_CODE, 2, N_CONTENT, dark_bg=True)

    # S3 — GEOMETRIC split panel (mental model)
    s = _blank(prs)
    draw_split_panel(s,
        left={
            "label": "錯誤心智模型：盒子",
            "style": "wrong",
            "cells": [
                {"text": "a = [1, 2, 3]"},
                {"text": "b = a"},
                {"text": "兩個獨立盒子 ✗"},
            ],
            "note": "以為 b = a 是複製",
        },
        right={
            "label": "正確心智模型：名字綁定",
            "cells": [
                {"text": "a ─┐"},
                {"text": "b ─┴→  [1, 2, 3]", "highlight": True},
                {"text": "reference（同一物件）"},
            ],
            "note": "b = a 不是複製；是多一個名字指同一個物件。",
        },
        title="變數不是盒子，是繫在物件上的一條線",
    )
    add_source(s, "Python Language Reference §3.1 Objects, values and types")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # S4 — TABLE container comparison
    s = _blank(prs)
    draw_editorial_table(s,
        header=["容器", "可變", "有序", "可重複", "索引方式", "資料世界場景"],
        rows=[
            ["list", "✓", "✓", "✓", "整數索引", "一欄原始資料"],
            ["tuple", "✗", "✓", "✓", "整數索引", "df.shape 回傳 (列, 欄)"],
            ["dict", "✓", "✓ (Py3.7+)", "鍵唯一", "鍵索引", "一筆記錄 / JSON 對映"],
            ["set", "✓", "✗", "✗", "無索引", "去重 / 成員檢查"],
        ],
        col_widths=[0.7, 0.5, 0.5, 0.7, 0.9, 1.8],
        title="容器四寶對照：可變性、索引、重複",
    )
    add_textbox(s, T.MARGIN_X, Inches(5.6),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "選容器先問三件事：要不要改？在不在意順序？允不允許重複？",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "Python Data Model §3.1, PEP 468 (dict ordering)")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # S5 — RISK-MITIGATION mutability trio
    s = _blank(prs)
    draw_risk_mitigation(s,
        risks=[
            "def f(x, acc=[])：下次呼叫 acc 還帶著上次的資料",
            "Jupyter 前一個 cell 改了 list；後面 cell 默默讀到被污染狀態",
            "d[[1,2]] = 'x' 直接 TypeError：list 不能當 dict key",
        ],
        mitigations=[
            "根因：可變物件 + 被共享的繫結（預設值、全域狀態、雜湊需求）",
            "緩解 1：預設參數改寫 acc=None，進函式首行 if acc is None: acc = []",
            "緩解 2：Notebook 養成 Restart & Run All 習慣",
            "緩解 3：需要當 key 時改用 tuple 或 frozenset",
        ],
        risks_title="症狀（三處都中同一個陷阱）",
        miti_title="根因 + 緩解",
        title="可變性三病同源：預設參數 / Notebook 污染 / dict key",
        summary="三個症狀同一個病：分不清「物件」與「綁定」。",
    )
    add_source(s, "Python FAQ — Why did changing list 'y' also change list 'x'?")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # S6 — GEOMETRIC three building blocks
    s = _blank(prs)
    draw_three_blocks_flow(s,
        blocks=[
            {"heading": "順序（Sequence）",
             "items": ["step 1", "step 2", "step 3"]},
            {"heading": "分支（Branch）",
             "items": ["if amount > 1000:",
                      "  label = 'high'",
                      "else:",
                      "  label = 'standard'"]},
            {"heading": "迴圈（Loop）",
             "items": ["for row in rows:",
                      "  process(row)",
                      "done"]},
        ],
        title="控制流只有三塊積木：順序、分支、迴圈",
        bottom_note="if / for / list comprehension 都是這三塊的組合；語法糖不改本質。",
    )
    add_source(s, "Böhm & Jacopini 1966 結構化程式定理（教學改寫）")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # S7 — ASK
    s = _blank(prs)
    draw_ask_page(s,
        "你寫的函式，別人看簽章就該知道怎麼用嗎？",
        data_card={
            "label": "帶 type hint 的 PyPI 套件占比",
            "stat": "12% → 68%",
            "caption": "2018 vs 2024，成長 5.7 倍",
        },
    )
    add_source(s, "PyPI Top-5000 packages type-annotation scan, 2024")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # S8 — BEFORE/AFTER with code
    s = _blank(prs)
    add_title(s, "函式契約：加一行 type hint，把「猜」變成「讀」")
    # BEFORE
    draw_code_panel(s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.3),
        label="BEFORE：無契約，讀者靠猜",
        code='def clean(x):\n    return x.replace("$", "").replace(",", "")',
        bullets=[
            "回傳型態？不知道",
            "若傳 int 會 AttributeError；傳 None 會 crash",
            "IDE 無自動補全 / 靜態檢查無效",
        ],
        label_dark=False,
    )
    # AFTER
    draw_code_panel(s,
        x=T.MARGIN_X, y=Inches(3.9),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.6),
        label="AFTER：簽章即契約",
        code=('def clean(x: str) -> float:\n'
              '    """Convert \'$1,200\' -> 1200.0"""\n'
              '    return float(x.replace("$", "").replace(",", ""))'),
        bullets=[
            "輸入輸出型態一眼清楚",
            "mypy / IDE 可靜態檢查",
            "docstring 補足語意，可由 Sphinx 生成文件",
        ],
        label_dark=True,
    )
    add_source(s, "PEP 484 Type Hints; PEP 257 Docstring")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # S9 — BEFORE/AFTER with code (anti-pattern)
    s = _blank(prs)
    add_title(s, "clean_amount anti-pattern：回 0.0 會讓錯誤變成無聲的髒資料")
    draw_code_panel(s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.4),
        label="ANTI-PATTERN：失敗回 0.0",
        code=('def clean_amount(raw: str) -> float:\n'
              '    try:\n'
              '        return float(raw.replace("$", "").replace(",", ""))\n'
              '    except ValueError:\n'
              '        return 0.0   # ← 無聲吞錯'),
        bullets=[
            "真實平均 $842",
            "被 0.0 污染後平均 $516",
            "誤差 -39%，且無警報",
        ],
        label_dark=False,
    )
    draw_code_panel(s,
        x=T.MARGIN_X, y=Inches(4.0),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.3),
        label="FIX：失敗回 None，讓錯誤浮出",
        code=('def clean_amount(raw: str) -> float | None:\n'
              '    try:\n'
              '        return float(raw.replace("$", "").replace(",", ""))\n'
              '    except (ValueError, AttributeError):\n'
              '        return None  # 讓後續 .dropna() / .isna() 接手'),
        bullets=[
            "0.0 會污染平均、相關係數",
            "None 是顯性缺失，pandas 可處理",
            ".dropna() / .isna() 接手",
        ],
        label_dark=True,
    )
    add_source(s, "Hillard 2020, Practices of the Python Pro §4")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # S10 — CHART grouped bar
    s = _blank(prs)
    add_title(s, "Notebook 能不能 Restart & Run All 跑完？成功率決定可重現性")
    chart_png = charts.grouped_bar_m1s10()
    s.shapes.add_picture(str(chart_png), Inches(1.3), Inches(1.3),
                          width=Inches(10.8))
    draw_inverted_thesis_box(s, "Notebook 不能重跑就不是分析，是工藝品。",
                              y=6.1, width=9.0)
    add_source(s, "Pimentel et al., MSR 2019 (n=1.4M notebooks)")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # S11 — MATRIX 2×2 test cost vs value
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "assert 於函式入口 / 出口",
         "sub": "低成本 × 高價值\n一行防呆，抓 70% 前置條件錯誤",
         "highlight": True},
        {"text": "pytest 完整測試套件",
         "sub": "高成本 × 高價值\n（M5 教）"},
        {"text": "print(x) 肉眼檢查",
         "sub": "低成本 × 低價值\n僅適合快速 debug"},
        {"text": "手刻 end-to-end mock",
         "sub": "高成本 × 低價值\n除非複雜整合才值得"},
    ], title="assert 是最便宜的測試：2×2 測試成本 vs 價值")
    add_textbox(s, T.MARGIN_X, Inches(6.3),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
                "橫軸：寫測試成本 低 → 高    縱軸：抓到 bug 的價值 低 → 高",
                font_size=T.FONT_SOURCE, color=T.GRAY_MID,
                align=PP_ALIGN.CENTER)
    add_source(s, "M1 課堂歸納；測試金字塔取自 Mike Cohn 2009")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # S12 — VS Excel vs Python
    s = _blank(prs)
    draw_vs_two_col(s,
        left_title="Excel 流程",
        right_title="Python 流程",
        left_items=[
            "下載 CSV",
            "手動篩選",
            "手動刪行",
            "改欄位公式",
            "複製到 PPT",
            "下週新資料 → 全部重來",
        ],
        right_items=[
            "下載 CSV",
            "python clean.py",
            "python analyze.py",
            "",
            "",
            "下週新資料 → 只換一行檔名",
        ],
        title="Excel vs Python：兩種工作流的可重現性對照",
        summary="操作留痕 = 可重現；可重現 = 可審計 = 可交付。",
        delta="可重現性\n全有 / 全無",
    )
    add_source(s, "Wilson et al., Good Enough Practices, PLOS 2017")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # S13 — TABLE four data pollution types
    s = _blank(prs)
    draw_editorial_table(s,
        header=["污染類型", "判斷訊號", "Python 第一動作", "不該做什麼"],
        rows=[
            ["缺失 (Missing)", "NaN / 空字串 / 可疑 0",
             "df.isna().sum() 再決策",
             "直接 fillna(0) 會污染平均"],
            ["重複 (Duplicate)", "完全重複 或 邏輯重複（拼法不一）",
             "drop_duplicates / str.lower().strip()",
             "假設「看起來不同」就真的不同"],
            ["離群 (Outlier)", "IQR 之外 / 業務常識之外",
             "先 describe() + boxplot 檢視",
             "未問業務就直接砍"],
            ["型態 (Type)", "dtype=object 該是數字",
             "to_numeric(errors='coerce')",
             "int(x) 硬轉，遇髒字就 crash"],
        ],
        col_widths=[1.0, 1.4, 1.5, 1.3],
        title="資料清理四類污染與 Python 對應動作",
    )
    add_textbox(s, T.MARGIN_X, Inches(5.6),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "四類 MECE；每類第一個動作 ≤ 一行程式。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "Wickham 2014, Tidy Data; pandas User Guide §10")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # S14 — GEOMETRIC concentric (simplified as 3 horizontal bands)
    s = _blank(prs)
    draw_concentric_zones(s,
        zones=[
            {"label": "Python 語言核心",
             "items": "syntax / types / control flow — M1–M2 主場",
             "highlight": True},
            {"label": "資料分析層",
             "items": "pandas · numpy · matplotlib · requests · pytest — M3–M4 主場"},
            {"label": "ML / 系統 / AI 層",
             "items": "scikit-learn · PyTorch · FastAPI · PySpark · LangChain · Docker SDK — M5–M7 主場",
             "dashed": True},
        ],
        title="import 與套件生態：三層同心結構",
    )
    add_textbox(s, T.MARGIN_X, Inches(6.15),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "import 的本質：把別人的圈納入你的圈。語言核心小，生態大。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "PyPI top downloads 2024 Q4 (pypistats.org)")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # S15 — PYRAMID hierarchy close
    s = _blank(prs)
    draw_thesis_hierarchy(s,
        blocks=[
            {"heading": "資料思維三問",
             "items": [
                 "這份資料的 schema 是什麼？（列 / 欄 / 型態 / 鍵）",
                 "哪一類髒了？（缺失 / 重複 / 離群 / 型態）",
                 "這個分析明天還跑得出來嗎？（可重現性）",
             ]},
            {"heading": "Python 語法三寶",
             "items": [
                 "容器：選 list / dict / tuple / set 前先問可變、順序、重複",
                 "流程：順序 / 分支 / 迴圈三積木足以表達任何清理邏輯",
                 "函式：簽章即契約（type hint + docstring + assert）",
             ]},
        ],
        title="M1 收束：資料思維三問 + 語法三寶",
        thesis="語法是手段，可信是目的。M2 要讓這些函式長出骨架。",
    )
    add_source(s, "M1 module synthesis")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # S16 — SILENT closing
    s = _blank(prs)
    draw_silent_page(s, "會寫 Python ≠ 會做資料；\n兩者之間差一個「可信」。")
    add_footer(s, MODULE_CODE, 16, N_CONTENT, dark_bg=True)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
