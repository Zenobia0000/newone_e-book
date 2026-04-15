"""M6 deck — 15 content slides + cover + copyright page.

Module: M6 · 計算機組織與作業系統
Governing thought: Performance is physics respected; every abstraction leaks.

Required chart functions (Agent-E to add to charts.py; DO NOT edit here):
- charts.bar_chart_m6s2_dtype_ram(outfile=None) -> Path
    Horizontal bar chart, linear scale.
    data: 4 rows (label, MB):
        ("int64 × 1M 列", 8),
        ("float64 × 1M 列", 8),
        ("category × 1M 列（100 類）", 2),
        ("object (str) × 1M 列", 900)
    figsize ≈ (11.5, 4.0). End-of-bar value labels in #1B5E3F bold:
    "8 MB (基準)", "8 MB", "2 MB (−75%)", "~900 MB (×100 倍)".
    X-axis "RAM 佔用 (MB)" ticks 0/200/400/600/800/1000.

- charts.bar_chart_m6s3_memory_hierarchy(outfile=None) -> Path
    Horizontal bar chart, LOG scale (x from 1 to 1e9 ns).
    data: 6 rows:
        ("L1 cache", 1),
        ("L2 cache", 4),
        ("L3 cache", 15),
        ("RAM (DRAM)", 100),
        ("NVMe SSD", 100_000),
        ("Network RTT（同機房）", 1_000_000)
    End labels: "≈ 1 ns", "≈ 4 ns", "≈ 15 ns",
    "≈ 100 ns (×100 於 L1)", "≈ 100 μs (×10^5)", "≈ 1 ms (×10^6)".
    figsize ≈ (11.5, 4.2).

- charts.bar_chart_m6s9_io_evolution(outfile=None) -> Path
    Horizontal bar chart, LOG scale (x from 1K to 10M).
    data: 5 rows:
        ("1983 · fd + select()", 1_024),
        ("1997 · poll()", 10_000),
        ("2002 · epoll (Linux)", 1_000_000),
        ("2003 · kqueue (BSD/macOS)", 1_000_000),
        ("2019 · io_uring (Linux 5.1+)", 10_000_000)
    End labels: "1,024 (硬上限 FD_SETSIZE)", "~10K (O(n) 掃描)",
    "~1M (O(1) 事件驅動)", "~1M", "~10M (無系統呼叫、零拷貝)".
    figsize ≈ (11.5, 4.2).

- charts.before_after_m6s10_dtype(outfile=None) -> Path
    Two side-by-side vertical bar charts sharing Y axis (0-1000 MB linear).
    Left: "BEFORE：df['city'] as object" — one bar at 900 MB;
          top label "900 MB (基準)"; caption "dtype=object / ~80 bytes × 1M + 指標".
    Right: "AFTER：df['city'].astype('category')" — bar at 2 MB;
          top label "2 MB (−99.8%)"; caption "dtype=category / 100 類別 + int8".
    Grey vertical divider between panels. figsize ≈ (11.0, 4.5).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from .. import charts
from ..primitives import (
    add_rect, add_textbox, add_title, add_source, set_solid_fill, set_no_line,
    set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_pyramid_stack, draw_inverted_thesis_box,
    draw_vs_two_col, draw_three_blocks_flow, draw_split_panel,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "M6"
MODULE_TITLE = "計算機組織與作業系統"
MODULE_SUBTITLE = "Python 數據分析與 AI 基礎 · M6"
TIME_MIN = 45
N_CONTENT = 15


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_m6(output_path, image_registry=None):
    """Build M6 deck. image_registry is optional placeholder collector."""
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE, TIME_MIN, N_CONTENT)

    # S1 — ASK · 你的訓練為什麼慢，是算不完還是在等搬運？
    s = _blank(prs)
    draw_ask_page(
        s,
        "你的訓練為什麼慢，是算不完還是在等搬運？",
        data_card={
            "label": "單 iteration 時間拆解",
            "stat": "80%",
            "caption": "GPU 等資料，只 20% 在真算（num_workers=0 示例）",
        },
    )
    add_source(s, "本課 DataLoader profile 示例（num_workers=0）")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # S2 — CHART · DataFrame 一個 object 欄，RAM 瞬間多吃一個零
    s = _blank(prs)
    add_title(s, "DataFrame 一個 object 欄，RAM 瞬間多吃一個零")
    chart_png = charts.bar_chart_m6s2_dtype_ram()
    s.shapes.add_picture(str(chart_png), Inches(0.8), Inches(1.3),
                         width=Inches(11.7))
    add_textbox(s, T.MARGIN_X, Inches(5.6),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "同樣 1M 列，差一百倍。",
                font_size=T.FONT_BODY, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "pandas memory_usage(deep=True) 實測；1M rows 英文字串平均 80 bytes")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # S3 — CHART · 記憶體階層延遲（對數刻度）
    s = _blank(prs)
    add_title(s, "記憶體階層：同一件事在不同層級快或慢 10^6 倍")
    chart_png = charts.bar_chart_m6s3_memory_hierarchy()
    s.shapes.add_picture(str(chart_png), Inches(0.8), Inches(1.3),
                         width=Inches(11.7))
    add_textbox(s, T.MARGIN_X, Inches(5.6),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "從 L1 到網路，每一層 ×100 到 ×10^6。",
                font_size=T.FONT_BODY, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "Jeff Dean Numbers Every Programmer Should Know / Intel Skylake-X 官方 cache latency")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # S4 — TABLE · 把奈秒放大成人類時間
    s = _blank(prs)
    draw_editorial_table(s,
        header=["層級", "實際延遲", "放大 10^9 倍後的人類尺度"],
        rows=[
            ["L1 cache", "1 ns", "1 秒"],
            ["L2 cache", "4 ns", "4 秒"],
            ["RAM", "100 ns", "1 分 40 秒"],
            ["NVMe SSD", "100 μs", "28 小時"],
            ["同機房網路", "1 ms", "11.6 天"],
            ["跨大陸網路", "150 ms", "4.75 年"],
        ],
        col_widths=[1.0, 1.0, 2.0],
        title="把奈秒放大成人類時間：一杯水 vs 一趟火車",
    )
    draw_inverted_thesis_box(
        s, "CPU 等 RAM 是喘口氣，CPU 等網路是等下一次生日。",
        y=6.0, width=10.5)
    add_source(s, "Jeff Dean 2009 / Peter Norvig 放大比喻整理")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # S5 — VS · CPU 幾個強核心 vs GPU 幾千個弱核心
    s = _blank(prs)
    draw_vs_two_col(s,
        left_title="CPU（8 cores, ~3 GHz）",
        right_title="GPU（~10,000 CUDA cores, ~1.5 GHz）",
        left_items=[
            "核心數：8（典型 desktop）",
            "單核時脈：約 3 GHz",
            "強項：分支、邏輯、控制流",
            "cache 大、分支預測強",
            "什麼都能做，數量少",
        ],
        right_items=[
            "核心數：約 10,000 CUDA cores",
            "單核時脈：約 1.5 GHz",
            "強項：同一種運算 × 極大量",
            "SIMT、向量化友善",
            "只做同一件事，但同時做",
        ],
        title="CPU 幾個強核心 vs GPU 幾千個弱核心",
        summary="差別不是誰快，而是數量。矩陣乘法是同一件事做一兆次，GPU 贏在這裡。",
        delta="8 vs 10,000",
    )
    add_source(s, "Intel i9-13900K / NVIDIA RTX 4090 官方規格")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # S6 — GEOMETRIC-DIAGRAM · 虛擬記憶體三層翻譯
    s = _blank(prs)
    add_title(s, "虛擬記憶體三層翻譯：VA → Page Table → PA")
    # Layer 1 — Virtual Address (top band)
    l1_y = Inches(1.3)
    l1 = add_rect(s, Inches(2.5), l1_y, Inches(8.3), Inches(0.9))
    set_no_fill(l1); set_line(l1, T.PRIMARY, 1.2)
    add_textbox(s, Inches(2.5), l1_y, Inches(8.3), Inches(0.9),
                "Virtual Address  0x7fff_a1b2_c000",
                font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                family=T.FONT_MONO)
    add_textbox(s, Inches(10.9), l1_y, Inches(2.2), Inches(0.9),
                "每個 process 自有\n2^48 位址空間",
                font_size=T.FONT_SMALL, color=T.CHARCOAL,
                anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    add_textbox(s, Inches(0.3), l1_y, Inches(2.1), Inches(0.9),
                "程式看到的",
                font_size=T.FONT_SMALL, color=T.PRIMARY, bold=True,
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # Arrow down
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(6.65), Inches(2.3),
                                   Inches(6.65), Inches(2.8))
    set_line(conn, T.GRAY_MID, 1.2)

    # Layer 2 — three boxes: Page Table / TLB Cache / Page Fault Handler
    l2_y = Inches(2.85)
    box_w = Inches(2.5); gap = Inches(0.3)
    xs = [Inches(2.5) + i * (box_w + gap) for i in range(3)]
    labels = [("Page Table", False),
              ("TLB Cache\n64–1536 entries", True),
              ("Page Fault Handler", False)]
    for xi, (lab, emph) in zip(xs, labels):
        r = add_rect(s, xi, l2_y, box_w, Inches(1.0))
        set_no_fill(r)
        set_line(r, T.PRIMARY, 2.0 if emph else 1.0)
        add_textbox(s, xi, l2_y, box_w, Inches(1.0),
                    lab, font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                    line_spacing=1.2)
    add_textbox(s, Inches(0.3), l2_y, Inches(2.1), Inches(1.0),
                "OS 的翻譯表",
                font_size=T.FONT_SMALL, color=T.PRIMARY, bold=True,
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(10.9), l2_y, Inches(2.2), Inches(1.0),
                "TLB hit ≈ 1 ns\nTLB miss ≈ 100 ns\nPage fault ≈ 100 μs–10 ms",
                font_size=T.FONT_SMALL, color=T.CHARCOAL,
                anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

    # Arrow down to Layer 3
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(6.65), Inches(3.95),
                                   Inches(6.65), Inches(4.45))
    set_line(conn, T.GRAY_MID, 1.2)

    # Layer 3 — RAM + Swap
    l3_y = Inches(4.5)
    b1 = add_rect(s, Inches(3.5), l3_y, Inches(3.2), Inches(1.0))
    set_no_fill(b1); set_line(b1, T.PRIMARY, 1.0)
    add_textbox(s, Inches(3.5), l3_y, Inches(3.2), Inches(1.0),
                "RAM 物理頁 4 KB",
                font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    b2 = add_rect(s, Inches(6.9), l3_y, Inches(3.2), Inches(1.0))
    set_no_fill(b2); set_line(b2, T.PRIMARY, 1.0)
    add_textbox(s, Inches(6.9), l3_y, Inches(3.2), Inches(1.0),
                "Swap on SSD 4 KB",
                font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(0.3), l3_y, Inches(3.1), Inches(1.0),
                "實體硬體",
                font_size=T.FONT_SMALL, color=T.PRIMARY, bold=True,
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(3.5), l3_y + Inches(1.05), Inches(6.6), Inches(0.4),
                "RAM ≈ 100 ns　／　Swap ≈ 100 μs（×1000）",
                font_size=T.FONT_SMALL, color=T.CHARCOAL,
                align=PP_ALIGN.CENTER)

    draw_inverted_thesis_box(s, "你寫的位址，不是真的位址。", y=6.4, width=9.0)
    add_source(s, "Bryant & O'Hallaron CS:APP Ch.9 / Linux x86_64 4-level paging")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # S7 — TABLE · Process 與 Thread
    s = _blank(prs)
    draw_editorial_table(s,
        header=["維度", "Process", "Thread"],
        rows=[
            ["記憶體空間", "獨立（各自 virtual address space）", "共用（同一 process 內）"],
            ["建立成本", "高（~1–10 ms，複製 page table）", "低（~10 μs，共用資源）"],
            ["資料共享", "需 IPC / shared memory", "直接讀寫同一變數"],
            ["崩潰隔離", "一個掛不影響別人", "一個掛整個 process 死"],
            ["Python GIL 影響", "不受限（各自一把）", "CPython 同時只一個跑 bytecode"],
        ],
        col_widths=[1.0, 1.9, 1.6],
        title="Process 與 Thread：五個維度一眼對照",
    )
    draw_inverted_thesis_box(
        s, "CPU 密集選 multiprocessing，I/O 密集選 threading；選錯就是白忙。",
        y=6.0, width=11.0)
    add_source(s, "Python docs / Linux clone(2) / PEP 703 (GIL)")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # S8 — TABLE · fork / exec / spawn
    s = _blank(prs)
    draw_editorial_table(s,
        header=["方式", "機制", "平台", "建立延遲"],
        rows=[
            ["fork", "複製父 process 整個位址空間（copy-on-write）",
             "Linux / macOS ✓ / Windows ✗", "~1 ms"],
            ["exec", "取代當前 process 影像、保留 PID",
             "全平台", "~2 ms"],
            ["spawn", "全新 process + 重新 import 所有模組",
             "全平台（Windows 預設）", "~50–200 ms"],
            ["forkserver", "預啟一台乾淨伺服器分發 fork",
             "Linux / macOS", "~2 ms（暖機後）"],
        ],
        col_widths=[0.7, 2.2, 1.6, 0.9],
        title="fork / exec / spawn 三種產生 process 的方式，Windows 只有一種能用",
    )
    draw_inverted_thesis_box(
        s,
        "Windows 沒有 fork。if __name__ == '__main__': 不是裝飾，是 spawn 必需的防無限迴圈。",
        y=6.0, width=11.5)
    add_source(s, "Python multiprocessing docs / POSIX fork(2) / Windows CreateProcess")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # S9 — CHART · I/O 模型演進（對數）
    s = _blank(prs)
    add_title(s, "I/O 模型演進：select → epoll → io_uring，吞吐量一次一個數量級")
    chart_png = charts.bar_chart_m6s9_io_evolution()
    s.shapes.add_picture(str(chart_png), Inches(0.8), Inches(1.3),
                         width=Inches(11.7))
    add_textbox(s, T.MARGIN_X, Inches(5.6),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "36 年，單執行緒連線上限 ×10,000。",
                font_size=T.FONT_BODY, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "Linux kernel docs / Jens Axboe io_uring 2019 / Kegel C10K 1999")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # S10 — BEFORE/AFTER · object vs category
    s = _blank(prs)
    add_title(s, "同一份 1M 列表格，object vs category 差 450 倍")
    chart_png = charts.before_after_m6s10_dtype()
    s.shapes.add_picture(str(chart_png), Inches(1.0), Inches(1.3),
                         width=Inches(11.3))
    draw_inverted_thesis_box(
        s, "改一個 dtype，省下 99.8% RAM。別讓 pandas 幫你吃記憶體。",
        y=6.0, width=10.5)
    add_source(s, "pandas memory_usage(deep=True) 實測、1M rows 100 unique cities")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # S11 — GEOMETRIC-DIAGRAM · DataLoader 流水線時序對比
    s = _blank(prs)
    add_title(s, "DataLoader 流水線：num_workers=4 讓 GPU 不再空等")

    # Top track: num_workers=0
    label_x = Inches(0.4); label_w = Inches(2.0)
    track_x = Inches(2.5); track_w = Inches(9.5)
    t1_y = Inches(1.6); track_h = Inches(0.7)
    add_textbox(s, label_x, t1_y, label_w, track_h,
                "num_workers=0\n（單 process）",
                font_size=T.FONT_SMALL, color=T.PRIMARY, bold=True,
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=1.2)
    # 4 cycles: 800ms CPU (gray) + 200ms GPU (green) — ratio 4:1
    seg_total = track_w
    unit = seg_total / (4 * (4 + 1))  # 4 cycles of 5 units
    cx = track_x
    for _ in range(4):
        # CPU (gray, 4 units)
        cpu = add_rect(s, cx, t1_y, unit * 4, track_h)
        set_solid_fill(cpu, T.LIGHT_GRAY); set_no_line(cpu)
        cx += unit * 4
        gpu = add_rect(s, cx, t1_y, unit * 1, track_h)
        set_solid_fill(gpu, T.PRIMARY); set_no_line(gpu)
        cx += unit * 1
    add_textbox(s, track_x, t1_y + track_h + Inches(0.05), track_w, Inches(0.3),
                "灰＝CPU 讀檔+解碼 800 ms　│　綠＝GPU 計算 200 ms　→　GPU 利用率 ≈ 20%",
                font_size=T.FONT_SMALL, color=T.CHARCOAL,
                align=PP_ALIGN.CENTER)

    # Divider
    div_y = Inches(3.3)
    div = add_rect(s, Inches(0.4), div_y, T.SLIDE_W - Inches(0.8), Inches(0.02))
    set_solid_fill(div, T.GRAY_MID); set_no_line(div)

    # Bottom track: num_workers=4
    t2_y = Inches(3.8)
    add_textbox(s, label_x, t2_y, label_w, track_h,
                "num_workers=4\n（四個子 process）",
                font_size=T.FONT_SMALL, color=T.PRIMARY, bold=True,
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=1.2)
    # 4 parallel pre-fetch lanes (small gray stripes) at top
    lane_h = Inches(0.12)
    for i in range(4):
        ly = t2_y - Inches(0.6) + i * Inches(0.12)
        lane = add_rect(s, track_x, ly, unit * 4, lane_h)
        set_solid_fill(lane, T.LIGHT_GRAY); set_no_line(lane)
    # Main axis: 800ms gray prelude then continuous GPU green
    warm = add_rect(s, track_x, t2_y, unit * 4, track_h)
    set_solid_fill(warm, T.LIGHT_GRAY); set_no_line(warm)
    gpu_cont = add_rect(s, track_x + unit * 4, t2_y,
                         track_w - unit * 4, track_h)
    set_solid_fill(gpu_cont, T.PRIMARY); set_no_line(gpu_cont)
    add_textbox(s, track_x, t2_y + track_h + Inches(0.05), track_w, Inches(0.3),
                "四個 worker 在背景預取下一批　→　GPU 利用率 ≈ 95%",
                font_size=T.FONT_SMALL, color=T.CHARCOAL,
                align=PP_ALIGN.CENTER)

    draw_inverted_thesis_box(
        s, "四個 worker 把等待疊成背景；你只改了一個參數。",
        y=6.0, width=10.5)
    add_source(s, "PyTorch DataLoader docs / 本課 ImageFolder ResNet18 batch=32")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # S12 — MATRIX · RAM OOM vs VRAM OOM 四象限
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        # Top-left: RAM + 緩增
        {"text": "RAM 緩增：進程 RSS 持續上升",
         "sub": "原因：循環內建立新 DataFrame 未釋放 / 閉包抓參考\n"
                "第一動作：tracemalloc 定位、明確 del"},
        # Top-right: VRAM + 緩增
        {"text": "VRAM 緩增：nvidia-smi 每輪 +幾十 MB",
         "sub": "原因：loss.item() 漏呼叫 / retain_graph=True / 中間張量沒 detach\n"
                "第一動作：loss = loss.item() ／ .detach()"},
        # Bottom-left: RAM + 分配失敗
        {"text": "RAM 分配失敗：MemoryError: Unable to allocate（高發 80%）",
         "sub": "原因：pd.concat 大表 / object dtype / 一次 read_csv 全檔\n"
                "第一動作：改 chunksize、astype('category')、del + gc.collect()",
         "highlight": True},
        # Bottom-right: VRAM + 分配失敗
        {"text": "VRAM 分配失敗：CUDA out of memory on allocation（高發 80%）",
         "sub": "原因：batch_size 太大 / 模型太大 / 沒 zero_grad\n"
                "第一動作：降 batch_size、加 optimizer.zero_grad()、torch.cuda.empty_cache()",
         "highlight": True},
    ], title="RAM OOM vs VRAM OOM：四個象限、四個第一動作")
    add_textbox(s, T.MARGIN_X, Inches(6.4),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
                "橫軸：發生位置 CPU RAM ↔ GPU VRAM　　縱軸：症狀 緩增 ↔ 分配失敗",
                font_size=T.FONT_SOURCE, color=T.GRAY_MID,
                align=PP_ALIGN.CENTER)
    add_source(s, "PyTorch CUDA memory docs / pandas Memory usage / tracemalloc docs")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # S13 — PYRAMID · 五個延遲量級（橫向 flow chain 承載）
    s = _blank(prs)
    draw_flow_chain(s,
        nodes=[
            {"label": "1  L1 ≈ 1 ns", "sub": "ALU 指令級"},
            {"label": "2  RAM ≈ 100 ns", "sub": "×100 於 L1"},
            {"label": "3  SSD ≈ 100 μs", "sub": "×1000 於 RAM"},
            {"label": "4  同機房網路 ≈ 1 ms", "sub": "×10 於 SSD"},
            {"label": "5  跨洲網路 ≈ 150 ms", "sub": "×150 於同機房",
             "highlight": True},
        ],
        title="記住五個延遲量級，你就有系統直覺",
        y=3.0,
    )
    draw_inverted_thesis_box(
        s, "把這五個數字刻在腦裡，你的效能直覺比九成工程師準。",
        y=5.5, width=11.0)
    add_source(s, "Jeff Dean 2009 updated by Colin Scott 2020 / 本課精簡裁版")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # S14 — SILENT · Performance is physics respected
    s = _blank(prs)
    draw_silent_page(s, "Performance is physics respected.")
    add_footer(s, MODULE_CODE, 14, N_CONTENT, dark_bg=True)

    # S15 — SILENT closing
    s = _blank(prs)
    draw_silent_page(s, "每一層抽象 都會漏。\n到了 M7，你會再次感謝它。")
    add_footer(s, MODULE_CODE, 15, N_CONTENT, dark_bg=True)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
