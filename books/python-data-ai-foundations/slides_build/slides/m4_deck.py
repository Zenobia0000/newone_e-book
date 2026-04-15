"""M4 deck — 16 content slides + cover + copyright page.

Editorial-strict, Chart-first. 配色 #1B5E3F 深綠 / #333333 炭灰 / #D3D3D3 淺灰 / 白，
禁止紅 / 黃 / 橙 / 粉 / 淺藍 (G8)。

Required chart functions (to be implemented separately in charts.py — parallel edits forbidden):

  - datasaurus_dozen_m4s3(): 3x4 散佈圖, Datasaurus Dozen 十二形狀，
      同一組摘要 (μx≈54.3, μy≈47.8, ρ≈-0.06)。figsize=(11.5, 6.0), 純綠圓點。
      Source: Matejka & Fitzmaurice 2017.

  - before_after_truncated_axis_m4s6(): 左右對照長條，兩組同值 49% / 51%。
      左圖 Y 軸 48–52%，右圖 Y 軸 0–100%。figsize=(11.5, 4.8)，純綠柱。
      Source: 本課自擬示例 (G5 圖表說謊示範)。

  - bar_vs_box_m4s9(): 左長條 (A/B 班皆 75)、右箱形圖 (A 中位 76 IQR 8 離群 0;
      B 中位 78 IQR 14 離群 4)。figsize=(11.5, 4.8)，純綠。
      Source: 本課自擬成績資料。

  - ci_shrink_loglog_m4s10(): 對數 X 軸 n=10..10000, Y 軸 ±% 0–30。
      三節點 (30, ±9%), (300, ±3%), (3000, ±1%)。figsize=(11.5, 5.2)。
      Source: 二項分布常態近似估算。

These helpers must be added to charts.py by the Charts agent. This deck assumes
`getattr(charts, <name>, None)` — if missing, the slide falls back to a textual
placeholder panel so the deck still renders.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from .. import charts
from ..primitives import (
    add_rect, add_textbox, add_title, add_source, set_solid_fill, set_no_line,
    set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_pyramid_stack, draw_inverted_thesis_box,
    draw_vs_two_col, draw_risk_mitigation, draw_grid,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "M4"
MODULE_TITLE = "EDA、視覺化與統計直覺"
MODULE_SUBTITLE = "Python 數據分析與 AI 基礎 · M4"
TIME_MIN = 45
N_CONTENT = 16


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _chart_or_placeholder(slide, fn_name: str, x, y, w, h, fallback_text: str):
    """Insert chart PNG if the chart function exists in charts.py;
    otherwise draw an editorial placeholder rectangle with the spec text.
    """
    fn = getattr(charts, fn_name, None)
    if fn is not None:
        png = fn()
        slide.shapes.add_picture(str(png), x, y, width=w)
        return
    rect = add_rect(slide, x, y, w, h)
    set_no_fill(rect)
    set_line(rect, T.PRIMARY, 1.0, dash=True)
    add_textbox(
        slide, x + Inches(0.3), y + Inches(0.3),
        w - Inches(0.6), h - Inches(0.6),
        f"[CHART PENDING: {fn_name}]\n{fallback_text}",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4,
    )


def build_m4(output_path, image_registry=None):
    """Build M4 deck. image_registry kept for API parity (M4 has no PHOTO)."""
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE, TIME_MIN, N_CONTENT)

    # S1 — ASK · 你敢只看平均 4.0 星就下單嗎？
    s = _blank(prs)
    draw_ask_page(
        s,
        "平均 4.0 = 安全嗎？",
        data_card={
            "label": "同樣 mean=4.0，分佈不同",
            "stat": "4.0★",
            "caption": "整齊好評 vs 兩極地雷，摘要看不出差別",
        },
    )
    add_source(s, "本課自擬示例")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # S2 — VS · 四組資料統計量完全相同，散佈卻完全不同 (Anscombe)
    s = _blank(prs)
    add_title(s, "數字一樣 圖形不一樣：Anscombe 四組資料")
    # Left: editorial TABLE of Anscombe summary stats
    draw_editorial_table(s,
        header=["Set", "mean(x)", "var(x)", "mean(y)", "var(y)", "corr"],
        rows=[
            ["I",   "9.0", "11.0", "7.5", "4.1", "0.82"],
            ["II",  "9.0", "11.0", "7.5", "4.1", "0.82"],
            ["III", "9.0", "11.0", "7.5", "4.1", "0.82"],
            ["IV",  "9.0", "11.0", "7.5", "4.1", "0.82"],
        ],
        col_widths=[0.5, 0.8, 0.8, 0.8, 0.8, 0.8],
        top=1.3,
    )
    # Centred "≠" + bottom summary
    add_textbox(s, T.MARGIN_X, Inches(4.6),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
                "摘要相同  ≠  真相相同",
                font_size=Pt(24), color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER)
    add_textbox(s, T.MARGIN_X, Inches(5.3),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.9),
                "II 是曲線 · III 有離群點 · IV 根本只剩一個極端值在撐。只看表格會以為是同一個故事；畫出來才看見。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL,
                align=PP_ALIGN.CENTER, line_spacing=1.4)
    add_source(s, "Anscombe 1973")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # S3 — CHART · Datasaurus Dozen 3×4
    s = _blank(prs)
    add_title(s, "十二種形狀 同一組摘要：Datasaurus Dozen")
    _chart_or_placeholder(
        s, "datasaurus_dozen_m4s3",
        Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.0),
        "3×4 散佈圖：Dino/Star/X/Bullseye/Circle/Dots/High/Slant-up/Slant-down/V/H/Wide\n"
        "全體 μx 54.3 · μy 47.8 · ρ -0.06",
    )
    add_textbox(s, T.MARGIN_X, Inches(6.3),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "從恐龍到星星到 X，統計摘要完全一樣。不畫分佈就下結論，你不知道自己正在看恐龍。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "Matejka & Fitzmaurice 2017")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # S4 — GEOMETRIC-DIAGRAM · 四類問題決策樹 (用 grid 4 葉呈現 MECE)
    s = _blank(prs)
    add_title(s, "先問要回答什麼 再挑圖：四類問題、四種圖")
    # Root node
    root_w = Inches(4.0)
    root_x = (T.SLIDE_W - root_w) / 2
    root = add_rect(s, root_x, Inches(1.3), root_w, Inches(0.7))
    set_no_fill(root)
    set_line(root, T.PRIMARY, 1.5)
    add_textbox(s, root_x, Inches(1.3), root_w, Inches(0.7),
                "我想回答什麼？",
                font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Four leaf boxes
    draw_grid(s, rows=1, cols=4, cells=[
        {"label": "分佈", "sub": "直方圖 / 箱形圖", "note": "例：薪水分佈"},
        {"label": "比較", "sub": "長條圖 / 分組箱形圖", "note": "例：A/B 班成績"},
        {"label": "關聯", "sub": "散佈圖 / 熱力圖", "note": "例：身高 × 體重"},
        {"label": "時序", "sub": "折線圖 / 區域圖", "note": "例：月營收走勢"},
    ], top=2.6, bottom=2.0)
    add_textbox(s, T.MARGIN_X, Inches(5.5),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
                "四類問題 四種圖",
                font_size=Pt(22), color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER)
    add_textbox(s, T.MARGIN_X, Inches(6.1),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
                "圖表不是美工，是回答問題的工具。問題先鎖定，圖才選得對。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL,
                align=PP_ALIGN.CENTER)
    add_source(s, "本課改寫自 Cleveland 1985")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # S5 — PYRAMID · EDA 七步法 (使用 flow_chain 呈現七節 + 倒掛框)
    s = _blank(prs)
    add_title(s, "EDA 七步法：終點是一句可驗證的假設")
    draw_flow_chain(s, nodes=[
        {"label": "0 品質", "caption": "dtypes/NA"},
        {"label": "1 全貌", "caption": "describe"},
        {"label": "2 單欄", "caption": "hist/box"},
        {"label": "3 雙欄", "caption": "scatter/corr"},
        {"label": "4 多欄", "caption": "pairplot/PCA"},
        {"label": "5 異常", "caption": "IQR/Z"},
        {"label": "6 假設", "caption": "H0/H1", "highlight": True},
    ], y=3.0)
    draw_inverted_thesis_box(s, "EDA 的終點不是漂亮圖，是一句可驗證的假設。", y=6.0)
    add_source(s, "Tukey 1977 / 本課改寫")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # S6 — BEFORE/AFTER · Y 軸截斷
    s = _blank(prs)
    add_title(s, "同一份資料，截斷 Y 軸就能把 2% 畫成懸崖")
    _chart_or_placeholder(
        s, "before_after_truncated_axis_m4s6",
        Inches(0.9), Inches(1.3), Inches(11.5), Inches(4.6),
        "左 BEFORE Y=48–52% (視覺差近 5 倍) · 右 AFTER Y=0–100% (幾乎等高)\n"
        "兩柱 49% (基準) / 51% (+2pp)",
    )
    add_textbox(s, T.MARGIN_X, Inches(6.1),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
                "Y 軸從哪開始 決定真相",
                font_size=Pt(22), color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "本課自擬示例")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # S7 — TABLE · 圖表說謊四連擊 vs 四個拆穿動作
    s = _blank(prs)
    draw_editorial_table(s,
        header=["#", "常見說謊招式", "拆穿這招的第一個動作"],
        rows=[
            ["1", "Y 軸截斷", "看軸起點，是否從 0（或合理基準）起算"],
            ["2", "雙 Y 軸疊折線", "兩軸單位 / 量級是否可比，不可比就拆成兩張"],
            ["3", "挑起訖日期", "擴大時間窗，看完整區間是否同方向"],
            ["4", "彩虹 / 連續色分類", "類別變數應用離散調色盤，非連續漸層"],
        ],
        col_widths=[0.3, 1.6, 2.5],
        title="圖表說謊四連擊 vs 四個拆穿動作",
    )
    draw_inverted_thesis_box(s, "先看軸，再看色，最後才看數字。", y=5.8, width=9.0)
    add_source(s, "Cairo 2016, The Truthful Art / Few 2012")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # S8 — MATRIX · 色盲友善配色 2×2
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        # Top-left: low distinguishable, high colorblind-safe
        {"text": "單色深淺階",
         "sub": "深綠 → 白四階\n類別少時可用；類別多會混"},
        # Top-right: both high — Okabe-Ito (G8 合規版，去黃橙)
        {"text": "Okabe-Ito 裁版  ✓ 建議預設",
         "sub": "#1B5E3F 深綠 · #0072B2 靛 · #9467BD 深紫 · #333333 深灰\n（剔除黃 #F0E442 與橙 #E69F00 以守 G8 禁色）",
         "highlight": True},
        # Bottom-left: both low
        {"text": "紅 ✗ 綠 ✗ 橘 ✗",
         "sub": "紅綠對立 + 橙：色盲觀眾 8% 看不出\n類別、可讀性同時失敗"},
        # Bottom-right: high distinguishable, low friendly
        {"text": "彩虹連續漸層",
         "sub": "連續色用在類別：\n順序被誤讀、色盲觀眾失訊"},
    ], title="色盲友善配色的 2×2：可讀性 × 可區分")
    add_textbox(s, T.MARGIN_X, Inches(6.35),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.35),
                "縱軸：色盲友善度 低 → 高    橫軸：單色明暗可區分性 低 → 高",
                font_size=T.FONT_SOURCE, color=T.GRAY_MID,
                align=PP_ALIGN.CENTER)
    add_source(s, "Okabe & Ito 2008 裁版（剔除黃橙以守 G8）")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # S9 — VS · 平均 vs 分佈（長條 vs 箱形）
    s = _blank(prs)
    add_title(s, "平均騙你的時候，箱形圖是抓包工具")
    _chart_or_placeholder(
        s, "bar_vs_box_m4s9",
        Inches(0.9), Inches(1.3), Inches(11.5), Inches(4.6),
        "左：只看平均——A/B 班皆 75\n"
        "右：看分佈——A 中位 76 IQR 8 離群 0 · B 中位 78 IQR 14 離群 4 人",
    )
    add_textbox(s, T.MARGIN_X, Inches(6.1),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
                "平均一樣 分佈差很多",
                font_size=Pt(22), color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "本課自擬成績資料")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # S10 — CHART · n 變大，不確定性收斂
    s = _blank(prs)
    add_title(s, "n 變大 誤差變小：樣本數是信心的單位")
    _chart_or_placeholder(
        s, "ci_shrink_loglog_m4s10",
        Inches(0.8), Inches(1.3), Inches(11.7), Inches(4.8),
        "X 對數 10→10,000 · Y 線性 0–30%\n"
        "標註：n=30 → ±9% · n=300 → ±3% · n=3000 → ±1%",
    )
    add_textbox(s, T.MARGIN_X, Inches(6.3),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "五個朋友說好吃 vs 五百個說好吃：差的不是數字，是你敢不敢推薦給客戶。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "二項分布常態近似估算")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # S11 — PYRAMID · p-value 三個「不」(使用三卡 matrix 1x3 + 倒掛框)
    s = _blank(prs)
    add_title(s, "三個不：p-value 不是魔法")
    draw_matrix(s, rows=1, cols=3, cells=[
        {"text": "① 不告訴你差異有多大",
         "sub": "effect size 另算\n（Cohen's d / 差異百分比）"},
        {"text": "② 不告訴你差異在業務上重不重要",
         "sub": "業務判斷另做\n（成本 / 規模 / 風險三問）"},
        {"text": "③ 不告訴你有沒有因果",
         "sub": "需實驗設計另證\n（隨機分派 / RCT）"},
    ], top=1.4, bottom=2.5)
    draw_inverted_thesis_box(
        s, "p 小不等於重要，p 小不等於有因果。p 只說「這不太像巧合」。",
        y=5.8, width=11.0)
    add_source(s, "ASA Statement on p-values, 2016")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # S12 — RISK-MITIGATION · A/B 四坑 × 四緩解
    s = _blank(prs)
    draw_risk_mitigation(s,
        risks=[
            "Peeking — 每天看一次就想停，偽陽性暴增",
            "SRM — 分組比例偏離 50/50，分流系統有 bug",
            "MDE 未設 — 沒定義最小可見差異，結果無從判讀",
            "Multiple testing — 同時測十個指標，總有一個假贏",
        ],
        mitigations=[
            "預設樣本量跑完再看，或用 sequential testing",
            "每日 chi-square SRM 檢查，偏離立刻停",
            "開跑前寫下 MDE 與所需樣本量（power ≥ 0.8）",
            "Bonferroni / FDR 校正或預設主指標",
        ],
        risks_title="四坑 Pitfalls",
        miti_title="四緩解 Mitigations",
        title="A/B 測試四個坑 × 四個緩解",
        summary="每個坑都有解；四對八招是 A/B 的最小裝備。",
    )
    add_source(s, "Kohavi et al. 2020, Trustworthy Online Controlled Experiments")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # S13 — TABLE · 誠實度檢查清單
    s = _blank(prs)
    draw_editorial_table(s,
        header=["☐", "自查項", "不過關的後果"],
        rows=[
            ["☐", "Y 軸是否從 0 或合理基準起算？", "2% 被畫成懸崖，主管誤判"],
            ["☐", "是否挑了對自己有利的時間窗？", "看似改善，實為季節性"],
            ["☐", "顏色是否靠紅綠分類？", "8% 觀眾無法解讀"],
            ["☐", "是否只秀平均、沒秀分佈？", "離群值 / 雙峰被隱藏"],
            ["☐", "標題是結論句還是目錄句？", "「銷售分析」≠「北區下滑 12%」"],
            ["☐", "是否標樣本數 n 與資料期間？", "無法評估可靠度"],
            ["☐", "p 值旁有沒有 effect size？", "小差異 + 大樣本 = 假重要"],
        ],
        col_widths=[0.3, 2.0, 2.1],
        title="報告送出前的誠實度檢查清單",
    )
    draw_inverted_thesis_box(s, "這七格都要打勾，再按寄送。", y=5.95, width=9.0)
    add_source(s, "本課整合 Cairo 2016 + Kohavi 2020")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # S14 — ASK · 你手上最近一張圖，Y 軸從幾開始？
    s = _blank(prs)
    draw_ask_page(
        s,
        "你的 Y 軸 從幾開始？",
        data_card={
            "label": "回呼 Slide 6",
            "stat": "0 ?",
            "caption": "翻開手邊最近一份報告的第一張圖，Y 軸從幾開始？若不是 0，你有理由嗎？",
        },
    )
    add_source(s, "回到 Slide 6")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # S15 — SILENT · 金句主張
    s = _blank(prs)
    draw_silent_page(s, "一張圖 一個主張\n要證據")
    add_footer(s, MODULE_CODE, 15, N_CONTENT, dark_bg=True)

    # S16 — SILENT · 倫理收尾（白底版本）
    s = _blank(prs)
    # 白底 + 深綠字 (SILENT 的白底版本：手刻以保留 Editorial 收尾紀律)
    add_textbox(
        s, Inches(0.8), Inches(2.6),
        T.SLIDE_W - Inches(1.6), Inches(2.2),
        "誠實 是唯一的底線",
        font_size=Pt(44), color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3,
    )
    add_source(
        s,
        "延伸閱讀：Cairo《The Truthful Art》 · Kohavi 2020 · Okabe-Ito palette · ASA 2016 Statement",
    )
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
