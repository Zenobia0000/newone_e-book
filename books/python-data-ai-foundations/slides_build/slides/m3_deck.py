"""M3 deck — 16 content slides + cover + copyright page.

Chart functions required in charts.py (document only — do not edit charts.py here):

  hbar_log_m3s2(outfile: Path = None) -> Path
    - Horizontal bar chart, log-scale X axis (1 .. 10,000 ms).
    - Two bars: `Python for 迴圈` 1800 ms (baseline), `NumPy 向量化 (np.sum)` 12 ms (-99.3%).
    - Both bars solid PRIMARY (#1B5E3F), no border, no shadow.
    - X label: "執行時間 (ms)，對數刻度"; no gridlines; direct value labels at bar tips.
    - Subtle note bottom-right: "資料：1e6 float64 / M2 Pro 單執行緒".
    - figsize ~ (10.5, 4.2), 220 dpi, white facecolor.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from .. import charts
from ..primitives import (
    add_rect, add_textbox, add_title, add_source, set_solid_fill, set_no_line,
    set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_pyramid_stack, draw_inverted_thesis_box,
    draw_vs_two_col, draw_risk_mitigation, draw_dual_track, draw_grid,
    draw_code_panel, draw_split_panel, draw_three_blocks_flow,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "M3"
MODULE_TITLE = "NumPy 與 pandas"
MODULE_SUBTITLE = "Python 數據分析與 AI 基礎 · M3"
TIME_MIN = 40
N_CONTENT = 16


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_m3(output_path, image_registry=None):
    """Build M3 deck.

    image_registry: optional list collector for image placeholder metadata.
                    M3 has no PHOTO prototypes — registry unused but accepted
                    for build-pipeline symmetry with build_m0 / build_m1.
    """
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE, TIME_MIN, N_CONTENT)

    # S1 — ASK (Hook)
    s = _blank(prs)
    draw_ask_page(
        s,
        "100 萬列資料，你打算用 for 迴圈處理嗎？",
        data_card={
            "label": "1e6 float64 加總 · 單執行緒",
            "stat": "≈ 150×",
            "caption": "for 迴圈 1,800 ms  →  向量化 12 ms",
        },
    )
    add_source(s, "本課自擬 numpy 1e6 float64 加總微基準")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # S2 — CHART (horizontal bar, log scale)
    s = _blank(prs)
    add_title(s, "100 萬列加總：for 迴圈 1,800 ms，向量化 12 ms，差 150 倍")
    chart_png = charts.hbar_log_m3s2()
    s.shapes.add_picture(str(chart_png), Inches(0.8), Inches(1.3),
                          width=Inches(11.7))
    add_textbox(s, T.MARGIN_X, Inches(5.6),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "向量化省 99.3% 時間",
                font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "本課 benchmark，%.repeat=7 取中位 / 資料 1e6 float64 / 機器 M2 Pro 單執行緒")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # S3 — GEOMETRIC-DIAGRAM: list 散落指標 vs ndarray 連續記憶體
    s = _blank(prs)
    draw_split_panel(s,
        left={
            "label": "Python list：8 個 float",
            "style": "wrong",
            "cells": [
                {"text": "[ ptr ptr ptr ptr ptr ptr ptr ptr ]"},
                {"text": "↓  ↓  散落的 PyFloat 物件  ↓  ↓"},
                {"text": "記憶體：約 224 bytes"},
            ],
            "note": "每取一值要跳一次記憶體，CPU cache 不友善。",
        },
        right={
            "label": "NumPy ndarray：8 個 float64",
            "cells": [
                {"text": "■■■■■■■■", "highlight": True},
                {"text": "連續 8×8 = 64 bytes"},
                {"text": "記憶體：64 bytes (-71%)"},
            ],
            "note": "連續記憶體，CPU 一次抓一排——向量化的物理基礎。",
        },
        title="ndarray 是一塊連續記憶體，list 是一串指標",
    )
    add_source(s, "CPython 3.12 / NumPy 2.0 記憶體量測")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # S4 — GEOMETRIC-DIAGRAM: Broadcasting 四步驟
    s = _blank(prs)
    draw_flow_chain(s, nodes=[
        {"label": "Step 1 · 對齊",
         "caption": "(3,4) 與 (4,)\n右對齊比較"},
        {"label": "Step 2 · 補 1",
         "caption": "(4,) → (1,4)\n前補 1"},
        {"label": "Step 3 · 虛擬擴張",
         "caption": "(1,4) → (3,4)\n不真複製"},
        {"label": "Step 4 · 運算",
         "caption": "逐元素相加\n得 (3,4)", "highlight": True},
    ], title="Broadcasting 四步驟：從右對齊，補 1 擴張，不複製",
       y=3.2)
    draw_inverted_thesis_box(s,
        "規則：從右對齊，維度相同或其一為 1 才能廣播。",
        y=6.0, width=10.5)
    add_source(s, "NumPy 2.0 官方文件 Broadcasting rules")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # S5 — TABLE: DataFrame 三層結構
    s = _blank(prs)
    draw_editorial_table(s,
        header=["index", "employee_id", "hire_date", "salary", "is_manager"],
        rows=[
            ["0", "E001", "2020-03-15", "68000.0", "True"],
            ["1", "E002", "2019-07-01", "92000.0", "False"],
            ["2", "E003", "2021-11-20", "55000.0", "False"],
            ["3", "E004", "2018-01-08", "NaN", "True"],
            ["dtype", "int64 / object", "datetime64[ns]", "float64", "bool"],
        ],
        col_widths=[0.6, 1.0, 1.1, 0.9, 0.9],
        title="DataFrame 三層結構：index / columns / dtype，每欄各自有型別",
    )
    add_textbox(s, T.MARGIN_X, Inches(5.6),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
                "index 是身份證不是資料欄；salary 出現 NaN → 整欄升級 float64。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "pandas 2.2 官方文件")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # S6 — BEFORE/AFTER: view vs copy
    s = _blank(prs)
    draw_split_panel(s,
        left={
            "label": "BEFORE · view（共享記憶體）",
            "cells": [
                {"text": "s = df['salary']"},
                {"text": "id(s._values) == id(df['salary']._values)", "highlight": True},
                {"text": "寫 s 會改到 df"},
            ],
            "note": "同一個物件，兩個名字——改一邊另一邊跟著改。",
        },
        right={
            "label": "AFTER · copy（獨立記憶體）",
            "cells": [
                {"text": "s = df['salary'].copy()"},
                {"text": "id(s._values) ≠ id(df['salary']._values)"},
                {"text": "寫 s 不影響 df"},
            ],
            "note": "彼此獨立；id() 是你的偵測器。",
        },
        title="同一欄切片：view 共享記憶體，copy 各自獨立",
    )
    add_source(s, "pandas 2.2 indexing-copy-view")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # S7 — TABLE: pandas 1.5 / 2.0 / 2.2
    s = _blank(prs)
    draw_editorial_table(s,
        header=["特性", "pandas 1.5 (2022)", "pandas 2.0 (2023)", "pandas 2.2 (2024)"],
        rows=[
            ["鏈式賦值語意", "未定義（有時成功）", "CoW opt-in", "CoW 預設啟用"],
            ["底層儲存", "僅 NumPy", "NumPy 或 Arrow", "NumPy 或 Arrow"],
            ["字串欄位記憶體", "基準", "-35% (Arrow)", "-35% (Arrow)"],
            ["整數欄 + NaN", "升級為 float64", "可用 Int64 原生可空", "可用 Int64 原生可空"],
            ["SettingWithCopyWarning", "常見", "仍出現", "CoW 下大幅減少"],
            ["典型資料集總記憶體", "1.0× 基準", "0.65× (-35%)", "0.60× (-40%)"],
        ],
        col_widths=[1.3, 1.2, 1.2, 1.2],
        title="pandas 版本三代對照：1.5 / 2.0 / 2.2 的關鍵差異",
    )
    draw_inverted_thesis_box(s,
        "2.2 起 CoW 預設啟用——你的舊程式碼可能需要檢查。",
        y=6.15, width=10.5)
    add_source(s, "pandas release notes 1.5 / 2.0 / 2.2")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # S8 — GEOMETRIC-DIAGRAM: SettingWithCopyWarning 根因
    s = _blank(prs)
    draw_three_blocks_flow(s,
        blocks=[
            {"heading": "Step 1 · 先 filter",
             "items": ["df[df['age'] > 30]",
                      "→ 可能回傳 view",
                      "→ 也可能 copy",
                      "pandas 無法保證"]},
            {"heading": "Step 2 · 再 ['salary'] = 0",
             "items": ["__setitem__",
                      "若前步是 copy",
                      "寫入被丟棄",
                      "pandas 發出警告"]},
            {"heading": "結果 · 未定義行為",
             "items": ["SettingWithCopyWarning",
                      "有時成功、有時失敗",
                      "解法：.loc 一次存取",
                      ".loc[mask, 'salary'] = 0"]},
        ],
        title="SettingWithCopyWarning 根因：一條有歧義的存取路徑",
        bottom_note="解法只有一招：把兩步路徑合併成一次 .loc 存取。",
    )
    add_source(s, "pandas 官方 Indexing and Selecting Data")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # S9 — GEOMETRIC-DIAGRAM: groupby split → apply → combine
    s = _blank(prs)
    draw_three_blocks_flow(s,
        blocks=[
            {"heading": "Split · 切開",
             "items": ["df 6 列 region / revenue",
                      "group N (2 列)",
                      "group S (2 列)",
                      "group E (2 列)"]},
            {"heading": "Apply · 各算",
             "items": ["每組獨立 sum()",
                      "N: 3,200",
                      "S: 2,100",
                      "E: 4,500"]},
            {"heading": "Combine · 合回",
             "items": ["兩欄 TABLE",
                      "region / revenue_sum",
                      "N 3,200 · S 2,100 · E 4,500",
                      "回到一張表"]},
        ],
        title="groupby 三步驟：split → apply → combine",
        bottom_note="df.groupby('region')['revenue'].sum() —— 業務語言即 groupby 語法。",
    )
    add_source(s, "Wickham 2011 Split-Apply-Combine")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # S10 — BEFORE/AFTER: merge 兩窄表 → 一寬表
    s = _blank(prs)
    add_title(s, "merge：兩張窄表 → 一張寬表，鍵對齊決定一切")
    # BEFORE: two small tables side by side
    add_textbox(s, T.MARGIN_X, Inches(1.2),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.35),
                "BEFORE · on='customer_id'",
                font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True)
    # Left table orders
    left_x = T.MARGIN_X
    left_w = Inches(5.8)
    _mini_table(s, left_x, Inches(1.6), left_w,
                header=["order_id", "customer_id", "amount"],
                rows=[["O101", "C01", "520"],
                      ["O102", "C02", "890"],
                      ["O103", "C03", "340"],
                      ["O104", "C99", "210"]],
                caption="orders (4 列)")
    # Right table customers
    right_x = T.MARGIN_X + left_w + Inches(0.3)
    right_w = T.SLIDE_W - 2 * T.MARGIN_X - left_w - Inches(0.3)
    _mini_table(s, right_x, Inches(1.6), right_w,
                header=["customer_id", "name", "city"],
                rows=[["C01", "Alice", "Taipei"],
                      ["C02", "Bob", "Kaohsiung"],
                      ["C03", "Cara", "Taichung"]],
                caption="customers (3 列)")
    # AFTER: combined table
    add_textbox(s, T.MARGIN_X, Inches(4.3),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.35),
                "AFTER · 合併 5 欄",
                font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True)
    _mini_table(s, T.MARGIN_X, Inches(4.7), T.SLIDE_W - 2 * T.MARGIN_X,
                header=["order_id", "customer_id", "amount", "name", "city"],
                rows=[["O101", "C01", "520", "Alice", "Taipei"],
                      ["O102", "C02", "890", "Bob", "Kaohsiung"],
                      ["O103", "C03", "340", "Cara", "Taichung"],
                      ["O104", "C99", "210", "NaN", "NaN"]],
                caption="左表 4 列保留、右表缺鍵者為 NaN")
    add_source(s, "pandas 2.2 merge how='left'")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # S11 — BEFORE/AFTER: pivot 長表 → 寬表
    s = _blank(prs)
    add_title(s, "pivot：長表轉寬表，一列多指標變欄位")
    add_textbox(s, T.MARGIN_X, Inches(1.2),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.35),
                "BEFORE · long 9 列",
                font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True)
    _mini_table(s, T.MARGIN_X, Inches(1.6), Inches(5.8),
                header=["region", "quarter", "revenue"],
                rows=[["N", "Q1", "1200"], ["N", "Q2", "1500"], ["N", "Q3", "1800"],
                      ["S", "Q1", "900"], ["S", "Q2", "1100"], ["S", "Q3", "1300"],
                      ["E", "Q1", "1600"], ["E", "Q2", "1900"], ["E", "Q3", "2100"]],
                caption="region × quarter × revenue")
    add_textbox(s, T.MARGIN_X + Inches(6.2), Inches(3.2),
                Inches(6.0), Inches(0.5),
                "df.pivot(index='region',\n   columns='quarter', values='revenue')",
                font_size=T.FONT_SMALL, color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER)
    _mini_table(s, T.MARGIN_X + Inches(6.3), Inches(4.3),
                T.SLIDE_W - 2 * T.MARGIN_X - Inches(6.3),
                header=["region", "Q1", "Q2", "Q3"],
                rows=[["N", "1200", "1500", "1800"],
                      ["S", "900", "1100", "1300"],
                      ["E", "1600", "1900", "2100"]],
                caption="AFTER · wide 3 列（Q1/Q2/Q3 由 quarter 值攤開）")
    add_source(s, "pandas 2.2 reshape pivot")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # S12 — MATRIX 2×2: 小中大資料 × SQL 偏好
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "DuckDB + pandas",
         "sub": "資料 < 1 GB · SQL 風\nSQL 語法、in-process、零伺服器"},
        {"text": "DuckDB 獨立",
         "sub": "1–100 GB · SQL 風\nGB 到百 GB 單機 SQL"},
        {"text": "pandas 2.2（本課主場）",
         "sub": "資料 < 1 GB · Python 風\nin-memory、生態最成熟"},
        {"text": "Polars 1.x（建議演進）",
         "sub": "1–100 GB · Python 風\nRust · lazy 評估 · 5–20× pandas",
         "highlight": True},
    ], title="小中大資料 × SQL 偏好，四象限選工具")
    add_textbox(s, T.MARGIN_X, Inches(6.35),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
                "橫軸：資料量 < 1 GB ↔ 1–100 GB    縱軸：團隊 SQL 偏好 Python API ↔ SQL",
                font_size=T.FONT_SOURCE, color=T.GRAY_MID,
                align=PP_ALIGN.CENTER)
    add_source(s, "Polars 1.0 公告 2024 / DuckDB 1.0 release")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # S13 — VS: SQL ↔ pandas
    s = _blank(prs)
    add_title(s, "SQL ↔ pandas：同一個業務問題，兩種語言")
    draw_code_panel(s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=(T.SLIDE_W - 2 * T.MARGIN_X - Inches(0.3)) / 2, h=Inches(3.6),
        label="SQL",
        code=('SELECT region,\n'
              '       COUNT(*)     AS n,\n'
              '       AVG(revenue) AS avg_rev\n'
              'FROM   orders\n'
              'WHERE  revenue > 1000\n'
              'GROUP  BY region\n'
              'ORDER  BY avg_rev DESC\n'
              'LIMIT  3;'),
        bullets=["WHERE → 過濾",
                 "GROUP BY + AVG → 聚合",
                 "ORDER BY + LIMIT → 排序取頂"],
        label_dark=True,
    )
    right_x = T.MARGIN_X + (T.SLIDE_W - 2 * T.MARGIN_X - Inches(0.3)) / 2 + Inches(0.3)
    draw_code_panel(s,
        x=right_x, y=Inches(1.3),
        w=(T.SLIDE_W - 2 * T.MARGIN_X - Inches(0.3)) / 2, h=Inches(3.6),
        label="pandas",
        code=('(orders\n'
              "  .query('revenue > 1000')\n"
              "  .groupby('region')\n"
              "  .agg(n=('revenue', 'size'),\n"
              "       avg_rev=('revenue', 'mean'))\n"
              "  .sort_values('avg_rev',\n"
              "               ascending=False)\n"
              '  .head(3))'),
        bullets=[".query() ≡ WHERE",
                 ".groupby().agg() ≡ GROUP BY",
                 ".sort_values().head() ≡ ORDER BY + LIMIT"],
        label_dark=True,
    )
    # Mapping mini-table
    _mini_table(s, T.MARGIN_X + Inches(1.5), Inches(5.2),
                T.SLIDE_W - 2 * T.MARGIN_X - Inches(3.0),
                header=["SQL 子句", "pandas 對應", "順序"],
                rows=[["WHERE", ".query() 或 bool indexing", "1"],
                      ["GROUP BY + agg", ".groupby().agg()", "2"],
                      ["ORDER BY + LIMIT", ".sort_values().head()", "3"]],
                caption=None)
    add_source(s, "本課整理，對應 pandas 2.2 / ANSI SQL")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # S14 — PYRAMID compression pack: 4 cards + inverted thesis
    s = _blank(prs)
    draw_grid(s, rows=1, cols=4, cells=[
        {"label": "① ndarray",
         "sub": "連續記憶體、同型別\n-71% vs list"},
        {"label": "② broadcasting",
         "sub": "從右對齊、補 1、不複製\n(X - X.mean(0)) 一行"},
        {"label": "③ groupby",
         "sub": "split → apply → combine\n業務問題的通用骨架"},
        {"label": "④ pandas 2.0",
         "sub": "CoW + Arrow\n-40% 記憶體", "highlight": True},
    ], title="M3 壓縮包：四件事，一句話",
       top=1.4, bottom=2.0)
    draw_inverted_thesis_box(s,
        "向量化不是為了快，是為了用對的語言描述計算。",
        y=6.0, width=11.0)
    add_source(s, "本模組 Slide 2–13 收束")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # S15 — ASK (Feel bridge)
    s = _blank(prs)
    draw_ask_page(
        s,
        "你現在手上那段 for 迴圈，能不能寫成一行？",
        data_card={
            "label": "典型轉譯",
            "stat": "N → 1",
            "caption": "for i in range(len(arr)):\n→  arr.sum()",
        },
    )
    add_source(s, "回到 Slide 2")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # S16 — SILENT (Feel close)
    s = _blank(prs)
    draw_silent_page(s, "向量化\n是對的語言，不只是快的語言。")
    add_footer(s, MODULE_CODE, 16, N_CONTENT, dark_bg=True)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path


# ---------------------------------------------------------------------------
# Local helpers (scoped to m3_deck; mini table for BEFORE/AFTER composites).
# ---------------------------------------------------------------------------

def _mini_table(slide, x, y, w, header, rows, caption=None):
    """Editorial mini table: top/bottom borders, green header, alt-row banding."""
    ncols = len(header)
    col_w = w / ncols
    header_h = Inches(0.36)
    row_h = Inches(0.32)

    # Top border
    top_line = add_rect(slide, x, y, w, Inches(0.015))
    set_solid_fill(top_line, T.PRIMARY)
    set_no_line(top_line)

    # Header band
    hdr_y = y + Inches(0.015)
    hdr = add_rect(slide, x, hdr_y, w, header_h)
    set_solid_fill(hdr, T.PRIMARY)
    set_no_line(hdr)
    for i, h in enumerate(header):
        add_textbox(slide, x + i * col_w + Inches(0.06), hdr_y,
                    col_w - Inches(0.12), header_h,
                    h, font_size=T.FONT_SMALL, color=T.WHITE, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)

    cy = hdr_y + header_h
    for r, row in enumerate(rows):
        if r % 2 == 1:
            band = add_rect(slide, x, cy, w, row_h)
            set_solid_fill(band, T.TABLE_ALT)
            set_no_line(band)
        for i, val in enumerate(row):
            add_textbox(slide, x + i * col_w + Inches(0.06), cy,
                        col_w - Inches(0.12), row_h,
                        str(val), font_size=T.FONT_SMALL, color=T.CHARCOAL,
                        anchor=MSO_ANCHOR.MIDDLE)
        cy += row_h

    bot_line = add_rect(slide, x, cy, w, Inches(0.015))
    set_solid_fill(bot_line, T.PRIMARY)
    set_no_line(bot_line)

    if caption:
        add_textbox(slide, x, cy + Inches(0.05), w, Inches(0.3),
                    caption, font_size=T.FONT_SOURCE, color=T.GRAY_MID,
                    italic=True)
