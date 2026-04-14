"""M0 stability test suite.

Five categories of invariants:
  A. Schema — 00_skeleton.yaml + 04_image_placeholders.yaml parse & structure
  B. Build — `python -m slides_build.build --module M0` produces .pptx cleanly
  C. Structure — slide count, cover + copyright page, logo per slide, footer format
  D. Visual discipline — color palette locked, no forbidden content
  E. Content — skeleton thesis appears in deck; course-scope consistency (24h × 8 blocks)

Run:
  cd books/python-data-ai-foundations
  python -m slides_build.tests.test_m0
"""
from __future__ import annotations

import sys
import re
import subprocess
from pathlib import Path
from typing import Iterable

import yaml
from pptx import Presentation
from pptx.util import Emu

from .. import theme as T


MODULE_CODE = "M0"
CHAPTERS = Path(__file__).resolve().parents[2] / "chapters" / f"{MODULE_CODE}_annotated"
SKELETON = CHAPTERS / "00_skeleton.yaml"
PLACEHOLDERS = CHAPTERS / "04_image_placeholders.yaml"
SLIDES_DESIGN = CHAPTERS / "02_slides_design.md"
MVK = CHAPTERS / "05_mvk.md"

OUTPUT_DIR = T.OUTPUT_DIR
PPTX_PATH = OUTPUT_DIR / "M0_開場_Python與AI系統全景.pptx"
AUDIT_YAML = OUTPUT_DIR / "_image_placeholders_M0.yaml"

ALLOWED_COLORS = {
    (0x1B, 0x5E, 0x3F),  # PRIMARY
    (0x33, 0x33, 0x33),  # CHARCOAL
    (0x80, 0x80, 0x80),  # GRAY_MID
    (0xD3, 0xD3, 0xD3),  # LIGHT_GRAY
    (0xF0, 0xF0, 0xF0),  # TABLE_ALT
    (0xF7, 0xF7, 0xF7),  # placeholder bg
    (0xFF, 0xFF, 0xFF),  # WHITE
    (0, 0, 0),           # occasional default black (ignored but tolerated)
}

# For palette audit we scan fill colors only; text colors default to theme
FORBIDDEN_TOKENS = [
    "M8", "M9 ", "M10", "M11", "M12", "M13", "M14",
    "M15", "M16", "M17", "M18", "M19", "M20",
    "stock photo",   # Editorial禁令
]
# Patterns using regex to avoid substring false positives
FORBIDDEN_PATTERNS = [
    (r"(?<!\d)4 小時", "leftover '4 小時' (should be '24 小時')"),
]

REQUIRED_PROTOTYPES = {
    "SILENT", "ASK", "CHART", "MATRIX", "GEOMETRIC-DIAGRAM",
    "BEFORE-AFTER", "PYRAMID", "TABLE", "RISK-MITIGATION", "VS", "PHOTO",
}


# ---------------- helpers ----------------

class Reporter:
    def __init__(self):
        self.passed = 0
        self.failed: list[tuple[str, str]] = []

    def ok(self, name: str):
        self.passed += 1
        print(f"  \033[92m✓\033[0m {name}")

    def fail(self, name: str, reason: str):
        self.failed.append((name, reason))
        print(f"  \033[91m✗\033[0m {name}")
        for line in reason.splitlines():
            print(f"      {line}")

    def summary(self) -> bool:
        total = self.passed + len(self.failed)
        if self.failed:
            print(f"\n{len(self.failed)}/{total} FAILED")
            return False
        print(f"\n{self.passed}/{total} PASSED")
        return True


def _slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text:
                        parts.append(run.text)
    return "\n".join(parts)


def _slide_text_all(prs: Presentation) -> str:
    return "\n\n".join(_slide_text(s) for s in prs.slides)


def _shape_has_picture(slide) -> bool:
    return any(
        shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
        for shape in slide.shapes
    )


def _collect_solid_fill_rgb(prs: Presentation) -> set[tuple[int, int, int]]:
    colors = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if shape.fill.type == 1:  # MSO_FILL.SOLID
                    rgb = shape.fill.fore_color.rgb
                    colors.add((rgb[0], rgb[1], rgb[2]))
            except Exception:
                continue
    return colors


# ---------------- categories ----------------

def test_schema(r: Reporter):
    print("\n[A] Schema validation")
    # Skeleton
    try:
        skel = yaml.safe_load(SKELETON.read_text(encoding="utf-8"))
        assert skel["module"] == MODULE_CODE
        assert skel.get("governing_thought")
        assert skel.get("narrative_arc") and set(skel["narrative_arc"]) == {
            "hook", "tension", "reveal", "ground", "feel",
        }
        slides = skel.get("slides", [])
        assert len(slides) == 15, f"expected 15 slides, got {len(slides)}"
        for s in slides:
            assert s["index"] in range(1, 16)
            assert isinstance(s["thesis"], str) and len(s["thesis"]) > 5
            assert s["prototype"] in REQUIRED_PROTOTYPES
        r.ok(f"00_skeleton.yaml parses · 15 slides · narrative_arc complete")
    except Exception as e:
        r.fail("00_skeleton.yaml structure", str(e))
        return

    # Narrative arc covers every slide exactly once
    arc = skel["narrative_arc"]
    covered = []
    for seg in arc.values():
        covered.extend(seg)
    if sorted(covered) == list(range(1, 16)):
        r.ok("narrative_arc covers slides 1-15 with no gaps / dupes")
    else:
        r.fail("narrative_arc coverage", f"got sorted: {sorted(covered)}")

    # Prototype diversity (≥7 per §5 Layer B)
    protos = {s["prototype"] for s in slides}
    if len(protos) >= 7:
        r.ok(f"prototype diversity ≥ 7 (actual: {len(protos)})")
    else:
        r.fail("prototype diversity", f"only {len(protos)} used: {protos}")

    # No consecutive same prototype
    bad = []
    for i in range(len(slides) - 1):
        if slides[i]["prototype"] == slides[i + 1]["prototype"]:
            bad.append((slides[i]["index"], slides[i + 1]["index"]))
    if not bad:
        r.ok("no two adjacent slides share a prototype")
    else:
        r.fail("adjacent prototype check", f"dupes at: {bad}")

    # Placeholders yaml
    try:
        ph = yaml.safe_load(PLACEHOLDERS.read_text(encoding="utf-8"))
        assert ph["module"] == MODULE_CODE
        assert ph["total"] == len(ph.get("placeholders", []))
        for p in ph["placeholders"]:
            assert p["status"] in {"pending", "resolved", "skipped"}
            bb = p["bounding_box"]
            assert all(k in bb for k in ("x", "y", "w", "h"))
        r.ok(f"04_image_placeholders.yaml parses · {ph['total']} placeholder(s)")
    except Exception as e:
        r.fail("04_image_placeholders.yaml structure", str(e))


def test_build(r: Reporter):
    print("\n[B] Build pipeline")
    # Remove stale output to force fresh build
    if PPTX_PATH.exists():
        PPTX_PATH.unlink()
    if AUDIT_YAML.exists():
        AUDIT_YAML.unlink()

    proc = subprocess.run(
        [sys.executable, "-m", "slides_build.build", "--module", MODULE_CODE],
        cwd=str(Path(__file__).resolve().parents[2]),
        capture_output=True, text=True,
    )
    if proc.returncode != 0:
        r.fail("build exit code", f"stderr:\n{proc.stderr}")
        return
    r.ok("build completes with exit code 0")

    if not PPTX_PATH.exists():
        r.fail("pptx output", f"expected at {PPTX_PATH}")
        return
    r.ok(f"pptx produced ({PPTX_PATH.stat().st_size} bytes)")

    if not AUDIT_YAML.exists():
        r.fail("audit yaml", f"expected at {AUDIT_YAML}")
    else:
        r.ok(f"audit yaml emitted at {AUDIT_YAML.name}")


def test_structure(r: Reporter, prs: Presentation):
    print("\n[C] Structural integrity")
    n = len(prs.slides)
    expected = 1 + 15 + 1  # cover + 15 content + copyright
    if n == expected:
        r.ok(f"total slide count = {n} (cover + 15 content + copyright)")
    else:
        r.fail("slide count", f"expected {expected}, got {n}")

    # Cover (slide 0)
    cover_text = _slide_text(prs.slides[0])
    if "桑尼資料科學" in cover_text and "M0" in cover_text and "開場" in cover_text:
        r.ok("cover has brand + module code + title")
    else:
        r.fail("cover content", f"text snapshot:\n{cover_text[:200]}")

    if _shape_has_picture(prs.slides[0]):
        r.ok("cover contains logo image")
    else:
        r.fail("cover logo", "no picture shape on cover")

    # 15 content slides
    logo_missing = []
    footer_missing = []
    for i in range(1, 16):
        slide = prs.slides[i]
        if not _shape_has_picture(slide):
            logo_missing.append(i)
        txt = _slide_text(slide)
        if f"M0 · {i}/15" not in txt:
            footer_missing.append(i)
    if not logo_missing:
        r.ok("every content slide has logo")
    else:
        r.fail("per-slide logo", f"missing on: {logo_missing}")

    if not footer_missing:
        r.ok("every content slide has footer 'M0 · N/15'")
    else:
        r.fail("per-slide footer", f"missing or wrong on: {footer_missing}")

    # Copyright page
    cr_text = _slide_text(prs.slides[-1])
    if "版權" in cr_text and "All rights reserved" in cr_text:
        r.ok("copyright page present with full wording")
    else:
        r.fail("copyright page", f"snapshot:\n{cr_text[:200]}")


def test_visual_discipline(r: Reporter, prs: Presentation):
    print("\n[D] Visual discipline (§10.2 contract)")
    colors = _collect_solid_fill_rgb(prs)
    unexpected = colors - ALLOWED_COLORS
    if not unexpected:
        r.ok(f"fill palette locked to 7 allowed colors (found: {len(colors)})")
    else:
        r.fail(
            "palette lock",
            f"unexpected RGB fills: {sorted(unexpected)}\n"
            f"allowed: {sorted(ALLOWED_COLORS)}",
        )

    # Scan all text for forbidden tokens + patterns
    all_text = _slide_text_all(prs)
    hits = [tok for tok in FORBIDDEN_TOKENS if tok in all_text]
    pattern_hits = [
        label for pat, label in FORBIDDEN_PATTERNS
        if re.search(pat, all_text)
    ]
    if not hits and not pattern_hits:
        r.ok("no forbidden tokens / patterns (stale module refs, '4 小時', stock photo)")
    else:
        r.fail("forbidden tokens", f"literal: {hits} · patterns: {pattern_hits}")


def test_content_consistency(r: Reporter, prs: Presentation):
    print("\n[E] Content consistency")
    all_text = _slide_text_all(prs)

    # Each skeleton thesis should surface somewhere in deck (fuzzy: first 10 chars)
    skel = yaml.safe_load(SKELETON.read_text(encoding="utf-8"))
    missing = []
    for s in skel["slides"]:
        # Thesis key substring (first line up to punctuation)
        key = re.split(r"[；，。？！]", s["thesis"])[0][:14]
        if key and key not in all_text:
            missing.append((s["index"], s["prototype"], key))
    if not missing:
        r.ok("all 15 skeleton thesis keys surface in deck")
    else:
        r.fail(
            "thesis coverage",
            "\n".join(f"slide {i} ({p}): key '{k}' not found" for i, p, k in missing),
        )

    # Course scope: "24 小時" must appear
    occurrences = all_text.count("24 小時")
    if occurrences >= 3:
        r.ok(f"'24 小時' appears {occurrences}× (S2/S9/S10 + optional)")
    else:
        r.fail("24h label present", f"only {occurrences} occurrence(s)")

    # S10 八格 hours sum to 24 — take FIRST occurrence per module
    # (caption also contains M{x}·{h}h; de-dup by module keeps only grid cells)
    s10 = _slide_text(prs.slides[10])
    seen: dict[int, int] = {}
    for m, h in re.findall(r"M(\d)\s*·\s*(\d)h", s10):
        seen.setdefault(int(m), int(h))
    if set(seen) == set(range(8)):
        total = sum(seen.values())
        if total == 24:
            r.ok(f"S10 八格 M0–M7 小時數加總 = 24h "
                 f"({sorted(seen.items())})")
        else:
            r.fail("S10 總時數", f"加總 = {total}h, expected 24h, 分布 {seen}")
    else:
        r.fail(
            "S10 八格模組標籤",
            f"找到模組 {sorted(seen)}, expected {{0..7}}",
        )

    # S2/S9 must NOT say "4 小時" (leftover from earlier erroneous rewrite)
    s2 = _slide_text(prs.slides[2])
    s9 = _slide_text(prs.slides[9])
    for name, txt in [("S2", s2), ("S9", s9)]:
        if "24 小時" in txt and "4 小時押" not in txt:
            r.ok(f"{name} uses '24 小時' correctly")
        elif "24 小時" not in txt:
            r.fail(f"{name} 24h label", f"no '24 小時' found")


def test_placeholder_system(r: Reporter, prs: Presentation):
    print("\n[F] Placeholder system (S14)")
    audit = yaml.safe_load(AUDIT_YAML.read_text(encoding="utf-8"))
    if audit["total"] == 3 and audit["pending"] == 3:
        r.ok("audit YAML: 3 pending placeholders")
    else:
        r.fail("audit counts", f"total={audit['total']}, pending={audit['pending']}")

    # Each placeholder should have required fields
    required = {"id", "slot_name", "description", "source_url",
                "recommended_size", "bounding_box", "status"}
    for p in audit["placeholders"]:
        miss = required - set(p)
        if miss:
            r.fail(f"placeholder {p.get('id', '?')} fields", f"missing: {miss}")
            return
    r.ok("all placeholders have required schema fields")

    # S14 rendering: should contain '待補真圖' × 3
    s14_text = _slide_text(prs.slides[14])
    todo_count = s14_text.count("待補真圖")
    if todo_count == 3:
        r.ok("S14 renders 3 '待補真圖' placeholder headers")
    else:
        r.fail("S14 placeholder count", f"expected 3, got {todo_count}")

    # URL hints present
    urls = ["pandas.pydata.org", "scikit-learn.org", "pytorch.org"]
    missing_urls = [u for u in urls if u not in s14_text]
    if not missing_urls:
        r.ok("S14 shows all 3 documentation URLs")
    else:
        r.fail("S14 URL hints", f"missing: {missing_urls}")


def test_pdf_renders(r: Reporter):
    print("\n[G] PDF render (integration)")
    if not PPTX_PATH.exists():
        r.fail("pdf prerequisite", "pptx not found")
        return

    proc = subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf",
         "--outdir", str(OUTPUT_DIR), str(PPTX_PATH)],
        capture_output=True, text=True, timeout=60,
    )
    pdf_path = PPTX_PATH.with_suffix(".pdf")
    if proc.returncode == 0 and pdf_path.exists():
        r.ok(f"libreoffice renders PDF cleanly ({pdf_path.stat().st_size} bytes)")
    else:
        r.fail(
            "libreoffice render",
            f"rc={proc.returncode}\nstderr: {proc.stderr[:500]}",
        )


# ---------------- main ----------------

def main() -> int:
    r = Reporter()
    print(f"═══ M0 stability test ═══")
    print(f"chapters: {CHAPTERS}")
    print(f"pptx:     {PPTX_PATH}")

    test_schema(r)
    test_build(r)

    if not PPTX_PATH.exists():
        print("\nFATAL: pptx build failed — skipping downstream tests")
        return 1

    prs = Presentation(str(PPTX_PATH))
    test_structure(r, prs)
    test_visual_discipline(r, prs)
    test_content_consistency(r, prs)
    test_placeholder_system(r, prs)
    test_pdf_renders(r)

    return 0 if r.summary() else 1


if __name__ == "__main__":
    sys.exit(main())
