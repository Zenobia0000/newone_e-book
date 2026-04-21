"""Cover, footer, copyright page. All decks share these."""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from . import theme as T
from .primitives import (
    add_rect, add_textbox, set_solid_fill, set_no_line, set_text_font
)


def add_cover_slide(prs, module_code: str, module_title: str, subtitle: str,
                    time_min: int, slide_count: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full dark green background
    bg = add_rect(slide, 0, 0, T.SLIDE_W, T.SLIDE_H)
    set_solid_fill(bg, T.PRIMARY)
    set_no_line(bg)

    # Top band: brand text, small
    brand_tb = add_textbox(
        slide,
        T.MARGIN_X, T.MARGIN_Y,
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        f"{T.BRAND_NAME_ZH} · {T.BRAND_NAME_EN}",
        font_size=T.FONT_BODY, color=T.WHITE, align=PP_ALIGN.LEFT, bold=False,
    )

    # Center logo
    logo_w = T.LOGO_COVER_SIZE
    logo_x = (T.SLIDE_W - logo_w) / 2
    logo_y = Inches(1.6)
    slide.shapes.add_picture(str(T.LOGO_PATH), logo_x, logo_y, width=logo_w, height=logo_w)

    # Module code (big)
    code_y = logo_y + logo_w + Inches(0.3)
    add_textbox(
        slide,
        Inches(0), code_y, T.SLIDE_W, Inches(0.9),
        module_code,
        font_size=T.FONT_HERO, color=T.WHITE, align=PP_ALIGN.CENTER, bold=True,
    )

    # Module title
    add_textbox(
        slide,
        Inches(0), code_y + Inches(0.9), T.SLIDE_W, Inches(0.7),
        module_title,
        font_size=T.FONT_COVER_TITLE, color=T.WHITE, align=PP_ALIGN.CENTER, bold=True,
    )

    # Subtitle line
    add_textbox(
        slide,
        Inches(0), code_y + Inches(1.6), T.SLIDE_W, Inches(0.5),
        subtitle,
        font_size=T.FONT_COVER_SUB, color=T.WHITE, align=PP_ALIGN.CENTER, bold=False,
    )

    # Meta line
    meta = f"時長 {time_min} 分鐘 · 共 {slide_count} 張"
    add_textbox(
        slide,
        Inches(0), code_y + Inches(2.15), T.SLIDE_W, Inches(0.4),
        meta,
        font_size=T.FONT_BODY, color=T.LIGHT_GRAY, align=PP_ALIGN.CENTER, bold=False,
    )

    # Footer: short copyright + brand website placeholder
    foot_y = T.SLIDE_H - Inches(0.5)
    add_textbox(
        slide,
        T.MARGIN_X, foot_y, T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        f"{T.COPYRIGHT_SHORT} · 未經授權不得轉載",
        font_size=T.FONT_SMALL, color=T.LIGHT_GRAY, align=PP_ALIGN.CENTER, bold=False,
    )

    return slide


def add_footer(slide, module_code: str, slide_no: int, total: int, dark_bg: bool = False):
    """Add right-bottom logo + left-bottom copyright/pagination to a content slide."""
    fg = T.WHITE if dark_bg else T.CHARCOAL
    sub = T.LIGHT_GRAY if dark_bg else T.GRAY_MID

    # Left-bottom: copyright + pagination
    foot_y = T.SLIDE_H - Inches(0.35)
    add_textbox(
        slide,
        T.MARGIN_X, foot_y, Inches(6), Inches(0.25),
        f"{T.COPYRIGHT_SHORT} · {module_code} · {slide_no}/{total}",
        font_size=T.FONT_SOURCE, color=sub, align=PP_ALIGN.LEFT, bold=False,
    )

    # Right-bottom: logo
    logo_w = T.LOGO_CORNER_SIZE
    logo_x = T.SLIDE_W - logo_w - Inches(0.35)
    logo_y = T.SLIDE_H - logo_w - Inches(0.2)
    slide.shapes.add_picture(str(T.LOGO_PATH), logo_x, logo_y, width=logo_w, height=logo_w)


def add_copyright_slide(prs, module_code: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background, dark green top band
    band = add_rect(slide, 0, 0, T.SLIDE_W, Inches(0.7))
    set_solid_fill(band, T.PRIMARY)
    set_no_line(band)

    add_textbox(
        slide,
        T.MARGIN_X, Inches(0.15), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        f"版權與授權聲明 · {module_code} · Copyright & Licensing Notice",
        font_size=T.FONT_SUBTITLE, color=T.WHITE, align=PP_ALIGN.LEFT, bold=True,
    )

    # Full copyright body
    add_textbox(
        slide,
        Inches(1.2), Inches(1.4), T.SLIDE_W - Inches(2.4), Inches(5.2),
        T.COPYRIGHT_FULL,
        font_size=T.FONT_BODY, color=T.CHARCOAL, align=PP_ALIGN.LEFT, bold=False,
        line_spacing=1.45,
    )

    # Small logo bottom-right
    logo_w = Inches(0.7)
    slide.shapes.add_picture(
        str(T.LOGO_PATH),
        T.SLIDE_W - logo_w - Inches(0.6),
        T.SLIDE_H - logo_w - Inches(0.5),
        width=logo_w, height=logo_w,
    )

    # Bottom brand bar
    add_textbox(
        slide,
        T.MARGIN_X, T.SLIDE_H - Inches(0.4), Inches(8), Inches(0.3),
        f"{T.BRAND_NAME_ZH} · {T.BRAND_NAME_EN}",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, align=PP_ALIGN.LEFT, bold=False,
    )

    return slide
