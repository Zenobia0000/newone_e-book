"""Ch08 deck — 機率統計與大數據處理分析（完整版）
24 content slides + cover + copyright.

三大段落：
  起（S1-S4）  — 定錨 + 全覽 + 描述統計
  承（S5-S10） — Z-score / IQR / 機率分佈 / 假設檢定
  轉（S11-S18）— PCA / 聚類 / ROC / 串流 / 差分隱私
  合（S19-S24）— CIFAR-10 / 可視化 / Python 陷阱 / 練習 / 結語

評鑑範圍：中級 L221 + L222 + L223 + L224（科目2）
含 Python 程式題投影片
Aligned to chapters/Ch08_機率統計與大數據處理分析/02_slides_design.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source, add_rect,
    set_solid_fill, set_no_fill, set_line, set_no_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_vs_two_col, draw_pyramid_stack,
    draw_three_blocks_flow, draw_image_placeholder,
    draw_inverted_thesis_box, draw_grid, draw_code_panel,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch08"
MODULE_TITLE = "機率統計與大數據處理分析"
MODULE_SUBTITLE = "Z-score × IQR × PCA × ROC——用統計語言解讀大數據結果"
TIME_MIN = 150
N_CONTENT = 24


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_ch08(output_path, image_registry=None):
    """Build Ch08 deck; 24 content slides + cover + copyright."""
    prs = _new_prs()

    # ── Cover ──
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───── S1 · SILENT · 起 ─────
    s = _blank(prs)
    draw_silent_page(
        s,
        "統計和大數據不是兩件事——\n考試考的是你能否用統計語言\n解讀大數據結果。",
    )
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ───── S2 · ASK · 起 ─────
    s = _blank(prs)
    draw_ask_page(
        s,
        "看到 p<0.05，\n你知道該高興還是該懷疑？",
        data_card={
            "label": "科目2 第一題",
            "stat": "Q1",
            "caption": "就考檢定方法選擇",
        },
    )
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───── S3 · MATRIX · 起 ─────
    s = _blank(prs)
    draw_matrix(
        s, rows=2, cols=2,
        cells=[
            {"text": "L221 統計基礎",
             "sub": "描述統計 / 機率分佈\n假設檢定\nQ1 Q7 Q10", "highlight": True},
            {"text": "L222 數據處理",
             "sub": "ETL / 儲存 / 串流\nQ12 批次 vs 串流", "highlight": False},
            {"text": "L223 大數據分析",
             "sub": "PCA / 聚類 / ROC\n可視化\nQ2 Q9", "highlight": False},
            {"text": "L224 ML+AI+隱私",
             "sub": "ML 結合大數據\n差分隱私 / Python\nQ8 Q15", "highlight": True},
        ],
        title="科目2 四大評鑑範圍 · 八道樣題一覽",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "Python 程式題散佈在各區塊，約佔 25%。計算題集中在 L221 和 L223。",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───── S4 · TABLE · 起 ─────
    s = _blank(prs)
    draw_editorial_table(
        s,
        title="描述統計三維度",
        header=["維度", "指標", "定義", "考試重點"],
        rows=[
            ["集中趨勢", "Mean/Median/Mode", "中心在哪", "右偏：Mean>Median>Mode"],
            ["離散程度", "Range/IQR/Var/SD", "散多開", "SD 大 = 不穩定"],
            ["分佈型態", "Skewness/Kurtosis", "長什麼樣", "偏態方向 vs 尾巴厚度"],
        ],
        top=1.3,
        col_widths=[1.2, 1.8, 1.5, 2.5],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.0), T.SLIDE_W - 2 * T.MARGIN_X, Inches(1.5),
        "常考情境：「某班平均分 80、標準差 20」vs「平均分 80、標準差 5」\n"
        "→ 哪班成績更穩定？→ 標準差小的那班。\n\n"
        "右偏分佈排序：Mean > Median > Mode（尾巴在右邊拉高平均值）",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───── S5 · MECHANISM_FLOW · 承 ─────
    s = _blank(prs)
    draw_flow_chain(
        s,
        title="Z-score 計算流程（Q10）",
        nodes=[
            {"label": "原始值 X", "sub": "X = 88"},
            {"label": "減平均值", "sub": "X - mu\n88 - 72 = 16"},
            {"label": "除以 SD", "sub": "/ sigma\n16 / 8 = 2.0"},
            {"label": "Z-score", "sub": "Z = 2.0\n前 2.5%", "highlight": True},
        ],
        y=2.2,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.8), T.SLIDE_W - 2 * T.MARGIN_X, Inches(2.0),
        "Z = (X - mu) / sigma = 「距離平均值幾個標準差」\n\n"
        "68-95-99.7 法則：1 個 SD 內 68% / 2 個 SD 內 95% / 3 個 SD 內 99.7%\n"
        "|Z| > 3 通常視為異常值",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───── S6 · CODE · 承 ─────
    s = _blank(prs)
    add_title(s, "Python Z-score 計算（程式題考法）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.4),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.5),
        label="Z-score — numpy + scipy",
        code=(
            "import numpy as np\n"
            "from scipy import stats\n"
            "\n"
            "data = np.array([65, 70, 72, 75, 88])\n"
            "mean = np.mean(data)        # 74.0\n"
            "std  = np.std(data, ddof=1)  # 8.34\n"
            "z    = (data - mean) / std\n"
            "\n"
            "# 或直接用 scipy\n"
            "z2 = stats.zscore(data, ddof=1)\n"
            "outliers = data[np.abs(z2) > 3]"
        ),
        bullets=[
            "np.mean() 算均值",
            "np.std(ddof=1) = 樣本 SD",
            "np.std() 預設是母體 SD（除以 N）",
            "ddof=1 → 除以 N-1 → 樣本",
            "scipy.stats.zscore() 一步到位",
            "|Z| > 3 標記異常值",
        ],
    )
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───── S7 · TABLE · 承 ─────
    s = _blank(prs)
    draw_editorial_table(
        s,
        title="IQR 異常值判斷（Q7）",
        header=["步驟", "計算", "範例"],
        rows=[
            ["排序資料", "由小到大排列", "10, 15, 20, 25, 30, 35, 40, 45, 50"],
            ["找 Q1, Q3", "第 25/75 百分位", "Q1=20, Q3=40"],
            ["算 IQR", "Q3 - Q1", "IQR = 40 - 20 = 20"],
            ["下界", "Q1 - 1.5 × IQR", "20 - 30 = -10"],
            ["上界", "Q3 + 1.5 × IQR", "40 + 30 = 70"],
            ["判斷", "超出上下界 = 異常", "65 正常 / 75 異常"],
        ],
        top=1.3,
        col_widths=[1.0, 2.0, 3.0],
    )
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───── S8 · PITFALL · 承 ─────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        title="IQR 忘記乘 1.5——科目2 最常見錯誤",
        left_title="錯：直接用 IQR",
        right_title="對：乘以 1.5",
        left_items=[
            "IQR = 20",
            "下界 = Q1 - 20 = 0",
            "上界 = Q3 + 20 = 60",
            "範圍太窄，誤判太多",
        ],
        right_items=[
            "IQR = 20",
            "下界 = Q1 - 1.5×20 = -10",
            "上界 = Q3 + 1.5×20 = 70",
            "1.5 = 箱型圖鬍鬚長度倍率",
        ],
        summary="記憶法：鬍鬚長度 = 箱子的 1.5 倍。mild outlier 用 1.5，extreme 用 3。",
    )
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───── S9 · MATRIX · 承 ─────
    s = _blank(prs)
    draw_three_blocks_flow(
        s,
        title="三種機率分佈——場景決定選哪個",
        blocks=[
            {"heading": "常態分佈",
             "items": [
                 "連續 + 對稱鐘形",
                 "參數：mu, sigma",
                 "身高/成績/測量誤差",
             ]},
            {"heading": "二項分佈",
             "items": [
                 "固定 n 次試驗",
                 "參數：n, p",
                 "瑕疵品數/擲硬幣",
             ]},
            {"heading": "Poisson 分佈",
             "items": [
                 "單位時間事件計數",
                 "參數：lambda",
                 "客服來電/網頁點擊",
             ]},
        ],
        bottom_note="「次數 + 固定時間」→ Poisson。「n 次試驗 + 成功/失敗」→ 二項。",
    )
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───── S10 · VS · 承 ─────
    s = _blank(prs)
    add_title(s, "假設檢定決策樹（Q1）")
    # Use a flow chain to represent the decision tree
    draw_flow_chain(
        s,
        nodes=[
            {"label": "比較什麼？", "sub": "根節點", "highlight": True},
            {"label": "兩組均值", "sub": "→ t 檢定\n獨立 / 配對"},
            {"label": "三組+均值", "sub": "→ F 檢定\n(ANOVA)"},
            {"label": "類別頻次", "sub": "→ 卡方檢定\nchi-square"},
        ],
        y=2.2,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.8), T.SLIDE_W - 2 * T.MARGIN_X, Inches(2.0),
        "Q1：給場景，選檢定方法。規則：\n"
        "• 比兩組均值 → t 檢定\n"
        "• 比三組以上均值 → F 檢定（ANOVA）——不能做三次 t 檢定（多重比較膨脹 Type I error）\n"
        "• 比類別頻次 → 卡方檢定",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───── S11 · SILENT · 轉 ─────
    s = _blank(prs)
    draw_silent_page(
        s,
        "統計是語言，大數據分析方法是工具箱。\n接下來把工具拿出來用：\nPCA 降維、ROC 評估、聚類分析。",
    )
    add_footer(s, MODULE_CODE, 11, N_CONTENT, dark_bg=True)

    # ───── S12 · MECHANISM_FLOW · 轉 ─────
    s = _blank(prs)
    draw_flow_chain(
        s,
        title="PCA 降維四步驟（Q9）",
        nodes=[
            {"label": "標準化", "sub": "去量綱\nStandardScaler"},
            {"label": "協方差矩陣", "sub": "計算特徵間\n共變關係"},
            {"label": "特徵值分解", "sub": "找變異\n最大方向"},
            {"label": "選前 k 個", "sub": "累積解釋率\n≥ 85%", "highlight": True},
        ],
        y=2.2,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.8), T.SLIDE_W - 2 * T.MARGIN_X, Inches(2.0),
        "PCA = 找「資料變異最大的方向」——降維不是丟資料，是壓縮資訊。\n"
        "PCA 是無監督的——不看標籤。\n"
        "explained_variance_ratio_：第一個主成分解釋 40%、第二 25%、第三 15%，累積 80%。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───── S13 · CODE · 轉 ─────
    s = _blank(prs)
    add_title(s, "Python PCA 程式碼（程式題考法）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.4),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.5),
        label="PCA — sklearn",
        code=(
            "from sklearn.preprocessing import StandardScaler\n"
            "from sklearn.decomposition import PCA\n"
            "\n"
            "# 1. 標準化\n"
            "scaler = StandardScaler()\n"
            "X_scaled = scaler.fit_transform(X)\n"
            "\n"
            "# 2. PCA 降維\n"
            "pca = PCA(n_components=3)\n"
            "X_pca = pca.fit_transform(X_scaled)\n"
            "\n"
            "# 3. 解釋變異量\n"
            "print(pca.explained_variance_ratio_)\n"
            "print(X_pca.shape)  # (100, 3)"
        ),
        bullets=[
            "先 StandardScaler（PCA 對量綱敏感）",
            "PCA(n_components=3) → 保留 3 個主成分",
            "fit_transform = fit + transform",
            "shape: (100,5) → (100,3)",
            "explained_variance_ratio_ 加總看累積",
            "不標準化 → 被大數值特徵主導",
        ],
    )
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───── S14 · VS · 轉 ─────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        title="聚類分析：K-means vs DBSCAN",
        left_title="K-means",
        right_title="DBSCAN",
        left_items=[
            "預設 k 個群心",
            "迭代分配 + 更新群心",
            "假設群是球型（凸形）",
            "需要指定 k（群數）",
            "適合圓形群",
        ],
        right_items=[
            "基於密度可達性",
            "任意形狀群集",
            "自動偵測噪聲點",
            "需調 eps + min_samples",
            "適合月牙/環形群",
        ],
        summary="K-means 把月牙切兩半，DBSCAN 能正確識別。場景決定選哪個。",
    )
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───── S15 · MECHANISM_FLOW · 轉 ─────
    s = _blank(prs)
    add_title(s, "ROC 曲線與 AUC（Q2）")
    # Use editorial table for the confusion matrix + explanation
    draw_editorial_table(
        s,
        header=["", "預測正", "預測負"],
        rows=[
            ["實際正", "TP（真陽）", "FN（假陰）"],
            ["實際負", "FP（假陽）", "TN（真陰）"],
        ],
        top=1.3,
        col_widths=[1.5, 2.0, 2.0],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(3.5), T.SLIDE_W - 2 * T.MARGIN_X, Inches(3.0),
        "TPR (Recall) = TP / (TP + FN)    FPR = FP / (FP + TN)\n\n"
        "ROC 曲線 = 不同 threshold 下的 (FPR, TPR) 連線\n"
        "AUC = 曲線下面積\n\n"
        "AUC = 0.5 → 對角線 = 隨機猜 = 沒用\n"
        "AUC = 0.9+ → 很好\n"
        "AUC = 1.0 → 完美分類器\n\n"
        "ROC 曲線越靠左上角越好。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───── S16 · CODE · 轉 ─────
    s = _blank(prs)
    add_title(s, "Python ROC 曲線（程式題考法）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.4),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.5),
        label="ROC — sklearn",
        code=(
            "from sklearn.metrics import (\n"
            "    roc_curve, roc_auc_score\n"
            ")\n"
            "import matplotlib.pyplot as plt\n"
            "\n"
            "fpr, tpr, thresholds = roc_curve(\n"
            "    y_true, y_score\n"
            ")\n"
            "auc = roc_auc_score(y_true, y_score)\n"
            "\n"
            "plt.plot(fpr, tpr, label=f'AUC={auc:.2f}')\n"
            "plt.plot([0,1],[0,1],'--', color='gray')\n"
            "plt.xlabel('FPR')\n"
            "plt.ylabel('TPR')\n"
            "plt.legend()\n"
            "plt.show()"
        ),
        bullets=[
            "roc_curve() 回傳 fpr, tpr, thresholds",
            "roc_auc_score() 直接算 AUC",
            "虛線 = 隨機猜的基準線",
            "考題：哪條曲線的模型較好？",
            "→ 離左上角較近、AUC 較大",
        ],
    )
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───── S17 · TABLE · 轉 ─────
    s = _blank(prs)
    draw_editorial_table(
        s,
        title="批次處理 vs 串流處理（Q12）",
        header=["面向", "批次處理", "串流處理"],
        rows=[
            ["資料特性", "固定時間的完整資料", "持續流入的即時資料"],
            ["延遲", "高（分鐘~小時）", "低（毫秒~秒）"],
            ["工具", "Hadoop / Spark", "Kafka / Flink\nSpark Streaming"],
            ["適用場景", "日報表 / 歷史分析\n排程作業", "即時監控 / 推薦\n詐欺偵測"],
            ["關鍵字", "排程 / 完整資料集", "即時 / 低延遲 / 持續流入"],
        ],
        top=1.3,
        col_widths=[1.2, 2.5, 2.5],
    )
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───── S18 · VS · 轉 ─────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        title="差分隱私：epsilon 越小越安全（Q8）",
        left_title="epsilon 小 (e.g. 0.1)",
        right_title="epsilon 大 (e.g. 10)",
        left_items=[
            "噪音大",
            "隱私保護強",
            "資料可用性低",
            "查詢結果不太準確",
        ],
        right_items=[
            "噪音小",
            "隱私保護弱",
            "資料可用性高",
            "查詢結果較準確",
        ],
        summary="epsilon 是隱私和可用性的調節旋鈕。Q8：epsilon 減小 → 隱私增強 + 可用性降低。",
    )
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ───── S19 · CODE · 合 ─────
    s = _blank(prs)
    add_title(s, "Python CIFAR-10（Q15）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.4),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.5),
        label="CIFAR-10 — keras",
        code=(
            "from keras.datasets import cifar10\n"
            "\n"
            "(x_train, y_train), (x_test, y_test) = \\\n"
            "    cifar10.load_data()\n"
            "\n"
            "print(x_train.shape)\n"
            "# (50000, 32, 32, 3)\n"
            "# 50000 張 32x32 RGB 彩色影像\n"
            "\n"
            "# 正規化：[0,255] → [0,1]\n"
            "x_train = x_train / 255.0\n"
            "x_test  = x_test  / 255.0\n"
            "\n"
            "# 為什麼除 255？\n"
            "# 神經網路在小數值範圍收斂更快"
        ),
        bullets=[
            "shape (50000, 32, 32, 3)",
            "50000 張 / 32x32 px / 3 通道 RGB",
            "10 類：飛機、汽車、鳥...",
            "除 255.0 → 正規化到 [0,1]",
            "用 255.0 不是 255（避免整數除法）",
            "考題問 shape 或「某行 code 的作用」",
        ],
    )
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───── S20 · TABLE · 合 ─────
    s = _blank(prs)
    draw_editorial_table(
        s,
        title="可視化圖表選擇指南",
        header=["資料類型", "推薦圖表", "次選", "避免"],
        rows=[
            ["連續 × 連續", "散佈圖 scatter", "氣泡圖", "圓餅圖"],
            ["類別 × 連續", "箱型圖 box", "長條圖 bar", "折線圖"],
            ["分佈", "直方圖 histogram", "密度圖 KDE", "散佈圖"],
            ["時間序列", "折線圖 line", "面積圖", "圓餅圖"],
            ["類別占比", "圓餅圖 / 長條圖", "樹狀圖", "散佈圖"],
        ],
        top=1.3,
        col_widths=[1.5, 2.0, 1.5, 1.5],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.5), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "口訣：兩個連續 → scatter / 類別比較 → bar/box / 分佈 → histogram / 時間 → line",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ───── S21 · PITFALL · 合 ─────
    s = _blank(prs)
    add_title(s, "Python 程式題四大陷阱")
    draw_grid(
        s, rows=2, cols=2,
        cells=[
            {"label": "shape 搞混",
             "sub": "(50000,32,32,3)\n50000 是樣本數不是維度"},
            {"label": "axis 搞混",
             "sub": "axis=0 沿列(向下)\nnp.mean(X, axis=0) = 每個特徵均值"},
            {"label": "dtype 沒轉",
             "sub": "整數 / 255 → 0\n要用 / 255.0 → 浮點數"},
            {"label": "random_state",
             "sub": "train_test_split 不設\n→ 每次結果不同\n→ 不可重現"},
        ],
        title="",
        top=1.4,
        bottom=1.2,
        caption="這四個是科目2 Python 題最常見的出題點。",
    )
    add_footer(s, MODULE_CODE, 21, N_CONTENT)

    # ───── S22 · PRACTICE · 合 ─────
    s = _blank(prs)
    add_title(s, "Practice · 科目2 計算題三連發")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 22 / 24 · 五分鐘計時",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(10.8), Inches(3.5),
        "Q1  Z-score：mu=100, sigma=15, X=130\n"
        "    → Z = ?\n\n"
        "Q2  IQR：Q1=35, Q3=65，判斷 X=110 是否異常\n"
        "    → IQR = ?  上界 = ?  結論 = ?\n\n"
        "Q3  檢定選擇：\n"
        "    (a) 比較兩組薪資差異 → ?\n"
        "    (b) 比較三家門市銷售差異 → ?",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.8), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.8),
        "答案：Q1 Z=2.0（高於平均 2 個 SD）· Q2 IQR=30, 上界=65+45=110, "
        "110 不超過上界 → 不算異常 · Q3 (a) t 檢定 (b) F 檢定(ANOVA)",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
        line_spacing=1.4,
    )
    add_footer(s, MODULE_CODE, 22, N_CONTENT)

    # ───── S23 · PYRAMID · 合 ─────
    s = _blank(prs)
    draw_pyramid_stack(
        s,
        title="Ch08 收束：三層能力堆疊",
        layers=[
            {"name": "進階應用",
             "caption": "ML + 差分隱私 + Python 程式（L224）"},
            {"name": "大數據分析",
             "caption": "PCA / ROC / 聚類 / 串流 / 可視化（L222+L223）"},
            {"name": "統計基礎",
             "caption": "描述統計 / 機率分佈 / 假設檢定（L221）"},
        ],
        thesis="統計是基礎，分析是工具，Python 是驗證。",
    )
    add_footer(s, MODULE_CODE, 23, N_CONTENT)

    # ───── S24 · SILENT · 合 ─────
    s = _blank(prs)
    draw_silent_page(
        s,
        "用統計語言解讀大數據結果——\n這是科目2 的核心能力。\n下一章 Ch09：機器學習與深度學習。",
    )
    add_footer(s, MODULE_CODE, 24, N_CONTENT, dark_bg=True)

    # ── Copyright ──
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
