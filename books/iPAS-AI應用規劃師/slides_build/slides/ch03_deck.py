"""Ch03 deck — 機器學習概念與鑑別式/生成式 AI（完整版）
22 content slides + cover + copyright.

三大 Part：
  Part A — ML 基礎與學習類型（S1-S12）
  Part B — 深度學習與生成式家族（S13-S15）
  Part C — 鑑別式 vs 生成式（S16-S22）

受眾：iPAS 初級考生，高中程度即可。
Aligned to chapters/Ch03_機器學習概念與鑑別式生成式AI/02_slides_design.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source, add_rect,
    set_solid_fill, set_no_fill, set_line, set_no_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_vs_two_col, draw_pyramid_stack,
    draw_three_blocks_flow, draw_image_placeholder,
    draw_inverted_thesis_box,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch03"
MODULE_TITLE = "機器學習概念與鑑別式/生成式 AI"
MODULE_SUBTITLE = "ML 家族圖 \u00d7 四種學習類型 \u00d7 鑑別 vs 生成\u2014\u2014三層過濾答案自動浮現"
TIME_MIN = 120
N_CONTENT = 22


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _draw_bridge_note(slide, text, y_inches=6.6):
    add_textbox(
        slide, T.MARGIN_X, Inches(y_inches),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        text,
        font_size=T.FONT_SMALL, color=T.PRIMARY, bold=True, italic=True,
        align=PP_ALIGN.CENTER,
    )


def _draw_pitfall_card(slide, title_text, wrong, right, why=""):
    """PITFALL layout: title + wrong (red-bordered) + right (green-bordered) + optional why."""
    add_title(slide, title_text)
    col_w = Inches(5.5)
    gap = Inches(0.5)
    total = col_w * 2 + gap
    left_x = (T.SLIDE_W - total) / 2
    right_x = left_x + col_w + gap
    top = Inches(1.5)
    col_h = Inches(3.5)

    # Wrong box
    w_box = add_rect(slide, left_x, top, col_w, col_h)
    set_no_fill(w_box)
    set_line(w_box, T.ACCENT_WARM, 2.0)
    add_textbox(
        slide, left_x, top, col_w, Inches(0.5),
        "WRONG",
        font_size=T.FONT_BODY, color=T.ACCENT_WARM, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide, left_x + Inches(0.25), top + Inches(0.6),
        col_w - Inches(0.5), col_h - Inches(0.7),
        wrong,
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )

    # Right box
    r_box = add_rect(slide, right_x, top, col_w, col_h)
    set_no_fill(r_box)
    set_line(r_box, T.PRIMARY, 2.0)
    add_textbox(
        slide, right_x, top, col_w, Inches(0.5),
        "RIGHT",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide, right_x + Inches(0.25), top + Inches(0.6),
        col_w - Inches(0.5), col_h - Inches(0.7),
        right,
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )

    if why:
        why_y = top + col_h + Inches(0.2)
        why_box = add_rect(slide, left_x, why_y, total, Inches(0.6))
        set_solid_fill(why_box, T.PRIMARY)
        set_no_line(why_box)
        add_textbox(
            slide, left_x + Inches(0.3), why_y,
            total - Inches(0.6), Inches(0.6),
            why,
            font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )


def build_ch03(output_path, image_registry=None):
    """Build Ch03 deck; 22 content slides + cover + copyright."""
    prs = _new_prs()

    # -- Cover --
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ===== S1 . MOTIVATION -- ML 家族佔考試四分之一 =====
    s = _blank(prs)
    draw_ask_page(
        s,
        "ML 家族有七八個成員\u2014\u2014\n考試考的不是你會不會寫，\n是你分不分得清誰是誰。",
        data_card={
            "label": "L113 + L114 合計",
            "stat": "~25-30%",
            "caption": "初級科目一佔比\n機器學習基本原理 + 鑑別式 vs 生成式",
        },
    )
    add_source(s, "iPAS 歷屆考古題分析 \u00b7 初級科目一")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # ===== S2 . ASK -- 選錯學習類型就全盤皆錯 =====
    s = _blank(prs)
    draw_ask_page(
        s,
        "客戶分群用 supervised？\n推薦系統用 unsupervised？\n場景描述換一個詞，答案就不同。",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.2), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "「分群」= unsupervised（沒有預設答案）·「流失預測」= supervised（有 label）",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ===== S3 . SILENT -- 立論 =====
    s = _blank(prs)
    draw_silent_page(
        s,
        "ML 不是一種技術，是一個家族。\n考試不考你寫程式，\n考你能不能分辨誰適合做什麼。",
    )
    add_footer(s, MODULE_CODE, 3, N_CONTENT, dark_bg=True)

    # ===== S4 . CONCEPT-CARD -- ML 的定義與目的 =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="傳統程式",
        right_title="機器學習",
        left_items=[
            "人寫規則 \u2192 電腦照做",
            "輸入：規則 + 資料",
            "輸出：答案",
            "維護成本隨規則數量上升",
        ],
        right_items=[
            "人給資料 \u2192 電腦自己學",
            "輸入：資料 + 答案",
            "輸出：規則（模型）",
            "規律自動從資料中浮現",
        ],
        title="ML 的定義：傳統程式 vs 機器學習",
        summary="ML 的核心目的：泛化\u2014\u2014用舊資料學到的規律，去預測新資料",
    )
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ===== S5 . CONCEPT-CARD -- 模型訓練與泛化 =====
    s = _blank(prs)
    add_title(s, "模型訓練與泛化：五步驟流程")
    draw_flow_chain(
        s,
        nodes=[
            {"label": "資料收集", "sub": "蒐集歷史資料"},
            {"label": "特徵工程", "sub": "挑選有用特徵"},
            {"label": "模型訓練", "sub": "從資料學規律", "highlight": True},
            {"label": "驗證評估", "sub": "測試泛化能力"},
            {"label": "部署上線", "sub": "服務新資料"},
        ],
        y=2.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.3), T.SLIDE_W - 2 * T.MARGIN_X, Inches(1.5),
        "\u2022 訓練 = 從歷史資料學規律\n"
        "\u2022 泛化 = 規律能用在新資料\n"
        "\u2022 Overfitting = 背答案\u2014\u2014訓練很好但遇到新題就崩",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    warn_y = Inches(6.0)
    warn = add_rect(s, T.MARGIN_X, warn_y,
                    T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5))
    set_solid_fill(warn, T.PRIMARY)
    set_no_line(warn)
    add_textbox(
        s, T.MARGIN_X + Inches(0.3), warn_y,
        T.SLIDE_W - 2 * T.MARGIN_X - Inches(0.6), Inches(0.5),
        "考題：泛化能力不足的模型具有什麼特性？\u2192 Overfitting",
        font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ===== S6 . PITFALL (P1) -- 深度學習不等於機器學習 =====
    s = _blank(prs)
    _draw_pitfall_card(
        s,
        "P1 深度學習不等於機器學習",
        wrong="「深度學習就是機器學習的另一個說法」",
        right="「深度學習是機器學習的子集（一種特殊方法）」\n\nAI \u2283 ML \u2283 DL 三層同心圓",
        why="就像柯基是狗的一種\u2014\u2014你不會說所有狗都是柯基",
    )
    draw_image_placeholder(
        s, Inches(4.5), Inches(1.5), Inches(4.5), Inches(3.0),
        slot_name="AI \u2283 ML \u2283 DL 同心圓",
        description="三層同心圓：最外圈 AI，中圈 ML，內圈 DL。",
        size_hint="1000\u00d7800 px",
        placeholder_id="Ch03_S6_ai_ml_dl_circles",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ===== S7 . CONCEPT-CARD -- 監督式學習 =====
    s = _blank(prs)
    add_title(s, "監督式學習 Supervised Learning")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
        "定義：資料有 label（標準答案）\u00b7 模型學 feature \u2192 label 的對應關係",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    draw_vs_two_col(
        s,
        left_title="Regression \u8ff4\u6b78",
        right_title="Classification \u5206\u985e",
        left_items=[
            "Label 是連續數字",
            "房價、銷售額、溫度",
            "輸出：一個數值",
        ],
        right_items=[
            "Label 是離散類別",
            "流失/不流失、垃圾/正常",
            "輸出：一個類別",
        ],
        summary="口訣：有答案 = supervised \u00b7 答案是數字 = regression \u00b7 答案是類別 = classification",
    )
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ===== S8 . CHECKPOINT -- Regression vs Classification 場景判斷 =====
    s = _blank(prs)
    add_title(s, "Check Point \u00b7 Regression vs Classification 場景判斷")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 8 / 22 \u00b7 不看投影片作答",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(10.8), Inches(3.5),
        "Q1  預測明天最高溫度 \u2192 Regression（連續數字）\n\n"
        "Q2  判斷貸款是否違約 \u2192 Classification（是/否）\n\n"
        "Q3  預測下個月銷售額 \u2192 Regression（連續數字）\n\n"
        "Q4  偵測信用卡詐騙交易 \u2192 Classification（詐騙/正常）\n\n"
        "Q5  預測客戶滿意度 1-5 分 \u2192 看定義：連續=Reg\uff1b離散等級=Cls",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.0), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "陷阱：Q5 兩種都可以\u2014\u2014看題目怎麼定義。「預測分數」\u2192 Reg\uff1b「預測等級」\u2192 Cls",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ===== S9 . CONCEPT-CARD -- 非監督式學習 =====
    s = _blank(prs)
    add_title(s, "非監督式學習 Unsupervised Learning")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
        "定義：資料沒有 label \u00b7 模型自己找資料中的結構與模式",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "分群 Clustering",
             "items": [
                 "K-Means / DBSCAN",
                 "客戶分群",
                 "市場區隔",
             ]},
            {"heading": "降維 Dim. Reduction",
             "items": [
                 "PCA / t-SNE",
                 "資料視覺化",
                 "特徵壓縮",
             ]},
            {"heading": "異常偵測",
             "items": [
                 "Isolation Forest",
                 "設備故障偵測",
                 "網路入侵偵測",
             ]},
        ],
        bottom_note="口訣：沒答案 = unsupervised \u00b7 找結構 = 分群 \u00b7 找異類 = 異常偵測",
    )
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ===== S10 . PITFALL (P2) -- 非監督式不代表不需要資料 =====
    s = _blank(prs)
    _draw_pitfall_card(
        s,
        "P2 非監督式不代表不需要資料",
        wrong="「非監督式學習不需要訓練資料」",
        right="「非監督式不需要 label，\n但仍然需要大量高品質的資料」",
        why="分群要有足夠資料才能找到有意義的群組結構\uff1b降維要有足夠維度才值得降",
    )
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ===== S11 . CONCEPT-CARD -- 半監督式與強化式學習 =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="半監督式 Semi-supervised",
        right_title="強化式 Reinforcement",
        left_items=[
            "少量有 label + 大量無 label",
            "標注成本高（醫療影像、法律文件）",
            "用少量標注訓練初始模型",
            "再用未標注資料擴展",
            "關鍵字：「借力」",
        ],
        right_items=[
            "Agent 在環境中試錯",
            "透過獎勵/懲罰學最佳策略",
            "AlphaGo / 自動駕駛",
            "推薦系統 / 機器人控制",
            "關鍵字：「試錯」",
        ],
        title="半監督式 vs 強化式學習",
        summary="半監督「借力」\u2014\u2014少量 label + 大量未標注。強化式「試錯」\u2014\u2014Agent + 環境 + 獎勵。",
    )
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ===== S12 . CHECKPOINT -- 四種學習類型配對場景 =====
    s = _blank(prs)
    add_title(s, "Check Point \u00b7 四種學習類型配對場景")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 12 / 22 \u00b7 不看投影片作答",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    draw_editorial_table(
        s,
        header=["場景", "學習類型", "判斷依據"],
        rows=[
            ["客戶分群做精準行銷", "非監督式", "沒有預設答案"],
            ["垃圾郵件偵測（大量標注）", "監督式", "有 label"],
            ["醫療影像（200 張標注 + 10000 未標注）", "半監督式", "label 少 + 未標注多"],
            ["自動駕駛車輛路徑規劃", "強化式", "Agent + 環境互動"],
        ],
        top=1.8,
        col_widths=[3, 1.5, 2],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.5), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "場景 C 最易錯\u2014\u2014有一些 label 但不多 + 大量未標注 = 半監督式",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ===== S13 . CONCEPT-CARD -- 深度學習 =====
    s = _blank(prs)
    add_title(s, "深度學習：自動學特徵的多層神經網路")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(3.8),
        "\u2022 多層神經網路（輸入層 \u2192 隱藏層 x N \u2192 輸出層）\n\n"
        "\u2022 優勢：自動從原始資料提取特徵\n"
        "  不需手動特徵工程\n\n"
        "\u2022 代價：需要大量資料 + 大量運算資源（GPU）\n\n"
        "\u2022 擅長：影像（CNN）、序列（RNN/Transformer）、\n"
        "  生成（GAN/GPT）\n\n"
        "\u2022 DL 是 ML 的子集，專門處理非結構化資料",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(4.0),
        slot_name="簡化神經網路架構圖",
        description="輸入層 \u2192 隱藏層 x N \u2192 輸出層，箭頭標示自動特徵提取。",
        size_hint="1200\u00d7900 px",
        placeholder_id="Ch03_S13_nn_architecture",
        registry=image_registry,
    )
    _draw_bridge_note(s, "考題：DL 跟傳統 ML 最大差異？\u2192 自動特徵提取")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ===== S14 . PITFALL (P3) -- 生成式 AI 不只有 ChatGPT =====
    s = _blank(prs)
    add_title(s, "P3 生成式 AI 不只有 ChatGPT")
    _draw_pitfall_card(
        s,
        "P3 生成式 AI 不只有 ChatGPT",
        wrong="「生成式 AI 就是 ChatGPT」",
        right="「ChatGPT 是生成式 AI 家族的一員（LLM），\n但家族還有很多成員」",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.6), T.SLIDE_W - 2 * T.MARGIN_X, Inches(1.0),
        "GAN 生成對抗網路 \u00b7 VAE 變分自編碼器 \u00b7 Diffusion 擴散模型 \u00b7 LLM 大語言模型\n"
        "四個家族成員都是生成式\u2014\u2014考題問「何者為生成式 AI？」時都要認得",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        align=PP_ALIGN.CENTER,
        line_spacing=1.45,
    )
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ===== S15 . CONCEPT-CARD -- 多模態學習 =====
    s = _blank(prs)
    add_title(s, "多模態學習：同時處理多種資料類型")
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "文字 Text",
             "items": [
                 "NLP / LLM",
                 "GPT-4 / Gemini",
             ]},
            {"heading": "圖片 Image",
             "items": [
                 "CV / CNN",
                 "CLIP / DALL-E",
             ]},
            {"heading": "語音 Audio",
             "items": [
                 "ASR / TTS",
                 "Whisper",
             ]},
        ],
        bottom_note="",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.5), T.SLIDE_W - 2 * T.MARGIN_X, Inches(1.0),
        "多模態 = 同時吃多種資料。GPT-4V 傳照片問問題 = 文字 + 圖片融合。\n"
        "考題：「何者為多模態 AI 應用？」\u2192 找同時涉及兩種以上資料類型的選項",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        align=PP_ALIGN.CENTER,
        line_spacing=1.45,
    )
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ===== S16 . CONCEPT-CARD -- 鑑別式 AI =====
    s = _blank(prs)
    add_title(s, "鑑別式 AI：學會分辨的專家")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
        "鑑別式 AI 學的是「邊界」\u2014\u2014給一個輸入，判斷它屬於哪一類 \u00b7 學 P(Y|X)",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(2.2), Inches(6.5), Inches(3.5),
        "代表模型：\n"
        "\u2022 Logistic Regression\n"
        "\u2022 SVM\n"
        "\u2022 Random Forest\n"
        "\u2022 CNN（分類用途）\n\n"
        "應用：\n"
        "\u2022 垃圾郵件分類\n"
        "\u2022 人臉辨識\n"
        "\u2022 醫療影像判讀\n"
        "\u2022 信用評分",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    draw_inverted_thesis_box(s, "一句話：看了東西做判斷", y=6.0)
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ===== S17 . CONCEPT-CARD -- 生成式 AI =====
    s = _blank(prs)
    add_title(s, "生成式 AI：學會創造的藝術家")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
        "生成式 AI 學的是「分佈」\u2014\u2014學會資料整體長什麼樣，就能生成新的 \u00b7 學 P(X) 或 P(X|Y)",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(2.2), Inches(6.5), Inches(3.5),
        "代表模型：\n"
        "\u2022 GAN \u00b7 VAE \u00b7 Diffusion\n"
        "\u2022 Transformer (GPT) \u00b7 LLM\n\n"
        "應用：\n"
        "\u2022 文字生成 / 圖片生成\n"
        "\u2022 程式碼生成 / 語音合成\n"
        "\u2022 影片生成 / 藥物分子設計",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    draw_inverted_thesis_box(s, "一句話：憑條件創造新東西", y=6.0)
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ===== S18 . TABLE -- 鑑別式 vs 生成式對比表 =====
    s = _blank(prs)
    add_title(s, "鑑別式 vs 生成式：一張表搞定考題")
    draw_editorial_table(
        s,
        header=["比較維度", "鑑別式 AI", "生成式 AI"],
        rows=[
            ["核心目標", "畫邊界，做判斷", "學分佈，做創造"],
            ["數學直覺", "P(Y|X)", "P(X) 或 P(X|Y)"],
            ["輸入\u2192輸出", "資料\u2192類別/分數", "條件\u2192新內容"],
            ["代表模型", "LogReg/SVM/RF/CNN", "GAN/VAE/Diffusion/GPT"],
            ["典型應用", "分類/辨識/偵測", "生成/創作/合成"],
            ["一句話", "看了東西做判斷", "憑條件創造新東西"],
        ],
        top=1.3,
        col_widths=[1.5, 2.5, 2.5],
    )
    _draw_bridge_note(s, "鑑別式 = 法官（做判斷）\u00b7 生成式 = 藝術家（做創造）")
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ===== S19 . CHECKPOINT -- 鑑別式 vs 生成式場景判斷 =====
    s = _blank(prs)
    add_title(s, "Check Point \u00b7 鑑別式 vs 生成式場景判斷")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 19 / 22 \u00b7 不看投影片作答",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(10.8), Inches(3.5),
        "Q1  信用卡詐騙偵測 \u2192 鑑別式（判斷是否詐騙）\n\n"
        "Q2  自動生成行銷文案 \u2192 生成式（創造新文字）\n\n"
        "Q3  手機人臉解鎖 \u2192 鑑別式（判斷是否為機主）\n\n"
        "Q4  AI 繪圖（文字\u2192圖片）\u2192 生成式（創造新圖片）\n\n"
        "Q5  語音轉文字 \u2192 鑑別式（辨識語音內容）",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.0), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "陷阱：Q5 語音轉文字是辨識任務 \u2260 生成。真正的生成式語音任務是 TTS（語音合成）。",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ===== S20 . TABLE -- CV / 語音 / 生成技術應用場域 =====
    s = _blank(prs)
    add_title(s, "CV\u3001語音\u3001生成技術的應用場域總覽")
    draw_editorial_table(
        s,
        header=["電腦視覺 CV", "語音技術", "生成技術"],
        rows=[
            ["人臉辨識", "語音助理（Siri/Alexa）", "文字生成（ChatGPT）"],
            ["自動駕駛（物件偵測）", "語音轉文字 ASR", "圖片生成（Midjourney）"],
            ["醫療影像判讀", "聲紋辨識", "程式碼生成（Copilot）"],
            ["工廠品質檢測", "語音合成 TTS", "藥物分子設計"],
            ["OCR 文字辨識", "情緒語音分析", "數據增強/合成資料"],
        ],
        top=1.3,
        col_widths=[2, 2, 2],
    )
    _draw_bridge_note(s, "考題問場景配技術\u2014\u2014這張表就是你的配對卡")
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ===== S21 . PYRAMID -- ML 家族分類三層記憶法 =====
    s = _blank(prs)
    draw_pyramid_stack(
        s,
        title="Ch03 收束：ML 家族分類三層記憶法",
        layers=[
            {"name": "第三層：判斷 vs 創造",
             "caption": "鑑別式 / 生成式"},
            {"name": "第二層：任務類型",
             "caption": "Regression / Classification / Clustering / 降維 / 試錯"},
            {"name": "第一層：有沒有 label",
             "caption": "Supervised / Unsupervised / Semi-supervised / RL"},
        ],
        thesis="三層過濾，答案自動浮現。",
    )
    add_footer(s, MODULE_CODE, 21, N_CONTENT)

    # ===== S22 . MOTIVATION -- 銜接 Ch04 =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="本章完成",
        right_title="Ch04 預告",
        left_items=[
            "ML 家族全貌",
            "四種學習類型",
            "鑑別式 vs 生成式",
            "CV / 語音 / 生成應用場域",
        ],
        right_items=[
            "生成式 AI 工具實戰",
            "提示工程 RTFC",
            "RAG vs 微調",
            "No Code / Low Code",
        ],
        title="認識了 ML 家族，接下來學最熱門的生成式工具",
        summary="知道原理之後，來看怎麼用\u2014\u2014這正是 L121+L122 在考的。",
    )
    add_footer(s, MODULE_CODE, 22, N_CONTENT)

    # -- Copyright --
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
