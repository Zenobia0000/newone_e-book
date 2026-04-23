"""Ch05 deck — 生成式 AI 導入評估與風險管理（完整版）
20 content slides + cover + copyright.

三大 Part：
  Part A — 導入評估（S1-S8）
  Part B — 導入規劃與 AI Agent（S9-S12）
  Part C — 風險管理（S13-S20）

受眾：iPAS AI 應用規劃師初級/中級備考。
Aligned to chapters/Ch05_生成式AI導入評估與風險管理/02_slides_design.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source, add_rect,
    set_solid_fill, set_no_fill, set_line, set_no_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_vs_two_col, draw_pyramid_stack,
    draw_three_blocks_flow, draw_image_placeholder,
    draw_inverted_thesis_box, draw_grid,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch05"
MODULE_TITLE = "生成式 AI 導入評估與風險管理"
MODULE_SUBTITLE = "三維評估 \u00d7 五階段導入 \u00d7 風險四層策略\u2014\u2014規劃師的核心判斷力"
TIME_MIN = 120
N_CONTENT = 20


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_ch05(output_path, image_registry=None):
    """Build Ch05 deck; 20 content slides + cover + copyright."""
    prs = _new_prs()

    # == Cover ==
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ===== S1 . SILENT . 起 =====
    s = _blank(prs)
    draw_silent_page(
        s,
        "導入 AI 不是裝一個工具，\n是做一個決策。\n考試考的是你的規劃判斷力。",
    )
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ===== S2 . ASK . 起 =====
    s = _blank(prs)
    draw_ask_page(
        s,
        "你的公司說要導入 AI，\n第一步該做什麼？",
        data_card={
            "label": "McKinsey 2024",
            "stat": "87%",
            "caption": "企業 AI 專案\n未達預期效益",
        },
    )
    add_source(s, "McKinsey State of AI 2024")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ===== S3 . MATRIX 2x2 . 起 — 三維評估 =====
    s = _blank(prs)
    draw_matrix(
        s, rows=1, cols=3,
        title="導入評估三維度——缺一不可",
        cells=[
            {"text": "技術可行性",
             "sub": "資料品質\n模型成熟度\n基礎設施",
             "highlight": True},
            {"text": "商業價值",
             "sub": "ROI 預估\n市場競爭力\n客戶需求",
             "highlight": False},
            {"text": "組織準備度",
             "sub": "人才技能\n文化接受度\n管理支持",
             "highlight": False},
        ],
        top=1.4, bottom=2.5,
    )
    draw_inverted_thesis_box(
        s,
        "三維交集 = Go / No-Go 決策點。考題最愛出「以下何者不屬於導入評估範疇」。",
        y=5.8,
    )
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ===== S4 . TABLE . 起 — 效能指標 =====
    s = _blank(prs)
    add_title(s, "效能指標：不是越高越好，要看場景選對的那一個")
    draw_editorial_table(
        s,
        header=["指標", "適用任務", "公式 / 定義", "使用場景", "考試陷阱"],
        rows=[
            ["Accuracy", "均衡分類", "正確數/總數", "類別均衡", "不均衡時失效"],
            ["Precision", "分類", "TP/(TP+FP)", "垃圾郵件(不誤殺)", "高≠好(看場景)"],
            ["Recall", "分類", "TP/(TP+FN)", "醫療診斷(不漏診)", "與Precision互斥"],
            ["F1", "分類", "2PR/(P+R)", "不均衡分類", "是調和平均"],
            ["Latency", "推論", "回應時間", "即時系統", "與精度 trade-off"],
            ["Cost/query", "API 調用", "每次呼叫成本", "大規模部署", "含隱性成本"],
        ],
        top=1.3,
        col_widths=[1.2, 1.0, 1.5, 1.5, 1.3],
    )
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ===== S5 . VS . 承 — 自建 vs API vs 開源微調 =====
    s = _blank(prs)
    draw_three_blocks_flow(
        s,
        title="三條路：自建模型 / API 調用 / 開源微調",
        blocks=[
            {"heading": "自建模型",
             "items": [
                 "成本：最高",
                 "彈性：最高",
                 "時間：最長",
                 "門檻：需 ML 團隊",
                 "隱私：資料不出門",
             ]},
            {"heading": "API 調用",
             "items": [
                 "成本：最低(初期)",
                 "彈性：受限供應商",
                 "時間：最快",
                 "門檻：低",
                 "隱私：資料送出去",
             ]},
            {"heading": "開源微調",
             "items": [
                 "成本：中(需GPU)",
                 "彈性：中高",
                 "時間：中",
                 "門檻：需ML基礎",
                 "隱私：可控",
             ]},
        ],
        bottom_note="考題給公司條件，要你選最適合的方案。",
    )
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ===== S6 . TABLE . 承 — ROI / TCO =====
    s = _blank(prs)
    add_title(s, "ROI 看報酬率，TCO 看總花費——方向不同，考題不同")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(0.6),
        "ROI = (收益 - 成本) / 成本",
        font_size=Pt(18), color=T.PRIMARY, bold=True,
        family=T.FONT_MONO,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(2.2), Inches(6.5), Inches(1.5),
        "範例：AI 客服省下 500 萬人力，建置花 200 萬\n"
        "ROI = (500 - 200) / 200 = 150%",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(3.8), Inches(12), Inches(0.5),
        "TCO（總擁有成本）= 冰山圖",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    draw_editorial_table(
        s,
        header=["水面上（可見）", "水面下（隱性）"],
        rows=[
            ["硬體 / 軟體 / 雲端", "資料標註 / 清洗"],
            ["初期開發人力", "員工訓練 / 變更管理"],
            ["", "法遵 / 稽核 / 退場成本"],
        ],
        top=4.4,
        col_widths=[1, 1],
    )
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ===== S7 . MECHANISM_FLOW . 承 — 五階段 =====
    s = _blank(prs)
    add_title(s, "《AI 導入指引》五階段——背起來就是送分題")
    draw_flow_chain(
        s,
        nodes=[
            {"label": "需求確認", "sub": "需求規格書",
             "caption": "PM / 使用者"},
            {"label": "可行性評估", "sub": "評估報告",
             "caption": "規劃師"},
            {"label": "PoC", "sub": "PoC 結果",
             "caption": "ML 工程師"},
            {"label": "系統開發", "sub": "系統文件",
             "caption": "開發團隊"},
            {"label": "部署維運", "sub": "SLA",
             "caption": "維運團隊", "highlight": True},
        ],
        y=2.2,
    )
    draw_inverted_thesis_box(
        s,
        "需求確認一定在最前面——不是可行性評估。每階段的順序和產出是高頻考題。",
        y=5.8,
    )
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ===== S8 . PITFALL . 承 — 三維 vs 只做技術 =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        title="PITFALL：只做技術評估就決定導入？",
        left_title="錯誤做法",
        right_title="正確做法",
        left_items=[
            "只做技術評估就決定導入",
            "上線後員工不會用",
            "老闆不買單",
            "專案失敗",
        ],
        right_items=[
            "技術 + 商業 + 組織 三維評估",
            "確認 Go / No-Go",
            "漸進導入",
            "成功擴展",
        ],
        summary="選項裡只列技術面看起來很專業，但缺了商業/組織就不完整。三維缺一不可。",
    )
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ===== S9 . MATRIX . 轉 — 導入規劃五面向 =====
    s = _blank(prs)
    draw_matrix(
        s, rows=2, cols=3,
        title="導入規劃五面向——系統思考",
        cells=[
            {"text": "目標設置",
             "sub": "SMART 原則\n量化 KPI",
             "highlight": True},
            {"text": "資源分配",
             "sub": "人力/預算/時間\n資料/基礎設施",
             "highlight": False},
            {"text": "導入策略",
             "sub": "漸進式 vs 全面式\nPilot → Scale",
             "highlight": False},
            {"text": "因應措施",
             "sub": "風險預案\n變更管理",
             "highlight": False},
            {"text": "測試",
             "sub": "功能/壓力/UAT\n驗收標準",
             "highlight": False},
            {"text": "",
             "sub": "",
             "highlight": False},
        ],
        top=1.4, bottom=1.8,
    )
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ===== S10 . MECHANISM_FLOW . 轉 — AI Agent 決策樹 =====
    s = _blank(prs)
    add_title(s, "AI Agent：不只是更強的 ChatGPT")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(3.5),
        "決策流程：\n"
        "需要自主多步驟決策？\n"
        "  → Yes → AI Agent\n"
        "    → 評估自主權邊界\n"
        "    → 設計監控機制\n"
        "  → No → 傳統 AI 工具\n"
        "    → 評估 API / 自建 / 微調\n\n"
        "AI Agent 核心循環：\n"
        "  感知 → 推理 → 行動 → 環境回饋",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(4.0),
        slot_name="AI Agent 架構圖",
        description="感知 → 推理 → 行動 → 環境回饋 循環圖。標示自主權邊界和人類介入點。",
        size_hint="1200\u00d7900 px",
        placeholder_id="Ch05_S10_agent_loop",
        registry=image_registry,
    )
    draw_inverted_thesis_box(
        s,
        "導入 Agent 關鍵：自主權邊界、人類何時介入、出錯怎麼回滾。",
        y=5.8,
    )
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ===== S11 . VS . 轉 — 傳統 AI 工具 vs AI Agent =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        title="傳統 AI 工具 vs AI Agent",
        left_title="傳統 AI 工具",
        right_title="AI Agent",
        left_items=[
            "被動回應",
            "單次執行",
            "明確輸入輸出",
            "可預測性高",
            "適合固定流程",
        ],
        right_items=[
            "主動規劃",
            "多步驟決策",
            "動態環境互動",
            "行為可能超預期",
            "適合動態複雜任務",
        ],
        summary="考試會問：什麼時候該用 Agent、什麼時候傳統工具就夠了。",
    )
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ===== S12 . SILENT . 轉 — 過渡到風險管理 =====
    s = _blank(prs)
    draw_silent_page(
        s,
        "導入是規劃問題，\n風險管理是存活問題。\n接下來：AI 專案可能死在哪裡。",
    )
    add_footer(s, MODULE_CODE, 12, N_CONTENT, dark_bg=True)

    # ===== S13 . MATRIX . 轉 — 三類風險 =====
    s = _blank(prs)
    draw_three_blocks_flow(
        s,
        title="GenAI 風險三維度——不只是資安",
        blocks=[
            {"heading": "倫理風險",
             "items": [
                 "偏見 / 歧視",
                 "透明度不足",
                 "問責不明",
                 "Deepfake 濫用",
             ]},
            {"heading": "資安隱私",
             "items": [
                 "資料外洩",
                 "模型竊取",
                 "PII 暴露",
                 "Prompt Injection",
             ]},
            {"heading": "合規性",
             "items": [
                 "個資法",
                 "著作權法",
                 "AI 基本法草案",
                 "EU AI Act",
             ]},
        ],
        bottom_note="考生最常見盲區：把風險管理等於資安。三個維度都要準備。",
    )
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ===== S14 . TABLE . 轉 — 風險影響矩陣 (3x3) =====
    s = _blank(prs)
    draw_grid(
        s, rows=3, cols=3,
        title="風險影響矩陣：機率 \u00d7 衝擊 = 應對策略",
        cells=[
            # row 0: high impact
            {"label": "緩解", "sub": "高衝擊/低機率", "highlight": False},
            {"label": "迴避/緩解", "sub": "高衝擊/中機率", "highlight": True},
            {"label": "迴避", "sub": "高衝擊/高機率", "highlight": True},
            # row 1: mid impact
            {"label": "監控", "sub": "中衝擊/低機率", "highlight": False},
            {"label": "緩解", "sub": "中衝擊/中機率", "highlight": False},
            {"label": "迴避/緩解", "sub": "中衝擊/高機率", "highlight": True},
            # row 2: low impact
            {"label": "接受", "sub": "低衝擊/低機率", "highlight": False},
            {"label": "監控", "sub": "低衝擊/中機率", "highlight": False},
            {"label": "轉移", "sub": "低衝擊/高機率", "highlight": False},
        ],
        top=1.4, bottom=1.2,
        caption="高機率高衝擊 = 迴避或緩解。不要反射性選「全部迴避」。",
    )
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ===== S15 . PITFALL . 轉 — 只做資安 =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        title="PITFALL：風險管理 = 買防火牆 + 加密？",
        left_title="錯誤做法",
        right_title="正確做法",
        left_items=[
            "風險管理 = 買防火牆 + 加密",
            "倫理爭議爆發",
            "無法回應",
            "品牌受損",
        ],
        right_items=[
            "資安 + 倫理審查 + 合規檢核",
            "三管齊下",
            "全面防護",
            "有備無患",
        ],
        summary="2024 年有公司 AI 客服回答歧視性內容上新聞——問題不在資安，在倫理審查沒做。",
    )
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ===== S16 . PYRAMID . 合 — 風險緩解四層策略 =====
    s = _blank(prs)
    draw_pyramid_stack(
        s,
        title="風險緩解四層策略：從預防到復原",
        layers=[
            {"name": "復原",
             "caption": "影響評估 + 改善措施"},
            {"name": "回應",
             "caption": "事件應變計畫 + 通報機制"},
            {"name": "偵測",
             "caption": "持續監控 + 定期稽核"},
            {"name": "預防",
             "caption": "倫理審查 + 資料脫敏 + 存取控制"},
        ],
        thesis="考題問「最優先採取的策略」，答案通常是預防。",
    )
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ===== S17 . EXAM_DRILL . 合 — 導入評估/規劃考題 =====
    s = _blank(prs)
    add_title(s, "Exam Drill \u00b7 導入評估 / 規劃情境題")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "每題 90 秒，先自己選再看解析",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(0.8), Inches(1.9), Inches(11.5), Inches(4.0),
        "Q1  以下何者不屬於 AI 導入評估的三維度？\n"
        "    (A) 技術可行性 (B) 商業價值 (C) 組織準備度 (D) 模型準確度\n\n"
        "Q2  經濟部《AI 導入指引》的第一個階段是？\n"
        "    (A) 可行性評估 (B) 需求確認 (C) PoC (D) 系統開發\n\n"
        "Q3  TCO 不包含以下哪一項？\n"
        "    (A) 資料標註成本 (B) 員工訓練成本 (C) 競爭對手分析 (D) 退場成本",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "答案：Q1(D) 模型準確度是效能指標不是評估維度 \u00b7 Q2(B) 需求確認在最前面 \u00b7 Q3(C) 競爭對手分析不屬於 TCO",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ===== S18 . EXAM_DRILL . 合 — 風險管理考題 =====
    s = _blank(prs)
    add_title(s, "Exam Drill \u00b7 風險管理判斷題")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "L3 風險管理 \u00b7 情境判斷題",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(0.8), Inches(1.9), Inches(11.5), Inches(4.0),
        "Q1  AI 客服系統回答帶有性別歧視的內容，這屬於哪類風險？\n"
        "    (A) 資安風險 (B) 倫理風險 (C) 合規風險 (D) 技術風險\n\n"
        "Q2  某風險「發生機率高、衝擊程度高」，應採取什麼策略？\n"
        "    (A) 接受 (B) 監控 (C) 迴避或緩解 (D) 轉移\n\n"
        "Q3  以下何者屬於風險管理的「預防」層級措施？\n"
        "    (A) 事件通報機制 (B) 影響評估 (C) 設計階段倫理審查 (D) 持續監控",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "答案：Q1(B) 偏見歧視=倫理風險 \u00b7 Q2(C) 高機率高衝擊=迴避/緩解 \u00b7 Q3(C) 設計階段=預防",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ===== S19 . PRACTICE_PROMPT . 合 =====
    s = _blank(prs)
    add_title(s, "Practice \u00b7 綜合情境題")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "5 分鐘 \u00b7 能寫出來才算真的會",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.0), Inches(2.0), Inches(11.0), Inches(4.0),
        "情境：一間 200 人的中型製造業要導入 AI 品檢系統。\n\n"
        "請完成：\n"
        "(1) 三維評估摘要（技術/商業/組織各一句話）\n\n"
        "(2) 導入規劃五面向各一句話\n"
        "    （目標設置/資源分配/導入策略/因應措施/測試）\n\n"
        "(3) 列出前三大風險及緩解措施",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.3), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "這題整合了 L1/L2/L3 三個範疇，就是考試的綜合情境題型。",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ===== S20 . SILENT . 合 — 結語 =====
    s = _blank(prs)
    draw_silent_page(
        s,
        "考試考的不是你會不會用 AI，\n是你敢不敢為這個決策負責。\n"
        "導入評估看判斷力，風險管理看責任感\n"
        "\u2014\u2014這就是規劃師的核心能力。\n"
        "下一章：NLP / CV / 多模態技術辨識。",
    )
    add_footer(s, MODULE_CODE, 20, N_CONTENT, dark_bg=True)

    # == Copyright ==
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
