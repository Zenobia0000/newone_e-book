"""Ch01 deck -- AI 概念、分類與治理政策（完整版）
20 content slides + cover + copyright.

三大 Part：
  Part A -- AI 定義與分類（S1-S10）
  Part B -- 治理框架（S11-S16）
  Part C -- 倫理與收束（S17-S20）

受眾：iPAS 初級考生，所有案例走「考題情境」路線。
Aligned to chapters/Ch01_AI概念分類與治理政策/02_slides_design.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source, add_rect,
    set_solid_fill, set_no_fill, set_line, set_no_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_vs_two_col, draw_pyramid_stack,
    draw_three_blocks_flow, draw_image_placeholder,
    draw_inverted_thesis_box,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch01"
MODULE_TITLE = "AI 概念、分類與治理政策"
MODULE_SUBTITLE = "PRA 三要素 × 三三四分類 × 禁嚴透自——iPAS 基礎分拿穩"
TIME_MIN = 120
N_CONTENT = 20


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _draw_bridge_note(slide, text, y_inches=6.6):
    add_textbox(
        slide, T.MARGIN_X, Inches(y_inches),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        text,
        font_size=T.FONT_SMALL, color=T.PRIMARY, bold=True, italic=True,
        align=PP_ALIGN.CENTER,
    )


def build_ch01(output_path, image_registry=None):
    """Build Ch01 deck; 20 content slides + cover + copyright."""
    prs = _new_prs()

    # == Cover ==
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ===== S1 . MOTIVATION -- iPAS 考什麼？L111 佔分地圖 =====
    s = _blank(prs)
    draw_ask_page(
        s,
        "初級科目一：AI 定義分類 + AI 治理\n≈ 25-30% 考題",
        data_card={
            "label": "L111 佔分",
            "stat": "~15 題",
            "caption": "AI 定義分類 + 治理\n搞定這章 = 穩拿基本分",
        },
    )
    add_source(s, "iPAS 歷屆考古題分析 · L111 佔分估算")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # ===== S2 . CONCEPT-CARD -- AI 的定義：感知、推理、行動三要素 =====
    s = _blank(prs)
    add_title(s, "AI 的定義：感知、推理、行動三要素")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(0.5),
        "記憶口訣：PRA -- Perceive, Reason, Act",
        font_size=Pt(16), color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(2.0), Inches(6.5), Inches(3.5),
        "感知（Perceive）：接收環境資訊（影像、聲音、文字、數據）\n\n"
        "推理（Reason）：基於資訊做判斷或預測（不是固定規則）\n\n"
        "行動（Act）：根據推理結果採取回應\n\n"
        "判斷反射：三個都有 = AI / 缺推理 = 自動化",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(4.0),
        slot_name="PRA 三角形圖",
        description="三角形三頂點：感知 Perceive / 推理 Reason / 行動 Act，中心標示 AI。",
        size_hint="1200x900 px",
        placeholder_id="Ch01_S2_pra_triangle",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ===== S3 . EXAMPLE-I/O -- 考題實戰：判斷哪個是 AI =====
    s = _blank(prs)
    add_title(s, "考題實戰：判斷哪個是 AI")
    draw_editorial_table(
        s,
        header=["場景", "感知", "推理", "行動", "判定"],
        rows=[
            ["自動販賣機", "V", "X", "V", "不是 AI"],
            ["掃地機器人（路徑規劃型）", "V", "V", "V", "是 AI"],
            ["Netflix 推薦系統", "V", "V", "V", "是 AI"],
            ["銀行自動扣繳", "X", "X", "V", "不是 AI"],
        ],
        top=1.3,
        col_widths=[3.0, 1.0, 1.0, 1.0, 1.5],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.5), T.SLIDE_W - 2 * T.MARGIN_X, Inches(1.5),
        "考題反射 = 先找「推理」，有推理再看其他\n\n"
        "自動販賣機只是 if-else，不推理 → 不是 AI\n"
        "掃地機器人用演算法規劃路徑 → 是 AI\n"
        "銀行自動扣繳是固定日期扣固定金額 → 純自動化",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ===== S4 . CONCEPT-CARD -- AI 分類（一）：依能力分級 ANI / AGI / ASI =====
    s = _blank(prs)
    draw_pyramid_stack(
        s,
        title="AI 分類（一）：依能力分級 ANI / AGI / ASI",
        layers=[
            {"name": "ASI 超級 AI",
             "caption": "超越所有人類智力 → 理論假設"},
            {"name": "AGI 通用 AI",
             "caption": "能做任何人類智力任務 → 目前不存在"},
            {"name": "ANI 弱 AI / 窄 AI",
             "caption": "只能做一件事 → Siri、AlphaGo、ChatGPT"},
        ],
        thesis="考試重點：目前所有 AI 都是 ANI，AGI 尚未實現。",
    )
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ===== S5 . PITFALL (P1, P2) -- AI 分類混淆陷阱 =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="常見錯誤",
        right_title="正確認知",
        left_items=[
            "RPA 流程自動化是 ANI",
            "深度學習 = 非監督學習",
            "ChatGPT 是 AGI",
        ],
        right_items=[
            "RPA 不是 AI，根本不在分類裡",
            "分類維度不同，深度學習可搭配任何學習方式",
            "ChatGPT 是 ANI（只做語言任務）",
        ],
        title="PITFALL：AI 分類混淆陷阱",
        summary="解題反射：先確認「這是不是 AI」→ 再確認「問的是哪個分類維度」",
    )
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ===== S6 . CHECKPOINT -- AI 定義與分類快速檢核 =====
    s = _blank(prs)
    add_title(s, "Check Point . AI 定義與分類快速檢核")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 6 / 20 . 不看投影片作答",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(10.8), Inches(3.5),
        "Q1  Siri 屬於 ANI / AGI / ASI？\n\n"
        "Q2  自動扣款系統是不是 AI？\n\n"
        "Q3  「用大量有標籤資料訓練模型做分類」是哪種學習方式？",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.0), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "答案：Q1 ANI . Q2 不是（缺推理） . Q3 監督式學習",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ===== S7 . CONCEPT-CARD -- AI 分類（二）：依學習方式 =====
    s = _blank(prs)
    add_title(s, "AI 分類（二）：依學習方式")
    draw_editorial_table(
        s,
        header=["", "監督式 Supervised", "非監督式 Unsupervised", "強化式 Reinforcement"],
        rows=[
            ["資料", "有標籤", "無標籤", "有獎懲訊號"],
            ["任務", "分類/回歸", "分群/降維", "決策/控制"],
            ["實例", "垃圾郵件分類", "客戶分群", "AlphaGo 對弈"],
        ],
        top=1.3,
        col_widths=[1.2, 2.5, 2.5, 2.5],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.2), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "口訣：「有標籤找答案、沒標籤找結構、有獎懲學策略」",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ===== S8 . PRACTICE -- 學習方式場景配對 =====
    s = _blank(prs)
    add_title(s, "Practice . 學習方式場景配對")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 8 / 20 . 30 秒自己先配",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.0), Inches(1.8), Inches(11.0), Inches(4.5),
        "場景 1：垃圾郵件分類（有歷史標記）→ 監督式\n\n"
        "場景 2：商場客流分群（無預設群組）→ 非監督式\n\n"
        "場景 3：圍棋對弈學習 → 強化式\n\n"
        "場景 4：醫療影像判讀（有醫師標記）→ 監督式\n\n"
        "場景 5：推薦系統（根據點擊回饋優化）→ 強化式",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.3), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "關鍵判斷：有標籤→監督 / 找結構→非監督 / 有獎懲回饋→強化",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ===== S9 . PITFALL (P2) -- 技術路線 vs 學習方式的交叉混淆 =====
    s = _blank(prs)
    add_title(s, "PITFALL：技術路線 vs 學習方式的交叉混淆")
    draw_editorial_table(
        s,
        header=["", "規則式", "統計 ML", "深度學習", "生成式 AI"],
        rows=[
            ["監督式", "", "線性回歸", "CNN 圖片分類", "有監督微調"],
            ["非監督式", "", "K-Means", "Autoencoder", "GPT 預訓練"],
            ["強化式", "", "", "AlphaGo", "RLHF"],
        ],
        top=1.3,
        col_widths=[1.5, 1.5, 1.8, 2.0, 2.0],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.5), T.SLIDE_W - 2 * T.MARGIN_X, Inches(1.5),
        "技術路線和學習方式是兩個獨立維度，可以自由組合\n\n"
        "錯誤示範：「深度學習一定是非監督」 X\n"
        "正確認知：「深度學習可搭配任何學習方式」 V",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ===== S10 . CHECKPOINT -- AI 分類總整理 =====
    s = _blank(prs)
    add_title(s, "Check Point . AI 分類總整理：三三四")
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "能力分級（3 級）",
             "items": [
                 "ANI 弱 AI",
                 "AGI 通用 AI",
                 "ASI 超級 AI",
             ]},
            {"heading": "學習方式（3 種）",
             "items": [
                 "監督式",
                 "非監督式",
                 "強化式",
             ]},
            {"heading": "技術路線（4 條）",
             "items": [
                 "規則式",
                 "統計 ML",
                 "深度學習",
                 "生成式 AI",
             ]},
        ],
        bottom_note="考試記憶法：「三三四」→ 3 級能力 / 3 種學習 / 4 條技術路線",
    )
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ===== S11 . CONCEPT-CARD -- 歐盟 AI Act 四級風險分類 =====
    s = _blank(prs)
    draw_pyramid_stack(
        s,
        title="歐盟 AI Act 四級風險分類",
        layers=[
            {"name": "不可接受風險（禁止）",
             "caption": "社會信用評分、即時遠端生物辨識"},
            {"name": "高風險（嚴格監管）",
             "caption": "履歷篩選、信用評分、醫療診斷"},
            {"name": "有限風險（透明義務）",
             "caption": "聊天機器人須告知、Deepfake 須標示"},
            {"name": "最小風險（自由使用）",
             "caption": "垃圾郵件過濾、遊戲 AI、推薦系統"},
        ],
        thesis="記憶口訣：「禁嚴透自」→ 禁止 / 嚴格 / 透明 / 自由",
    )
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ===== S12 . PRACTICE -- 歐盟 AI Act 風險等級配對 =====
    s = _blank(prs)
    add_title(s, "Practice . 歐盟 AI Act 風險等級配對")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 12 / 20 . 45 秒自己先配",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.0), Inches(1.8), Inches(11.0), Inches(4.5),
        "場景 1：政府用 AI 對公民做社會信用評分 → 不可接受\n\n"
        "場景 2：企業用 AI 自動篩選履歷 → 高風險\n\n"
        "場景 3：電商用 AI 推薦商品 → 最小風險\n\n"
        "場景 4：客服聊天機器人 → 有限風險\n\n"
        "場景 5：AI 自動駕駛系統 → 高風險\n\n"
        "場景 6：AI 生成 Deepfake 影片 → 有限風險（須標示）",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ===== S13 . PITFALL (P3) -- 風險等級記反：不可接受 vs 高風險 =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="常見錯誤",
        right_title="正確認知",
        left_items=[
            "社會信用評分 → 高風險（嚴格監管就好）",
            "履歷篩選 AI → 不可接受（太嚴了）",
        ],
        right_items=[
            "社會信用評分 → 不可接受（直接禁止）",
            "履歷篩選 AI → 高風險（嚴格監管但允許）",
        ],
        title="PITFALL：不可接受 vs 高風險 記反",
        summary="判斷標準：能不能透過管控降低風險？不能→禁止；可以→高風險",
    )
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ===== S14 . CHECKPOINT -- 歐盟 AI Act 情境速判 =====
    s = _blank(prs)
    add_title(s, "Check Point . 歐盟 AI Act 情境速判")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 14 / 20 . 每題 15 秒",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(10.8), Inches(3.5),
        "Q1  警方在廣場用 AI 即時辨識所有路人身份\n"
        "     → 不可接受\n\n"
        "Q2  保險公司用 AI 自動核保\n"
        "     → 高風險\n\n"
        "Q3  新聞網站用 AI 摘要文章\n"
        "     → 最小風險\n\n"
        "Q4  AI 語音助手未告知用戶它是 AI\n"
        "     → 違反有限風險的透明義務",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ===== S15 . CONCEPT-CARD -- 台灣 AI 治理：公部門手冊 vs 金管會指引 =====
    s = _blank(prs)
    add_title(s, "台灣 AI 治理：公部門手冊 vs 金管會指引")
    draw_editorial_table(
        s,
        header=["", "數位部《公部門 AI 手冊》", "金管會《金融業 AI 指引》"],
        rows=[
            ["適用範圍", "行政院及所屬機關", "金融業（銀行/保險/證券）"],
            ["核心原則", "人為監督、透明可解釋\n隱私保護、安全可靠", "公平性、可課責性\n透明性、客戶權益保護"],
            ["性質", "參考指引（非法律強制）", "行政指導（具監理效力）"],
            ["考試重點", "「誰用？」→ 政府機關", "「誰用？」→ 金融業"],
        ],
        top=1.3,
        col_widths=[1.5, 3.5, 3.5],
    )
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ===== S16 . PRACTICE -- 治理框架適用情境 =====
    s = _blank(prs)
    add_title(s, "Practice . 治理框架適用情境")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 16 / 20 . 快速配對",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.0), Inches(1.8), Inches(11.0), Inches(3.8),
        "情境 A：某銀行用 AI 審核貸款 → 金管會指引\n\n"
        "情境 B：某市政府用 AI 分配社福資源 → 公部門手冊\n\n"
        "情境 C：某電商用 AI 推薦商品 → 目前無專門框架（一般法規）\n\n"
        "情境 D：某保險公司用 AI 自動核保 → 金管會指引",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    # Warning box
    warn_y = Inches(5.6)
    warn_box = add_rect(slide=s, x=T.MARGIN_X, y=warn_y,
                        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(0.6))
    set_solid_fill(warn_box, T.PRIMARY)
    set_no_line(warn_box)
    add_textbox(
        s, T.MARGIN_X + Inches(0.3), warn_y,
        T.SLIDE_W - 2 * T.MARGIN_X - Inches(0.6), Inches(0.6),
        "陷阱提示：不是所有 AI 應用都有對應的治理框架！",
        font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ===== S17 . CONCEPT-CARD -- AI 治理三大支柱 =====
    s = _blank(prs)
    add_title(s, "AI 治理三大支柱")
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "透明性 Transparency",
             "items": [
                 "使用者知道 AI 在做決定",
                 "AI 決策邏輯可被解釋",
                 "關鍵字：黑箱、不解釋",
             ]},
            {"heading": "可課責性 Accountability",
             "items": [
                 "AI 出錯有明確責任歸屬",
                 "有人工覆核與申訴機制",
                 "關鍵字：沒人負責",
             ]},
            {"heading": "公平性 Fairness",
             "items": [
                 "不因種族/性別/年齡歧視",
                 "訓練資料無偏見",
                 "關鍵字：偏見、歧視",
             ]},
        ],
        bottom_note="情境題解題法：看到場景 → 判斷違反哪根柱子 → 選對應改善措施",
    )
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ===== S18 . CHECKPOINT -- AI 倫理情境判斷 =====
    s = _blank(prs)
    add_title(s, "Check Point . AI 倫理情境判斷")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 18 / 20 . 三題練到反射級",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(10.8), Inches(3.5),
        "Q1  AI 面試系統對女性求職者評分偏低\n"
        "     → 違反公平性\n\n"
        "Q2  AI 醫療診斷不提供任何決策理由\n"
        "     → 違反透明性\n\n"
        "Q3  AI 信貸系統拒絕申請，客戶無法申訴也找不到負責人\n"
        "     → 違反可課責性",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.0), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "考試反射：歧視→公平 / 不告知→透明 / 無人負責→可課責",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ===== S19 . PYRAMID -- 全章收束：AI 知識 + 治理金字塔 =====
    s = _blank(prs)
    draw_pyramid_stack(
        s,
        title="Ch01 收束：AI 知識 + 治理金字塔",
        layers=[
            {"name": "AI 倫理三支柱",
             "caption": "怎麼管 AI → 透明/可課責/公平"},
            {"name": "國際與台灣治理框架",
             "caption": "誰在管 AI → 歐盟 AI Act / 台灣指引"},
            {"name": "AI 分類三軸線",
             "caption": "AI 有哪些種類 → 三三四"},
            {"name": "AI 定義與三要素",
             "caption": "什麼是 AI → PRA 感知推理行動"},
        ],
        thesis="定義佔基礎分、治理佔情境分——從底到頂都會考。",
    )
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ===== S20 . MOTIVATION -- 銜接 Ch02：AI 的燃料是資料 =====
    s = _blank(prs)
    draw_silent_page(
        s,
        "AI 定義 V / AI 分類 V / 治理框架 V / 倫理原則 V\n"
        "下一站：Ch02 資料是 AI 的燃料",
    )
    add_footer(s, MODULE_CODE, 20, N_CONTENT, dark_bg=True)

    # == Copyright ==
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
