"""Ch10 deck — 建模調參、評估與 ML 治理（完整版）
22 content slides + cover + copyright.

三大 Part：
  Part A — 特徵工程（S1-S8）
  Part B — 評估指標與調參（S9-S17）
  Part C — ML 治理與收束（S18-S22）

受眾：iPAS 中級科目3 考生，調好評好管好導向。
Aligned to chapters/Ch10_建模調參評估與ML治理/02_slides_design.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source, add_rect,
    set_solid_fill, set_no_fill, set_line, set_no_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_vs_two_col, draw_pyramid_stack,
    draw_three_blocks_flow, draw_image_placeholder,
    draw_inverted_thesis_box, draw_code_panel,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch10"
MODULE_TITLE = "建模調參、評估與 ML 治理"
MODULE_SUBTITLE = "特徵工程 × 評估指標 × Loss Curve 診斷 × GDPR——調好評好管好"
TIME_MIN = 120
N_CONTENT = 22


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _draw_bridge_note(slide, text, y_inches=6.6):
    add_textbox(
        slide, T.MARGIN_X, Inches(y_inches),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        text,
        font_size=T.FONT_SMALL, color=T.PRIMARY, bold=True, italic=True,
        align=PP_ALIGN.CENTER,
    )


def build_ch10(output_path, image_registry=None):
    """Build Ch10 deck; 22 content slides + cover + copyright."""
    prs = _new_prs()

    # == Cover ==
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───── S1 · MOTIVATION — 模型做出來只是開始 ─────
    s = _blank(prs)
    draw_ask_page(
        s,
        "科目3 的另一半考題不是考演算法，\n"
        "而是考你會不會調參、會不會評估、懂不懂治理。",
        data_card={
            "label": "歷屆分析",
            "stat": "~40-50%",
            "caption": "調參+評估+治理佔比\nL233 + L234 約 20-25 題",
        },
    )
    add_source(s, "iPAS 歷屆考古題分析 · 科目3 樣題")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # ───── S2 · ASK — Accuracy 95% 真的可以上線嗎 ─────
    s = _blank(prs)
    draw_ask_page(
        s,
        "你的模型 Accuracy 95%\n——真的可以上線嗎？",
        data_card={
            "label": "三個隱藏地雷",
            "stat": "?",
            "caption": "資料 95% 同一類？\nGDPR 不合規？\n模型有偏見？",
        },
    )
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───── S3 · SILENT — 立論 ─────
    s = _blank(prs)
    draw_silent_page(
        s,
        "模型做出來只是開始。\n調好、評好、管好才是考試重點。",
    )
    add_footer(s, MODULE_CODE, 3, N_CONTENT, dark_bg=True)

    # ───── S4 · CONCEPT-CARD — 特徵工程三步驟 ─────
    s = _blank(prs)
    add_title(s, "特徵工程三步驟：轉換 / 萃取 / 選擇")
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "特徵轉換",
             "items": [
                 "One-Hot Encoding",
                 "標準化/正規化",
                 "Log 轉換（偏態分佈）",
             ]},
            {"heading": "特徵萃取",
             "items": [
                 "多項式/交互特徵",
                 "TF-IDF（文字→數值）",
                 "PCA（降維=萃取）",
             ]},
            {"heading": "特徵選擇",
             "items": [
                 "相關係數篩選",
                 "Feature Importance",
                 "L1 Regularization（自動歸零）",
             ]},
        ],
        bottom_note="轉換=格式化 / 萃取=創造新特徵 / 選擇=去蕪存菁",
    )
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───── S5 · CONCEPT-CARD — 標準化 vs 正規化 ─────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="標準化 Z-score",
        right_title="正規化 Min-Max",
        left_items=[
            "z = (x - mean) / std",
            "結果：均值 0、標準差 1",
            "適用：SVM / KNN / 神經網路",
            "不受 outlier 影響太大",
        ],
        right_items=[
            "x' = (x - min) / (max - min)",
            "結果：壓縮到 [0, 1]",
            "適用：影像像素 / NN 輸入",
            "對 outlier 敏感",
        ],
        title="標準化(Z-score) vs 正規化(Min-Max)",
        summary="口訣：Z 標 M 正（Z-score = 標準化、Min-Max = 正規化）",
    )
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───── S6 · CONCEPT-CARD — 標註品質與訓練前配置 ─────
    s = _blank(prs)
    add_title(s, "標註品質 = 模型天花板")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "Garbage In, Garbage Out — 再好的演算法也救不了爛標註",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    draw_flow_chain(
        s,
        title="",
        nodes=[
            {"label": "標註品質", "sub": "Cohen's Kappa\nInter-Annotator", "highlight": True},
            {"label": "Loss Function", "sub": "MSE / CE"},
            {"label": "優化器", "sub": "SGD / Adam"},
            {"label": "學習率", "sub": "起手 0.001"},
        ],
        y=3.0,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.5), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "樣題 Q5：標註品質對模型的影響 · 標註不一致 → 模型學到噪音",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───── S7 · PITFALL — 標準化和正規化搞混 ─────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="錯誤做法",
        right_title="正確做法",
        left_items=[
            "把 Min-Max 叫做標準化",
            "把 Z-score 叫做正規化",
            "中文翻譯混亂 → 搞混",
        ],
        right_items=[
            "Z-score = 標準化 (Standardization)",
            "Min-Max = 正規化 (Normalization)",
            "記公式更保險",
            "口訣：Z 標 M 正",
        ],
        title="PITFALL：標準化和正規化搞混",
        summary="解決方案：不記中文名稱，記公式。看到 mean/std → 標準化；看到 min/max → 正規化",
    )
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───── S8 · CHECKPOINT — 特徵工程分類 ─────
    s = _blank(prs)
    add_title(s, "Check Point · 特徵工程分類")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 8 / 22 · 以下操作屬於轉換、萃取、還是選擇？",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    draw_editorial_table(
        s,
        header=["操作", "分類", "理由"],
        rows=[
            ["One-Hot Encoding", "轉換", "格式化類別特徵"],
            ["TF-IDF", "萃取", "從文字創造數值特徵"],
            ["PCA 降維", "萃取", "從多特徵萃取主成分"],
            ["相關係數篩選", "選擇", "根據相關性挑特徵"],
            ["多項式特徵", "萃取", "x → x, x^2, x^3"],
            ["L1 Regularization", "選擇", "自動歸零不重要特徵"],
        ],
        top=2.0,
        col_widths=[2.5, 1.2, 3.0],
    )
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───── S9 · CONCEPT-CARD — 分類評估指標全家福 ─────
    s = _blank(prs)
    add_title(s, "分類評估指標全家福")
    draw_editorial_table(
        s,
        header=["", "預測正", "預測負"],
        rows=[
            ["實際正", "TP（抓到）", "FN（漏掉）"],
            ["實際負", "FP（誤報）", "TN（正確排除）"],
        ],
        top=1.3,
        col_widths=[1.5, 2.5, 2.5],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(3.3), T.SLIDE_W - 2 * T.MARGIN_X, Inches(3.0),
        "Accuracy = (TP + TN) / 全部 → 整體正確率\n"
        "Precision = TP / (TP + FP) → 預測為正的有多少是對的\n"
        "Recall = TP / (TP + FN) → 真正為正的你抓到多少\n"
        "F1 = 2 x Precision x Recall / (P + R) → 調和平均\n\n"
        "心法：漏掉嚴重追 Recall / 誤報嚴重追 Precision / 兩者都重要看 F1",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───── S10 · CONCEPT-CARD — AUC-ROC 與 R² ─────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="AUC-ROC（分類）",
        right_title="R²（迴歸）",
        left_items=[
            "ROC：X 軸 FPR / Y 軸 TPR",
            "AUC = 曲線下面積",
            "0.5 = 隨機猜 / 1.0 = 完美",
            "不受閾值影響",
        ],
        right_items=[
            "R^2 = 1 - (SS_res / SS_tot)",
            "模型解釋了多少 % 變異",
            "R^2 = 1 完美 / 0 = 猜平均",
            "迴歸專用（不是分類！）",
        ],
        title="AUC-ROC vs R²：分類迴歸不要混用",
        summary="分類用 AUC / F1，迴歸用 R^2 / MSE，不要混用。",
    )
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───── S11 · PITFALL — Imbalanced Data 上用 Accuracy ─────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="錯誤做法",
        right_title="正確做法",
        left_items=[
            "詐騙偵測 10,000 筆僅 20 筆詐騙",
            "全猜「不是詐騙」",
            "Accuracy = 99.8%",
            "報告「模型很準」",
        ],
        right_items=[
            "Precision = 0/0 = undefined",
            "Recall = 0/20 = 0%",
            "一筆詐騙都沒抓到！",
            "應看 Precision/Recall/F1/AUC-PR",
        ],
        title="PITFALL：Imbalanced Data 上用 Accuracy",
        summary="樣題 Q10：Accuracy 計算——但更重要的是知道什麼時候 Accuracy 沒意義",
    )
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───── S12 · CHECKPOINT — 指標選用判斷 ─────
    s = _blank(prs)
    add_title(s, "Check Point · 指標選用判斷")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 12 / 22 · 選出最適合的評估指標",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    draw_editorial_table(
        s,
        header=["場景", "優先指標", "理由"],
        rows=[
            ["癌症篩檢", "高 Recall", "不能漏診比誤診重要"],
            ["垃圾郵件過濾", "高 Precision", "不能把正常信當垃圾"],
            ["信用評分", "AUC / F1", "整體表現+不受閾值影響"],
            ["房價預測", "R^2 / MSE", "迴歸任務（不是 Accuracy！）"],
        ],
        top=2.0,
        col_widths=[2.0, 1.8, 3.5],
    )
    _draw_bridge_note(s, "心法：漏掉嚴重追 Recall，誤報嚴重追 Precision，兩者都重要看 F1")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───── S13 · CONCEPT-CARD — Loss Function 與 Loss Curve ─────
    s = _blank(prs)
    add_title(s, "Loss Function 與 Loss Curve 三種模式")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(1.5),
        "Loss Function = 模型成績單的反面（越低越好）\n"
        "  迴歸：MSE = SUM(y-y')^2 / n\n"
        "  分類：Cross-Entropy = -SUM[y log(p) + (1-y) log(1-p)]",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "A 正常收斂",
             "items": [
                 "train/val 都下降",
                 "最終接近",
                 "繼續訓練或微調",
             ]},
            {"heading": "B Overfitting",
             "items": [
                 "train 持續下降",
                 "val 先降後升",
                 "解法：Early Stopping/Dropout",
             ]},
            {"heading": "C Underfitting",
             "items": [
                 "train/val 都很高",
                 "下降很慢",
                 "解法：增加模型複雜度",
             ]},
        ],
    )
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───── S14 · CONCEPT-CARD — 調參策略 ─────
    s = _blank(prs)
    add_title(s, "調參策略：四大參數 x 影響方向")
    draw_editorial_table(
        s,
        header=["參數", "調大", "調小", "常見值"],
        rows=[
            ["Learning Rate", "震盪不收斂", "收斂太慢", "0.001(Adam) / 0.01(SGD)"],
            ["Batch Size", "穩定但泛化差", "噪音大但泛化好", "32 / 64 / 128 / 256"],
            ["Regularization", "降 overfit 可能 underfit", "可能 overfit", "Dropout 0.2-0.5"],
            ["Early Stopping", "——", "——", "val loss N epoch 沒改善就停"],
        ],
        top=1.3,
        col_widths=[1.8, 2.2, 2.2, 3.0],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.5), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "Early Stopping = 最簡單有效的防 overfitting 手段",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───── S15 · PITFALL — 不看 Loss Curve 就調參 ─────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="錯誤做法",
        right_title="正確做法",
        left_items=[
            "不畫 Loss Curve",
            "直接跑 GridSearch",
            "三小時後得到一組參數",
            "但不知道為什麼好",
        ],
        right_items=[
            "先畫 Loss Curve",
            "判斷 overfit 還是 underfit",
            "針對性調參",
            "先診斷再治療",
        ],
        title="PITFALL：不看 Loss Curve 就調參",
        summary="Overfit → 加 regularization / Underfit → 增加模型複雜度",
    )
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───── S16 · CHECKPOINT — Loss Curve 診斷練習 ─────
    s = _blank(prs)
    add_title(s, "Check Point · Loss Curve 診斷練習")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 16 / 22 · 看圖判斷模型狀態",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    draw_editorial_table(
        s,
        header=["圖", "train loss", "val loss", "診斷", "解法"],
        rows=[
            ["A", "穩定下降", "穩定下降且接近", "正常收斂", "繼續或微調"],
            ["B", "持續下降", "先降後升", "Overfitting", "Early Stop/Dropout"],
            ["C", "很高下降慢", "很高下降慢", "Underfitting", "加複雜度/加特徵"],
        ],
        top=2.0,
        col_widths=[0.6, 1.5, 1.8, 1.5, 2.0],
    )
    _draw_bridge_note(s, "val loss 上升的轉折點 = Early Stopping 該停的地方")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───── S17 · CONCEPT-CARD — 處理 Class Imbalance ─────
    s = _blank(prs)
    add_title(s, "處理 Class Imbalance 三策略")
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "Oversampling",
             "items": [
                 "Random Oversampling",
                 "SMOTE（合成新樣本）",
                 "不丟資料但可能 overfit",
             ]},
            {"heading": "Undersampling",
             "items": [
                 "Random Undersampling",
                 "快但丟資料",
                 "可能丟失重要資訊",
             ]},
            {"heading": "演算法層面",
             "items": [
                 "class_weight='balanced'",
                 "調整分類閾值",
                 "調損失函數權重",
             ]},
        ],
    )
    # Warning box
    warn_y = Inches(5.8)
    warn_box = add_rect(slide=s, x=T.MARGIN_X, y=warn_y,
                        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(0.6))
    set_solid_fill(warn_box, T.PRIMARY)
    set_no_line(warn_box)
    add_textbox(
        s, T.MARGIN_X + Inches(0.3), warn_y,
        T.SLIDE_W - 2 * T.MARGIN_X - Inches(0.6), Inches(0.6),
        "致命注意：重抽樣只能在 train set 內做！在 split 前做 = data leakage",
        font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───── S18 · CONCEPT-CARD — GDPR 被遺忘權 ─────
    s = _blank(prs)
    add_title(s, "GDPR 被遺忘權：對 ML 的影響")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), T.SLIDE_W - 2 * T.MARGIN_X, Inches(4.5),
        "GDPR（General Data Protection Regulation）= 歐盟資料保護法規\n\n"
        "被遺忘權（Right to Erasure）：\n"
        "  用戶有權要求刪除其個人資料\n"
        "  組織必須在合理時間內完成刪除\n"
        "  包含：資料庫記錄 + 備份 + 衍生資料\n\n"
        "對 ML 的影響：\n"
        "  模型可能需要重新訓練（移除該用戶的訓練資料後）\n"
        "  需要建立 data lineage（資料血統追蹤）\n"
        "  需要記錄哪些資料用於訓練哪個模型\n\n"
        "樣題 Q12：GDPR 被遺忘權的具體要求",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ───── S19 · CONCEPT-CARD — 演算法偏誤與公平性 ─────
    s = _blank(prs)
    add_title(s, "演算法偏誤與公平性")
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "偏誤三大來源",
             "items": [
                 "訓練資料偏差",
                 "特徵選擇偏差（敏感特徵）",
                 "標籤偏差（歷史偏見）",
             ]},
            {"heading": "公平性指標",
             "items": [
                 "Demographic Parity",
                 "Equalized Odds",
                 "Predictive Parity",
             ]},
            {"heading": "調整策略",
             "items": [
                 "移除敏感特徵",
                 "重新平衡訓練資料",
                 "後處理校正閾值",
             ]},
        ],
        bottom_note="演算法不是故意歧視——是從有偏見的資料中學到了偏見",
    )
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───── S20 · PRACTICE — sklearn 程式排序題 ─────
    s = _blank(prs)
    add_title(s, "Practice · sklearn 程式排序題")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 20 / 22 · 樣題 Q15 題型",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.8),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.0),
        label="打亂順序 → 正確排列",
        code=(
            "# 正確順序：B → G → H → D → E → A → F → C\n"
            "\n"
            "(B) from sklearn.ensemble import RandomForestClassifier\n"
            "(G) from sklearn.model_selection import train_test_split\n"
            "(H) from sklearn.metrics import classification_report\n"
            "(D) X_train, X_test, y_train, y_test = train_test_split(...)\n"
            "(E) model = RandomForestClassifier(n_estimators=100)\n"
            "(A) model.fit(X_train, y_train)\n"
            "(F) y_pred = model.predict(X_test)\n"
            "(C) print(classification_report(y_test, y_pred))"
        ),
        bullets=[
            "import 永遠在最前面",
            "evaluate 永遠在最後面",
            "fit 一定在 predict 前面",
            "流程：import→切→建→訓→預→評",
        ],
    )
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ───── S21 · PYRAMID — 建模三層功夫 ─────
    s = _blank(prs)
    draw_pyramid_stack(
        s,
        title="Ch10 收束：建模三層功夫",
        layers=[
            {"name": "治理到位",
             "caption": "合規+公平 — GDPR/偏見/可解釋性"},
            {"name": "評估指標選對",
             "caption": "不要用錯指標騙自己 — P/R/F1/AUC"},
            {"name": "特徵工程做好",
             "caption": "垃圾進垃圾出 — 轉換/萃取/選擇"},
        ],
        thesis="底層不穩，上面全白做。",
    )
    add_footer(s, MODULE_CODE, 21, N_CONTENT)

    # ───── S22 · MOTIVATION — 課程總結 ─────
    s = _blank(prs)
    draw_ask_page(
        s,
        "從 Ch01 AI 概念到 Ch10 ML 治理\n"
        "——你已完成 iPAS 中級科目3 完整備考。",
        data_card={
            "label": "核心金句",
            "stat": "2句",
            "caption": "演算法考選型不考推導\n指標要會選不只會算",
        },
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "考試不只考技術，更考判斷力。祝你考試順利！",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 22, N_CONTENT)

    # == Copyright ==
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
