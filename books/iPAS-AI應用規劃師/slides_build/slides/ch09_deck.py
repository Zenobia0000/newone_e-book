"""Ch09 deck — 機器學習演算法與深度學習（完整版）
22 content slides + cover + copyright.

三大 Part：
  Part A — ML 演算法選型（S1-S12）
  Part B — 深度學習架構與激活函數（S13-S18）
  Part C — 實作與收束（S19-S22）

受眾：iPAS 中級科目3 考生，選型判斷力導向。
Aligned to chapters/Ch09_機器學習演算法與深度學習/02_slides_design.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source, add_rect,
    set_solid_fill, set_no_fill, set_line, set_no_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_vs_two_col, draw_pyramid_stack,
    draw_three_blocks_flow, draw_image_placeholder,
    draw_inverted_thesis_box, draw_code_panel,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch09"
MODULE_TITLE = "機器學習演算法與深度學習"
MODULE_SUBTITLE = "ML 全景地圖 × CNN 三架構 × 激活函數四天王——選型判斷力"
TIME_MIN = 120
N_CONTENT = 22


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _draw_bridge_note(slide, text, y_inches=6.6):
    add_textbox(
        slide, T.MARGIN_X, Inches(y_inches),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        text,
        font_size=T.FONT_SMALL, color=T.PRIMARY, bold=True, italic=True,
        align=PP_ALIGN.CENTER,
    )


def build_ch09(output_path, image_registry=None):
    """Build Ch09 deck; 22 content slides + cover + copyright."""
    prs = _new_prs()

    # == Cover ==
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───── S1 · MOTIVATION — 科目3 演算法題佔比最高 ─────
    s = _blank(prs)
    draw_ask_page(
        s,
        "科目3 有一半的題目跟演算法有關\n——但它不考推導，考的是選型判斷力。",
        data_card={
            "label": "歷屆分析",
            "stat": "~50%",
            "caption": "演算法+架構辨識佔比\nL231 + L232 約 25 題",
        },
    )
    add_source(s, "iPAS 歷屆考古題分析 · 科目3 樣題")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # ───── S2 · ASK — 10 種演算法怎麼選 ─────
    s = _blank(prs)
    draw_ask_page(
        s,
        "Linear Regression / Logistic Regression / Decision Tree\n"
        "Random Forest / SVM / KNN / K-means / DBSCAN\n"
        "PCA / CNN / RNN / Transformer...\n"
        "這麼多，考試到底考什麼？",
    )
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───── S3 · SILENT — 立論 ─────
    s = _blank(prs)
    draw_silent_page(
        s,
        "演算法不用全背，但必須知道每個的定位。\n"
        "考試考的是選型判斷力，不是公式記憶力。",
    )
    add_footer(s, MODULE_CODE, 3, N_CONTENT, dark_bg=True)

    # ───── S4 · CONCEPT-CARD — ML 演算法全景地圖 ─────
    s = _blank(prs)
    add_title(s, "ML 演算法全景地圖")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(4.5),
        "監督式學習\n"
        "  迴歸：Linear Regression\n"
        "  分類：Logistic Regression / Decision Tree\n"
        "        Random Forest / SVM / KNN\n\n"
        "非監督式學習\n"
        "  分群：K-means / DBSCAN\n"
        "  降維：PCA\n\n"
        "深度學習\n"
        "  影像：CNN（VGG / Inception / ResNet）\n"
        "  序列：RNN / LSTM / Transformer",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(4.5),
        slot_name="ML 演算法全景樹狀圖",
        description="樹狀心智圖：監督(迴歸/分類) / 非監督(分群/降維) / 深度學習(CNN/RNN/Transformer)。",
        size_hint="1200x1000 px",
        placeholder_id="Ch09_S4_ml_landscape",
        registry=image_registry,
    )
    _draw_bridge_note(s, "先記住大分類，後面一個一個填細節。")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───── S5 · CONCEPT-CARD — Linear Regression + R² ─────
    s = _blank(prs)
    add_title(s, "Linear Regression + R² 判定係數")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(0.6),
        "y = wx + b    R² = 1 - (SS_res / SS_tot)",
        font_size=Pt(18), color=T.PRIMARY, bold=True,
        family=T.FONT_MONO,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(2.2), Inches(6.5), Inches(3.5),
        "Linear Regression：假設輸入與輸出為線性關係\n\n"
        "R²（判定係數）= 模型解釋了多少百分比的變異\n"
        "  R² = 1 → 完美預測\n"
        "  R² = 0 → 跟猜平均值一樣\n"
        "  R² < 0 → 比猜平均值還爛\n\n"
        "考點：R² 越接近 1 越好",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(3.8),
        slot_name="散點圖 + 迴歸線 + R²",
        description="左側散點圖疊加迴歸直線，右側標示 R² 數值意義。",
        size_hint="1200x900 px",
        placeholder_id="Ch09_S5_linear_r2",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───── S6 · CONCEPT-CARD — LR / DT / RF 三欄比較 ─────
    s = _blank(prs)
    add_title(s, "Logistic Regression / Decision Tree / Random Forest")
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "Logistic Regression",
             "items": [
                 "Sigmoid → 輸出 0-1 機率",
                 "線性決策邊界",
                 "適合近似線性的場景",
             ]},
            {"heading": "Decision Tree",
             "items": [
                 "if-else 規則分裂",
                 "高解釋性",
                 "容易 overfit",
             ]},
            {"heading": "Random Forest",
             "items": [
                 "多棵 DT 投票（Bagging）",
                 "降低 variance",
                 "穩定的 baseline 首選",
             ]},
        ],
        bottom_note="選型口訣：資料量小/要解釋 → LR；規則可視化 → DT；要穩定 → RF",
    )
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───── S7 · CONCEPT-CARD — SVM 與 KNN ─────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="SVM 支持向量機",
        right_title="KNN K-近鄰",
        left_items=[
            "找最大間隔的超平面",
            "Kernel Trick 可處理非線性",
            "高維資料有效",
            "大資料量慢",
        ],
        right_items=[
            "用最近 K 個鄰居投票",
            "不需要訓練（懶學習）",
            "預測慢",
            "維度詛咒嚴重",
        ],
        title="SVM vs KNN",
        summary="共同點：兩者都對特徵縮放敏感——考試會考。",
    )
    draw_image_placeholder(
        s, Inches(5.0), Inches(6.2), Inches(3.2), Inches(0.5),
        slot_name="SVM 最大間隔 / KNN 鄰居示意",
        description="左邊 SVM 超平面+間隔，右邊 KNN 鄰居投票。",
        size_hint="800x200 px",
        placeholder_id="Ch09_S7_svm_knn",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───── S8 · CHECKPOINT — 演算法選型快問快答 ─────
    s = _blank(prs)
    add_title(s, "Check Point · 演算法選型快問快答")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 8 / 22 · 不看投影片作答",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(10.8), Inches(3.5),
        "Q1  電商客戶分群 → ?\n\n"
        "Q2  房價預測 → ?\n\n"
        "Q3  垃圾郵件分類 → ?\n\n"
        "Q4  人臉辨識 → ?\n\n"
        "Q5  工廠設備異常偵測（無標註）→ ?",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.8), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.7),
        "Q1 K-means/DBSCAN(非監督) · Q2 Linear Regression(迴歸)\n"
        "Q3 LR/RF(分類) · Q4 CNN(深度學習) · Q5 DBSCAN/Isolation Forest(非監督)",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───── S9 · CONCEPT-CARD — K-means vs DBSCAN ─────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="K-means",
        right_title="DBSCAN",
        left_items=[
            "需預先指定 K",
            "假設球狀群集",
            "對 outlier 敏感",
            "計算快",
        ],
        right_items=[
            "不需指定 K",
            "能找任意形狀群集",
            "自動標記噪音點",
            "需設定 eps 和 min_samples",
        ],
        title="K-means vs DBSCAN",
        summary="知道群數+球狀 → K-means；不知群數+不規則 → DBSCAN",
    )
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───── S10 · CONCEPT-CARD — PCA 降維 ─────
    s = _blank(prs)
    add_title(s, "PCA 降維：線性降維的標準工具")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(4.0),
        "原理：找資料變異量最大的方向（第一主成分）\n"
        "      然後找正交的第二大方向...\n\n"
        "保留 95% 累積變異通常就夠\n\n"
        "考試必記：\n"
        "  PCA 是線性降維\n"
        "  需要先標準化\n"
        "  主成分彼此正交\n"
        "  不適合非線性結構\n\n"
        "用途：特徵太多時降維 / 資料視覺化（降到 2D/3D）",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(4.0),
        slot_name="PCA 投影 + 累積變異圖",
        description="左側 PCA 投影示意，右側累積變異解釋曲線（保留 95% 處劃虛線）。",
        size_hint="1200x900 px",
        placeholder_id="Ch09_S10_pca",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───── S11 · PITFALL — MapReduce 不是 ML 但會考 ─────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="錯誤做法",
        right_title="正確做法",
        left_items=[
            "跳過 MapReduce",
            "覺得不是 ML 範圍",
            "結果樣題 Q1 直接考",
        ],
        right_items=[
            "理解 Map → Shuffle → Reduce",
            "Map：拆解+處理",
            "Reduce：合併結果",
            "屬於 L231 演算法效率範疇",
        ],
        title="PITFALL：MapReduce 不是 ML 但會考",
        summary="範例：計算全校平均分 → Map(每班算總分和人數) → Reduce(合併除以總人數)",
    )
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───── S12 · CHECKPOINT — 演算法特性配對 ─────
    s = _blank(prs)
    add_title(s, "Check Point · 演算法特性配對")
    draw_editorial_table(
        s,
        header=["特性", "對應演算法"],
        rows=[
            ["不需指定群數", "DBSCAN"],
            ["需要 Kernel Trick 處理非線性", "SVM"],
            ["輸出 0-1 之間的機率值", "Logistic Regression"],
            ["對特徵縮放不敏感", "Decision Tree / Random Forest"],
            ["線性降維、需要標準化", "PCA"],
            ["懶學習、不建模型", "KNN"],
        ],
        top=1.3,
        col_widths=[3.0, 3.0],
    )
    _draw_bridge_note(s, "tree-based 不看距離 → 不需縮放；SVM/KNN 看距離 → 必須縮放")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───── S13 · CONCEPT-CARD — CNN 經典架構：VGG / Inception / ResNet ─────
    s = _blank(prs)
    add_title(s, "CNN 三大經典架構：VGG / Inception / ResNet")
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "VGG",
             "items": [
                 "3x3 小卷積核",
                 "簡單堆疊加深",
                 "16/19 層",
                 "概念直覺但參數多",
             ]},
            {"heading": "Inception (GoogLeNet)",
             "items": [
                 "1x1/3x3/5x5 卷積核並行",
                 "加寬不加深",
                 "多尺度特徵擷取",
             ]},
            {"heading": "ResNet",
             "items": [
                 "殘差連接(skip connection)",
                 "解決梯度消失",
                 "可訓練 100+ 層",
             ]},
        ],
        bottom_note="考試口訣：VGG 加深 / Inception 加寬 / ResNet 跳接",
    )
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───── S14 · CONCEPT-CARD — R-CNN / RNN / Transformer ─────
    s = _blank(prs)
    add_title(s, "R-CNN / RNN / Transformer 任務配對")
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "R-CNN 系列",
             "items": [
                 "影像物件偵測",
                 "Region Proposal + CNN 分類",
                 "R-CNN → Fast → Faster",
             ]},
            {"heading": "RNN / LSTM",
             "items": [
                 "序列資料處理",
                 "RNN 有梯度消失問題",
                 "LSTM 用 gate 機制解決",
             ]},
            {"heading": "Transformer",
             "items": [
                 "Self-Attention 機制",
                 "可平行處理（比 RNN 快）",
                 "NLP + CV 主流架構",
             ]},
        ],
        bottom_note="影像分類→CNN / 物件偵測→R-CNN / 序列→RNN|Transformer / NLP→Transformer",
    )
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───── S15 · CONCEPT-CARD — 激活函數四大天王 ─────
    s = _blank(prs)
    add_title(s, "激活函數四大天王")
    draw_editorial_table(
        s,
        header=["函數", "公式 / 範圍", "使用場景", "注意"],
        rows=[
            ["ReLU", "max(0, x)\n[0, +inf)", "隱藏層首選", "Dead ReLU"],
            ["Sigmoid", "1/(1+e^-x)\n[0, 1]", "二元分類輸出層", "梯度消失"],
            ["Softmax", "多個機率加總=1", "多元分類輸出層", "——"],
            ["Tanh", "(e^x-e^-x)/(e^x+e^-x)\n[-1, 1]", "中心化版 Sigmoid", "梯度消失"],
        ],
        top=1.3,
        col_widths=[1.2, 2.5, 2.0, 1.3],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.5), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "必記：隱藏層 → ReLU / 二分類輸出 → Sigmoid / 多分類輸出 → Softmax",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───── S16 · CHECKPOINT — 深度學習架構與激活函數配對 ─────
    s = _blank(prs)
    add_title(s, "Check Point · 深度學習架構與激活函數配對")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 16 / 22 · 不看投影片作答",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(10.8), Inches(3.5),
        "Q1  CNN 影像分類（10 類）最後一層用什麼激活函數？\n\n"
        "Q2  二元情感分析最後一層用什麼？\n\n"
        "Q3  為什麼隱藏層不用 Sigmoid？\n\n"
        "Q4  哪個 CNN 架構的核心設計是「加寬不加深」？\n\n"
        "Q5  Transformer 為什麼能取代 RNN？",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.8), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.7),
        "Q1 Softmax · Q2 Sigmoid · Q3 梯度消失\n"
        "Q4 Inception · Q5 Self-Attention 可平行處理",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───── S17 · CONCEPT-CARD — 前向傳播：矩陣乘法 ─────
    s = _blank(prs)
    add_title(s, "前向傳播 = 矩陣乘法：Z = X W + b")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(0.6),
        "Z = X . W + b  →  A = activation(Z)",
        font_size=Pt(18), color=T.PRIMARY, bold=True,
        family=T.FONT_MONO,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(2.2), Inches(6.5), Inches(3.5),
        "維度規則：(batch, input) x (input, output) = (batch, output)\n\n"
        "手算範例：\n"
        "  X = [[1, 2, 3]]        # (1, 3)\n"
        "  W = [[0.1],[0.2],[0.3]] # (3, 1)\n"
        "  b = [0.1]\n"
        "  Z = 1x0.1 + 2x0.2 + 3x0.3 + 0.1 = 1.5\n"
        "  A = ReLU(1.5) = 1.5\n\n"
        "考點：維度對齊 + 加偏差 + 過激活函數",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.4,
        family=T.FONT_MONO,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(4.0),
        slot_name="前向傳播矩陣運算圖",
        description="輸入 X(1,3) x 權重 W(3,1) + 偏差 b → 過 ReLU 得到輸出。",
        size_hint="1200x900 px",
        placeholder_id="Ch09_S17_forward_pass",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───── S18 · CONCEPT-CARD — 優化器：SGD / Adam / Adagrad ─────
    s = _blank(prs)
    add_title(s, "優化器三兄弟：SGD / Adagrad / Adam")
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "SGD",
             "items": [
                 "基本款·學習率固定",
                 "容易卡 local minimum",
                 "需手動調 lr",
             ]},
            {"heading": "Adagrad",
             "items": [
                 "自動調整學習率",
                 "稀疏特徵友善",
                 "學習率持續衰減→後期停止",
             ]},
            {"heading": "Adam",
             "items": [
                 "Momentum + 自適應 lr",
                 "收斂快",
                 "最常用的預設選擇",
             ]},
        ],
        bottom_note="考試口訣：SGD 手動調 / Adagrad 會衰減 / Adam 最常用（樣題 Q14）",
    )
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ───── S19 · PRACTICE — Python 程式題模擬 ─────
    s = _blank(prs)
    add_title(s, "Practice · Python sklearn 程式題")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 19 / 22 · 科目3 約 25% 程式閱讀題",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.8),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.0),
        label="sklearn Random Forest 標準流程",
        code=(
            "from sklearn.ensemble import RandomForestClassifier\n"
            "from sklearn.model_selection import train_test_split\n"
            "from sklearn.metrics import accuracy_score\n"
            "\n"
            "X_train, X_test, y_train, y_test = train_test_split(\n"
            "    X, y, test_size=0.2, random_state=42)\n"
            "model = RandomForestClassifier(n_estimators=100)\n"
            "model.fit(X_train, y_train)\n"
            "y_pred = model.predict(X_test)\n"
            "print(accuracy_score(y_test, y_pred))"
        ),
        bullets=[
            "test_size=0.2 → 20% 測試集",
            "n_estimators=100 → 100 棵決策樹",
            "迴歸任務改用 RandomForestRegressor",
            "流程：切資料→建模→訓練→預測→評估",
        ],
    )
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───── S20 · CONCEPT-CARD — Overfitting 辨識與處理 ─────
    s = _blank(prs)
    add_title(s, "Overfitting 辨識與處理（跨章考點）")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
        "辨識：訓練集表現好 + 測試集表現差 = Overfitting",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "正則化",
             "items": [
                 "L1 (Lasso) → 稀疏",
                 "L2 (Ridge) → 小權重",
                 "Dropout → 隨機關閉神經元",
             ]},
            {"heading": "資料策略",
             "items": [
                 "更多訓練資料",
                 "Data Augmentation",
                 "Cross-Validation",
             ]},
            {"heading": "訓練策略",
             "items": [
                 "Early Stopping",
                 "簡化模型（減層/減特徵）",
                 "降低模型容量",
             ]},
        ],
    )
    _draw_bridge_note(s, "樣題 Q3 直接考 · Ch10 會用 Loss Curve 進一步診斷")
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ───── S21 · PYRAMID — 演算法選型三層思維 ─────
    s = _blank(prs)
    draw_pyramid_stack(
        s,
        title="Ch09 收束：演算法選型三層思維",
        layers=[
            {"name": "效能需求",
             "caption": "速度/精度/可解釋性 → 決定調參方向"},
            {"name": "資料特性",
             "caption": "維度/量級/分布形狀 → 決定具體演算法"},
            {"name": "問題類型",
             "caption": "監督/非監督 x 迴歸/分類/分群 → 決定大類"},
        ],
        thesis="從下往上選，不要反過來。",
    )
    add_footer(s, MODULE_CODE, 21, N_CONTENT)

    # ───── S22 · MOTIVATION — 銜接 Ch10 ─────
    s = _blank(prs)
    draw_ask_page(
        s,
        "模型做出來只是開始\n——調好、評好、管好才是考試重點。",
        data_card={
            "label": "下一章",
            "stat": "Ch10",
            "caption": "建模調參評估與 ML 治理\nL233 + L234",
        },
    )
    add_footer(s, MODULE_CODE, 22, N_CONTENT)

    # == Copyright ==
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
