"""Ch00 deck — AI 數學基礎：從零開始（完整版）
22 content slides + cover + copyright.

三大 Part：
  Part A — 機率統計（S1-S9）
  Part B — 線性代數（S10-S16）
  Part C — 最佳化（S17-S22）

受眾：沒有學過 AI 的人，高中數學為前提。
Aligned to chapters/Ch00_AI數學基礎_從零開始/02_slides_design.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source, add_rect,
    set_solid_fill, set_no_fill, set_line, set_no_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_vs_two_col, draw_pyramid_stack,
    draw_three_blocks_flow, draw_image_placeholder,
    draw_inverted_thesis_box,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch00"
MODULE_TITLE = "AI 數學基礎：從零開始"
MODULE_SUBTITLE = "機率統計 × 線性代數 × 最佳化——讀懂 AI 公式的眼鏡"
TIME_MIN = 170
N_CONTENT = 22


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _draw_bridge_note(slide, text, y_inches=6.6):
    add_textbox(
        slide, T.MARGIN_X, Inches(y_inches),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        text,
        font_size=T.FONT_SMALL, color=T.PRIMARY, bold=True, italic=True,
        align=PP_ALIGN.CENTER,
    )


def _draw_formula_card(slide, title_text, formula, explanation, ai_app,
                       image_placeholder=None, image_registry=None):
    """概念公式卡：左側公式+文字，右側可選圖片佔位。"""
    add_title(slide, title_text)
    left_w = Inches(6.5) if image_placeholder else T.SLIDE_W - 2 * T.MARGIN_X

    # Formula box
    add_textbox(
        slide, T.MARGIN_X, Inches(1.4), left_w, Inches(0.7),
        formula,
        font_size=Pt(18), color=T.PRIMARY, bold=True,
        family=T.FONT_MONO,
    )
    # Explanation
    add_textbox(
        slide, T.MARGIN_X, Inches(2.2), left_w, Inches(2.8),
        explanation,
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    # AI application
    ai_box = add_rect(slide, T.MARGIN_X, Inches(5.2), left_w, Inches(0.6))
    set_solid_fill(ai_box, T.PRIMARY)
    set_no_line(ai_box)
    add_textbox(
        slide, T.MARGIN_X + Inches(0.2), Inches(5.2),
        left_w - Inches(0.4), Inches(0.6),
        f"AI 應用：{ai_app}",
        font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def _draw_pitfall_two_col(slide, title_text, left_items, right_items,
                          bottom_note=""):
    """PITFALL 型：2x2 表格或對比式。"""
    add_title(slide, title_text)
    col_w = Inches(5.5)
    gap = Inches(0.5)
    total = col_w * 2 + gap
    left_x = (T.SLIDE_W - total) / 2
    right_x = left_x + col_w + gap
    top = Inches(1.5)
    col_h = Inches(4.3)

    for x, items in [(left_x, left_items), (right_x, right_items)]:
        body = add_rect(slide, x, top, col_w, col_h)
        set_no_fill(body)
        set_line(body, T.PRIMARY, 1.0)
        item_text = "\n".join(items)
        add_textbox(
            slide, x + Inches(0.25), top + Inches(0.2),
            col_w - Inches(0.5), col_h - Inches(0.4),
            item_text,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            line_spacing=1.45,
        )

    if bottom_note:
        add_textbox(
            slide, T.MARGIN_X, top + col_h + Inches(0.15),
            T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
            bottom_note,
            font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
            align=PP_ALIGN.CENTER,
        )


def build_ch00(output_path, image_registry=None):
    """Build Ch00 deck; 22 content slides + cover + copyright."""
    prs = _new_prs()

    # ── Cover ──
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───── S1 · MOTIVATION — 你不需要成為數學家 ─────
    s = _blank(prs)
    draw_ask_page(
        s,
        "你不需要成為數學家，\n但你需要看懂 AI 用數學在做什麼。",
        data_card={
            "label": "iPAS 中級考試",
            "stat": "~30%",
            "caption": "公式題佔比\n但只需高中數學程度",
        },
    )
    add_source(s, "iPAS 歷屆考古題分析 · 高中數學起步")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # ───── S2 · SILENT — 三把數學工具 ─────
    s = _blank(prs)
    draw_silent_page(
        s,
        "AI 用機率衡量不確定性\n用線性代數搬運資料\n用最佳化找最佳答案",
    )
    add_footer(s, MODULE_CODE, 2, N_CONTENT, dark_bg=True)

    # ───── S3 · CONCEPT-CARD — 隨機變數：離散 vs 連續 ─────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="離散隨機變數",
        right_title="連續隨機變數",
        left_items=[
            "值可列舉（骰子 1-6、硬幣正反）",
            "P(X = x)，所有加總 = 1",
            "長條圖（每個值有明確機率）",
            "AI 應用：分類模型的輸出 = 離散",
        ],
        right_items=[
            "值不可列舉（身高 170.32...）",
            "用密度函數 f(x)",
            "P(a < X < b) = 曲線下面積",
            "AI 應用：迴歸模型的輸出 = 連續",
        ],
        title="隨機變數：離散 vs 連續",
        summary="分類 → 離散 / 迴歸 → 連續。這個區分貫穿整個 AI。",
    )
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───── S4 · MECHANISM-FLOW — 條件機率 P(A|B) ─────
    s = _blank(prs)
    add_title(s, "條件機率 P(A|B)——方向很重要")
    # Formula
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(0.7),
        "P(A|B) = P(A ∩ B) / P(B)",
        font_size=Pt(22), color=T.PRIMARY, bold=True,
        family=T.FONT_MONO,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(2.2), Inches(6.5), Inches(2.5),
        "「已知 B 發生了，A 也發生的機率是多少？」\n\n"
        "• B 已經發生 → 你的世界縮小到 B 這個圓圈\n"
        "• A 在 B 裡面佔多少比例 = P(A|B)\n\n"
        "例：P(有病|陽性) ≠ P(陽性|有病)\n"
        "    方向反過來是完全不同的數字",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    # Warning box
    warn_y = Inches(5.0)
    warn_box = add_rect(slide=s, x=T.MARGIN_X, y=warn_y,
                        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(0.6))
    set_solid_fill(warn_box, T.PRIMARY)
    set_no_line(warn_box)
    add_textbox(
        s, T.MARGIN_X + Inches(0.3), warn_y,
        T.SLIDE_W - 2 * T.MARGIN_X - Inches(0.6), Inches(0.6),
        "考試陷阱：P(A|B) ≠ P(B|A)！方向不能反。",
        font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(3.5),
        slot_name="Venn 圖：條件機率",
        description="兩個重疊圓圈 A 和 B，重疊部分標示 A∩B，箭頭指向 P(A|B) = 重疊面積 / B 的面積。",
        size_hint="1200×800 px",
        placeholder_id="Ch00_S4_venn_cond_prob",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───── S5 · EXAMPLE-I/O — 貝氏定理 + 醫療篩檢 ─────
    s = _blank(prs)
    add_title(s, "貝氏定理：為什麼 COVID 要做二次篩檢")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(12), Inches(0.6),
        "P(A|B) = P(B|A) × P(A) / P(B)",
        font_size=Pt(18), color=T.PRIMARY, bold=True,
        family=T.FONT_MONO,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(2.1), Inches(7.0), Inches(3.5),
        "設定：盛行率 1%、靈敏度 99%、特異度 95%\n\n"
        "1000 人起始：\n"
        "  → 10 有病 → 9.9 陽性 / 0.1 陰性\n"
        "  → 990 沒病 → 49.5 陽性 / 940.5 陰性\n\n"
        "總陽性 = 9.9 + 49.5 = 59.4 人\n"
        "P(有病|陽性) = 9.9 / 59.4 ≈ 17%",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    # Conclusion box
    conc_y = Inches(5.8)
    conc = add_rect(s, T.MARGIN_X, conc_y,
                    T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6))
    set_solid_fill(conc, T.PRIMARY)
    set_no_line(conc)
    add_textbox(
        s, T.MARGIN_X + Inches(0.3), conc_y,
        T.SLIDE_W - 2 * T.MARGIN_X - Inches(0.6), Inches(0.6),
        "檢測很準(99%) 但陽性結果只有 17% 是真的——因為盛行率太低",
        font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    draw_image_placeholder(
        s, Inches(8.0), Inches(1.8), Inches(4.8), Inches(3.8),
        slot_name="貝氏定理樹狀圖",
        description="樹狀圖：1000 人分支到有病/沒病，再分支到陽性/陰性，標示各數字。",
        size_hint="1200×900 px",
        placeholder_id="Ch00_S5_bayes_tree",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───── S6 · CONCEPT-CARD — E(X) 與 Var(X) ─────
    s = _blank(prs)
    add_title(s, "期望值 E(X) 與變異數 Var(X)")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(3.5),
        "E(X) = Σ xᵢ × P(xᵢ) = 加權平均\n"
        "  例：骰子 E(X) = (1+2+3+4+5+6)/6 = 3.5\n"
        "  AI 應用：損失函數 = 最小化預測誤差的期望值\n\n"
        "Var(X) = E[(X - μ)²] = 離散程度\n"
        "  AI 應用：模型 variance 高 = overfitting\n\n"
        "直覺：E(X) = 中心位置 / Var(X) = 寬度",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(3.8),
        slot_name="兩個常態分佈比較",
        description="兩個常態分佈重疊：同均值不同寬度。窄的標示 Var 小=穩定，寬的標示 Var 大=不穩定。",
        size_hint="1200×900 px",
        placeholder_id="Ch00_S6_var_comparison",
        registry=image_registry,
    )
    _draw_bridge_note(s, "E(X) = 中心 / Var(X) = 寬度——這兩個概念會出現在後面每一章")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───── S7 · CONCEPT-CARD — 中央極限定理 CLT ─────
    s = _blank(prs)
    add_title(s, "中央極限定理 CLT：為什麼假設檢定成立")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(3.8),
        "• 母體什麼分佈不重要\n\n"
        "• 只要 n ≥ 30，抽樣平均數的分佈 → 常態分佈\n\n"
        "• 這就是假設檢定能用常態分佈的原因\n\n"
        "• 考試不考 CLT 推導\n"
        "  只考「為什麼假設檢定可以假設常態分佈？」\n"
        "  → 答案：中央極限定理",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(4.5),
        slot_name="CLT 視覺化",
        description="三行：均勻/指數/雙峰分佈 → 各自的抽樣均值分佈皆趨近常態。",
        size_hint="1200×1000 px",
        placeholder_id="Ch00_S7_clt_visual",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───── S8 · PITFALL — Type I/II 錯誤 ─────
    s = _blank(prs)
    add_title(s, "Type I / Type II 錯誤 + α + p-value")
    draw_editorial_table(
        s,
        header=["", "H0 為真（沒事）", "H0 為假（有事）"],
        rows=[
            ["拒絕 H0", "Type I (α)\n喊狼來了", "正確拒絕\n抓到壞人"],
            ["不拒絕 H0", "正確\n平安無事", "Type II (β)\n放走壞人"],
        ],
        top=1.3,
        col_widths=[1.5, 2.5, 2.5],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(3.8), T.SLIDE_W - 2 * T.MARGIN_X, Inches(2.0),
        "• Type I = 喊「狼來了」但沒有狼（假陽性）\n"
        "• Type II = 狼來了但你沒喊（假陰性）\n"
        "• α = 容忍 Type I 的上限（通常 0.05）\n"
        "• p-value < α → 拒絕 H0\n"
        "• 注意：不拒絕 ≠ 接受——只是沒有足夠證據拒絕",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    _draw_bridge_note(s, "→ Ch08 會用假設檢定做完整的統計分析實戰")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───── S9 · CHECKPOINT — Part A 小測驗 ─────
    s = _blank(prs)
    add_title(s, "Check Point · Part A 機率統計小測驗")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 9 / 22 · 不看投影片作答",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(10.8), Inches(4.0),
        "Q1  P(A|B) = 0.6，P(B|A) 一定也是 0.6 嗎？\n\n"
        "Q2  某疾病盛行率 0.1%，篩檢陽性結果可靠嗎？\n\n"
        "Q3  p-value = 0.03，α = 0.05，你的結論？",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.0), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "答案：Q1 否（方向不同）· Q2 不可靠（盛行率太低→假陽性多）· Q3 p<α→拒絕 H0",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───── S10 · CONCEPT-CARD — 向量 = 特徵清單 ─────
    s = _blank(prs)
    add_title(s, "向量 = 特徵清單：每筆資料都是一個向量")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(3.5),
        "x = [x₁, x₂, ..., x_d]\n\n"
        "一個樣本的 d 個特徵 = d 維向量\n\n"
        "例：房屋 = [坪數, 樓層, 屋齡] = [30, 5, 10]\n\n"
        "AI 應用：ML 把每一筆資料變成向量",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(4.0),
        slot_name="向量 2D 示意圖",
        description="左側試算表一行資料。右側 2D 座標平面，向量 [3,4] 畫成從原點出發的箭頭。",
        size_hint="1200×900 px",
        placeholder_id="Ch00_S10_vector_visual",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───── S11 · MECHANISM-FLOW — 內積 = 相似度 ─────
    s = _blank(prs)
    add_title(s, "內積 = 相似度：a · b = |a||b| cos θ")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(0.6),
        "a · b = Σ aᵢ × bᵢ = |a| |b| cos θ",
        font_size=Pt(18), color=T.PRIMARY, bold=True,
        family=T.FONT_MONO,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(2.2), Inches(6.5), Inches(3.0),
        "• cos θ ≈ 1 → 方向相同 → 很相似\n"
        "• cos θ ≈ 0 → 垂直 → 不相關\n"
        "• cos θ ≈ -1 → 方向相反 → 很不像\n\n"
        "手算：[1, 2, 3] · [4, 5, 6]\n"
        "     = 4 + 10 + 18 = 32\n\n"
        "AI 應用：線性回歸 ŷ = w · x + b = 內積",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(3.8),
        slot_name="內積與夾角示意圖",
        description="2D 視覺：兩個向量 a 和 b，標示夾角 theta，虛線標示投影。",
        size_hint="1200×900 px",
        placeholder_id="Ch00_S11_dot_product",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───── S12 · CONCEPT-CARD — L2 範數 ─────
    s = _blank(prs)
    add_title(s, "L2 範數 = 向量長度（畢氏定理推廣）")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(0.6),
        "||x||₂ = √(x₁² + x₂² + ... + x_d²)",
        font_size=Pt(18), color=T.PRIMARY, bold=True,
        family=T.FONT_MONO,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(2.2), Inches(6.5), Inches(2.5),
        "例：||[3, 4]||₂ = √(9 + 16) = √25 = 5\n\n"
        "就是畢氏定理的高維推廣\n\n"
        "→ 連結 Part C：\n"
        "  L2 正則化 penalty = ||w||₂²\n"
        "  = 限制權重的長度不能太大",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(3.5),
        slot_name="L2 範數視覺：3-4-5 三角形",
        description="2D 視覺：向量 [3,4] 標示長度=5，畢氏定理直角三角形。",
        size_hint="1200×800 px",
        placeholder_id="Ch00_S12_l2_norm",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───── S13 · CONCEPT-CARD — 矩陣 = 一批向量 ─────
    s = _blank(prs)
    add_title(s, "矩陣 = 一批向量：試算表就是矩陣")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(3.5),
        "n 筆資料 × d 個特徵 = shape(n, d) 矩陣\n\n"
        "• 每行 = 一個樣本向量\n"
        "• 每列 = 一個特徵維度\n\n"
        "Python：X.shape = (100, 5)\n"
        "  → 100 筆資料、5 個特徵\n\n"
        "記住：第一個數字是幾筆(n)，第二個是幾個特徵(d)",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(4.0),
        slot_name="Excel → 矩陣 對應圖",
        description="左側 Excel 試算表風格（5 行 3 列，行=樣本、列=特徵）。右側對應矩陣符號 X(5,3)。",
        size_hint="1200×900 px",
        placeholder_id="Ch00_S13_matrix_excel",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───── S14 · PRACTICE — 矩陣乘法手算 ─────
    s = _blank(prs)
    add_title(s, "Practice · 矩陣乘法手算")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 14 / 22 · 跟著算一遍",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(1.8), Inches(7.0), Inches(0.5),
        "維度規則：(m, n) × (n, p) = (m, p)    中間的 n 必須一樣！",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, Inches(1.0), Inches(2.5), Inches(11.0), Inches(3.5),
        "A = [[1,2,3],[4,5,6]]  shape(2,3)\n"
        "B = [[7,8],[9,10],[11,12]]  shape(3,2)\n\n"
        "C = A × B = shape(2,2)\n"
        "  C[0,0] = 1×7 + 2×9 + 3×11 = 58\n"
        "  C[0,1] = 1×8 + 2×10 + 3×12 = 64\n"
        "  C[1,0] = 4×7 + 5×9 + 6×11 = 139\n"
        "  C[1,1] = 4×8 + 5×10 + 6×12 = 154",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        family=T.FONT_MONO,
        line_spacing=1.4,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "考試天花板：只考 2×3 × 3×2，不會考更大的。",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───── S15 · MECHANISM-FLOW — 前向傳播 Wx + b ─────
    s = _blank(prs)
    add_title(s, "前向傳播 = 矩陣乘法：z = Wx + b")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(0.6),
        "z = W × x + b → a = activation(z)",
        font_size=Pt(18), color=T.PRIMARY, bold=True,
        family=T.FONT_MONO,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(2.2), Inches(6.5), Inches(3.0),
        "維度：輸入(batch, d_in) × 權重(d_in, d_out) = 輸出(batch, d_out)\n\n"
        "例：3 個輸入特徵 → 2 個輸出神經元\n"
        "  W 的 shape = (3, 2)\n"
        "  x 的 shape = (1, 3)\n\n"
        "→ Ch09 S17 會完整手算",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(4.0),
        slot_name="神經網路單層前向傳播",
        description="左邊輸入節點 x(3 維)，中間權重矩陣 W(3,2)，右邊輸出 z(2 維)，加偏差 b(2 維)。",
        size_hint="1200×900 px",
        placeholder_id="Ch00_S15_forward_pass",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───── S16 · CONCEPT-CARD — 特徵值/特徵向量 = PCA ─────
    s = _blank(prs)
    add_title(s, "特徵值/特徵向量 = PCA 降維的原理")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(0.6),
        "Av = λv",
        font_size=Pt(22), color=T.PRIMARY, bold=True,
        family=T.FONT_MONO,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(2.2), Inches(6.5), Inches(3.3),
        "矩陣 A 作用在特徵向量 v\n"
        "  → 只改變長度(λ)，不改變方向\n\n"
        "• λ = 拉伸倍率 = 特徵值\n"
        "• v = 被拉伸的方向 = 特徵向量\n\n"
        "PCA = 找協方差矩陣最大特徵值的方向\n"
        "    = 資料變異最大的方向\n\n"
        "考試不考計算，只考概念",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(4.0),
        slot_name="PCA 橢圓資料雲",
        description="2D 橢圓形資料雲，兩個箭頭標示第一主成分(長軸)和第二主成分(短軸)。",
        size_hint="1200×900 px",
        placeholder_id="Ch00_S16_pca_ellipse",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───── S17 · CONCEPT-CARD — MSE vs Cross-Entropy ─────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="MSE → 迴歸",
        right_title="Cross-Entropy → 分類",
        left_items=[
            "MSE = (1/n) Σ(y - ŷ)²",
            "直覺：預測偏差的平方平均",
            "預測房價差 20 萬 → 懲罰 400 萬",
            "差越多懲罰越重",
        ],
        right_items=[
            "CE = -Σ yᵢ log(ŷᵢ)",
            "直覺：預測機率偏離真實標籤",
            "log 函數放大錯誤",
            "機率越偏離，懲罰越重",
        ],
        title="損失函數：MSE vs Cross-Entropy",
        summary="口訣：迴歸用 MSE，分類用 Cross-Entropy。考試直接考。",
    )
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───── S18 · MECHANISM-FLOW — 導數/偏微分 → 梯度 ─────
    s = _blank(prs)
    add_title(s, "導數 → 偏微分 → 梯度：最陡下坡方向")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(3.8),
        "導數 = 某點的斜率 = f'(x)（高中複習）\n\n"
        "偏微分 = 只看一個方向的斜率\n"
        "  ∂L/∂w₁ = 固定其他變數，只看 w₁ 的變化率\n\n"
        "梯度 = 所有偏微分組成的向量\n"
        "  ∇L = [∂L/∂w₁, ∂L/∂w₂, ...]\n\n"
        "• 梯度指向最陡上坡方向\n"
        "• 沿負梯度走 = 最陡下坡 = 梯度下降",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(4.5),
        slot_name="梯度視覺化",
        description="左側 2D 曲線 f(x)=x² 標示斜率線。右側 3D 曲面標示兩方向偏微分。底部梯度向量箭頭。",
        size_hint="1200×1000 px",
        placeholder_id="Ch00_S18_gradient_visual",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ───── S19 · PITFALL — Learning Rate 三面板 ─────
    s = _blank(prs)
    add_title(s, "Learning Rate：太大震盪、太小龜速、剛好收斂")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "更新規則：w = w - lr × ∇L",
        font_size=Pt(16), color=T.PRIMARY, bold=True,
        family=T.FONT_MONO,
    )
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "lr 太大",
             "items": [
                 "loss 曲線上下震盪",
                 "跳過最低點",
                 "不收斂",
             ]},
            {"heading": "lr 太小",
             "items": [
                 "loss 曲線緩慢下降",
                 "要訓練很久",
                 "可能卡在局部最小",
             ]},
            {"heading": "lr 剛好",
             "items": [
                 "loss 穩定下降",
                 "到達最低點",
                 "收斂！",
             ]},
        ],
    )
    _draw_bridge_note(s, "→ Ch09：Adam / Adagrad 就是在自動調 lr")
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───── S20 · PRACTICE — 鏈式法則 + 凸性 + L1 vs L2 ─────
    s = _blank(prs)
    add_title(s, "鏈式法則 · 凸性 · L1 vs L2 正則化")
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "鏈式法則",
             "items": [
                 "dL/dw = dL/dz × dz/dw",
                 "反向傳播的數學基礎",
                 "考試只考概念",
             ]},
            {"heading": "凸 vs 非凸",
             "items": [
                 "碗形(凸) = 全域最佳",
                 "崎嶇(非凸) = 局部最小",
                 "線性回歸=凸 / NN=非凸",
             ]},
            {"heading": "L1 vs L2 正則化",
             "items": [
                 "L1(Lasso) → 鑽石→稀疏",
                 "L2(Ridge) → 圓形→小權重",
                 "考試直接考差異",
             ]},
        ],
    )
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ───── S21 · PYRAMID — 三層數學能力金字塔 ─────
    s = _blank(prs)
    draw_pyramid_stack(
        s,
        title="Ch00 收束：三層數學能力金字塔",
        layers=[
            {"name": "最佳化",
             "caption": "找到最佳答案 → Ch10 調參評估"},
            {"name": "線性代數",
             "caption": "搬運與組合資料 → Ch09 前向傳播"},
            {"name": "機率統計",
             "caption": "衡量不確定性 → Ch08 假設檢定"},
        ],
        thesis="三層疊起來 = AI 的數學引擎。",
    )
    add_footer(s, MODULE_CODE, 21, N_CONTENT)

    # ───── S22 · SILENT — 數學工具已裝備 ─────
    s = _blank(prs)
    draw_silent_page(
        s,
        "你已經擁有讀懂 AI 數學的眼鏡。\n下一站：Ch08 統計 · Ch09 建模 · Ch10 調參",
    )
    add_footer(s, MODULE_CODE, 22, N_CONTENT, dark_bg=True)

    # ── Copyright ──
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
