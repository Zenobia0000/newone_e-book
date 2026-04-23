"""Ch07 deck — AI 導入規劃與系統部署（完整版）
20 content slides + cover + copyright.

三大段落：
  起（S1-S4）  — 定錨 + 全覽 + 生命週期
  承（S5-S10） — 需求分析 / 方案選擇 / 模型比較 / 數據準備
  轉（S11-S17）— 風險管理 / 資安 / 部署架構 / 測試四層
  合（S18-S20）— 實戰練習 / 金字塔 / 結語

評鑑範圍：中級 L212 + L213
Aligned to chapters/Ch07_AI導入規劃與系統部署/02_slides_design.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source, add_rect,
    set_solid_fill, set_no_fill, set_line, set_no_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_vs_two_col, draw_pyramid_stack,
    draw_three_blocks_flow, draw_image_placeholder,
    draw_inverted_thesis_box, draw_grid,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch07"
MODULE_TITLE = "AI 導入規劃與系統部署"
MODULE_SUBTITLE = "五階段生命週期 × 風險四象限 × 測試四層——知道坑在哪才能繞過去"
TIME_MIN = 120
N_CONTENT = 20


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_ch07(output_path, image_registry=None):
    """Build Ch07 deck; 20 content slides + cover + copyright."""
    prs = _new_prs()

    # ── Cover ──
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───── S1 · SILENT · 起 ─────
    s = _blank(prs)
    draw_silent_page(
        s,
        "從選模型到上線，中間有一百個坑。\n考試考的是你知不知道坑在哪。",
    )
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ───── S2 · ASK · 起 ─────
    s = _blank(prs)
    draw_ask_page(
        s,
        "你公司的 AI 專案，卡在哪一步？",
        data_card={
            "label": "VentureBeat 2024",
            "stat": "87%",
            "caption": "AI 專案未能從概念驗證\n進入生產部署",
        },
    )
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───── S3 · MATRIX · 起 ─────
    s = _blank(prs)
    draw_matrix(
        s, rows=2, cols=3,
        cells=[
            {"text": "L212 導入評估",
             "sub": "需求分析 / 方案設計\nQ5 AutoML 趨勢", "highlight": True},
            {"text": "L212 規劃管理",
             "sub": "可行性評估 / ROI\n優先級排序", "highlight": False},
            {"text": "L212 風險管理",
             "sub": "四象限 / 資安三威脅\nQ9 供應鏈攻擊", "highlight": False},
            {"text": "L213 數據準備",
             "sub": "EDA / 清洗 / 特徵工程\n佔專案 60-80% 時間", "highlight": False},
            {"text": "L213 模型選擇",
             "sub": "多維度比較\n精度/速度/可解釋性", "highlight": False},
            {"text": "L213 系統部署",
             "sub": "雲端/地端/Edge\nQ15 整合測試", "highlight": True},
        ],
        title="L212 + L213 全覽：科目1 的重頭戲",
    )
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───── S4 · MECHANISM_FLOW · 起 ─────
    s = _blank(prs)
    draw_flow_chain(
        s,
        title="AI 導入五階段生命週期",
        nodes=[
            {"label": "需求分析", "sub": "痛點量化\nROI 評估", "highlight": False},
            {"label": "方案設計", "sub": "自建/採購/AutoML", "highlight": False},
            {"label": "數據準備", "sub": "EDA / 清洗\n特徵工程", "highlight": False},
            {"label": "模型建構", "sub": "選型 / 訓練\n驗證", "highlight": False},
            {"label": "系統部署", "sub": "容器化 / CI/CD\n監控", "highlight": True},
        ],
        y=2.8,
        branch={"from_index": 4, "label": "監控更新", "sub": "模型漂移 → 迴圈", "above": True},
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.5), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "AI 導入不是線性的——是有回饋迴圈的流程。考試用情境題測你知不知道問題屬於哪個階段。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───── S5 · TABLE · 承 ─────
    s = _blank(prs)
    draw_editorial_table(
        s,
        title="需求分析四步法",
        header=["步驟", "做什麼", "怎麼做", "考試怎麼考"],
        rows=[
            ["痛點識別", "量化問題", "「人工目檢錯誤率 15%」", "情境 → 問題定義"],
            ["可行性評估", "資料/技術/團隊", "三個都要過關", "「最不適合導入 AI」"],
            ["ROI 估算", "成本 vs 效益", "導入成本 + 維護成本", "「投資報酬率最高」"],
            ["優先級排序", "影響力×可行性", "矩陣排序", "「應先導入哪個」"],
        ],
        top=1.3,
        col_widths=[1.2, 1.5, 2.0, 2.0],
    )
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───── S6 · VS · 承 ─────
    s = _blank(prs)
    draw_three_blocks_flow(
        s,
        title="方案三條路：自建 vs 採購 vs AutoML（Q5）",
        blocks=[
            {"heading": "自建模型",
             "items": [
                 "客製化高",
                 "成本高、需專業團隊",
                 "適合核心競爭力場景",
             ]},
            {"heading": "採購方案",
             "items": [
                 "快速導入、成本可控",
                 "客製化受限",
                 "適合標準化需求",
             ]},
            {"heading": "AutoML",
             "items": [
                 "降低建模門檻",
                 "自動化特徵/模型/超參",
                 "問題定義仍需人工",
             ]},
        ],
        bottom_note="Q5：AutoML 降低了建模門檻，但沒有降低規劃門檻。",
    )
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───── S7 · PITFALL · 承 ─────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        title="模型選擇 ≠ 只看 Accuracy",
        left_title="錯：只看 accuracy",
        right_title="對：多維度比較",
        left_items=[
            "模型 A accuracy=95% → 選 A",
            "忽略推論時間 10 秒",
            "忽略模型大小 2GB",
            "忽略可解釋性為零",
        ],
        right_items=[
            "同時比較 accuracy / 速度",
            "考慮 inference time / cost",
            "考慮 interpretability",
            "場景需即時 → 速度優先",
        ],
        summary="準確率是必要條件，不是充分條件。",
    )
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───── S8 · TABLE · 承 ─────
    s = _blank(prs)
    draw_editorial_table(
        s,
        title="模型選型速查表",
        header=["模型", "精度", "速度", "可解釋性", "適用場景"],
        rows=[
            ["線性迴歸", "★★", "★★★★★", "★★★★★", "簡單預測/基線"],
            ["決策樹", "★★★", "★★★★", "★★★★★", "小資料/需解釋"],
            ["隨機森林", "★★★★", "★★★", "★★★", "結構化資料"],
            ["XGBoost", "★★★★★", "★★★", "★★★", "競賽/高精度"],
            ["SVM", "★★★", "★★★", "★★", "小樣本分類"],
            ["CNN", "★★★★★", "★★", "★★", "影像辨識"],
            ["Transformer", "★★★★★", "★★", "★", "文本/序列"],
        ],
        top=1.3,
        col_widths=[1.5, 1.0, 1.0, 1.0, 2.0],
    )
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───── S9 · MECHANISM_FLOW · 承 ─────
    s = _blank(prs)
    draw_flow_chain(
        s,
        title="數據準備流程（佔專案 60-80% 時間）",
        nodes=[
            {"label": "數據收集", "sub": "內部/外部\n公開資料集"},
            {"label": "EDA", "sub": "探索性分析\n看資料長相"},
            {"label": "清洗處理", "sub": "缺值/異常/去重\n格式統一"},
            {"label": "特徵工程", "sub": "衍生特徵\n編碼/標準化"},
            {"label": "特徵選擇", "sub": "降維/篩選\n去除冗餘"},
        ],
        y=2.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.0), T.SLIDE_W - 2 * T.MARGIN_X, Inches(1.2),
        "清洗重點：MCAR 可直接刪、MAR 用模型填補、MNAR 需額外調查\n"
        "標準化 vs 正規化：標準化(Z-score)去量綱、正規化(Min-Max)壓到[0,1]",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───── S10 · PITFALL · 承 ─────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        title="垃圾進垃圾出——不做 EDA 就是盲人開車",
        left_title="錯：拿到資料直接訓練",
        right_title="對：先做 EDA 再訓練",
        left_items=[
            "10 萬筆資料直接訓練",
            "Accuracy 只有 60%",
            "不知道 30% 欄位有缺值",
            "不知道標籤 95:5 不平衡",
        ],
        right_items=[
            "EDA 發現缺值 + 不平衡",
            "處理缺值 + 重新抽樣",
            "去除 5% 重複資料",
            "同模型 accuracy → 85%",
        ],
        summary="考試問「下一步該做什麼」——答案通常是 EDA。",
    )
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───── S11 · SILENT · 轉 ─────
    s = _blank(prs)
    draw_silent_page(
        s,
        "規劃做完了，數據準備好了，\n接下來是兩個考試重災區：\n風險管理 + 系統部署。",
    )
    add_footer(s, MODULE_CODE, 11, N_CONTENT, dark_bg=True)

    # ───── S12 · MATRIX · 轉 ─────
    s = _blank(prs)
    draw_matrix(
        s, rows=2, cols=2,
        cells=[
            {"text": "低機率 × 高影響",
             "sub": "法規風險\nGDPR 違規 / AI 法案\n罰款巨大但發生率低", "highlight": False},
            {"text": "高機率 × 高影響",
             "sub": "資安 / 倫理風險\n資料外洩 / 偏見歧視\n出事上新聞", "highlight": True},
            {"text": "低機率 × 低影響",
             "sub": "環境風險\n硬體故障 / 斷網\n可快速修復", "highlight": False},
            {"text": "高機率 × 中影響",
             "sub": "技術風險\n模型漂移 / 過擬合\n可透過重訓修復", "highlight": False},
        ],
        title="風險四象限：按優先級分配資源",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.0), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "考試：「以下哪個風險最應該優先處理？」→ 高影響的優先，同影響看機率。",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───── S13 · TABLE · 轉 ─────
    s = _blank(prs)
    draw_editorial_table(
        s,
        title="資安三大威脅（Q9）",
        header=["威脅", "攻擊方式", "防範措施", "考試出法"],
        rows=[
            ["供應鏈攻擊(Q9)", "在第三方元件埋雷\n預訓練模型/套件",
             "SBOM / 可信任儲存庫\n安全掃描", "「如何防範？」"],
            ["對抗攻擊", "輸入加微小擾動\n騙過模型判斷",
             "對抗訓練\n輸入驗證", "「哪種攻擊？」"],
            ["資料投毒", "訓練資料混入\n惡意樣本",
             "資料審核\n異常偵測", "「後果是什麼？」"],
        ],
        top=1.3,
        col_widths=[1.5, 2.0, 2.0, 1.5],
    )
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───── S14 · VS · 轉 ─────
    s = _blank(prs)
    draw_three_blocks_flow(
        s,
        title="部署架構：雲端 vs 地端 vs Edge",
        blocks=[
            {"heading": "雲端部署",
             "items": [
                 "彈性擴縮 / 按量付費",
                 "資料需上傳",
                 "適合快速上線 MVP",
             ]},
            {"heading": "地端部署",
             "items": [
                 "低延遲 / 資料不出場",
                 "維護成本高",
                 "適合醫療/金融",
             ]},
            {"heading": "Edge 部署",
             "items": [
                 "超低延遲 / 離線運行",
                 "模型大小受限",
                 "適合工廠產線",
             ]},
        ],
        bottom_note="決策樹：資料敏感？→ 地端。需離線？→ Edge。其他 → 雲端。",
    )
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───── S15 · MECHANISM_FLOW · 轉 ─────
    s = _blank(prs)
    draw_flow_chain(
        s,
        title="MLOps：模型部署工程管線",
        nodes=[
            {"label": "開發環境", "sub": "Jupyter\nVSCode"},
            {"label": "Docker", "sub": "容器封裝\n環境一致性"},
            {"label": "K8s", "sub": "編排調度\n自動擴縮"},
            {"label": "CI/CD", "sub": "自動測試\n自動部署"},
            {"label": "生產環境", "sub": "監控告警\n模型版本"},
        ],
        y=2.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.0), T.SLIDE_W - 2 * T.MARGIN_X, Inches(1.0),
        "Docker = 環境一致性 + 可移植性 / K8s = 自動擴縮容 / CI/CD = 程式碼推送後自動部署\n"
        "考試：「容器化的主要優勢？」→ 環境一致性和可移植性。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───── S16 · TABLE · 轉 ─────
    s = _blank(prs)
    draw_editorial_table(
        s,
        title="系統測試四層（Q15）",
        header=["測試類型", "測什麼", "何時做", "常見考法"],
        rows=[
            ["單元測試", "單一模組功能", "整合前", "「測試範圍最小？」"],
            ["整合測試(Q15)", "子系統介面\n資料流正確性", "單元測試後", "「最優先驗什麼？」\n→ 介面正確性"],
            ["壓力測試", "高負載效能\n回應時間", "整合測試後", "「何時做壓力測試？」"],
            ["A/B 測試", "新舊版本比較\n使用者體驗", "上線時", "「灰度發布用什麼？」"],
        ],
        top=1.3,
        col_widths=[1.5, 2.0, 1.5, 2.0],
    )
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───── S17 · PITFALL · 轉 ─────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        title="子系統各自 OK ≠ 合起來 OK",
        left_title="錯：單元測試通過就上線",
        right_title="對：完整測試流程",
        left_items=[
            "每個子系統各自通過測試",
            "直接上線",
            "使用者回報 JSON 格式不一致",
            "推論服務 OK 但串前端爆了",
        ],
        right_items=[
            "子系統各自通過單元測試",
            "→ 整合測試驗證介面",
            "→ 壓力測試驗效能",
            "→ 灰度發布 → 全量上線",
        ],
        summary="整合測試存在的理由：測「串起來」有沒有問題。",
    )
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───── S18 · PRACTICE · 合 ─────
    s = _blank(prs)
    add_title(s, "Practice · 製造業 AI 瑕疵檢測導入規劃")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 18 / 20 · 五分鐘情境練習",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    draw_grid(
        s, rows=1, cols=5,
        cells=[
            {"label": "需求分析", "sub": "人工目檢痛點？"},
            {"label": "方案選擇", "sub": "CNN + Edge？\n雲端 API？"},
            {"label": "風險評估", "sub": "模型漂移\n對抗攻擊"},
            {"label": "部署架構", "sub": "Edge\n(即時+離線)"},
            {"label": "測試計畫", "sub": "整合 → 壓力\n→ A/B"},
        ],
        title="",
        top=2.2,
        bottom=2.0,
        caption="用本章框架解一道完整情境題——包含 L212 + L213 所有考點。",
    )
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ───── S19 · PYRAMID · 合 ─────
    s = _blank(prs)
    draw_pyramid_stack(
        s,
        title="Ch07 收束：五層金字塔",
        layers=[
            {"name": "系統部署",
             "caption": "容器化 / CI/CD / 測試四層"},
            {"name": "風險管理",
             "caption": "四象限 / 資安三威脅"},
            {"name": "模型選擇",
             "caption": "多維度比較 / 選型速查"},
            {"name": "數據準備",
             "caption": "EDA / 清洗 / 特徵工程"},
            {"name": "需求分析",
             "caption": "痛點量化 / ROI / 優先級"},
        ],
        thesis="底層沒做好，上面全部白搭——這是 Ch07 的核心邏輯。",
    )
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───── S20 · SILENT · 合 ─────
    s = _blank(prs)
    draw_silent_page(
        s,
        "知道坑在哪，才能繞過去——\n這就是規劃師的價值。\n下一章 Ch08：統計和大數據。",
    )
    add_footer(s, MODULE_CODE, 20, N_CONTENT, dark_bg=True)

    # ── Copyright ──
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
