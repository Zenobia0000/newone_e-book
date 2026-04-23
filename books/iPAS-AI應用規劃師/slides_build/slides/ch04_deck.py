"""Ch04 deck — 生成式 AI 工具、提示工程與 RAG（完整版）
22 content slides + cover + copyright.

三大 Part：
  Part A — No Code / Low Code / GenAI 工具地圖（S1-S12）
  Part B — 提示工程（S13-S16）
  Part C — RAG / 微調 / 整合（S17-S22）

受眾：iPAS 初級考生，無技術背景亦可。
Aligned to chapters/Ch04_生成式AI工具提示工程與RAG/02_slides_design.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source, add_rect,
    set_solid_fill, set_no_fill, set_line, set_no_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_vs_two_col, draw_pyramid_stack,
    draw_three_blocks_flow, draw_image_placeholder,
    draw_inverted_thesis_box,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch04"
MODULE_TITLE = "生成式 AI 工具、提示工程與 RAG"
MODULE_SUBTITLE = "四象限工具地圖 \u00d7 RTFC 提示框架 \u00d7 RAG vs 微調\u2014\u2014工具選擇判斷力"
TIME_MIN = 120
N_CONTENT = 22


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _draw_bridge_note(slide, text, y_inches=6.6):
    add_textbox(
        slide, T.MARGIN_X, Inches(y_inches),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        text,
        font_size=T.FONT_SMALL, color=T.PRIMARY, bold=True, italic=True,
        align=PP_ALIGN.CENTER,
    )


def _draw_pitfall_card(slide, title_text, wrong, right, why=""):
    """PITFALL layout: title + wrong (red-bordered) + right (green-bordered) + optional why."""
    add_title(slide, title_text)
    col_w = Inches(5.5)
    gap = Inches(0.5)
    total = col_w * 2 + gap
    left_x = (T.SLIDE_W - total) / 2
    right_x = left_x + col_w + gap
    top = Inches(1.5)
    col_h = Inches(3.5)

    # Wrong box
    w_box = add_rect(slide, left_x, top, col_w, col_h)
    set_no_fill(w_box)
    set_line(w_box, RGBColor(0xC6, 0x28, 0x28), 2.0)
    add_textbox(
        slide, left_x, top, col_w, Inches(0.5),
        "WRONG",
        font_size=T.FONT_BODY, color=RGBColor(0xC6, 0x28, 0x28), bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide, left_x + Inches(0.25), top + Inches(0.6),
        col_w - Inches(0.5), col_h - Inches(0.7),
        wrong,
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )

    # Right box
    r_box = add_rect(slide, right_x, top, col_w, col_h)
    set_no_fill(r_box)
    set_line(r_box, T.PRIMARY, 2.0)
    add_textbox(
        slide, right_x, top, col_w, Inches(0.5),
        "RIGHT",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide, right_x + Inches(0.25), top + Inches(0.6),
        col_w - Inches(0.5), col_h - Inches(0.7),
        right,
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )

    if why:
        why_y = top + col_h + Inches(0.2)
        why_box = add_rect(slide, left_x, why_y, total, Inches(0.6))
        set_solid_fill(why_box, T.PRIMARY)
        set_no_line(why_box)
        add_textbox(
            slide, left_x + Inches(0.3), why_y,
            total - Inches(0.6), Inches(0.6),
            why,
            font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )


def build_ch04(output_path, image_registry=None):
    """Build Ch04 deck; 22 content slides + cover + copyright."""
    prs = _new_prs()

    # -- Cover --
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ===== S1 . MOTIVATION -- GenAI 工具題佔科目二最大比重 =====
    s = _blank(prs)
    draw_ask_page(
        s,
        "你每天都在用 ChatGPT\u2014\u2014\n但 Copilot Studio 和 GitHub Copilot 差在哪，\n你答得出來嗎？",
        data_card={
            "label": "L121 + L122 合計",
            "stat": "~30-40%",
            "caption": "初級科目二佔比\nNo Code/Low Code + GenAI 工具與提示工程",
        },
    )
    add_source(s, "iPAS 歷屆考古題分析 \u00b7 初級科目二")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # ===== S2 . ASK -- 工具那麼多，考的是選擇判斷力 =====
    s = _blank(prs)
    draw_ask_page(
        s,
        "ChatGPT\u3001Gemini\u3001Midjourney\u3001\nGitHub Copilot\u3001Cursor\u2014\u2014\n都是 AI 工具，但定位天差地遠。",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.2), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "用 Midjourney 寫程式？用 GitHub Copilot 做流程自動化？那就是選錯工具。",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ===== S3 . SILENT -- 立論 =====
    s = _blank(prs)
    draw_silent_page(
        s,
        "工具會用不難，\n難的是知道什麼時候用、怎麼用對。\n這正是考試在考的。",
    )
    add_footer(s, MODULE_CODE, 3, N_CONTENT, dark_bg=True)

    # ===== S4 . CONCEPT-CARD -- No Code =====
    s = _blank(prs)
    add_title(s, "No Code：完全不寫程式的開發方式")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
        "No Code = 用拖拉視覺化介面完成應用開發 \u00b7 目標用戶：非技術人員",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    draw_vs_two_col(
        s,
        left_title="優勢",
        right_title="限制",
        left_items=[
            "上手快 / 成本低",
            "迭代快 / 不需技術人員",
            "代表：Bubble / Airtable / Zapier",
        ],
        right_items=[
            "客製化程度低",
            "複雜邏輯難實現",
            "平台綁定（vendor lock-in）",
        ],
        summary="No Code 像 Canva\u2014\u2014不需要會 Photoshop 也能做圖，但天花板較低",
    )
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ===== S5 . CONCEPT-CARD -- Low Code =====
    s = _blank(prs)
    add_title(s, "Low Code：少量程式碼的加速開發")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
        "Low Code = 視覺化開發為主，必要時寫少量程式碼 \u00b7 目標：公民開發者 Citizen Developer",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    draw_vs_two_col(
        s,
        left_title="優勢",
        right_title="限制",
        left_items=[
            "比傳統開發快 3-5 倍",
            "有一定客製化能力",
            "企業級整合 / 適合內部應用",
            "代表：Power Platform / OutSystems",
        ],
        right_items=[
            "仍需基礎技術能力",
            "複雜整合有挑戰",
            "平台相依性",
            "學習曲線中等",
        ],
        summary="Low Code 是加速器\u2014\u2014視覺化搞定八成，剩下兩成用少量程式碼補",
    )
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ===== S6 . PITFALL (P1) -- No Code 和 Low Code 不是同義詞 =====
    s = _blank(prs)
    _draw_pitfall_card(
        s,
        "P1 No Code 和 Low Code 不是同義詞",
        wrong="「No Code 和 Low Code 基本上是同一種東西」",
        right="「No Code 完全不寫程式碼\uff1b\nLow Code 以視覺化為主但需要少量程式碼」",
    )
    draw_editorial_table(
        s,
        header=["", "No Code", "Low Code"],
        rows=[
            ["程式碼", "零", "少量"],
            ["目標用戶", "非技術人員", "有基礎技術者"],
            ["客製化", "低", "中"],
            ["複雜度上限", "低", "中高"],
        ],
        top=5.3,
        col_widths=[1.5, 2, 2],
    )
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ===== S7 . CONCEPT-CARD -- GenAI 工具全景：四大類工具地圖 =====
    s = _blank(prs)
    add_title(s, "GenAI 工具全景：四象限工具地圖")
    draw_matrix(
        s, rows=2, cols=2,
        cells=[
            {"text": "文字生成",
             "sub": "ChatGPT / Gemini / Claude\n通用對話、寫作、分析",
             "highlight": False},
            {"text": "圖片生成",
             "sub": "Midjourney / DALL-E / Stable Diffusion\n設計、創意、視覺",
             "highlight": False},
            {"text": "程式輔助",
             "sub": "GitHub Copilot / Cursor\n寫程式、補全、重構",
             "highlight": False},
            {"text": "流程自動化",
             "sub": "Copilot Studio / Power Automate\n聊天機器人、工作流程",
             "highlight": False},
        ],
        title="",
        top=1.3,
    )
    _draw_bridge_note(s, "場景決定象限，象限決定工具")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ===== S8 . CHECKPOINT -- No Code/Low Code/GenAI 場景配對 =====
    s = _blank(prs)
    add_title(s, "Check Point \u00b7 No Code / Low Code / GenAI 場景配對")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 8 / 22 \u00b7 不看投影片作答",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    draw_editorial_table(
        s,
        header=["場景", "工具/平台", "判斷依據"],
        rows=[
            ["行銷人員做線上問卷", "No Code (Airtable/Typeform)", "非技術、簡單需求"],
            ["工程師在 IDE 中補全程式碼", "GitHub Copilot", "程式輔助象限"],
            ["業務建客服聊天機器人", "Copilot Studio (Low Code)", "流程自動化、不寫程式"],
            ["PM 做內部審批 App", "Power Apps (Low Code)", "內部應用、基礎技術"],
        ],
        top=1.8,
        col_widths=[2.5, 2.5, 2],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.5), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "陷阱：「建聊天機器人」\u2260 選 ChatGPT\u2014\u2014ChatGPT 是通用對話，客製化機器人用 Copilot Studio",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ===== S9 . CONCEPT-CARD -- ChatGPT vs OpenAI API =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="ChatGPT \u2014 產品",
        right_title="OpenAI API \u2014 程式介面",
        left_items=[
            "面向終端用戶",
            "網頁 / App 介面",
            "直接對話",
            "適合個人使用",
        ],
        right_items=[
            "面向開發者",
            "程式碼呼叫",
            "system prompt + 參數控制",
            "適合嵌入產品",
        ],
        title="ChatGPT vs OpenAI API",
        summary="ChatGPT 是成品，API 是積木\u2014\u2014一個直接用，一個拿來蓋東西",
    )
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ===== S10 . PITFALL (P4) -- Copilot 家族定位大不同 =====
    s = _blank(prs)
    add_title(s, "P4 Copilot 家族定位大不同")
    _draw_pitfall_card(
        s,
        "P4 Copilot 家族定位大不同",
        wrong="「Copilot Studio 是用來輔助寫程式碼的工具」",
        right="「Copilot Studio 是 Low Code 流程自動化平台，\n用來建聊天機器人和工作流程」",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.5), T.SLIDE_W - 2 * T.MARGIN_X, Inches(1.2),
        "GitHub Copilot \u2192 IDE 外掛 \u00b7 寫程式用\n"
        "VS Code Copilot \u2192 GitHub Copilot 在 VS Code 中的實現\n"
        "Copilot Studio \u2192 Low Code 平台 \u00b7 聊天機器人 / 流程\n"
        "Microsoft 365 Copilot \u2192 Office 助手 \u00b7 Word/Excel/PPT",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        align=PP_ALIGN.CENTER,
        line_spacing=1.45,
    )
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ===== S11 . CONCEPT-CARD -- Midjourney / Gemini / 其他工具 =====
    s = _blank(prs)
    add_title(s, "Midjourney / Gemini / 其他工具定位")
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "Midjourney",
             "items": [
                 "圖片生成專精",
                 "文字描述 \u2192 圖片",
                 "藝術風格控制",
                 "平台：Discord",
             ]},
            {"heading": "Gemini (Google)",
             "items": [
                 "多模態 AI",
                 "文字 + 圖片 + 程式",
                 "Google 生態整合",
                 "長上下文",
             ]},
            {"heading": "其他",
             "items": [
                 "DALL-E (OpenAI 圖片)",
                 "Stable Diffusion (開源)",
                 "Claude (Anthropic)",
                 "Llama (Meta 開源)",
             ]},
        ],
        bottom_note="每個工具都有自己的生態位\u2014\u2014考試考你知不知道它的定位",
    )
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ===== S12 . CHECKPOINT -- GenAI 工具配對 =====
    s = _blank(prs)
    add_title(s, "Check Point \u00b7 GenAI 工具配對：給場景選工具")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 12 / 22 \u00b7 不看投影片作答",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(10.8), Inches(3.5),
        "Q1  設計師要生成產品概念插圖 \u2192 Midjourney / DALL-E\n\n"
        "Q2  工程師要在 VS Code 中補全 Python \u2192 GitHub Copilot\n\n"
        "Q3  客服部門要建 FAQ 自動回覆機器人 \u2192 Copilot Studio\n\n"
        "Q4  研究員要分析 100 頁 PDF \u2192 ChatGPT / Gemini\n\n"
        "Q5  開發者要在 App 中嵌入 AI 問答 \u2192 OpenAI API / Gemini API",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.0), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "判斷：先定象限（文字/圖片/程式/流程），再選具體工具",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ===== S13 . CONCEPT-CARD -- 提示工程 RTFC 四要素 =====
    s = _blank(prs)
    add_title(s, "提示工程：RTFC 四要素")
    draw_matrix(
        s, rows=2, cols=2,
        cells=[
            {"text": "R = Role \u89d2\u8272",
             "sub": "「你是一位...」\n設定 AI 的角色和專業背景",
             "highlight": True},
            {"text": "T = Task \u4efb\u52d9",
             "sub": "「請你...」\n明確說明要 AI 做什麼",
             "highlight": False},
            {"text": "F = Format \u683c\u5f0f",
             "sub": "「以...格式呈現」\n指定輸出格式（條列/表格/段落）",
             "highlight": False},
            {"text": "C = Constraint \u9650\u5236",
             "sub": "「注意.../不要...」\n設定邊界和限制條件",
             "highlight": True},
        ],
        title="",
        top=1.3,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.8), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.8),
        "範例：你是一位資深行銷專家(R)，請撰寫 500 字產品介紹(T)，\n使用條列式重點(F)，避免使用專業術語(C)。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        align=PP_ALIGN.CENTER,
        line_spacing=1.4,
    )
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ===== S14 . PITFALL (P2) -- 提示工程不是把問題打長一點 =====
    s = _blank(prs)
    _draw_pitfall_card(
        s,
        "P2 提示工程不是把問題打長一點",
        wrong="「幫我寫一篇文章，要很好很專業很詳細，\n"
              "內容要豐富，字數要多一點，語氣要正式。」\n\n"
              "\u2192 含糊、沒有具體指示",
        right="「你是一位資深科技記者(R)。請撰寫 800 字\n"
              "文章介紹 RAG 技術(T)，三段落：定義/流程/\n"
              "應用(F)。讀者為非技術管理者，避免程式碼(C)。」\n\n"
              "\u2192 結構清晰、具體可執行",
        why="好的提示不是字多，是有結構。",
    )
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ===== S15 . CONCEPT-CARD -- Zero-shot / Few-shot / CoT =====
    s = _blank(prs)
    add_title(s, "提示工程三大技巧：Zero-shot / Few-shot / CoT")
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "Zero-shot",
             "items": [
                 "不給範例直接問",
                 "「請翻譯以下句子...」",
                 "適合：簡單明確任務",
             ]},
            {"heading": "Few-shot",
             "items": [
                 "給 2-3 個範例",
                 "「範例一...範例二...現在請...」",
                 "適合：特定格式/風格",
             ]},
            {"heading": "CoT (Chain-of-Thought)",
             "items": [
                 "要求一步步推理",
                 "「請一步步思考...」",
                 "適合：推理/計算/邏輯",
             ]},
        ],
        bottom_note="Zero-shot 最快 / Few-shot 最穩 / CoT 最適合推理",
    )
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ===== S16 . PRACTICE -- 判斷提示好壞 =====
    s = _blank(prs)
    add_title(s, "Practice \u00b7 判斷提示好壞")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 16 / 22 \u00b7 30 秒自己判斷",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.0), Inches(1.8), Inches(11.0), Inches(4.0),
        "組一\n"
        "  A「幫我寫一封信」\n"
        "  B「你是人資主管，請撰寫 300 字新人歡迎信，語氣親切正式，含報到日期」\n"
        "  \u2192 B 好（有 RTFC）\n\n"
        "組二\n"
        "  A「這個數學題怎麼解？」\n"
        "  B「請一步步解這個數學題：...首先計算...然後...」\n"
        "  \u2192 B 好（用了 CoT）\n\n"
        "組三\n"
        "  A「翻譯以下內容」\n"
        "  B「範例：Hello\u2192你好 / Thank you\u2192謝謝 / 請翻譯：Goodbye」\n"
        "  \u2192 B 好（用了 Few-shot）",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.35,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "好提示的判斷標準：有結構、有具體性、有適當的技巧",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ===== S17 . CONCEPT-CARD -- RAG =====
    s = _blank(prs)
    add_title(s, "RAG：檢索增強生成")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(0.8),
        "LLM 的兩大痛點：\n"
        "\u2022 知識過時（訓練截止日後的事不知道）\n"
        "\u2022 幻覺（不知道就亂編）",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    draw_flow_chain(
        s,
        nodes=[
            {"label": "Retrieve", "sub": "從知識庫搜尋\n相關段落", "highlight": False},
            {"label": "Augment", "sub": "段落+問題\n組合增強提示", "highlight": True},
            {"label": "Generate", "sub": "LLM 根據增強\n提示生成回答", "highlight": False},
        ],
        title="",
        y=3.2,
    )
    draw_inverted_thesis_box(
        s,
        "RAG = 帶小抄考試。知識庫是小抄，LLM 是考生。不用全背，帶著小抄就能答對。",
        y=5.5,
    )
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ===== S18 . CONCEPT-CARD -- RAG vs 微調 =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="RAG 檢索增強生成",
        right_title="微調 Fine-tuning",
        left_items=[
            "不改模型，外掛知識庫",
            "知識可即時更新",
            "成本低（不用重新訓練）",
            "適合：知識會變動的場景",
            "類比：帶小抄考試",
        ],
        right_items=[
            "改模型權重",
            "學會新風格/領域/行為",
            "成本高（需要訓練資料+GPU）",
            "適合：改變模型核心行為",
            "類比：讀完整本書再考",
        ],
        title="RAG vs 微調 (Fine-tuning)",
        summary="知識會更新 \u2192 RAG\u3002行為要改變 \u2192 微調\u3002",
    )
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ===== S19 . CHECKPOINT -- RAG vs 微調場景判斷 =====
    s = _blank(prs)
    add_title(s, "Check Point \u00b7 RAG vs 微調場景判斷")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 19 / 22 \u00b7 不看投影片作答",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    draw_editorial_table(
        s,
        header=["場景", "選擇", "判斷依據"],
        rows=[
            ["公司內部文件問答，文件每月更新", "RAG", "知識常更新"],
            ["讓 AI 用台語風格回覆客戶", "微調", "改變語言風格/行為"],
            ["法律事務所查詢最新法規判例", "RAG", "即時檢索最新資訊"],
            ["醫療 AI 學會醫學術語和診斷邏輯", "微調", "學習專業領域行為"],
        ],
        top=1.8,
        col_widths=[3, 1, 2],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.5), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "口訣：知識更新 \u2192 RAG \u00b7 行為改變 \u2192 微調",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ===== S20 . CONCEPT-CARD -- AI 工具整合：組合拳思維 =====
    s = _blank(prs)
    add_title(s, "AI 工具整合：組合拳思維")
    draw_flow_chain(
        s,
        nodes=[
            {"label": "需求分析", "sub": "定義問題"},
            {"label": "AI 初稿", "sub": "ChatGPT\n生成文案"},
            {"label": "AI 配圖", "sub": "Midjourney\n生成視覺"},
            {"label": "AI 程式", "sub": "Copilot\n寫程式碼"},
            {"label": "自動化", "sub": "Power Automate\n串接流程"},
        ],
        y=2.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.5), T.SLIDE_W - 2 * T.MARGIN_X, Inches(1.5),
        "\u2022 實務中很少只用一個 AI 工具\n"
        "\u2022 考試考的是整合思維\u2014\u2014知道每個環節該用什麼工具\n"
        "\u2022 AI 工具是輔助，最後一關永遠是人工審核",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    warn_y = Inches(6.2)
    warn = add_rect(s, T.MARGIN_X, warn_y,
                    T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5))
    set_solid_fill(warn, T.PRIMARY)
    set_no_line(warn)
    add_textbox(
        s, T.MARGIN_X + Inches(0.3), warn_y,
        T.SLIDE_W - 2 * T.MARGIN_X - Inches(0.6), Inches(0.5),
        "GenAI 限制：幻覺 / 偏見 / 隱私\u2014\u2014人工審核不能省",
        font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ===== S21 . PYRAMID -- GenAI 工具選擇三層判斷 =====
    s = _blank(prs)
    draw_pyramid_stack(
        s,
        title="Ch04 收束：GenAI 工具選擇三層判斷",
        layers=[
            {"name": "第三層：知識需求",
             "caption": "通用 LLM / RAG（知識更新）/ 微調（行為改變）"},
            {"name": "第二層：技術能力",
             "caption": "No Code（非技術）/ Low Code（基礎技術）/ Pro Code（開發者）"},
            {"name": "第一層：任務類型",
             "caption": "文字生成 / 圖片生成 / 程式輔助 / 流程自動化"},
        ],
        thesis="三層過濾，工具自動浮現。",
    )
    add_footer(s, MODULE_CODE, 21, N_CONTENT)

    # ===== S22 . MOTIVATION -- 銜接 Ch05 =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="本章完成",
        right_title="Ch05 預告",
        left_items=[
            "No Code / Low Code",
            "GenAI 工具地圖（四象限）",
            "提示工程 RTFC + 三大技巧",
            "RAG vs 微調",
        ],
        right_items=[
            "GenAI 導入評估",
            "幻覺與偏見風險",
            "倫理與法規",
            "成本效益分析",
        ],
        title="知道工具怎麼用，接下來看導入的風險",
        summary="工具能用不代表該用\u2014\u2014導入 GenAI 有哪些風險？怎麼評估？下一章告訴你。",
    )
    add_footer(s, MODULE_CODE, 22, N_CONTENT)

    # -- Copyright --
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
