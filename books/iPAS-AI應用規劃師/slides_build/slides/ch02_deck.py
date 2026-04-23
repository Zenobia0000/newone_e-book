"""Ch02 deck -- 資料處理分析與隱私安全（完整版）
18 content slides + cover + copyright.

三大 Part：
  Part A -- 資料概念（S1-S6）
  Part B -- 資料處理與特徵工程（S7-S14）
  Part C -- 隱私安全與收束（S15-S18）

受眾：iPAS 初級考生，所有案例走「考題情境」路線。
Aligned to chapters/Ch02_資料處理分析與隱私安全/02_slides_design.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source, add_rect,
    set_solid_fill, set_no_fill, set_line, set_no_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_vs_two_col, draw_pyramid_stack,
    draw_three_blocks_flow, draw_image_placeholder,
    draw_inverted_thesis_box,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch02"
MODULE_TITLE = "資料處理分析與隱私安全"
MODULE_SUBTITLE = "5V 量速多真值 × 五階段流程 × CIA 三原則——資料是 AI 的燃料"
TIME_MIN = 90
N_CONTENT = 18


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _draw_bridge_note(slide, text, y_inches=6.6):
    add_textbox(
        slide, T.MARGIN_X, Inches(y_inches),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        text,
        font_size=T.FONT_SMALL, color=T.PRIMARY, bold=True, italic=True,
        align=PP_ALIGN.CENTER,
    )


def build_ch02(output_path, image_registry=None):
    """Build Ch02 deck; 18 content slides + cover + copyright."""
    prs = _new_prs()

    # == Cover ==
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ===== S1 . MOTIVATION -- L112 考什麼？資料題佔分地圖 =====
    s = _blank(prs)
    draw_ask_page(
        s,
        "L112：資料概念 + 處理流程 + 隱私安全\n≈ 20-25% 考題",
        data_card={
            "label": "L112 佔分",
            "stat": "~12 題",
            "caption": "資料是 AI 的燃料\n燃料搞懂，L112 穩拿",
        },
    )
    add_source(s, "iPAS 歷屆考古題分析 · L112 佔分估算")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # ===== S2 . CONCEPT-CARD -- 大數據 5V 特性 =====
    s = _blank(prs)
    add_title(s, "大數據 5V 特性")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(0.5),
        "記憶口訣：「量速多真值」",
        font_size=Pt(16), color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(2.0), Inches(6.5), Inches(4.0),
        "Volume（量）：資料量大到傳統工具處理不了 → TB/PB 級\n\n"
        "Velocity（速）：資料產生速度快，需即時處理 → 串流資料\n\n"
        "Variety（多樣）：資料格式多樣 → 文字/圖片/影音/log\n\n"
        "Veracity（真實性）：資料品質有疑慮 → 雜訊、缺失、不一致\n\n"
        "Value（價值）：從海量資料萃取商業洞見 → 最終目的",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(4.5),
        slot_name="5V 五角形圖",
        description="五角形五頂點各標一個 V，每個 V 配圖示和一句話定義。",
        size_hint="1200x1000 px",
        placeholder_id="Ch02_S2_5v_pentagon",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ===== S3 . CONCEPT-CARD -- 資料型態三分法 =====
    s = _blank(prs)
    add_title(s, "資料型態三分法")
    draw_editorial_table(
        s,
        header=["", "結構化", "半結構化", "非結構化"],
        rows=[
            ["特徵", "固定 schema、欄位明確", "有標籤/標記、schema 彈性", "無固定格式"],
            ["格式", "SQL 表格、CSV、Excel", "JSON、XML、HTML", "圖片、影音、文字"],
            ["儲存", "關聯式資料庫", "NoSQL / 文件型資料庫", "物件儲存 / 資料湖"],
        ],
        top=1.3,
        col_widths=[1.2, 2.8, 2.8, 2.8],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.2), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "判斷口訣：有固定欄位→結構 / 有標籤但彈性→半結構 / 什麼都沒有→非結構",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ===== S4 . PRACTICE -- 5V 配對 + 資料型態辨識 =====
    s = _blank(prs)
    add_title(s, "Practice . 5V 配對 + 資料型態辨識")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 4 / 18 . 六個場景快速配對",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.0), Inches(1.8), Inches(11.0), Inches(4.5),
        "場景 1：社群平台每秒產生數千則貼文 → Velocity\n\n"
        "場景 2：電商資料庫有 10 億筆交易紀錄 → Volume\n\n"
        "場景 3：感測器收集的資料中有 15% 是雜訊 → Veracity\n\n"
        "場景 4：醫院 X 光片影像 → 非結構化\n\n"
        "場景 5：API 回傳的 JSON 資料 → 半結構化\n\n"
        "場景 6：公司 ERP 系統的訂單表格 → 結構化",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ===== S5 . PITFALL (P1, P2) -- 5V 混淆 + 半結構化陷阱 =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="常見錯誤",
        right_title="正確認知",
        left_items=[
            "資料有雜訊代表沒有 Value",
            "JSON 是結構化資料",
            "Excel 是半結構化",
        ],
        right_items=[
            "Veracity 是品質問題，Value 是價值問題，兩者獨立",
            "JSON 有標籤但 schema 彈性，是半結構化",
            "Excel 表格有固定欄位，是結構化",
        ],
        title="PITFALL：5V 混淆 + 半結構化陷阱",
        summary="判斷反射：先問「有沒有固定 schema？」→ 有=結構 / 有標籤但彈性=半結構 / 都沒有=非結構",
    )
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ===== S6 . CHECKPOINT -- 大數據與資料型態檢核 =====
    s = _blank(prs)
    add_title(s, "Check Point . 大數據與資料型態檢核")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 6 / 18 . 不看投影片作答",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(10.8), Inches(3.5),
        "Q1  5V 中哪個 V 關注「資料品質與可信度」？\n\n"
        "Q2  CSV 檔案屬於什麼資料型態？\n\n"
        "Q3  YouTube 上的影片屬於什麼資料型態？",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.0), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "答案：Q1 Veracity . Q2 結構化 . Q3 非結構化",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ===== S7 . CONCEPT-CARD -- 資料處理流程五階段 =====
    s = _blank(prs)
    add_title(s, "資料處理流程五階段")
    draw_flow_chain(
        s,
        title="",
        nodes=[
            {"label": "收集", "sub": "爬蟲/API/問卷"},
            {"label": "清理", "sub": "缺失/重複/異常"},
            {"label": "轉換", "sub": "編碼/標準化"},
            {"label": "分析", "sub": "統計/建模"},
            {"label": "呈現", "sub": "圖表/儀表板"},
        ],
        y=2.0,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(3.8), T.SLIDE_W - 2 * T.MARGIN_X, Inches(2.5),
        "① 收集：爬蟲、API、問卷、感測器、資料庫匯出\n"
        "② 清理：處理缺失值、移除重複、修正異常值、格式統一\n"
        "③ 轉換：特徵工程、One-hot/Label 編碼、標準化/正規化\n"
        "④ 分析：描述統計、視覺化探索、相關性分析、模型建構\n"
        "⑤ 呈現：圖表、儀表板、報告、摘要\n\n"
        "比喻：買食材→洗切→調味→下鍋→擺盤",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ===== S8 . PRACTICE -- 流程階段配對 =====
    s = _blank(prs)
    add_title(s, "Practice . 流程階段配對")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 8 / 18 . 六個操作配五個階段",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.0), Inches(1.8), Inches(11.0), Inches(4.5),
        "操作 1：用爬蟲抓取新聞網站資料 → 收集\n\n"
        "操作 2：刪除重複的客戶紀錄 → 清理\n\n"
        "操作 3：將類別欄位做 One-hot 編碼 → 轉換\n\n"
        "操作 4：計算各欄位的平均數和標準差 → 分析\n\n"
        "操作 5：製作互動式 dashboard → 呈現\n\n"
        "操作 6：將年齡欄位中的「-5」替換為中位數 → 清理（異常值處理）",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    _draw_bridge_note(s, "分界線：資料品質問題→清理；資料格式問題→轉換")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ===== S9 . PITFALL (P3) -- 資料清理不只是刪 NULL =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="常見錯誤",
        right_title="正確認知",
        left_items=[
            "資料清理 = 刪除缺失值",
            "只會 dropna() 就夠了",
        ],
        right_items=[
            "清理 = 缺失值 + 異常值 + 重複值 + 格式統一 + 型態轉換",
            "缺失值處理：刪除 / 填均值中位數 / 模型預測填補",
            "異常值處理：刪除 / 截斷 capping / 替換合理值",
        ],
        title="PITFALL：資料清理不只是刪 NULL",
        summary="考題常問：「以下哪些操作屬於資料清理？」→ 五項全選",
    )
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ===== S10 . CHECKPOINT -- 資料處理流程情境判斷 =====
    s = _blank(prs)
    add_title(s, "Check Point . 資料處理流程情境判斷")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 10 / 18 . 每題選擇所屬階段",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(10.8), Inches(3.5),
        "Q1  將「男/女」欄位轉換為 0/1\n"
        "     → 轉換（編碼）\n\n"
        "Q2  發現某筆訂單金額為 -999 並修正\n"
        "     → 清理（異常值）\n\n"
        "Q3  用 matplotlib 繪製銷售趨勢圖\n"
        "     → 呈現\n\n"
        "Q4  透過 API 取得天氣資料\n"
        "     → 收集",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.0), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "注意：編碼是把乾淨資料變成模型能用的格式，不是修資料品質問題",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ===== S11 . CONCEPT-CARD -- 特徵工程四大操作 =====
    s = _blank(prs)
    add_title(s, "特徵工程四大操作")
    draw_matrix(
        s,
        rows=2, cols=2,
        cells=[
            {"text": "特徵選取 Selection",
             "sub": "挑出有用特徵，丟掉無關的\n相關性分析、特徵重要性排序",
             "highlight": False},
            {"text": "特徵建立 Creation",
             "sub": "從現有特徵組合出新特徵\n從「生日」算出「年齡」",
             "highlight": False},
            {"text": "特徵轉換 Transformation",
             "sub": "改變特徵的表示方式\nOne-hot 編碼、Label 編碼",
             "highlight": False},
            {"text": "特徵縮放 Scaling",
             "sub": "統一不同特徵的尺度\n標準化 Z-score、正規化 Min-Max",
             "highlight": True},
        ],
        title="",
        top=1.3,
    )
    _draw_bridge_note(s, "特徵工程 = 幫模型準備最好消化的食材")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ===== S12 . EXAMPLE-I/O -- 標準化 vs 正規化：考題實戰 =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="標準化 Z-score",
        right_title="正規化 Min-Max",
        left_items=[
            "公式：(X - mean) / std",
            "[10,20,30] → [-1, 0, 1]",
            "結果以 0 為中心，無固定範圍",
            "適用：常態分布、SVM/PCA",
        ],
        right_items=[
            "公式：(X - min) / (max - min)",
            "[10,20,30] → [0, 0.5, 1]",
            "結果縮放到 [0, 1]",
            "適用：固定範圍需求、NN/KNN",
        ],
        title="標準化 vs 正規化：考題實戰",
        summary="考試不考計算，考的是概念和使用時機。",
    )
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ===== S13 . PITFALL (P4) -- 標準化 vs 正規化分不清 =====
    s = _blank(prs)
    add_title(s, "PITFALL：標準化 vs 正規化分不清")
    draw_editorial_table(
        s,
        header=["", "標準化 Standardization", "正規化 Normalization"],
        rows=[
            ["公式", "(X-mean)/std", "(X-min)/(max-min)"],
            ["結果範圍", "無固定範圍（以 0 為中心）", "[0, 1]"],
            ["適用情境", "常態分布資料、SVM/PCA", "固定範圍需求、NN/KNN"],
        ],
        top=1.3,
        col_widths=[1.5, 3.5, 3.5],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.0), T.SLIDE_W - 2 * T.MARGIN_X, Inches(2.0),
        "錯：「Min-Max 是標準化」\n"
        "對：「Min-Max 是正規化（Normalization），Z-score 才是標準化（Standardization）」\n\n"
        "記憶法：「標=Z、正=MinMax」\n"
        "標準化的「標」長得像 Z → 標=Z-score\n"
        "正規化的「正」像 0 到 1 之間的正數 → 正=Min-Max",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ===== S14 . CHECKPOINT -- 特徵工程與縮放 =====
    s = _blank(prs)
    add_title(s, "Check Point . 特徵工程與縮放")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 14 / 18 . 每題 15 秒",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(10.8), Inches(3.5),
        "Q1  年齡 0-80、薪水 20000-200000，要統一尺度→用什麼？\n"
        "     → 特徵縮放（標準化或正規化）\n\n"
        "Q2  從日期欄位取出「是否為週末」→ 哪種特徵操作？\n"
        "     → 特徵建立\n\n"
        "Q3  刪除與目標變數相關性低於 0.1 的欄位\n"
        "     → 特徵選取",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.0), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "如果你答對了三題，特徵工程的概念已經到位",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ===== S15 . CONCEPT-CARD -- 資料隱私：個資法 + 去識別化技術 =====
    s = _blank(prs)
    add_title(s, "資料隱私：個資法 + 去識別化技術")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.0), Inches(1.5),
        "個資定義：可直接或間接識別特定個人的資料\n\n"
        "個資法核心原則：\n"
        "  告知義務、目的限制、最小蒐集、安全維護",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    draw_editorial_table(
        s,
        header=["技術", "說明", "可逆性"],
        rows=[
            ["匿名化 Anonymization", "刪除所有可識別資訊", "不可逆"],
            ["假名化 Pseudonymization", "用代號替換，保留對照表", "可逆"],
            ["資料遮蔽 Data Masking", "部分隱藏，如 A***1234", "部分可逆"],
        ],
        top=3.8,
        col_widths=[2.5, 4.0, 1.5],
    )
    _draw_bridge_note(s, "考試重點：匿名化是不可逆的，假名化是可逆的")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ===== S16 . PRACTICE -- 隱私安全情境判斷 =====
    s = _blank(prs)
    add_title(s, "Practice . 隱私安全情境判斷")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 16 / 18 . 五個情境配對",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.0), Inches(1.8), Inches(11.0), Inches(4.5),
        "情境 1：醫院將病歷中姓名替換為隨機編號，保留對照表\n"
        "         → 假名化\n\n"
        "情境 2：研究機構將資料中所有可識別欄位永久刪除\n"
        "         → 匿名化\n\n"
        "情境 3：電商收集用戶位置資訊但未告知用戶\n"
        "         → 違反告知義務\n\n"
        "情境 4：公司將員工身分證字號顯示為 A1234***89\n"
        "         → 資料遮蔽\n\n"
        "情境 5：行銷部門將客戶資料分享給未經授權的第三方\n"
        "         → 違反目的限制 + 安全維護",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.35,
    )
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ===== S17 . CONCEPT-CARD -- 資料安全 CIA 三原則 =====
    s = _blank(prs)
    add_title(s, "資料安全 CIA 三原則")
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "Confidentiality 機密性",
             "items": [
                 "只有授權的人能存取",
                 "違反：資料庫密碼外洩",
                 "關鍵字：未授權存取",
             ]},
            {"heading": "Integrity 完整性",
             "items": [
                 "資料未被非法竄改",
                 "違反：駭客修改交易金額",
                 "注意：不是「資料完不完整」",
             ]},
            {"heading": "Availability 可用性",
             "items": [
                 "需要時能正常存取",
                 "違反：DDoS 導致癱瘓",
                 "關鍵字：系統不可用",
             ]},
        ],
        bottom_note="記憶口訣：CIA = 機完用（機密/完整/可用）",
    )
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ===== S18 . PYRAMID -- 全章收束：資料生命週期金字塔 =====
    s = _blank(prs)
    draw_pyramid_stack(
        s,
        title="Ch02 收束：資料生命週期金字塔",
        layers=[
            {"name": "資料隱私安全",
             "caption": "個資法 + 去識別化 + CIA"},
            {"name": "資料處理分析",
             "caption": "五階段 + 特徵工程 + 縮放"},
            {"name": "資料概念",
             "caption": "5V + 型態三分法"},
        ],
        thesis="考試反射：5V 配場景、型態看 schema、流程看操作、隱私看法規",
    )
    _draw_bridge_note(s, "→ Ch03 機器學習——有了乾淨的資料，下一步就是建模")
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # == Copyright ==
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
