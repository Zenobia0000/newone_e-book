"""Ch06 deck — AI 技術應用：NLP / CV / 多模態（完整版）
22 content slides + cover + copyright.

三大 Part：
  Part A — NLP（S1-S8）
  Part B — CV + Transformer（S9-S13）
  Part C — 多模態 + 收束（S14-S22）

受眾：iPAS AI 應用規劃師中級備考。
Aligned to chapters/Ch06_AI技術應用_NLP_CV_多模態/02_slides_design.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source, add_rect,
    set_solid_fill, set_no_fill, set_line, set_no_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_vs_two_col, draw_pyramid_stack,
    draw_three_blocks_flow, draw_image_placeholder,
    draw_inverted_thesis_box, draw_grid,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch06"
MODULE_TITLE = "AI \u6280\u8853\u61c9\u7528\uff1aNLP / CV / \u591a\u6a21\u614b"
MODULE_SUBTITLE = "NLP Pipeline \u00d7 CV \u56db\u5927\u4efb\u52d9 \u00d7 Transformer \u00d7 \u591a\u6a21\u614b\u878d\u5408\u2014\u2014\u6280\u8853\u9078\u64c7\u529b"
TIME_MIN = 120
N_CONTENT = 22


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_ch06(output_path, image_registry=None):
    """Build Ch06 deck; 22 content slides + cover + copyright."""
    prs = _new_prs()

    # == Cover ==
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ===== S1 . SILENT . 起 =====
    s = _blank(prs)
    draw_silent_page(
        s,
        "中級科目一的核心：\n你能否在 NLP、CV、GenAI、多模態\n之間選對技術。",
    )
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ===== S2 . ASK . 起 =====
    s = _blank(prs)
    draw_ask_page(
        s,
        "什麼時候該用 AI，\n什麼時候寫個 if-else 就夠了？",
    )
    # Right side: two-column comparison via textbox
    col_x = Inches(8.0)
    col_w = Inches(4.8)
    add_textbox(
        s, col_x, Inches(4.0), col_w, Inches(3.0),
        "AI 適用：模式複雜、資料量大、規則難窮舉\n"
        "Rule-based：規則明確、場景有限、可解釋性要求高",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ===== S3 . MECHANISM_FLOW . 起 — NLP Pipeline =====
    s = _blank(prs)
    add_title(s, "NLP Pipeline：把文字變成數字，讓機器能處理")
    draw_flow_chain(
        s,
        nodes=[
            {"label": "原始文字", "sub": "清洗/正規化"},
            {"label": "Tokenization", "sub": "切成最小單位"},
            {"label": "向量化", "sub": "W2V/TF-IDF/Emb"},
            {"label": "下游任務", "sub": "分類/NER/情緒/摘要",
             "highlight": True},
        ],
        y=2.5,
    )
    draw_inverted_thesis_box(
        s,
        "NLP 全流程 = 文字 → 數字 → 任務。理解 pipeline 就能應對大部分考題。",
        y=5.8,
    )
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ===== S4 . TABLE . 起 — Tokenization 三切法 =====
    s = _blank(prs)
    add_title(s, "Tokenization 三種切法——BPE 是現在的主流")
    draw_editorial_table(
        s,
        header=["切法", "粒度", "詞彙量", "OOV 問題", "序列長度", "代表系統"],
        rows=[
            ["字元級", "最細", "極小", "無", "很長", "早期 NLP"],
            ["詞級", "最粗", "極大", "嚴重", "短", "傳統分詞"],
            ["子詞級(BPE)", "折衷", "適中", "極少", "適中", "GPT / BERT"],
        ],
        top=1.3,
        col_widths=[1.2, 0.8, 0.8, 0.8, 0.8, 1.2],
    )
    draw_inverted_thesis_box(
        s,
        "Q2 考 Tokenization 定義——「把文字切成模型能處理的最小單位」。",
        y=5.5,
    )
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ===== S5 . VS . 承 — Word2Vec vs TF-IDF =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        title="Word2Vec vs TF-IDF——本質不同",
        left_title="Word2Vec",
        right_title="TF-IDF",
        left_items=[
            "稠密向量（~300 維）",
            "捕捉語義關係",
            "king - man + woman \u2248 queen",
            "不看文件頻率",
            "關鍵字：語義/向量/類比",
        ],
        right_items=[
            "稀疏向量（維度=詞彙量）",
            "衡量詞重要度",
            "TF \u00d7 IDF",
            "不懂語義",
            "關鍵字：頻率/重要度",
        ],
        summary="看到「頻率/重要度」→ TF-IDF；看到「向量/語義/類比」→ Word2Vec。",
    )
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ===== S6 . TABLE . 承 — NLP 下游任務 =====
    s = _blank(prs)
    add_title(s, "NLP 下游任務四大天王——考試就考配對")
    draw_editorial_table(
        s,
        header=["NLP 任務", "定義", "典型應用", "相關技術"],
        rows=[
            ["情緒分析", "判斷正/負面", "客戶評論分析", "BERT / RoBERTa"],
            ["NER", "抽取實體名稱", "合約甲乙方抽取", "CRF / BiLSTM-CRF"],
            ["文本分類", "歸類別", "郵件分類/主題分類", "CNN / Transformer"],
            ["文本摘要", "壓縮重點", "新聞摘要", "Seq2Seq / T5"],
        ],
        top=1.3,
        col_widths=[1.2, 1.2, 1.5, 1.5],
    )
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ===== S7 . PITFALL . 承 — Word2Vec 常見混淆 =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        title="PITFALL：Word2Vec 是衡量詞頻重要性的方法？",
        left_title="常見錯誤 \u2717",
        right_title="正確理解 \u2713",
        left_items=[
            "Word2Vec 是衡量詞頻重要性",
            "Word2Vec 是稀疏向量",
            "Word2Vec 看文件頻率",
        ],
        right_items=[
            "Word2Vec 將詞映射為稠密向量",
            "捕捉語義關係",
            "king-man+woman\u2248queen",
        ],
        summary="考題最愛把 Word2Vec 和 TF-IDF 的描述交叉放在選項裡。",
    )
    draw_image_placeholder(
        s, Inches(4.5), Inches(5.6), Inches(4.5), Inches(1.5),
        slot_name="Word2Vec 向量空間示意圖",
        description="king/queen/man/woman 四點位置關係圖。",
        size_hint="900\u00d7300 px",
        placeholder_id="Ch06_S7_w2v_analogy",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ===== S8 . SILENT . 承 — 過渡到 CV =====
    s = _blank(prs)
    draw_silent_page(
        s,
        "從文字到圖像\u2014\u2014CV 的世界觀完全不同。\n"
        "NLP 處理的是序列，\nCV 處理的是像素矩陣。",
    )
    add_footer(s, MODULE_CODE, 8, N_CONTENT, dark_bg=True)

    # ===== S9 . MATRIX . 轉 — CV 四大任務 =====
    s = _blank(prs)
    draw_matrix(
        s, rows=2, cols=2,
        title="CV 四大任務：粒度 \u00d7 是否分個體",
        cells=[
            {"text": "圖像分類",
             "sub": "整張圖 + 單標籤\n貓/狗/車",
             "highlight": False},
            {"text": "物件偵測",
             "sub": "整張圖 + 多 BBox\nYOLO / Faster R-CNN",
             "highlight": False},
            {"text": "語義分割",
             "sub": "像素級 + 不分個體\n所有貓同色",
             "highlight": False},
            {"text": "實例分割",
             "sub": "像素級 + 分個體\n每隻貓不同色",
             "highlight": True},
        ],
        top=1.4, bottom=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "橫軸：整張圖 / 像素級。縱軸：單一類別 / 多個個體。Q7 考的就是這個矩陣。",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ===== S10 . MECHANISM_FLOW . 轉 — CNN → ResNet → ViT =====
    s = _blank(prs)
    add_title(s, "CV 演進：CNN → ResNet → ViT")
    draw_flow_chain(
        s,
        nodes=[
            {"label": "CNN 基礎", "sub": "LeNet / AlexNet",
             "caption": "卷積提取局部特徵"},
            {"label": "深層 CNN", "sub": "VGG / ResNet",
             "caption": "殘差連接解決深度",
             "highlight": True},
            {"label": "ViT", "sub": "Vision Transformer",
             "caption": "圖像切 patch 當 token"},
        ],
        y=2.5,
    )
    draw_inverted_thesis_box(
        s,
        "ViT 是理解「為什麼 Transformer 能跨領域」的關鍵——跨領域遷移。",
        y=5.5,
    )
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ===== S11 . PITFALL . 轉 — 語義分割 vs 實例分割 =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        title="語義分割 vs 實例分割\u2014\u2014這張圖記住，考試就不會錯",
        left_title="語義分割",
        right_title="實例分割",
        left_items=[
            "兩隻貓都標為橙色（同類同色）",
            "只分類別，不分個體",
            "語義 = 類別語義",
            "適合：道路/天空分割",
        ],
        right_items=[
            "貓A 藍色、貓B 紅色（同類不同色）",
            "分類別且分個體",
            "實例 = 個體",
            "適合：自動駕駛辨識前方各車",
        ],
        summary="語義不分個體，實例分個體。自動駕駛需要知道「前面有三輛車各在哪」→ 實例分割。",
    )
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ===== S12 . MECHANISM_FLOW . 轉 — Transformer 架構 =====
    s = _blank(prs)
    add_title(s, "Transformer 核心：Self-Attention")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.5), Inches(0.6),
        "Attention(Q,K,V) = softmax(QK\u1d40/\u221ad_k)V",
        font_size=Pt(16), color=T.PRIMARY, bold=True,
        family=T.FONT_MONO,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(2.2), Inches(6.5), Inches(3.5),
        "簡化流程：\n"
        "Input Embedding + Positional Encoding\n"
        "  \u2192 Multi-Head Self-Attention\n"
        "  \u2192 Feed Forward\n"
        "  \u2192 Output\n\n"
        "Self-Attention 機制：\n"
        "\u2022 每個 token \u2192 Q（查詢）、K（鍵）、V（值）\n"
        "\u2022 Q \u00d7 K = 相似度 \u2192 注意力權重\n"
        "\u2022 權重 \u00d7 V = 加權輸出\n"
        "\u2022 位置編碼補上順序資訊",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    draw_image_placeholder(
        s, Inches(7.5), Inches(1.3), Inches(5.2), Inches(4.5),
        slot_name="Transformer 架構簡化圖",
        description="左側：完整 Transformer Block 流程。右側：Self-Attention Q/K/V 計算放大圖。",
        size_hint="1200\u00d71000 px",
        placeholder_id="Ch06_S12_transformer_arch",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ===== S13 . TABLE . 轉 — 模型壓縮三劍客 =====
    s = _blank(prs)
    add_title(s, "模型壓縮三劍客：蒸餾 / 剪枝 / 量化")
    draw_editorial_table(
        s,
        header=["壓縮方法", "原理", "優點", "缺點", "精度損失", "速度提升"],
        rows=[
            ["知識蒸餾", "大模型教小模型", "結構靈活", "需重新訓練", "中", "高"],
            ["剪枝", "砍掉不重要連接", "保持架構", "需微調", "低-中", "中"],
            ["量化", "降低數字精度", "簡單快速", "精度下降", "低", "高"],
        ],
        top=1.3,
        col_widths=[1.2, 1.5, 1.0, 1.0, 0.8, 0.8],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.0), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "量化階梯：FP32 \u2192 FP16 \u2192 INT8 \u2192 INT4",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    draw_inverted_thesis_box(
        s,
        "記憶法：蒸「教」、剪「砍」、量「降」。Q8 考模型壓縮原理。",
        y=5.0,
    )
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ===== S14 . MATRIX . 轉 — 三種融合策略 =====
    s = _blank(prs)
    draw_three_blocks_flow(
        s,
        title="多模態融合三策略：什麼時候合併？",
        blocks=[
            {"heading": "早期融合",
             "items": [
                 "特徵層合併",
                 "模態交互充分",
                 "維度爆炸",
                 "適合模態高度相關",
             ]},
            {"heading": "晚期融合",
             "items": [
                 "決策層合併",
                 "模組化/各自優化",
                 "錯過低層交互",
                 "適合模態可獨立處理",
             ]},
            {"heading": "注意力機制融合",
             "items": [
                 "Cross-Attention",
                 "折衷方案",
                 "目前最先進",
                 "計算量適中",
             ]},
        ],
        bottom_note="Q11 考 Early Fusion。融合位置 = pipeline 前端/後端/中間。",
    )
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ===== S15 . VS . 轉 — 早期 vs 晚期融合 =====
    s = _blank(prs)
    draw_vs_two_col(
        s,
        title="早期融合 vs 晚期融合",
        left_title="早期融合（特徵層）",
        right_title="晚期融合（決策層）",
        left_items=[
            "在特徵層合併",
            "模態交互深",
            "維度高 / 計算量大",
            "適合模態高度相關",
            "例：醫療影像 + 病歷",
        ],
        right_items=[
            "在決策層合併",
            "模組化 / 各自優化",
            "錯過低層交互",
            "適合模態可獨立處理",
            "例：語音辨識 + 文字分類",
        ],
        summary="注意力機制融合（Cross-Attention）= 折衷方案。選擇取決於模態相關程度。",
    )
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ===== S16 . MATRIX . 合 — 多模態醫療案例 =====
    s = _blank(prs)
    add_title(s, "多模態醫療案例全景圖")
    draw_flow_chain(
        s,
        nodes=[
            {"label": "醫學影像", "sub": "X光/MRI",
             "caption": "ViT 處理"},
            {"label": "電子病歷", "sub": "文字",
             "caption": "BERT 處理"},
            {"label": "基因數據", "sub": "序列",
             "caption": "1D-CNN 處理"},
            {"label": "Cross-Attention\n融合", "sub": "三路匯合",
             "highlight": True,
             "caption": ""},
            {"label": "輔助診斷", "sub": "輸出結果",
             "caption": "最終預測"},
        ],
        y=2.2,
    )
    draw_inverted_thesis_box(
        s,
        "Q12 考這個。影像用 ViT、文字用 BERT、基因用 1D-CNN、融合用 Cross-Attention。",
        y=5.5,
    )
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ===== S17 . PYRAMID . 合 — 四層技術金字塔 =====
    s = _blank(prs)
    draw_pyramid_stack(
        s,
        title="Ch06 收束：四層技術金字塔",
        layers=[
            {"name": "多模態",
             "caption": "融合應用 \u2192 技術選擇力"},
            {"name": "GenAI / Transformer",
             "caption": "跨領域引擎"},
            {"name": "CV",
             "caption": "圖像理解基礎"},
            {"name": "NLP",
             "caption": "文字處理基礎"},
        ],
        thesis="NLP 和 CV 是基礎，Transformer 是引擎，多模態是整合。考試難題在頂端。",
    )
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ===== S18 . EXAM_DRILL . 合 — NLP 考題 =====
    s = _blank(prs)
    add_title(s, "Exam Drill \u00b7 NLP 技術辨識題")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "每題 60 秒 \u00b7 NLP 題訣竅：抓選項裡的關鍵字",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(0.8), Inches(1.9), Inches(11.5), Inches(4.0),
        "Q1  Tokenization 的定義是？\n"
        "    (A) 把文字轉成向量 (B) 把文字切成最小處理單位\n"
        "    (C) 計算詞頻 (D) 語義分析\n\n"
        "Q2  Word2Vec 的核心特性是？\n"
        "    (A) 計算詞頻重要度 (B) 將詞映射為稠密向量並捕捉語義\n"
        "    (C) 文件分類 (D) 稀疏矩陣分解\n\n"
        "Q3  從合約中抽取甲乙方名稱，應用哪個 NLP 任務？\n"
        "    (A) 情緒分析 (B) 文本分類 (C) NER (D) 文本摘要",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "答案：Q1(B) 切成最小單位 \u00b7 Q2(B) 稠密向量+語義 \u00b7 Q3(C) NER=抽取實體名稱",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ===== S19 . EXAM_DRILL . 合 — CV + 壓縮考題 =====
    s = _blank(prs)
    add_title(s, "Exam Drill \u00b7 CV 技術配對 + 模型壓縮")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "CV 配對看「輸出格式」 \u00b7 壓縮題看「原理」",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(0.8), Inches(1.9), Inches(11.5), Inches(4.0),
        "Q1  自動駕駛需辨識前方每輛車的獨立位置，應用哪個 CV 任務？\n"
        "    (A) 圖像分類 (B) 物件偵測 (C) 語義分割 (D) 實例分割\n\n"
        "Q2  「用大模型訓練出的知識指導小模型」是哪種壓縮方法？\n"
        "    (A) 量化 (B) 剪枝 (C) 知識蒸餾 (D) 降維\n\n"
        "Q3  將模型從 FP32 轉為 INT8 是哪種壓縮方法？\n"
        "    (A) 知識蒸餾 (B) 剪枝 (C) 量化 (D) 正則化",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "答案：Q1(D) 分個體=實例分割 \u00b7 Q2(C) 大教小=蒸餾 \u00b7 Q3(C) 降精度=量化",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ===== S20 . EXAM_DRILL . 合 — 多模態 + Transformer =====
    s = _blank(prs)
    add_title(s, "Exam Drill \u00b7 多模態融合 + Transformer")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "中級壓軸題型",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(0.8), Inches(1.9), Inches(11.5), Inches(4.0),
        "Q1  Early Fusion 是在 pipeline 的哪個位置合併？\n"
        "    (A) 輸入層 (B) 特徵層 (C) 決策層 (D) 輸出層\n\n"
        "Q2  Transformer 的核心機制是？\n"
        "    (A) 卷積 (B) 循環 (C) Self-Attention (D) 池化\n\n"
        "Q3  醫療 AI 系統中，X光影像應選用什麼技術處理？\n"
        "    (A) BERT (B) Word2Vec (C) ViT (D) TF-IDF",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "答案：Q1(B) 特徵層=Early \u00b7 Q2(C) Self-Attention \u00b7 Q3(C) 影像=ViT",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ===== S21 . PRACTICE_PROMPT . 合 =====
    s = _blank(prs)
    add_title(s, "Practice \u00b7 綜合情境題")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "5 分鐘 \u00b7 這就是一道完整的中級綜合題",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.0), Inches(2.0), Inches(11.0), Inches(4.0),
        "情境：一家醫院要建 AI 輔助診斷系統，\n"
        "輸入包含 X光影像、病歷文字、血液檢驗數值。\n\n"
        "請回答：\n"
        "(1) 每種輸入該用什麼技術處理？\n\n"
        "(2) 融合策略選哪種？為什麼？\n\n"
        "(3) 如果要部署到邊緣裝置，該用什麼壓縮技術？",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.3), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "這題把 CV、NLP、多模態融合、模型壓縮全串在一起。",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 21, N_CONTENT)

    # ===== S22 . SILENT . 合 — 結語 =====
    s = _blank(prs)
    draw_silent_page(
        s,
        "技術不是越新越好，是選對的那一個。\n"
        "NLP、CV、GenAI、多模態\u2014\u2014\n"
        "考的不是你背了多少名詞，\n"
        "是你能不能在情境中選出最適合的技術。\n"
        "下一章：AI 導入規劃與系統部署。",
    )
    add_footer(s, MODULE_CODE, 22, N_CONTENT, dark_bg=True)

    # == Copyright ==
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
