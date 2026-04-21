"""Editorial-strict shape primitives. All slide builders compose these."""
from typing import Optional, Sequence
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

from . import theme as T


# ---------- low-level helpers ----------

def set_solid_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_no_fill(shape):
    shape.fill.background()


def set_line(shape, color: RGBColor, width_pt: float = 1.0, dash: bool = False):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)
    if dash:
        ln = shape.line._get_or_add_ln()
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dash")


def set_no_line(shape):
    shape.line.fill.background()


def set_text_font(run, size=None, color=None, bold=None, family=None, italic=None):
    if size is not None:
        run.font.size = size
    if color is not None:
        run.font.color.rgb = color
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if family is not None:
        run.font.name = family
    else:
        run.font.name = T.FONT_FAMILY


def _apply_para(para, font_size, color, bold, align, family, line_spacing):
    if align is not None:
        para.alignment = align
    if line_spacing is not None:
        para.line_spacing = line_spacing
    for run in para.runs:
        set_text_font(run, size=font_size, color=color, bold=bold, family=family)


def add_rect(slide, x, y, w, h, shape_type=MSO_SHAPE.RECTANGLE):
    return slide.shapes.add_shape(shape_type, x, y, w, h)


def add_textbox(slide, x, y, w, h, text,
                font_size=None, color=None, bold=False, italic=False,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                family=None, line_spacing=None, word_wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        set_text_font(run, size=font_size, color=color, bold=bold,
                      italic=italic, family=family)
        if align is not None:
            p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
    return tb


def add_title(slide, text, y=None, color=None, size=None):
    y = y if y is not None else Inches(0.4)
    color = color if color is not None else T.PRIMARY
    size = size if size is not None else T.FONT_SUBTITLE
    return add_textbox(
        slide, T.MARGIN_X, y, T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.7),
        text, font_size=size, color=color, bold=True,
        align=PP_ALIGN.LEFT,
    )


def add_source(slide, text, y=None):
    y = y if y is not None else T.SLIDE_H - Inches(0.7)
    return add_textbox(
        slide, T.MARGIN_X, y, T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.25),
        f"Source: {text}" if not text.lower().startswith("source") else text,
        font_size=T.FONT_SOURCE, color=T.GRAY_MID, align=PP_ALIGN.RIGHT, bold=False,
    )


# ---------- composite builders ----------

def draw_silent_page(slide, thesis: str):
    """SILENT: dark green full bleed, white hero statement centered."""
    bg = add_rect(slide, 0, 0, T.SLIDE_W, T.SLIDE_H)
    set_solid_fill(bg, T.PRIMARY)
    set_no_line(bg)

    add_textbox(
        slide, Inches(1.0), Inches(2.8), T.SLIDE_W - Inches(2.0), Inches(2.0),
        thesis,
        font_size=T.FONT_HERO, color=T.WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.3,
    )

    # Thin accent line under thesis
    line = add_rect(slide, Inches(6.1), Inches(4.9), Inches(1.1), Inches(0.015))
    set_solid_fill(line, T.WHITE)
    set_no_line(line)


def draw_ask_page(slide, question: str, data_card: Optional[dict] = None):
    """ASK: white background, big Socratic question top-third, data card bottom-right."""
    # Background white — default, no action needed

    add_textbox(
        slide, Inches(1.0), Inches(1.6), T.SLIDE_W - Inches(2.0), Inches(2.0),
        question,
        font_size=Pt(30), color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.35,
    )

    if data_card:
        card_w = Inches(3.2)
        card_h = Inches(2.0)
        cx = T.SLIDE_W - card_w - Inches(0.8)
        cy = Inches(4.7)
        card = add_rect(slide, cx, cy, card_w, card_h)
        set_no_fill(card)
        set_line(card, T.PRIMARY, 1.0)

        add_textbox(
            slide, cx + Inches(0.2), cy + Inches(0.15), card_w - Inches(0.4), Inches(0.3),
            data_card.get("label", ""),
            font_size=T.FONT_BODY, color=T.CHARCOAL, bold=False,
        )
        add_textbox(
            slide, cx + Inches(0.2), cy + Inches(0.5), card_w - Inches(0.4), Inches(0.9),
            data_card.get("stat", ""),
            font_size=Pt(44), color=T.PRIMARY, bold=True,
        )
        add_textbox(
            slide, cx + Inches(0.2), cy + Inches(1.4), card_w - Inches(0.4), Inches(0.5),
            data_card.get("caption", ""),
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            line_spacing=1.2,
        )


def draw_matrix(slide, rows: int, cols: int, cells: Sequence[dict],
                title: str = "", top: float = 1.2, bottom: float = 0.9):
    """
    cells: list of dicts with keys: text, sub (optional), highlight (bool)
           Ordered row-major (left-to-right, top-to-bottom).
    """
    if title:
        add_title(slide, title)

    top_in = Inches(top)
    bottom_in = Inches(bottom)
    total_h = T.SLIDE_H - top_in - bottom_in
    total_w = T.SLIDE_W - 2 * T.MARGIN_X

    gap = Inches(0.1)
    cell_w = (total_w - gap * (cols - 1)) / cols
    cell_h = (total_h - gap * (rows - 1)) / rows

    for idx, cell in enumerate(cells):
        r, c = divmod(idx, cols)
        x = T.MARGIN_X + c * (cell_w + gap)
        y = top_in + r * (cell_h + gap)
        rect = add_rect(slide, x, y, cell_w, cell_h)
        if cell.get("highlight"):
            set_solid_fill(rect, T.PRIMARY)
            set_no_line(rect)
            txt_color = T.WHITE
            sub_color = T.LIGHT_GRAY
        else:
            set_no_fill(rect)
            set_line(rect, T.PRIMARY, 1.0)
            txt_color = T.PRIMARY
            sub_color = T.CHARCOAL

        add_textbox(
            slide, x + Inches(0.2), y + Inches(0.2),
            cell_w - Inches(0.4), Inches(0.5),
            cell.get("text", ""),
            font_size=T.FONT_BODY, color=txt_color, bold=True,
            anchor=MSO_ANCHOR.TOP,
        )
        if cell.get("sub"):
            add_textbox(
                slide, x + Inches(0.2), y + Inches(0.75),
                cell_w - Inches(0.4), cell_h - Inches(0.95),
                cell["sub"],
                font_size=T.FONT_SMALL, color=sub_color, bold=False,
                line_spacing=1.3,
            )


def draw_editorial_table(slide, header: Sequence[str], rows: Sequence[Sequence[str]],
                         title: str = "", top: float = 1.2,
                         col_widths: Optional[Sequence[float]] = None):
    if title:
        add_title(slide, title)

    top_in = Inches(top)
    total_w = T.SLIDE_W - 2 * T.MARGIN_X
    ncols = len(header)
    if col_widths is None:
        col_widths_in = [total_w / ncols for _ in range(ncols)]
    else:
        total_ratio = sum(col_widths)
        col_widths_in = [total_w * (w / total_ratio) for w in col_widths]

    header_h = Inches(0.5)
    row_h = Inches(0.55)

    # Top border line
    top_line = add_rect(slide, T.MARGIN_X, top_in, total_w, Inches(0.02))
    set_solid_fill(top_line, T.PRIMARY)
    set_no_line(top_line)

    # Header band
    hdr_y = top_in + Inches(0.02)
    hdr_band = add_rect(slide, T.MARGIN_X, hdr_y, total_w, header_h)
    set_solid_fill(hdr_band, T.PRIMARY)
    set_no_line(hdr_band)

    x = T.MARGIN_X
    for i, h in enumerate(header):
        add_textbox(
            slide, x + Inches(0.1), hdr_y, col_widths_in[i] - Inches(0.2), header_h,
            h, font_size=T.FONT_BODY, color=T.WHITE, bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        x += col_widths_in[i]

    # Data rows
    cy = hdr_y + header_h
    for r, row in enumerate(rows):
        if r % 2 == 1:
            band = add_rect(slide, T.MARGIN_X, cy, total_w, row_h)
            set_solid_fill(band, T.TABLE_ALT)
            set_no_line(band)
        x = T.MARGIN_X
        for i, val in enumerate(row):
            add_textbox(
                slide, x + Inches(0.12), cy, col_widths_in[i] - Inches(0.24), row_h,
                val, font_size=T.FONT_CAPTION, color=T.CHARCOAL,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            x += col_widths_in[i]
        cy += row_h

    # Bottom border
    bot_line = add_rect(slide, T.MARGIN_X, cy, total_w, Inches(0.02))
    set_solid_fill(bot_line, T.PRIMARY)
    set_no_line(bot_line)
    return cy + Inches(0.05)


def draw_flow_chain(slide, nodes: Sequence[dict], title: str = "",
                    y: float = 3.0, branch: Optional[dict] = None):
    """
    nodes: list of dicts {label, sub, highlight}
    branch (optional): {from_index, label, sub, above: bool}
    """
    if title:
        add_title(slide, title)

    n = len(nodes)
    total_w = T.SLIDE_W - 2 * T.MARGIN_X
    gap = Inches(0.5)
    node_w = (total_w - gap * (n - 1)) / n
    node_h = Inches(1.0)
    y_in = Inches(y)

    centers = []
    for i, node in enumerate(nodes):
        x = T.MARGIN_X + i * (node_w + gap)
        rect = add_rect(slide, x, y_in, node_w, node_h)
        if node.get("highlight"):
            set_solid_fill(rect, T.PRIMARY)
            set_no_line(rect)
            txt_color = T.WHITE
            sub_color = T.LIGHT_GRAY
        else:
            set_no_fill(rect)
            set_line(rect, T.PRIMARY, 1.0)
            txt_color = T.PRIMARY
            sub_color = T.CHARCOAL

        add_textbox(
            slide, x, y_in + Inches(0.15), node_w, Inches(0.4),
            node.get("label", ""),
            font_size=T.FONT_BODY, color=txt_color, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        if node.get("sub"):
            add_textbox(
                slide, x, y_in + Inches(0.55), node_w, Inches(0.4),
                node["sub"],
                font_size=T.FONT_SMALL, color=sub_color,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )

        # Definition/caption below box
        if node.get("caption"):
            add_textbox(
                slide, x, y_in + node_h + Inches(0.1), node_w, Inches(0.4),
                node["caption"],
                font_size=T.FONT_SMALL, color=T.CHARCOAL,
                align=PP_ALIGN.CENTER,
            )

        centers.append((x + node_w, y_in + node_h / 2))

    # Arrows between nodes
    for i in range(n - 1):
        x1, yc = centers[i]
        x2 = T.MARGIN_X + (i + 1) * (node_w + gap)
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, yc, x2, yc)
        set_line(conn, T.PRIMARY, 1.5)
        # Add arrow head via XML
        ln = conn.line._get_or_add_ln()
        tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
        tailEnd.set("type", "triangle")
        tailEnd.set("w", "med")
        tailEnd.set("h", "med")

    # Branch node
    if branch:
        idx = branch["from_index"]
        src_x = T.MARGIN_X + idx * (node_w + gap) + node_w / 2
        src_y = y_in if branch.get("above") else y_in + node_h
        b_y = y_in - Inches(1.4) if branch.get("above") else y_in + Inches(1.5)
        b_x = src_x - node_w / 2
        # Dashed connector
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, src_x, src_y,
                                           src_x, b_y + (node_h / 2 if branch.get("above") else 0))
        set_line(conn, T.GRAY_MID, 1.0, dash=True)

        brect = add_rect(slide, b_x, b_y, node_w, Inches(0.7))
        set_no_fill(brect)
        set_line(brect, T.GRAY_MID, 1.0, dash=True)
        add_textbox(
            slide, b_x, b_y, node_w, Inches(0.7),
            branch.get("label", ""),
            font_size=T.FONT_BODY, color=T.CHARCOAL, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        if branch.get("sub"):
            add_textbox(
                slide, b_x, b_y + Inches(0.75), node_w, Inches(0.3),
                branch["sub"],
                font_size=T.FONT_SMALL, color=T.GRAY_MID,
                align=PP_ALIGN.CENTER,
            )


def draw_pyramid_stack(slide, layers: Sequence[dict], thesis: str = "",
                       title: str = "", cross_cuts: Optional[Sequence[str]] = None):
    """
    Four-layer stacked rectangles + optional left/right vertical cross-cut bands.
    layers: [{name, caption}] top-to-bottom
    cross_cuts: [left_label, right_label]
    """
    if title:
        add_title(slide, title)

    # Layout: stack occupies middle 55% width, center
    stack_w = Inches(5.5)
    stack_x = (T.SLIDE_W - stack_w) / 2
    stack_y = Inches(1.4)
    layer_h = Inches(0.7)

    # Widths vary: top thin, bottom thick — but keep rectangles equal for clarity
    widths = [stack_w * 0.6, stack_w * 0.72, stack_w * 0.86, stack_w * 1.0]

    # Reserve space for cross-cut bands on both sides (§10.2 Evaluation/Governance)
    cc_reserve = Inches(0.9) if "__needs_cc_space__" else Inches(0)
    # Side-caption starts AFTER the right cross-cut band to avoid overlap
    side_cap_x = stack_x + stack_w + Inches(1.1)
    side_cap_w = T.SLIDE_W - side_cap_x - T.MARGIN_X

    for i, layer in enumerate(layers):
        w = widths[i]
        x = stack_x + (stack_w - w) / 2
        y = stack_y + i * layer_h
        rect = add_rect(slide, x, y, w, layer_h)
        set_solid_fill(rect, T.PRIMARY)
        set_no_line(rect)

        add_textbox(
            slide, x, y, w, layer_h,
            layer.get("name", ""),
            font_size=T.FONT_BODY, color=T.WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        if layer.get("caption"):
            add_textbox(
                slide, side_cap_x, y,
                side_cap_w, layer_h,
                layer["caption"],
                font_size=T.FONT_SMALL, color=T.CHARCOAL,
                anchor=MSO_ANCHOR.MIDDLE,
            )

    total_h = layer_h * len(layers)

    # Cross cuts — widened so vertical labels don't wrap
    if cross_cuts:
        cc_w = Inches(0.8)
        left_x_cc = stack_x - cc_w - Inches(0.15)
        right_x_cc = stack_x + stack_w + Inches(0.15)
        for cc_x, label in zip([left_x_cc, right_x_cc], cross_cuts):
            rect = add_rect(slide, cc_x, stack_y, cc_w, total_h)
            set_no_fill(rect)
            set_line(rect, T.PRIMARY, 1.0)
            # Render label as stacked single characters to avoid horizontal wrap
            stacked = "\n".join(list(label))
            add_textbox(
                slide, cc_x, stack_y, cc_w, total_h,
                stacked,
                font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=1.2,
            )

    if thesis:
        draw_inverted_thesis_box(slide, thesis, y=6.0)


def draw_inverted_thesis_box(slide, text: str, y: float = 6.0, width: float = 10.0):
    """G11 — dark solid box with white text, bottom of slide."""
    w = Inches(width)
    x = (T.SLIDE_W - w) / 2
    y_in = Inches(y)
    h = Inches(0.8)
    rect = add_rect(slide, x, y_in, w, h)
    set_solid_fill(rect, T.PRIMARY)
    set_no_line(rect)
    add_textbox(
        slide, x + Inches(0.3), y_in, w - Inches(0.6), h,
        text,
        font_size=T.FONT_BODY, color=T.WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.3,
    )


def draw_vs_two_col(slide, left_title: str, right_title: str,
                    left_items: Sequence[str], right_items: Sequence[str],
                    title: str = "", summary: str = "",
                    delta: str = ""):
    if title:
        add_title(slide, title)

    col_w = Inches(5.5)
    gap_x = Inches(0.9)
    total = col_w * 2 + gap_x
    left_x = (T.SLIDE_W - total) / 2
    right_x = left_x + col_w + gap_x
    top = Inches(1.5)
    col_h = Inches(4.3)

    # Center VS + optional delta badge (Guidance layer)
    vs_w = Inches(0.9)
    vs_x = left_x + col_w + (gap_x - vs_w) / 2
    add_textbox(
        slide, vs_x, top + col_h / 2 - Inches(0.7), vs_w, Inches(0.5),
        "VS",
        font_size=Pt(26), color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    for x, t_text, items in [(left_x, left_title, left_items),
                              (right_x, right_title, right_items)]:
        hdr = add_rect(slide, x, top, col_w, Inches(0.5))
        set_solid_fill(hdr, T.PRIMARY)
        set_no_line(hdr)
        add_textbox(
            slide, x, top, col_w, Inches(0.5),
            t_text,
            font_size=T.FONT_BODY, color=T.WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        body = add_rect(slide, x, top + Inches(0.5), col_w, col_h - Inches(0.5))
        set_no_fill(body)
        set_line(body, T.PRIMARY, 1.0)

        # Items
        item_text = "\n".join(f"• {it}" for it in items)
        add_textbox(
            slide, x + Inches(0.25), top + Inches(0.7),
            col_w - Inches(0.5), col_h - Inches(0.8),
            item_text,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            line_spacing=1.5,
        )

    # Delta badge in the center gap (Guidance layer)
    if delta:
        badge_w = Inches(1.4)
        badge_h = Inches(0.8)
        badge_x = left_x + col_w + (gap_x - badge_w) / 2
        badge_y = top + col_h / 2 - Inches(0.1)
        draw_delta_badge(slide, badge_x, badge_y, delta,
                         w=badge_w, h=badge_h, inverted=True, pill=True)

    if summary:
        add_textbox(
            slide, T.MARGIN_X, top + col_h + Inches(0.15),
            T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
            summary,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            align=PP_ALIGN.CENTER,
        )


def draw_risk_mitigation(slide, risks: Sequence[str], mitigations: Sequence[str],
                         title: str = "", summary: str = "",
                         risks_title: str = "Risks（真風險）",
                         miti_title: str = "Mitigations（課程設計）"):
    if title:
        add_title(slide, title)

    col_w = Inches(5.5)
    gap_x = Inches(0.4)
    total = col_w * 2 + gap_x
    left_x = (T.SLIDE_W - total) / 2
    right_x = left_x + col_w + gap_x
    top = Inches(1.5)
    col_h = Inches(4.3)

    for x, t_text, items in [(left_x, risks_title, risks),
                              (right_x, miti_title, mitigations)]:
        outer = add_rect(slide, x, top, col_w, col_h)
        set_no_fill(outer)
        set_line(outer, T.PRIMARY, 1.5)

        hdr = add_rect(slide, x, top, col_w, Inches(0.5))
        set_solid_fill(hdr, T.PRIMARY)
        set_no_line(hdr)
        add_textbox(
            slide, x, top, col_w, Inches(0.5),
            t_text,
            font_size=T.FONT_BODY, color=T.WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        item_text = "\n".join(f"• {it}" for it in items)
        add_textbox(
            slide, x + Inches(0.25), top + Inches(0.7),
            col_w - Inches(0.5), col_h - Inches(0.9),
            item_text,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            line_spacing=1.55,
        )

    if summary:
        add_textbox(
            slide, T.MARGIN_X, top + col_h + Inches(0.2),
            T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
            summary,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            align=PP_ALIGN.CENTER, italic=False,
        )


def draw_dual_track(slide, track_a_label: str, track_b_label: str,
                    a_nodes: Sequence[str], b_nodes: Sequence[str],
                    bridges: Sequence[tuple],
                    terminal_label: str = "",
                    title: str = ""):
    """Two parallel tracks with nodes; bridges are (a_idx, b_idx, label)."""
    if title:
        add_title(slide, title)

    top = Inches(1.5)
    bottom = Inches(5.5)
    track_a_y = top + Inches(0.3)
    track_b_y = bottom - Inches(0.3)

    left_pad = Inches(2.3)
    right_pad = Inches(2.3)
    start_x = left_pad
    end_x = T.SLIDE_W - right_pad

    # Two horizontal lines
    for y in [track_a_y, track_b_y]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, y, end_x, y)
        set_line(line, T.PRIMARY, 2.0)

    # Labels
    add_textbox(
        slide, Inches(0.4), track_a_y - Inches(0.2), left_pad - Inches(0.4), Inches(0.4),
        track_a_label,
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide, Inches(0.4), track_b_y - Inches(0.2), left_pad - Inches(0.4), Inches(0.4),
        track_b_label,
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
    )

    n = max(len(a_nodes), len(b_nodes))
    span = end_x - start_x
    xs = [start_x + int(span * (i / (n - 1))) for i in range(n)]

    # Nodes on track A
    for i, label in enumerate(a_nodes):
        cx = xs[i]
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Emu(60000), track_a_y - Emu(60000),
                                      Emu(120000), Emu(120000))
        set_solid_fill(dot, T.PRIMARY)
        set_no_line(dot)
        add_textbox(
            slide, cx - Inches(0.9), track_a_y - Inches(0.7), Inches(1.8), Inches(0.4),
            label,
            font_size=T.FONT_SMALL, color=T.CHARCOAL,
            align=PP_ALIGN.CENTER,
        )

    for i, label in enumerate(b_nodes):
        cx = xs[i]
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Emu(60000), track_b_y - Emu(60000),
                                      Emu(120000), Emu(120000))
        set_solid_fill(dot, T.PRIMARY)
        set_no_line(dot)
        add_textbox(
            slide, cx - Inches(0.9), track_b_y + Inches(0.25), Inches(1.8), Inches(0.4),
            label,
            font_size=T.FONT_SMALL, color=T.CHARCOAL,
            align=PP_ALIGN.CENTER,
        )

    # Bridges
    for a_i, b_i, blabel in bridges:
        cx_a = xs[a_i]
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, cx_a, track_a_y, cx_a, track_b_y)
        set_line(conn, T.LIGHT_GRAY, 1.0, dash=True)
        if blabel:
            add_textbox(
                slide, cx_a - Inches(0.9), (track_a_y + track_b_y) / 2 - Inches(0.2),
                Inches(1.8), Inches(0.4),
                blabel,
                font_size=T.FONT_SOURCE, color=T.GRAY_MID,
                align=PP_ALIGN.CENTER,
            )

    # Terminal box
    if terminal_label:
        tw = Inches(2.0)
        th = Inches(0.8)
        tx = end_x + Inches(0.15)
        ty = (track_a_y + track_b_y) / 2 - th / 2
        rect = add_rect(slide, tx, ty, tw, th)
        set_solid_fill(rect, T.PRIMARY)
        set_no_line(rect)
        add_textbox(
            slide, tx, ty, tw, th,
            terminal_label,
            font_size=T.FONT_SMALL, color=T.WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )


def draw_grid(slide, rows: int, cols: int, cells: Sequence[dict],
              title: str = "", top: float = 1.2, bottom: float = 0.9,
              caption: str = ""):
    """Like matrix but for S10-style milestone grid; cell can have label + sub + highlight + dashed."""
    if title:
        add_title(slide, title)

    top_in = Inches(top)
    bottom_in = Inches(bottom)
    total_h = T.SLIDE_H - top_in - bottom_in - (Inches(0.5) if caption else Inches(0))
    total_w = T.SLIDE_W - 2 * T.MARGIN_X

    gap = Inches(0.15)
    cell_w = (total_w - gap * (cols - 1)) / cols
    cell_h = (total_h - gap * (rows - 1)) / rows

    for idx, cell in enumerate(cells):
        r, c = divmod(idx, cols)
        x = T.MARGIN_X + c * (cell_w + gap)
        y = top_in + r * (cell_h + gap)
        rect = add_rect(slide, x, y, cell_w, cell_h)

        if cell.get("highlight"):
            set_solid_fill(rect, T.PRIMARY)
            set_no_line(rect)
            label_color = T.WHITE
            sub_color = T.LIGHT_GRAY
        else:
            set_no_fill(rect)
            set_line(rect, T.CHARCOAL if cell.get("dashed") else T.PRIMARY, 1.0,
                     dash=cell.get("dashed", False))
            label_color = T.PRIMARY
            sub_color = T.CHARCOAL

        add_textbox(
            slide, x, y + Inches(0.1), cell_w, Inches(0.5),
            cell.get("label", ""),
            font_size=T.FONT_BODY, color=label_color, bold=True,
            align=PP_ALIGN.CENTER,
        )
        if cell.get("sub"):
            add_textbox(
                slide, x, y + Inches(0.55), cell_w, Inches(0.35),
                cell["sub"],
                font_size=T.FONT_SMALL, color=sub_color,
                align=PP_ALIGN.CENTER,
            )
        if cell.get("note"):
            add_textbox(
                slide, x, y + cell_h - Inches(0.5), cell_w, Inches(0.4),
                cell["note"],
                font_size=T.FONT_SOURCE, color=sub_color, italic=True,
                align=PP_ALIGN.CENTER,
            )

    if caption:
        add_textbox(
            slide, T.MARGIN_X, T.SLIDE_H - bottom_in + Inches(0.05),
            T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
            caption,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            align=PP_ALIGN.CENTER,
        )


def draw_code_panel(slide, x, y, w, h, label, code, bullets,
                     label_dark: bool = True):
    """Code comparison panel for BEFORE/AFTER slides.
    Left side: label + code block; Right side: bullet notes."""
    code_w = w * 0.60
    note_w = w * 0.36
    gap = w * 0.04

    # Label tag
    tag_h = Inches(0.35)
    tag = add_rect(slide, x, y, code_w, tag_h)
    set_solid_fill(tag, T.PRIMARY if label_dark else T.LIGHT_GRAY)
    set_no_line(tag)
    add_textbox(
        slide, x + Inches(0.15), y, code_w - Inches(0.3), tag_h,
        label,
        font_size=T.FONT_CAPTION,
        color=T.WHITE if label_dark else T.CHARCOAL,
        bold=True, anchor=MSO_ANCHOR.MIDDLE,
    )

    # Code box
    code_y = y + tag_h
    code_h = h - tag_h
    code_box = add_rect(slide, x, code_y, code_w, code_h)
    set_solid_fill(code_box, RGBColor(0xF7, 0xF7, 0xF7))
    set_line(code_box, T.PRIMARY, 1.0)
    code_tb = slide.shapes.add_textbox(
        x + Inches(0.15), code_y + Inches(0.1),
        code_w - Inches(0.3), code_h - Inches(0.2)
    )
    tf = code_tb.text_frame
    tf.word_wrap = True
    lines = code.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = T.FONT_MONO
        run.font.size = Pt(11)
        run.font.color.rgb = T.CHARCOAL
        p.line_spacing = 1.2

    # Bullet notes
    bx = x + code_w + gap
    bullet_text = "\n".join(f"• {b}" for b in bullets)
    add_textbox(
        slide, bx, code_y, note_w, code_h,
        bullet_text,
        font_size=T.FONT_SMALL,
        color=T.PRIMARY if label_dark else T.CHARCOAL,
        line_spacing=1.45, anchor=MSO_ANCHOR.TOP,
    )


def draw_split_panel(slide, left: dict, right: dict, title: str = ""):
    """Two-column split panel for GEOMETRIC comparison (e.g. wrong vs right mental model).
    left/right dicts: {label, cells: [{box_text, lines: [...]}, ...], note, style: 'wrong'|'right'}
    Simplified structural rendering without full fidelity."""
    if title:
        add_title(slide, title)

    col_w = (T.SLIDE_W - 2 * T.MARGIN_X - Inches(0.3)) / 2
    top = Inches(1.4)
    col_h = Inches(4.8)

    left_x = T.MARGIN_X
    right_x = T.MARGIN_X + col_w + Inches(0.3)

    # Divider
    div = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        T.MARGIN_X + col_w + Inches(0.15), top,
        T.MARGIN_X + col_w + Inches(0.15), top + col_h,
    )
    set_line(div, T.PRIMARY, 1.0)

    for x, col in [(left_x, left), (right_x, right)]:
        # Section label
        label_color = T.GRAY_MID if col.get("style") == "wrong" else T.PRIMARY
        add_textbox(
            slide, x, top, col_w, Inches(0.4),
            col.get("label", ""),
            font_size=T.FONT_BODY, color=label_color, bold=True,
        )

        # Cells
        y = top + Inches(0.55)
        for cell in col.get("cells", []):
            box_w = Inches(2.5)
            box_h = Inches(0.7)
            cx = x + (col_w - box_w) / 2
            rect = add_rect(slide, cx, y, box_w, box_h)
            if col.get("style") == "wrong":
                set_no_fill(rect)
                set_line(rect, T.LIGHT_GRAY, 1.0, dash=True)
                t_color = T.CHARCOAL
            elif cell.get("highlight"):
                set_solid_fill(rect, T.PRIMARY)
                set_no_line(rect)
                t_color = T.WHITE
            else:
                set_no_fill(rect)
                set_line(rect, T.PRIMARY, 1.0)
                t_color = T.PRIMARY
            add_textbox(
                slide, cx, y, box_w, box_h,
                cell.get("text", ""),
                font_size=T.FONT_SMALL, color=t_color, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )
            y += box_h + Inches(0.2)

        if col.get("note"):
            add_textbox(
                slide, x, top + col_h - Inches(0.5), col_w, Inches(0.4),
                col["note"],
                font_size=T.FONT_SMALL, color=label_color, italic=True,
                align=PP_ALIGN.CENTER,
            )


def draw_concentric_zones(slide, zones: Sequence[dict], title: str = ""):
    """Simplified concentric structure: three horizontal bands (not true circles)."""
    if title:
        add_title(slide, title)

    top = Inches(1.3)
    total_h = Inches(4.6)
    band_h = total_h / len(zones)

    for i, zone in enumerate(zones):
        y = top + i * band_h
        w = T.SLIDE_W - 2 * T.MARGIN_X
        inset = Inches(i * 0.0)  # Could indent, but keep flat
        bx = T.MARGIN_X + inset
        bw = w - 2 * inset
        rect = add_rect(slide, bx, y + Inches(0.05), bw, band_h - Inches(0.1))
        if zone.get("highlight"):
            set_solid_fill(rect, T.PRIMARY)
            set_no_line(rect)
            t_color = T.WHITE
            sub_color = T.LIGHT_GRAY
        else:
            set_no_fill(rect)
            set_line(rect, T.PRIMARY, 1.5,
                     dash=bool(zone.get("dashed", False)))
            t_color = T.PRIMARY
            sub_color = T.CHARCOAL

        add_textbox(
            slide, bx + Inches(0.3), y + Inches(0.1),
            Inches(3.0), band_h - Inches(0.2),
            zone.get("label", ""),
            font_size=T.FONT_BODY, color=t_color, bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide, bx + Inches(3.5), y + Inches(0.1),
            bw - Inches(3.8), band_h - Inches(0.2),
            zone.get("items", ""),
            font_size=T.FONT_SMALL, color=sub_color,
            anchor=MSO_ANCHOR.MIDDLE,
        )


def draw_thesis_hierarchy(slide, blocks: Sequence[dict], title: str = "",
                           thesis: str = ""):
    """Two-column hierarchical bullet layout for M1-S15 closing PYRAMID."""
    if title:
        add_title(slide, title)

    col_w = (T.SLIDE_W - 2 * T.MARGIN_X - Inches(0.5)) / 2
    top = Inches(1.4)
    col_h = Inches(4.3)

    for i, block in enumerate(blocks[:2]):
        x = T.MARGIN_X + i * (col_w + Inches(0.5))
        # Title band
        band = add_rect(slide, x, top, col_w, Inches(0.5))
        set_solid_fill(band, T.PRIMARY)
        set_no_line(band)
        add_textbox(
            slide, x, top, col_w, Inches(0.5),
            block.get("heading", ""),
            font_size=T.FONT_BODY, color=T.WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        body_y = top + Inches(0.6)
        item_text = "\n\n".join(f"• {it}" for it in block.get("items", []))
        add_textbox(
            slide, x + Inches(0.25), body_y,
            col_w - Inches(0.5), col_h - Inches(0.7),
            item_text,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            line_spacing=1.4,
        )

    if thesis:
        draw_inverted_thesis_box(slide, thesis, y=6.0, width=9.0)


def draw_three_blocks_flow(slide, blocks: Sequence[dict], title: str = "",
                           bottom_note: str = ""):
    """Three horizontal blocks for S6-style 'control flow three building blocks'.
    Each block: {heading, items: [str, ...]} where items are steps."""
    if title:
        add_title(slide, title)

    n = len(blocks)
    total_w = T.SLIDE_W - 2 * T.MARGIN_X
    gap = Inches(0.25)
    block_w = (total_w - gap * (n - 1)) / n
    top = Inches(1.5)
    block_h = Inches(4.3)

    for i, block in enumerate(blocks):
        x = T.MARGIN_X + i * (block_w + gap)
        outer = add_rect(slide, x, top, block_w, block_h)
        set_no_fill(outer)
        set_line(outer, T.PRIMARY, 1.0, dash=True)

        # Heading
        add_textbox(
            slide, x, top + Inches(0.15), block_w, Inches(0.4),
            block.get("heading", ""),
            font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
            align=PP_ALIGN.CENTER,
        )

        # Items rendered as vertical chain of small boxes
        items = block.get("items", [])
        step_h = Inches(0.55)
        cy = top + Inches(0.7)
        for j, item in enumerate(items):
            bw = block_w - Inches(0.5)
            bx = x + Inches(0.25)
            rect = add_rect(slide, bx, cy, bw, step_h)
            set_no_fill(rect)
            set_line(rect, T.PRIMARY, 1.0)
            add_textbox(
                slide, bx, cy, bw, step_h,
                item,
                font_size=T.FONT_SMALL, color=T.CHARCOAL,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                family=T.FONT_MONO if item.startswith("`") or "=" in item else None,
            )
            cy += step_h + Inches(0.15)

    if bottom_note:
        add_textbox(
            slide, T.MARGIN_X, top + block_h + Inches(0.15),
            T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
            bottom_note,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            align=PP_ALIGN.CENTER,
        )


# ---------- Guidance Layer primitives (§ 10.7) ----------

def draw_delta_badge(slide, x, y, value: str, *,
                     w=Inches(1.6), h=Inches(0.7),
                     inverted: bool = True, pill: bool = True):
    """Variation / delta callout badge — '+23pp' / '-39%' / '3–5×'.
    inverted=True: dark green fill, white text (primary emphasis).
    inverted=False: white fill, dark green border+text (secondary).
    pill=True: rounded corners."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if pill else MSO_SHAPE.RECTANGLE
    badge = slide.shapes.add_shape(shape_type, x, y, w, h)
    if inverted:
        set_solid_fill(badge, T.PRIMARY)
        set_no_line(badge)
        fg = T.WHITE
    else:
        badge.fill.solid()
        badge.fill.fore_color.rgb = T.WHITE
        set_line(badge, T.PRIMARY, 1.5)
        fg = T.PRIMARY

    tf = badge.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = value
    set_text_font(run, size=Pt(22), color=fg, bold=True, family=T.FONT_FAMILY)
    return badge


def draw_emphasis_pill(slide, x, y, w, h, text: str, *,
                       font_size=None, inverted: bool = True):
    """Rounded pill for short call-out text. Similar to delta_badge but
    generalised for words, not just numbers."""
    font_size = font_size or T.FONT_BODY
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if inverted:
        set_solid_fill(pill, T.PRIMARY)
        set_no_line(pill)
        fg = T.WHITE
    else:
        pill.fill.solid()
        pill.fill.fore_color.rgb = T.WHITE
        set_line(pill, T.PRIMARY, 1.5)
        fg = T.PRIMARY
    tf = pill.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_text_font(run, size=font_size, color=fg, bold=True)
    return pill


def draw_callout_arrow(slide, from_xy, to_xy, *,
                       note: str = "", note_at: str = "from",
                       dashed: bool = False, width_pt: float = 1.5,
                       note_offset=(Inches(0.1), Inches(-0.3)),
                       note_width=Inches(2.3)):
    """Connector arrow from from_xy → to_xy with optional short note.
    note_at: 'from' (anchor note near start) or 'to' (near end)."""
    fx, fy = from_xy
    tx, ty = to_xy
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, fx, fy, tx, ty)
    set_line(conn, T.PRIMARY, width_pt, dash=dashed)
    # Arrow head
    ln = conn.line._get_or_add_ln()
    tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
    tailEnd.set("type", "triangle")
    tailEnd.set("w", "med")
    tailEnd.set("h", "med")

    if note:
        nx = (fx if note_at == "from" else tx) + note_offset[0]
        ny = (fy if note_at == "from" else ty) + note_offset[1]
        add_textbox(
            slide, nx, ny, note_width, Inches(0.4),
            note,
            font_size=T.FONT_SMALL, color=T.PRIMARY, bold=True,
        )
    return conn


def draw_dividing_rule(slide, *, axis: str = "vertical",
                       at=Inches(6.5), length=Inches(4.5),
                       offset=Inches(1.5), label: str = "",
                       kind: str = "time_split"):
    """Vertical or horizontal dividing dashed line with semantic label.
    kind ∈ time_split / category_split / threshold — label placement convention.
    `at` is the axis position (x for vertical, y for horizontal),
    `offset` is the start position on the other axis, `length` is the span."""
    if axis == "vertical":
        x1 = x2 = at
        y1 = offset
        y2 = offset + length
    else:
        y1 = y2 = at
        x1 = offset
        x2 = offset + length

    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    set_line(conn, T.GRAY_MID, 1.2, dash=True)

    if label:
        if axis == "vertical":
            lx = at - Inches(1.5)
            ly = offset - Inches(0.35)
            lw = Inches(3.0)
        else:
            lx = offset + Inches(0.2)
            ly = at - Inches(0.35)
            lw = Inches(4.0)
        add_textbox(
            slide, lx, ly, lw, Inches(0.3),
            label,
            font_size=T.FONT_SOURCE, color=T.GRAY_MID, italic=True,
            align=PP_ALIGN.CENTER,
        )
    return conn


def draw_highlight_ring(slide, x, y, w, h, *, width_pt: float = 2.0):
    """Draw a non-filled rectangle (or rounded rect) ring around a target area
    to emphasise 'look here'. Caller provides the target's bounding box."""
    ring = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_no_fill(ring)
    set_line(ring, T.PRIMARY, width_pt)
    return ring


def draw_image_placeholder(slide, x, y, w, h, *,
                           slot_name: str = "真圖",
                           description: str = "",
                           url_hint: str = "",
                           placeholder_id: Optional[str] = None,
                           registry: Optional[list] = None,
                           size_hint: str = ""):
    """Reserve a bounding box for a real image that the author will paste in later.

    §7 of 投影片生成 SOP — 統一外觀：
      - 外框：深綠 1.5pt 實線
      - 底色：#F7F7F7 淺灰
      - 中央：「[ 待補真圖：{slot_name} ]」12pt 灰斜體
      - 下方：description（12pt 炭灰）+ url_hint（10pt 灰）+ size_hint
      - 右下角標：「TODO」8pt 深綠小字（editorial 紀律：不使用紅色）

    若傳入 registry (list)，會把佔位資訊追加進去，供 build.py 最後寫出
    `output/_image_placeholders.yaml`。
    """
    # 背景
    bg = add_rect(slide, x, y, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF7, 0xF7, 0xF7)
    set_line(bg, T.PRIMARY, 1.5)

    # 中央標題
    header_h = Inches(0.5)
    add_textbox(
        slide, x, y + h * 0.25, w, header_h,
        f"[ 待補真圖 · {slot_name} ]",
        font_size=T.FONT_BODY, color=T.GRAY_MID, italic=True, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # description
    if description:
        add_textbox(
            slide, x + Inches(0.2), y + h * 0.45,
            w - Inches(0.4), Inches(0.9),
            description,
            font_size=T.FONT_SMALL, color=T.CHARCOAL,
            align=PP_ALIGN.CENTER, line_spacing=1.35,
        )

    # URL hint
    if url_hint:
        add_textbox(
            slide, x + Inches(0.2), y + h * 0.72,
            w - Inches(0.4), Inches(0.35),
            url_hint,
            font_size=T.FONT_SOURCE, color=T.PRIMARY,
            align=PP_ALIGN.CENTER,
        )

    # Size hint
    if size_hint:
        add_textbox(
            slide, x + Inches(0.2), y + h * 0.85,
            w - Inches(0.4), Inches(0.3),
            f"建議尺寸：{size_hint}",
            font_size=T.FONT_SOURCE, color=T.GRAY_MID,
            align=PP_ALIGN.CENTER,
        )

    # TODO 角標（editorial-strict：深綠不用紅色）
    todo_w = Inches(0.5)
    todo_h = Inches(0.22)
    todo_x = x + w - todo_w - Inches(0.08)
    todo_y = y + h - todo_h - Inches(0.08)
    badge = add_rect(slide, todo_x, todo_y, todo_w, todo_h)
    set_solid_fill(badge, T.PRIMARY)
    set_no_line(badge)
    add_textbox(
        slide, todo_x, todo_y, todo_w, todo_h,
        "TODO",
        font_size=T.FONT_SOURCE, color=T.WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # Register for audit
    if registry is not None:
        # Convert Emu to inches for human-readable YAML
        from pptx.util import Emu
        def _in(v): return round(v / 914400, 3)  # EMU per inch
        registry.append({
            "id": placeholder_id or f"ph_{len(registry)+1}",
            "slot_name": slot_name,
            "description": description,
            "source_url": url_hint,
            "recommended_size": size_hint,
            "bounding_box": {"x": _in(x), "y": _in(y),
                             "w": _in(w), "h": _in(h)},
            "status": "pending",
            "asset_path": None,
        })


# ---------- existing (unchanged below) ----------

def draw_photo_placeholder_triptych(slide, items: Sequence[dict],
                                     title: str = "", thesis: str = ""):
    """S14 — three-column placeholder with 工具名 + URL."""
    if title:
        add_title(slide, title)

    n = len(items)
    gap = Inches(0.2)
    total_w = T.SLIDE_W - 2 * T.MARGIN_X
    col_w = (total_w - gap * (n - 1)) / n
    top = Inches(1.5)
    col_h = Inches(4.2)

    for i, item in enumerate(items):
        x = T.MARGIN_X + i * (col_w + gap)
        frame = add_rect(slide, x, top, col_w, col_h)
        set_no_fill(frame)
        set_line(frame, T.PRIMARY, 2.0)

        # Upper half — placeholder text
        add_textbox(
            slide, x, top, col_w, col_h * 0.55,
            "[ 官方文件截圖待補 ]",
            font_size=T.FONT_BODY, color=T.LIGHT_GRAY, italic=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Lower half — tool name + URL
        add_textbox(
            slide, x, top + col_h * 0.55, col_w, Inches(0.45),
            item.get("tool", ""),
            font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, x + Inches(0.1), top + col_h * 0.55 + Inches(0.5),
            col_w - Inches(0.2), Inches(0.5),
            item.get("url", ""),
            font_size=T.FONT_SMALL, color=T.CHARCOAL,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, x + Inches(0.1), top + col_h * 0.55 + Inches(1.05),
            col_w - Inches(0.2), Inches(0.5),
            item.get("depth", ""),
            font_size=T.FONT_SMALL, color=T.GRAY_MID,
            align=PP_ALIGN.CENTER,
        )

    if thesis:
        add_textbox(
            slide, T.MARGIN_X, top + col_h + Inches(0.3),
            T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
            thesis,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
            align=PP_ALIGN.CENTER,
            line_spacing=1.3,
        )
