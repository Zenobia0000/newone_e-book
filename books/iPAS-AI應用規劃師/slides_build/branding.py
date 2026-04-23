"""Cover, footer, copyright page — Uber-inspired black/white design."""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from . import theme as T
from .primitives import (
    add_rect, add_textbox, set_solid_fill, set_no_line, set_text_font
)


def add_cover_slide(prs, module_code: str, module_title: str, subtitle: str,
                    time_min: int, slide_count: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full black background
    bg = add_rect(slide, 0, 0, T.SLIDE_W, T.SLIDE_H)
    set_solid_fill(bg, T.PRIMARY)
    set_no_line(bg)

    # Left-aligned layout — Uber confidence

    # Small logo top-left
    logo_w = T.LOGO_COVER_SIZE
    try:
        slide.shapes.add_picture(
            str(T.LOGO_PATH), T.MARGIN_X, Inches(0.6),
            width=logo_w, height=logo_w)
    except Exception:
        pass  # logo file may not exist

    # Module code — massive, left-aligned
    code_y = Inches(3.0)
    add_textbox(
        slide,
        T.MARGIN_X, code_y, T.SLIDE_W - 2 * T.MARGIN_X, Inches(1.0),
        module_code,
        font_size=Pt(64), color=T.WHITE, align=PP_ALIGN.LEFT, bold=True,
    )

    # Thin white rule
    rule_y = code_y + Inches(1.1)
    rule = add_rect(slide, T.MARGIN_X, rule_y,
                    Inches(4.0), Inches(0.02))
    set_solid_fill(rule, T.WHITE)
    set_no_line(rule)

    # Module title
    add_textbox(
        slide,
        T.MARGIN_X, rule_y + Inches(0.2),
        T.SLIDE_W * 0.7, Inches(0.8),
        module_title,
        font_size=T.FONT_COVER_TITLE, color=T.WHITE,
        align=PP_ALIGN.LEFT, bold=True,
    )

    # Subtitle
    add_textbox(
        slide,
        T.MARGIN_X, rule_y + Inches(1.0),
        T.SLIDE_W * 0.7, Inches(0.5),
        subtitle,
        font_size=T.FONT_COVER_SUB, color=T.GRAY_MID,
        align=PP_ALIGN.LEFT, bold=False,
    )

    # Meta line — bottom left
    meta = f"{time_min} min · {slide_count} slides"
    add_textbox(
        slide,
        T.MARGIN_X, T.SLIDE_H - Inches(1.0),
        Inches(4), Inches(0.3),
        meta,
        font_size=T.FONT_BODY, color=T.GRAY_MID,
        align=PP_ALIGN.LEFT, bold=False,
    )

    # Brand name — bottom left
    add_textbox(
        slide,
        T.MARGIN_X, T.SLIDE_H - Inches(0.55),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        f"{T.BRAND_NAME_ZH} · {T.BRAND_NAME_EN}",
        font_size=T.FONT_SMALL, color=T.GRAY_MID,
        align=PP_ALIGN.LEFT, bold=False,
    )

    return slide


def add_footer(slide, module_code: str, slide_no: int, total: int,
               dark_bg: bool = False):
    """Minimal footer: left pagination, right logo."""
    fg = T.WHITE if dark_bg else T.GRAY_MID

    foot_y = T.SLIDE_H - Inches(0.35)
    add_textbox(
        slide,
        T.MARGIN_X, foot_y, Inches(6), Inches(0.25),
        f"{module_code} · {slide_no}/{total}",
        font_size=T.FONT_SOURCE, color=fg, align=PP_ALIGN.LEFT, bold=False,
    )

    # Right-bottom: logo
    logo_w = T.LOGO_CORNER_SIZE
    logo_x = T.SLIDE_W - logo_w - Inches(0.4)
    logo_y = T.SLIDE_H - logo_w - Inches(0.15)
    try:
        slide.shapes.add_picture(
            str(T.LOGO_PATH), logo_x, logo_y,
            width=logo_w, height=logo_w)
    except Exception:
        pass


def add_copyright_slide(prs, module_code: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Black top band
    band = add_rect(slide, 0, 0, T.SLIDE_W, Inches(0.7))
    set_solid_fill(band, T.PRIMARY)
    set_no_line(band)

    add_textbox(
        slide,
        T.MARGIN_X, Inches(0.15), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        f"Copyright · {module_code}",
        font_size=T.FONT_SUBTITLE, color=T.WHITE,
        align=PP_ALIGN.LEFT, bold=True,
    )

    # Copyright body
    add_textbox(
        slide,
        Inches(1.2), Inches(1.4), T.SLIDE_W - Inches(2.4), Inches(5.2),
        T.COPYRIGHT_FULL,
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        align=PP_ALIGN.LEFT, bold=False,
        line_spacing=1.45,
    )

    # Small logo bottom-right
    logo_w = Inches(0.6)
    try:
        slide.shapes.add_picture(
            str(T.LOGO_PATH),
            T.SLIDE_W - logo_w - Inches(0.6),
            T.SLIDE_H - logo_w - Inches(0.5),
            width=logo_w, height=logo_w,
        )
    except Exception:
        pass

    # Bottom brand bar
    add_textbox(
        slide,
        T.MARGIN_X, T.SLIDE_H - Inches(0.4), Inches(8), Inches(0.3),
        f"{T.BRAND_NAME_ZH} · {T.BRAND_NAME_EN}",
        font_size=T.FONT_SMALL, color=T.GRAY_MID,
        align=PP_ALIGN.LEFT, bold=False,
    )

    return slide
