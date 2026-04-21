"""F3 deck — OOP 核心觀念與實例化
21 content slides + cover + copyright.

Governing thought:
    腳本解決問題，類別解決規模。
    OOP 不是炫技，是把散落的狀態收進一個
    可重用、可測試、可組合的殼。

Aligned to chapters/F3_OOP核心觀念與實例化/00_skeleton.yaml
  · 5 Learning Objectives × 5 Common Pitfalls
  · Teaching-track primitives: ASK / VS / MATRIX / SILENT / CONCEPT-CARD /
    CODE / EXAMPLE-I/O / CHECKPOINT / PITFALL / TABLE / PRACTICE / PYRAMID
  · 2 bridging slides: S18 → F4, S20 → F5
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_rect, add_textbox, add_title, add_source,
    set_solid_fill, set_no_line, set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_vs_two_col, draw_code_panel, draw_thesis_hierarchy,
    draw_flow_chain, draw_image_placeholder, draw_pyramid_stack,
    draw_emphasis_pill, draw_three_blocks_flow,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "F3"
MODULE_TITLE = "OOP 核心觀念與實例化"
MODULE_SUBTITLE = "腳本解決問題，類別解決規模"
TIME_MIN = 90
N_CONTENT = 21

RED_ERROR = RGBColor(0xC6, 0x28, 0x28)
GREEN_OK = RGBColor(0x2E, 0x7D, 0x32)


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _draw_pitfall(slide, title, wrong_label, wrong_code, right_label,
                  right_code, why):
    """PITFALL 頁：左紅錯 / 右綠對 / 下方 why 條。"""
    add_title(slide, title)
    top = Inches(1.4)
    col_h = Inches(4.2)
    col_w = Inches(5.8)
    gap = Inches(0.3)
    total = col_w * 2 + gap
    left_x = (T.SLIDE_W - total) / 2
    right_x = left_x + col_w + gap

    for x, label, code, accent, mark in [
        (left_x, wrong_label, wrong_code, RED_ERROR, "✗"),
        (right_x, right_label, right_code, GREEN_OK, "✓"),
    ]:
        hdr = add_rect(slide, x, top, col_w, Inches(0.5))
        set_solid_fill(hdr, accent)
        set_no_line(hdr)
        add_textbox(
            slide, x + Inches(0.2), top, col_w - Inches(0.4), Inches(0.5),
            f"{mark}  {label}",
            font_size=T.FONT_BODY, color=T.WHITE, bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        body_y = top + Inches(0.5)
        body_h = col_h - Inches(0.5)
        box = add_rect(slide, x, body_y, col_w, body_h)
        set_solid_fill(box, RGBColor(0xF9, 0xF9, 0xF9))
        set_line(box, accent, 1.2)
        add_textbox(
            slide, x + Inches(0.25), body_y + Inches(0.2),
            col_w - Inches(0.5), body_h - Inches(0.4),
            code,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            family=T.FONT_MONO, line_spacing=1.45,
        )

    why_y = top + col_h + Inches(0.2)
    why_box = add_rect(slide, T.MARGIN_X, why_y,
                       T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.55))
    set_solid_fill(why_box, T.PRIMARY)
    set_no_line(why_box)
    add_textbox(
        slide, T.MARGIN_X + Inches(0.3), why_y,
        T.SLIDE_W - 2 * T.MARGIN_X - Inches(0.6), Inches(0.55),
        f"為什麼：{why}",
        font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def build_ch04(output_path, image_registry=None):
    """Build F3 deck; 21 content slides."""
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── S1 · ASK — 複製貼上幾次了 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "為什麼同一段清洗邏輯，\n你已經複製貼上第四次？",
        data_card={
            "label": "資料工程師日常",
            "stat": "61%",
            "caption": "每週維護 ≥ 3 條\n結構相似的管線",
        },
    )
    add_source(s, "Kaggle State of Data Science 2024 (n=26,000)")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # ───────── S2 · VS — 腳本式 vs 類別式 ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="腳本式 (Procedural)",
        right_title="類別式 (Object-Oriented)",
        left_items=[
            "5 個散落的函式",
            "10 個全域變數穿插",
            "狀態住在 global namespace",
            "改一處，擔心另外三處",
            "單元測試很難寫",
        ],
        right_items=[
            "1 個 DataPipeline 類別",
            "狀態收進 self.xxx",
            "行為封裝成方法",
            "實例化 N 次 = 管線複製 N 條",
            "測試只要 pipe = DataPipeline(); assert ...",
        ],
        title="腳本式 vs 類別式：差的不是行數，是狀態的去處",
        summary="兩邊都 100 行，但維護成本差 10 倍——因為狀態的可見性不同。",
        delta="狀態\n的歸屬",
    )
    add_source(s, "F3 課堂歸納 · Gamma et al., Design Patterns §Encapsulation")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── S3 · MATRIX 1×3 — 三大效益 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=1, cols=3, cells=[
        {"text": "可重用 (Reusable)",
         "sub": "清洗邏輯寫一次\n套到 N 個資料集\n實例化 = 複製管線",
         "highlight": True},
        {"text": "可測試 (Testable)",
         "sub": "小單元獨立驗證\npipe.add_step() 可 mock\n單元測試 < 10 行"},
        {"text": "可組合 (Composable)",
         "sub": "管線各步驟可替換\n換資料源不改邏輯\nF5 method chaining 基礎"},
    ], title="資料工程採用 OOP 的三大效益：背後都指向降耦")
    add_source(s, "Fowler, Refactoring §Encapsulate Record · Beck, TDD by Example")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───────── S4 · SILENT — 立論 ─────────
    s = _blank(prs)
    draw_silent_page(s, "腳本解決問題，\n類別解決規模。")
    add_footer(s, MODULE_CODE, 4, N_CONTENT, dark_bg=True)

    # ───────── S5 · CONCEPT-CARD — Class vs Object 比喻 ─────────
    s = _blank(prs)
    add_title(s, "Class 是建築藍圖，Object 是蓋出來的房子")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.0), Inches(0.5),
        "類比（僅用一次）",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(1.95), Inches(6.0), Inches(4.2),
        "一份 class 定義\n可以實例化成多個 object。\n\n"
        "每個 object 各自有\n獨立的屬性值 (instance attribute)，\n"
        "但結構 (方法與屬性名) 由 class 決定。\n\n"
        "→ Class：定義資料與行為的模板\n"
        "→ Object：執行期真正持有狀態的實體\n\n"
        "術語對照：Class / Instance / Attribute / Method",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2), Inches(6.0), Inches(0.5),
        "記憶點：一份 class → N 個獨立 instance。",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True, italic=True,
    )
    draw_image_placeholder(
        s,
        x=Inches(7.0), y=Inches(1.3),
        w=Inches(5.7), h=Inches(5.0),
        slot_name="Class 定義 vs Instance 實體示意",
        description="左：一份 class 原始碼\n右：三個 instance 記憶體方塊\n(同結構，各自 name / steps 獨立值)",
        url_hint="https://docs.python.org/3/tutorial/classes.html",
        placeholder_id="F3_S05_class_vs_instance",
        registry=image_registry,
        size_hint="1400×1200 px",
    )
    add_source(s, "Python Tutorial §9 Classes")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── S6 · CODE — 最小類別 ─────────
    s = _blank(prs)
    add_title(s, "最小類別：三個關鍵字 + 兩個方法，比 5 個散落函式好用")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="class DataPipeline — 一個最小可用雛形",
        code=(
            'class DataPipeline:\n'
            '    def __init__(self, name):\n'
            '        self.name = name       # instance attribute：每個實例獨立\n'
            '        self.steps = []        # 可變物件放 __init__！\n'
            '\n'
            '    def add_step(self, step):\n'
            '        self.steps.append(step)\n'
            '\n'
            '# 實例化\n'
            'pipe = DataPipeline("ETL_v1")\n'
            'pipe.add_step("read_csv")\n'
            'pipe.add_step("drop_nulls")\n'
            'print(pipe.steps)   # [\'read_csv\', \'drop_nulls\']'
        ),
        bullets=[
            "class：宣告類別",
            "__init__：初始化鉤子",
            "self：當前實例參照",
            "add_step：行為 (method)",
            "pipe：instance",
        ],
        label_dark=True,
    )
    add_source(s, "Python Tutorial §9.3 A First Look at Classes")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───────── S7 · EXAMPLE-I/O — 實例化三欄走位 ─────────
    s = _blank(prs)
    draw_three_blocks_flow(
        s,
        title="實例化那一行，Python 在背後走三步",
        blocks=[
            {"heading": "① 輸入 (呼叫)",
             "items": [
                 '`pipe = DataPipeline("ETL_v1")`',
                 "使用者只寫這一行",
                 "看起來像呼叫函式",
             ]},
            {"heading": "② 過程 (直譯器動作)",
             "items": [
                 "建立空 instance obj",
                 '呼叫 __init__(self=obj, ...)',
                 "執行 self.name = ...",
             ]},
            {"heading": "③ 產出 (綁名字)",
             "items": [
                 "obj 綁到 pipe",
                 "pipe.name == 'ETL_v1'",
                 "pipe.steps == []",
             ]},
        ],
        bottom_note="三步背熟 → 任何 class 的實例化都一樣流程；不是黑魔法，是固定劇本。",
    )
    add_source(s, "Python Data Model §3.3.1 Basic customization")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── S8 · CODE — __init__ + self 解剖 ─────────
    s = _blank(prs)
    add_title(s, "__init__ 是初始化鉤子，self 是當前實例參照——都是慣例，不是關鍵字")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="__init__ 與 self 的實際語意",
        code=(
            'class DataPipeline:\n'
            '    def __init__(self, name):     # self 由 Python 自動傳入\n'
            '        self.name = name          # 寫入當前實例的屬性\n'
            '\n'
            '# self 是 PEP 8 慣例，理論上可改名 (但別改)：\n'
            'class Bad:\n'
            '    def __init__(me, x):          # 可跑但沒人看得懂\n'
            '        me.x = x\n'
            '\n'
            '# 方法呼叫時 Python 做的事：\n'
            '# pipe.add_step("x")\n'
            '# → DataPipeline.add_step(pipe, "x")   ← 實例被塞進第一參數\n'
        ),
        bullets=[
            "__init__ 不是建構子\n是「初始化鉤子」",
            "self 由 Python 自動傳入\n你不手動傳",
            "慣例叫 self (PEP 8)\n改名合法但違規",
            "省略 self → TypeError\n(見 S14)",
        ],
        label_dark=True,
    )
    add_source(s, "Python Data Model §3.3.1 · PEP 8 §Function and method arguments")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── S9 · CHECKPOINT — self 三問 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=1, cols=3, cells=[
        {"text": "Q1 · 誰傳進來？",
         "sub": "Python 自動把\n當前實例塞進第一個參數\n你不用手動傳",
         "highlight": True},
        {"text": "Q2 · 指到哪？",
         "sub": "指到那一顆 instance\nself.name = ... 就是\n寫入這顆物件的屬性"},
        {"text": "Q3 · 不寫會怎樣？",
         "sub": "TypeError\n或被誤判為 classmethod\n(未進階前不要碰)"},
    ], title="Check Point · self 三問：背熟，OOP 一半功課就完成")
    add_source(s, "Python Tutorial §9.3.2 Class and Instance Variables")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── S10 · CODE — 兩實例獨立驗證 ─────────
    s = _blank(prs)
    add_title(s, "兩個實例各自的 steps，互不污染——記憶體才是真相")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.8), h=Inches(5.0),
        slot_name="兩實例各自 steps 的記憶體示意",
        description="兩個 instance 方塊：p1 與 p2\n各自指向獨立的 steps list\n(reference 箭頭標示)",
        url_hint="https://docs.python.org/3/reference/datamodel.html",
        placeholder_id="F3_S10_memory_model",
        registry=image_registry,
        size_hint="1400×1200 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.7), y=Inches(1.3),
        w=Inches(6.0), h=Inches(5.0),
        label="驗證兩實例獨立",
        code=(
            'p1 = DataPipeline("ETL_A")\n'
            'p2 = DataPipeline("ETL_B")\n'
            '\n'
            'p1.add_step("read_csv")\n'
            'p2.add_step("read_json")\n'
            '\n'
            'print(p1.steps)\n'
            '# [\'read_csv\']\n'
            'print(p2.steps)\n'
            '# [\'read_json\']\n'
            '\n'
            'print(p1.steps is p2.steps)\n'
            '# False  ← 不同記憶體'
        ),
        bullets=[
            "每次 __init__ 執行\n都建一個新 list",
            "p1.steps / p2.steps\n指向不同物件",
            "這就是實例獨立的根基",
        ],
        label_dark=True,
    )
    add_source(s, "Python Reference §3.1 Objects, values, and types")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── S11 · PITFALL (P3) — 可變 Class Attribute ─────────
    s = _blank(prs)
    _draw_pitfall(
        s,
        title="陷阱 · items = [] 放類別頂端 → 所有實例共享",
        wrong_label="Class Attribute 放可變物件",
        wrong_code="class Bad:\n"
                   "    items = []         # 類別頂端\n\n"
                   "a, b = Bad(), Bad()\n"
                   "a.items.append('x')\n\n"
                   "print(b.items)\n"
                   "→ ['x']  ← b 也被污染！\n"
                   "# 所有實例共享同一 list",
        right_label="可變物件一律放 __init__",
        right_code="class Good:\n"
                   "    def __init__(self):\n"
                   "        self.items = []  # 每實例獨立\n\n"
                   "a, b = Good(), Good()\n"
                   "a.items.append('x')\n\n"
                   "print(b.items)\n"
                   "→ []  ← b 不受影響",
        why="類別頂端只在 class 定義時執行一次；__init__ 每次實例化都重建 → list/dict/set 必須放 __init__",
    )
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── S12 · TABLE — Class Attribute 紅綠燈 ─────────
    s = _blank(prs)
    draw_editorial_table(
        s,
        header=["用途", "範例", "是否 OK", "備註"],
        rows=[
            ["常數 (不可變)",
             'PI = 3.14159',
             "✓ OK",
             "int / float / str / tuple 皆安全"],
            ["預設設定字串",
             'DEFAULT_ENCODING = "utf-8"',
             "✓ OK",
             "字串不可變，共享無副作用"],
            ["版本號 / 類別識別碼",
             'VERSION = "1.0"',
             "✓ OK",
             "配 classmethod 讀取更整齊"],
            ["可變容器 (list)",
             'items = []',
             "✗ 禁",
             "永遠是 bug，放 __init__"],
            ["可變容器 (dict)",
             'cache = {}',
             "✗ 禁",
             "實例間共用快取 → 資料污染"],
            ["可變容器 (set)",
             'seen = set()',
             "✗ 禁",
             "同上，放 __init__"],
        ],
        col_widths=[1.6, 2.2, 1.0, 2.2],
        title="Class Attribute 的紅綠燈：只給不可變物件用",
    )
    add_source(s, "PEP 8 §Designing for Inheritance · Ruff RUF012")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── S13 · CHECKPOINT — 中段驗收 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=1, cols=3, cells=[
        {"text": "Q1 · 可變放哪？",
         "sub": "self.cache = {}\n放 __init__ 還是\n類別頂端？",
         "highlight": True},
        {"text": "Q2 · 共享會怎樣？",
         "sub": "class Bad: items=[]\n兩個實例 append()\n結果 shape 長怎樣？"},
        {"text": "Q3 · 哪些可以放頂端？",
         "sub": "VERSION='1.0'\nPI=3.14\nDEFAULT_CONF={}\n三選二，哪個不行？"},
    ], title="Check Point · 中段驗收：答得出來再往下 (卡住回 S10-S12)")
    add_source(s, "F3 課程節奏：每 5 張插一次 Checkpoint")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── S14 · PITFALL (P4) — 省略 self ─────────
    s = _blank(prs)
    _draw_pitfall(
        s,
        title="陷阱 · 省略 self 或亂改名",
        wrong_label="忘了寫 self",
        wrong_code="class Pipeline:\n"
                   "    def __init__(name):\n"
                   "        name = name\n\n"
                   "p = Pipeline('ETL')\n"
                   "# TypeError:\n"
                   "# __init__() takes 1\n"
                   "# positional argument\n"
                   "# but 2 were given",
        right_label="第一個參數永遠是 self",
        right_code="class Pipeline:\n"
                   "    def __init__(self, name):\n"
                   "        self.name = name\n\n"
                   "p = Pipeline('ETL')\n"
                   "print(p.name)\n"
                   "→ 'ETL'\n"
                   "# self 是慣例，別亂改",
        why="Python 呼叫方法時會自動把實例塞進第一個參數；沒 self 接 → 參數數量對不上 → TypeError",
    )
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── S15 · CONCEPT-CARD — 何時「不」該用 OOP ─────────
    s = _blank(prs)
    add_title(s, "反例 · 何時不該用 OOP：不是每個腳本都需要 class")
    # 左欄：不該用的三種情境
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.0), Inches(0.5),
        "✗ 這些情境，function 就夠了",
        font_size=T.FONT_BODY, color=RED_ERROR, bold=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(1.95), Inches(6.0), Inches(4.2),
        "① 一次性腳本\n"
        "   只跑一次、不會被別人 import\n"
        "   例：一次性匯出當月報表的 50 行程式\n\n"
        "② 純函式式轉換\n"
        "   輸入 → 輸出，無需保存狀態\n"
        "   例：df.apply(lambda x: x.strip())\n\n"
        "③ 單一資料結構 + 少量操作\n"
        "   dict / dataclass / NamedTuple 更輕量\n"
        "   例：使用者設定、設定檔結構",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    # 右欄：該用的三種情境
    add_textbox(
        s, Inches(7.0), Inches(1.4), Inches(5.7), Inches(0.5),
        "✓ 這些情境，OOP 會降低成本",
        font_size=T.FONT_BODY, color=GREEN_OK, bold=True,
    )
    add_textbox(
        s, Inches(7.0), Inches(1.95), Inches(5.7), Inches(4.2),
        "① 狀態需要跨方法共享\n"
        "   清洗管線記住 schema / 錯誤清單\n\n"
        "② 同一組邏輯套到多個資料集\n"
        "   三家客戶格式略不同\n"
        "   → 實例化三次，傳不同 config\n\n"
        "③ 需要 method chaining / 可替換步驟\n"
        "   cleaner.read().clean().export()\n"
        "   (F5 的目標形狀)",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "判斷準則：有沒有「跨方法共享的狀態」＋「重複套用的需求」？兩個都有 → 用 class；只有一個 → function 即可。",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "Martin, Clean Code §10 Classes · Python dataclasses docs")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── S16 · CODE — DataPipeline 雛形 + chaining 伏筆 ─────────
    s = _blank(prs)
    add_title(s, "DataPipeline 雛形：F3 的終點，F5 的起點")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="完整雛形：狀態 + 行為 + method chaining 伏筆",
        code=(
            'class DataPipeline:\n'
            '    """A minimal reusable pipeline skeleton."""\n'
            '\n'
            '    def __init__(self, name):\n'
            '        self.name = name\n'
            '        self.steps = []\n'
            '\n'
            '    def add_step(self, step):\n'
            '        self.steps.append(step)\n'
            '        return self              # F4/F5 chaining 伏筆\n'
            '\n'
            '    def run(self, data):\n'
            '        for step in self.steps:\n'
            '            data = step(data)\n'
            '        return data\n'
            '\n'
            '# 使用：\n'
            'pipe = DataPipeline("ETL_v1")\n'
            'pipe.add_step(str.strip).add_step(str.lower)\n'
            'pipe.run("  HELLO ")   # \'hello\''
        ),
        bullets=[
            "狀態：name / steps",
            "行為：add_step / run",
            "return self →\nmethod chaining",
            "F4 加繼承 + 魔術方法",
            "F5 擴成 DataCleaner",
        ],
        label_dark=True,
    )
    add_source(s, "scikit-learn Pipeline 設計啟發 · Martin, Clean Code §10 Classes")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── S17 · PRACTICE-PROMPT — Student 類別 ─────────
    s = _blank(prs)
    add_title(s, "練習時間 · 5 分鐘 · 各自動手")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(11), Inches(0.4),
        "🟡 核心題 · 難度：容易 · 目標：鞏固 LO3 / LO4 / LO5",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(1.85), Inches(11.0), Inches(4.3),
        "情境：你拿到三家客戶的銷售資料，格式略有不同。\n"
        "         要計算每家的平均客單價。用 function 還是 class？\n\n"
        "   任務：寫一個 Student 類別 (模擬三家客戶 → 三個 instance)\n\n"
        "     屬性：\n"
        "       · name   (字串)\n"
        "       · scores (list，可變 → 放哪裡？)\n\n"
        "     方法：\n"
        "       · add_score(score)   加一筆成績\n"
        "       · average()          回傳平均分數\n\n"
        "   驗收：建立兩個 Student，各自 add_score，互不污染\n"
        "   思考題：若改成三家客戶各自的『平均客單價』，這個類別改幾行就能重用？",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        family=T.FONT_MONO, line_spacing=1.4,
    )
    draw_emphasis_pill(
        s, Inches(4.0), Inches(6.35), Inches(5.3), Inches(0.5),
        "Think · Pair · Share — 5 分鐘後對答案",
        inverted=True,
    )
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───────── S18 · SILENT — 銜接 F4 預告 ─────────
    s = _blank(prs)
    draw_silent_page(s, "懂了藍圖，\n下一步是學會美化它。")
    add_textbox(
        s, T.MARGIN_X, Inches(5.8), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
        "F4 · 封裝 (property) ｜ 繼承 (is-a 關係) ｜ 魔術方法 (__repr__ / __eq__ / __len__)",
        font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 18, N_CONTENT, dark_bg=True)

    # ───────── S19 · PYRAMID — 收束 ─────────
    s = _blank(prs)
    draw_thesis_hierarchy(
        s,
        blocks=[
            {"heading": "四個關鍵詞 (硬記)",
             "items": [
                 "Class = 定義模板",
                 "Object / Instance = 執行期的實體",
                 "__init__ = 實例化時自動執行的初始化鉤子",
                 "self = 指向當前 instance 的第一參數",
             ]},
            {"heading": "三章銜接線",
             "items": [
                 "F3 · class + __init__ + self (今日)",
                 "F4 · 加封裝 / 繼承 / 魔術方法",
                 "F5 · DataCleaner().read().clean().export()",
             ]},
        ],
        title="F3 收束：地基不穩，上面再漂亮都會塌",
        thesis="先搞懂狀態住哪，才配談設計——F4 延伸封裝、繼承與魔術方法；F5 擴成實戰 DataCleaner。",
    )
    add_source(s, "F3 module synthesis · F4 / F5 銜接預告")
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───────── S20 · CONCEPT-CARD — 銜接 F5 DataCleaner 預告 ─────────
    s = _blank(prs)
    add_title(s, "銜接 F5：今日的 DataPipeline 會長大成 DataCleaner")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.0), Inches(0.5),
        "F3 今日產出（雛形）",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(1.95), Inches(6.0), Inches(4.2),
        'class DataPipeline:\n'
        '    def __init__(self, name):\n'
        '        self.name = name\n'
        '        self.steps = []\n\n'
        '    def add_step(self, step):\n'
        '        self.steps.append(step)\n'
        '        return self\n\n'
        '    def run(self, data):\n'
        '        for step in self.steps:\n'
        '            data = step(data)\n'
        '        return data',
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        family=T.FONT_MONO, line_spacing=1.4,
    )
    add_textbox(
        s, Inches(7.0), Inches(1.4), Inches(5.7), Inches(0.5),
        "F5 成熟產出（實戰）",
        font_size=T.FONT_BODY, color=GREEN_OK, bold=True,
    )
    add_textbox(
        s, Inches(7.0), Inches(1.95), Inches(5.7), Inches(4.2),
        'cleaner = DataCleaner("sales.csv")\n'
        'result = (cleaner\n'
        '    .read()\n'
        '    .drop_duplicates()\n'
        '    .fill_missing("price", 0)\n'
        '    .clean_column_names()\n'
        '    .export("clean.parquet"))\n\n'
        '# 今日的 return self\n'
        '# 就是明天的 method chaining\n'
        '# 今日的 self.steps = []\n'
        '# 就是明天的 self._errors = []',
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        family=T.FONT_MONO, line_spacing=1.4,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "沒有 F3 的三個關鍵詞，F5 的 method chaining 都只是抄語法；會了 F3，F5 只是把規模做大。",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "F5 DataCleaner 設計 · pandas.DataFrame method chaining pattern")
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ───────── S21 · SILENT — 過渡 ─────────
    s = _blank(prs)
    draw_silent_page(s, "先搞懂狀態住哪，\n才配談設計。")
    add_footer(s, MODULE_CODE, 21, N_CONTENT, dark_bg=True)

    # Copyright
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
