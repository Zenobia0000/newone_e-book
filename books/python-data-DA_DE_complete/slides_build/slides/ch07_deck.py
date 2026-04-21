"""F5 deck — OOP × Pandas 整合實戰（6 小時基礎段收斂章）.

22 content slides + cover + copyright page.

Governing thought:
    6 小時基礎的終點，不是『你會 OOP』——
    是你能把 F1-F4 + Pandas 組成一個可鏈式呼叫的 DataCleaner。

Aligned to chapters/F5_OOP_Pandas整合實戰/00_skeleton.yaml
    · Teaching-track archetypes：SILENT / ASK / VS / MATRIX / CONCEPT / CODE /
      MECHANISM-FLOW / PITFALL / PRACTICE-PROMPT / TABLE
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source,
    draw_silent_page, draw_ask_page, draw_matrix,
    draw_vs_two_col, draw_code_panel, draw_flow_chain,
    draw_inverted_thesis_box, draw_editorial_table,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "F5"
MODULE_TITLE = "OOP × Pandas 整合實戰"
MODULE_SUBTITLE = "把 F1-F4 + Pandas 組成一條會跑的 DataCleaner 管線"
TIME_MIN = 90
N_CONTENT = 22


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_ch07(output_path, image_registry=None):
    """Build F5 deck — foundation-closing OOP × Pandas chapter."""
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── F01 · SILENT ─────────
    s = _blank(prs)
    draw_silent_page(
        s,
        "6 小時基礎的終點，不是「你會 OOP」。\n"
        "是你能把 F1-F4 + Pandas 組成一條會跑的管線。",
    )
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ───────── F02 · ASK ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "每月一份髒 CSV，你會交 5 個零散 function，\n"
        "還是一個能 .validate().clean().export() 串起來的類別？",
        data_card={
            "label": "工程現場觀察",
            "stat": "3×",
            "caption": "寫成 class 的程式\n6 個月後還能用的機率\n是零散腳本的 3 倍",
        },
    )
    add_source(s, "工程可維護性經驗法則 · F3-F4 OOP 章回顧")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── F03 · VS · 零散 function vs 類別管線 ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="零散 function 版（腳本思維）",
        right_title="類別管線版（工程思維）",
        left_items=[
            "df = load(path); df = validate(df)",
            "df = clean(df, 'drop'); df = apply(df, 'revenue', f)",
            "狀態在每個變數名上，誰被改過要用人腦追",
            "錯誤發生時，沒有統一的『錯誤邊界』可 catch",
            "改一個欄位規則，5 個 function 都要翻一遍",
        ],
        right_items=[
            "cleaner = DataCleaner(path)",
            ".validate().clean_missing_values().apply_custom_transform(...)",
            "狀態收在 self.df，整條管線共享同一份資料",
            "raise DataValidationError → 上游一個 except 就能處理",
            "改規則只改一個 method，其他呼叫方零影響",
        ],
        title="同一份 raw.csv，兩種寫法的工程差距在哪",
        summary="差距不在『能不能跑』，在『6 個月後能不能改』。",
    )
    add_source(s, "F3-F4 OOP 設計原則 · 資料管線封裝模式")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───────── F04 · MATRIX · 狀態 vs 動作 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "狀態（Instance Attribute）",
         "sub": "self.path / self.df\n「這個物件記得什麼」\n整個生命週期都要存的",
         "highlight": True},
        {"text": "動作（Method）",
         "sub": "validate / clean / export\n「這個物件會做什麼」\n對 self.df 的改造流程",
         "highlight": True},
        {"text": "可鏈式（return self）",
         "sub": "改完 self.df 後 return self\n讓 a().b().c() 成立\nMethod Chaining 的核心"},
        {"text": "會 raise 的（Exception）",
         "sub": "業務規則違反時 raise\n自訂 DataValidationError\n給上游精準的錯誤語意"},
    ], title="類別設計先分四格：狀態 / 動作 / 鏈式 / 例外")
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "先分清楚誰該存起來、誰該做事，class 才會乾淨。",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID, align=PP_ALIGN.CENTER,
    )
    add_source(s, "F3 OOP 核心觀念 · F4 封裝繼承 · Linus 資料結構優先原則")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── F05 · CODE · 類別介面草圖 ─────────
    s = _blank(prs)
    add_title(s, "先寫介面、再寫實作：DataCleaner 6 個方法簽名")
    draw_code_panel(
        s, x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="介面 = 你對未來自己的承諾",
        code=(
            'from pathlib import Path\n'
            'import pandas as pd\n'
            '\n'
            'class DataCleaner:\n'
            '    def __init__(self, data_path: str): ...\n'
            '    def validate(self) -> "DataCleaner": ...\n'
            '    def clean_missing_values(\n'
            '        self, strategy: str = "drop",\n'
            '    ) -> "DataCleaner": ...\n'
            '    def apply_custom_transform(\n'
            '        self, column: str, func,\n'
            '    ) -> "DataCleaner": ...\n'
            '    def generate_eda_report(\n'
            '        self, out_path: str,\n'
            '    ) -> "DataCleaner": ...\n'
            '    def export_data(self, out_path: str) -> "DataCleaner": ...'
        ),
        bullets=[
            "六個方法：\n一進（init）、五動（其他）",
            "每個 method 都 return\nself → 才能鏈式",
            "type hint 讓 IDE\n能提示下一步怎麼接",
            "先定介面、再填實作 —\nAPI 設計紀律",
        ],
        label_dark=True,
    )
    add_source(s, "DataCleaner API 設計 · PEP 484 type hint")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── F06 · CODE · DataValidationError ─────────
    s = _blank(prs)
    add_title(s, "自訂 Exception：讓錯誤有語意、能被精準 catch")
    draw_code_panel(
        s, x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="一個 Exception class 就夠 — 別過度設計",
        code=(
            'class DataValidationError(Exception):\n'
            '    """資料驗證失敗時拋出"""\n'
            '    def __init__(self, column: str, reason: str):\n'
            '        self.column = column\n'
            '        self.reason = reason\n'
            '        super().__init__(f"[{column}] {reason}")\n'
            '\n'
            '# 使用：\n'
            'if (df["revenue"] < 0).any():\n'
            '    raise DataValidationError(\n'
            '        "revenue", "出現負值",\n'
            '    )\n'
            '\n'
            '# 捕捉（上游可精準處理）：\n'
            'try:\n'
            '    cleaner.validate()\n'
            'except DataValidationError as e:\n'
            '    log.warning(f"欄位 {e.column} 失敗：{e.reason}")'
        ),
        bullets=[
            "繼承 Exception —\n這就是 F4 的繼承實戰",
            "column + reason 兩個欄位 —\n錯誤帶語意",
            "上游能 catch\nDataValidationError\n而非所有 Exception",
            "raise Exception('xxx') 是\n業餘寫法，抓不到語意",
        ],
        label_dark=True,
    )
    add_source(s, "Python Docs §Built-in Exceptions · F4 繼承章回顧")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───────── F07 · MECHANISM-FLOW · 六段管線 ─────────
    s = _blank(prs)
    draw_flow_chain(
        s,
        title="DataCleaner 六段管線：__init__ → validate → clean → apply → eda → export",
        y=3.2,
        nodes=[
            {"label": "① __init__",
             "sub": "讀檔 + 路徑檢查",
             "caption": "fail fast"},
            {"label": "② validate",
             "sub": "欄位 + 值域\n業務規則",
             "caption": "DataValidationError",
             "highlight": True},
            {"label": "③ clean",
             "sub": "dropna / fillna\nPandas 缺失值",
             "caption": "return self"},
            {"label": "④ apply",
             "sub": "Lambda × apply\n新增衍生欄",
             "caption": "F2 × Pandas",
             "highlight": True},
            {"label": "⑤ eda",
             "sub": "Matplotlib\n自動 4 圖",
             "caption": "報表輸出"},
            {"label": "⑥ export",
             "sub": "to_csv\n收尾交付",
             "caption": "clean.csv"},
        ],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "每段都對應一個你學過的概念 —— 這章是在把地基接起來，不是學新東西。",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "DataCleaner pipeline 設計 · F1-F4 + Pandas 能力整合")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── F08 · CODE · EXAMPLE I/O ─────────
    s = _blank(prs)
    add_title(s, "EXAMPLE I/O：raw.csv 進、clean.csv + eda.png 出")
    draw_code_panel(
        s, x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="先看輸入 / 輸出長什麼樣，再看程式碼",
        code=(
            '# 輸入 raw.csv（髒資料典型樣貌）\n'
            'order_id,date,category,revenue,qty\n'
            '1001,2026-01-03,A,1200,3\n'
            '1002,2026-01-04,B,,2          ← revenue 缺值\n'
            '1003,2026-01-05,A,-50,1       ← revenue 負值（業務規則違反）\n'
            '1004,2026-01-05,C,880,NaN\n'
            '\n'
            '# 期望輸出（兩個檔）\n'
            '# ① clean.csv — 缺值補 mean、revenue 做 *1.05、無負值\n'
            'order_id,date,category,revenue,qty,revenue_transformed\n'
            '1001,2026-01-03,A,1200,3,1260.0\n'
            '1004,2026-01-05,C,880.0,2.0,924.0\n'
            '\n'
            '# ② eda.png — 2x2 子圖：4 個數值欄各一張 histogram'
        ),
        bullets=[
            "髒資料三種典型：\n缺值 / 異常值 / 型別錯",
            "validate 要把\n『revenue < 0』擋掉",
            "clean 用 fillna(mean)\n把 1002、1004 補齊",
            "apply 產 _transformed 欄\n不蓋原欄（可追溯）",
        ],
        label_dark=True,
    )
    add_source(s, "DataCleaner 範例資料集 · DA/DE 課程教材")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── F09 · CODE · __init__ ─────────
    s = _blank(prs)
    add_title(s, "Step 1 · __init__：讀檔 + 路徑檢查（fail fast）")
    draw_code_panel(
        s, x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="錯就現在炸 — 別讓錯誤延後到下游才爆｜呼應 F3",
        code=(
            'import pandas as pd\n'
            'from pathlib import Path\n'
            '\n'
            'class DataCleaner:\n'
            '    def __init__(self, data_path: str):\n'
            '        self.path = Path(data_path)\n'
            '        if not self.path.exists():\n'
            '            raise FileNotFoundError(\n'
            '                f"找不到檔案：{self.path}",\n'
            '            )\n'
            '        self.df = pd.read_csv(self.path)\n'
            '        print(\n'
            '            f"已載入 {self.path.name}，"\n'
            '            f"shape={self.df.shape}",\n'
            '        )'
        ),
        bullets=[
            "F3 __init__ 的實戰 —\n不是只會 self.x = x",
            "pathlib.Path 比\nos.path 乾淨（全域守則）",
            "fail fast：檔案不存在\n立刻炸，不要讀到一半",
            "print shape 是禮貌 —\n讓使用者確認讀對了",
        ],
        label_dark=True,
    )
    add_source(s, "F3 OOP __init__ · pathlib 最佳實踐 · fail fast 原則")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── F10 · CODE · validate ─────────
    s = _blank(prs)
    add_title(s, "Step 2 · validate：業務規則寫成 if，不通過就 raise")
    draw_code_panel(
        s, x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="validate 的責任：把業務規則變成程式｜呼應 F4 + S2",
        code=(
            '    def validate(self) -> "DataCleaner":\n'
            '        if self.df.empty:\n'
            '            raise DataValidationError(\n'
            '                "df", "資料為空",\n'
            '            )\n'
            '        if "revenue" in self.df.columns:\n'
            '            if (self.df["revenue"] < 0).any():\n'
            '                raise DataValidationError(\n'
            '                    "revenue", "出現負值",\n'
            '                )\n'
            '        return self   # ← 鏈式的靈魂'
        ),
        bullets=[
            "先檢查整體（empty）\n再檢查欄位",
            "Pandas 布林運算：\n(col < 0).any()",
            "違反規則 raise\nDataValidationError",
            "return self —\n少了就不能鏈式",
        ],
        label_dark=True,
    )
    add_source(s, "pandas §DataFrame.any · 業務驗證設計模式")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── F11 · CODE · clean_missing_values ─────────
    s = _blank(prs)
    add_title(s, "Step 3 · clean_missing_values：策略化 + return self")
    draw_code_panel(
        s, x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="strategy 參數讓方法『可選』而非『寫死』｜呼應 S2",
        code=(
            '    def clean_missing_values(\n'
            '        self, strategy: str = "drop",\n'
            '    ) -> "DataCleaner":\n'
            '        if strategy == "drop":\n'
            '            self.df = self.df.dropna()\n'
            '        elif strategy == "mean":\n'
            '            self.df = self.df.fillna(\n'
            '                self.df.mean(numeric_only=True),\n'
            '            )\n'
            '        else:\n'
            '            raise ValueError(\n'
            '                f"未知策略：{strategy}",\n'
            '            )\n'
            '        return self'
        ),
        bullets=[
            "default='drop' —\n最常用的當預設",
            "dropna / fillna 是\nPandas 的兩招通吃",
            "numeric_only=True\n避開字串欄",
            "未知策略就 raise —\n不要偷偷當成預設",
        ],
        label_dark=True,
    )
    add_source(s, "pandas §dropna §fillna · 策略模式（Strategy Pattern）")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── F12 · CODE · apply_custom_transform ─────────
    s = _blank(prs)
    add_title(s, "Step 4 · apply_custom_transform：Lambda × Pandas.apply")
    draw_code_panel(
        s, x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="F2 的 Lambda 餵進 Pandas.apply — 一行產新欄｜呼應 F2 + S3",
        code=(
            '    def apply_custom_transform(\n'
            '        self, column: str, func,\n'
            '    ) -> "DataCleaner":\n'
            '        if column not in self.df.columns:\n'
            '            raise DataValidationError(\n'
            '                column, "欄位不存在",\n'
            '            )\n'
            '        self.df[f"{column}_transformed"] = (\n'
            '            self.df[column].apply(func)\n'
            '        )\n'
            '        return self\n'
            '\n'
            '# 使用：\n'
            'cleaner.apply_custom_transform(\n'
            '    "revenue", lambda x: x * 1.05,\n'
            ')'
        ),
        bullets=[
            "func 是參數 —\n把行為注入物件",
            "F2 Lambda 的實戰場 —\n不是玩具範例",
            "新欄名 col_transformed —\n不蓋掉原欄",
            "欄位不存在先 raise —\n不要靜默吞錯",
        ],
        label_dark=True,
    )
    add_source(s, "F2 Lambda 章回顧 · pandas §Series.apply")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── F13 · CODE · generate_eda_report ─────────
    s = _blank(prs)
    add_title(s, "Step 5 · generate_eda_report：Matplotlib 自動產 4 圖")
    draw_code_panel(
        s, x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="EDA 不是手畫 — 是 class 的一個 method｜呼應 S5",
        code=(
            '    def generate_eda_report(\n'
            '        self, out_path: str = "eda.png",\n'
            '    ) -> "DataCleaner":\n'
            '        import matplotlib.pyplot as plt\n'
            '        numeric = self.df.select_dtypes(include="number")\n'
            '        fig, axes = plt.subplots(2, 2, figsize=(10, 8))\n'
            '        for ax, col in zip(\n'
            '            axes.flat, numeric.columns[:4],\n'
            '        ):\n'
            '            ax.hist(numeric[col].dropna(), bins=30)\n'
            '            ax.set_title(col)\n'
            '        fig.tight_layout()\n'
            '        fig.savefig(out_path, dpi=150)\n'
            '        plt.close(fig)\n'
            '        return self'
        ),
        bullets=[
            "select_dtypes 只挑數值欄 —\n字串欄不畫 hist",
            "2x2 共 4 張 subplots —\n一頁看完",
            "plt.close 釋放記憶體 —\n批次跑不會爆",
            "import 放 method 內 —\n這章先求能跑",
        ],
        label_dark=True,
    )
    add_source(s, "matplotlib §subplots · EDA 自動化流水線")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── F14 · CODE · export + 完整鏈式 ─────────
    s = _blank(prs)
    add_title(s, "Step 6 · export_data + 完整鏈式：5 行完成整條管線")
    draw_code_panel(
        s, x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="這就是你用 6 小時基礎換來的產物",
        code=(
            '# 匯出方法（最後一塊拼圖）\n'
            '    def export_data(self, out_path: str) -> "DataCleaner":\n'
            '        self.df.to_csv(out_path, index=False)\n'
            '        print(f"匯出 {len(self.df)} 列到 {out_path}")\n'
            '        return self\n'
            '\n'
            '# 使用：一段 5 行完成 ETL + EDA\n'
            '(DataCleaner("raw.csv")\n'
            '    .validate()\n'
            '    .clean_missing_values(strategy="mean")\n'
            '    .apply_custom_transform("revenue", lambda x: x * 1.05)\n'
            '    .generate_eda_report("report.png")\n'
            '    .export_data("clean.csv"))'
        ),
        bullets=[
            "每個 . 都是一個\nreturn self 接起來",
            "讀起來像英文：\n驗證→清洗→變換→報表→匯出",
            "6 個月後自己回來看\n還看得懂 —— 這才算過關",
            "這就是 Method\nChaining 的商業價值",
        ],
        label_dark=True,
    )
    add_source(s, "Method Chaining 經典實作 · jQuery / Pandas 皆此模式")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── F15 · PITFALL · 三大踩雷 ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="業餘寫法（P1 / P2 / P4）",
        right_title="資深紀律（這樣才過關）",
        left_items=[
            "P1 · method 忘記 return self → 第二層 chain 就變 None",
            "P2 · __init__ 不驗證 → 錯誤延後到 clean 才爆，難 debug",
            "P4 · 30 行腳本硬包成 class → 過度設計，沒重用價值",
            "用字面字串 raise Exception('壞了') → 上游抓不到語意",
        ],
        right_items=[
            "每個改動 self.df 的 method 結尾都 return self",
            "__init__ 就 fail fast — 錯就立刻炸，訊息明確",
            "class 是為了重用與擴充 — 用兩次以上才包",
            "自訂 DataValidationError — 帶 column/reason 語意",
        ],
        title="踩雷總表 1/2：return self / fail fast / 別過度設計",
        summary="class 不是拿來看起來專業，是拿來讓 6 個月後的你還能改。",
    )
    add_source(s, "工程反模式觀察 · Linus 實用主義（不解決假想威脅）")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── F16 · PITFALL · P3/P5 進階踩雷 ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="進階錯誤心態（P3 / P5）",
        right_title="能帶著走的紀律",
        left_items=[
            "P3 · 一口氣定 5 個 Exception：DataMissingError / DataNegError / ...",
            "    → 分類過細，使用者不知道該 catch 哪個",
            "P5 · 只會照抄 DataCleaner 範本，換一份資料就卡住",
            "    → 改個欄位名、加個規則就不知從哪下手",
        ],
        right_items=[
            "一個 DataValidationError 就夠 —— column + reason 兩欄攜帶語意",
            "    raise DataValidationError('revenue', '出現負值')",
            "把範本當框架：validate 加 1 條 / apply 換 1 個 lambda / eda 多 1 圖",
            "    改得動才是真的會；背得熟只是會念",
        ],
        title="踩雷總表 2/2：Exception 不要濫用 / 範本要能改造",
        summary="P3：錯誤分類要粗；P5：範本要能改 —— 這兩點過了，才算真正學會 F5。",
    )
    add_source(s, "F4 繼承章 · 教學遷移學習觀察")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── F17 · PRACTICE-PROMPT · 改造練習 ─────────
    s = _blank(prs)
    add_title(s, "核心練習：把 DataCleaner 改造成你自己的領域（5 分鐘）")
    draw_matrix(s, rows=1, cols=3, top=1.4, bottom=1.8, cells=[
        {"text": "任務 1 · 加一個 method",
         "sub": "挑一個你常做的資料動作\n（去重 / 標準化 / 類別編碼）\n\n"
                "規格：\n回傳 self、有 docstring、\n至少 1 個參數",
         "highlight": True},
        {"text": "任務 2 · 加一條驗證",
         "sub": "在 validate() 內\n加一條你領域的業務規則\n\n"
                "範例：\n「age 不能 > 150」\n「email 必須含 @」\n"
                "raise DataValidationError",
         "highlight": False},
        {"text": "任務 3 · 加一張 EDA 圖",
         "sub": "在 generate_eda_report\n加第 5 張子圖\n\n"
                "範例：\n類別欄 value_counts 長條圖\n或相關係數 heatmap",
         "highlight": False},
    ])
    draw_inverted_thesis_box(
        s,
        "改得動，才是真的會 —— 背範本不算，改一次才是自己的。",
        y=6.1, width=11.0,
    )
    add_source(s, "F5 核心練習設計 · 遷移學習（transfer learning）教學法")
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───────── F18 · SILENT · 段落收束 ─────────
    s = _blank(prs)
    draw_silent_page(
        s,
        "能改別人的範本，才算真的會 OOP。\n"
        "能把自己的資料套進去，才算真的過關。",
    )
    add_footer(s, MODULE_CODE, 18, N_CONTENT, dark_bg=True)

    # ───────── F19 · MATRIX · F1-F5 + S1-S6 能力盤點 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=3, cells=[
        {"text": "基礎 · F1 系統直覺",
         "sub": "OS / RAM / I/O 三關鍵字\n看得懂效能從哪來\n來源：F1"},
        {"text": "基礎 · F2 資料結構",
         "sub": "list / dict / Lambda /\nComprehension / Generator\n來源：F2"},
        {"text": "基礎 · F3-F4 OOP",
         "sub": "類別 / __init__ / 繼承 /\n魔術方法 / 封裝\n來源：F3-F4",
         "highlight": True},
        {"text": "基礎 · F5 整合（今天）",
         "sub": "OOP × Pandas\n可鏈式 DataCleaner\n來源：F5",
         "highlight": True},
        {"text": "實戰 · S1-S4 資料處理",
         "sub": "NumPy 向量化 / Pandas I/O /\ngroupby×merge×pivot /\n時序 EDA"},
        {"text": "實戰 · S5-S6 視覺化",
         "sub": "matplotlib 五圖 /\nPlotly dashboard.html /\nCapstone 交付"},
    ], title="F1–F5 + S1–S6 能力盤點：11 段地基，少一段上面就不穩")
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "你今天能寫 DataCleaner，是因為前 4 章都穩了 —— 這就是基礎段的收斂。",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID, align=PP_ALIGN.CENTER,
    )
    add_source(s, "DA/DE 12h 課程能力地圖 · F1-F5 + S1-S6 合成")
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───────── F20 · TABLE · 下一步五條路線 ─────────
    s = _blank(prs)
    draw_editorial_table(
        s,
        title="下一站五條路線：挑一條走 3 個月，比五條都碰過更值錢",
        header=["方向", "推薦起點", "什麼時候該選它"],
        rows=[
            ["資料庫整合", "SQL + SQLAlchemy",
             "工作開始接觸資料庫、CSV 已經餵不飽你了"],
            ["機器學習", "scikit-learn 入門",
             "清洗好的資料想預測未來，不只描述現在"],
            ["大規模資料", "Polars / DuckDB",
             "Pandas 讀 GB 級檔案開始慢 / OOM 時"],
            ["自動化排程", "Airflow / Prefect / n8n",
             "手動跑 DataCleaner 已經煩到想自動化"],
            ["雲端部署", "AWS S3 + Lambda / GCP BigQuery",
             "管線要共享、要排程、要可被其他系統呼叫"],
        ],
        col_widths=[1.2, 1.8, 3.0],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "五條都是獨立主題，別同時開五個坑 —— 選一條，3 個月後你會是那個主題的可靠同事。",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "DA/DE 職涯路徑觀察 · 業界主流工具棧（2026）")
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ───────── F21 · CODE · 結業三句忠告 ─────────
    s = _blank(prs)
    add_title(s, "結業三句忠告：帶著離開教室的紀律")
    draw_code_panel(
        s, x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="Linus 風格 · 專業版三件事",
        code=(
            '# 1) 跑得起來不算數 —— 6 個月後看得懂才算數\n'
            '#    return self / type hint / docstring / 好的變數名\n'
            '#    是給未來的你，不是給現在的自己。\n'
            '\n'
            '# 2) 最強的優化是「不要做」\n'
            '#    在向量化、在 Method Chaining、在自訂 Exception 之前——\n'
            '#    先問：這一步能不能跳過？\n'
            '#    能砍的程式比能寫的程式更值錢。\n'
            '\n'
            '# 3) OOP 不是萬靈丹\n'
            '#    30 行的腳本不需要 class；\n'
            '#    一次性的分析 notebook 不需要 class；\n'
            '#    class 是為了「會被呼叫第二次以上」才包。'
        ),
        bullets=[
            "紀律 1：寫給未來自己 —\ntype hint、docstring 別省",
            "紀律 2：最強優化是不做 —\n能砍比能寫重要",
            "紀律 3：OOP 不濫用 —\n重用 2 次以上才包",
            "這三句內化了，\n下一份工作直接用",
        ],
        label_dark=True,
    )
    add_source(s, "Linus Torvalds 實用主義精神 · F1-F5 全程收斂")
    add_footer(s, MODULE_CODE, 21, N_CONTENT)

    # ───────── F22 · SILENT 結業 ─────────
    s = _blank(prs)
    draw_silent_page(
        s,
        "6 小時基礎段到此收束。\n"
        "你今天寫出的 DataCleaner，就是 F1-F5 的結業證書。\n"
        "下一站 —— 選一條路，3 個月後見。",
    )
    add_footer(s, MODULE_CODE, 22, N_CONTENT, dark_bg=True)

    # Copyright
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
