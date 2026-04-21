"""Ch03 deck — 自訂函式與特殊 Python 函式
20 content slides + cover + copyright page.

Governing thought：
    Lambda / map / filter / Comprehension 不是炫技糖衣，
    是為 Pandas apply 鋪跑道。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_rect, add_textbox, add_title, add_source,
    set_solid_fill, set_no_line, set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_vs_two_col, draw_code_panel, draw_thesis_hierarchy,
    draw_inverted_thesis_box, draw_image_placeholder,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "Ch03"
MODULE_TITLE = "自訂函式與特殊 Python 函式"
MODULE_SUBTITLE = "會宣告，就少寫 80% 迴圈"
TIME_MIN = 120
N_CONTENT = 20


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_ch03(output_path, image_registry=None):
    """Build Ch03 deck.

    image_registry: optional list collecting image placeholder metadata
                    for downstream YAML emission.
    """
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── S1 · SILENT ─────────
    s = _blank(prs)
    draw_silent_page(s, "會宣告，\n就少寫 80% 迴圈。")
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ───────── S2 · ASK ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "為什麼資深資料工程師讀到 for 迴圈就皺眉？",
        data_card={
            "label": "Pandas 使用者日常寫法",
            "stat": "73%",
            "caption": "日常用 apply(lambda)；\nComprehension 使用率 89%",
        },
    )
    add_source(s, "Real Python Pandas Usage Survey 2024 (n=3,200)")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── S3 · MATRIX 2×3：函式參數五種 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=3, cells=[
        {"text": "位置參數 (positional)",
         "sub": "def f(x, y):\n必填，順序敏感",
         "highlight": True},
        {"text": "關鍵字參數 (keyword)",
         "sub": "f(x=1, y=2)\n呼叫端明示語意"},
        {"text": "預設參數 (default)",
         "sub": "def f(x=10):\n可省略，但別給可變物件"},
        {"text": "*args",
         "sub": "收剩下的位置參數\n進函式是 tuple"},
        {"text": "**kwargs",
         "sub": "收剩下的關鍵字參數\n進函式是 dict"},
        {"text": "/ 與 * 分隔（進階）",
         "sub": "PEP 570 · 強制位置 / 強制關鍵字\n日常較少用"},
    ], title="函式參數五種分工：由嚴到鬆")
    add_source(s, "Python Language Reference §8.6 · PEP 3102 · PEP 570")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───────── S4 · CODE：*args / **kwargs 實戰 ─────────
    s = _blank(prs)
    add_title(s, "*args 收散彈，**kwargs 收規則 — 彈性 API 的兩招")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.6),
        label="資料驗證函式：欄位清單 + 規則 dict",
        code=(
            'def validate(row, *required_fields, **rules):\n'
            '    for f in required_fields:\n'
            '        if f not in row:\n'
            '            return False\n'
            '    for k, rule in rules.items():\n'
            '        if not rule(row.get(k)):\n'
            '            return False\n'
            '    return True\n'
            '\n'
            'validate(row, "id", "email",\n'
            '         age=lambda x: 0 < x < 120,\n'
            '         email=lambda s: "@" in s)'
        ),
        bullets=[
            "*args 進來是 tuple",
            "**kwargs 進來是 dict",
            "規則用 lambda 從呼叫端注入",
            "這正是 Ch10 DataCleaner 的雛形",
        ],
        label_dark=True,
    )
    add_source(s, "PEP 3102 Keyword-Only Arguments")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── S5 · VS-CODE：可變預設值陷阱 ─────────
    s = _blank(prs)
    add_title(s, "可變預設值：Python 最常被當面試考的陷阱")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.4),
        label="BEFORE：bag=[] 共用同一個 list",
        code=(
            'def add(item, bag=[]):\n'
            '    bag.append(item)\n'
            '    return bag\n'
            '\n'
            'add(1)  # [1]\n'
            'add(2)  # [1, 2]  ← 不是 [2]！'
        ),
        bullets=[
            "預設值在『def 那一刻』建立一次",
            "之後呼叫終身共用",
            "Ruff B006 / pylint 會告警",
        ],
        label_dark=False,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(3.9),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.4),
        label="AFTER：None 哨兵 + 首行初始化",
        code=(
            'def add(item, bag=None):\n'
            '    if bag is None:\n'
            '        bag = []\n'
            '    bag.append(item)\n'
            '    return bag'
        ),
        bullets=[
            "None 不可變，可當哨兵",
            "每次呼叫都重建新 list",
            "同樣適用於 dict / set 預設值",
        ],
        label_dark=True,
    )
    add_source(s, "PEP 8 · Ruff rule B006 Mutable default argument")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── S6 · CODE + IMAGE：LEGB 作用域 ─────────
    s = _blank(prs)
    add_title(s, "LEGB：Python 找名字的四層，由內往外")
    # 左圖 placeholder
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.3), h=Inches(5.0),
        slot_name="LEGB 四層作用域示意",
        description="四層巢狀方框：Local → Enclosing → Global → Built-in\n由內往外查找箭頭",
        url_hint="https://docs.python.org/3/tutorial/classes.html#python-scopes-and-namespaces",
        placeholder_id="Ch03_S06_LEGB",
        registry=image_registry,
        size_hint="1280×1200 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.2), y=Inches(1.3),
        w=Inches(6.5), h=Inches(5.0),
        label="LEGB 查找順序示範",
        code=(
            'x = "Global"\n'
            '\n'
            'def outer():\n'
            '    x = "Enclosing"\n'
            '    def inner():\n'
            '        x = "Local"\n'
            '        print(x)   # Local\n'
            '    inner()\n'
            '    print(x)       # Enclosing\n'
            '\n'
            'outer()\n'
            'print(x)           # Global'
        ),
        bullets=[
            "由內往外找",
            "找不到 → NameError",
            "UnboundLocalError = 在 Local 層被指派過",
        ],
        label_dark=True,
    )
    add_source(s, "Python Tutorial §9.2 Python Scopes and Namespaces")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───────── S7 · TABLE：global / nonlocal 代價 ─────────
    s = _blank(prs)
    draw_editorial_table(
        s,
        header=["關鍵字", "作用", "合理場景", "為什麼少用"],
        rows=[
            ["global x", "從函式內改全域名字",
             "module flag / logging 設定",
             "呼叫端無法預測狀態，難 mock"],
            ["nonlocal x", "改外層函式的名字",
             "closure 累加器 / decorator",
             "只有在真寫 closure 才該出現"],
            ["（盡量避免）", "大多改用回傳值即可",
             "回傳值 > 共用狀態",
             "Pipeline 可重跑 > 寫得短"],
        ],
        col_widths=[1.3, 1.6, 1.8, 1.8],
        title="能不用就不用：global / nonlocal 的真正代價",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
        "看到 global 就該問：真的沒辦法用回傳值嗎？99% 時候有。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "PEP 3104 nonlocal · Beazley, Python Essential Reference §6")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── S8 · ASK：Lambda 存在的理由 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "如果 Lambda 只能寫一行，\n它能解決什麼樣的問題？",
        data_card={
            "label": "Lambda 設計意圖",
            "stat": "≤ 1 行",
            "caption": "當函式 / 臨時 / 即用即丟\n塞進別人要 callable 的地方",
        },
    )
    add_source(s, "Guido van Rossum, The History of Python: Origins of Lambda")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── S9 · VS-CODE：def vs lambda 邊界 ─────────
    s = _blank(prs)
    add_title(s, "def 與 lambda 的邊界：一行是線，過線就 def")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.2),
        label="LAMBDA OK：單一 expression、即用即丟",
        code=(
            'square = lambda x: x ** 2\n'
            '\n'
            'sorted(users, key=lambda u: u["age"])\n'
            'df["price"].apply(lambda x: x * 0.9)'
        ),
        bullets=[
            "單一 expression，沒有 return",
            "傳給 HOF 當 callback 最合適",
            "即用即丟，不需要名字",
        ],
        label_dark=False,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(3.7),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.6),
        label="過線該 def：有分支、有多行、要測試",
        code=(
            '# BAD：lambda row: row["p"]*0.9 if row["p"]>100 else row["p"]\n'
            '\n'
            'def discount(row):\n'
            '    """High-tier 9 折；其他原價。"""\n'
            '    if row["p"] > 100:\n'
            '        return row["p"] * 0.9\n'
            '    return row["p"]'
        ),
        bullets=[
            "可讀性 > 一行主義",
            "def 才能加 docstring、寫單元測試",
            "有分支 / try / 多行 → def",
        ],
        label_dark=True,
    )
    add_source(s, "PEP 8 §Programming Recommendations")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── S10 · CODE：map / filter / sorted 三件套 ─────────
    s = _blank(prs)
    add_title(s, "三件套：告訴 Python『做什麼』，不寫『怎麼做』")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="map · filter · sorted(key=) — Pandas apply 的祖先",
        code=(
            'prices = [120, 80, 200, 50]\n'
            '\n'
            '# map：全部套函式\n'
            'discounted = list(map(lambda x: x * 0.9, prices))\n'
            '\n'
            '# filter：保留符合條件\n'
            'bigs = list(filter(lambda x: x > 100, prices))\n'
            '\n'
            '# sorted(key=)：自訂排序鍵\n'
            'users = [{"age": 30}, {"age": 22}, {"age": 45}]\n'
            'by_age = sorted(users, key=lambda u: u["age"])'
        ),
        bullets=[
            "map / filter 回傳 iterator",
            "要 list() 才物化看內容",
            "sorted 是穩定排序",
            "Ch08 的 df.apply / sort_values\n就是同概念搬到表格",
        ],
        label_dark=True,
    )
    add_source(s, "Python Built-in Functions §map, filter, sorted")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── S11 · VS：map+filter vs List Comprehension ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="map + filter（函數式遺產）",
        right_title="List Comprehension（Pythonic）",
        left_items=[
            "list(map(lambda x: x*0.9,",
            "     filter(lambda x: x>100, prices)))",
            "要 list() 物化",
            "巢狀 lambda 難讀",
            "懂原理可用，但不推薦",
        ],
        right_items=[
            "[x*0.9 for x in prices if x > 100]",
            "讀起來就是需求本身",
            "CPython 內部優化，通常更快",
            "Guido 本人公開推薦的寫法",
            "Code review 會過的那版",
        ],
        title="兩種寫法，一個是 Python 偏愛的母語",
        summary="寫得出 map+filter 代表你懂原理；寫 Comprehension 代表你能過 review。",
        delta="可讀性\n單向勝",
    )
    add_source(s, "Guido van Rossum, python-dev archive 2005 · PEP 202")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── S12 · CODE：List Comprehension 三段式 ─────────
    s = _blank(prs)
    add_title(s, "List Comprehension 三格：輸出 / 來源 / 條件")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="三格骨架：讀時從 for 開始讀",
        code=(
            '# [  expr    for  x  in  iterable   if  cond  ]\n'
            '#    ↑輸出        ↑來源              ↑可選條件\n'
            '\n'
            'squares     = [x**2 for x in range(10)]\n'
            'big_squares = [x**2 for x in range(10) if x > 5]\n'
            '\n'
            '# 多行排版：expr / for / if 各一行\n'
            'cleaned_names = [\n'
            '    n.strip().lower()\n'
            '    for n in raw_names\n'
            '    if n and n.strip()\n'
            ']'
        ),
        bullets=[
            "輸出可以是任何 expression",
            "條件可連鎖：if a if b",
            "三秒講完 = 這行 comprehension\n在做什麼",
            "超過 80 字元換行排版",
        ],
        label_dark=True,
    )
    add_source(s, "PEP 202 List Comprehensions")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── S13 · CODE：Dict / Set / Generator Comprehension ─────────
    s = _blank(prs)
    add_title(s, "同骨架，換外框：四兄弟同根同源")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="{} + 冒號 → dict；{} 無冒號 → set；() → generator",
        code=(
            '# Dict Comprehension\n'
            'name_age = {u["name"]: u["age"] for u in users}\n'
            'inverse  = {v: k for k, v in name_age.items()}\n'
            '\n'
            '# Set Comprehension（自動去重）\n'
            'domains = {email.split("@")[1] for email in emails}\n'
            '\n'
            '# Generator Expression（懶求值，不物化）\n'
            'total = sum(x**2 for x in range(10**7))'
        ),
        bullets=[
            "方括號 [] → list",
            "大括號 + 冒號 → dict",
            "大括號無冒號 → set",
            "圓括號 () → generator",
            "看括號就知道結果型態",
        ],
        label_dark=True,
    )
    add_source(s, "PEP 274 Dict · PEP 289 Generator Expressions")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── S14 · VS-CODE：巢狀 Comprehension 界線 ─────────
    s = _blank(prs)
    add_title(s, "兩層可讀，三層就拆：可讀性優先於單行主義")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.2),
        label="OK：兩層扁平化，無巢狀條件",
        code=(
            'matrix = [[1, 2, 3], [4, 5, 6]]\n'
            'flat = [x for row in matrix for x in row]\n'
            '# → [1, 2, 3, 4, 5, 6]'
        ),
        bullets=[
            "讀順序：for 由左而右 = 由外而內",
            "兩層且無 if 仍可讀",
        ],
        label_dark=False,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(3.7),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.6),
        label="NOT OK：三層 + 多條件，拆回 for",
        code=(
            '# 不要這樣寫\n'
            'result = [f(x, y, z)\n'
            '          for x in xs for y in ys for z in zs\n'
            '          if x > 0 if y < z if g(x, y)]\n'
            '\n'
            '# 改寫：可加 log、可 breakpoint、可除錯\n'
            'result = []\n'
            'for x in xs:\n'
            '    for y in ys:\n'
            '        ...'
        ),
        bullets=[
            "三層 + 多條件 → 拆回 for",
            "for 可加 log、可 breakpoint",
            "PEP 20：寫給人看的優先",
        ],
        label_dark=True,
    )
    add_source(s, "The Zen of Python (PEP 20) · Readability counts.")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── S15 · MATRIX 2×2：選型四象限 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "List Comprehension",
         "sub": "小量 × 簡單邏輯\n90% 日常場景首選\n簡潔、快、Pythonic",
         "highlight": True},
        {"text": "Generator Expression",
         "sub": "大量 × 簡單邏輯\n讀不進記憶體 / 只跑一遍\n用 () 省 memory"},
        {"text": "def + 普通 for",
         "sub": "小量 × 複雜邏輯\n要 log / 例外處理\n可讀優先"},
        {"text": "def + yield + pipeline",
         "sub": "大量 × 複雜邏輯\n串流處理\nCh06 例外 + Ch10 整合"},
    ], title="四種寫法的選型：資料量 × 可讀性需求")
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "橫軸：邏輯複雜度 低 → 高    縱軸：資料量 小 → 大",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "Ch03 課堂歸納")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── S16 · TABLE：時間複雜度直覺 ─────────
    s = _blank(prs)
    draw_editorial_table(
        s,
        header=["寫法", "複雜度", "相對速度", "何時該警覺"],
        rows=[
            ["for x in xs: ...（單層）", "O(N)", "基準 1×",
             "百萬級以下都沒事"],
            ["[f(x) for x in xs]", "O(N)", "~0.7×",
             "CPython 內部優化，通常更快"],
            ["巢狀 for x/for y", "O(N×M)", "N×M ×",
             "兩邊千級就要小心"],
            ["x in some_list", "O(N) 每次", "慢",
             "改 set / dict 變 O(1)"],
            ["找到就 break", "O(k), k ≤ N", "視情況",
             "比全掃快，尤其 early-exit"],
        ],
        col_widths=[1.6, 1.1, 1.1, 2.2],
        title="時間複雜度直覺：線性是命，巢狀是罪，早 break 是救贖",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.6),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "三直覺：單層 OK、巢狀看乘積、查找用 set / dict。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "CPython source · Beazley, Python Essential Reference §Perf")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── S17 · VS-CODE：List vs Generator ─────────
    s = _blank(prs)
    add_title(s, "差一個括號，記憶體差 1000 倍：List vs Generator")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.3),
        label="LIST：[] 立刻物化全部",
        code=(
            'sq = [x**2 for x in range(10_000_000)]\n'
            'total = sum(sq)\n'
            '# 常駐記憶體：~ 350 MB'
        ),
        bullets=[
            "立刻算完全部、存起來",
            "可多次迭代、可索引",
            "佔約 350 MB",
        ],
        label_dark=False,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(3.8),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.5),
        label="GENERATOR：() 懶求值，整段流",
        code=(
            'sq = (x**2 for x in range(10_000_000))\n'
            'total = sum(sq)\n'
            '# 常駐記憶體：~ 128 B'
        ),
        bullets=[
            "懶求值，yield one by one",
            "只能迭代一次",
            "常駐僅 ~ 128 B",
            "Pandas chunked read = 同概念",
        ],
        label_dark=True,
    )
    add_source(s, "PEP 289 Generator Expressions · sys.getsizeof 實測")
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───────── S18 · PHOTO-CODE：預告 Pandas apply ─────────
    s = _blank(prs)
    add_title(s, "Lambda 的真實棲地：Pandas apply（Ch08 主場）")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.8), h=Inches(5.0),
        slot_name="pandas.DataFrame.apply 官方文件",
        description="pandas.pydata.org DataFrame.apply\n文件首屏（含 URL 列 + 函式簽章）",
        url_hint="https://pandas.pydata.org/docs/reference/api/pandas.DataFrame.apply.html",
        placeholder_id="Ch03_S18_pandas_apply",
        registry=image_registry,
        size_hint="1400×1200 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.7), y=Inches(1.3),
        w=Inches(6.0), h=Inches(5.0),
        label="Ch03 → Ch08：同一個概念換載體",
        code=(
            '# Ch03 你學的 map\n'
            'list(map(lambda x: x * 0.9, prices))\n'
            '\n'
            '# Ch08 你天天寫的 apply\n'
            'df["price"] = df["price"].apply(\n'
            '    lambda x: x * 0.9\n'
            ')\n'
            '\n'
            'df["tier"] = df["amount"].apply(\n'
            '    lambda x: "high" if x > 1000 else "low"\n'
            ')'
        ),
        bullets=[
            "apply = 把 lambda 套到每列",
            "與 map 概念完全一致",
            "Ch03 練好的 lambda 功力直接可用",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Function application")
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ───────── S19 · PYRAMID：收束 ─────────
    s = _blank(prs)
    draw_thesis_hierarchy(
        s,
        blocks=[
            {"heading": "Python 特殊機制三件套",
             "items": [
                 "Lambda：即用即丟，一行為限",
                 "map / filter / sorted(key=)：聲明式處理 iterable",
                 "Comprehension：[expr for x in iter if cond]",
             ]},
            {"heading": "設計紀律",
             "items": [
                 "可變預設值 → None + 首行初始化",
                 "global / nonlocal 能不用就不用",
                 "巢狀 Comprehension 超過兩層就拆回 for",
             ]},
        ],
        title="Ch03 收束：參數四件 + Lambda 三件套 + Comprehension 三格",
        thesis="Ch04 起進入 OOP：把這些函式技巧升級為『物件方法』，為 Ch10 DataCleaner 奠基。",
    )
    add_source(s, "Ch03 module synthesis")
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───────── S20 · SILENT：收尾 ─────────
    s = _blank(prs)
    draw_silent_page(s, "好程式不是寫得多，\n是說得準。")
    add_footer(s, MODULE_CODE, 20, N_CONTENT, dark_bg=True)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
