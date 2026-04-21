"""F1 deck — 計算機概論與 OS 角色（完整版 v1.1）
21 content slides + cover + copyright.

Governing thought:
    資料處理的瓶頸永遠在硬體——RAM 有多大、
    資料在哪一層跑，決定你的程式會不會卡。

受眾：僅有 Python 基礎的非理工資訊背景學生（商管、人文、行銷、
      行政、設計、醫護等轉職或自學者）。所以本章大量用生活類比
      （廚房、便當、倉庫、物流），不走硬體術語路線。

Aligned to chapters/F1_計算機概論與OS角色/00_skeleton.yaml
  · 6 Learning Objectives × 5 Common Pitfalls
  · Teaching-track 7 原型：MOTIVATION / CONCEPT / MECHANISM / EXAMPLE / PITFALL / PRACTICE / CHECKPOINT

三件核心 take-away：
  1. RAM 是工作檯（→ 銜接 S2 chunksize / S4 大資料）
  2. I/O 很慢（→ 銜接 S1 向量化 / S2 read_csv）
  3. 路徑要用 /（→ 銜接 S2 pathlib）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source, add_rect,
    set_solid_fill, set_no_fill, set_line, set_no_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_code_panel, draw_vs_two_col, draw_pyramid_stack,
    draw_thesis_hierarchy, draw_three_blocks_flow, draw_image_placeholder,
    draw_delta_badge, draw_emphasis_pill, draw_inverted_thesis_box,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "F1"
MODULE_TITLE = "計算機概論與 OS 角色"
MODULE_SUBTITLE = "廚房三角色 × OS 管委會——資料處理的底層心智模型"
TIME_MIN = 90
N_CONTENT = 21

# Pitfall 頁用的警示色（僅 PITFALL 頁限定）
RED_ERROR = RGBColor(0xC6, 0x28, 0x28)
GREEN_OK = RGBColor(0x2E, 0x7D, 0x32)


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _draw_three_io(slide, title, input_text, process_text, output_text,
                   bottom_note=""):
    """EXAMPLE-I/O 三欄版：Input | Process | Output。"""
    add_title(slide, title)
    top = Inches(1.5)
    col_h = Inches(4.2)
    total_w = T.SLIDE_W - 2 * T.MARGIN_X
    gap = Inches(0.2)
    col_w = (total_w - gap * 2) / 3
    labels = ["Input", "Process", "Output"]
    bodies = [input_text, process_text, output_text]
    for i, (lbl, body) in enumerate(zip(labels, bodies)):
        x = T.MARGIN_X + i * (col_w + gap)
        hdr = add_rect(slide, x, top, col_w, Inches(0.45))
        set_solid_fill(hdr, T.PRIMARY)
        set_no_line(hdr)
        add_textbox(
            slide, x, top, col_w, Inches(0.45),
            lbl,
            font_size=T.FONT_BODY, color=T.WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        body_y = top + Inches(0.45)
        body_h = col_h - Inches(0.45)
        box = add_rect(slide, x, body_y, col_w, body_h)
        set_no_fill(box)
        set_line(box, T.PRIMARY, 1.0)
        add_textbox(
            slide, x + Inches(0.2), body_y + Inches(0.15),
            col_w - Inches(0.4), body_h - Inches(0.3),
            body,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            family=T.FONT_MONO, line_spacing=1.35,
        )
    for i in range(2):
        x = T.MARGIN_X + (i + 1) * col_w + i * gap + gap * 0.1
        add_textbox(
            slide, x, top + col_h / 2 - Inches(0.25),
            gap * 0.8, Inches(0.5),
            "→",
            font_size=Pt(22), color=T.PRIMARY, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
    if bottom_note:
        add_textbox(
            slide, T.MARGIN_X, top + col_h + Inches(0.25),
            T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
            bottom_note,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            align=PP_ALIGN.CENTER, italic=True,
        )


def _draw_pitfall(slide, title, wrong_label, wrong_code, right_label,
                  right_code, why):
    """PITFALL 頁：左紅錯 / 右綠對 / 下方 why。"""
    add_title(slide, title)
    top = Inches(1.4)
    col_h = Inches(4.2)
    col_w = Inches(5.8)
    gap = Inches(0.3)
    total = col_w * 2 + gap
    left_x = (T.SLIDE_W - total) / 2
    right_x = left_x + col_w + gap

    for x, label, code, accent, mark in [
        (left_x, wrong_label, wrong_code, RED_ERROR, "✗"),
        (right_x, right_label, right_code, GREEN_OK, "✓"),
    ]:
        hdr = add_rect(slide, x, top, col_w, Inches(0.5))
        set_solid_fill(hdr, accent)
        set_no_line(hdr)
        add_textbox(
            slide, x + Inches(0.2), top, col_w - Inches(0.4), Inches(0.5),
            f"{mark}  {label}",
            font_size=T.FONT_BODY, color=T.WHITE, bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        body_y = top + Inches(0.5)
        body_h = col_h - Inches(0.5)
        box = add_rect(slide, x, body_y, col_w, body_h)
        set_solid_fill(box, RGBColor(0xF9, 0xF9, 0xF9))
        set_line(box, accent, 1.2)
        add_textbox(
            slide, x + Inches(0.25), body_y + Inches(0.2),
            col_w - Inches(0.5), body_h - Inches(0.4),
            code,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            family=T.FONT_MONO, line_spacing=1.45,
        )

    why_y = top + col_h + Inches(0.2)
    why_box = add_rect(slide, T.MARGIN_X, why_y,
                       T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.55))
    set_solid_fill(why_box, T.PRIMARY)
    set_no_line(why_box)
    add_textbox(
        slide, T.MARGIN_X + Inches(0.3), why_y,
        T.SLIDE_W - 2 * T.MARGIN_X - Inches(0.6), Inches(0.55),
        f"為什麼：{why}",
        font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def _draw_bridge_note(slide, text, y_inches=6.6):
    """底部橫幅：銜接下一章的提示 pill。"""
    add_textbox(
        slide, T.MARGIN_X, Inches(y_inches),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        text,
        font_size=T.FONT_SMALL, color=T.PRIMARY, bold=True, italic=True,
        align=PP_ALIGN.CENTER,
    )


def build_ch01(output_path, image_registry=None):
    """Build F1 deck; 21 content slides + cover + copyright."""
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── F1 · MOTIVATION — Excel 卡 3 分鐘 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "你打開 5 萬筆訂單的 Excel——\n電腦卡 3 分鐘才回應你。",
        data_card={
            "label": "這不是你電腦爛",
            "stat": "RAM 在抗議",
            "caption": "記憶體不夠\n系統在偷渡硬碟",
        },
    )
    add_source(s, "非理工學員最常見的痛點 · 行銷 / 財務 / 行政每月都遇到")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # ───────── F2 · ASK — 同一行程式，為什麼差這麼多 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "同一行 pd.read_csv()——\n為什麼有時秒回、有時凍結？",
        data_card={
            "label": "線索",
            "stat": "RAM",
            "caption": "不是程式問題\n是『工作檯』不夠大",
        },
    )
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── F3 · SILENT — 一句話立論 ─────────
    s = _blank(prs)
    draw_silent_page(
        s,
        "寫得快的程式，是懂硬體的人寫的。\n不是語法炫的人寫的。",
    )
    add_footer(s, MODULE_CODE, 3, N_CONTENT, dark_bg=True)

    # ───────── F4 · CONCEPT-CARD — 廚房三角色（類比主軸）─────────
    s = _blank(prs)
    add_title(s, "電腦裡的廚房：廚師 / 工作檯 / 倉庫")
    left_x = T.MARGIN_X
    left_w = Inches(6.0)
    add_textbox(
        s, left_x, Inches(1.4), left_w, Inches(0.5),
        "一個類比，記三件事",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, left_x, Inches(1.95), left_w, Inches(3.9),
        "• CPU — 廚師\n"
        "  真正做事的人；越多核＝越多廚師\n\n"
        "• RAM — 工作檯\n"
        "  料要放到檯面上才能切；檯面大小決定一次能做多少菜\n\n"
        "• Storage（SSD/HDD）— 倉庫\n"
        "  容量大但遠；廚師看不到倉庫，料要先『叫上來』",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    add_textbox(
        s, left_x, Inches(5.95), left_w, Inches(0.55),
        "記憶點：料不搬到工作檯，廚師碰不到 → 這就是「讀檔」在做的事。",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True, italic=True,
    )
    draw_image_placeholder(
        s, Inches(7.0), Inches(1.3), Inches(5.7), Inches(4.8),
        slot_name="廚房三角色示意圖",
        description="廚師（CPU）在工作檯（RAM）前切菜；旁邊大型倉庫（Storage）標示『遠、慢、容量大』；手推車箭頭代表資料搬運。避免硬體照片，用插畫。",
        url_hint="自製插畫或參考 Von Neumann 架構圖改造成廚房隱喻",
        size_hint="1400×1120 px",
        placeholder_id="F1_S4_kitchen_trio",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── F5 · EXAMPLE-I/O — 一行 read_csv 的旅程 ─────────
    s = _blank(prs)
    _draw_three_io(
        s,
        "範例：一行 read_csv() 背後，OS 幫你跑了幾步",
        input_text="# 你只寫這一行\n\n"
                   "df = pd.read_csv(\n"
                   "  'sales.csv'\n"
                   ")",
        process_text="Python 直譯器\n"
                     "   ↓ 翻譯\n"
                     "OS System Call\n"
                     "   ↓ 「幫我開檔」\n"
                     "檔案系統找到檔\n"
                     "   ↓\n"
                     "I/O 控制器讀磁碟\n"
                     "   ↓\n"
                     "資料搬到 RAM",
        output_text="DataFrame\n"
                    "  躺在 RAM 裡\n\n"
                    "廚師（CPU）\n"
                    "  才能開始做菜\n\n"
                    "（塞不下→OOM）",
        bottom_note="你寫一行，OS 幫你跑七步——這就是為什麼 OS 是主角，不是配角。",
    )
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── F6 · CHECKPOINT — 廚房角色連連看 ─────────
    s = _blank(prs)
    add_title(s, "Check Point · 角色連連看（30 秒）")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 6 / 21 · 不看投影片作答",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.3),
        "請把左邊的廚房角色配到右邊的電腦零件：\n\n"
        "   廚師     ←→    ? \n"
        "   工作檯   ←→    ? \n"
        "   倉庫     ←→    ? \n\n"
        "延伸題：\n"
        "  「檯面放不下料」在電腦上對應到什麼錯誤訊息？\n"
        "  「叫貨要等 5 分鐘」對應到什麼瓶頸？",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.35), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "答案：廚師=CPU · 工作檯=RAM · 倉庫=Storage · 檯面滿=OOM · 叫貨慢=I/O Bound",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───────── F7 · CONCEPT-CARD — Bit/Byte 量級（生活對照）─────────
    s = _blank(prs)
    add_title(s, "Bit / Byte 量級：你每天都在用的單位")
    draw_editorial_table(
        s,
        header=["單位", "大小", "生活對照", "資料場景"],
        rows=[
            ["1 Byte", "8 Bit", "1 個英文字母", "—"],
            ["1 KB", "≈1 千 Byte", "半頁 Word 文字", "設定檔 / 小 JSON"],
            ["1 MB", "≈1 百萬 Byte", "1 張手機照片 / 1 分鐘 MP3", "小型 CSV"],
            ["1 GB", "≈10 億 Byte", "1 部標清電影 / 250 首歌", "桌機 RAM 單位 / 中型資料"],
            ["1 TB", "≈1 兆 Byte", "250 部 HD 電影", "企業日誌 / 小型資料倉儲"],
            ["1 PB", "≈1000 TB", "—", "Netflix / 雲端倉儲級"],
        ],
        top=1.3,
        col_widths=[1.2, 1.7, 3.0, 2.3],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.1), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "唯一要背的數字：1 個 float64 = 8 Bytes。算 RAM 就靠它。",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── F8 · CHECKPOINT — 量級估算 ─────────
    s = _blank(prs)
    add_title(s, "Check Point · RAM 能裝多少筆資料")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(4), Inches(0.4),
        "進度 8 / 21 · 心算不查 Google",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.5),
        "Q1  1 億筆 float64 佔多少 RAM？\n"
        "       提示：1 個 float64 = 8 Bytes\n\n"
        "Q2  1000 萬列 × 50 欄（全 float64）呢？\n"
        "       這種 DataFrame 能在 8GB 筆電上活嗎？\n\n"
        "Q3  如果把 float64 改成 float32，記憶體降幾成？\n"
        "       （這就是 S2 會教的 dtype downcast）",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.3), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "答案：Q1≈800MB · Q2≈4GB（勉強但危險）· Q3 降 50%。",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── F9 · CONCEPT-CARD — RAM 是工作檯（1000× 差距）─────────
    s = _blank(prs)
    add_title(s, "為什麼資料一定要『先搬到 RAM』")
    draw_vs_two_col(
        s,
        left_title="工作檯（RAM）的現場",
        right_title="倉庫（Storage）的距離",
        left_items=[
            "廚師（CPU）只看得到這裡",
            "拿料要 1 秒（奈秒級 ns）",
            "容量小（4GB / 8GB / 16GB）",
            "揮發性：關機就空",
        ],
        right_items=[
            "廚師碰不到，要叫貨",
            "拿料要 1000 秒（毫秒級 ms）",
            "容量大（TB 級）",
            "永久保存：關機還在",
        ],
        summary="所以 pd.read_csv() 在做的事，本質上就是「叫貨到工作檯」——它在搬家。",
        delta="1000×",
    )
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── F10 · PITFALL (P1) — OOM (→ 銜接 S2) ─────────
    s = _blank(prs)
    _draw_pitfall(
        s,
        title="陷阱 · 一口氣讀 5GB CSV → 系統凍結",
        wrong_label="整檔塞 RAM",
        wrong_code="import pandas as pd\n\n"
                   "df = pd.read_csv(\n"
                   "  'big_5gb.csv'\n"
                   ")\n\n"
                   "# MemoryError\n"
                   "# 或系統卡死 10 分鐘\n"
                   "# （Swap 被觸發：\n"
                   "#  拿硬碟假裝 RAM）",
        right_label="分批串流 / 降 dtype",
        right_code="for chunk in pd.read_csv(\n"
                   "    'big_5gb.csv',\n"
                   "    chunksize=100_000,\n"
                   "    dtype={'id': 'int32'}\n"
                   "):\n"
                   "    process(chunk)\n\n"
                   "# 每次只吃 ~100MB\n"
                   "# 系統呼吸順暢",
        why="RAM 塞爆就觸發 Swap（拿硬碟假裝 RAM），系統會卡到你以為當機",
    )
    _draw_bridge_note(s, "→ 銜接 S2：chunksize、usecols、dtype downcast 三招組合拳")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── F11 · CONCEPT-CARD — I/O 很慢（核心觀念）─────────
    s = _blank(prs)
    add_title(s, "為什麼『讀檔』比『算』慢 100 倍")
    draw_vs_two_col(
        s,
        left_title="CPU 算一次（ns 級）",
        right_title="硬碟讀一次（ms 級）",
        left_items=[
            "加一次 = 約 1 奈秒（0.000000001 秒）",
            "一秒可以算幾十億次",
            "像廚師切菜：一秒切一百刀",
            "這一層很快，不是瓶頸",
        ],
        right_items=[
            "讀一次 SSD = 約 0.1 毫秒",
            "一秒只能讀幾萬次",
            "像倉庫叫貨：一趟 5 分鐘",
            "這一層才是瓶頸所在",
        ],
        summary="所以你的 Pandas 卡住，90% 是在等讀檔，不是在等計算。",
        delta="100-1000×",
    )
    _draw_bridge_note(s, "→ 銜接 S1：Python 層運算也慢 50×，向量化把事情交給 C 層")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── F12 · MECHANISM-FLOW — OS 三大職責 ─────────
    s = _blank(prs)
    draw_flow_chain(
        s,
        nodes=[
            {"label": "記憶體管理", "sub": "RAM / Swap / OOM",
             "caption": "誰先、誰後、誰被 kill"},
            {"label": "檔案系統", "sub": "路徑 / 權限 / I/O",
             "caption": "Windows \\ vs Linux /", "highlight": True},
            {"label": "行程排程", "sub": "Process / Thread / GIL",
             "caption": "多核怎麼分工"},
        ],
        title="OS 三大職責：你碰不到硬體，都由 OS 代理",
        y=2.4,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.3), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
        "口訣：Python 只是租客，OS 才是管委會。",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── F13 · PITFALL (P2) — 路徑跨平台 ─────────
    s = _blank(prs)
    _draw_pitfall(
        s,
        title="陷阱 · 路徑字串：Windows 能跑，Linux 爆炸",
        wrong_label="硬拼反斜線字串",
        wrong_code="# 在 Windows 寫的\n"
                   "path = 'data\\\\2024\\\\x.csv'\n\n"
                   "df = pd.read_csv(path)\n\n"
                   "# 部署到雲端 Linux：\n"
                   "# FileNotFoundError\n"
                   "# 原因：\\ 在 Linux\n"
                   "# 不是分隔符",
        right_label="用 pathlib.Path",
        right_code="from pathlib import Path\n\n"
                   "path = Path('data') \\\n"
                   "     / '2024' / 'x.csv'\n\n"
                   "df = pd.read_csv(path)\n\n"
                   "# Windows / Linux / Mac\n"
                   "# 三平台自動正確",
        why="DA/DE 工作幾乎都在 Linux 雲端跑；第一次部署就會被反斜線咬",
    )
    _draw_bridge_note(s, "→ 銜接 S2：pathlib 會是你讀每一個 CSV / Parquet 的起手式")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── F14 · PRACTICE — 估算你的筆電能讀多大 ─────────
    s = _blank(prs)
    add_title(s, "Practice · 估算你自己的筆電能讀多大 CSV")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 14 / 21 · 1 分鐘心算",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(10.8), Inches(4.2),
        "① 你的筆電 RAM 幾 GB？（常見：8 / 16 / 32）\n\n"
        "② 安全額度 = RAM × 0.3\n"
        "       為什麼 ×0.3？OS 要、Chrome 要、Pandas 載入會膨脹 2–3×\n\n"
        "③ 算一筆資料大小：欄數 × 8 Bytes（假設都是 float64）\n\n"
        "④ 安全筆數上限 = 安全額度 ÷ 一筆大小\n\n"
        "   例：16GB RAM，50 欄 → 安全額度 ≈ 4.8GB ÷ (50×8) ≈ 1200 萬列",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    _draw_bridge_note(s, "超過這個數？就需要 S2 的 chunksize 上場了")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── F15 · CONCEPT-CARD — I/O vs CPU Bound ─────────
    s = _blank(prs)
    add_title(s, "I/O Bound vs CPU Bound：選錯策略就白工")
    draw_vs_two_col(
        s,
        left_title="I/O Bound（瓶頸在『等』）",
        right_title="CPU Bound（瓶頸在『算』）",
        left_items=[
            "特徵：CPU 閒著、在等網路/磁碟",
            "典型：爬蟲、讀大檔、API 呼叫",
            "策略：async / threading",
            "多開行程沒用（還在等）",
        ],
        right_items=[
            "特徵：CPU 100%、風扇狂轉",
            "典型：矩陣運算、特徵工程",
            "策略：向量化 / multiprocessing",
            "多開行程真的會快",
        ],
        summary="判斷法：打開工作管理員，看 CPU 飆滿還是網路/硬碟燈閃爍。",
    )
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── F16 · PRACTICE-PROMPT — 判斷三個任務 ─────────
    s = _blank(prs)
    add_title(s, "Practice · 這三個任務該用什麼策略？")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 16 / 21 · 2 分鐘小組討論",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(10.8), Inches(4.2),
        "① 爬 1000 個商品頁面抓價格\n"
        "     → I/O or CPU？用什麼？\n\n"
        "② 10 億筆數值做 groupby 加總\n"
        "     → I/O or CPU？用什麼？\n\n"
        "③ 把一個 5GB CSV 讀進 Pandas\n"
        "     → I/O or CPU？用什麼？",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.3), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "提示：① I/O→async · ② CPU→向量化（S1） · ③ I/O→chunksize 串流（S2）。",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── F17 · PITFALL (P3) — 用錯武器 ─────────
    s = _blank(prs)
    _draw_pitfall(
        s,
        title="陷阱 · I/O Bound 上 multiprocessing → 白工",
        wrong_label="開 8 個行程爬網頁",
        wrong_code="from multiprocessing import Pool\n\n"
                   "with Pool(8) as p:\n"
                   "    results = p.map(\n"
                   "      fetch_url, urls\n"
                   "    )\n\n"
                   "# 跟單行程差不多慢\n"
                   "# 因為 8 個行程\n"
                   "# 全部在等網路",
        right_label="用 async / threading",
        right_code="import asyncio, aiohttp\n\n"
                   "async def main():\n"
                   "    async with aiohttp.\\\n"
                   "       ClientSession() as s:\n"
                   "        tasks = [fetch(s,u)\n"
                   "                 for u in urls]\n"
                   "        return await asyncio.\\\n"
                   "               gather(*tasks)\n\n"
                   "# 1000 個請求並行等",
        why="等網路時 CPU 本來就閒——多開 CPU 沒意義，要的是並行『等』",
    )
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───────── F18 · CHECKPOINT — 三核心回顧 ─────────
    s = _blank(prs)
    add_title(s, "Check Point · 離開教室前的三題")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 18 / 21 · 用自己的話回答",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(10.8), Inches(4.2),
        "Q1  用一句話說：RAM 在整台電腦裡扮演什麼角色？\n"
        "       （提示：類比 + 電腦術語各一遍）\n\n"
        "Q2  為什麼『讀一個大檔』比『算 10 億次加法』慢？\n"
        "       （提示：ns vs ms）\n\n"
        "Q3  你在 Windows 寫的程式要上雲，路徑該怎麼寫？\n"
        "       （提示：不是字串）",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.3), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "能答出這三題，本章就畢業了。答不出來的，回去翻 F6 / F11 / F13。",
        font_size=T.FONT_SMALL, color=T.PRIMARY, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ───────── F19 · CONCEPT-CARD — 三件帶走的反射 ─────────
    s = _blank(prs)
    draw_three_blocks_flow(
        s,
        title="本章三件帶走的事（會出現在後面每一章）",
        blocks=[
            {"heading": "① RAM 是工作檯",
             "items": [
                 "檔案大小 ≤ RAM × 0.3",
                 "超過就要分批 chunk",
                 "→ 銜接 S2 chunksize",
                 "→ 銜接 S4 大資料 EDA",
             ]},
            {"heading": "② I/O 很慢",
             "items": [
                 "讀檔比算術慢 100×",
                 "能少讀一次就少讀",
                 "→ 銜接 S1 向量化",
                 "→ 銜接 S2 read_csv",
             ]},
            {"heading": "③ 路徑用 /",
             "items": [
                 "pathlib.Path 而非字串",
                 "Windows 上也用 /",
                 "雲端 Linux 不咬你",
                 "→ 銜接 S2 pathlib 起手式",
             ]},
        ],
    )
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───────── F20 · PYRAMID — 收束 ─────────
    s = _blank(prs)
    draw_pyramid_stack(
        s,
        title="F1 收束：從硬體到 Python 的全景",
        layers=[
            {"name": "Python 程式",
             "caption": "你寫的那一行"},
            {"name": "作業系統 OS",
             "caption": "記憶體 / 檔案 / 行程的管委會"},
            {"name": "硬體 CPU + RAM + Storage",
             "caption": "真正做事的地方（廚師 / 工作檯 / 倉庫）"},
            {"name": "Bit / Byte",
             "caption": "一切資料的本體"},
        ],
        thesis="懂硬體 → 懂 OS → 才寫得動 Python 資料處理。",
    )
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ───────── F21 · 銜接 F2 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "RAM 省一半的秘訣——\n不在硬體，在你選哪個資料結構。",
        data_card={
            "label": "下一章 F2",
            "stat": "Python\n資料結構",
            "caption": "list vs tuple\nvs dict vs generator\n記憶體差幾倍？",
        },
    )
    _draw_bridge_note(s, "F2：Python 核心與資料結構深化——把本章的 RAM 觀念落地到程式寫法")
    add_footer(s, MODULE_CODE, 21, N_CONTENT)

    # Copyright
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
