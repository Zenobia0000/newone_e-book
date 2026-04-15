"""S4 deck — 時間序列與 EDA 實戰
23 content slides + cover + copyright (120 min, teaching track).

Governing thought:
    商業問題都有時間軸 —— 學會把時間變成一條軸、把 EDA 變成一種直覺，
    才能產出報表。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_rect, add_textbox, add_title, add_source,
    set_solid_fill, set_no_line, set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_vs_two_col, draw_code_panel, draw_thesis_hierarchy,
    draw_inverted_thesis_box, draw_image_placeholder,
    draw_flow_chain, draw_pyramid_stack, draw_three_blocks_flow,
    draw_emphasis_pill,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "S4"
MODULE_TITLE = "時間序列與 EDA 實戰"
MODULE_SUBTITLE = "從 .dt / resample / rolling 到一張月度經營報表"
TIME_MIN = 120
N_CONTENT = 23


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_s4(output_path, image_registry=None):
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── S1 · SILENT · 開場金句 ─────────
    s = _blank(prs)
    draw_silent_page(s, "數字不會說話，\n時間軸與 EDA 會。")
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ───────── S2 · MOTIVATION · 老闆四問 ─────────
    s = _blank(prs)
    add_title(s, "老闆四問：每一題都帶著時間維度")
    # Left: pain story
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(5.8), Inches(0.5),
        "真實場景：週一早會",
        font_size=T.FONT_SUBTITLE, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(1.9), Inches(5.8), Inches(4.2),
        "老闆看著儀表板，隨口丟四個問題 —\n"
        "你回得出「讓我查一下」還是「這張圖」？\n\n"
        "不會處理時間軸的 DA，\n"
        "永遠停在「等我回去跑一下」。",
        font_size=T.FONT_BODY, color=T.CHARCOAL, line_spacing=1.6,
    )
    # Right: 4 questions matrix
    q_labels = [
        ("Q1", "上週 vs 上月差多少？"),
        ("Q2", "哪一天最好賣？"),
        ("Q3", "哪款在衰？"),
        ("Q4", "下個月預估？"),
    ]
    qx = Inches(7.0)
    qy = Inches(1.3)
    qw = Inches(5.7)
    qh = Inches(1.1)
    gap = Inches(0.15)
    for i, (tag, q) in enumerate(q_labels):
        y = qy + i * (qh + gap)
        rect = add_rect(s, qx, y, qw, qh)
        set_no_fill(rect)
        set_line(rect, T.PRIMARY, 1.0)
        add_textbox(
            s, qx + Inches(0.25), y + Inches(0.1), Inches(0.8), qh - Inches(0.2),
            tag, font_size=T.FONT_SUBTITLE, color=T.PRIMARY, bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            s, qx + Inches(1.0), y + Inches(0.1), qw - Inches(1.2), qh - Inches(0.2),
            q, font_size=T.FONT_BODY, color=T.CHARCOAL,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    add_source(s, "改編自真實零售電商週會訪談 · 2024")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── S3 · ASK · 入門提問 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "你拿到一份訂單 CSV，\n第一件事是 head() 還是……？",
        data_card={
            "label": "90% 新手會直接 head()",
            "stat": "正解",
            "caption": "parse_dates + set_index\n+ sort_index —— 三連招",
        },
    )
    add_source(s, "Kaggle Learn Time Series Track 常見新手訪談")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───────── S4 · CONCEPT-CARD · 學習目標 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "LO1 · .dt accessor",
         "sub": "能用 .dt 拆年/月/日/星期\n向量化版 Excel 日期函式",
         "highlight": True},
        {"text": "LO2 · resample / groupby",
         "sub": "能分辨兩條路徑的差異\n並選對情境使用"},
        {"text": "LO3 · rolling",
         "sub": "能用移動視窗平滑資料\n解讀趨勢而非噪音",
         "highlight": True},
        {"text": "LO4+5 · EDA 三板斧 + 報表",
         "sub": "describe / value_counts / corr\n→ 產出月度經營報表"},
    ], title="本節結束後你能：寫出一張月度經營報表")
    add_source(s, "Bloom Taxonomy · Apply 層動詞：寫出 / 分辨 / 產出")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── S5 · CONCEPT-CARD · datetime64 與 .dt ─────────
    s = _blank(prs)
    add_title(s, "datetime64[ns]：pandas 的時間原語")
    # Left definition
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(5.8), Inches(0.5),
        "定義（一句話）",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(1.9), Inches(5.8), Inches(1.6),
        "datetime64[ns] 是 pandas 的時間資料型別；\n"
        ".dt 是它的 accessor，把 Timestamp 當成「有結構的欄位」來拆。",
        font_size=T.FONT_BODY, color=T.CHARCOAL, line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(3.8), Inches(5.8), Inches(0.4),
        "類比（但不同的是）",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.3), Inches(5.8), Inches(2.0),
        "像 Excel 的 YEAR() / MONTH() / WEEKDAY()，\n"
        "但不同的是：一行 code 一次拆四欄、百萬列秒算，\n"
        "且回傳是 Series 能直接當新欄位接下去用。",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True, line_spacing=1.5,
    )
    # Right code
    draw_code_panel(
        s,
        x=Inches(6.8), y=Inches(1.3),
        w=Inches(5.9), h=Inches(5.0),
        label="datetime64 長這樣",
        code=(
            'df["order_date"].dtype\n'
            '# datetime64[ns]\n'
            '\n'
            'ts = df["order_date"].iloc[0]\n'
            '# Timestamp("2024-03-15 10:23:00")\n'
            '\n'
            'ts.year, ts.month, ts.day\n'
            '# (2024, 3, 15)\n'
            '\n'
            '# 批次版（整個 Series）\n'
            'df["order_date"].dt.year\n'
            'df["order_date"].dt.day_name()'
        ),
        bullets=[
            "dtype 對了，.dt 才會動",
            "Timestamp 是單點；\nSeries.dt 是批次",
            ".dt 是向量化通道",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Time series / date functionality §Time/Date components")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── S6 · MECHANISM-FLOW · 時序標準開場 4 步 ─────────
    s = _blank(prs)
    draw_flow_chain(s, nodes=[
        {"label": "read_csv",
         "sub": "parse_dates=[...]",
         "caption": "讀檔就轉型"},
        {"label": "to_datetime",
         "sub": "補轉 / 補齊時區",
         "caption": "漏網之魚處理"},
        {"label": "set_index",
         "sub": "order_date → Index",
         "caption": "開啟時序 API", "highlight": True},
        {"label": "sort_index",
         "sub": "由舊到新",
         "caption": "resample 才會對"},
    ], title="時序版 hello world：4 步驟，少一步下游全錯",
       y=2.5)
    add_textbox(
        s, T.MARGIN_X, Inches(5.6),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.8),
        "每次拿到帶時間的 CSV，閉著眼睛先跑這四步 —— 之後所有 .dt / resample / rolling 才能動。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "pandas docs §IO · §Indexing with sort_index · Wes McKinney 3e §11.1")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───────── S7 · EXAMPLE-I/O · .dt 四連拆 ─────────
    s = _blank(prs)
    add_title(s, "範例：.dt 四連拆 — Excel 四個公式 vs pandas 四行")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="Input: order_date 一欄 → Output: 四欄衍生特徵",
        code=(
            '# Input\n'
            'df = pd.read_csv("orders.csv", parse_dates=["order_date"])\n'
            '\n'
            '# Process — 四連拆\n'
            'df["year"]    = df["order_date"].dt.year\n'
            'df["month"]   = df["order_date"].dt.month\n'
            'df["day"]     = df["order_date"].dt.day\n'
            'df["weekday"] = df["order_date"].dt.day_name()   # Monday..Sunday\n'
            '\n'
            '# Output（head）\n'
            '#   order_date   year  month  day  weekday\n'
            '#   2024-03-15   2024   3     15   Friday\n'
            '#   2024-03-16   2024   3     16   Saturday'
        ),
        bullets=[
            "一次四欄 = 向量化",
            "day_name() 傳英文；\n中文需 .dt.day_name(locale=...)",
            "is_weekend = dayofweek\n.isin([5,6])",
            "Excel 要四公式 + 下拉",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §dt accessor · PyData Best Practices 2024")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── S8 · CONCEPT-CARD · resample = 時間版 groupby ─────────
    s = _blank(prs)
    add_title(s, "resample：時間版的 groupby（專門認 DatetimeIndex）")
    # Left: one-sentence definition
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(5.9), Inches(0.5),
        "一句話定義",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(1.9), Inches(5.9), Inches(1.6),
        "resample(rule) = 把時間軸切成固定粒度的 bucket，\n"
        "再對每 bucket 做聚合（sum / mean / count…）。",
        font_size=T.FONT_BODY, color=T.CHARCOAL, line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(3.8), Inches(5.9), Inches(0.5),
        "前提（硬性）",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.3), Inches(5.9), Inches(2.0),
        "Index 必須是 DatetimeIndex；否則：\n"
        "TypeError: Only valid with DatetimeIndex\n\n"
        "忘了 set_index → resample 直接罷工。",
        font_size=T.FONT_BODY, color=T.CHARCOAL, line_spacing=1.5,
    )
    # Right: image placeholder (bucketing illustration)
    draw_image_placeholder(
        s,
        x=Inches(6.8), y=Inches(1.4),
        w=Inches(5.9), h=Inches(4.8),
        slot_name="resample 時間分桶示意",
        description=(
            "每日 Timestamp 串流（左側散點）\n"
            "→ 三個月 bucket（右側方塊）\n"
            "中間箭頭標 .resample('ME').sum()"
        ),
        url_hint="",
        placeholder_id="S4_S08_resample_bucketing",
        registry=image_registry,
        size_hint="1400×1000 px",
    )
    add_source(s, "pandas User Guide §Resampling · McKinney 3e §11.6")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── S9 · MECHANISM-FLOW · D / W / ME ─────────
    s = _blank(prs)
    draw_three_blocks_flow(s, blocks=[
        {"heading": "D — 每日",
         "items": ["rule = 'D'",
                   "`.resample('D').sum()`",
                   "日營收 / 日單量",
                   "最細粒度"]},
        {"heading": "W — 每週",
         "items": ["rule = 'W'",
                   "`.resample('W').mean()`",
                   "每週日結束為界",
                   "週報表常用"]},
        {"heading": "ME — 每月末",
         "items": ["rule = 'ME' (2.2+)",
                   "`.resample('ME').sum()`",
                   "舊寫法 'M' 被 deprecation",
                   "也有 QE / YE"]},
    ], title="三種粒度：D / W / ME —— 選對規則比寫對 code 重要",
       bottom_note="同理：QE（季末）、YE（年末）、H（小時）、min（分鐘）—— rule 是 pandas 官方字母表。")
    add_source(s, "pandas docs §Offset aliases · Release Notes 2.2 §Deprecations")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── S10 · EXAMPLE-I/O · resample 月營收 ─────────
    s = _blank(prs)
    add_title(s, "範例：resample('ME') — 日訂單 → 月營收（3 行）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="Input: 日訂單 (order_date 為 Index) → Output: 月營收 Series",
        code=(
            '# Input — 確保 index 是 DatetimeIndex 且已排序\n'
            'df = (pd.read_csv("orders.csv", parse_dates=["order_date"])\n'
            '        .set_index("order_date")\n'
            '        .sort_index())\n'
            '\n'
            '# Process — 一句話\n'
            'monthly_revenue = df["amount"].resample("ME").sum()\n'
            '\n'
            '# Output\n'
            '# order_date\n'
            '# 2024-01-31    1,240,500\n'
            '# 2024-02-29    1,388,200\n'
            '# 2024-03-31    1,510,800\n'
            '# Freq: ME, Name: amount, dtype: int64'
        ),
        bullets=[
            "三行出結果",
            "sum / mean / count / agg\n皆可接在 resample 後",
            "Index 是 Month-End Timestamp",
            "缺月會自動補 NaN\n（groupby 不會）",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Resampling §Upsampling/Downsampling")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── S11 · PITFALL · 沒 sort_index ─────────
    s = _blank(prs)
    add_title(s, "常見錯誤：忘了 sort_index —— 不噴錯、但給你錯數字")
    # Left wrong
    left_x = T.MARGIN_X
    right_x = T.MARGIN_X + Inches(6.2)
    top_y = Inches(1.3)
    box_h = Inches(4.3)
    # Wrong box
    wrong = add_rect(s, left_x, top_y, Inches(5.9), box_h)
    set_no_fill(wrong)
    set_line(wrong, T.GRAY_MID, 1.5)
    add_textbox(
        s, left_x, top_y + Inches(0.1), Inches(5.9), Inches(0.5),
        "✗ 錯誤寫法",
        font_size=T.FONT_BODY, color=T.CHARCOAL, bold=True, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s, left_x + Inches(0.3), top_y + Inches(0.7),
        Inches(5.3), box_h - Inches(1.0),
        '# 讀完直接 resample\n'
        'df = pd.read_csv("orders.csv",\n'
        '                 parse_dates=["order_date"])\n'
        'df = df.set_index("order_date")\n'
        '\n'
        '# 沒 sort_index()\n'
        'df["amount"].resample("ME").sum()\n'
        '# → 不會噴錯\n'
        '# → 但月界切錯，數字全歪\n'
        '# → review 時才發現',
        font_size=Pt(11), color=T.CHARCOAL, family=T.FONT_MONO, line_spacing=1.3,
    )
    # Right correct
    right = add_rect(s, right_x, top_y, Inches(5.9), box_h)
    set_no_fill(right)
    set_line(right, T.PRIMARY, 1.5)
    add_textbox(
        s, right_x, top_y + Inches(0.1), Inches(5.9), Inches(0.5),
        "✓ 正確寫法",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s, right_x + Inches(0.3), top_y + Inches(0.7),
        Inches(5.3), box_h - Inches(1.0),
        '# 標準 4 步到位\n'
        'df = (pd.read_csv("orders.csv",\n'
        '        parse_dates=["order_date"])\n'
        '        .set_index("order_date")\n'
        '        .sort_index())\n'
        '\n'
        'df["amount"].resample("ME").sum()\n'
        '# ✓ 月界正確\n'
        '# ✓ 可重現\n'
        '# ✓ 下游 rolling 也安全',
        font_size=Pt(11), color=T.CHARCOAL, family=T.FONT_MONO, line_spacing=1.3,
    )
    # Bottom explanation
    add_textbox(
        s, T.MARGIN_X, Inches(5.8),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.8),
        "為什麼：resample 依 Index 順序切桶；亂序時會用遇到的第一個時間當起點，結果每桶資料混亂 —— "
        "pandas 不檢查、不警告，是「靜默 bug」之王。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, line_spacing=1.4,
    )
    add_source(s, "pandas docs §DataFrame.sort_index · Stack Overflow 年度 pandas bug top 10")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── S12 · CONCEPT-CARD · rolling ─────────
    s = _blank(prs)
    add_title(s, "rolling：移動視窗 —— 平滑雜訊、看趨勢")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(5.9), Inches(0.5),
        "一句話定義",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(1.9), Inches(5.9), Inches(1.6),
        "rolling(window) 把相鄰 N 筆當一個視窗、\n"
        "視窗滑過整個序列，對每視窗做聚合。",
        font_size=T.FONT_BODY, color=T.CHARCOAL, line_spacing=1.5,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(3.8), Inches(5.9), Inches(0.5),
        "類比（但不同的是）",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(4.3), Inches(5.9), Inches(2.0),
        "像股票的 MA5 / MA20 —— 但不同的是：\n"
        "rolling 輸出是同長度 Series（前 N-1 列為 NaN），\n"
        "可直接當 df 新欄位接下去用，不限金融。",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True, line_spacing=1.5,
    )
    # Right illustration (text-based since no external image)
    ill_x = Inches(6.8)
    ill_y = Inches(1.4)
    add_textbox(
        s, ill_x, ill_y, Inches(5.9), Inches(0.5),
        "視窗示意（window = 3）",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True, align=PP_ALIGN.CENTER,
    )
    # Row of boxes representing days
    box_w = Inches(0.7)
    box_h = Inches(0.7)
    start_x = ill_x + Inches(0.2)
    row_y = ill_y + Inches(0.9)
    for i in range(7):
        bx = start_x + i * (box_w + Inches(0.05))
        r = add_rect(s, bx, row_y, box_w, box_h)
        set_no_fill(r)
        set_line(r, T.PRIMARY, 1.0)
        add_textbox(
            s, bx, row_y, box_w, box_h,
            f"D{i+1}", font_size=T.FONT_SMALL, color=T.CHARCOAL,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
    # Window highlight over D3-D5
    win = add_rect(s, start_x + 2 * (box_w + Inches(0.05)) - Inches(0.05),
                   row_y - Inches(0.1),
                   3 * box_w + 2 * Inches(0.05) + Inches(0.1),
                   box_h + Inches(0.2))
    set_no_fill(win)
    set_line(win, T.PRIMARY, 2.5)
    add_textbox(
        s, ill_x, row_y + Inches(1.1), Inches(5.9), Inches(0.5),
        "rolling(3).mean() → (D3+D4+D5)/3 貼在 D5 位置",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s, ill_x, row_y + Inches(2.0), Inches(5.9), Inches(1.8),
        "• 前 window-1 列 = NaN（D1 / D2 無前兩筆）\n"
        "• min_periods=1 可強制不產 NaN\n"
        "  但前幾個值會失真\n"
        "• window 小 → 保細節；大 → 看趨勢",
        font_size=T.FONT_SMALL, color=T.CHARCOAL, line_spacing=1.5,
    )
    add_source(s, "pandas User Guide §Window operations §Rolling")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── S13 · CHART · window 大小對比 ─────────
    s = _blank(prs)
    add_title(s, "window 大小 vs 平滑度：看多遠由你決定")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.4),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.6),
        slot_name="rolling window 三層對比圖",
        description=(
            "三條線疊加折線圖：\n"
            "(1) 原始日資料（細線、雜訊明顯）\n"
            "(2) rolling(7).mean()（中粗、週節律仍可見）\n"
            "(3) rolling(30).mean()（粗線、長期趨勢）"
        ),
        url_hint="",
        placeholder_id="S4_S13_rolling_windows",
        registry=image_registry,
        size_hint="1800×900 px",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.15),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "選 window：報表週期相關 → 7 或 30；平滑決策 → 與下游動作週期一致（週報用 7、月報用 30）。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "Hyndman & Athanasopoulos, Forecasting: Principles and Practice 3e §3.3")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── S14 · EXAMPLE-I/O · rolling(7).mean() ─────────
    s = _blank(prs)
    add_title(s, "範例：rolling(7).mean() —— 鋸齒日線變成趨勢線")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="Input: 日營收 Series → Output: 7 日移動平均欄",
        code=(
            '# Input\n'
            'daily = df["amount"].resample("D").sum()\n'
            '\n'
            '# Process — 三種寫法對比\n'
            'daily.rolling(3).mean()                 # 週節律仍有鋸齒\n'
            'daily.rolling(7).mean()                 # 週節律被壓平\n'
            'daily.rolling(7, min_periods=1).mean()  # 前 6 列不 NaN 但失真\n'
            '\n'
            '# 接回 df 當新欄位\n'
            'daily_df = daily.to_frame("revenue")\n'
            'daily_df["ma7"]  = daily_df["revenue"].rolling(7).mean()\n'
            'daily_df["ma30"] = daily_df["revenue"].rolling(30).mean()\n'
            '\n'
            '# Output（畫圖前記得 dropna()）\n'
            'daily_df.dropna().plot(y=["revenue", "ma7", "ma30"])'
        ),
        bullets=[
            "畫圖前先 dropna()\n否則線會斷",
            "min_periods=1 保長度\n但要提醒前段失真",
            "ma7 + ma30 疊圖是\n經典趨勢解讀版面",
            "rolling 可串 .std() .max()",
        ],
        label_dark=True,
    )
    add_source(s, "pandas docs §Window operations §rolling.mean · McKinney 3e §11.7")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── S15 · CONCEPT-CARD · EDA 三板斧 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "describe()",
         "sub": "數值欄分佈：\ncount / mean / std / min / Q1 / Q2 / Q3 / max\n一行看完「長相」",
         "highlight": True},
        {"text": "value_counts()",
         "sub": "類別欄頻次：\n降序排列、預設排除 NaN\n配 normalize=True 看比例"},
        {"text": "corr()",
         "sub": "數值欄兩兩相關：\n只看 Pearson 線性相關\n類別欄要 crosstab / 卡方",
         "highlight": True},
        {"text": "三招口訣",
         "sub": "拿到新資料先跑這三件套\n→ 再決定清洗 / 建模策略\n「先認識它，再動手」"},
    ], title="EDA 三板斧：describe / value_counts / corr")
    add_source(s, "John W. Tukey, Exploratory Data Analysis (1977) · pandas User Guide §Descriptive statistics")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── S16 · EXAMPLE-I/O · 三板斧實戰 ─────────
    s = _blank(prs)
    add_title(s, "範例：EDA 三板斧一次跑 —— 從資料到一句結論")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(6.2), h=Inches(5.0),
        label="三行 code，三個角度理解資料",
        code=(
            '# 1. describe — 數值欄分佈\n'
            'df[["amount", "qty", "discount"]].describe()\n'
            '\n'
            '# 2. value_counts — 類別欄頻次\n'
            'df["category"].value_counts(normalize=True).head()\n'
            '# 3C  0.42\n'
            '# 電子  0.31\n'
            '# 家用  0.18\n'
            '\n'
            '# 3. corr — 數值欄相關\n'
            'df[["amount","qty","discount","age"]].corr()\n'
            '#            amount   qty  discount   age\n'
            '# amount      1.00   0.78    -0.31   0.12\n'
            '# qty         0.78   1.00    -0.24   0.08'
        ),
        bullets=[
            "describe 先 .T 看\n多欄比較更直覺",
            "value_counts + normalize\n= 比例分佈",
            "corr 對角線恆 1.0",
        ],
        label_dark=True,
    )
    draw_image_placeholder(
        s,
        x=Inches(7.0), y=Inches(1.4),
        w=Inches(5.7), h=Inches(4.6),
        slot_name="corr heatmap 截圖",
        description=(
            "df.corr() 產出的 4×4 熱力圖\n"
            "(amount / qty / discount / age)\n"
            "對角線 1.0，主體深綠→白漸層\n"
            "seaborn heatmap 風格"
        ),
        url_hint="",
        placeholder_id="S4_S16_corr_heatmap",
        registry=image_registry,
        size_hint="1200×1000 px",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "結論口語版：『amount 與 qty 強正相關（0.78）—— 買得多、金額也高，符合直覺。』",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "pandas docs §DataFrame.describe / .value_counts / .corr")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── S17 · PITFALL · corr 解讀陷阱 ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="✗ 只丟數字",
        right_title="✓ 三段式口語解讀",
        left_items=[
            "「amount 和 qty 的 corr 是 0.65」",
            "「你自己看相關係數表」",
            "老闆：所以呢？",
            "類別欄也丟 corr → 結果是 NaN 或 0",
            "看到 0.2 就說「弱相關」→\n但 p-value 沒驗",
        ],
        right_items=[
            "「amount 與 qty 中度正相關（0.65）」",
            "「買越多、金額也越高，\n但不完全線性」",
            "「符合直覺，不需額外動作」",
            "類別欄改用 crosstab / 卡方檢定",
            "門檻：>0.7 強 / 0.3–0.7 中 / <0.3 弱",
        ],
        title="常見錯誤：相關係數只丟數字、忽略類別欄",
        summary="相關不等於因果；三段式：『欄 A 與欄 B』+『強度』+『商業解讀』—— 才是老闆要聽的。",
        delta="0.65",
    )
    add_source(s, "Hadley Wickham, R for Data Science 2e §EDA · Kahneman, Thinking Fast and Slow §Illusion of causality")
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───────── S18 · PITFALL · 'M' vs 'ME' 與四錯並陳 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "✗ resample('M')",
         "sub": "pandas 2.2+ FutureWarning\n→ 改 'ME'（Month End）\n同理：'Q'→'QE' / 'Y'→'YE'"},
        {"text": "✗ AttributeError: .dt",
         "sub": "read_csv 沒帶 parse_dates\n→ order_date 是 object\n→ 補 pd.to_datetime() 或\nparse_dates=['order_date']"},
        {"text": "✗ 星期字母序",
         "sub": "groupby('weekday') 後畫圖\n→ Friday 排最前\n→ .reindex([Mon..Sun]) 固定順序"},
        {"text": "✗ to_period('M') 型別",
         "sub": "Period 物件寫 CSV / merge 會出錯\n→ 畫圖前先 .astype(str)\n或保持 Timestamp"},
    ], title="四個時序常見錯誤：看過一次就不會再中")
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "每個錯誤都有具體症狀、具體修法 —— 記不住沒關係，這四張卡收進 MVK。",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "pandas 2.2 Release Notes §Deprecations · Stack Overflow pandas top errors")
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ───────── S19 · EXAMPLE-I/O · 月度經營報表 walkthrough ─────────
    s = _blank(prs)
    add_title(s, "整合實戰：月度經營報表 —— 交付給 S5 的 monthly_revenue.csv")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="一個 pipeline 產六欄報表：給老闆看、給 S5 畫圖、給 S6 建模",
        code=(
            '# 1. 標準開場\n'
            'df = (pd.read_csv("orders.csv", parse_dates=["order_date"])\n'
            '        .set_index("order_date").sort_index())\n'
            '\n'
            '# 2. 月度 agg —— 多欄不同聚合\n'
            'monthly = df.resample("ME").agg(\n'
            '    revenue=("amount", "sum"),\n'
            '    orders =("order_id", "count"),\n'
            '    aov    =("amount", "mean"),\n'
            '    unique_customers=("customer_id", "nunique"),\n'
            ')\n'
            '\n'
            '# 3. 加衍生欄：月成長率 + 7 日平均回填\n'
            'monthly["mom_growth"] = monthly["revenue"].pct_change()\n'
            'monthly["ma3_revenue"] = monthly["revenue"].rolling(3).mean()\n'
            '\n'
            '# 4. 輸出（index 保留為 Month-End Timestamp）\n'
            'monthly.to_csv("monthly_revenue.csv")'
        ),
        bullets=[
            "named agg：結果欄名乾淨",
            "mom_growth 第一列 NaN\n是對的（沒前一月）",
            "要不要補 0 看下游：\n畫圖 → dropna\nKPI → 顯示 N/A",
            "這張表就是 S5 的輸入",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Named aggregation · 對應 01_outline.md §2 60~70 分鐘段落")
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───────── S20 · PRACTICE-PROMPT · 三層練習 ─────────
    s = _blank(prs)
    add_title(s, "練習時間（10 分鐘）：送分 / 核心 / 挑戰")
    # Three cards
    card_w = Inches(3.9)
    card_h = Inches(4.5)
    card_y = Inches(1.4)
    gap_x = Inches(0.2)
    cards = [
        {"tag": "送分 · 3 min",
         "title": "算每月營收",
         "body": "用 resample('ME') + sum\n把 orders.csv 的 amount\n聚合成月營收 Series。\n\n關鍵字：\nparse_dates / set_index\nsort_index / resample"},
        {"tag": "核心 · 5 min",
         "title": "加上 7 日 rolling",
         "body": "在每日營收上\n疊一條 rolling(7).mean()\n並畫出雙線圖。\n\n關鍵字：\nresample('D') / rolling\ndropna / plot"},
        {"tag": "挑戰 · 10 min",
         "title": "月度經營報表 + 週一平均",
         "body": "產出 S19 的六欄報表；\n再用 groupby('weekday')\n算每個星期幾的平均營收，\n並固定 Mon-Sun 排序。\n\n對應 §5-Q1"},
    ]
    for i, c in enumerate(cards):
        x = T.MARGIN_X + i * (card_w + gap_x)
        rect = add_rect(s, x, card_y, card_w, card_h)
        if i == 2:
            set_solid_fill(rect, T.PRIMARY)
            set_no_line(rect)
            tag_color = T.LIGHT_GRAY
            title_color = T.WHITE
            body_color = T.LIGHT_GRAY
        else:
            set_no_fill(rect)
            set_line(rect, T.PRIMARY, 1.2)
            tag_color = T.GRAY_MID
            title_color = T.PRIMARY
            body_color = T.CHARCOAL
        add_textbox(
            s, x + Inches(0.25), card_y + Inches(0.2),
            card_w - Inches(0.5), Inches(0.4),
            c["tag"], font_size=T.FONT_SMALL, color=tag_color, bold=True,
        )
        add_textbox(
            s, x + Inches(0.25), card_y + Inches(0.65),
            card_w - Inches(0.5), Inches(0.7),
            c["title"], font_size=T.FONT_SUBTITLE, color=title_color, bold=True,
        )
        add_textbox(
            s, x + Inches(0.25), card_y + Inches(1.5),
            card_w - Inches(0.5), card_h - Inches(1.7),
            c["body"], font_size=T.FONT_CAPTION, color=body_color,
            line_spacing=1.55,
        )
    add_textbox(
        s, T.MARGIN_X, Inches(6.1),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "與同伴討論：挑戰題第一個月的 mom_growth 會是 NaN —— 要不要補 0？",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "對應 01_outline.md §5 討論提問 Q1、Q3")
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ───────── S21 · CHECKPOINT · 選用決策 ─────────
    s = _blank(prs)
    add_title(s, "Check Point：選對工具比寫對 code 重要")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "下列情境，你會用 resample / groupby / rolling 哪一個？",
        font_size=T.FONT_SUBTITLE, color=T.PRIMARY, bold=True,
    )
    scenarios = [
        ("Q1", "每月營收總和（要補齊缺月）",       "resample('ME').sum()"),
        ("Q2", "每個星期幾的平均單量",              "groupby(weekday).mean() + reindex"),
        ("Q3", "平滑每日營收看趨勢",                "rolling(7).mean()"),
        ("Q4", "每季各品類的訂單數",                "groupby([pd.Grouper('QE'), 'category'])"),
    ]
    qy = Inches(2.0)
    qh = Inches(0.85)
    for i, (tag, q, a) in enumerate(scenarios):
        y = qy + i * (qh + Inches(0.15))
        add_textbox(
            s, T.MARGIN_X, y, Inches(0.7), qh,
            tag, font_size=T.FONT_SUBTITLE, color=T.PRIMARY, bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            s, T.MARGIN_X + Inches(0.9), y, Inches(6.2), qh,
            q, font_size=T.FONT_BODY, color=T.CHARCOAL,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Hidden answer box (right side, light fill)
        ans = add_rect(s, T.MARGIN_X + Inches(7.4), y, Inches(4.8), qh)
        set_no_fill(ans)
        set_line(ans, T.LIGHT_GRAY, 1.0, dash=True)
        add_textbox(
            s, T.MARGIN_X + Inches(7.55), y, Inches(4.5), qh,
            a, font_size=T.FONT_SMALL, color=T.GRAY_MID,
            family=T.FONT_MONO, anchor=MSO_ANCHOR.MIDDLE,
        )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "右側虛線答案演講時揭示 —— 重點不是背 API，是看到情境能分類。",
        font_size=T.FONT_CAPTION, color=T.GRAY_MID, italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "教學型守則 T10 · 每 5 張必有 Checkpoint")
    add_footer(s, MODULE_CODE, 21, N_CONTENT)

    # ───────── S22 · PYRAMID · 收束三層 ─────────
    s = _blank(prs)
    draw_pyramid_stack(s, layers=[
        {"name": "Next", "caption": "S5：把這些報表畫成圖 —— 折線、長條、散點、熱圖、盒鬚"},
        {"name": "Why",  "caption": "因為：商業問題都有時間維度；老闆四問全部關時間"},
        {"name": "What", "caption": "本節四工具：.dt / resample / rolling / EDA 三板斧"},
    ], title="收束：What / Why / Next",
       thesis="時間是軸 · EDA 是直覺 · 報表是產出 —— 這是 DA 工作的骨架。")
    add_source(s, "教學型守則 T12 · 收尾三件套")
    add_footer(s, MODULE_CODE, 22, N_CONTENT)

    # ───────── S23 · SILENT · 下一節預告 ─────────
    s = _blank(prs)
    draw_silent_page(s, "S5 預告：\n數字不會說話，圖會。")
    add_footer(s, MODULE_CODE, 23, N_CONTENT, dark_bg=True)

    # Copyright page
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
