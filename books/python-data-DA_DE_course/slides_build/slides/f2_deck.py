"""F2 deck — Python 核心與資料結構深化.

21 content slides + cover + copyright page.
Teaching-track; palette: 黑 / 灰 / 深綠 (#1B5E3F) + 警示紅 (#C62828) + 成功綠 (#2E7D32).

受眾：只有 Python 基礎的非理工背景學員（商管/人文/行銷/行政/設計/醫護）。
  - 不用 CS 術語（hash table / O(1) / pointer / invariant / stack-heap）
  - 大量使用生活類比：聯絡人、字典查字、集合去重、逐頁翻書、便利貼

Governing thought:
    會選容器，代表你能寫對；
    會用 generator，代表你能寫大——
    為 Pandas / ETL 打下兩個根基。

Aligned to chapters/F2_Python核心與資料結構深化/00_skeleton.yaml
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source, add_rect,
    set_solid_fill, set_no_fill, set_line, set_no_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_code_panel, draw_vs_two_col,
    draw_inverted_thesis_box, draw_emphasis_pill, draw_pyramid_stack,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "F2"
MODULE_TITLE = "Python 核心與資料結構深化"
MODULE_SUBTITLE = "容器選用 × 可變性 × Iterator/Generator——為 Pandas 與 ETL 鋪路"
TIME_MIN = 90
N_CONTENT = 21

RED_ERROR = RGBColor(0xC6, 0x28, 0x28)
GREEN_OK = RGBColor(0x2E, 0x7D, 0x32)


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_f2(output_path, image_registry=None):
    """Build F2 deck.

    image_registry: optional list collecting image placeholder metadata
    (unused in F2 — no real images required).
    """
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ---- S1 MOTIVATION — Pandas 踩的坑，根在 Python 基礎 ----
    s = _blank(prs)
    draw_ask_page(
        s,
        "你在 Pandas 裡踩過的那些『改了 A 結果 B 也變』『一個警告看不懂』，\n根源其實不在 Pandas——在 Python 基礎。",
        data_card={
            "label": "課程樣本觀察",
            "stat": "60%",
            "caption": "Pandas 初學者的錯誤，回溯到容器選錯 / 複製不對 / 迭代用錯",
        },
    )
    add_source(s, "本課程 2023–2025 企業訓練 code review 樣本 n=120")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # ---- S2 ASK — 10GB log 要怎麼讀 ----
    s = _blank(prs)
    draw_ask_page(
        s,
        "要在 10GB 的 log 檔裡挑出所有 ERROR 行，\n你會怎麼寫？整個檔案讀進來還是怎樣？",
        data_card={
            "label": "笨方法 vs 聰明方法",
            "stat": "~200×",
            "caption": "全部讀進來 → 電腦直接卡死；逐頁翻 → 只用 < 50 MB",
        },
    )
    add_source(s, "本課程實測 Python 3.11 / 16GB 機器")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ---- S3 SILENT — 一句話立論 ----
    s = _blank(prs)
    draw_silent_page(
        s,
        "今天只教兩件事：\n會選容器 → 能寫對；會用 generator → 能寫大。",
    )
    add_footer(s, MODULE_CODE, 3, N_CONTENT, dark_bg=True)

    # ---- S4 CONCEPT-CARD — 四容器生活類比 ----
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "List  [ ]  =  聯絡人清單",
         "sub": "有順序、可以重複\n第 1 個、第 2 個、第 3 個\n加一筆到最後很快"},
        {"text": "Tuple ( )  =  身分證欄位",
         "sub": "固定欄位、不能改\n(姓, 名, 生日)\n保證不被亂動"},
        {"text": "Dict  { k: v }  =  字典查字",
         "sub": "用『關鍵字』直接翻到答案\nID → 名字、欄位 → 值\n不用翻整本"},
        {"text": "Set   { }  =  收藏夾/去重袋",
         "sub": "自動去重、不管順序\n訂單編號去重、標籤集合\n判斷『有沒有』超快",
         "highlight": True},
    ], title="四種容器 = 四種生活場景（選錯，後面會付兩種代價：效能 + 正確性）")
    add_source(s, "Python 官方文件 docs.python.org/3/library/stdtypes.html")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ---- S5 MECHANISM-FLOW — 30 秒決策樹 ----
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "Q1：要保留順序、允許重複？",
         "sub": "→ 用 List\n例：每天的營業額紀錄\n（1 號、2 號、3 號… 逐筆）"},
        {"text": "Q2：只在乎『有沒有』、要去重？",
         "sub": "→ 用 Set\n例：所有下過單的客戶編號\n（同一人下 10 次只算 1 次）"},
        {"text": "Q3：要用『關鍵字』快速找到答案？",
         "sub": "→ 用 Dict\n例：員工編號 → 姓名\n（查 E1024 直接跳答案）"},
        {"text": "Q4：欄位固定、不希望被改？",
         "sub": "→ 用 Tuple\n例：(縣市, 鄉鎮, 郵遞區號)\n寫好就鎖死，也能當 Dict 的 key",
         "highlight": True},
    ], title="30 秒容器選用決策：照順序問這四個問題，選到就停")
    add_textbox(s, T.MARGIN_X, Inches(6.4),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
                "實務分布（本課程 code review 樣本）：List 80% · Dict 15% · Set 4% · Tuple 1%",
                font_size=T.FONT_SOURCE, color=T.GRAY_MID,
                align=PP_ALIGN.CENTER)
    add_source(s, "本課程 2023–2025 企業訓練樣本 n=120")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ---- S6 EXAMPLE-I/O — 四容器關鍵操作對照表 ----
    s = _blank(prs)
    draw_editorial_table(s,
        header=["容器", "生活場景", "關鍵操作", "查一筆資料的速度"],
        rows=[
            ["List",  "逐筆紀錄（有順序）",
             "append / [i] / 切片",
             "要從頭翻到尾（資料量大會慢）"],
            ["Tuple", "固定欄位（不會變）",
             "解構 a,b,c = t",
             "同 List，但保證不被改"],
            ["Dict",  "用關鍵字查答案",
             "d[key] / .get() / in",
             "直接翻到（資料量多大都一樣快）"],
            ["Set",   "去重、判斷有沒有",
             "add / & | - / in",
             "直接翻到（同 Dict）"],
        ],
        col_widths=[0.9, 1.8, 1.6, 2.5],
        title="查一筆資料：List 是翻整本，Dict / Set 是直接翻索引",
    )
    add_textbox(s, T.MARGIN_X, Inches(5.6),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
                "常見題：「這個 ID 有沒有下過單？」——10 萬筆時，List 要翻幾秒、Set 幾乎 0 秒。",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER)
    add_source(s, "Python 官方 Time Complexity 文件")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ---- S7 EXAMPLE-I/O — 商業場景 + 銜接 S2 DataFrame ----
    s = _blank(prs)
    add_title(s, "商業場景實例——這些結構就是下週 Pandas DataFrame 的前身")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.4),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.4),
        label="EXAMPLE · 三個日常商業任務",
        code=(
            "# 1) 訂單編號去重（同客戶下多次）—— 用 Set\n"
            "orders = ['A001', 'A002', 'A001', 'A003', 'A002']\n"
            "unique = set(orders)          # {'A001', 'A002', 'A003'}\n"
            "\n"
            "# 2) 員工 ID → 姓名 對照 —— 用 Dict\n"
            "emp = {'E1024': '王小明', 'E1025': '林小華'}\n"
            "emp['E1024']                  # '王小明'\n"
            "\n"
            "# 3) 一筆客戶資料（欄位固定）—— 用 Dict（一列）\n"
            "row = {'id': 'E1024', 'name': '王小明', 'city': '台北'}\n"
            "# 很多列堆起來 → 就是下週 pandas 的 DataFrame！\n"
            "rows = [row, {'id': 'E1025', 'name': '林小華', 'city': '新北'}]"
        ),
        bullets=[
            "set → 「去重 / 判斷有沒有」的萬用解",
            "dict → 「用關鍵字查答案」的萬用解",
            "一個 dict = 資料表的一列（row）",
            "一串 dict = 資料表（DataFrame 前身）",
            "→ 下週 S2：`pd.DataFrame(rows)` 直接變表",
        ],
    )
    add_source(s, "銜接 S2 Pandas I/O · DataFrame 的本質")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ---- S8 CONCEPT-CARD — 便利貼心智模型 ----
    s = _blank(prs)
    draw_flow_chain(s, nodes=[
        {"label": "name", "caption": "便利貼\n（變數名）"},
        {"label": "→ 指向", "caption": "貼在……"},
        {"label": "object", "caption": "真正的東西\n（在記憶體裡）", "highlight": True},
    ], title="變數不是盒子，是便利貼——貼在某個『東西』上",
       y=3.0)
    draw_inverted_thesis_box(
        s,
        "`b = a` 不是影印一份，是再貼一張便利貼到同一個東西上。改它 → 兩邊都變。",
        y=5.9, width=11.0,
    )
    add_source(s, "生活類比版 · 對應 Python data model")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ---- S9 CONCEPT-CARD — 可變/不可變 × 能否當 key ----
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "會變的東西 · 不能當 key",
         "sub": "list · dict · set\n內容會變 → Python 記不住它的指紋\n當 key 會直接報錯"},
        {"text": "不會變的東西 · 能當 key",
         "sub": "數字 · 字串 · tuple\n寫好就鎖死 → 有固定指紋\n適合當『關鍵字』",
         "highlight": True},
        {"text": "會變的 · 只能放在『答案』那邊",
         "sub": "`d[key] = [商品1, 商品2]`\nvalue 想怎麼變都可以\n這是最常用的情況"},
        {"text": "不會變的 · 也能放答案",
         "sub": "`d[key] = 100` 或 `'VIP'`\n計數器、標籤、狀態值\n都很常見"},
    ], title="會變 ≠ 能當 key — 因為 Python 記不住一個會變的東西的『指紋』")
    add_textbox(s, T.MARGIN_X, Inches(6.4),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
                "橫軸：會不會變   ·   縱軸：要當關鍵字（key）還是答案（value）",
                font_size=T.FONT_SOURCE, color=T.GRAY_MID,
                align=PP_ALIGN.CENTER)
    add_source(s, "生活類比版 · Python hashable 規則")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ---- S10 PITFALL — [[]] * 3 ----
    s = _blank(prs)
    add_title(s, "`[[]] * 3` 不是建三個空 list——是同一個 list 被貼了三張便利貼")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.4),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.4),
        label="PITFALL · 新手第一坑",
        code=(
            "# 你以為：三個獨立的空籃子\n"
            "a = [[]] * 3\n"
            "a[0].append(1)              # 只想放進第 1 個籃子\n"
            "print(a)\n"
            "# 結果：[[1], [1], [1]]   ← 三格全中，嚇到\n"
            "\n"
            "# 正解：真的建三個獨立的\n"
            "a = [[] for _ in range(3)]  # 逐一生一個\n"
            "a[0].append(1)\n"
            "print(a)   # [[1], [], []]  ← 這才對"
        ),
        bullets=[
            "`*` 複製的是『便利貼』，不是後面那個東西",
            "三張便利貼全貼在同一個 list 上",
            "改任何一格，三格都跟著變（必然，不是 bug）",
            "心法：要 n 個獨立的 → 用 comprehension",
            "記口訣：`[[] for _ in range(n)]` 逐一生",
        ],
    )
    add_source(s, "Python FAQ · 'How do I make a list of lists?'")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ---- S11 PITFALL — 淺深拷貝 ----
    s = _blank(prs)
    draw_vs_two_col(s,
        left_title="淺拷貝（shallow）= 只複製外殼",
        right_title="深拷貝（deep）= 連裡面一起複製",
        left_items=[
            "`list.copy()` · `x[:]` · `list(x)`",
            "只複製最外層，裡面還是共用",
            "像複製 Excel 分頁——改分頁 A 的儲存格 A1 不會影響 B",
            "但巢狀 list 改『內層』→ 原本也會中",
            "扁平資料沒事，巢狀就中毒",
        ],
        right_items=[
            "`copy.deepcopy(x)`",
            "從外殼到葉子整份重新複製一遍",
            "改到哪層都不影響原本",
            "代價：時間與記憶體多一份",
            "不確定時先 deepcopy 保命",
        ],
        title="淺拷貝複製外殼，深拷貝連骨頭——巢狀結構只有後者安全",
        summary="這就是你未來看到 Pandas 那個 `SettingWithCopyWarning` 的根源——系統在提醒你：你改的到底是本尊還是分身？",
        delta="只差一層",
    )
    add_source(s, "Python 標準函式庫 copy module · 對應 Pandas SettingWithCopy")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ---- S12 CONCEPT-CARD — Iterator 生活版 ----
    s = _blank(prs)
    draw_flow_chain(s, nodes=[
        {"label": "for 問：", "caption": "還有嗎？"},
        {"label": "容器答：", "caption": "有，這是下一個"},
        {"label": "for 處理：", "caption": "做事"},
        {"label": "容器答：", "caption": "沒了（停止）", "highlight": True},
    ], title="for 迴圈其實是在『跟容器要下一個』——直到它說『沒了』",
       y=3.2)
    draw_inverted_thesis_box(
        s,
        "你不用寫 while 跟計數器——Python 的 for 已經幫你問到底、接收到『沒了』就自動停。",
        y=5.9, width=11.0,
    )
    add_source(s, "生活類比版 · 對應 Iterator protocol")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ---- S13 MECHANISM-FLOW — yield 暫停/續行 ----
    s = _blank(prs)
    add_title(s, "yield 不是 return——它讓函式『暫停、記住現場、下次從這裡繼續演』")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.4),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.4),
        label="MECHANISM · 把函式變成連續劇",
        code=(
            "def counter():\n"
            "    print('第一集開始')\n"
            "    yield 1          # ← 本集結束，暫停在這\n"
            "    print('第二集開播')\n"
            "    yield 2\n"
            "    print('最終回')\n"
            "\n"
            "g = counter()        # 此時什麼都沒播（只是『訂閱』）\n"
            "next(g)   # 第一集開始  → 拿到 1\n"
            "next(g)   # 第二集開播  → 拿到 2\n"
            "next(g)   # 最終回      → 結束（StopIteration）"
        ),
        bullets=[
            "return 一去不回；yield 下次會從斷點接著演",
            "呼叫 `g = counter()` 不會真的執行——只是訂閱",
            "每次 yield 凍住『場景 + 進度』，像劇集暫停",
            "記憶體只裝『當下這一集』，不裝整季",
            "→ 這就是下頁『逐頁翻書』的引擎",
        ],
    )
    add_source(s, "PEP 255 · Simple Generators（生活類比版）")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ---- S14 EXAMPLE-I/O — [ ] vs ( ) ----
    s = _blank(prs)
    draw_vs_two_col(s,
        left_title="[ ] 一次算完塞 RAM",
        right_title="( ) 要一個算一個",
        left_items=[
            "`[x*2 for x in range(10**8)]`",
            "一次算完 1 億個，整排塞進記憶體",
            "記憶體：~ 3.2 GB（電腦會哀號）",
            "適用：資料小、要反覆讀",
            "特徵：可以數 len、可以 [i] 取",
        ],
        right_items=[
            "`(x*2 for x in range(10**8))`",
            "只算當下那一個，下一個才繼續",
            "記憶體：~ 200 bytes（幾乎不佔）",
            "適用：資料大、只需讀一遍",
            "特徵：沒有 len、只能從頭跑到尾",
        ],
        title="括號不同、命運不同：[ ] 立刻算完全部，( ) 要一個算一個",
        summary="要存結果用 [ ]；要處理很大的資料用 ( )——同樣的寫法，記憶體差兩個數量級。",
        delta="~ 200× RAM",
    )
    add_source(s, "本課程實測 Python 3.11 · tracemalloc 觀測")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ---- S15 EXAMPLE-I/O — 10GB log 三行 generator ----
    s = _blank(prs)
    add_title(s, "10GB log 挑 ERROR 行——三行 generator 解決（『逐頁翻書』）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.4),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.4),
        label="EXAMPLE · 不把整本搬進 RAM",
        code=(
            "def read_errors(path):\n"
            "    with open(path, encoding='utf-8') as f:\n"
            "        for line in f:              # 檔案本身就是『一頁一頁』\n"
            "            if 'ERROR' in line:\n"
            "                yield line.rstrip() # 翻到一頁就給出去\n"
            "\n"
            "# 使用端：完全不用關心檔案多大\n"
            "for err in read_errors('/var/log/app.log'):\n"
            "    send_to_alert(err)              # 處理完這頁，再翻下一頁"
        ),
        bullets=[
            "open() 回傳的 file 天生就能『一頁一頁讀』",
            "for line in f：一次只拿一行，RAM 幾乎沒壓力",
            "yield 讓呼叫端也『要一筆給一筆』",
            "組合起來：10GB 流過，RAM 用不到 50 MB",
            "→ 下一張：這招在 Pandas 裡叫 chunksize",
        ],
    )
    add_source(s, "本課程實戰範例")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ---- S16 EXAMPLE-I/O — 銜接 S2 chunksize ----
    s = _blank(prs)
    add_title(s, "這就是下週 `pd.read_csv(chunksize=...)` 背後的魔法")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.4),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.4),
        label="BRIDGE · F2 → S2",
        code=(
            "# 下週 S2 你會寫：處理 5GB CSV 也不會爆\n"
            "import pandas as pd\n"
            "\n"
            "# chunksize=100_000 → 回傳『generator』，不是整個 DataFrame\n"
            "chunks = pd.read_csv('orders_5GB.csv', chunksize=100_000)\n"
            "\n"
            "for chunk in chunks:            # ← for + generator = 逐頁\n"
            "    result = chunk[chunk['amount'] > 1000]  # 處理這塊\n"
            "    result.to_csv('out.csv', mode='a', header=False)\n"
            "    # 這塊處理完 → 自動丟掉 → 下一塊才載入\n"
            "\n"
            "# 記憶體只裝當下這 10 萬列，不是 5GB 全部"
        ),
        bullets=[
            "Pandas 的 chunksize 回傳的就是 generator",
            "跟今天教的 `yield` 是同一套引擎",
            "for chunk in chunks：跟今天 for line in f 一模一樣",
            "→ 處理大資料的標準招式：讀一塊、算一塊、寫一塊",
            "今天你懂 yield = 下週你懂 chunksize",
        ],
    )
    add_source(s, "Pandas 官方文件 · read_csv(chunksize) — 銜接 S2")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ---- S17 PITFALL — generator 只能迭代一次 ----
    s = _blank(prs)
    add_title(s, "generator 只能『翻』一次——翻完就沒了，第二次跑是空的")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.4),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(4.4),
        label="PITFALL · 最後一坑",
        code=(
            "g = (x*2 for x in range(5))\n"
            "\n"
            "print(sum(g))      # 20   ✅ 第一次翻完了\n"
            "print(list(g))     # []   ← 空的！已經翻到最後一頁\n"
            "print(max(g))      # ValueError: 空的沒辦法取最大值\n"
            "\n"
            "# 要重用，兩條路：\n"
            "# (A) 重新建一個\n"
            "data = list(range(5))\n"
            "print(sum(x*2 for x in data))\n"
            "print(max(x*2 for x in data))\n"
            "\n"
            "# (B) 一次存成 list（資料小時）\n"
            "nums = [x*2 for x in range(5)]"
        ),
        bullets=[
            "generator 像一本書：翻完就到最後一頁了",
            "沒有『回到第一頁』這件事",
            "常見地雷：把同一個 generator 傳給多個函式",
            "除錯訣竅：想多次用 → 先 list() 存下來",
            "資料量允許就 list；不允許就重新建",
        ],
    )
    add_source(s, "Python 官方文件 · Generator iterator semantics")
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ---- S18 PRACTICE-PROMPT — 3 分鐘練習 ----
    s = _blank(prs)
    add_title(s, "練習時間 · 3 分鐘 · Think-Pair-Share")
    draw_matrix(s, rows=1, cols=3, cells=[
        {"text": "① 訂單去重",
         "sub": "`orders = ['A1','A2','A1','A3']`\n\n任務：去掉重複的訂單編號\n\n提示：哪個容器自動去重？"},
        {"text": "② ID 對照",
         "sub": "`emp_id = ['E1','E2']`\n`name   = ['王','林']`\n\n任務：用 E1 查到『王』\n\n提示：哪個容器是『查字典』？"},
        {"text": "③ 逐行讀檔",
         "sub": "給一個 log.txt（很大）\n\n任務：寫一個 generator\n只回傳含 'ERROR' 的行\n\n提示：`yield` + `for line in f`",
         "highlight": True},
    ], title="用今天三招解三題（每題一行就夠）")
    add_textbox(s, T.MARGIN_X, Inches(6.0),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.8),
                "3 分鐘後對答案：① `set(orders)`  ② `dict(zip(emp_id, name))['E1']`  ③ 見 S15 三行 yield",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER, line_spacing=1.3)
    add_source(s, "本課程 F2 課堂練習")
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ---- S19 CHECKPOINT — 三題快問 ----
    s = _blank(prs)
    draw_matrix(s, rows=1, cols=3, cells=[
        {"text": "Q1 · 容器選用",
         "sub": "要記錄 10 萬筆訂單 ID，\n之後頻繁檢查『某 ID 有沒有下過單』，\n選哪個容器？為什麼？"},
        {"text": "Q2 · 淺深拷貝",
         "sub": "`a = [[1,2],[3,4]]`\n`b = a.copy()`\n`b[0][0] = 99`\n\na 變成什麼？"},
        {"text": "Q3 · Generator",
         "sub": "`g = (x for x in range(3))`\n先 `list(g)` 再 `list(g)`\n\n兩次各印什麼？",
         "highlight": True},
    ], title="30 秒三題快問——跟得上嗎？")
    add_textbox(s, T.MARGIN_X, Inches(6.0),
                T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.8),
                "答：Q1 Set（判斷有沒有超快） · Q2 `[[99,2],[3,4]]`（淺拷貝共用內層） · Q3 `[0,1,2]` 然後 `[]`",
                font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER, line_spacing=1.3)
    add_source(s, "本課程 F2 checkpoint")
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ---- S20 PYRAMID — 三層 take-away ----
    s = _blank(prs)
    draw_pyramid_stack(s,
        layers=[
            {"name": "思維：寫大",
             "caption": "用 generator 逐頁翻，不把整本搬進 RAM"},
            {"name": "心智：寫對",
             "caption": "變數是便利貼；會變的東西不能當 key"},
            {"name": "工具：選容器",
             "caption": "List 排序 · Dict 查找 · Set 去重 · Tuple 鎖死"},
            {"name": "地基：Python 基礎",
             "caption": "變數、容器、for、函式——課前已備"},
        ],
        thesis="這三層會了——下週 S1 NumPy / S2 Pandas，你會發現它們就長在這塊地基上。",
        title="F2 三層 take-away：從基礎 → 選容器 → 懂可變 → 用 generator",
    )
    add_source(s, "本課程 F2 收束")
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ---- S21 SILENT — 收束 ----
    s = _blank(prs)
    draw_silent_page(
        s,
        "會選容器 → 能寫對；會用 generator → 能寫大。\n\n下一站：S2 Pandas —— DataFrame 就是一排 dict + chunksize。",
    )
    add_footer(s, MODULE_CODE, 21, N_CONTENT, dark_bg=True)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
