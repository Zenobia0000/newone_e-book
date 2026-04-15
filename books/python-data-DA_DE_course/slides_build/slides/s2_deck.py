"""S2 deck — Pandas I/O 與資料清理（2 小時 · 23 張內容投影片）

Governing thought:
    資料科學家 80% 時間在清資料——把『拿到新檔案』的 SOP 肌肉記憶化，
    讓清理從憑感覺變成可複現的流水線。

教學型 7 原型配置：
    MOTIVATION (S1,S2) → CONCEPT-CARD (S3-S15) → MECHANISM-FLOW (S6,S12,S17)
    → EXAMPLE-I/O (S18,S19) → PITFALL (S8,S11,S16,S20) → PRACTICE (S21)
    → CHECKPOINT (S22,S23)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_rect, add_textbox, add_title, add_source,
    set_solid_fill, set_no_line, set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_vs_two_col, draw_code_panel, draw_flow_chain,
    draw_three_blocks_flow, draw_thesis_hierarchy,
    draw_inverted_thesis_box, draw_image_placeholder,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "S2"
MODULE_TITLE = "Pandas I/O 與資料清理"
MODULE_SUBTITLE = "拿到新檔案的 SOP：讀 → 看 → 清 → 驗 → 存"
TIME_MIN = 120
N_CONTENT = 23


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_s2(output_path, image_registry=None):
    prs = _new_prs()

    # ───────── Cover ─────────
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── S1 · SILENT · Motivation ─────────
    s = _blank(prs)
    draw_silent_page(s, "資料科學家 80% 時間在清資料。\n這一節，把 SOP 刻進肌肉。")
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ───────── S2 · ASK · Hook ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "拿到一份新 CSV，\n你的第一個指令是什麼？",
        data_card={
            "label": "Kaggle 2024 ML & DS Survey",
            "stat": "80%",
            "caption": "資料工作者自承\n清理與探索吃掉大半工時",
        },
    )
    add_source(s, "Kaggle, State of Machine Learning and Data Science 2024")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── S3 · MATRIX 2×3 · 三件套 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=3, cells=[
        {"text": "Series",
         "sub": "一維帶 Index 的陣列\nNumPy + 標籤\n單欄資料的化身",
         "highlight": True},
        {"text": "DataFrame",
         "sub": "二維表\n多個 Series 共用同一 Index\n類似 SQL Table"},
        {"text": "Index",
         "sub": "資料的「名字」\n≠ 欄位、≠ row number\n對齊與合併的根據"},
        {"text": "對齊（alignment）",
         "sub": "兩 Series 相加\n自動依 Index 對齊\n缺位 → NaN",
         "highlight": True},
        {"text": "Index 不可變",
         "sub": "建立後不能 in-place 改\n要用 reindex / rename"},
        {"text": "set / reset_index",
         "sub": "在 Index 與欄位之間\n互換的兩道門"},
    ], title="DataFrame / Series / Index：Pandas 的三件套心智模型")
    add_source(s, "pandas User Guide §Intro to data structures · McKinney, Python for Data Analysis 3e §5")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───────── S4 · CODE · Index 對齊 ─────────
    s = _blank(prs)
    add_title(s, "Index 不是欄位——它是對齊與合併的根據")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.5), h=Inches(5.0),
        slot_name="Index 對齊示意圖",
        description=(
            "s1=[A,B,C] 與 s2=[B,C,D]\n"
            "中間 + 號 → 結果 [A=NaN, B=合, C=合, D=NaN]\n"
            "強調「對齊靠名字，不是位置」"
        ),
        url_hint="https://pandas.pydata.org/docs/user_guide/dsintro.html",
        placeholder_id="S2_S04_index_alignment",
        registry=image_registry,
        size_hint="1280×720 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.4), y=Inches(1.3),
        w=Inches(6.3), h=Inches(5.0),
        label="標籤對齊：pandas 的隱藏外掛",
        code=(
            'import pandas as pd\n'
            '\n'
            's1 = pd.Series([10, 20, 30],\n'
            '               index=["A","B","C"])\n'
            's2 = pd.Series([1, 2, 3],\n'
            '               index=["B","C","D"])\n'
            '\n'
            's1 + s2\n'
            '# A     NaN  ← 只在 s1\n'
            '# B    21.0\n'
            '# C    32.0\n'
            '# D     NaN  ← 只在 s2'
        ),
        bullets=[
            "對齊靠 Index 名字\n不是位置",
            "Excel 沒這個能力",
            "merge / join 的底層原理",
            "看懂這張 = 懂 50% pandas",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Intro §Alignment")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── S5 · CODE · 建立 + 五件事 SOP ─────────
    s = _blank(prs)
    add_title(s, "建立 DataFrame 的三條路 + 看資料的五件事 SOP")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="拿到任何新資料，先用這五行把它看穿",
        code=(
            'import pandas as pd\n'
            '\n'
            '# 三條建立路徑\n'
            'df1 = pd.DataFrame({"name":["A","B"], "age":[20,30]})\n'
            'df2 = pd.DataFrame([{"name":"A","age":20}, {"name":"B","age":30}])\n'
            'df  = pd.read_csv("orders_raw.csv", encoding="utf-8-sig")\n'
            '\n'
            '# 五件事 SOP —— senior 和 junior 的差別\n'
            'df.shape              # (rows, cols) 規模感\n'
            'df.head()             # 前 5 列：長相\n'
            'df.info()             # 型別 + 缺失，debug 第一站\n'
            'df.describe()         # 數值欄分佈\n'
            'df.isna().sum()       # 各欄缺值數\n'
        ),
        bullets=[
            "別急著 query\n先跑五件事",
            "object dtype\n八成是字串或混型別",
            "describe 只看數值欄\n類別欄用 value_counts",
            "isna 是缺值的第一道\ndebug 入口",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Essential basic functionality · PyData best practices")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── S6 · FLOW · read_csv 四關鍵參數 ─────────
    s = _blank(prs)
    draw_flow_chain(s, nodes=[
        {"label": "encoding", "sub": "utf-8-sig / big5", "highlight": True},
        {"label": "sep", "sub": "',' / '\\t' / ';'"},
        {"label": "dtype", "sub": "釘住欄位型別", "highlight": True},
        {"label": "parse_dates", "sub": "日期一開始\n就是 datetime64"},
    ], title="read_csv 一行背後的四個關鍵參數", y=3.0)
    add_textbox(
        s, T.MARGIN_X, Inches(5.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.7),
        "這四個參數決定後面要不要重來——一次讀對，比之後 astype 補救便宜十倍。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "pandas User Guide §IO tools (text, CSV, HDF5, …)")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───────── S7 · CODE · 四參數實戰 ─────────
    s = _blank(prs)
    add_title(s, "encoding / sep / dtype / parse_dates：讓讀檔一次對")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="一次把四個關鍵參數都指定好",
        code=(
            'df = pd.read_csv(\n'
            '    "orders_raw.csv",\n'
            '    encoding="utf-8-sig",       # Windows 中文檔常帶 BOM\n'
            '    sep=",",                     # TSV 改 "\\t"\n'
            '    dtype={                      # 釘住，後面不用 astype\n'
            '        "order_id":   "string",\n'
            '        "customer_id":"string",\n'
            '        "qty":        "Int64",   # 可為 NaN 的整數\n'
            '    },\n'
            '    parse_dates=["order_date"],  # 一次變 datetime64\n'
            '    na_values=["-", "N/A", "?", ""],\n'
            ')\n'
        ),
        bullets=[
            "utf-8-sig 吃 BOM；\nbig5 給舊系統",
            "Int64（大寫）是\n可為 NaN 的整數",
            "parse_dates 省掉\n一次 to_datetime",
            "na_values 在讀檔\n就把「-」「N/A」\n統一變 NaN",
        ],
        label_dark=True,
    )
    add_source(s, "pandas docs §pandas.read_csv · §Nullable integer data type")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── S8 · PITFALL · 讀 CSV 沒 encoding ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="❌ 常見寫法：中文變亂碼",
        right_title="✅ 推薦寫法：讀前先想 encoding",
        left_items=[
            'pd.read_csv("orders.csv")',
            "預設 utf-8，Windows 匯出常帶 BOM → 首欄名變 '\\ufefforder_id'",
            "遇 big5 舊系統直接 UnicodeDecodeError",
            "在下游用 rename 補救是治標不治本",
        ],
        right_items=[
            'pd.read_csv("orders.csv", encoding="utf-8-sig")',
            "utf-8-sig：吃掉 BOM、同時相容純 utf-8",
            '遇舊系統：encoding="big5" 或 "cp950"',
            "源頭解決，後面程式碼都乾淨",
        ],
        title="讀 CSV 不指定 encoding——中文檔的第一個坑",
        summary="Windows 匯出的 CSV 幾乎都帶 BOM；utf-8-sig 是安全預設。",
    )
    add_source(s, "pandas docs §read_csv encoding · Python Docs §codecs")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── S9 · MATRIX 2×2 · loc vs iloc ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "df.loc[label]",
         "sub": "用 Index / 欄名 / 布林取\n包頭包尾（SQL 式）\n業務邏輯 90% 用這個",
         "highlight": True},
        {"text": "df.iloc[pos]",
         "sub": "用整數位置取\n包頭不包尾（Python list 式）\n對資料順序敏感"},
        {"text": "df[df.col > x]",
         "sub": "布林遮罩\n最 Pythonic\n多條件用 & | 加括號",
         "highlight": True},
        {"text": 'df.query("col > @x")',
         "sub": "字串條件、可讀性最高\n可用 @var 帶外部變數\n複雜邏輯不會一坨括號"},
    ], title="loc / iloc / 布林 / query：取資料的四把鑰匙")
    add_source(s, "pandas User Guide §Indexing and selecting data")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── S10 · CODE · loc vs iloc 切片差異 ─────────
    s = _blank(prs)
    add_title(s, "同樣的 0:3，loc 包後、iloc 不包後——一組 demo 看清楚")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.3),
        label="iloc — 用整數位置（包頭不包尾，與 list 一致）",
        code=(
            'df.iloc[0:3, 0:2]\n'
            '# 取位置 0、1、2 三列、位置 0、1 兩欄\n'
            '# end 不含 → 與 list[0:3] 同邏輯'
        ),
        bullets=[
            "對資料順序敏感",
            "排序後同樣寫法\n結果就變",
            "ETL / 隨機抽樣適用",
        ],
        label_dark=False,
    )
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(3.7),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(2.5),
        label="loc — 用標籤（包頭包尾，像 SQL BETWEEN）",
        code=(
            'df = df.set_index("order_id")\n'
            'df.loc["A001":"A003", "name":"amount"]\n'
            '# 取 Index 為 A001~A003 三列（含 A003）\n'
            '# 欄位 name 到 amount（含 amount）\n'
            '\n'
            'df.loc[df["age"] > 30, ["name","city"]]\n'
            '# 布林條件 + 欄位選取 一句到位'
        ),
        bullets=[
            "業務邏輯 90% 用 loc",
            "可放布林、標籤、切片",
            "set_index 後 loc[0]\n與 iloc[0] 完全不同列",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Indexing — Different choices for indexing")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── S11 · PITFALL · 布林用 iloc 會爆 ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="❌ 布林用 iloc：直接 raise",
        right_title="✅ 布林用 loc 或直接 []",
        left_items=[
            'df.iloc[df["age"] > 18]',
            "→ IndexError: iLocation based boolean indexing on an integer type is not available",
            "iloc 只吃整數位置，布林 Series 不吃",
            "新手最常踩的第三名 bug",
        ],
        right_items=[
            'df.loc[df["age"] > 18]',
            'df[df["age"] > 18]    # 等價簡寫',
            "條件用 loc / []、位置用 iloc",
            "寫錯時退回去想：我用的是『條件』還是『位置』？",
        ],
        title="布林條件要用 loc，寫成 iloc 直接 raise",
        summary="loc 吃（標籤 / 布林 / 切片）三種、iloc 只吃整數位置。",
    )
    add_source(s, "pandas docs §Boolean indexing · GitHub issues 常見求助 Top 10")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── S12 · FLOW · 清理四大手法 ─────────
    s = _blank(prs)
    draw_flow_chain(s, nodes=[
        {"label": "① 欄名標準化", "sub": "strip / lower\nreplace ' ' → '_'", "highlight": True},
        {"label": "② 型別轉換", "sub": "to_datetime\nto_numeric\nerrors='coerce'"},
        {"label": "③ 缺值處理", "sub": "依欄位意義\ndropna / fillna", "highlight": True},
        {"label": "④ 重複移除", "sub": "drop_duplicates\nsubset=['id']"},
    ], title="清理四大手法：順序不要亂", y=3.0)
    add_textbox(
        s, T.MARGIN_X, Inches(5.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.7),
        "順序有邏輯：先正名（key 才對得上）→ 再定型別（NaN 才正確）→ 再補缺 → 最後去重。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "Wickham, Tidy Data 2014 · pandas User Guide §Working with missing data")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── S13 · CODE · 欄名標準化 ─────────
    s = _blank(prs)
    add_title(s, "欄名標準化：三行下去，未來程式碼不會再崩潰")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="一次把 Customer ID / customer_id / CustomerID 統一",
        code=(
            '# BEFORE ：欄名花式寫法滿天飛\n'
            'df.columns\n'
            '# ["Order ID", " Customer ID", "Order_Date", "amount "]\n'
            '\n'
            '# AFTER ：三行標準化\n'
            'df.columns = (\n'
            '    df.columns\n'
            '      .str.strip()                  # 去前後空白\n'
            '      .str.lower()                  # 全部小寫\n'
            '      .str.replace(" ", "_")        # 空白換底線\n'
            ')\n'
            '# ["order_id", "customer_id", "order_date", "amount"]\n'
            '\n'
            '# 特殊情境：重複欄名\n'
            'df = df.rename(columns={"amount.1": "amount_usd"})'
        ),
        bullets=[
            "「免費」的投資\n成本三行、收益整份專案",
            "S2 之後所有 merge、\ngroupby 才不會出意外",
            "遇到中文欄名：\n先 rename 成英文\n再標準化",
            "重複欄名用 rename\n個別處理",
        ],
        label_dark=True,
    )
    add_source(s, "pandas docs §Working with text data · Data engineering best practices")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── S14 · CODE · 型別轉換 + errors='coerce' ─────────
    s = _blank(prs)
    add_title(s, "型別轉換的救命符：errors='coerce'")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="coerce：把不合法值悄悄變 NaT / NaN，而不是整批爆炸",
        code=(
            '# 1. 日期：遇到 "2023/13/45" 會爆 → coerce 後變 NaT\n'
            'df["order_date"] = pd.to_datetime(\n'
            '    df["order_date"], errors="coerce")\n'
            '\n'
            '# 2. 金額：字串 "$1,355" 不能直接 astype(float)\n'
            'df["amount"] = (\n'
            '    df["amount"]\n'
            '      .str.replace("$", "", regex=False)\n'
            '      .str.replace(",", "", regex=False)\n'
            ')\n'
            'df["amount"] = pd.to_numeric(df["amount"], errors="coerce")\n'
            '\n'
            '# 3. 驗證 ：coerce 後務必回頭看缺值數\n'
            'df[["order_date", "amount"]].isna().sum()\n'
            '# 若突然暴增 → 原始資料有奇怪格式，回頭檢查'
        ),
        bullets=[
            "errors='coerce' 是救命符：\n處理髒資料不爆炸",
            "副作用：『失敗悄悄發生』\n清完務必檢查缺值",
            "金額 / 百分比 / 電話\n都是字串偽裝的數字",
            "不要用 astype(float) \n硬轉髒欄位",
        ],
        label_dark=True,
    )
    add_source(s, "pandas docs §to_datetime · §to_numeric errors parameter")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── S15 · THREE BLOCKS · 缺值三策略 ─────────
    s = _blank(prs)
    draw_three_blocks_flow(s, blocks=[
        {"heading": "訂單日期缺失",
         "items": [
             "無日期 = 無法時序分析",
             "直接 dropna",
             'df.dropna(subset=["order_date"])',
             "關鍵欄缺失就丟",
         ]},
        {"heading": "金額缺失",
         "items": [
             "數值欄、有分佈",
             "中位數比平均穩",
             'df["amount"].fillna(\n  df["amount"].median())',
             "別無腦 fillna(0)",
         ]},
        {"heading": "類別缺失",
         "items": [
             "類別欄、可能含意義",
             "補 'Unknown' 保留訊號",
             'df["channel"].fillna("Unknown")',
             "別無腦 dropna",
         ]},
    ], title="缺值處理沒有標準答案——依欄位意義決定",
        bottom_note="策略選擇要問三件事：這欄是什麼？缺多少？下游誰會用？")
    add_source(s, "pandas User Guide §Working with missing data · Hadley Wickham, Tidy Data")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── S16 · PITFALL · inplace / append / dropna 地雷 ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="❌ 三個 Review 會被打槍的寫法",
        right_title="✅ 改寫：語意清楚、版本安全",
        left_items=[
            'df.drop("col", axis=1, inplace=True)',
            "→ inplace 新版不建議、讓人搞不清變數狀態",
            "df = df.append(other)",
            "→ append 已 deprecated（pandas 2.0 移除）",
            'df.dropna()',
            "→ 無 subset 會砍到不想砍的欄，資料爆縮",
        ],
        right_items=[
            'df = df.drop("col", axis=1)',
            "→ 賦值寫法，可讀、可鏈式",
            'df = pd.concat([df, other], ignore_index=True)',
            "→ concat 是新版標準",
            'df.dropna(subset=["order_date"])',
            "→ 只砍關鍵欄缺失的列",
        ],
        title="三個 pandas 2.x 時代要改掉的經典錯",
        summary="賦值 > inplace、concat 取代 append、dropna 一定加 subset。",
    )
    add_source(s, "pandas 2.0 Release Notes · Deprecation warnings")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── S17 · FLOW · 七步 ETL ─────────
    s = _blank(prs)
    draw_flow_chain(s, nodes=[
        {"label": "① 讀檔", "sub": "read_csv\n+ encoding", "highlight": True},
        {"label": "② 五件事", "sub": "shape / head /\ninfo / describe /\nisna"},
        {"label": "③ 欄名", "sub": "strip / lower\nreplace"},
        {"label": "④ 型別", "sub": "to_datetime\nto_numeric\ncoerce", "highlight": True},
        {"label": "⑤ 缺值", "sub": "dropna / fillna\n依欄位"},
        {"label": "⑥ 去重", "sub": "drop_duplicates"},
        {"label": "⑦ 輸出", "sub": "to_csv\nindex=False", "highlight": True},
    ], title="orders_raw → orders_clean：七步微型 ETL", y=3.0)
    add_textbox(
        s, T.MARGIN_X, Inches(5.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.7),
        "記住這七步的順序——未來任何髒檔，你都能複製這條流水線。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "本課程 Notebook M2_Pandas_Basic/S2_pandas_io_cleaning.ipynb")
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───────── S18 · CODE · 七步一次跑完 ─────────
    s = _blank(prs)
    add_title(s, "七步一次跑完：從 orders_raw.csv 到 orders_clean.csv")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="這 18 行 = Notebook Case，等一下練習請封裝成 clean_orders(path)",
        code=(
            '# ① 讀檔\n'
            'df = pd.read_csv("orders_raw.csv", encoding="utf-8-sig",\n'
            '                 na_values=["-", "N/A", "?", ""])\n'
            '# ② 五件事（實務省略，debug 時跑）\n'
            '# df.shape; df.head(); df.info(); df.describe(); df.isna().sum()\n'
            '\n'
            '# ③ 欄名標準化\n'
            'df.columns = df.columns.str.strip().str.lower().str.replace(" ", "_")\n'
            '\n'
            '# ④ 型別轉換\n'
            'df["order_date"] = pd.to_datetime(df["order_date"], errors="coerce")\n'
            'df["amount"] = (df["amount"].str.replace("$", "", regex=False)\n'
            '                              .str.replace(",", "", regex=False))\n'
            'df["amount"] = pd.to_numeric(df["amount"], errors="coerce")\n'
            '\n'
            '# ⑤ 缺值處理\n'
            'df = df.dropna(subset=["order_date"])\n'
            'df["amount"] = df["amount"].fillna(df["amount"].median())\n'
            'df["channel"] = df["channel"].fillna("Unknown")\n'
            '\n'
            '# ⑥ 去重 & ⑦ 輸出\n'
            'df = df.drop_duplicates(subset=["order_id"])\n'
            'df.to_csv("orders_clean.csv", index=False, encoding="utf-8-sig")'
        ),
        bullets=[
            "一次看完七步流水線",
            "每一行對應前張流程圖",
            "下一張用數字驗證\n到底做了什麼",
            "練習：封裝成函式\n+ type hints + docstring",
        ],
        label_dark=True,
    )
    add_source(s, "本課程 Notebook Case · clean_orders reference implementation")
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ───────── S19 · TABLE · 清理前 vs 清理後 ─────────
    s = _blank(prs)
    draw_editorial_table(
        s,
        header=["指標", "清理前 (raw)", "清理後 (clean)", "變化"],
        rows=[
            ["總列數", "10,842", "10,315", "−527 列（重複+無效日期）"],
            ["dtypes 正確率", "3 / 7 欄", "7 / 7 欄", "+4 欄（日期、金額、ID）"],
            ["order_date 缺失", "411", "0", "dropna subset"],
            ["amount 缺失", "156", "0", "fillna median"],
            ["channel 缺失", "89", "0", "fillna 'Unknown'"],
            ["重複 order_id", "318", "0", "drop_duplicates"],
        ],
        col_widths=[1.4, 1.6, 1.6, 2.4],
        title="清理前 vs 清理後：用數字驗證我們做了什麼",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.0),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "清理完務必出這張表——不是憑感覺說『我清好了』，是用數字證明。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "本課程 Notebook 驗證儲存格 · Data quality reporting pattern")
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───────── S20 · PITFALL · 六個錯誤綜合 ─────────
    s = _blank(prs)
    draw_editorial_table(
        s,
        header=["常犯錯誤", "為什麼錯", "正確寫法"],
        rows=[
            ['pd.to_datetime(df["date"])',
             "遇到 '2023/13/45' 整批 raise",
             'pd.to_datetime(df["date"], errors="coerce")'],
            ['df["amount"].astype(float)',
             "'$1,355' 字串無法轉 float",
             '.str.replace("$","").str.replace(",","") + to_numeric'],
            ['df.iloc[df["age"] > 18]',
             "iloc 不吃布林 Series → IndexError",
             'df.loc[df["age"] > 18]'],
            ['df.drop(..., inplace=True)',
             "新版不建議、變數狀態混亂",
             'df = df.drop(...)'],
            ['df.append(other)',
             "pandas 2.0 已移除",
             'pd.concat([df, other], ignore_index=True)'],
            ['pd.read_csv("orders.csv")',
             "Windows 中文檔變亂碼",
             'encoding="utf-8-sig"（或 big5）'],
        ],
        col_widths=[2.0, 2.4, 2.6],
        title="六個常犯錯誤一次看清（印下來貼在螢幕邊）",
    )
    add_source(s, "GitHub issues 彙整 · pandas 2.0 Migration Guide")
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ───────── S21 · ASK · 練習題 + 討論題 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "把剛剛 18 行封裝成 clean_orders(path) ——\n並回答三個討論題",
        data_card={
            "label": "課堂練習（40 min）",
            "stat": "2 + 3",
            "caption": "兩題動手 + 三題討論\n🟡 跑通 · 🔴 封裝",
        },
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.5),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(1.5),
        ("討論題 ①　100 萬筆訂單、5% 缺 amount / 3% 缺 order_date / 1% 缺 customer_id，\n"
         "　　　　　你會怎麼決定各欄位策略？依據是什麼？\n"
         "討論題 ②　loc 和 iloc 看起來很像，為什麼 Pandas 故意設計兩個？只保留一個會有什麼問題？\n"
         "討論題 ③　清完的資料要不要覆蓋原始檔？團隊協作時哪種做法比較安全？"),
        font_size=T.FONT_SMALL, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    add_source(s, "teacher_notes §5 Discussion Prompts")
    add_footer(s, MODULE_CODE, 21, N_CONTENT)

    # ───────── S22 · CHECKPOINT · MVK 回顧 ─────────
    s = _blank(prs)
    draw_thesis_hierarchy(
        s,
        blocks=[
            {"heading": "概念 MVK（腦裡要有）",
             "items": [
                 "三件套：DataFrame / Series / Index",
                 "五件事 SOP：shape / head / info / describe / isna",
                 "loc vs iloc：標籤包尾、位置不包尾",
                 "清理四大手法順序：欄名 → 型別 → 缺值 → 重複",
             ]},
            {"heading": "動作 MVK（手要會打）",
             "items": [
                 'read_csv(encoding=, dtype=, parse_dates=, na_values=)',
                 'columns.str.strip().lower().replace(" ", "_")',
                 'to_datetime(..., errors="coerce") · to_numeric(...)',
                 'dropna(subset=[...]) · drop_duplicates(subset=[...])',
             ]},
        ],
        title="S2 MVK 回顧：腦裡有觀念、手上有八招",
        thesis="資料清理不是憑手感——是按順序做完這八件事。",
    )
    add_source(s, "本節 05_mvk.md · 對應 teacher_notes §1 Learning Objectives")
    add_footer(s, MODULE_CODE, 22, N_CONTENT)

    # ───────── S23 · SILENT · 下節預告 ─────────
    s = _blank(prs)
    draw_silent_page(s, "下節 S3 —— 把乾淨資料變成洞察：\ngroupby · merge · pivot。")
    add_footer(s, MODULE_CODE, 23, N_CONTENT, dark_bg=True)

    # ───────── Copyright ─────────
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
