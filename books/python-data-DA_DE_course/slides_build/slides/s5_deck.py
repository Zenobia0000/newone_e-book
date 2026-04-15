"""S5 deck — 視覺化精華：5 種必懂圖
24 content slides + cover + copyright page (S5 · 120 min).

Governing thought:
    5 張圖吃掉 90% DA 場景 —
    問題決定圖：趨勢/比較/關聯/分布/矩陣，
    seaborn 打底、matplotlib 微調、subplots 封裝儀表板。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_rect, add_textbox, add_title, add_source,
    set_solid_fill, set_no_line, set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_vs_two_col, draw_code_panel, draw_thesis_hierarchy,
    draw_inverted_thesis_box, draw_image_placeholder,
    draw_flow_chain, draw_risk_mitigation, draw_delta_badge,
    draw_emphasis_pill,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "S5"
MODULE_TITLE = "視覺化精華：5 種必懂圖"
MODULE_SUBTITLE = "seaborn 打底 × matplotlib 微調 × 2×3 儀表板交付"
TIME_MIN = 120
N_CONTENT = 24


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_s5(output_path, image_registry=None):
    """Build S5 deck.

    image_registry: optional list collecting image placeholder metadata
                    for downstream YAML emission.
    """
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── S1 · SILENT ─────────
    s = _blank(prs)
    draw_silent_page(s, "數字不會說話。\n圖會。")
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ───────── S2 · ASK ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "老闆看完表格\n還是問『所以誰最好？』",
        data_card={
            "label": "閱讀時間對比（12 列銷售資料）",
            "stat": "4s vs 37s",
            "caption": "圖 4 秒解完 · 表格要 37 秒\n老闆時間貴、要的是訊號不是資料",
        },
    )
    add_source(s, "Few, Show Me the Numbers 2e §Visual Perception")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── S3 · MATRIX 2×3 — 五問題 → 五圖 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=3, cells=[
        {"text": "趨勢 · 時間",
         "sub": "→ 折線 lineplot\nx=時間、y=量\nmarker='o' 看點",
         "highlight": True},
        {"text": "比較 · 誰大",
         "sub": "→ 長條 barplot\n先 sort_values\n橫排讀長標籤",
         "highlight": True},
        {"text": "關聯 · 兩數值",
         "sub": "→ 散佈 scatterplot\nalpha 降透明\nhue 上第三維"},
        {"text": "分布 · 離群",
         "sub": "→ 箱型 boxplot\n中位 / IQR / 離群\n三組以上才划算"},
        {"text": "矩陣 · 兩類別",
         "sub": "→ 熱力 heatmap\nannot=True\nfmt=',.0f'"},
        {"text": "心法",
         "sub": "問題決定圖\n不是品味決定\n先問再畫",
         "highlight": True},
    ], title="五個問題 → 五種圖（選圖思維）")
    add_source(s, "Tukey EDA 1977 · Knaflic Storytelling with Data · seaborn gallery")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───────── S4 · VS — 表格 vs 折線 ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="表格（給你讀的）",
        right_title="折線（給老闆看的）",
        left_items=[
            "12 列 × 2 欄數字",
            "要心算比大小",
            "趨勢藏在腦補裡",
            "平均閱讀時間 37 秒",
            "3 秒後注意力流失",
        ],
        right_items=[
            "一條線、一眼到底",
            "斜率 = 成長訊號",
            "峰谷直接顯影",
            "平均閱讀時間 4 秒",
            "下一個問題馬上問",
        ],
        title="同一份資料 · 兩種呈現 · 十倍效率差",
        summary="不是圖比較好看 — 是圖把「比較」這個認知動作交給了視覺皮層。",
        delta="9×",
    )
    add_source(s, "Cleveland & McGill 1984, Graphical Perception JASA")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── S5 · CODE — 環境三件套 ─────────
    s = _blank(prs)
    add_title(s, "環境三件套：notebook 首行寫一次、之後不再煩")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="S5 所有範例共用這 6 行 — 字型／負號／網格一次搞定",
        code=(
            'import pandas as pd\n'
            'import matplotlib.pyplot as plt\n'
            'import seaborn as sns\n'
            '\n'
            '# 1. seaborn 主題（配色 + 網格）\n'
            'sns.set_theme(style="whitegrid", palette="deep")\n'
            '\n'
            '# 2. 中文字型（Win / Mac / Linux 擇一）\n'
            'plt.rcParams["font.sans-serif"] = ["Microsoft JhengHei"]\n'
            '#   Mac   → ["PingFang TC"]\n'
            '#   Linux → ["Noto Sans CJK TC"]\n'
            '\n'
            '# 3. 修正負號變方塊\n'
            'plt.rcParams["axes.unicode_minus"] = False'
        ),
        bullets=[
            "sns.set_theme\n= 網格 + 字級 + 配色 三合一",
            "rcParams 是全域設定\n寫一次全 notebook 受益",
            "unicode_minus=False\n擋掉負號方塊",
            "Windows → JhengHei\nMac → PingFang TC",
        ],
        label_dark=True,
    )
    add_source(s, "seaborn docs §set_theme · matplotlib FAQ §Working with text")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── S6 · SILENT — 圖 1 折線 ─────────
    s = _blank(prs)
    draw_silent_page(s, "圖 1 · 折線。\n看的是趨勢。")
    add_footer(s, MODULE_CODE, 6, N_CONTENT, dark_bg=True)

    # ───────── S7 · IMAGE+CODE — 折線 ─────────
    s = _blank(prs)
    add_title(s, "sns.lineplot：月度趨勢一行畫")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.8), h=Inches(5.0),
        slot_name="折線圖 · 月度趨勢",
        description=(
            "深綠折線 + 空心圓 marker\n"
            "x 軸 = 月份、y 軸 = 銷售額（千元）\n"
            "可疊 2023 灰色對照線"
        ),
        url_hint="seaborn gallery §lineplot",
        placeholder_id="S5_ph01_line_monthly_trend",
        registry=image_registry,
        size_hint="1200×720 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.7), y=Inches(1.3),
        w=Inches(6.0), h=Inches(5.0),
        label="時間／有序序列 — lineplot 是直覺反應",
        code=(
            'fig, ax = plt.subplots(figsize=(8, 5))\n'
            '\n'
            '# 單條線\n'
            'sns.lineplot(data=df, x="month", y="sales",\n'
            '             marker="o", ax=ax,\n'
            '             color="#1B5E3F")\n'
            '\n'
            '# 多條線 → hue 自動分色 + legend\n'
            'sns.lineplot(data=df, x="month", y="sales",\n'
            '             hue="year", marker="o", ax=ax)\n'
            '\n'
            'ax.set_title("月度銷售趨勢（2023 vs 2024）")\n'
            'ax.set_xlabel("月份"); ax.set_ylabel("金額 (千元)")'
        ),
        bullets=[
            "marker='o' 強制顯點\n離散月份看得清楚",
            "hue 自動分色\n不用 for 迴圈",
            "ax= 明確指定\n儀表板才不畫錯格",
            "x 密時 rotation=45",
        ],
        label_dark=True,
    )
    add_source(s, "seaborn docs §lineplot · McKinney PDA 3e §9.4")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── S8 · PITFALL 折線 ─────────
    s = _blank(prs)
    draw_risk_mitigation(
        s,
        title="折線常見坑：三個讓圖說謊的錯誤",
        risks=[
            "類別資料硬套折線（如地區 x 軸）→ 錯覺趨勢",
            "x 軸日期太密 → 標籤重疊看不清",
            "忘了 marker → 離散點看不出",
            "多條線用同色 → 疊圖難辨識",
        ],
        mitigations=[
            "類別改 barplot；折線留給「有序」",
            "plt.xticks(rotation=45) + tight_layout",
            "一律 marker='o'（養成習慣）",
            "hue=分類欄，交給 seaborn 自動配色",
        ],
        summary="折線的定義：x 軸是「有序」。地區、產品類是類別 — 那不是折線的地盤。",
    )
    add_source(s, "seaborn docs §lineplot Common pitfalls · Tukey 1977")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── S9 · IMAGE+CODE — 長條 ─────────
    s = _blank(prs)
    add_title(s, "sns.barplot：排序 + 數字標註，老闆 2 秒懂")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.8), h=Inches(5.0),
        slot_name="長條圖 · 地區排名（已排序）",
        description=(
            "橫向長條、由大到小、深綠色\n"
            "每條末端 plt.text 標千分位金額\n"
            "y 軸 = 地區、x 軸 = 金額"
        ),
        url_hint="seaborn gallery §barplot",
        placeholder_id="S5_ph02_bar_region_sorted",
        registry=image_registry,
        size_hint="1200×720 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.7), y=Inches(1.3),
        w=Inches(6.0), h=Inches(5.0),
        label="比較 / 排名 — 排序是長條的靈魂",
        code=(
            '# 1. 先排序（靈魂一步）\n'
            'agg = (df.groupby("region")["sales"].sum()\n'
            '         .sort_values(ascending=True)\n'
            '         .reset_index())\n'
            '\n'
            'fig, ax = plt.subplots(figsize=(8, 5))\n'
            '\n'
            '# 2. seaborn 0.13+ : hue + legend=False\n'
            'sns.barplot(data=agg, x="sales", y="region",\n'
            '            hue="region", legend=False,\n'
            '            palette="Greens_r", ax=ax)\n'
            '\n'
            '# 3. 數字標註\n'
            'for i, v in enumerate(agg["sales"]):\n'
            '    ax.text(v, i, f" {v:,.0f}", va="center")\n'
            '\n'
            'ax.set_title("Q4 各地區銷售排名")'
        ),
        bullets=[
            "sort_values 先於畫圖\n否則按字母順序",
            "橫排（y=類別）\n適合長中文標籤",
            "seaborn 0.13+ palette\n必搭 hue+legend=False",
            "標註讓老闆不用讀軸",
        ],
        label_dark=True,
    )
    add_source(s, "seaborn 0.13 release notes · Knaflic 2015 §Bar charts")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── S10 · PITFALL 長條 ─────────
    s = _blank(prs)
    draw_risk_mitigation(
        s,
        title="長條常見坑：未排序是老闆最大的怒點",
        risks=[
            "直接 groupby().sum() 畫圖 → 按字母排，誰第一看不出",
            "seaborn 0.13+ 直接傳 palette → FutureWarning",
            "地區名稱長 → 標籤重疊糊掉",
            "多組長條不加 hue → 每組只看到同色",
        ],
        mitigations=[
            ".sort_values(ascending=False) 放在 barplot 前",
            "加 hue=x_col, legend=False（官方推薦寫法）",
            "改 x=數值、y=類別 橫排；或 xticks rotation=30",
            "hue='subgroup'，seaborn 自動分色 + legend",
        ],
        summary="長條圖 = 排名圖。未排序的長條圖 = 沒有結論的簡報。",
    )
    add_source(s, "seaborn 0.13 release · S5 teacher_notes §Common Pitfalls")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── S11 · IMAGE+CODE — 散佈 ─────────
    s = _blank(prs)
    add_title(s, "sns.scatterplot：hue 上色 + legend 外移")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.8), h=Inches(5.0),
        slot_name="散佈圖 · 單價 × 數量（依品類 hue）",
        description=(
            "alpha=0.6、hue=category(4 類)\n"
            "legend 右側外移（不擋資料）\n"
            "可見品類聚落 + 離群點"
        ),
        url_hint="seaborn gallery §scatterplot",
        placeholder_id="S5_ph03_scatter_hue_category",
        registry=image_registry,
        size_hint="1200×720 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.7), y=Inches(1.3),
        w=Inches(6.0), h=Inches(5.0),
        label="兩連續變數關係 — alpha + hue 是兩大利器",
        code=(
            'fig, ax = plt.subplots(figsize=(8, 5))\n'
            '\n'
            'sns.scatterplot(\n'
            '    data=df,\n'
            '    x="unit_price", y="quantity",\n'
            '    hue="category",       # 第三維：類別\n'
            '    alpha=0.6,            # 重疊降透明\n'
            '    s=60,                 # 點大小\n'
            '    ax=ax,\n'
            ')\n'
            '\n'
            '# legend 外移（不擋資料）\n'
            'ax.legend(title="品類",\n'
            '          bbox_to_anchor=(1.02, 1),\n'
            '          loc="upper left",\n'
            '          frameon=False)\n'
            '\n'
            'ax.set_title("單價 vs 數量（依品類）")'
        ),
        bullets=[
            "alpha=0.4~0.6\n重疊看密度",
            "hue 自動上色\n比手動 for 簡單",
            "bbox_to_anchor=(1.02,1)\n= legend 右側外移",
            "size= 可做第四維\n但別超過 3 個維度",
        ],
        label_dark=True,
    )
    add_source(s, "seaborn docs §scatterplot · Wilke Fundamentals ch12")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── S12 · PITFALL 散佈 ─────────
    s = _blank(prs)
    draw_risk_mitigation(
        s,
        title="散佈常見坑：彩虹 legend 與一坨點",
        risks=[
            "大量點重疊不用 alpha → 看起來只有一坨",
            "legend 擠在圖內 → 擋到資料點",
            "hue 類別 > 7 → 彩虹災難、顏色難辨",
            "把『無關聯』誤讀成『負關聯』",
        ],
        mitigations=[
            "alpha=0.4~0.6（密度即訊號）",
            "bbox_to_anchor=(1.02, 1) 外移",
            "改用 facet (sns.FacetGrid) 分面",
            "看趨勢線：sns.regplot 疊上回歸",
        ],
        summary="散佈圖回答的是「關聯」。類別太多時，把 hue 換成 facet — 一類一格。",
    )
    add_source(s, "seaborn docs §scatterplot · Cleveland 1993 Visualizing Data")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── S13 · IMAGE+CODE — 箱型 ─────────
    s = _blank(prs)
    add_title(s, "sns.boxplot：中位數 / IQR / 離群一次解剖")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.8), h=Inches(5.0),
        slot_name="箱型圖 · 各地區訂單金額",
        description=(
            "三地區並排箱型\n"
            "中位數線、IQR 盒、whisker、離群點\n"
            "y 軸 = 訂單金額"
        ),
        url_hint="seaborn gallery §boxplot",
        placeholder_id="S5_ph04_box_by_region",
        registry=image_registry,
        size_hint="1200×720 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.7), y=Inches(1.3),
        w=Inches(6.0), h=Inches(5.0),
        label="多組分布比較 — 先會看再會畫",
        code=(
            'fig, ax = plt.subplots(figsize=(8, 5))\n'
            '\n'
            'sns.boxplot(data=df, x="region", y="amount",\n'
            '            hue="region", legend=False,\n'
            '            palette="Greens", ax=ax)\n'
            '\n'
            'ax.set_title("各地區訂單金額分布")\n'
            'ax.set_ylabel("金額 (元)")\n'
            '\n'
            '# ─── 怎麼看一個箱子 ──────────────\n'
            '#   盒子       = IQR (Q1 ~ Q3，中間 50%)\n'
            '#   中間橫線   = 中位數\n'
            '#   鬚 whisker = 1.5 × IQR 範圍內極值\n'
            '#   鬚外的點   = 離群值（不是錯，是訊號）'
        ),
        bullets=[
            "三組以上才體現價值\n兩組改 hist 更直覺",
            "中位數 > 平均\n抗離群干擾",
            "離群點 ≠ 錯誤\n先問為什麼",
            "vert=False 橫排\n適合長標籤",
        ],
        label_dark=True,
    )
    add_source(s, "Tukey 1977 §Box-and-Whisker · seaborn docs §boxplot")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── S14 · PITFALL 箱型 ─────────
    s = _blank(prs)
    draw_risk_mitigation(
        s,
        title="箱型常見坑：離群點不是垃圾、兩組不划算",
        risks=[
            "只比兩組時用 box → 資訊量不如直方",
            "把離群點當錯誤、直接砍掉",
            "y 軸起點不 0 → 看起來分布被壓扁",
            "不看 hue 就畫 → 全部混成一箱",
        ],
        mitigations=[
            "兩組 → sns.histplot(..., hue=...) 更易讀",
            "先問離群原因（VIP / 異常 / Bug）再決定",
            "ax.set_ylim(0, ...) 或保留自動 + 加註",
            "x=分組欄、y=數值欄，分開才有比較意義",
        ],
        summary="離群點是訊號：可能是 VIP、可能是詐騙、可能是鍵盤打錯。先看再說，別先砍。",
    )
    add_source(s, "Tukey 1977 · Wickham & Stryjewski 2011 §40 years of boxplots")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── S15 · IMAGE+CODE — 熱力 ─────────
    s = _blank(prs)
    add_title(s, "sns.heatmap：矩陣型視覺（S3 pivot 的圖形版）")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.8), h=Inches(5.0),
        slot_name="熱力圖 · 地區 × 品類 總額",
        description=(
            "annot=True、fmt=',.0f'\n"
            "cmap='Greens'（單向：低 → 高）\n"
            "列=地區、行=品類，標題：Q4 總額矩陣"
        ),
        url_hint="seaborn gallery §heatmap",
        placeholder_id="S5_ph05_heatmap_region_x_category",
        registry=image_registry,
        size_hint="1200×720 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.7), y=Inches(1.3),
        w=Inches(6.0), h=Inches(5.0),
        label="兩個類別交叉的數值 — heatmap 一張搞定",
        code=(
            '# 1. 先做 pivot（S3 的知識派上用場）\n'
            'pivot = df.pivot_table(\n'
            '    index="region", columns="category",\n'
            '    values="amount", aggfunc="sum"\n'
            ')\n'
            '\n'
            'fig, ax = plt.subplots(figsize=(8, 5))\n'
            '\n'
            '# 2. 一行熱力圖\n'
            'sns.heatmap(\n'
            '    pivot,\n'
            '    annot=True,           # 格內寫數字\n'
            '    fmt=",.0f",           # 千分位整數（擋科學記號）\n'
            '    cmap="Greens",        # 單向：低 → 高\n'
            '    linewidths=0.5,\n'
            '    ax=ax,\n'
            ')\n'
            '\n'
            'ax.set_title("Q4 地區 × 品類 總額矩陣")'
        ),
        bullets=[
            "annot=True 數字補上\n不用切去讀色階",
            "fmt=',.0f' 擋科學記號",
            "單向值用 sequential\n(Greens / Blues / viridis)",
            "正負值才用 diverging\n(RdBu / coolwarm)",
        ],
        label_dark=True,
    )
    add_source(s, "seaborn docs §heatmap · Wilke ch16 Visualizing amounts")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── S16 · PITFALL 熱力 ─────────
    s = _blank(prs)
    draw_risk_mitigation(
        s,
        title="熱力常見坑：格子太擠、cmap 選錯",
        risks=[
            "類別 > 15 → 格子擠到看不見數字",
            "忘了 fmt → 數字跑科學記號（1.2e+06）",
            "單向值誤用 diverging cmap（RdBu）",
            "沒補 annot → 讀者要對色階猜數字",
        ],
        mitigations=[
            "先 .nlargest(10) 過濾，或切分面",
            "永遠加 fmt=',.0f' / ',.1f' / ',.0%'",
            "低→高用 Greens / viridis；有正負才 RdBu",
            "annot=True 是交付預設（除非純視覺探索）",
        ],
        summary="熱力圖的工作是「讓讀者不用查色階」。annot + fmt 是最低交付標配。",
    )
    add_source(s, "seaborn docs §heatmap · matplotlib §Choosing colormaps")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── S17 · TABLE — 五圖速查 ─────────
    s = _blank(prs)
    draw_editorial_table(
        s,
        header=["問題", "圖種", "一行 API", "關鍵參數 / 陷阱"],
        rows=[
            ["趨勢 / 時間",
             "折線 lineplot",
             "sns.lineplot(x, y, marker='o', hue)",
             "marker / hue；x 必須有序"],
            ["比較 / 排名",
             "長條 barplot",
             "sns.barplot(x, y, hue=x, legend=False)",
             "sort_values 先行"],
            ["關聯 / 兩變數",
             "散佈 scatterplot",
             "sns.scatterplot(x, y, hue, alpha=0.6)",
             "alpha / legend bbox"],
            ["分布 / 離群",
             "箱型 boxplot",
             "sns.boxplot(x, y, hue=x, legend=False)",
             "三組以上才划算"],
            ["矩陣 / 兩類別",
             "熱力 heatmap",
             "sns.heatmap(pivot, annot=True, fmt=',.0f')",
             "fmt 擋科學記號；cmap 選對"],
        ],
        col_widths=[1.6, 1.8, 4.0, 2.8],
        title="五圖速查：問題 → 圖 → API → 陷阱",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.6),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "這張截圖放工作桌布。遇到新需求 30 秒內挑好圖。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "S5 module synthesis · seaborn gallery overview")
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───────── S18 · CODE — subplots 骨架 ─────────
    s = _blank(prs)
    add_title(s, "plt.subplots(2, 3)：axes[i, j] 儀表板骨架")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="所有儀表板的共同模板 — 只有中間 6 行會變",
        code=(
            '# 2 列 × 3 行 = 6 張子圖\n'
            'fig, axes = plt.subplots(2, 3, figsize=(15, 8))\n'
            '\n'
            '# axes 是 2D numpy array，用 [i, j] 取\n'
            'sns.lineplot(..., ax=axes[0, 0])     # 左上：趨勢\n'
            'sns.barplot (..., ax=axes[0, 1])     # 中上：排名\n'
            'sns.scatterplot(..., ax=axes[0, 2])  # 右上：關聯\n'
            'sns.boxplot(..., ax=axes[1, 0])      # 左下：分布\n'
            'sns.heatmap(..., ax=axes[1, 1])      # 中下：矩陣\n'
            'axes[1, 2].text(0.5, 0.5, "KPI", ...)# 右下：KPI 文字\n'
            '\n'
            '# 收束三件套\n'
            'fig.suptitle("Q4 Sales Dashboard", fontsize=16)\n'
            'fig.tight_layout()\n'
            'fig.savefig("dashboard.png", dpi=300, bbox_inches="tight")\n'
            'plt.show()    # ← savefig 永遠在 show 前'
        ),
        bullets=[
            "axes[i, j] 統一寫法\n不要混 axes[i][j]",
            "ax= 明確指定\nseaborn 才畫對格",
            "suptitle 總標題\ntight_layout 防撞",
            "savefig → show 順序\n反了 = 存檔空白",
        ],
        label_dark=True,
    )
    add_source(s, "matplotlib docs §Arranging Axes · §tight_layout guide")
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ───────── S19 · IMAGE+CODE — Dashboard ─────────
    s = _blank(prs)
    add_title(s, "Q4 Sales Dashboard：履歷級 2×3 儀表板")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.8), h=Inches(5.0),
        slot_name="Q4 Sales Dashboard 2×3",
        description=(
            "六格整合：\n"
            "(0,0) 月趨勢  (0,1) 地區排名  (0,2) 品類散佈\n"
            "(1,0) 地區箱型  (1,1) 熱力矩陣  (1,2) KPI 區塊\n"
            "suptitle = Q4 Sales Dashboard"
        ),
        url_hint="自產 png，dpi=300",
        placeholder_id="S5_ph06_q4_dashboard_2x3",
        registry=image_registry,
        size_hint="1800×1100 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.7), y=Inches(1.3),
        w=Inches(6.0), h=Inches(5.0),
        label="六格整合範例（把前面 5 圖塞進 2×3）",
        code=(
            'fig, axes = plt.subplots(2, 3, figsize=(15, 8))\n'
            '\n'
            '# 上列：訊號\n'
            'sns.lineplot(data=df_m, x="month", y="sales",\n'
            '             marker="o", ax=axes[0, 0])\n'
            'axes[0, 0].set_title("月度趨勢")\n'
            '\n'
            'sns.barplot(data=rank, x="sales", y="region",\n'
            '            hue="region", legend=False,\n'
            '            palette="Greens_r", ax=axes[0, 1])\n'
            'axes[0, 1].set_title("地區排名")\n'
            '\n'
            'sns.scatterplot(data=df, x="unit_price", y="qty",\n'
            '                hue="category", alpha=0.6,\n'
            '                ax=axes[0, 2])\n'
            '\n'
            '# 下列：分布 / 矩陣 / KPI\n'
            'sns.boxplot (..., ax=axes[1, 0])\n'
            'sns.heatmap(pivot, annot=True, fmt=",.0f",\n'
            '            cmap="Greens", ax=axes[1, 1])\n'
            '\n'
            'axes[1, 2].axis("off")\n'
            'axes[1, 2].text(0.5, 0.5,\n'
            '  "Q4 總額\\n$1.23M\\n+18% YoY",\n'
            '  ha="center", va="center", fontsize=20)\n'
            '\n'
            'fig.suptitle("Q4 Sales Dashboard", fontsize=16)\n'
            'fig.tight_layout()'
        ),
        bullets=[
            "上列看訊號\n下列看細節",
            "KPI 格用 text\n不是每格都要畫圖",
            "suptitle 總結主題\n讀者 3 秒抓重點",
            "這張截圖就是履歷",
        ],
        label_dark=True,
    )
    add_source(s, "S5 Capstone · Knaflic 2015 §Dashboard design")
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───────── S20 · PITFALL Dashboard ─────────
    s = _blank(prs)
    draw_risk_mitigation(
        s,
        title="儀表板常見坑：axes 索引 / tight_layout / savefig 順序",
        risks=[
            "axes[0][1] 與 axes[0, 1] 混用 → debug 無解",
            "1 列或 1 行時 axes 變 1D → [i, j] 會爆",
            "沒 tight_layout → 標題、軸標撞版",
            "plt.show() 在 savefig 前 → 存檔是空圖",
            "六格塞太滿、每格都重裝飾 → 認知過載",
        ],
        mitigations=[
            "全程統一 axes[i, j]（Linus：明確 > 隱式）",
            "subplots(1, 3) → axes[j]；可用 axes = np.atleast_2d(axes)",
            "fig.tight_layout() 是儀表板的最後一步",
            "順序固定：tight_layout → savefig → show",
            "每格只做一件事；裝飾留給 suptitle",
        ],
        summary="儀表板不是「畫更多」，是「挑對 6 件事、整齊地講一個故事」。",
    )
    add_source(s, "S5 teacher_notes §Common Pitfalls · matplotlib tight_layout")
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ───────── S21 · VS seaborn vs matplotlib ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="seaborn（打底）",
        right_title="matplotlib（微調）",
        left_items=[
            "一行畫完 90% 場景",
            "統計類型圖內建（box / violin / reg）",
            "hue / palette 自動分色",
            "吃 DataFrame 欄名",
            "新手→產出可交付圖的最短路徑",
        ],
        right_items=[
            "客製座標軸、雙 y 軸、annotation",
            "spines / ticks / locator 精修",
            "subplots / gridspec 複雜排版",
            "savefig / backend / rcParams",
            "seaborn 最後的一哩路",
        ],
        title="seaborn vs matplotlib：分工、不是對立",
        summary="seaborn 打 90% 底，matplotlib 做最後 10% 精修 — 用 ax= 參數把兩者銜接起來。",
        delta="90/10",
    )
    add_source(s, "seaborn docs §Overview · matplotlib lifecycle guide")
    add_footer(s, MODULE_CODE, 21, N_CONTENT)

    # ───────── S22 · PRACTICE PROMPT ─────────
    s = _blank(prs)
    add_title(s, "課堂練習：Electronics 迷你儀表板（2×2 · 10 min）")

    # Prompt panel
    prompt_box = add_rect(s, T.MARGIN_X, Inches(1.3),
                          T.SLIDE_W - 2 * T.MARGIN_X, Inches(2.5))
    set_no_fill(prompt_box)
    set_line(prompt_box, T.PRIMARY, 1.5)
    add_textbox(
        s, T.MARGIN_X + Inches(0.1), Inches(1.3),
        Inches(3.0), Inches(0.4),
        "  PROMPT",
        font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
    )
    # Header strip
    hdr = add_rect(s, T.MARGIN_X, Inches(1.3),
                   Inches(1.6), Inches(0.35))
    set_solid_fill(hdr, T.PRIMARY)
    set_no_line(hdr)
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(1.6), Inches(0.35),
        "PROMPT",
        font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        s, T.MARGIN_X + Inches(0.3), Inches(1.8),
        T.SLIDE_W - 2 * T.MARGIN_X - Inches(0.6), Inches(2.0),
        "用 df[df['category']=='Electronics'] 子集，畫一張 2×2 儀表板：\n"
        "  (0,0) 月度趨勢（lineplot）   (0,1) Top 5 商品排名（barplot，sort！）\n"
        "  (1,0) 單價 vs 數量散佈      (1,1) 各地區金額箱型\n"
        "要求：中文字型、suptitle=\"Electronics 2024 Q4\"、savefig dpi=300。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, line_spacing=1.5,
    )

    # Rubric panel
    rubric_y = Inches(4.0)
    add_textbox(
        s, T.MARGIN_X, rubric_y, Inches(4.0), Inches(0.4),
        "評分 Rubric（自評 / 互評）",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, T.MARGIN_X, rubric_y + Inches(0.45),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(2.0),
        "• 選圖正確（4 格各對應：趨勢／排名／關聯／分布）\n"
        "• 長條圖有排序、有數字標註\n"
        "• 中文標題 / 軸標 無方塊字\n"
        "• suptitle + tight_layout 有執行\n"
        "• savefig 順序正確、輸出檔可交付",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, line_spacing=1.5,
    )
    add_source(s, "S5 teacher_notes §Discussion Prompts · §Practice")
    add_footer(s, MODULE_CODE, 22, N_CONTENT)

    # ───────── S23 · FLOW 選圖工作流 ─────────
    s = _blank(prs)
    draw_flow_chain(
        s,
        nodes=[
            {"label": "問題", "sub": "看什麼？",
             "caption": "趨勢 / 比較 / 關聯 / 分布 / 矩陣",
             "highlight": True},
            {"label": "圖種", "sub": "line / bar / scatter / box / heatmap",
             "caption": "5 選 1 · 30 秒內決定"},
            {"label": "排版", "sub": "subplots(r, c)",
             "caption": "單圖 or 儀表板",
             "highlight": True},
            {"label": "交付", "sub": "savefig dpi=300",
             "caption": "png / pdf / svg"},
        ],
        title="選圖工作流：問題 → 圖種 → 排版 → 交付",
        y=2.8,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.6),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.8),
        "拿到新需求 → 先問「看什麼」→ 30 秒選圖 → 10 分鐘畫完 → savefig 交付。\n"
        "這不是教條，是 DA 每天重複 10 次的例行公事。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER, line_spacing=1.4,
    )
    add_source(s, "S5 module synthesis · DA daily workflow")
    add_footer(s, MODULE_CODE, 23, N_CONTENT)

    # ───────── S24 · PYRAMID 收束 + S6 預告 ─────────
    s = _blank(prs)
    draw_thesis_hierarchy(
        s,
        blocks=[
            {"heading": "五圖速查（一輩子夠用）",
             "items": [
                 "趨勢 → sns.lineplot(x, y, marker='o')",
                 "比較 → sns.barplot(x, y, hue=x, legend=False) + sort",
                 "關聯 → sns.scatterplot(x, y, hue, alpha=0.6)",
                 "分布 → sns.boxplot(x, y) · 看中位 / IQR / 離群",
                 "矩陣 → sns.heatmap(pivot, annot=True, fmt=',.0f')",
             ]},
            {"heading": "四條紀律（帶回工作）",
             "items": [
                 "先 fig, ax = plt.subplots() 再畫",
                 "排序是長條圖的靈魂（sort_values 永遠先行）",
                 "seaborn 0.13+ palette 必搭 hue + legend=False",
                 "savefig 永遠在 plt.show() 之前",
             ]},
        ],
        title="S5 收束：5 圖 + 4 紀律",
        thesis="S6 把這 5 張圖換上 Plotly — 靜態變互動，讓老闆自己點、自己篩、自己下鑽。",
    )
    add_source(s, "S5 module synthesis · 銜接 S6 Plotly 互動與 Capstone")
    add_footer(s, MODULE_CODE, 24, N_CONTENT)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
