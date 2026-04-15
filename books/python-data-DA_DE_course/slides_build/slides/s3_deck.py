"""S3 deck — Pandas 轉換：groupby / merge / pivot
22 content slides + cover + copyright page (2 hr session).

Governing thought:
    groupby 是思考模型（Split→Apply→Combine），
    merge 是多表通行證（on + how），
    pivot 是交叉視角（index × columns × values）。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_rect, add_textbox, add_title, add_source,
    set_solid_fill, set_no_line, set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_vs_two_col, draw_code_panel, draw_flow_chain,
    draw_three_blocks_flow, draw_thesis_hierarchy, draw_image_placeholder,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "S3"
MODULE_TITLE = "Pandas 轉換：groupby / merge / pivot"
MODULE_SUBTITLE = "從單表查詢跨到多表分析 — 三把刀吃掉 80% 的商業提問"
TIME_MIN = 120
N_CONTENT = 22


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_s3(output_path, image_registry=None):
    """Build S3 deck.

    image_registry: optional list collecting image placeholder metadata
                    for downstream YAML emission.
    """
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── S1 · MOTIVATION · SILENT ─────────
    s = _blank(prs)
    draw_silent_page(
        s,
        "groupby 是思考模型，\nmerge 是通行證，\npivot 是視角。",
    )
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ───────── S2 · MOTIVATION · ASK ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "同樣問『各地區營收』，\ngroupby、pivot_table、SQL 三種寫法 ——\n背後的運算一樣嗎？你會怎麼選？",
        data_card={
            "label": "資料分析師日常 80% 提問",
            "stat": "3 刀",
            "caption": "groupby / merge / pivot\n組合起來能答多數商業問題",
        },
    )
    add_source(s, "S3 outline §5 Discussion Prompt Q1")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── S3 · CONCEPT-CARD · 條件篩選補強 ─────────
    s = _blank(prs)
    add_title(s, "條件篩選三把鑰匙：between / isin / 多條件括號")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="先把篩選寫乾淨，後面 groupby / merge 才不會一坨括號",
        code=(
            '# 1. between — 區間篩選（包前包後）\n'
            'df[df["amount"].between(100, 500)]\n'
            '\n'
            '# 2. isin — 類別多選（list / set / Series）\n'
            'df[df["region"].isin(["北區", "中區", "南區"])]\n'
            '\n'
            '# 3. 多條件 & / |（每個條件都要括號）\n'
            'df[(df["amount"] > 100) & (df["region"] == "北區")]\n'
            '\n'
            '# 4. query — 可讀性最高（可用 @var 帶外部變數）\n'
            'min_amt = 100\n'
            'df.query("amount > @min_amt and region == \\"北區\\" ")'
        ),
        bullets=[
            "between 省掉 >= 與 <= 兩個比較",
            "isin 是多個 == 的簡寫",
            "& | 優先級高於 > <\n沒括號一定出錯",
            "query 適合複雜條件\n也是 reviewer 最愛",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Indexing — Boolean indexing · §Indexing.query")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───────── S4 · PITFALL · 括號陷阱 ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="❌ 錯誤：少了括號 — TypeError",
        right_title="✅ 正確：每個比較都括起來",
        left_items=[
            "df[df.a > 0 & df.b < 10]",
            "& 優先級高於 > <",
            "先算 0 & df.b → 型別錯誤",
            "噴 TypeError 或結果錯亂",
            "新手最常踩的坑",
        ],
        right_items=[
            "df[(df.a > 0) & (df.b < 10)]",
            "每個比較都用括號包起來",
            "再用 & | 連接",
            "或改用 df.query() 一步到位",
            "Review 時一眼看得懂",
        ],
        title="多條件篩選：& 優先級陷阱 — 一個括號的差別",
        summary="記這句：pandas 的 & 是位元運算、優先級高於比較 — 每個條件必須括起來。",
    )
    add_source(s, "S3 outline §4 Pitfall #6 · pandas docs §Boolean indexing")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── S5 · MECHANISM-FLOW · Split-Apply-Combine ─────────
    s = _blank(prs)
    add_title(s, "groupby 心智模型：Split → Apply → Combine（三步驟秒解所有 groupby 題）")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.5), h=Inches(4.2),
        slot_name="Split-Apply-Combine 流程圖",
        description=(
            "完整 orders 表 → 依 region 切 3 組\n"
            "（北區 / 中區 / 南區）\n"
            "→ 每組對 amount 算 sum\n"
            "→ 合併回單一結果表（每 region 一列）"
        ),
        url_hint="https://pandas.pydata.org/docs/user_guide/groupby.html",
        placeholder_id="S3_S05_split_apply_combine",
        registry=image_registry,
        size_hint="1400×800 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.4), y=Inches(1.3),
        w=Inches(6.3), h=Inches(4.2),
        label="一句話 = split + apply + combine 三步驟",
        code=(
            '# 一句話完成三步驟\n'
            'df.groupby("region")["amount"].sum()\n'
            '\n'
            '# 拆開看：\n'
            '#   split   = 依 region 切成多組\n'
            '#   apply   = 每組對 amount 算 sum\n'
            '#   combine = 拼回單一結果（每 region 一列）\n'
            '\n'
            '# Wickham (2011) 的三步驟心智模型\n'
            '# 今天還是 pandas / dplyr / SQL 的核心'
        ),
        bullets=[
            "1956 年的概念，到今天還在用",
            "背概念、不背 API",
            "split = hash partition\napply = vectorised\ncombine = concat",
        ],
        label_dark=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.7),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "所有 groupby 題都能套這個三步驟 — 卡住時先問自己：我在 split 什麼？apply 什麼？",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "Wickham, The Split-Apply-Combine Strategy 2011 · pandas User Guide §Group by")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── S6 · CONCEPT-CARD · groupby Level 1 ─────────
    s = _blank(prs)
    add_title(s, "groupby 三層寫法 · Level 1：單欄 × 單函式（最常見）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="先熟練 Level 1 — 80% 日常報表就是這招",
        code=(
            '# 1. 單欄單聚合（最直覺）\n'
            'df.groupby("region")["amount"].sum()\n'
            'df.groupby("region")["amount"].mean()\n'
            'df.groupby("region")["amount"].count()\n'
            '\n'
            '# 2. 多鍵分組（產生 MultiIndex）\n'
            'df.groupby(["region", "category"])["amount"].sum()\n'
            '\n'
            '# 3. 排序看 Top N\n'
            '(df.groupby("region")["amount"].sum()\n'
            '   .sort_values(ascending=False)\n'
            '   .head(3))\n'
            '\n'
            '# 4. 攤平 MultiIndex（常用）\n'
            'result = df.groupby(["region", "category"])["amount"].sum()\n'
            'result.reset_index()     # 回到平面 DataFrame'
        ),
        bullets=[
            "先 groupby，再選欄、再聚合",
            "多鍵分組 → MultiIndex",
            "後續要 merge / 畫圖\n幾乎都要 reset_index",
            "sort_values + head\n是 Top N 的黃金組合",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Group by — Aggregation · McKinney 3e §10.2")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───────── S7 · CONCEPT-CARD · groupby Level 2 ─────────
    s = _blank(prs)
    add_title(s, "groupby 三層寫法 · Level 2：多函式聚合 .agg([...])")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="一次算多個統計 — 業務報表最愛",
        code=(
            '# 1. 單欄多函式（最常見進階用法）\n'
            'df.groupby("region")["amount"].agg(["count", "sum", "mean"])\n'
            '#          region  count     sum    mean\n'
            '#             北區    120  48,000   400\n'
            '\n'
            '# 2. 多欄不同函式（字典版）\n'
            'df.groupby("region").agg({\n'
            '    "amount":   "sum",\n'
            '    "qty":      "mean",\n'
            '    "order_id": "count",\n'
            '})\n'
            '\n'
            '# 3. 內建函式字串 vs numpy 函式 vs lambda\n'
            'df.groupby("region")["amount"].agg(\n'
            '    ["sum", "mean", lambda g: g.quantile(0.9)])'
        ),
        bullets=[
            "字串 'sum' 最快\n（走 pandas C 層）",
            "字典版：欄位 → 函式",
            "lambda 拿到整個 group\n（是 Series，不是單值）",
            "結果欄名可能很醜\n→ Level 3 解決",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Group by — Applying multiple functions")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── S8 · CONCEPT-CARD · groupby Level 3（具名聚合，推薦） ─────────
    s = _blank(prs)
    add_title(s, "groupby 三層寫法 · Level 3：具名聚合（推薦寫法，欄名乾淨）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="named aggregation — 結果欄名自己取，下游最好接",
        code=(
            '# 語法：new_col=(source_col, func)\n'
            'df.groupby("region").agg(\n'
            '    total=("amount", "sum"),\n'
            '    cnt=("order_id", "count"),\n'
            '    avg_qty=("qty", "mean"),\n'
            '    p90=("amount", lambda g: g.quantile(0.9)),\n'
            ')\n'
            '# region  total   cnt  avg_qty    p90\n'
            '# 北區    48000   120     2.5      780\n'
            '\n'
            '# 對比 Level 2：欄名變成 (amount, sum) 這種 tuple 很難接下游\n'
            '# Level 3 直接指定 total、cnt — 後續 merge / 畫圖超順\n'
            '\n'
            '# Bonus：搭配 reset_index() 就是一張報表\n'
            'report = (df.groupby("region")\n'
            '            .agg(total=("amount","sum"), cnt=("order_id","count"))\n'
            '            .reset_index())'
        ),
        bullets=[
            "Named aggregation =\n最推薦寫法",
            "結果是平面 DataFrame\n（attr 存取友善）",
            "欄名自己取 →\n下游 merge / 報表超順",
            "pandas 0.25+ 原生支援",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Named aggregation (since v0.25)")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── S9 · PITFALL · MultiIndex 忘 reset_index ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="❌ 忘了 reset_index — 後續 bug",
        right_title="✅ 加上 reset_index — 一路順暢",
        left_items=[
            "res = df.groupby(['region','cat']).sum()",
            "res.merge(other, on='region')",
            "→ KeyError: 'region'（它在 Index 不在欄位）",
            "畫 bar chart 也會出錯",
            "MultiIndex 的隱形坑",
        ],
        right_items=[
            "res = (df.groupby(['region','cat'])\n     .sum()\n     .reset_index())",
            "res.merge(other, on='region')  # OK",
            "也可用具名聚合 + reset_index 組合",
            "或 as_index=False 參數：",
            "df.groupby(..., as_index=False).sum()",
        ],
        title="groupby 後忘記 reset_index — MultiIndex 是新手第一大坑",
        summary="規則：groupby 結果只要要 merge / 畫圖 / 匯出 CSV，就先 reset_index()。",
    )
    add_source(s, "S3 outline §4 Pitfall #3 · pandas docs §Group by — as_index")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── S10 · EXAMPLE-IO · SQL 對照表 ─────────
    s = _blank(prs)
    draw_editorial_table(
        s,
        header=["SQL 子句", "Pandas 對應", "範例", "備註"],
        rows=[
            ["GROUP BY region",
             "df.groupby('region')",
             "df.groupby('region')",
             "分組鍵（可多欄）"],
            ["SELECT SUM(amount)",
             ".agg / .sum",
             "[...]['amount'].sum()",
             "選欄+聚合"],
            ["SELECT a, SUM(b), COUNT(*)",
             ".agg(具名聚合)",
             "agg(s=('b','sum'), c=('b','count'))",
             "最推薦寫法"],
            ["HAVING sum(amount) > 1000",
             "先 agg 再 query / 布林",
             ".query('total > 1000')",
             "Pandas 沒有 HAVING，\n事後篩"],
            ["ORDER BY total DESC",
             ".sort_values",
             ".sort_values('total', ascending=False)",
             "接 .head(N) = Top N"],
            ["JOIN on key",
             ".merge",
             "a.merge(b, on='key')",
             "S3 第二刀，下一段講"],
        ],
        col_widths=[2.6, 2.0, 3.4, 2.6],
        title="groupby = SQL 對照表：熟悉 SQL 就能把 Pandas 當 SQL 寫",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "Pandas 沒有 HAVING — 先 agg 完、再用 query / 布林篩結果。",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "pandas docs §Comparison with SQL · McKinney 3e §10.2")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── S11 · CHECKPOINT · 三種寫法 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "小測驗：把『各 region 總營收』\n用三種方式寫出來 ——\ngroupby().sum() / .agg() / pivot_table()\n三者的結果一樣嗎？有什麼差？",
        data_card={
            "label": "自我檢核",
            "stat": "3",
            "caption": "能寫三種 = 真的理解\n只會一種 = 還在背 API",
        },
    )
    add_source(s, "S3 outline §5 Discussion Prompt Q1")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── S12 · CONCEPT-CARD · merge 基礎 ─────────
    s = _blank(prs)
    add_title(s, "merge 基礎：on 與 how 兩個參數吃掉 90% 情境")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="merge = 依鍵合併兩表（對應 SQL JOIN）",
        code=(
            '# 1. 最基本：欄名相同 → on\n'
            'orders.merge(customers, on="customer_id")\n'
            '#        預設 how="inner"（取交集）\n'
            '\n'
            '# 2. 指定 how 四選一\n'
            'orders.merge(customers, on="customer_id", how="left")\n'
            '# 90% 實務用 how="left"：「從主表的角度去補資訊」\n'
            '\n'
            '# 3. debug 黃金公式 — merge 完一定做\n'
            'before = len(orders)\n'
            'merged = orders.merge(customers, on="customer_id", how="left")\n'
            'after  = len(merged)\n'
            'assert after == before, "鍵不唯一 → 筆數暴漲"\n'
            '\n'
            '# 4. indicator=True → 加 _merge 欄位幫 debug\n'
            'orders.merge(customers, on="customer_id",\n'
            '             how="outer", indicator=True)'
        ),
        bullets=[
            "預設 how='inner' —\n會默默丟掉沒匹配的",
            "90% 情境用 how='left'\n（保留主表全部）",
            "merge 完第一件事：\nlen() 比對前後",
            "indicator=True 是 debug 神器",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Merge, join, concatenate — merge")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── S13 · MECHANISM-FLOW · 四種 how ─────────
    s = _blank(prs)
    add_title(s, "merge 四種 how：inner / left / right / outer — 一張 Venn 圖說完")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.5), h=Inches(5.0),
        slot_name="merge how Venn 圖（四聯圖）",
        description=(
            "四個 Venn 圖水平排列：\n"
            "(1) inner — 交集塗色（左右都有）\n"
            "(2) left — 左圓全塗（保左全部，右補 NaN）\n"
            "(3) right — 右圓全塗\n"
            "(4) outer — 整體塗色（聯集，兩邊缺都 NaN）"
        ),
        url_hint="https://pandas.pydata.org/docs/user_guide/merging.html",
        placeholder_id="S3_S13_merge_how_venn",
        registry=image_registry,
        size_hint="1600×600 px",
    )
    draw_matrix(
        s, rows=2, cols=2,
        cells=[
            {"text": "inner（預設）",
             "sub": "只留兩表都有的 key\n會默默丟資料\n用前先確認"},
            {"text": "left（實務 90%）",
             "sub": "保留左表全部\n右表沒匹配 → NaN\n「主表補資訊」",
             "highlight": True},
            {"text": "right",
             "sub": "保留右表全部\n等價於對調 left / right\n實務少用"},
            {"text": "outer",
             "sub": "左右全部保留\n缺的位置 NaN\n搭配 indicator debug"},
        ],
        title="",
    )
    add_source(s, "S3 outline §3 Point 3 · pandas docs §Database-style joining")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── S14 · CONCEPT-CARD · left_on / right_on / suffixes ─────────
    s = _blank(prs)
    add_title(s, "欄名不一致、或兩表同名欄：left_on / right_on / suffixes")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="實務常見：orders.uid vs customers.customer_id",
        code=(
            '# 1. 欄名不一致：用 left_on / right_on\n'
            'orders.merge(\n'
            '    customers,\n'
            '    left_on="uid",           # orders 裡的欄\n'
            '    right_on="customer_id",  # customers 裡的欄\n'
            '    how="left",\n'
            ')\n'
            '# 結果會同時保留 uid 和 customer_id 兩欄（冗餘）\n'
            '# 建議事後 drop 一欄：.drop(columns="customer_id")\n'
            '\n'
            '# 2. 兩表有同名非鍵欄：用 suffixes 指定後綴\n'
            'orders.merge(\n'
            '    products,\n'
            '    on="product_id",\n'
            '    suffixes=("_order", "_prod"),   # 預設 ("_x","_y")\n'
            ')\n'
            '# price_order / price_prod — 比 _x / _y 清楚得多\n'
            '\n'
            '# 3. 也可以事前 rename 避開\n'
            'products2 = products.rename(columns={"price": "list_price"})\n'
            'orders.merge(products2, on="product_id")'
        ),
        bullets=[
            "欄名不一致 → left_on / right_on",
            "_x / _y 的來源 = suffixes 預設",
            "事前 rename 往往更清楚",
            "冗餘欄位記得 drop",
        ],
        label_dark=True,
    )
    add_source(s, "S3 outline §3 Point 4 · pandas docs §merge parameters")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── S15 · PITFALL · 筆數暴增 ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="❌ 沒檢查 key 唯一 — merge 後暴漲",
        right_title="✅ merge 前先驗 key + 事後比 len",
        left_items=[
            "orders 100 萬筆",
            "customers 1 萬筆，但 customer_id 重複了",
            "orders.merge(customers, on='customer_id')",
            "→ 150 萬筆（多出 50%！）",
            "報表數字全錯、沒人知道哪裡壞",
        ],
        right_items=[
            "# 前：驗 right 表 key 唯一",
            "dup = customers['customer_id'].duplicated().sum()",
            "assert dup == 0, f'key 重複 {dup} 筆'",
            "# 後：比 len 前後",
            "assert len(merged) == len(orders)  # left join 應相等",
            "# 或用 validate='m:1' 讓 pandas 自己檢查",
        ],
        title="merge 後筆數暴增 — 1:N 鍵不唯一（S3 第一名 bug）",
        summary="三行護身符：merge 前驗 duplicated、merge 時加 validate='m:1'、merge 後比 len。",
    )
    add_source(s, "S3 outline §4 Pitfall #2 · §5 Q3 · pandas docs §merge validate")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── S16 · CONCEPT-CARD · 三表 join 實戰 ─────────
    s = _blank(prs)
    add_title(s, "三表 join 實戰：orders × customers × products（鏈式 merge）")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="電商三表 → 寬表：後續 groupby / pivot 都吃這張",
        code=(
            '# 三表長相\n'
            '# orders    : order_id, customer_id, product_id, amount, qty, ts\n'
            '# customers : customer_id, name, region, tier(VIP/一般)\n'
            '# products  : product_id, category, price\n'
            '\n'
            '# 鏈式 merge — 從主表 orders 出發\n'
            'wide = (\n'
            '    orders\n'
            '      .merge(customers, on="customer_id", how="left",\n'
            '             validate="m:1")\n'
            '      .merge(products,  on="product_id",  how="left",\n'
            '             validate="m:1")\n'
            ')\n'
            '\n'
            '# 驗收（每次 merge 完都做）\n'
            'assert len(wide) == len(orders)\n'
            'print(wide.isna().sum())  # 看缺失：哪些 orders 沒對到客戶？\n'
            '\n'
            '# 這張 wide 就是後面 groupby / pivot 的料\n'
            'wide.groupby("region")["amount"].sum()\n'
            'wide.groupby(["region","category"])["amount"].sum()'
        ),
        bullets=[
            "永遠從「主表」出發\n鏈式 merge 下去",
            "validate='m:1'\n= 右表 key 必須唯一",
            "每次 merge 後\nassert len + isna",
            "結果是 wide 表\n後面 groupby 都吃它",
        ],
        label_dark=True,
    )
    add_source(s, "S3 outline §2 01:00-01:10 · LO5")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── S17 · CONCEPT-CARD · pivot_table ─────────
    s = _blank(prs)
    add_title(s, "pivot_table：列 × 欄的交叉分析矩陣（Excel 樞紐的 pandas 版）")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.5), h=Inches(4.8),
        slot_name="pivot_table 結構示意",
        description=(
            "左：long-format 原始表\n"
            "    (region, category, amount) × N rows\n"
            "右：wide-format pivot 結果矩陣\n"
            "    index=region 當列\n"
            "    columns=category 當欄\n"
            "    values=amount, aggfunc=sum 當儲存格\n"
            "    可加 margins=True 補總計列/欄"
        ),
        url_hint="https://pandas.pydata.org/docs/user_guide/reshaping.html",
        placeholder_id="S3_S17_pivot_table_schema",
        registry=image_registry,
        size_hint="1400×900 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.4), y=Inches(1.3),
        w=Inches(6.3), h=Inches(4.8),
        label="pivot_table = groupby(多鍵) 的寬表版",
        code=(
            '# 基本：列=region, 欄=category, 值=amount\n'
            'wide.pivot_table(\n'
            '    index="region",\n'
            '    columns="category",\n'
            '    values="amount",\n'
            '    aggfunc="sum",     # 一定要指定（預設 mean！）\n'
            '    fill_value=0,      # 缺的組合補 0\n'
            '    margins=True,      # 加總計列/欄\n'
            ')\n'
            '\n'
            '# 多個 aggfunc 一次算\n'
            'wide.pivot_table(\n'
            '    index="region", columns="category",\n'
            '    values="amount",\n'
            '    aggfunc=["sum", "mean", "count"])'
        ),
        bullets=[
            "想「行=A 列=B 值=m」\n就用 pivot_table",
            "aggfunc 預設 mean\n99% 情境要改成 sum",
            "margins=True → 加總行/列",
            "fill_value 避免一堆 NaN",
        ],
        label_dark=True,
    )
    add_source(s, "pandas User Guide §Reshaping and pivot tables")
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───────── S18 · PITFALL · aggfunc 預設 mean ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="❌ 用預設 aggfunc — 算的是平均不是總和",
        right_title="✅ 永遠明確指定 aggfunc",
        left_items=[
            "wide.pivot_table(index='region',",
            "    columns='category', values='amount')",
            "# aggfunc 省略 → 預設 'mean'",
            "→ 報表寫「總營收」但跑出平均值",
            "老闆問為什麼北區營收只有 400 元？",
        ],
        right_items=[
            "wide.pivot_table(index='region',",
            "    columns='category', values='amount',",
            "    aggfunc='sum',          # 明確指定",
            "    fill_value=0,           # 缺補 0",
            "    margins=True)           # 加總計",
        ],
        title="pivot_table 預設 aggfunc='mean' — 忘指定整張報表會錯",
        summary="pivot_table 每次寫都四件套：index / columns / values / aggfunc — 一個都不能省。",
    )
    add_source(s, "S3 outline §4 Pitfall #5 · pandas docs §pivot_table")
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ───────── S19 · EXAMPLE-IO · 電商三問 ─────────
    s = _blank(prs)
    draw_three_blocks_flow(
        s,
        blocks=[
            {"heading": "Q1: 各地區總營收",
             "items": [
                 "wide.groupby('region')",
                 "['amount'].sum()",
                 ".sort_values(",
                 "  ascending=False)",
                 "→ 一欄結果",
             ]},
            {"heading": "Q2: 各地區 Top3 商品",
             "items": [
                 "(wide.groupby(",
                 "  ['region','product_id'])",
                 "['amount'].sum()",
                 ".groupby('region')",
                 ".nlargest(3))",
             ]},
            {"heading": "Q3: VIP 貢獻佔比",
             "items": [
                 "tot = wide['amount'].sum()",
                 "vip = (wide.query(",
                 "  'tier==\"VIP\"')",
                 "['amount'].sum())",
                 "ratio = vip / tot",
             ]},
        ],
        title="電商三問實戰：同一張 wide 表回答三個商業問題",
        bottom_note=(
            "三題都走 Split→Apply→Combine — groupby 吃掉 Q1/Q2，query+sum 解 Q3。"
            "這就是 S3 要帶走的「組合拳」感。"
        ),
    )
    add_source(s, "S3 outline §2 01:00-01:10 · LO5")
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───────── S20 · CONCEPT-CARD · groupby vs pivot_table ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="groupby（long format）",
        right_title="pivot_table（wide format）",
        left_items=[
            "結果是 long 表（每組一列）",
            "適合：下游程式再吃、繼續 merge、畫圖",
            "語法：groupby(...).agg(...)",
            "多鍵 → MultiIndex（記得 reset_index）",
            "靈活度最高 — 可自訂函式、lambda",
        ],
        right_items=[
            "結果是 wide 矩陣（列 × 欄）",
            "適合：人看、貼 Excel、做儀表板",
            "語法：pivot_table(index, columns, values, aggfunc)",
            "margins=True 直接補總計",
            "本質 = groupby(多鍵).agg().unstack()",
        ],
        title="groupby vs pivot_table：同樣運算、兩種輸出形狀",
        summary=(
            "兩者可互換 — 「要給人看的交叉表」用 pivot_table；"
            "「要接下游程式」用 groupby。背後都是 Split-Apply-Combine。"
        ),
        delta="shape",
    )
    add_source(s, "S3 outline §3 Point 5 · §5 Q1")
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ───────── S21 · PRACTICE-PROMPT ─────────
    s = _blank(prs)
    draw_silent_page(
        s,
        "課堂練習（40 分鐘）\n\n🟡 品類營收 Top N\n🔴 RFM 粗估挑戰題\n\n先想清楚：split 什麼？apply 什麼？",
    )
    add_footer(s, MODULE_CODE, 21, N_CONTENT, dark_bg=True)

    # ───────── S22 · CHECKPOINT · 收束 ─────────
    s = _blank(prs)
    draw_thesis_hierarchy(
        s,
        blocks=[
            {"heading": "S3 三把刀（背下這三句就能上戰場）",
             "items": [
                 "groupby = 思考模型 Split→Apply→Combine（具名聚合最推薦）",
                 "merge = 多表通行證（on + how='left' 打 90%，記得 validate='m:1'）",
                 "pivot_table = 交叉視角（四件套：index / columns / values / aggfunc）",
             ]},
            {"heading": "四條紀律（可直接抄到個人 checklist）",
             "items": [
                 "多條件篩選：每個比較都括號；能用 query 就用 query",
                 "groupby 後要 merge / 畫圖，一律 reset_index()",
                 "merge 後第一件事：assert len() 與 isna().sum()",
                 "pivot_table 永遠明確寫 aggfunc，別靠預設 mean",
             ]},
        ],
        title="S3 收束：三把刀 + 四條紀律 — 下一節 S4 把這些結果畫出來",
        thesis="S4 進入時間序列與 EDA — 把 groupby / merge / pivot 的產物，變成看得見的洞察。",
    )
    add_source(s, "S3 module synthesis · 銜接 S4 時間序列與 EDA 實戰")
    add_footer(s, MODULE_CODE, 22, N_CONTENT)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
