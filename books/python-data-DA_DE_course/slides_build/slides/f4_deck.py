"""F4 deck — 封裝、繼承與魔術方法
20 content slides + cover + copyright page.

Governing thought:
    類別要守門、要分工、要像 Python 內建型別一樣自然——
    三件事學會就夠打 90% 的資料工程場景。

Aligned to chapters/F4_封裝繼承與魔術方法/00_skeleton.yaml
  · 4 Learning Objectives × 5 Common Pitfalls
  · Teaching-track: MOTIVATION / CONCEPT / MECHANISM / EXAMPLE / PITFALL / BRIDGE / CHECKPOINT / SUMMARY
  · 為 F5 OOP+Pandas 整合實戰的 DataCleaner 鋪路
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source, add_rect,
    set_solid_fill, set_no_fill, set_line, set_no_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_code_panel, draw_thesis_hierarchy, draw_image_placeholder,
    draw_flow_chain,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "F4"
MODULE_TITLE = "封裝、繼承與魔術方法"
MODULE_SUBTITLE = "類別要守門、要分工、要像 Python 內建型別一樣自然"
TIME_MIN = 90
N_CONTENT = 20

RED_ERROR = RGBColor(0xC6, 0x28, 0x28)
GREEN_OK = RGBColor(0x2E, 0x7D, 0x32)


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _draw_pitfall(slide, title, wrong_label, wrong_code, right_label,
                  right_code, why):
    """PITFALL 頁：左紅錯 / 右綠對 / 下方一句 why。"""
    add_title(slide, title)
    top = Inches(1.4)
    col_h = Inches(4.2)
    col_w = Inches(5.8)
    gap = Inches(0.3)
    total = col_w * 2 + gap
    left_x = (T.SLIDE_W - total) / 2
    right_x = left_x + col_w + gap

    for x, label, code, accent, mark in [
        (left_x, wrong_label, wrong_code, RED_ERROR, "✗"),
        (right_x, right_label, right_code, GREEN_OK, "✓"),
    ]:
        hdr = add_rect(slide, x, top, col_w, Inches(0.5))
        set_solid_fill(hdr, accent)
        set_no_line(hdr)
        add_textbox(
            slide, x + Inches(0.2), top, col_w - Inches(0.4), Inches(0.5),
            f"{mark}  {label}",
            font_size=T.FONT_BODY, color=T.WHITE, bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        body_y = top + Inches(0.5)
        body_h = col_h - Inches(0.5)
        box = add_rect(slide, x, body_y, col_w, body_h)
        set_solid_fill(box, RGBColor(0xF9, 0xF9, 0xF9))
        set_line(box, accent, 1.2)
        add_textbox(
            slide, x + Inches(0.25), body_y + Inches(0.2),
            col_w - Inches(0.5), body_h - Inches(0.4),
            code,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            family=T.FONT_MONO, line_spacing=1.45,
        )

    why_y = top + col_h + Inches(0.2)
    why_box = add_rect(slide, T.MARGIN_X, why_y,
                       T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.55))
    set_solid_fill(why_box, T.PRIMARY)
    set_no_line(why_box)
    add_textbox(
        slide, T.MARGIN_X + Inches(0.2), why_y,
        T.SLIDE_W - 2 * T.MARGIN_X - Inches(0.4), Inches(0.55),
        f"為什麼:{why}",
        font_size=T.FONT_BODY, color=T.WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def build_f4(output_path, image_registry=None):
    """Build F4 deck; 20 content slides."""
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── 1 · ASK — 為什麼資深工程師很少直接改欄位 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "為什麼資深工程師的 class，\n很少直接 obj.field = 改值？",
        data_card={
            "label": "code review 常見退件",
            "stat": "#1",
            "caption": "『直接改內部欄位』\n是第一名被打回原因",
        },
    )
    add_source(s, "Google Engineering Practices · Python Style Internal Review")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # ───────── 2 · SILENT — 立論 ─────────
    s = _blank(prs)
    draw_silent_page(s, "類別要守門，\n不是把變數塞在一起。")
    add_footer(s, MODULE_CODE, 2, N_CONTENT, dark_bg=True)

    # ───────── 3 · CONCEPT — 封裝三道門（使用時機） ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "x（公開欄位）",
         "sub": "時機:穩定、無驗證需求\n任何人都能讀寫\n代價:之後想改規則 → 改全世界"},
        {"text": "_x（弱封裝）",
         "sub": "時機:『內部用，但不強制』\n單底線 = 禮貌提醒\nlinter 警告、IDE 自動隱藏",
         "highlight": True},
        {"text": "__x（強封裝）",
         "sub": "時機:避免子類別意外覆蓋\n雙底線 → name mangling\n不是加密，是隔離"},
        {"text": "@property",
         "sub": "時機:需要驗證／計算／換實作\n對外看起來像欄位\n對內其實是方法",
         "highlight": True},
    ], title="封裝三道門:Python 沒有 private，一切靠使用時機判斷")
    add_source(s, "PEP 8 §Naming · Python Tutorial §9.6 Private Variables")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───────── 4 · EXAMPLE — _x / __x 實測 ─────────
    s = _blank(prs)
    add_title(s, "_x 與 __x 實測:name mangling 不是加密")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="單底線 vs 雙底線 — 兩種慣例的真實差別",
        code=(
            'class Account:\n'
            '    def __init__(self, balance):\n'
            '        self._hint    = balance   # 禮貌提醒\n'
            '        self.__secret = balance   # 真的改名\n'
            '\n'
            'a = Account(100)\n'
            'a._hint                 # OK:Python 不阻止\n'
            'a.__secret              # AttributeError\n'
            'a._Account__secret      # 仍可拿到，只是改名了'
        ),
        bullets=[
            "_x 只是約定，linter 會警告",
            "__x 被改寫為 _ClassName__x",
            "本意:避免子類別意外覆蓋",
            "不是加密，是把鑰匙藏在\n有名字的抽屜裡",
        ],
        label_dark=True,
    )
    add_source(s, "Python Tutorial §9.6 · Fluent Python Ch11")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── 5 · EXAMPLE — @property 真實用途 ─────────
    s = _blank(prs)
    add_title(s, "@property:對外是欄位，對內是方法")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="Temperature:Celsius 可寫、Fahrenheit 唯讀",
        code=(
            'class Temperature:\n'
            '    def __init__(self, c):\n'
            '        self._celsius = c\n'
            '\n'
            '    @property\n'
            '    def celsius(self):\n'
            '        return self._celsius\n'
            '\n'
            '    @celsius.setter\n'
            '    def celsius(self, value):\n'
            '        if value < -273.15:\n'
            '            raise ValueError("below absolute zero")\n'
            '        self._celsius = value\n'
            '\n'
            '    @property\n'
            '    def fahrenheit(self):\n'
            '        return self._celsius * 9 / 5 + 32'
        ),
        bullets=[
            "對外:t.celsius = 25、t.fahrenheit",
            "對內:可驗證、可計算、可換實作",
            "寫法像欄位，實際是方法",
            "未來換實作 → 呼叫端零改動",
        ],
        label_dark=True,
    )
    add_source(s, "Python Built-in Types §property · PEP 252")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── 6 · PITFALL (P1) — __x 不是加密 ─────────
    s = _blank(prs)
    _draw_pitfall(
        s,
        title="陷阱 · __x 不是加密，不要拿來裝敏感資料",
        wrong_label="以為 __x 是 private",
        wrong_code=(
            'class User:\n'
            '    def __init__(self, pw):\n'
            '        self.__password = pw\n'
            '\n'
            'u = User("s3cret")\n'
            'u.__password\n'
            '# AttributeError\n'
            '# 以為安全了？\n'
            'u._User__password\n'
            '# → "s3cret"  外人照樣拿到'
        ),
        right_label="敏感資料該走驗證 / 加密 / 不存",
        right_code=(
            'import hashlib\n'
            'class User:\n'
            '    def __init__(self, pw):\n'
            '        self._pw_hash = hashlib.sha256(\n'
            '            pw.encode()).hexdigest()\n'
            '\n'
            '    def verify(self, pw):\n'
            '        h = hashlib.sha256(\n'
            '            pw.encode()).hexdigest()\n'
            '        return h == self._pw_hash'
        ),
        why="__x 的目的是避免子類命名衝突，不是安全機制；敏感資料要 hash / 加密 / 外存",
    )
    add_source(s, "Python Data Model §3.1 · OWASP ASVS 2.4")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───────── 7 · ASK — 繼承觸發場景 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "要支援 CSV / JSON / Parquet 三種來源，\n"
        "你要寫三個獨立類別，還是一個家族？",
        data_card={
            "label": "重複邏輯比例",
            "stat": "70%",
            "caption": "讀檔 → 檢查 → 回 DataFrame\n70% 步驟完全一樣",
        },
    )
    add_source(s, "課堂實測:10 位學員獨立寫 3 個 Reader")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── 8 · CONCEPT — 繼承樹 ─────────
    s = _blank(prs)
    add_title(s, "DataReader 家族:共同點上移，差異點下放")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.3), h=Inches(5.0),
        slot_name="DataReader 繼承樹",
        description=(
            "根節點 DataReader（read 介面）\n"
            "向下分三支:CSVReader / JSONReader / ParquetReader\n"
            "連線下方標註 override read()"
        ),
        url_hint="",
        placeholder_id="F4_S08_inheritance_tree",
        registry=image_registry,
        size_hint="1280×900 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.2), y=Inches(1.3),
        w=Inches(6.5), h=Inches(5.0),
        label="基類定介面、子類實作差異",
        code=(
            'class DataReader:\n'
            '    def read(self, path):\n'
            '        raise NotImplementedError\n'
            '\n'
            'class CSVReader(DataReader):\n'
            '    def read(self, path):\n'
            '        return pd.read_csv(path)\n'
            '\n'
            'class JSONReader(DataReader):\n'
            '    def read(self, path):\n'
            '        return pd.read_json(path)'
        ),
        bullets=[
            "基類 = 合約（interface）",
            "子類 = 差異化實作",
            "共用邏輯（log / 檢查）\n寫在基類一次就好",
        ],
        label_dark=True,
    )
    add_source(s, "Gamma et al., Design Patterns §Template Method")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── 9 · EXAMPLE — super() 與 override ─────────
    s = _blank(prs)
    add_title(s, "super().__init__ 與方法覆寫:先做父親那份，再做自己的")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="CSVReader 完整繼承骨架",
        code=(
            'class DataReader:\n'
            '    def __init__(self, encoding="utf-8"):\n'
            '        self.encoding = encoding\n'
            '        self.rows_read = 0\n'
            '\n'
            '    def read(self, path):\n'
            '        raise NotImplementedError\n'
            '\n'
            'class CSVReader(DataReader):\n'
            '    def __init__(self, encoding="utf-8", sep=","):\n'
            '        super().__init__(encoding)   # 父類先設好\n'
            '        self.sep = sep\n'
            '\n'
            '    def read(self, path):            # override\n'
            '        df = pd.read_csv(path, sep=self.sep,\n'
            '                         encoding=self.encoding)\n'
            '        self.rows_read = len(df)\n'
            '        return df'
        ),
        bullets=[
            "super().__init__() 必做，\n否則父類屬性沒設",
            "override = 同名重寫，\n簽章可延伸",
            "子類可用 self.encoding\n直接存父類屬性",
            "漏 super → bug 兩個月後才爆",
        ],
        label_dark=True,
    )
    add_source(s, "Python Tutorial §9.5 Inheritance · PEP 3135")
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── 10 · PITFALL (P2) — 漏呼 super ─────────
    s = _blank(prs)
    _draw_pitfall(
        s,
        title="陷阱 · 子類 __init__ 忘記呼叫 super().__init__()",
        wrong_label="漏呼 super()",
        wrong_code=(
            'class CSVReader(DataReader):\n'
            '    def __init__(self, sep=","):\n'
            '        self.sep = sep  # 只設自己的\n'
            '\n'
            'r = CSVReader()\n'
            'r.read("a.csv")\n'
            '# AttributeError:\n'
            "# 'CSVReader' object has\n"
            "# no attribute 'encoding'"
        ),
        right_label="先 super()，再設自己的",
        right_code=(
            'class CSVReader(DataReader):\n'
            '    def __init__(self, encoding="utf-8",\n'
            '                 sep=","):\n'
            '        super().__init__(encoding)\n'
            '        self.sep = sep\n'
            '\n'
            'r = CSVReader()\n'
            'r.encoding   # "utf-8"  ✓\n'
            'r.sep        # ","      ✓'
        ),
        why="父類 __init__ 負責設定共用狀態；漏呼 = 父類屬性全部沒初始化，錯誤會延後兩個月才爆",
    )
    add_source(s, "Python Data Model §3.3.3 object.__init__")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── 11 · CONCEPT — IS-A vs HAS-A ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "繼承 · IS-A ✓",
         "sub": "CSVReader is a DataReader\n共用介面、共用狀態\n這是繼承的正當場景",
         "highlight": True},
        {"text": "繼承 · HAS-A ✗",
         "sub": "Order 不是 List\n即使它『裝了』items\n別讓 OrderList 誕生"},
        {"text": "組合 · HAS-A ✓",
         "sub": "Cleaner has a Logger\n傳進去、換掉、mock 都容易\n90% 的『擴充』都該用這個",
         "highlight": True},
        {"text": "本章不教",
         "sub": "多重繼承 / MRO / ABC\n單一繼承已覆蓋\n資料工程 90% 需求"},
    ], title="先問關係是 IS-A 還是 HAS-A:別用繼承解決組合問題")
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "實務鐵則:能用組合不用繼承。F5 DataCleaner 會再驗證一次。",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "GoF Design Patterns §Favor composition over inheritance")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── 12 · CONCEPT — 魔術方法四件套 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "__str__",
         "sub": "print(obj) / str(obj) 時呼叫\n給『人』看的字串\n最常用也最該先實作",
         "highlight": True},
        {"text": "__repr__",
         "sub": "REPL 顯示 / repr(obj)\n給『開發者』看\n最好能反貼 Python 重跑"},
        {"text": "__len__",
         "sub": "len(obj) 時呼叫\n有『數量感』才實作\n通常和 __iter__ 成對",
         "highlight": True},
        {"text": "__iter__",
         "sub": "for x in obj 時呼叫\n讓物件可被迭代\nF5 DataCleaner 會用上"},
    ], title="魔術方法四件套:Python 幫你自動呼叫的那些 dunder")
    add_source(s, "Python Data Model §3.3 Special method names")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── 13 · EXAMPLE — __str__ vs __repr__ ─────────
    s = _blank(prs)
    add_title(s, "__str__ vs __repr__:給人看 vs 給開發者看")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="Order 類:兩種字串表示同時定義",
        code=(
            'class Order:\n'
            '    def __init__(self, id_, amount):\n'
            '        self.id, self.amount = id_, amount\n'
            '\n'
            '    def __repr__(self):\n'
            '        return f"Order(id_={self.id!r}, amount={self.amount!r})"\n'
            '\n'
            '    def __str__(self):\n'
            '        return f"訂單 #{self.id}:NT$ {self.amount:,}"\n'
            '\n'
            'o = Order("A01", 12800)\n'
            'print(o)        # 訂單 #A01:NT$ 12,800\n'
            'repr(o)         # "Order(id_=\'A01\', amount=12800)"'
        ),
        bullets=[
            "__str__ 給使用者 / log",
            "__repr__ 給開發者 debug",
            "repr 理想上能反貼重跑",
            "沒定義 __str__ →\nprint 退而求其次用 __repr__",
        ],
        label_dark=True,
    )
    add_source(s, "Fluent Python §1 · Guido on str vs repr")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── 14 · EXAMPLE — __len__ / __iter__ ─────────
    s = _blank(prs)
    add_title(s, "__len__ 與 __iter__:讓類別用 len() 和 for 都自然")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="DataPipeline:同時支援 len() / for / chaining",
        code=(
            'class DataPipeline:\n'
            '    def __init__(self):\n'
            '        self._steps = []\n'
            '\n'
            '    def add(self, step):\n'
            '        self._steps.append(step)\n'
            '        return self               # 預告 Chaining\n'
            '\n'
            '    def __len__(self):\n'
            '        return len(self._steps)\n'
            '\n'
            '    def __iter__(self):\n'
            '        return iter(self._steps)\n'
            '\n'
            'pipe = DataPipeline()\n'
            'pipe.add(clean).add(normalize).add(export)\n'
            'len(pipe)           # 3\n'
            'for s in pipe: ...  # 自動逐步跑'
        ),
        bullets=[
            "有『數量感』→ 實作 __len__",
            "有『逐一處理』→ 實作 __iter__",
            "兩者通常成對出現",
            "集合概念的必備配件",
        ],
        label_dark=True,
    )
    add_source(s, "Python Data Model §3.3.6 Emulating container types")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── 15 · CONCEPT — Method Chaining ─────────
    s = _blank(prs)
    add_title(s, "Method Chaining:return self 讓資料流讀起來像一句話")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.8), h=Inches(4.8),
        slot_name="Method Chaining 流程圖",
        description=(
            "水平四節點:read → drop_na → normalize → export\n"
            "每節點上方標 return self\n"
            "下方虛線註:Chaining = 單一物件貫穿全流程"
        ),
        url_hint="",
        placeholder_id="F4_S15_chaining_flow",
        registry=image_registry,
        size_hint="1400×600 px",
    )
    draw_code_panel(
        s,
        x=Inches(6.7), y=Inches(1.3),
        w=Inches(6.0), h=Inches(4.8),
        label="DataCleaner 雛形:每個動作 return self",
        code=(
            'class DataCleaner:\n'
            '    def __init__(self, df):\n'
            '        self.df = df\n'
            '\n'
            '    def drop_na(self):\n'
            '        self.df = self.df.dropna()\n'
            '        return self\n'
            '\n'
            '    def normalize(self, cols):\n'
            '        z = (self.df[cols] - self.df[cols].mean()) \\\n'
            '            / self.df[cols].std()\n'
            '        self.df[cols] = z\n'
            '        return self\n'
            '\n'
            '    def export(self, path):\n'
            '        self.df.to_csv(path, index=False)\n'
            '        return self\n'
            '\n'
            '(DataCleaner(df)\n'
            ' .drop_na().normalize(["price"]).export("out.csv"))'
        ),
        bullets=[
            "每個方法 return self",
            "讀起來像句子:\n資料 → 去空 → 標準化 → 匯出",
            "程式結構 = 資料流結構",
            "這就是 F5 DataCleaner 的骨架",
        ],
        label_dark=True,
    )
    add_source(s, "Martin Fowler, DSL §Method Chaining")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── 16 · PITFALL (P5) — Chaining 忘記 return self ─────────
    s = _blank(prs)
    _draw_pitfall(
        s,
        title="陷阱 · Chaining 方法忘記 return self",
        wrong_label="忘寫 return self",
        wrong_code=(
            'class DataCleaner:\n'
            '    def drop_na(self):\n'
            '        self.df = self.df.dropna()\n'
            '        # 忘了 return self\n'
            '\n'
            '(DataCleaner(df)\n'
            ' .drop_na()       # 回傳 None\n'
            ' .normalize(...)  # AttributeError:\n'
            '                  # NoneType has no\n'
            '                  # attribute \'normalize\''
        ),
        right_label="每一個方法都 return self",
        right_code=(
            'class DataCleaner:\n'
            '    def drop_na(self):\n'
            '        self.df = self.df.dropna()\n'
            '        return self         # ← 必要\n'
            '\n'
            '    def normalize(self, cols):\n'
            '        ...\n'
            '        return self\n'
            '\n'
            '(DataCleaner(df)\n'
            ' .drop_na().normalize(["price"]))'
        ),
        why="Python 函式沒 return 就回 None；Chaining 鏈上任何一環回 None，下一步就 AttributeError",
    )
    add_source(s, "PEP 8 §Programming Recommendations · Fluent API")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── 17 · EXAMPLE — 四柱合體 ─────────
    s = _blank(prs)
    add_title(s, "F4 合體預告:封裝 + 繼承 + 魔術方法 + Chaining")
    draw_code_panel(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="BaseCleaner / PriceCleaner — F5 DataCleaner 的前身",
        code=(
            'class BaseCleaner:\n'
            '    def __init__(self, df):\n'
            '        self._df = df              # 封裝:底線提示\n'
            '\n'
            '    @property\n'
            '    def shape(self):               # property:暴露計算屬性\n'
            '        return self._df.shape\n'
            '\n'
            '    def __len__(self):             # 魔術方法:len() 回列數\n'
            '        return len(self._df)\n'
            '\n'
            '    def drop_na(self):\n'
            '        self._df = self._df.dropna()\n'
            '        return self                # Chaining\n'
            '\n'
            'class PriceCleaner(BaseCleaner):   # 繼承 + override\n'
            '    def drop_na(self):\n'
            '        super().drop_na()\n'
            '        self._df = self._df[self._df["price"] > 0]\n'
            '        return self'
        ),
        bullets=[
            "封裝:_df + @property.shape",
            "繼承:PriceCleaner 擴充基類",
            "魔術方法:len(cleaner) 直接可用",
            "Chaining:每個動作 return self",
            "四柱合體 = F5 起跑線",
        ],
        label_dark=True,
    )
    add_source(s, "F4 → F5 銜接草案")
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───────── 18 · BRIDGE — 銜接 F5 DataCleaner ─────────
    s = _blank(prs)
    add_title(s, "銜接 F5:四柱如何組成 DataCleaner")
    draw_flow_chain(
        s,
        nodes=[
            {"label": "封裝",
             "sub": "_df / @property\n外部只讀、內部可換"},
            {"label": "繼承",
             "sub": "BaseCleaner →\nPriceCleaner / LogCleaner"},
            {"label": "魔術方法",
             "sub": "__len__ / __iter__\nlen() / for in 直接可用"},
            {"label": "Chaining",
             "sub": "return self\n讀→清→正規化→匯出",
             "highlight": True},
        ],
        title="",
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.2),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(1.0),
        "F5 會把這四柱接成一條完整清理管線——\n"
        "今天你手上的每個零件，明天都會再裝回去一次。",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        align=PP_ALIGN.CENTER, line_spacing=1.4,
    )
    add_source(s, "F4 → F5 DataCleaner 銜接地圖")
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ───────── 19 · CHECKPOINT — 五題快問 ─────────
    s = _blank(prs)
    add_title(s, "Check Point · 五題快問（90 秒）")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(5), Inches(0.4),
        "進度 19 / 20 · 不看講義作答",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(1.9), Inches(10.8), Inches(4.8),
        "Q1  __x 真的是『私有』嗎？a._Class__x 拿得到嗎？\n\n"
        "Q2  子類 __init__ 忘記呼叫 super().__init__()，會發生什麼事？\n\n"
        "Q3  只寫 __str__ 不寫 __repr__，在 REPL 裡打物件名會看到什麼？\n\n"
        "Q4  Chaining 方法忘記 return self，下一次呼叫會拿到什麼？\n\n"
        "Q5  『訂單裝了三個商品』——該用繼承還是組合？為什麼？",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.45,
    )
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───────── 20 · PYRAMID — 收束 ─────────
    s = _blank(prs)
    draw_thesis_hierarchy(
        s,
        blocks=[
            {"heading": "三根柱子",
             "items": [
                 "封裝:_x / __x / @property — 誰能改、怎麼改",
                 "繼承:super + override — 共用上移、差異下放",
                 "魔術方法:__str__ / __repr__ / __len__ / __iter__",
             ]},
            {"heading": "設計紀律",
             "items": [
                 "能用組合（HAS-A）不用繼承（IS-A）",
                 "單一繼承已夠，多重繼承是最後手段",
                 "方法鏈要一路 return self 才不破壞",
             ]},
        ],
        title="F4 收束:守門 + 分工 + 自然 + Chaining = 可長可久的類別",
        thesis="F5 進入 OOP + Pandas 整合實戰 —— 用本章四柱打造 DataCleaner。",
    )
    add_source(s, "F4 module synthesis")
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # Copyright slide
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
