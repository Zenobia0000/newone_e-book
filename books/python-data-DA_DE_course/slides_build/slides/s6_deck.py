"""S6 deck — Plotly 互動 + 完整 Capstone（12 小時 DA/DE 課程最終收斂章）.

24 content slides + cover + copyright page.

Governing thought:
    12 小時的終點不是「你會用 Plotly」，
    是你能獨立把一份髒 CSV 變成老闆願意轉寄的 dashboard。
    這才是 DA 的商業價值。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_rect, add_textbox, add_title, add_source,
    set_solid_fill, set_no_line, set_no_fill, set_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_vs_two_col, draw_code_panel, draw_flow_chain,
    draw_inverted_thesis_box, draw_image_placeholder,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "S6"
MODULE_TITLE = "Plotly 互動 + 完整 Capstone"
MODULE_SUBTITLE = "把髒 CSV 變成老闆願意轉寄的 dashboard"
TIME_MIN = 120
N_CONTENT = 24


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_s6(output_path, image_registry=None):
    """Build S6 deck — course-closing Plotly + Capstone chapter."""
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── S01 · SILENT ─────────
    s = _blank(prs)
    draw_silent_page(
        s,
        "12 小時的終點，不是「你會用 Plotly」。\n"
        "是你能獨立把一份髒 CSV，變成老闆願意轉寄的 dashboard。",
    )
    add_footer(s, MODULE_CODE, 1, N_CONTENT, dark_bg=True)

    # ───────── S02 · ASK ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "如果現在要交給非技術主管一份「能自己點」的分析報告，\n"
        "你會交一個 .py 檔、一個 notebook、還是一份 dashboard.html？",
        data_card={
            "label": "DA 交付現場觀察",
            "stat": "80%",
            "caption": "主管最後轉寄給其他人的，\n是 HTML，不是 notebook",
        },
    )
    add_source(s, "業界 DA 交付模式彙整 · 以收件者視角倒推")
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── S03 · MATRIX 2x2 靜態 vs 互動 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "靜態圖（S5 · matplotlib）",
         "sub": "輸出 .png / .pdf\n一次定格、無 hover\n放投影片、報告最好用",
         "highlight": False},
        {"text": "互動圖（S6 · Plotly）",
         "sub": "輸出 .html\nhover / zoom / 下拉篩選\n郵件轉寄、內網公告最好用",
         "highlight": True},
        {"text": "檔案體積",
         "sub": "png 約 100 KB\nhtml 含 plotly.js 約 3 MB\n（用 CDN 可降到 50 KB）"},
        {"text": "場景選擇",
         "sub": "投影片 → 靜態\n自助分析 → 互動\n兩者並存，不是互斥"},
    ], title="靜態圖 vs 互動圖：四格場景對比")
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "不是誰取代誰，是場景選擇 — DA 的工具箱裡兩種都要有。",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID, align=PP_ALIGN.CENTER,
    )
    add_source(s, "Plotly Docs §Interactive vs static · S5 視覺化章回顧")
    add_footer(s, MODULE_CODE, 3, N_CONTENT)

    # ───────── S04 · VS · .png vs .html ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="matplotlib → .png",
        right_title="plotly → .html",
        left_items=[
            "fig.savefig('out.png', dpi=150)",
            "寄出：收件者看的是一張圖",
            "無法互動、無法探索細節",
            "投影片、報告正文的最佳選擇",
        ],
        right_items=[
            "fig.write_html('dashboard.html')",
            "寄出：收件者用瀏覽器點開即用",
            "hover / zoom / 下拉篩選全都在",
            "對方不需要裝 Python，這才是交付",
        ],
        title="交付差異：.png 是一張圖，.html 是一個產品",
        summary="本章焦點：讓你能交出第二種。",
    )
    add_source(s, "Plotly Docs §write_html · matplotlib §savefig")
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── S05 · MATRIX 2x2 express vs graph_objects ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=2, cells=[
        {"text": "plotly.express（px）",
         "sub": "高階一行 API\npx.line / px.bar / px.scatter ...\n快速原型、EDA 最佳",
         "highlight": True},
        {"text": "plotly.graph_objects（go）",
         "sub": "底層 trace 物件\ngo.Scatter / go.Bar / go.Pie\n可完全客製每一筆",
         "highlight": False},
        {"text": "subplot 支援",
         "sub": "px 無法直接組 subplot\nmake_subplots 只吃 go.*\n這是本章最常問的題"},
        {"text": "取捨策略",
         "sub": "能 px 就 px，\n要組儀表板時轉 go。\n90% EDA 用 px 就夠"},
    ], title="Plotly Express vs graph_objects：何時用哪個")
    add_source(s, "Plotly Docs §plotly.express §graph_objects")
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── S06 · CODE · px.line + px.bar ─────────
    s = _blank(prs)
    add_title(s, "px.line + px.bar：時序與類別（兩招打七成）")
    draw_code_panel(
        s, x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="一行畫一張 — Plotly Express 的承諾",
        code=(
            'import plotly.express as px\n'
            '\n'
            '# 時間序列 → line\n'
            'fig1 = px.line(\n'
            '    monthly_df, x="month", y="revenue",\n'
            '    title="月營收趨勢",\n'
            '    labels={"month": "月份", "revenue": "營收 (NT$)"},\n'
            ')\n'
            '\n'
            '# 類別比較 → bar\n'
            'fig2 = px.bar(\n'
            '    category_df, x="category", y="sales",\n'
            '    color="category", text_auto=True,\n'
            '    title="各品類銷售",\n'
            ')'
        ),
        bullets=[
            "時間選 line、類別選 bar —\n90% 的場景夠用",
            "labels 參數讓圖「自己會說中文」",
            "text_auto=True 直接把數字標上去",
            "色彩用 color 分群，\n不要手動指定 16 進位",
        ],
        label_dark=True,
    )
    add_source(s, "Plotly Express 官方教學 §Line / §Bar")
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───────── S07 · CODE · px.scatter + px.box ─────────
    s = _blank(prs)
    add_title(s, "px.scatter + px.box：關聯與分群分布")
    draw_code_panel(
        s, x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="找關聯、看分布 — DA 第二常用的兩招",
        code=(
            '# 雙變數關聯 → scatter + 趨勢線\n'
            'fig3 = px.scatter(\n'
            '    orders, x="qty", y="revenue",\n'
            '    color="category", size="profit",\n'
            '    trendline="ols",\n'
            '    hover_data=["order_id", "customer"],\n'
            ')\n'
            '\n'
            '# 分群分布 → box\n'
            'fig4 = px.box(\n'
            '    orders, x="category", y="unit_price",\n'
            '    color="region", points="outliers",\n'
            ')'
        ),
        bullets=[
            "trendline='ols' 一行畫出迴歸線",
            "size 讓第三維度用面積表達",
            "hover_data 只挑 2~3 欄商業關鍵，\n塞滿十欄會變資訊牆",
            "box 的 points='outliers' 只顯示離群點",
        ],
        label_dark=True,
    )
    add_source(s, "Plotly Express §Scatter §Box · statsmodels trendline")
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── S08 · CODE · px.histogram + px.pie ─────────
    s = _blank(prs)
    add_title(s, "px.histogram + px.pie：次數與占比")
    draw_code_panel(
        s, x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="單變數看分布、占比看 pie — 但超過 5 類就改 bar",
        code=(
            '# 次數分布 → histogram\n'
            'fig5 = px.histogram(\n'
            '    orders, x="unit_price",\n'
            '    nbins=30, color="region",\n'
            '    marginal="box",    # 頂部加 box\n'
            ')\n'
            '\n'
            '# 占比 → pie（或 donut）\n'
            'fig6 = px.pie(\n'
            '    category_df, names="category", values="revenue",\n'
            '    hole=0.45,           # 變成 donut 更現代\n'
            '    title="品類營收占比",\n'
            ')'
        ),
        bullets=[
            "histogram 的 marginal\n一行加上頂部 box / violin",
            "pie 的 hole 參數 → donut",
            "pie 超過 5 類就該改 bar —\n不是顯示得下，是讀不出差異",
            "pie 的 type 是 domain —\n下一頁就要踩到這個雷",
        ],
        label_dark=True,
    )
    add_source(s, "Plotly Express §Histogram §Pie · 可視化設計原則")
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── S09 · SILENT ─────────
    s = _blank(prs)
    draw_silent_page(
        s,
        "一行畫一張，是 Plotly Express 的承諾。\n"
        "複雜的組裝，才是 make_subplots 的舞台。",
    )
    add_footer(s, MODULE_CODE, 9, N_CONTENT, dark_bg=True)

    # ───────── S10 · PITFALL · domain vs xy ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="踩雷寫法（ValueError）",
        right_title="正確寫法（domain + xy 明標）",
        left_items=[
            "specs=[[{}, {}], [{}, {}]]",
            "fig.add_trace(go.Pie(...), row=2, col=2)",
            "→ ValueError: Trace type 'pie'",
            "   is not compatible with subplot type 'xy'",
        ],
        right_items=[
            "specs=[[{'type':'xy'},    {'type':'xy'}],",
            "       [{'type':'xy'},    {'type':'domain'}]]",
            "→ pie / sunburst 必須在 domain 格",
            "→ 其他圖用 xy（預設）",
        ],
        title="踩雷一：subplot 混 domain / xy，specs 寫錯",
        summary="這是本章最常出錯的點 — pie 家族是 domain 型，記起來。",
    )
    add_source(s, "Plotly Python §subplots §mixed subplot types")
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── S11 · CODE · make_subplots 2x2 骨架 ─────────
    s = _blank(prs)
    add_title(s, "make_subplots 2x2：儀表板骨架（熟到背得出來）")
    draw_code_panel(
        s, x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="三步驟：建骨架 → add_trace → update_layout",
        code=(
            'from plotly.subplots import make_subplots\n'
            'import plotly.graph_objects as go\n'
            '\n'
            'fig = make_subplots(\n'
            '    rows=2, cols=2,\n'
            '    subplot_titles=["月營收", "品類", "散點", "占比"],\n'
            '    specs=[[{"type": "xy"},     {"type": "xy"}],\n'
            '           [{"type": "xy"},     {"type": "domain"}]],\n'
            ')\n'
            'fig.add_trace(go.Scatter(x=m.month, y=m.revenue), row=1, col=1)\n'
            'fig.add_trace(go.Bar(x=c.category, y=c.sales),    row=1, col=2)\n'
            'fig.add_trace(go.Scatter(x=o.qty, y=o.revenue,\n'
            '                         mode="markers"),          row=2, col=1)\n'
            'fig.add_trace(go.Pie(labels=c.category,\n'
            '                     values=c.revenue, hole=.45),  row=2, col=2)\n'
            'fig.update_layout(height=800, title="Orders Dashboard")'
        ),
        bullets=[
            "subplot_titles 要簡短，\n中文 4~6 字最好讀",
            "row / col 從 1 開始（不是 0）",
            "不寫 row/col → 全部疊在 (1,1)",
            "最後 update_layout\n一次把標題 + 高度設好",
        ],
        label_dark=True,
    )
    add_source(s, "Plotly Python §make_subplots · graph_objects trace 類型")
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── S12 · PITFALL · px 不是 trace ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="踩雷寫法（TypeError）",
        right_title="正確寫法（兩種救場）",
        left_items=[
            "fig.add_trace(px.bar(df, ...), row=1, col=1)",
            "→ px.bar 回傳 Figure，不是 trace",
            "→ add_trace 只吃 go.* 或 dict",
            "新手 90% 第一次卡在這裡",
        ],
        right_items=[
            "方法 A：直接用 go.*",
            "fig.add_trace(go.Bar(x=..., y=...))",
            "方法 B：取出 px 的第一筆 trace",
            "fig.add_trace(px.bar(df, ...).data[0])",
        ],
        title="踩雷二：px 回傳 Figure、不是 trace",
        summary="記一句話：Figure 是成品，trace 是零件。subplot 要的是零件。",
    )
    add_source(s, "Plotly Python §Figure vs Trace · community FAQ")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── S13 · MECHANISM-FLOW · Capstone 五站 ─────────
    s = _blank(prs)
    draw_flow_chain(
        s,
        title="Capstone 全流程：一份髒 CSV → dashboard.html 五站",
        y=3.2,
        nodes=[
            {"label": "① 讀檔",
             "sub": "pd.read_csv",
             "caption": "orders_raw.csv"},
            {"label": "② 清洗",
             "sub": "to_datetime\n(errors='coerce')",
             "caption": "load_and_clean_orders",
             "highlight": True},
            {"label": "③ 合併 + 聚合",
             "sub": "merge + groupby",
             "caption": "monthly_df / category_df"},
            {"label": "④ 組裝",
             "sub": "make_subplots\n+ go.*",
             "caption": "2x2 dashboard",
             "highlight": True},
            {"label": "⑤ 交付",
             "sub": "write_html",
             "caption": "dashboard.html"},
        ],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.2),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.4),
        "不要 import 舊 notebook — 從頭寫一次，才驗證你真的會。",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "Capstone pipeline 設計 · S1-S5 能力整合總覽")
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── S14 · CODE · load_and_clean_orders ─────────
    s = _blank(prs)
    add_title(s, "Step 1-2：load_and_clean_orders（coerce 是靈魂）")
    draw_code_panel(
        s, x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="髒 CSV 的第一道防線 — 別讓爛日期噴錯炸掉整條管線",
        code=(
            'import pandas as pd\n'
            'from pathlib import Path\n'
            '\n'
            'def load_and_clean_orders(path: str) -> pd.DataFrame:\n'
            '    df = pd.read_csv(path)\n'
            '\n'
            '    # 壞日期 → NaT，而不是 raise\n'
            '    df["order_date"] = pd.to_datetime(\n'
            '        df["order_date"], errors="coerce",\n'
            '    )\n'
            '\n'
            '    # 必備欄位缺就丟\n'
            '    df = df.dropna(subset=["order_date", "revenue", "qty"])\n'
            '\n'
            '    # 型別校正\n'
            '    df["revenue"] = df["revenue"].astype(float)\n'
            '    df["qty"]     = df["qty"].astype(int)\n'
            '    return df.reset_index(drop=True)'
        ),
        bullets=[
            "errors='coerce' =\n壞資料變 NaT，不是炸",
            "dropna 只針對關鍵欄位 —\n不要 dropna 全表",
            "型別校正最後做，\n比 read_csv 的 dtype= 好改",
            "reset_index 是禮貌 —\n別把 index 留成洞",
        ],
        label_dark=True,
    )
    add_source(s, "pandas §to_datetime §errors='coerce' · Linus fail-soft 原則")
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── S15 · CODE · merge + KPIs ─────────
    s = _blank(prs)
    add_title(s, "Step 3：merge + KPIs（S3 能力回報）")
    draw_code_panel(
        s, x=T.MARGIN_X, y=Inches(1.3),
        w=T.SLIDE_W - 2 * T.MARGIN_X, h=Inches(5.0),
        label="把 S3 的 groupby × merge × pivot 在這裡一次用上",
        code=(
            '# 3-1 合併商品維度（拿到 category）\n'
            'orders = load_and_clean_orders("orders_raw.csv")\n'
            'products = pd.read_csv("products.csv")\n'
            'full = orders.merge(products, on="product_id", how="left")\n'
            '\n'
            '# 3-2 月營收（畫 line）\n'
            'monthly_df = (full\n'
            '    .assign(month=full.order_date.dt.to_period("M").astype(str))\n'
            '    .groupby("month", as_index=False)\n'
            '    .agg(revenue=("revenue", "sum"),\n'
            '         qty=("qty", "sum")))\n'
            '\n'
            '# 3-3 品類營收（畫 bar / pie）\n'
            'category_df = (full\n'
            '    .groupby("category", as_index=False)\n'
            '    .agg(revenue=("revenue", "sum"),\n'
            '         sales=("qty", "sum")))'
        ),
        bullets=[
            "merge how='left' —\n保留所有訂單",
            "to_period('M') 把日期壓成月",
            "agg 用 named tuple —\n欄名自己就是文件",
            "兩個小 df 分別餵四張圖，\n資料與視覺解耦",
        ],
        label_dark=True,
    )
    add_source(s, "S3 groupby+merge 回顧 · pandas §Period §named aggregation")
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── S16 · CODE · 4 張子圖組裝 ─────────
    s = _blank(prs)
    add_title(s, "Step 4：make_subplots 組裝 2x2 儀表板")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.3), h=Inches(5.0),
        slot_name="Plotly 2x2 dashboard 截圖",
        description=(
            "真實執行後截圖：\n"
            "左上月營收折線、右上品類長條、\n"
            "左下散點、右下 donut pie\n"
            "顯示 hover tooltip 示意"
        ),
        placeholder_id="S6_S16_subplot_screenshot",
        registry=image_registry,
        size_hint="1400×900 px",
    )
    draw_code_panel(
        s, x=Inches(6.2), y=Inches(1.3),
        w=Inches(6.5), h=Inches(5.0),
        label="把三個 df 灌進四格 — 儀表板活了",
        code=(
            'fig = make_subplots(\n'
            '    rows=2, cols=2,\n'
            '    subplot_titles=["月營收","品類銷售",\n'
            '                    "量-收關聯","品類占比"],\n'
            '    specs=[[{"type":"xy"},    {"type":"xy"}],\n'
            '           [{"type":"xy"},    {"type":"domain"}]],\n'
            ')\n'
            'fig.add_trace(go.Scatter(x=monthly_df.month,\n'
            '                         y=monthly_df.revenue,\n'
            '                         mode="lines+markers"), 1, 1)\n'
            'fig.add_trace(go.Bar(x=category_df.category,\n'
            '                     y=category_df.sales),     1, 2)\n'
            'fig.add_trace(go.Scatter(x=full.qty, y=full.revenue,\n'
            '                         mode="markers"),      2, 1)\n'
            'fig.add_trace(go.Pie(labels=category_df.category,\n'
            '                     values=category_df.revenue,\n'
            '                     hole=.45),                 2, 2)'
        ),
        bullets=[
            "右下是 pie → domain 格",
            "其他三格 xy 就可以",
            "row/col 每筆都要寫",
            "scatter mode='markers' 才是散點",
        ],
        label_dark=True,
    )
    add_source(s, "Plotly §make_subplots 實戰 · Capstone Step4")
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── S17 · CODE · write_html ─────────
    s = _blank(prs)
    add_title(s, "Step 5：write_html — 一行交付，可寄出的產品")
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.3),
        w=Inches(5.3), h=Inches(5.0),
        slot_name="dashboard.html 寄送示意",
        description=(
            "左側 email 附件 dashboard.html，\n"
            "右側瀏覽器打開的互動畫面，\n"
            "中間箭頭『一鍵寄出、點開即用』"
        ),
        placeholder_id="S6_S17_html_delivery",
        registry=image_registry,
        size_hint="1600×900 px",
    )
    draw_code_panel(
        s, x=Inches(6.2), y=Inches(1.3),
        w=Inches(6.5), h=Inches(5.0),
        label="路徑用絕對路徑、plotly.js 走 CDN — 給你最小交付檔",
        code=(
            'from pathlib import Path\n'
            '\n'
            'fig.update_layout(\n'
            '    height=800,\n'
            '    title="Orders Dashboard (2025-Q1)",\n'
            '    showlegend=False,\n'
            ')\n'
            '\n'
            'out_path = Path("dashboard.html").absolute()\n'
            'fig.write_html(\n'
            '    str(out_path),\n'
            '    include_plotlyjs="cdn",   # 檔案 3MB → 50KB\n'
            '    full_html=True,           # 獨立檔，可單獨開\n'
            ')\n'
            'print(f"已輸出：{out_path}")'
        ),
        bullets=[
            "absolute() → 瀏覽器\n不會找不到檔",
            "include_plotlyjs='cdn'\n檔案縮到 50KB",
            "full_html=True 可獨立開，\n也可用 False 嵌進網頁",
            "最後 print 路徑 —\n對學員最友善",
        ],
        label_dark=True,
    )
    add_source(s, "Plotly §write_html · HTML delivery best practices")
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───────── S18 · PITFALL · 路徑 + hover 爆版 ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="業餘寫法（開不了 / 讀不懂）",
        right_title="資深紀律（能寄、能讀）",
        left_items=[
            "write_html('dashboard.html')（相對路徑）",
            "→ 瀏覽器點 file:// 常常 404",
            "hover_data=list(df.columns)（十欄全塞）",
            "→ tooltip 變資訊牆、沒人看",
        ],
        right_items=[
            "write_html(Path(out).absolute())",
            "→ 或輸出到固定 reports/ 目錄",
            "hover_data=['order_id','customer','profit']",
            "→ 只留 2~3 欄最有商業意義的",
        ],
        title="踩雷三：write_html 路徑 + hover_data 爆版",
        summary="便利性的設計，才是專業 — 為收件者思考，不是為自己。",
    )
    add_source(s, "DA 交付紀律 · Plotly hover_data 設計原則")
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ───────── S19 · TABLE · Grading Rubric ─────────
    s = _blank(prs)
    draw_editorial_table(
        s,
        header=["面向", "權重", "達標關鍵", "拿 A 的條件"],
        rows=[
            ["技術完成度", "40%",
             "dashboard.html 能開啟、4 張子圖都有資料",
             "流程無 warning、輸出 <100KB"],
            ["洞察品質", "30%",
             "至少 1 個具體數字結論",
             "指出商業可行動項，非空話"],
            ["視覺呈現", "20%",
             "subplot_titles / hover 有商業語意",
             "色彩、排版讓非技術主管秒懂"],
            ["口語表達", "10%",
             "30 秒內講清「發現 + 行動」",
             "有數字、有建議、有下一步"],
        ],
        title="評分 Rubric：技術 40 / 洞察 30 / 視覺 20 / 口語 10",
        col_widths=[1.1, 0.6, 2.2, 2.1],
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "學員最常拿 B 是因為沒有「具體數字結論」 — 把 30% 拿滿，A 就穩了。",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID, align=PP_ALIGN.CENTER,
    )
    add_source(s, "Capstone Rubric · 業界 DA 交付評估整理")
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───────── S20 · PRACTICE-PROMPT · 30 秒三問 ─────────
    s = _blank(prs)
    add_title(s, "30 秒成果發表三問：發現、下一張圖、應用場景")
    draw_matrix(s, rows=1, cols=3, top=1.4, bottom=1.8, cells=[
        {"text": "問題 1 · 發現",
         "sub": "這份 dashboard 最關鍵的發現是什麼？\n\n"
                "規格：\n1 個具體數字 + 1 句商業解讀\n\n"
                "範例：\n「3 月營收 -18%，由 A 品類獨自拖累」",
         "highlight": True},
        {"text": "問題 2 · 下一張圖",
         "sub": "再給你 1 小時，你會新增哪張圖？\n\n"
                "規格：\n指明圖型 + 要回答什麼問題\n\n"
                "範例：\n「客戶分群 heatmap，\n看高價值客戶的流失」",
         "highlight": False},
        {"text": "問題 3 · 應用場景",
         "sub": "這個流程在你未來工作會用在哪？\n\n"
                "規格：\n工作場景 + 每週頻率\n\n"
                "範例：\n「每週一 standup 前 30 分\n自動跑 KPI 報表」",
         "highlight": False},
    ])
    draw_inverted_thesis_box(
        s,
        "30 秒內把三題答完 — 這套模板會跟你走一輩子（未來 standup / review / 1-on-1 都用）。",
        y=6.1, width=11.0,
    )
    add_source(s, "Capstone 成果發表流程 · §4 teacher notes")
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ───────── S21 · VS · 迷思 vs 紀律 ─────────
    s = _blank(prs)
    draw_vs_two_col(
        s,
        left_title="業餘迷思（別這樣）",
        right_title="資深紀律（這樣才專業）",
        left_items=[
            "「跑得起來就算交了」",
            "「圖越多越厲害」",
            "「hover 越詳細越貼心」",
            "「notebook 直接丟給主管」",
        ],
        right_items=[
            "能寄出去、對方會點開才算交付",
            "4 張能講完故事，勝過 12 張碎片",
            "hover 只留商業關鍵，多一欄就是噪音",
            "交 dashboard.html，不交 .ipynb",
        ],
        title="常見迷思 vs 資深 DA 交付紀律",
        summary="DA 的終端使用者是非技術主管 — 為他們設計，不為自己的炫技。",
    )
    add_source(s, "業界 DA 交付觀察 · Linus 實用主義原則延伸")
    add_footer(s, MODULE_CODE, 21, N_CONTENT)

    # ───────── S22 · PYRAMID · S1-S6 能力盤點 ─────────
    s = _blank(prs)
    draw_matrix(s, rows=2, cols=3, cells=[
        {"text": "1. 向量化思維",
         "sub": "NumPy broadcasting / ufunc\n從 for-loop 到陣列運算\n來源：S1"},
        {"text": "2. 資料入口",
         "sub": "Pandas I/O、缺失值、dtype\n讀髒資料的第一道防線\n來源：S2"},
        {"text": "3. 資料塑形",
         "sub": "groupby / merge / pivot\n把原始表轉成分析表\n來源：S3",
         "highlight": True},
        {"text": "4. 時序與 EDA",
         "sub": "datetime index / resample\n看模式、看趨勢、看異常\n來源：S4",
         "highlight": True},
        {"text": "5. 視覺化語法",
         "sub": "matplotlib / seaborn\n五種必懂圖的靜態交付\n來源：S5",
         "highlight": True},
        {"text": "6. 互動 + 整合",
         "sub": "Plotly + make_subplots\ndashboard.html 交付\n來源：S6（今天）",
         "highlight": True},
    ], title="S1–S6 能力盤點：12 小時六層地基")
    add_textbox(
        s, T.MARGIN_X, Inches(6.3),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "少一層，上面就不穩 — 你今天能寫 dashboard，是因為前五週都穩了。",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID, align=PP_ALIGN.CENTER,
    )
    add_source(s, "DA 課程架構總覽 · S1-S6 能力地圖合成")
    add_footer(s, MODULE_CODE, 22, N_CONTENT)

    # ───────── S23 · TABLE · 下一站五條路線 ─────────
    s = _blank(prs)
    draw_image_placeholder(
        s,
        x=T.MARGIN_X, y=Inches(1.2),
        w=Inches(5.0), h=Inches(5.0),
        slot_name="下一站學習路徑地圖",
        description=(
            "中央『你現在（DA · dashboard.html）』\n"
            "五條放射狀箭頭指向：\n"
            "Streamlit / Dash / SQL / Cloud / ML"
        ),
        placeholder_id="S6_S23_next_steps_map",
        registry=image_registry,
        size_hint="1600×1000 px",
    )
    # Right side table
    tx = Inches(6.0)
    tw = T.SLIDE_W - tx - T.MARGIN_X
    hdr_h = Inches(0.4)
    row_h = Inches(0.78)

    add_title(s, "下一站：五條路線（挑一條走 3 個月）")

    hdr = add_rect(s, tx, Inches(1.2), tw, hdr_h)
    set_solid_fill(hdr, T.PRIMARY)
    set_no_line(hdr)
    add_textbox(
        s, tx + Inches(0.15), Inches(1.2), tw - Inches(0.3), hdr_h,
        "方向　|　起點　|　一句話",
        font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    paths = [
        ("Web App 化", "Streamlit", "把 notebook 變可點擊的網頁"),
        ("儀表板框架", "Dash（Plotly）", "生產級儀表板、有權限控管"),
        ("資料源升級", "SQL + SQLAlchemy", "把 CSV 換成資料庫才是現實"),
        ("雲端部署", "AWS / GCP", "讓管線每天自動跑、連結可分享"),
        ("進階分析", "scikit-learn", "從描述性走向預測性"),
    ]
    cy = Inches(1.2) + hdr_h
    for i, (d, api, one) in enumerate(paths):
        if i % 2 == 1:
            band = add_rect(s, tx, cy, tw, row_h)
            set_solid_fill(band, T.TABLE_ALT)
            set_no_line(band)
        add_textbox(
            s, tx + Inches(0.15), cy + Inches(0.05),
            tw - Inches(0.3), Inches(0.3),
            f"{d}　·　{api}",
            font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
        )
        add_textbox(
            s, tx + Inches(0.15), cy + Inches(0.38),
            tw - Inches(0.3), Inches(0.35),
            one,
            font_size=T.FONT_SMALL, color=T.CHARCOAL,
        )
        cy += row_h
    add_textbox(
        s, T.MARGIN_X, Inches(6.35),
        T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.3),
        "挑一條走完，比五條都碰過更值錢 — Linus 紀律。",
        font_size=T.FONT_SOURCE, color=T.GRAY_MID, align=PP_ALIGN.CENTER,
    )
    add_source(s, "2025 DA/DE 學習路徑地圖 · 業界現實整理")
    add_footer(s, MODULE_CODE, 23, N_CONTENT)

    # ───────── S24 · SILENT 結業 ─────────
    s = _blank(prs)
    draw_silent_page(
        s,
        "把髒 CSV 變成老闆願意轉寄的 dashboard —\n"
        "這，才是 DA 的商業價值。\n"
        "12 小時結束。從今天起，寫給非技術主管看。",
    )
    add_footer(s, MODULE_CODE, 24, N_CONTENT, dark_bg=True)

    # Copyright
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
