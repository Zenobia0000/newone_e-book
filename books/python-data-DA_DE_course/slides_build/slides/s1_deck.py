"""S1 deck — NumPy 與向量化思維
23 content slides + cover + copyright.

Governing thought:
    把 ndarray 當成一塊『有形狀的記憶體』來想，
    for-loop 就會自己退場——
    向量化不是加個 np.，是換一種語言想運算。

Aligned to chapters/S1_NumPy與向量化思維/01_outline.md
  · 5 Learning Objectives × 5 Common Pitfalls
  · Teaching-track: MOTIVATION / CONCEPT / MECHANISM / EXAMPLE / PITFALL / PRACTICE / CHECKPOINT
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .. import theme as T
from ..primitives import (
    add_textbox, add_title, add_source, add_rect,
    set_solid_fill, set_no_fill, set_line, set_no_line,
    draw_silent_page, draw_ask_page, draw_matrix, draw_editorial_table,
    draw_flow_chain, draw_code_panel, draw_vs_two_col, draw_pyramid_stack,
    draw_thesis_hierarchy, draw_three_blocks_flow, draw_image_placeholder,
    draw_delta_badge, draw_emphasis_pill, draw_inverted_thesis_box,
)
from ..branding import add_cover_slide, add_footer, add_copyright_slide


MODULE_CODE = "S1"
MODULE_TITLE = "NumPy 與向量化思維"
MODULE_SUBTITLE = "ndarray × broadcasting × 向量化——從 for-loop 畢業"
TIME_MIN = 120
N_CONTENT = 23

# Pitfall 頁用的警示色（只在 PITFALL 頁限定）
RED_ERROR = RGBColor(0xC6, 0x28, 0x28)
GREEN_OK = RGBColor(0x2E, 0x7D, 0x32)


def _new_prs():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _draw_three_io(slide, title, input_text, process_text, output_text,
                   bottom_note=""):
    """EXAMPLE-I/O 三欄版：Input | Process | Output 搭配箭頭指示。"""
    add_title(slide, title)
    top = Inches(1.5)
    col_h = Inches(4.2)
    total_w = T.SLIDE_W - 2 * T.MARGIN_X
    gap = Inches(0.2)
    col_w = (total_w - gap * 2) / 3
    labels = ["Input", "Process", "Output"]
    bodies = [input_text, process_text, output_text]
    for i, (lbl, body) in enumerate(zip(labels, bodies)):
        x = T.MARGIN_X + i * (col_w + gap)
        # header band
        hdr = add_rect(slide, x, top, col_w, Inches(0.45))
        set_solid_fill(hdr, T.PRIMARY)
        set_no_line(hdr)
        add_textbox(
            slide, x, top, col_w, Inches(0.45),
            lbl,
            font_size=T.FONT_BODY, color=T.WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        # body box
        body_y = top + Inches(0.45)
        body_h = col_h - Inches(0.45)
        box = add_rect(slide, x, body_y, col_w, body_h)
        set_no_fill(box)
        set_line(box, T.PRIMARY, 1.0)
        add_textbox(
            slide, x + Inches(0.2), body_y + Inches(0.15),
            col_w - Inches(0.4), body_h - Inches(0.3),
            body,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            family=T.FONT_MONO, line_spacing=1.35,
        )
    # arrows between columns
    for i in range(2):
        x = T.MARGIN_X + (i + 1) * col_w + i * gap + gap * 0.1
        add_textbox(
            slide, x, top + col_h / 2 - Inches(0.25),
            gap * 0.8, Inches(0.5),
            "→",
            font_size=Pt(22), color=T.PRIMARY, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
    if bottom_note:
        add_textbox(
            slide, T.MARGIN_X, top + col_h + Inches(0.25),
            T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
            bottom_note,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            align=PP_ALIGN.CENTER, italic=True,
        )


def _draw_pitfall(slide, title, wrong_label, wrong_code, right_label,
                  right_code, why):
    """PITFALL 頁：左紅錯 / 右綠對 / 下方一句 why。
    本頁是全 deck 唯一使用紅/綠警示色的頁型。"""
    add_title(slide, title)
    top = Inches(1.4)
    col_h = Inches(4.2)
    col_w = Inches(5.8)
    gap = Inches(0.3)
    total = col_w * 2 + gap
    left_x = (T.SLIDE_W - total) / 2
    right_x = left_x + col_w + gap

    for x, label, code, accent, mark in [
        (left_x, wrong_label, wrong_code, RED_ERROR, "✗"),
        (right_x, right_label, right_code, GREEN_OK, "✓"),
    ]:
        hdr = add_rect(slide, x, top, col_w, Inches(0.5))
        set_solid_fill(hdr, accent)
        set_no_line(hdr)
        add_textbox(
            slide, x + Inches(0.2), top, col_w - Inches(0.4), Inches(0.5),
            f"{mark}  {label}",
            font_size=T.FONT_BODY, color=T.WHITE, bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        body_y = top + Inches(0.5)
        body_h = col_h - Inches(0.5)
        box = add_rect(slide, x, body_y, col_w, body_h)
        set_solid_fill(box, RGBColor(0xF9, 0xF9, 0xF9))
        set_line(box, accent, 1.2)
        add_textbox(
            slide, x + Inches(0.25), body_y + Inches(0.2),
            col_w - Inches(0.5), body_h - Inches(0.4),
            code,
            font_size=T.FONT_CAPTION, color=T.CHARCOAL,
            family=T.FONT_MONO, line_spacing=1.45,
        )

    # why 條
    why_y = top + col_h + Inches(0.2)
    why_box = add_rect(slide, T.MARGIN_X, why_y,
                       T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.55))
    set_solid_fill(why_box, T.PRIMARY)
    set_no_line(why_box)
    add_textbox(
        slide, T.MARGIN_X + Inches(0.3), why_y,
        T.SLIDE_W - 2 * T.MARGIN_X - Inches(0.6), Inches(0.55),
        f"為什麼：{why}",
        font_size=T.FONT_CAPTION, color=T.WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def build_s1(output_path, image_registry=None):
    """Build S1 deck; 23 content slides."""
    prs = _new_prs()

    # Cover
    add_cover_slide(prs, MODULE_CODE, MODULE_TITLE, MODULE_SUBTITLE,
                    TIME_MIN, N_CONTENT)

    # ───────── S1 · MOTIVATION — 別再等 for-loop 跑完 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "你還在等 for-loop 跑完？",
        data_card={
            "label": "1,000 萬筆加總實測",
            "stat": "80×",
            "caption": "Python list 3.20s\nNumPy 0.04s",
        },
    )
    add_source(s, "timeit · 講師 laptop 實測 · i7 class CPU")
    add_footer(s, MODULE_CODE, 1, N_CONTENT)

    # ───────── S2 · ASK — 差在哪一層 ─────────
    s = _blank(prs)
    draw_ask_page(
        s,
        "同一台電腦、同一個 Python——\n為什麼差 80 倍？",
        data_card={
            "label": "線索",
            "stat": "不是 CPU",
            "caption": "是運算\n跑在哪一層",
        },
    )
    add_footer(s, MODULE_CODE, 2, N_CONTENT)

    # ───────── S3 · SILENT — 一句話立論 ─────────
    s = _blank(prs)
    draw_silent_page(
        s,
        "NumPy 不是『加個 np. 就好』。\n是換一種語言想運算。",
    )
    add_footer(s, MODULE_CODE, 3, N_CONTENT, dark_bg=True)

    # ───────── S4 · CONCEPT-CARD — ndarray 的本質 ─────────
    s = _blank(prs)
    add_title(s, "ndarray ＝ 連續記憶體 ＋ 一層 view")
    # 左文字欄
    left_x = T.MARGIN_X
    left_w = Inches(6.0)
    add_textbox(
        s, left_x, Inches(1.4), left_w, Inches(0.5),
        "三張身分證",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, left_x, Inches(1.95), left_w, Inches(3.8),
        "• .shape — 形狀（幾列幾行）\n"
        "  (2, 3) = 2 列 × 3 行\n\n"
        "• .dtype — 資料型別（同質！）\n"
        "  int64 / float64 / bool\n\n"
        "• .ndim — 維度（幾個軸）\n"
        "  1D 向量 / 2D 矩陣 / 3D 影像",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        line_spacing=1.4,
    )
    add_textbox(
        s, left_x, Inches(5.9), left_w, Inches(0.6),
        "記憶點：reshape 不搬資料，只換眼鏡。",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True, italic=True,
    )
    # 右圖佔位
    draw_image_placeholder(
        s, Inches(7.0), Inches(1.3), Inches(5.7), Inches(4.8),
        slot_name="ndarray 記憶體佈局",
        description="連續記憶體格子 + shape/strides 標籤示意",
        url_hint="numpy.org/doc/stable/reference/arrays.ndarray.html",
        size_hint="1400×1120 px",
        placeholder_id="S1_S4_ndarray_memory",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 4, N_CONTENT)

    # ───────── S5 · EXAMPLE-I/O — np.array() 建立 ─────────
    s = _blank(prs)
    _draw_three_io(
        s,
        "範例：從 list 到 ndarray",
        input_text="lst = [\n  [1, 2, 3],\n  [4, 5, 6]\n]\n\n# Python list\n# 異質、分散",
        process_text="import numpy as np\n\na = np.array(lst)\n\n# 一行搞定",
        output_text="a.shape\n→ (2, 3)\n\na.dtype\n→ int64\n\na.ndim\n→ 2",
        bottom_note="三個屬性是你以後每次 debug 的第一張牌。",
    )
    add_footer(s, MODULE_CODE, 5, N_CONTENT)

    # ───────── S6 · CHECKPOINT — shape 快問快答 ─────────
    s = _blank(prs)
    add_title(s, "Check Point · shape 怎麼讀")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(4), Inches(0.4),
        "進度 6 / 23 · 不看講義作答",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.5),
        "Q1  (3,)  是幾維？\n"
        "       它是向量、列向量、還是行向量？\n\n"
        "Q2  (3, 1) 與 (1, 3) 差在哪？\n"
        "       兩個乘起來 shape 變什麼？\n\n"
        "Q3  (2, 3, 4)  有幾個元素？\n"
        "       怎麼 reshape 成 2D 而不掉資料？",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.5,
    )
    add_footer(s, MODULE_CODE, 6, N_CONTENT)

    # ───────── S7 · MECHANISM-FLOW — 索引三件套 ─────────
    s = _blank(prs)
    draw_flow_chain(
        s,
        nodes=[
            {"label": "基本切片", "sub": "a[1:3]", "caption": "view · 零複製"},
            {"label": "布林遮罩", "sub": "a[a>0]", "caption": "篩選 · 像 WHERE", "highlight": True},
            {"label": "fancy indexing", "sub": "a[[0,2,5]]", "caption": "copy · 任意取"},
        ],
        title="三種索引，三種語意，三種速度",
        y=2.4,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.3), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.6),
        "口訣：切片是借鏡、遮罩像 WHERE、fancy 是點名。",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, MODULE_CODE, 7, N_CONTENT)

    # ───────── S8 · PITFALL (P2) — and/or vs &/| ─────────
    s = _blank(prs)
    _draw_pitfall(
        s,
        title="陷阱 · 布林遮罩：and 不能用，要用 &",
        wrong_label="用 and / or",
        wrong_code="a = np.array([1, 5, 8, 12, 20])\n\n"
                   "a[a>0 and a<10]\n\n"
                   "# ValueError:\n"
                   "# The truth value of an\n"
                   "# array ... is ambiguous",
        right_label="用 & | ~ 加括號",
        right_code="a = np.array([1, 5, 8, 12, 20])\n\n"
                   "a[(a>0) & (a<10)]\n\n"
                   "# array([1, 5, 8])\n"
                   "# 每塊括號都不能省！",
        why="and/or 期待單一 bool；遮罩是整個陣列，必須用元素級運算子 & | ~",
    )
    add_footer(s, MODULE_CODE, 8, N_CONTENT)

    # ───────── S9 · EXAMPLE-I/O — 布林遮罩篩分數 ─────────
    s = _blank(prs)
    _draw_three_io(
        s,
        "範例：篩出 60-100 分的成績",
        input_text="scores = np.array([\n"
                   "  45, 72, 88,\n"
                   "  51, 99, 100,\n"
                   "  38\n"
                   "])",
        process_text="mask = (scores >= 60) \\\n"
                     "     & (scores <= 100)\n\n"
                     "passed = scores[mask]\n\n"
                     "# 一行 = 四行 for-if",
        output_text="passed\n"
                    "→ array([\n"
                    "    72, 88,\n"
                    "    99, 100\n"
                    "])",
        bottom_note="for-if 要寫 4 行，布林遮罩一行——但括號要全。",
    )
    add_footer(s, MODULE_CODE, 9, N_CONTENT)

    # ───────── S10 · PITFALL (P3) — 切片是 view ─────────
    s = _blank(prs)
    _draw_pitfall(
        s,
        title="陷阱 · 切片是 view，不是 copy",
        wrong_label="以為是 copy",
        wrong_code="a = np.array([1, 2, 3, 4, 5])\n\n"
                   "b = a[:3]\n"
                   "b[0] = 99\n\n"
                   "print(a)\n"
                   "→ [99, 2, 3, 4, 5]\n"
                   "# a 也被改到！",
        right_label="要獨立副本就 .copy()",
        right_code="a = np.array([1, 2, 3, 4, 5])\n\n"
                   "b = a[:3].copy()\n"
                   "b[0] = 99\n\n"
                   "print(a)\n"
                   "→ [1, 2, 3, 4, 5]\n"
                   "# a 不動",
        why="切片共用記憶體是 NumPy 刻意設計（省 RAM）；跟 Python list 的 [:] 不一樣",
    )
    add_footer(s, MODULE_CODE, 10, N_CONTENT)

    # ───────── S11 · CONCEPT-CARD — 向量化的本質 ─────────
    s = _blank(prs)
    add_title(s, "向量化 ＝ 把運算整塊送進 C 層")
    draw_vs_two_col(
        s,
        left_title="for-loop（Python 層）",
        right_title="向量化（C 層）",
        left_items=[
            "每輪進直譯器",
            "每個值包 PyObject",
            "記憶體分散、cache miss",
            "天花板在 CPython",
        ],
        right_items=[
            "整塊運算一次送進 C",
            "連續記憶體、SIMD 友善",
            "NumPy 呼叫 BLAS/MKL",
            "天花板在硬體",
        ],
        summary="關鍵：向量化 ≠ 加 np.；是讓整塊運算「一次都不要進入 Python 迴圈」。",
        delta="50-80×",
    )
    add_footer(s, MODULE_CODE, 11, N_CONTENT)

    # ───────── S12 · CHART — 1e7 實測對比 ─────────
    s = _blank(prs)
    add_title(s, "1,000 萬筆加總實測 · list vs NumPy")
    draw_image_placeholder(
        s, Inches(1.3), Inches(1.3), Inches(10.7), Inches(4.2),
        slot_name="list vs NumPy bar chart",
        description="1e7 筆加總 bar chart：Python list 3.20s vs NumPy 0.04s",
        url_hint="自行 timeit 實測 screenshot",
        size_hint="2100×840 px",
        placeholder_id="S1_S12_perf_bar",
        registry=image_registry,
    )
    # 80× badge 右上
    draw_delta_badge(
        s, Inches(10.5), Inches(1.5), "80×",
        w=Inches(1.5), h=Inches(0.9), inverted=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(5.7), T.SLIDE_W - 2 * T.MARGIN_X, Inches(0.5),
        "而且 NumPy 那行，只有 12 個字。",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(s, "timeit · 講師 laptop i7 class · 單次測量")
    add_footer(s, MODULE_CODE, 12, N_CONTENT)

    # ───────── S13 · CONCEPT-CARD — Broadcasting 兩條規則 ─────────
    s = _blank(prs)
    add_title(s, "Broadcasting 只有兩條規則")
    add_textbox(
        s, T.MARGIN_X, Inches(1.4), Inches(6.0), Inches(0.5),
        "規則",
        font_size=T.FONT_BODY, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(1.95), Inches(6.0), Inches(4.0),
        "① 從右往左對齊每個維度\n"
        "    (3, 1)\n"
        "       (4,)    ← 從右對齊\n\n"
        "② 該維度相等，或其中一個是 1\n"
        "    → 可擴展；否則 ValueError\n\n"
        "    (3, 1) + (1, 4) → (3, 4) ✓\n"
        "    (3,)   + (4,)   → 報錯 ✗",
        font_size=T.FONT_CAPTION, color=T.CHARCOAL,
        family=T.FONT_MONO, line_spacing=1.45,
    )
    add_textbox(
        s, T.MARGIN_X, Inches(6.0), Inches(6.0), Inches(0.5),
        "debug 第一步：print(a.shape, b.shape)",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True, italic=True,
    )
    draw_image_placeholder(
        s, Inches(7.0), Inches(1.3), Inches(5.7), Inches(4.8),
        slot_name="broadcasting 對齊示意",
        description="(3,1) + (1,4) → (3,4) 網格擴展示意圖",
        url_hint="numpy.org/doc/stable/user/basics.broadcasting.html",
        size_hint="1400×1120 px",
        placeholder_id="S1_S13_broadcasting",
        registry=image_registry,
    )
    add_footer(s, MODULE_CODE, 13, N_CONTENT)

    # ───────── S14 · PITFALL (P4) — (3,) vs (3,1) ─────────
    s = _blank(prs)
    _draw_pitfall(
        s,
        title="陷阱 · shape (3,) 與 (3,1) 看起來像，結果不同",
        wrong_label="以為一樣",
        wrong_code="a = np.array([1, 2, 3])       # (3,)\n"
                   "b = np.array([[1],[2],[3]])   # (3,1)\n\n"
                   "a + b\n"
                   "# 期待：逐元素加\n"
                   "# 實際：broadcasting\n"
                   "#   → (3, 3) 矩陣！\n"
                   "# [[2,3,4],\n"
                   "#  [3,4,5],\n"
                   "#  [4,5,6]]",
        right_label="先確認 shape",
        right_code="print(a.shape, b.shape)\n"
                   "# (3,)  (3, 1)\n\n"
                   "# 想要逐元素 → shape 對齊\n"
                   "a + b.flatten()   # (3,)+(3,)\n"
                   "→ array([2, 4, 6])\n\n"
                   "# 想要矩陣 → 保留原 shape\n"
                   "a + b            # (3,3) broadcasting",
        why="(3,) 是一維、(3,1) 是二維列向量；從右對齊後擴展方向完全不同",
    )
    add_footer(s, MODULE_CODE, 14, N_CONTENT)

    # ───────── S15 · EXAMPLE-I/O — 電商總庫存價值 ─────────
    s = _blank(prs)
    _draw_three_io(
        s,
        "電商 Case 1 · 總庫存價值（SUMPRODUCT）",
        input_text="price = np.array(\n"
                   "  [199, 599,\n"
                   "   1299, 89])\n\n"
                   "stock = np.array(\n"
                   "  [120, 45,\n"
                   "   8, 300])",
        process_text="total = (price *\n"
                     "         stock).sum()\n\n"
                     "# 逐元素乘\n"
                     "# 再一次加總\n"
                     "# broadcasting (4,)+(4,)",
        output_text="total\n→ 77,215\n\n\n"
                    "# Excel SUMPRODUCT\n"
                    "# 一行搞定",
        bottom_note="招式名：逐元素乘 + sum = SUMPRODUCT",
    )
    add_footer(s, MODULE_CODE, 15, N_CONTENT)

    # ───────── S16 · EXAMPLE-I/O — 電商低庫存商品數 ─────────
    s = _blank(prs)
    _draw_three_io(
        s,
        "電商 Case 2 · 低庫存商品數（COUNT-IF）",
        input_text="stock = np.array(\n"
                   "  [120, 45,\n"
                   "   8, 300])\n\n"
                   "# 低庫存定義：\n"
                   "# stock < 10",
        process_text="low = (stock < 10).sum()\n\n"
                     "# True = 1\n"
                     "# False = 0\n"
                     "# 布林加總 = 計數",
        output_text="low\n→ 1\n\n\n"
                    "# 只有 1 個商品\n"
                    "# 需要補貨",
        bottom_note="招式名：布林遮罩 + sum = COUNT-IF",
    )
    add_footer(s, MODULE_CODE, 16, N_CONTENT)

    # ───────── S17 · PITFALL (P1) — for-loop 反模式 ─────────
    s = _blank(prs)
    _draw_pitfall(
        s,
        title="陷阱 · 看到 for + 數字運算，先停手問向量化",
        wrong_label="for-loop 版（慢 50×）",
        wrong_code="data = np.arange(10_000_000)\n\n"
                   "total = 0\n"
                   "for x in data:\n"
                   "    total += x * 2\n\n"
                   "# 3.2 秒\n"
                   "# 每輪都進 Python 直譯器",
        right_label="向量化版（一行、快、可讀）",
        right_code="data = np.arange(10_000_000)\n\n"
                   "total = (data * 2).sum()\n\n\n"
                   "# 0.04 秒\n"
                   "# 整塊運算一次送進 C",
        why="80% 的數字 for-loop 都可以向量化；反射：看到 for + 數字 → 先停手問自己",
    )
    add_footer(s, MODULE_CODE, 17, N_CONTENT)

    # ───────── S18 · CHECKPOINT — 中段驗收 ─────────
    s = _blank(prs)
    add_title(s, "Check Point · 中段驗收（30 秒）")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(6), Inches(0.4),
        "進度 18 / 23 · 卡住的回翻對應章節",
        font_size=T.FONT_SMALL, color=T.GRAY_MID, italic=True,
    )
    add_textbox(
        s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.5),
        "Q1   a[:3]  是 copy 還是 view？        （回翻 S10）\n\n"
        "Q2   (4,) * (4, 1)  結果 shape 是什麼？  （回翻 S13 / S14）\n\n"
        "Q3   data[(data>0) and (data<10)]\n"
        "        會發生什麼？                    （回翻 S8）",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        line_spacing=1.6,
    )
    add_footer(s, MODULE_CODE, 18, N_CONTENT)

    # ───────── S19 · PITFALL (P5) — 浮點 == 比較 ─────────
    s = _blank(prs)
    _draw_pitfall(
        s,
        title="陷阱 · 浮點不能用 == 比較",
        wrong_label="用 ==",
        wrong_code="0.1 + 0.2 == 0.3\n→ False\n\n\n"
                   "# 為什麼？\n"
                   "# 電腦二進位\n"
                   "# 無法精確表示 0.1\n"
                   "# 實際值：\n"
                   "# 0.30000000000000004",
        right_label="用 np.isclose()",
        right_code="np.isclose(0.1 + 0.2, 0.3)\n→ True\n\n\n"
                   "# 允許微小誤差\n"
                   "# 預設 rtol=1e-05\n"
                   "# 做金融/科學運算\n"
                   "# 一定要肌肉記憶",
        why="浮點本質是近似值；任何浮點相等比較都必須用 np.isclose (或 math.isclose)",
    )
    add_footer(s, MODULE_CODE, 19, N_CONTENT)

    # ───────── S20 · PRACTICE-PROMPT — 🟡 核心練習 ─────────
    s = _blank(prs)
    add_title(s, "練習時間 · 3 分鐘 · 各自動手")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(8), Inches(0.4),
        "🟡 核心題 · 難度：容易 · 目標：鞏固 LO5",
        font_size=T.FONT_CAPTION, color=T.PRIMARY, bold=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(11.0), Inches(4.0),
        "情境：電商商品資料\n\n"
        "   price    = np.array([199, 599, 1299, 89])\n"
        "   stock    = np.array([120,  45,    8, 300])\n\n"
        "用一行分別寫出：\n\n"
        "   ① 總庫存價值\n"
        "   ② 低庫存商品數（stock < 10）\n"
        "   ③ 平均售價",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        family=T.FONT_MONO, line_spacing=1.5,
    )
    draw_emphasis_pill(
        s, Inches(4.5), Inches(6.2), Inches(4.3), Inches(0.6),
        "Think · Pair · Share — 3 分鐘後對答案",
        inverted=True,
    )
    add_footer(s, MODULE_CODE, 20, N_CONTENT)

    # ───────── S21 · PRACTICE-PROMPT — 🔴 雙 11 折扣挑戰 ─────────
    s = _blank(prs)
    add_title(s, "挑戰題 · 10 分鐘 · Broadcasting 實戰")
    add_textbox(
        s, T.MARGIN_X, Inches(1.3), Inches(8), Inches(0.4),
        "🔴 挑戰題 · 難度：進階 · 目標：鞏固 LO4",
        font_size=T.FONT_CAPTION, color=RED_ERROR, bold=True,
    )
    add_textbox(
        s, Inches(1.2), Inches(2.0), Inches(11.0), Inches(4.2),
        "情境：雙 11 全品類折扣\n\n"
        "   discounts = np.array([0.9, 0.8, 0.95])      # (3,) 3C/服飾/食品\n"
        "   prices    = np.random.randint(100, 2000, (3, 5))  # (3, 5)\n\n"
        "任務：一行算出折後價矩陣（shape 應為 (3, 5)）\n\n"
        "提示：\n"
        "   · discounts 要 reshape 成哪個 shape 才能跟 prices 對齊？\n"
        "   · 從右對齊：(3, 5) vs (?, ?) → (3, 5) ✓",
        font_size=T.FONT_BODY, color=T.CHARCOAL,
        family=T.FONT_MONO, line_spacing=1.45,
    )
    draw_emphasis_pill(
        s, Inches(4.0), Inches(6.25), Inches(5.3), Inches(0.55),
        "卡住可以跟鄰座討論 · 10 分鐘後公布答案",
        inverted=False,
    )
    add_footer(s, MODULE_CODE, 21, N_CONTENT)

    # ───────── S22 · PYRAMID — 三層 takeaway ─────────
    s = _blank(prs)
    draw_pyramid_stack(
        s,
        title="收束：NumPy 的三層地基",
        layers=[
            {"name": "思維", "caption": "向量化：整塊運算一次離開 Python"},
            {"name": "語言", "caption": "broadcasting：從右對齊 + 維度為 1 擴展"},
            {"name": "地基", "caption": "ndarray：連續記憶體 + shape/dtype/ndim"},
            {"name": "肌肉", "caption": "5 個一行解 + 5 個必踩地雷"},
        ],
        thesis="會了這三層，Pandas 只是加了欄名的 ndarray。",
    )
    add_footer(s, MODULE_CODE, 22, N_CONTENT)

    # ───────── S23 · SILENT — 銜接下一節 ─────────
    s = _blank(prs)
    draw_silent_page(
        s,
        "學會 NumPy，\nPandas 只是加了欄名的 ndarray。",
    )
    add_footer(s, MODULE_CODE, 23, N_CONTENT, dark_bg=True)

    # Copyright
    add_copyright_slide(prs, MODULE_CODE)

    prs.save(str(output_path))
    return output_path
